from pathlib import Path
from zipfile import BadZipFile, ZIP_DEFLATED, ZipFile
from collections import defaultdict
from io import BytesIO
import calendar
import csv
from datetime import date, datetime, time, timedelta
import html
import json
import os
import secrets
import sqlite3
import string
import threading
import uuid
import re
from urllib.parse import urlencode

from flask import Flask, redirect, request, send_file, session, jsonify
from werkzeug.security import check_password_hash, generate_password_hash

import xlrd
from openpyxl import Workbook, load_workbook
from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

try:
    import psycopg
    from psycopg.rows import dict_row
except ImportError:
    psycopg = None
    dict_row = None

try:
    import anthropic as anthropic_sdk
except ImportError:
    anthropic_sdk = None

try:
    import pandas as pd
except ImportError:
    pd = None


H_ALIGN = {1: "left", 2: "center", 3: "right", 4: "fill", 5: "justify", 6: "centerContinuous", 7: "distributed"}
V_ALIGN = {0: "top", 1: "center", 2: "bottom", 3: "justify", 4: "distributed"}
NO_BORDER = Border(
    left=Side(border_style=None),
    right=Side(border_style=None),
    top=Side(border_style=None),
    bottom=Side(border_style=None),
)
ALLOWED_EXTENSIONS = {"xls", "xlsx", "csv"}
IMAGE_EXTENSIONS = {"png", "jpg", "jpeg", "gif", "bmp", "webp"}
DOC_EXTENSIONS = {"doc", "docx", "pdf", "txt", "rtf", "dat", "xml", "json", "ppt", "pptx"}
CHAT_ALLOWED_EXTENSIONS = ALLOWED_EXTENSIONS | IMAGE_EXTENSIONS | DOC_EXTENSIONS
MAX_UPLOAD_SIZE = 10 * 1024 * 1024
DATABASE_URL = os.environ.get("DATABASE_URL")


def clean(val):
    if isinstance(val, str):
        return val.replace("*", "").replace("?", "").strip()
    return val


def get_extension(filename):
    if not filename or "." not in filename:
        return ""
    return filename.rsplit(".", 1)[-1].lower()


def detect_excel_signature(file_storage):
    pos = file_storage.stream.tell()
    header = file_storage.stream.read(8)
    file_storage.stream.seek(pos)
    if header.startswith(b"PK\x03\x04"):
        return "xlsx"
    if header.startswith(b"\xD0\xCF\x11\xE0\xA1\xB1\x1A\xE1"):
        return "xls"
    return ""


def detect_excel_file_signature(path_text):
    try:
        with open(path_text, "rb") as handle:
            header = handle.read(8)
    except OSError:
        return ""
    if header.startswith(b"PK\x03\x04"):
        return "xlsx"
    if header.startswith(b"\xD0\xCF\x11\xE0\xA1\xB1\x1A\xE1"):
        return "xls"
    return ""


def validate_upload(file_storage):
    if not file_storage or file_storage.filename == "":
        return "missing", None

    ext = get_extension(file_storage.filename)
    if ext not in ALLOWED_EXTENSIONS:
        return "unsupported", None

    file_storage.stream.seek(0, os.SEEK_END)
    size = file_storage.stream.tell()
    file_storage.stream.seek(0)
    if size <= 0:
        return "empty", None
    if size > MAX_UPLOAD_SIZE:
        return "too_large", None

    if ext == "csv":
        return None, ext

    detected = detect_excel_signature(file_storage)
    if detected != ext:
        return "invalid_excel", None

    return None, ext


def validate_upload_for_chat(file_storage):
    """Validate file upload for AI chat — accepts all common business file formats."""
    if not file_storage or file_storage.filename == "":
        return "missing", None

    ext = get_extension(file_storage.filename)
    if ext not in CHAT_ALLOWED_EXTENSIONS:
        return "unsupported", None

    file_storage.stream.seek(0, os.SEEK_END)
    size = file_storage.stream.tell()
    file_storage.stream.seek(0)
    if size <= 0:
        return "empty", None
    if size > MAX_UPLOAD_SIZE:
        return "too_large", None

    # Images, docs, text: skip excel signature check
    if ext in IMAGE_EXTENSIONS or ext in DOC_EXTENSIONS or ext == "csv":
        return None, ext

    # Excel files: verify signature
    detected = detect_excel_signature(file_storage)
    if detected != ext:
        return "invalid_excel", None

    return None, ext


def analyze_image_for_chat(file_path, filename):
    """Analyze an image file: return basic metadata string."""
    try:
        size_bytes = os.path.getsize(file_path)
        size_str = f"{size_bytes / 1024:.1f} KB" if size_bytes < 1024 * 1024 else f"{size_bytes / (1024 * 1024):.1f} MB"
        # Try to get dimensions without PIL (read basic headers)
        width, height = None, None
        try:
            with open(file_path, "rb") as f:
                header = f.read(32)
                # PNG
                if header[:8] == b'\x89PNG\r\n\x1a\n':
                    width = int.from_bytes(header[16:20], 'big')
                    height = int.from_bytes(header[20:24], 'big')
                # JPEG
                elif header[:2] == b'\xff\xd8':
                    f.seek(0)
                    data = f.read()
                    i = 2
                    while i < len(data) - 9:
                        if data[i] == 0xFF:
                            marker = data[i + 1]
                            if marker in (0xC0, 0xC1, 0xC2):
                                height = int.from_bytes(data[i + 5:i + 7], 'big')
                                width = int.from_bytes(data[i + 7:i + 9], 'big')
                                break
                            elif marker == 0xD9:
                                break
                            else:
                                seg_len = int.from_bytes(data[i + 2:i + 4], 'big')
                                i += 2 + seg_len
                        else:
                            i += 1
                # GIF
                elif header[:6] in (b'GIF87a', b'GIF89a'):
                    width = int.from_bytes(header[6:8], 'little')
                    height = int.from_bytes(header[8:10], 'little')
        except Exception:
            pass

        dims_str = f"{width}x{height}" if width and height else "unknown dimensions"
        return f"Image file: {filename}\nSize: {size_str}\nDimensions: {dims_str}\nThis is a screenshot/image. The user likely wants you to see the structure or layout shown in the image."
    except Exception as exc:
        return f"Image file: {filename} (could not analyze: {exc})"


def analyze_document_for_chat(file_path, filename, extension):
    """Analyze document files (txt, dat, csv-like, pdf, docx, etc.) for AI chat context."""
    try:
        size_bytes = os.path.getsize(file_path)
        size_str = f"{size_bytes / 1024:.1f} KB" if size_bytes < 1024 * 1024 else f"{size_bytes / (1024 * 1024):.1f} MB"

        # Plain text / DAT / CSV-like files
        if extension in ("txt", "dat", "rtf"):
            lines = []
            for enc in ("utf-8", "utf-8-sig", "cp1255", "latin-1"):
                try:
                    with open(file_path, "r", encoding=enc) as f:
                        lines = [f.readline() for _ in range(20)]
                    lines = [l.rstrip("\n\r") for l in lines if l.strip()]
                    break
                except (UnicodeDecodeError, UnicodeError):
                    continue
            if not lines:
                return f"Document: {filename} ({size_str}). Could not read text content."
            # Detect if it's tabular (tab/comma/pipe separated)
            sample = "\n".join(lines[:10])
            sep = None
            if "\t" in lines[0]:
                sep = "tab"
            elif "|" in lines[0]:
                sep = "pipe (|)"
            elif lines[0].count(",") > 2:
                sep = "comma"
            summary = f"Text file: {filename} ({size_str})\n"
            summary += f"Lines in preview: {len(lines)}\n"
            if sep:
                summary += f"Appears to be tabular data separated by {sep}\n"
                summary += f"First line (likely headers): {lines[0]}\n"
            summary += f"Preview (first {min(len(lines), 10)} lines):\n"
            for line in lines[:10]:
                summary += f"  {line}\n"
            return summary

        # XML files
        if extension == "xml":
            try:
                with open(file_path, "r", encoding="utf-8") as f:
                    content = f.read(4000)
            except UnicodeDecodeError:
                with open(file_path, "r", encoding="latin-1") as f:
                    content = f.read(4000)
            summary = f"XML file: {filename} ({size_str})\n"
            summary += f"Preview (first 4000 chars):\n{content}\n"
            return summary

        # JSON files
        if extension == "json":
            try:
                with open(file_path, "r", encoding="utf-8") as f:
                    data = json.load(f)
                if isinstance(data, list) and len(data) > 0:
                    summary = f"JSON file: {filename} ({size_str})\n"
                    summary += f"Array with {len(data)} items\n"
                    if isinstance(data[0], dict):
                        keys = list(data[0].keys())
                        summary += f"Item keys: {', '.join(str(k) for k in keys)}\n"
                        summary += f"Sample item: {json.dumps(data[0], ensure_ascii=False)[:500]}\n"
                    return summary
                elif isinstance(data, dict):
                    summary = f"JSON file: {filename} ({size_str})\n"
                    summary += f"Object with keys: {', '.join(str(k) for k in list(data.keys())[:20])}\n"
                    return summary
            except Exception:
                pass
            return f"JSON file: {filename} ({size_str}). Could not parse structure."

        # PDF files
        if extension == "pdf":
            # Read first bytes to confirm it's a PDF
            try:
                with open(file_path, "rb") as f:
                    header = f.read(5)
                is_pdf = header == b"%PDF-"
            except Exception:
                is_pdf = False
            summary = f"PDF file: {filename} ({size_str})\n"
            if is_pdf:
                summary += "This is a valid PDF document. The user may want you to understand its structure or content.\n"
                summary += "Note: PDF text extraction is limited. Ask the user to describe the content or copy-paste relevant text."
            else:
                summary += "File header does not match PDF format."
            return summary

        # Word documents (.docx)
        if extension == "docx":
            # docx is a ZIP — try to extract text from word/document.xml
            try:
                from zipfile import ZipFile
                with ZipFile(file_path) as zf:
                    if "word/document.xml" in zf.namelist():
                        import re as _re
                        xml_content = zf.read("word/document.xml").decode("utf-8", errors="replace")
                        # Strip XML tags to get text
                        text = _re.sub(r"<[^>]+>", " ", xml_content)
                        text = " ".join(text.split())[:3000]
                        summary = f"Word document: {filename} ({size_str})\n"
                        summary += f"Text preview (first 3000 chars):\n{text}\n"
                        return summary
            except Exception:
                pass
            return f"Word document: {filename} ({size_str}). Could not extract text."

        # Old Word (.doc)
        if extension == "doc":
            return f"Word document (legacy .doc): {filename} ({size_str}). Cannot read .doc content directly. Ask the user to save as .docx or copy-paste the relevant text."

        # PowerPoint (.pptx)
        if extension == "pptx":
            try:
                from zipfile import ZipFile
                import re as _re
                with ZipFile(file_path) as zf:
                    texts = []
                    for name in sorted(zf.namelist()):
                        if name.startswith("ppt/slides/slide") and name.endswith(".xml"):
                            xml_content = zf.read(name).decode("utf-8", errors="replace")
                            slide_text = _re.sub(r"<[^>]+>", " ", xml_content)
                            slide_text = " ".join(slide_text.split())
                            if slide_text.strip():
                                texts.append(slide_text[:500])
                    summary = f"PowerPoint: {filename} ({size_str}), {len(texts)} slides\n"
                    for i, t in enumerate(texts[:5]):
                        summary += f"  Slide {i+1}: {t[:200]}\n"
                    return summary
            except Exception:
                pass
            return f"PowerPoint: {filename} ({size_str}). Could not extract slide text."

        # Old PowerPoint (.ppt)
        if extension == "ppt":
            return f"PowerPoint (legacy .ppt): {filename} ({size_str}). Cannot read .ppt content directly. Ask the user to save as .pptx or describe the content."

        return f"Document: {filename} ({size_str}, type: {extension})"
    except Exception as exc:
        return f"Document: {filename} (could not analyze: {exc})"


def idx_to_hex(cmap, idx):
    if idx in (0, 64, 65, 32767, None):
        return None
    rgb = cmap.get(idx)
    if rgb and None not in rgb:
        return "{:02X}{:02X}{:02X}".format(*rgb)
    return None


def widen_known_columns(ws):
    for col, extra in (("G", 4.5), ("S", 3.0), ("AU", 3.0)):
        current = ws.column_dimensions[col].width or 8.43
        ws.column_dimensions[col].width = current + extra


def process_legacy_xls(input_path, output_path):
    try:
        wb_in = xlrd.open_workbook(input_path, formatting_info=True)
        cmap = wb_in.colour_map
        preserve_formatting = True
    except (NotImplementedError, TypeError):
        wb_in = xlrd.open_workbook(input_path)
        cmap = {}
        preserve_formatting = False
    wb_out = Workbook()
    wb_out.remove(wb_out.active)

    for s_idx in range(wb_in.nsheets):
        ws_in = wb_in.sheet_by_index(s_idx)
        ws_out = wb_out.create_sheet(title=ws_in.name)
        ws_out.sheet_view.rightToLeft = True
        ws_out.sheet_view.showGridLines = False

        for c in range(ws_in.ncols):
            ltr = get_column_letter(c + 1)
            ci = ws_in.colinfo_map.get(c)
            width = (ci.width / 256.0) if (ci and ci.width) else 1.0
            ws_out.column_dimensions[ltr].width = max(width, 0.5)

        for r in range(ws_in.nrows):
            ri = ws_in.rowinfo_map.get(r)
            orig_pts = (ri.height / 20.0) if (ri and ri.height) else 12.75
            has_content = any(v != "" for v in ws_in.row_values(r))
            if has_content:
                new_h = orig_pts if orig_pts >= 25 else (18 if orig_pts >= 15 else 15)
            else:
                new_h = 2 if orig_pts <= 8 else (orig_pts if orig_pts <= 50 else 12)
            ws_out.row_dimensions[r + 1].height = new_h

        for r in range(ws_in.nrows):
            for c in range(ws_in.ncols):
                val = clean(ws_in.cell_value(r, c))
                cell = ws_out.cell(row=r + 1, column=c + 1, value=val)
                if preserve_formatting:
                    xf = wb_in.xf_list[ws_in.cell_xf_index(r, c)]
                    fi = wb_in.font_list[xf.font_index]
                    fc = idx_to_hex(cmap, fi.colour_index)
                    cell.font = Font(
                        bold=bool(fi.bold),
                        italic=bool(fi.italic),
                        size=fi.height / 20,
                        name=fi.name or "Arial",
                        color=fc or "000000",
                    )
                    bg = idx_to_hex(cmap, xf.background.pattern_colour_index)
                    if bg and xf.background.fill_pattern != 0:
                        cell.fill = PatternFill(fill_type="solid", fgColor=bg)
                    cell.alignment = Alignment(
                        horizontal=H_ALIGN.get(xf.alignment.hor_align, "general"),
                        vertical=V_ALIGN.get(xf.alignment.vert_align, "bottom"),
                        wrapText=bool(xf.alignment.text_wrapped),
                        shrinkToFit=bool(xf.alignment.shrink_to_fit),
                        readingOrder=2,
                    )
                else:
                    cell.alignment = Alignment(horizontal="right", vertical="center", readingOrder=2)
                cell.border = NO_BORDER

        if preserve_formatting:
            for r1, r2, c1, c2 in ws_in.merged_cells:
                if r2 > r1 or c2 > c1:
                    try:
                        ws_out.merge_cells(start_row=r1 + 1, start_column=c1 + 1, end_row=r2, end_column=c2)
                    except Exception:
                        pass

    for ws in wb_out.worksheets:
        widen_known_columns(ws)

    wb_out.save(output_path)


def process_xlsx(input_path, output_path):
    wb = load_workbook(input_path)
    for ws in wb.worksheets:
        ws.sheet_view.rightToLeft = True
        ws.sheet_view.showGridLines = False
        for row in ws.iter_rows():
            for cell in row:
                if isinstance(cell.value, str):
                    cell.value = clean(cell.value)
        widen_known_columns(ws)
    wb.save(output_path)


def process_spreadsheet(input_path, output_path, extension):
    actual_extension = detect_excel_file_signature(input_path) or extension
    if actual_extension == "xls":
        process_legacy_xls(input_path, output_path)
        return
    if actual_extension == "xlsx":
        process_xlsx(input_path, output_path)
        return
    raise ValueError("Unsupported file type")


def run_attendance_cleanup(input_path, output_path, extension, options=None):
    process_spreadsheet(input_path, output_path, extension)


PAYABLE_HOUR_LABELS = {"רגילות", "׳¨׳’׳™׳׳•׳×", "100%", "125%", "150%", "175%", "200%"}
REGULAR_PAYABLE_HOUR_LABELS = {"רגילות", "׳¨׳’׳™׳׳•׳×"}

FLAMINGO_PAYABLE_KEYWORDS = ("שעותלתשלום", "שעותמשולמות", "רגילות", "נוכחות")
FLAMINGO_RATE_KEYWORDS = ("תעריף", "שעה", "הערות", "rate")


def parse_numeric_rate(value):
    if value in ("", None):
        return None
    if isinstance(value, (int, float)):
        return float(value)
    text = str(value).strip().replace(",", ".")
    if not text:
        return None
    return float(text)


def parse_hours_value(value):
    if value in ("", None):
        return None
    if isinstance(value, (int, float)):
        numeric = float(value)
        if 0 <= numeric <= 1:
            return numeric * 24
        return numeric
    text = str(value).strip()
    if not text:
        return None
    if ":" in text:
        hours_text, minutes_text = text.split(":", 1)
        return int(hours_text) + (int(minutes_text) / 60.0)
    return float(text.replace(",", "."))


def format_hours(hours_value):
    if hours_value in (None, ""):
        return ""
    total_minutes = int(round(float(hours_value) * 60))
    hours, minutes = divmod(total_minutes, 60)
    return f"{hours:02d}:{minutes:02d}"


def parse_hours_or_zero(value):
    parsed = parse_hours_value(value)
    return 0.0 if parsed is None else parsed


def try_parse_hours_value(value):
    try:
        return parse_hours_value(value)
    except (ValueError, TypeError):
        return None


def parse_float_or_none(value):
    text = str(value).strip()
    if not text:
        return None
    return float(text.replace(",", "."))


def parse_int_or_none(value):
    parsed = parse_float_or_none(value)
    if parsed is None:
        return None
    return int(parsed)


def detect_matan_missing_header_row(sheet):
    best_row = 4 if sheet.nrows > 4 else 0
    best_score = -1
    target_tokens = {
        normalize_token("שם עובד"),
        normalize_token("מספר עובד"),
        normalize_token("חוסר"),
        normalize_token("ש.תקן"),
        normalize_token("ש.נוכחות"),
    }
    for row_index in range(min(sheet.nrows, 12)):
        score = 0
        for col_index in range(sheet.ncols):
            token = normalize_token(sheet.cell_value(row_index, col_index))
            if token in target_tokens:
                score += 1
        if score > best_score:
            best_score = score
            best_row = row_index
    return best_row


def normalize_token(text):
    value = str(text or "").strip().lower()
    return re.sub(r"[\s_\-\"'`]+", "", value)


def safe_sheet_title(title, fallback):
    cleaned = "".join(ch for ch in str(title) if ch not in '[]:*?/\\')[:31].strip()
    return cleaned or fallback


def get_sheet_cell(sheet, row_index, col_index, default=""):
    if row_index >= sheet.nrows or col_index >= sheet.ncols:
        return default
    value = sheet.cell_value(row_index, col_index)
    if isinstance(value, str):
        return value.strip()
    return value


def find_row_label_value(sheet, row_index, label):
    if row_index >= sheet.nrows:
        return ""
    values = [sheet.cell_value(row_index, c) for c in range(sheet.ncols)]
    for idx, value in enumerate(values):
        if str(value).strip() == label:
            for next_idx in range(idx + 1, len(values)):
                candidate = values[next_idx]
                if candidate not in ("", None):
                    return candidate
            return ""
    return ""


def find_row_label_value_with_offsets(sheet, row_index, label, offsets):
    if row_index >= sheet.nrows:
        return ""
    values = [sheet.cell_value(row_index, c) for c in range(sheet.ncols)]
    for idx, value in enumerate(values):
        if str(value).strip() == label:
            for offset in offsets:
                candidate_index = idx + offset
                if 0 <= candidate_index < len(values):
                    candidate = values[candidate_index]
                    if candidate not in ("", None):
                        return candidate
            return ""
    return ""


def extract_payable_hours(summary_sheet):
    totals = {}
    for row_index in range(summary_sheet.nrows):
        label = str(get_sheet_cell(summary_sheet, row_index, 9, "")).strip()
        if label in PAYABLE_HOUR_LABELS:
            totals[label] = parse_hours_value(get_sheet_cell(summary_sheet, row_index, 13, ""))
    regular_hours = [totals[label] for label in REGULAR_PAYABLE_HOUR_LABELS if totals.get(label) is not None]
    if regular_hours:
        return regular_hours[0], totals
    available = [value for value in totals.values() if value is not None]
    if not available:
        return None, totals
    return sum(available), totals


def extract_flamingo_worker_pair(detail_sheet, summary_sheet, workbook_kind, mapping, manual_hourly_rate_text=""):
    worker_name = stringify_excel_value(extract_flamingo_mapping_value(detail_sheet, summary_sheet, workbook_kind, mapping.get("worker_name_source")))
    worker_name = worker_name or get_flamingo_sheet_name(detail_sheet, workbook_kind)
    department = stringify_excel_value(extract_flamingo_mapping_value(detail_sheet, summary_sheet, workbook_kind, mapping.get("department_source")))
    rate_raw = extract_flamingo_mapping_value(detail_sheet, summary_sheet, workbook_kind, mapping.get("hourly_rate_source"))
    worker_number = stringify_excel_value(extract_flamingo_mapping_value(detail_sheet, summary_sheet, workbook_kind, mapping.get("worker_number_source")))
    id_number = stringify_excel_value(extract_flamingo_mapping_value(detail_sheet, summary_sheet, workbook_kind, mapping.get("id_number_source")))
    start_date = stringify_excel_value(extract_flamingo_mapping_value(detail_sheet, summary_sheet, workbook_kind, mapping.get("start_date_source")))
    attendance_hours = extract_flamingo_mapping_value(detail_sheet, summary_sheet, workbook_kind, mapping.get("attendance_hours_source"))
    standard_hours = extract_flamingo_mapping_value(detail_sheet, summary_sheet, workbook_kind, mapping.get("standard_hours_source"))
    missing_hours = extract_flamingo_mapping_value(detail_sheet, summary_sheet, workbook_kind, mapping.get("missing_hours_source"))
    notes = []
    status = "OK"

    hourly_rate = None
    used_manual_rate = False
    if str(manual_hourly_rate_text or "").strip():
        try:
            hourly_rate = parse_numeric_rate(manual_hourly_rate_text)
            used_manual_rate = True
            rate_raw = manual_hourly_rate_text
        except ValueError:
            hourly_rate = None
            status = "Invalid manual hourly rate"
            notes.append(f"התעריף השעתי הידני אינו תקין: {manual_hourly_rate_text}")
    elif rate_raw not in ("", None):
        try:
            hourly_rate = parse_numeric_rate(rate_raw)
        except ValueError:
            hourly_rate = None
            status = "Invalid hourly rate"
            notes.append(f"ערך התעריף השעתי אינו תקין: {rate_raw}")

    if hourly_rate is None and status == "OK":
        status = "Missing hourly rate"
        notes.append("לא נבחר שדה תעריף שעתי ולא הוזן תעריף ידני.")

    payable_hours_raw = extract_flamingo_mapping_value(detail_sheet, summary_sheet, workbook_kind, mapping.get("payable_hours_source"))
    try:
        payable_hours = parse_hours_value(payable_hours_raw)
    except ValueError:
        payable_hours = None
        if status == "OK":
            status = "Invalid payable hours"
        notes.append(f"ערך שעות התשלום אינו תקין: {payable_hours_raw}")

    payable_breakdown = {}
    summary_name = get_flamingo_sheet_name(summary_sheet, workbook_kind) if summary_sheet is not None else get_flamingo_sheet_name(detail_sheet, workbook_kind)
    if payable_hours is None and status == "OK":
        status = "Missing payable hours"
        notes.append("לא נבחר או לא זוהה שדה שעות לתשלום בפועל.")

    calculated_salary = None
    if status == "OK":
        calculated_salary = round(payable_hours * hourly_rate, 2)

    return {
        "worker_name": worker_name,
        "department": department,
        "worker_number": worker_number,
        "id_number": id_number,
        "start_date": start_date,
        "detail_sheet": get_flamingo_sheet_name(detail_sheet, workbook_kind),
        "summary_sheet": summary_name,
        "hourly_rate": hourly_rate,
        "hourly_rate_raw": rate_raw,
        "payable_hours": payable_hours,
        "payable_breakdown": payable_breakdown,
        "attendance_hours": parse_hours_value(attendance_hours) if attendance_hours not in ("", None) else None,
        "standard_hours": parse_hours_value(standard_hours) if standard_hours not in ("", None) else None,
        "missing_hours": parse_hours_value(missing_hours) if missing_hours not in ("", None) else None,
        "salary": calculated_salary,
        "status": status,
        "used_manual_rate": used_manual_rate,
        "notes": " | ".join(notes),
    }


def translate_flamingo_status(status):
    return {
        "OK": "תקין",
        "Missing hourly rate": "חסר תעריף שעתי",
        "Invalid hourly rate": "תעריף שעתי לא תקין",
        "Invalid manual hourly rate": "תעריף ידני לא תקין",
        "Missing payable hours": "חסרות שעות לתשלום",
        "Invalid payable hours": "שעות לתשלום לא תקינות",
        "Could not match summary sheet": "לא זוהה גיליון סיכום",
    }.get(status, status)


def get_flamingo_attention_action(status):
    if status in {"Missing hourly rate", "Invalid hourly rate", "Invalid manual hourly rate"}:
        return "יש לבחור שדה תעריף שעתי נכון או להזין תעריף ידני תקין."
    if status == "Could not match summary sheet":
        return "יש לבדוק את מבנה הדוח ולוודא שקיים אזור סיכום תקין לעובד."
    return "יש לבדוק את שדה שעות התשלום ולוודא שנבחר השדה הנכון."


def write_flamingo_summary_sheet(ws, worker_rows):
    ws.sheet_view.rightToLeft = True
    ws.sheet_view.showGridLines = False
    ws.title = safe_sheet_title("סיכום שכר", "סיכום שכר")
    ws.freeze_panes = "A12"

    successful_rows = [row for row in worker_rows if row["status"] == "OK"]
    total_workers = len(worker_rows)
    unresolved_workers = len([row for row in worker_rows if row["status"] != "OK"])
    total_hours = sum(row["payable_hours"] or 0 for row in successful_rows)
    total_salary = sum(row["salary"] or 0 for row in successful_rows)
    average_rate = (sum(row["hourly_rate"] or 0 for row in successful_rows) / len(successful_rows)) if successful_rows else 0

    metrics = [
        ("סה\"כ עובדים", total_workers, "DBEAFE"),
        ("חושבו בהצלחה", len(successful_rows), "DCFCE7"),
        ("דורשים טיפול", unresolved_workers, "FEE2E2"),
        ("סה\"כ שעות לתשלום", format_hours(total_hours), "FEF3C7"),
        ("סה\"כ שכר", round(total_salary, 2), "E9D5FF"),
        ("ממוצע תעריף שעתי", round(average_rate, 2), "FCE7F3"),
    ]

    ws.merge_cells("A1:D1")
    ws["A1"] = "דוח סיכום שכר"
    ws["A1"].font = Font(bold=True, size=18, color="0F172A")
    ws["A1"].fill = PatternFill(fill_type="solid", fgColor="BFDBFE")
    ws["A1"].alignment = Alignment(horizontal="center")
    ws["A2"] = "חישוב שכר אוטומטי מתוך דוח מפורט חודשי"
    ws["A2"].font = Font(italic=True, size=11, color="475569")
    ws["A2"].alignment = Alignment(horizontal="center")

    for index, (label, value, fill_color) in enumerate(metrics):
        start_col = 1 + (index % 3) * 4
        row = 4 + (index // 3) * 3
        ws.merge_cells(start_row=row, start_column=start_col, end_row=row, end_column=start_col + 1)
        ws.merge_cells(start_row=row + 1, start_column=start_col, end_row=row + 1, end_column=start_col + 1)
        label_cell = ws.cell(row=row, column=start_col, value=label)
        value_cell = ws.cell(row=row + 1, column=start_col, value=value)
        label_cell.font = Font(bold=True, color="334155")
        value_cell.font = Font(bold=True, size=14, color="0F172A")
        label_cell.fill = PatternFill(fill_type="solid", fgColor=fill_color)
        value_cell.fill = PatternFill(fill_type="solid", fgColor=fill_color)
        label_cell.alignment = Alignment(horizontal="center")
        value_cell.alignment = Alignment(horizontal="center")

    header_row = 11
    headers = [
        "שם עובד",
        "מספר עובד",
        "תעודת זהות",
        "מחלקה",
        "תחילת עבודה",
        "גיליון מפורט",
        "גיליון סיכום",
        "תעריף שעתי",
        "שעות לתשלום",
        "שכר מחושב",
        "סטטוס",
        "הערות",
    ]

    for col_index, header in enumerate(headers, start=1):
        cell = ws.cell(row=header_row, column=col_index, value=header)
        cell.font = Font(bold=True, color="FFFFFF")
        cell.fill = PatternFill(fill_type="solid", fgColor="1E3A8A")
        cell.alignment = Alignment(horizontal="center")
    ws.auto_filter.ref = f"A{header_row}:{get_column_letter(len(headers))}{header_row}"

    for row_index, worker in enumerate(worker_rows, start=header_row + 1):
        values = [
            worker["worker_name"],
            worker["worker_number"],
            worker["id_number"],
            worker["department"],
            worker["start_date"],
            worker["detail_sheet"],
            worker["summary_sheet"],
            worker["hourly_rate"],
            format_hours(worker["payable_hours"]),
            worker["salary"],
            translate_flamingo_status(worker["status"]),
            worker["notes"],
        ]
        for col_index, value in enumerate(values, start=1):
            cell = ws.cell(row=row_index, column=col_index, value=value)
            cell.alignment = Alignment(horizontal="right")
        fill_color = "ECFDF5" if worker["status"] == "OK" else "FEE2E2"
        for col_index in range(1, len(headers) + 1):
            ws.cell(row=row_index, column=col_index).fill = PatternFill(fill_type="solid", fgColor=fill_color)
        ws.cell(row=row_index, column=8).number_format = '0.00'
        ws.cell(row=row_index, column=10).number_format = '#,##0.00'

    widths = [22, 16, 18, 18, 14, 18, 18, 14, 14, 16, 22, 42]
    for col_index, width in enumerate(widths, start=1):
        ws.column_dimensions[get_column_letter(col_index)].width = width


def write_flamingo_attention_sheet(ws, worker_rows):
    ws.sheet_view.rightToLeft = True
    ws.sheet_view.showGridLines = False
    ws.title = safe_sheet_title("דורש טיפול", "דורש טיפול")
    ws.freeze_panes = "A2"

    headers = ["שם עובד", "מספר עובד", "תעודת זהות", "סוג תקלה", "תעריף שעתי", "שעות לתשלום", "פעולה מומלצת"]
    for col_index, header in enumerate(headers, start=1):
        cell = ws.cell(row=1, column=col_index, value=header)
        cell.font = Font(bold=True, color="FFFFFF")
        cell.fill = PatternFill(fill_type="solid", fgColor="B91C1C")
        cell.alignment = Alignment(horizontal="center")
    ws.auto_filter.ref = f"A1:{get_column_letter(len(headers))}1"

    issues = [row for row in worker_rows if row["status"] != "OK"]
    for row_index, worker in enumerate(issues, start=2):
        values = [
            worker["worker_name"],
            worker["worker_number"],
            worker["id_number"],
            translate_flamingo_status(worker["status"]),
            worker["hourly_rate_raw"],
            format_hours(worker["payable_hours"]),
            get_flamingo_attention_action(worker["status"]),
        ]
        for col_index, value in enumerate(values, start=1):
            cell = ws.cell(row=row_index, column=col_index, value=value)
            cell.fill = PatternFill(fill_type="solid", fgColor="FEF2F2")
            cell.alignment = Alignment(horizontal="right")

    widths = [22, 16, 18, 24, 14, 14, 60]
    for col_index, width in enumerate(widths, start=1):
        ws.column_dimensions[get_column_letter(col_index)].width = width


def write_flamingo_department_sheet(ws, worker_rows):
    ws.sheet_view.rightToLeft = True
    ws.sheet_view.showGridLines = False
    ws.title = safe_sheet_title("סיכום מחלקות", "סיכום מחלקות")
    ws.freeze_panes = "A2"

    headers = ["מחלקה", "עובדים", "חושבו בהצלחה", "שעות לתשלום", "שכר כולל"]
    for col_index, header in enumerate(headers, start=1):
        cell = ws.cell(row=1, column=col_index, value=header)
        cell.font = Font(bold=True, color="FFFFFF")
        cell.fill = PatternFill(fill_type="solid", fgColor="0F766E")
        cell.alignment = Alignment(horizontal="center")
    ws.auto_filter.ref = f"A1:{get_column_letter(len(headers))}1"

    department_totals = defaultdict(lambda: {"workers": 0, "calculated": 0, "hours": 0.0, "salary": 0.0})
    for worker in worker_rows:
        department = worker["department"] or "ללא מחלקה"
        bucket = department_totals[department]
        bucket["workers"] += 1
        if worker["status"] == "OK":
            bucket["calculated"] += 1
            bucket["hours"] += worker["payable_hours"] or 0
            bucket["salary"] += worker["salary"] or 0

    for row_index, (department, totals) in enumerate(sorted(department_totals.items()), start=2):
        values = [
            department,
            totals["workers"],
            totals["calculated"],
            format_hours(totals["hours"]),
            round(totals["salary"], 2),
        ]
        for col_index, value in enumerate(values, start=1):
            cell = ws.cell(row=row_index, column=col_index, value=value)
            cell.alignment = Alignment(horizontal="right")
            if row_index % 2 == 0:
                cell.fill = PatternFill(fill_type="solid", fgColor="F0FDFA")

    widths = [24, 12, 18, 16, 16]
    for col_index, width in enumerate(widths, start=1):
        ws.column_dimensions[get_column_letter(col_index)].width = width


def write_flamingo_top_earners_sheet(ws, worker_rows):
    ws.sheet_view.rightToLeft = True
    ws.sheet_view.showGridLines = False
    ws.title = safe_sheet_title("שכר גבוה", "שכר גבוה")
    ws.freeze_panes = "A2"

    headers = ["דירוג", "שם עובד", "תעודת זהות", "מחלקה", "תעריף שעתי", "שעות לתשלום", "שכר מחושב"]
    for col_index, header in enumerate(headers, start=1):
        cell = ws.cell(row=1, column=col_index, value=header)
        cell.font = Font(bold=True, color="FFFFFF")
        cell.fill = PatternFill(fill_type="solid", fgColor="7C3AED")
        cell.alignment = Alignment(horizontal="center")
    ws.auto_filter.ref = f"A1:{get_column_letter(len(headers))}1"

    ranked_workers = sorted(
        [row for row in worker_rows if row["status"] == "OK" and row["salary"] is not None],
        key=lambda row: row["salary"],
        reverse=True,
    )[:10]

    for row_index, worker in enumerate(ranked_workers, start=2):
        values = [
            row_index - 1,
            worker["worker_name"],
            worker["id_number"],
            worker["department"],
            worker["hourly_rate"],
            format_hours(worker["payable_hours"]),
            worker["salary"],
        ]
        for col_index, value in enumerate(values, start=1):
            cell = ws.cell(row=row_index, column=col_index, value=value)
            cell.alignment = Alignment(horizontal="right")
            if row_index % 2 == 0:
                cell.fill = PatternFill(fill_type="solid", fgColor="F5F3FF")
        ws.cell(row=row_index, column=5).number_format = '0.00'
        ws.cell(row=row_index, column=7).number_format = '#,##0.00'

    widths = [10, 22, 18, 20, 14, 14, 16]
    for col_index, width in enumerate(widths, start=1):
        ws.column_dimensions[get_column_letter(col_index)].width = width


def load_org_structure_csv(csv_path):
    records = {}
    unmatched = []
    with open(csv_path, "r", encoding="utf-8-sig", newline="") as handle:
        reader = csv.DictReader(handle)
        for row in reader:
            employee_number = (row.get("שכר") or row.get("׳©׳›׳¨") or "").strip()
            employee_id = (row.get("ת.ז") or row.get("׳×.׳–") or "").strip()
            entry = {
                "employee_number": employee_number,
                "id_number": employee_id,
                "employee_name": (row.get("שם עובד") or row.get("׳©׳ ׳¢׳•׳‘׳“") or "").strip(),
                "direct_manager": (row.get("מנהל ישיר") or row.get("׳׳ ׳”׳ ׳™׳©׳™׳¨") or "").strip(),
                "department": (row.get("מחלקה") or row.get("׳׳—׳׳§׳”") or "").strip(),
                "agreement_name": (row.get("שם הסכם") or row.get("׳©׳ ׳”׳¡׳›׳") or "").strip(),
                "agreement_number": (row.get("מס' הסכם") or row.get("׳׳¡' ׳”׳¡׳›׳") or "").strip(),
            }
            if employee_number:
                records[("number", employee_number)] = entry
            if employee_id:
                records[("id", employee_id)] = entry
            if not employee_number and not employee_id:
                unmatched.append(entry)
    return records, unmatched


def extract_org_mapping_value(raw_row, source):
    source_text = str(source or "").strip()
    if not source_text.startswith("header:"):
        return ""
    header = source_text.split(":", 1)[1]
    return str(raw_row.get(header) or "").strip()


def parse_org_hierarchy_csv(csv_path, mapping=None):
    rows = []
    with open(csv_path, "r", encoding="utf-8-sig", newline="") as handle:
        reader = csv.DictReader(handle)
        for raw in reader:
            if mapping:
                employee_name = extract_org_mapping_value(raw, mapping.get("employee_name_source"))
                employee_number = extract_org_mapping_value(raw, mapping.get("employee_number_source"))
                id_number = extract_org_mapping_value(raw, mapping.get("id_number_source"))
                passport_number = extract_org_mapping_value(raw, mapping.get("passport_number_source"))
                direct_manager = extract_org_mapping_value(raw, mapping.get("direct_manager_source"))
                manager_flag = extract_org_mapping_value(raw, mapping.get("manager_flag_source"))
                department = extract_org_mapping_value(raw, mapping.get("department_source"))
                email = extract_org_mapping_value(raw, mapping.get("email_source"))
                secondary_email = extract_org_mapping_value(raw, mapping.get("secondary_email_source"))
                app_access = extract_org_mapping_value(raw, mapping.get("app_access_source"))
                employment_percent = extract_org_mapping_value(raw, mapping.get("employment_percent_source"))
                agreement_number = extract_org_mapping_value(raw, mapping.get("agreement_number_source"))
                agreement_name = extract_org_mapping_value(raw, mapping.get("agreement_name_source"))
            else:
                employee_name = (raw.get("שם עובד") or raw.get("׳©׳ ׳¢׳•׳‘׳“") or "").strip()
                employee_number = (raw.get("שכר") or raw.get("׳©׳›׳¨") or "").strip()
                id_number = (raw.get("ת.ז") or raw.get("׳×.׳–") or "").strip()
                passport_number = ""
                direct_manager = (raw.get("מנהל ישיר") or raw.get("׳׳ ׳”׳ ׳™׳©׳™׳¨") or "").strip()
                manager_flag = (raw.get("מנהל") or "").strip()
                department = (raw.get("מחלקה") or raw.get("׳׳—׳׳§׳”") or "").strip()
                email = (raw.get("אימייל") or "").strip()
                secondary_email = (raw.get("אימייל נוסף") or "").strip()
                app_access = (raw.get("הרשאה לאפליקציה") or "").strip()
                employment_percent = (raw.get("אחוז משרה") or "").strip()
                agreement_number = (raw.get("מס' הסכם") or "").strip()
                agreement_name = (raw.get("שם הסכם") or "").strip()
            if not any([employee_name, employee_number, id_number, passport_number, direct_manager, department, email, secondary_email]):
                continue
            rows.append(
                {
                    "employee_name": employee_name,
                    "employee_number": employee_number,
                    "id_number": id_number or passport_number,
                    "direct_manager": direct_manager,
                    "is_manager": manager_flag in {"[+]", "+", "כן", "yes", "true", "1"},
                    "department": department,
                    "email": email,
                    "secondary_email": secondary_email,
                    "app_access": app_access,
                    "employment_percent": employment_percent,
                    "agreement_number": agreement_number,
                    "agreement_name": agreement_name,
                }
            )

    name_counts = defaultdict(int)
    id_counts = defaultdict(int)
    for row in rows:
        if row["employee_name"]:
            name_counts[row["employee_name"]] += 1
        if row["id_number"]:
            id_counts[row["id_number"]] += 1

    unique_name_map = {row["employee_name"]: row for row in rows if row["employee_name"] and name_counts[row["employee_name"]] == 1}
    children_map = defaultdict(list)
    exception_rows = []
    roots = []

    for row in rows:
        if not row["direct_manager"]:
            roots.append(row)
            exception_rows.append(
                {
                    "category": "צומת שורש",
                    "employee_name": row["employee_name"],
                    "employee_number": row["employee_number"],
                    "detail": "מנהל ישיר ריק",
                }
            )
            continue
        if row["direct_manager"] not in unique_name_map:
            roots.append(row)
            exception_rows.append(
                {
                    "category": "מנהל לא נמצא",
                    "employee_name": row["employee_name"],
                    "employee_number": row["employee_number"],
                    "detail": row["direct_manager"],
                }
            )
            continue
        children_map[row["direct_manager"]].append(row)

    for name, count in sorted(name_counts.items()):
        if count > 1:
            exception_rows.append(
                {
                    "category": "שם כפול",
                    "employee_name": name,
                    "employee_number": "",
                    "detail": f"נמצאו {count} שורות עם אותו שם עובד",
                }
            )
    for id_number, count in sorted(id_counts.items()):
        if id_number and count > 1:
            exception_rows.append(
                {
                    "category": "תעודת זהות כפולה",
                    "employee_name": "",
                    "employee_number": "",
                    "detail": f"תעודת הזהות {id_number} מופיעה {count} פעמים",
                }
            )

    visited = set()
    summary_rows = []
    tree_rows = []

    def walk(node, depth, root_name, root_key, stack):
        node_key = (node["employee_number"], node["id_number"], node["employee_name"])
        if node_key in stack:
            exception_rows.append(
                {
                    "category": "זוהתה לולאה",
                    "employee_name": node["employee_name"],
                    "employee_number": node["employee_number"],
                    "detail": "זוהתה לולאה בהיררכיה במהלך מעבר על מנהלים",
                }
            )
            return

        visited.add(node_key)
        summary_rows.append(
            {
                "employee_name": node["employee_name"],
                "employee_number": node["employee_number"],
                "id_number": node["id_number"],
                "direct_manager": node["direct_manager"],
                "is_manager": "כן" if node["is_manager"] or bool(children_map.get(node["employee_name"])) else "לא",
                "department": node["department"],
                "email": node["email"],
                "depth": depth,
                "root_name": root_name,
                "root_employee_number": root_key[0],
                "root_id_number": root_key[1],
            }
        )
        tree_rows.append(
            {
                "root_name": root_name,
                "depth": depth,
                "display_name": ("    " * depth) + node["employee_name"],
                "employee_number": node["employee_number"],
                "department": node["department"],
                "direct_manager": node["direct_manager"],
            }
        )
        next_stack = stack | {node_key}
        for child in sorted(children_map.get(node["employee_name"], []), key=lambda item: (item["employee_name"], item["employee_number"])):
            walk(child, depth + 1, root_name, root_key, next_stack)

    unique_roots = []
    seen_root_keys = set()
    for row in roots:
        root_key = (row["employee_number"], row["id_number"], row["employee_name"])
        if root_key not in seen_root_keys:
            unique_roots.append(row)
            seen_root_keys.add(root_key)

    for root in sorted(unique_roots, key=lambda item: (item["employee_name"], item["employee_number"])):
        walk(root, 0, root["employee_name"], (root["employee_number"], root["id_number"]), set())

    for row in rows:
        node_key = (row["employee_number"], row["id_number"], row["employee_name"])
        if node_key not in visited:
            exception_rows.append(
                {
                    "category": "צומת לא מקושר",
                    "employee_name": row["employee_name"],
                    "employee_number": row["employee_number"],
                    "detail": "השורה לא קושרה לאף צומת שורש",
                }
            )
            walk(row, 0, row["employee_name"], (row["employee_number"], row["id_number"]), set())

    stats = {
        "employee_count": len(rows),
        "root_count": len(unique_roots),
        "exception_count": len(exception_rows),
    }
    return summary_rows, tree_rows, exception_rows, stats


def write_org_hierarchy_summary(ws, summary_rows, stats, root_rows):
    ws.title = safe_sheet_title("סיכום מבנה ארגוני", "Org Summary")
    ws.sheet_view.rightToLeft = True
    ws.sheet_view.showGridLines = False
    ws.freeze_panes = "A11"

    ws["A1"] = "דוח מבנה ארגוני"
    ws["A1"].font = Font(bold=True, size=18, color="0F172A")
    ws["A1"].fill = PatternFill(fill_type="solid", fgColor="BFDBFE")

    direct_reports = defaultdict(int)
    department_counts = defaultdict(int)
    manager_rows = []
    for row in summary_rows:
        if row["direct_manager"]:
            direct_reports[row["direct_manager"]] += 1
        if row["department"]:
            department_counts[row["department"]] += 1
        if row["is_manager"] == "כן":
            manager_rows.append(row)

    metrics = [
        ("סה\"כ עובדים", stats["employee_count"], "DBEAFE"),
        ("סה\"כ שורשים (ללא מנהל ישיר או עם מנהל לא מזוהה)", stats["root_count"], "DCFCE7"),
        ("סה\"כ מנהלים", len(manager_rows), "EDE9FE"),
        ("סה\"כ חריגים", stats["exception_count"], "FEE2E2"),
    ]
    for idx, (label, value, fill_color) in enumerate(metrics, start=3):
        ws.cell(row=idx, column=1, value=label).font = Font(bold=True, color="334155")
        ws.cell(row=idx, column=2, value=value).font = Font(bold=True, color="0F172A")
        ws.cell(row=idx, column=1).fill = PatternFill(fill_type="solid", fgColor=fill_color)
        ws.cell(row=idx, column=2).fill = PatternFill(fill_type="solid", fgColor=fill_color)

    ws["D3"] = "שורשים שזוהו"
    ws["D3"].font = Font(bold=True, color="1E3A8A")
    ws["D3"].fill = PatternFill(fill_type="solid", fgColor="DBEAFE")
    ws["D4"] = ", ".join(root["employee_name"] for root in sorted(root_rows, key=lambda item: (item["employee_name"], item["employee_number"]))) or "אין"

    ws["D6"] = "רשימת מנהלים ומספר כפיפים ישירים"
    ws["D6"].font = Font(bold=True, color="1E3A8A")
    ws["D6"].fill = PatternFill(fill_type="solid", fgColor="EDE9FE")
    ws["D7"] = "מנהל"
    ws["E7"] = "כפיפים ישירים"
    for cell_ref in ("D7", "E7"):
        ws[cell_ref].font = Font(bold=True, color="FFFFFF")
        ws[cell_ref].fill = PatternFill(fill_type="solid", fgColor="7C3AED")
    for row_idx, row in enumerate(sorted(manager_rows, key=lambda item: (-direct_reports.get(item["employee_name"], 0), item["employee_name"])), start=8):
        ws.cell(row=row_idx, column=4, value=row["employee_name"])
        ws.cell(row=row_idx, column=5, value=direct_reports.get(row["employee_name"], 0))
        if row_idx % 2 == 0:
            ws.cell(row=row_idx, column=4).fill = PatternFill(fill_type="solid", fgColor="F5F3FF")
            ws.cell(row=row_idx, column=5).fill = PatternFill(fill_type="solid", fgColor="F5F3FF")

    dept_start_row = 8
    ws["G6"] = "רשימת מחלקות וכמות עובדים"
    ws["G6"].font = Font(bold=True, color="0F766E")
    ws["G6"].fill = PatternFill(fill_type="solid", fgColor="CCFBF1")
    ws["G7"] = "מחלקה"
    ws["H7"] = "כמות עובדים"
    for cell_ref in ("G7", "H7"):
        ws[cell_ref].font = Font(bold=True, color="FFFFFF")
        ws[cell_ref].fill = PatternFill(fill_type="solid", fgColor="0F766E")
    for row_idx, (department, count) in enumerate(sorted(department_counts.items(), key=lambda item: (-item[1], item[0])), start=dept_start_row):
        ws.cell(row=row_idx, column=7, value=department)
        ws.cell(row=row_idx, column=8, value=count)
        if row_idx % 2 == 0:
            ws.cell(row=row_idx, column=7).fill = PatternFill(fill_type="solid", fgColor="ECFDF5")
            ws.cell(row=row_idx, column=8).fill = PatternFill(fill_type="solid", fgColor="ECFDF5")

    headers = [
        "שם עובד",
        "מספר שכר",
        "תעודת זהות",
        "מנהל ישיר",
        "האם מנהל",
        "מחלקה",
        "אימייל",
        "עומק בעץ",
        "שם שורש",
    ]
    header_row = 10
    for col, header in enumerate(headers, start=1):
        cell = ws.cell(row=header_row, column=col, value=header)
        cell.font = Font(bold=True, color="FFFFFF")
        cell.fill = PatternFill(fill_type="solid", fgColor="1E3A8A")

    for row_idx, row in enumerate(sorted(summary_rows, key=lambda item: (item["root_name"], item["depth"], item["employee_name"])), start=header_row + 1):
        values = [
            row["employee_name"],
            row["employee_number"],
            row["id_number"],
            row["direct_manager"],
            row["is_manager"],
            row["department"],
            row["email"],
            row["depth"],
            row["root_name"],
        ]
        for col, value in enumerate(values, start=1):
            ws.cell(row=row_idx, column=col, value=value)
            if row_idx % 2 == 0:
                ws.cell(row=row_idx, column=col).fill = PatternFill(fill_type="solid", fgColor="F8FAFC")

    widths = [24, 14, 16, 22, 12, 24, 26, 10, 20]
    for col, width in enumerate(widths, start=1):
        ws.column_dimensions[get_column_letter(col)].width = width
    ws.column_dimensions["D"].width = 28
    ws.column_dimensions["E"].width = 14
    ws.column_dimensions["G"].width = 28
    ws.column_dimensions["H"].width = 14


def write_org_hierarchy_tree(ws, tree_rows):
    ws.title = safe_sheet_title("טבלת היררכיה", "Hierarchy Table")
    ws.sheet_view.rightToLeft = True
    ws.sheet_view.showGridLines = False
    ws.freeze_panes = "A5"

    ws["A1"] = "טבלת היררכיה ארגונית"
    ws["A1"].font = Font(bold=True, size=18, color="0F172A")
    ws["A1"].fill = PatternFill(fill_type="solid", fgColor="CCFBF1")
    ws["A2"] = "הטבלה מציגה נתיב דיווח ורמת עומק. ההזחה מסייעת לקריאה, אך אינה תרשים חזותי מלא."
    ws["A2"].font = Font(bold=True, color="0F766E")

    headers = ["שם עובד בהזחה לפי רמה", "רמה", "מספר שכר", "מחלקה", "מנהל ישיר"]
    for col, header in enumerate(headers, start=1):
        cell = ws.cell(row=4, column=col, value=header)
        cell.font = Font(bold=True, color="FFFFFF")
        cell.fill = PatternFill(fill_type="solid", fgColor="0F766E")

    row_idx = 5
    current_root = None
    for row in sorted(tree_rows, key=lambda item: (item["root_name"], item["depth"], item["display_name"])):
        if row["root_name"] != current_root:
            if current_root is not None:
                row_idx += 1
            ws.cell(row=row_idx, column=1, value=f"מקטע שורש: {row['root_name']}")
            ws.cell(row=row_idx, column=1).font = Font(bold=True, color="1E3A8A")
            ws.cell(row=row_idx, column=1).fill = PatternFill(fill_type="solid", fgColor="DBEAFE")
            current_root = row["root_name"]
            row_idx += 1
        values = [row["display_name"], row["depth"], row["employee_number"], row["department"], row["direct_manager"]]
        for col, value in enumerate(values, start=1):
            ws.cell(row=row_idx, column=col, value=value)
            if row_idx % 2 == 0:
                ws.cell(row=row_idx, column=col).fill = PatternFill(fill_type="solid", fgColor="ECFDF5")
        row_idx += 1

    widths = [40, 10, 14, 24, 22]
    for col, width in enumerate(widths, start=1):
        ws.column_dimensions[get_column_letter(col)].width = width


def write_org_hierarchy_exceptions(ws, exception_rows, root_rows):
    ws.title = safe_sheet_title("חריגים", "Exceptions")
    ws.sheet_view.rightToLeft = True
    ws.sheet_view.showGridLines = False
    ws.freeze_panes = "A2"

    headers = ["סוג חריג", "שם עובד", "מספר שכר", "פירוט"]
    for col, header in enumerate(headers, start=1):
        cell = ws.cell(row=1, column=col, value=header)
        cell.font = Font(bold=True, color="FFFFFF")
        cell.fill = PatternFill(fill_type="solid", fgColor="B91C1C")

    row_idx = 2
    for row in exception_rows:
        values = [row["category"], row["employee_name"], row["employee_number"], row["detail"]]
        for col, value in enumerate(values, start=1):
            ws.cell(row=row_idx, column=col, value=value)
            if row_idx % 2 == 0:
                ws.cell(row=row_idx, column=col).fill = PatternFill(fill_type="solid", fgColor="FEF2F2")
        row_idx += 1

    if row_idx == 2:
        ws.cell(row=row_idx, column=1, value="אין חריגים")
        row_idx += 1

    row_idx += 1
    ws.cell(row=row_idx, column=1, value="רשימת שורשים").font = Font(bold=True, color="1E3A8A")
    ws.cell(row=row_idx, column=1).fill = PatternFill(fill_type="solid", fgColor="DBEAFE")
    row_idx += 1
    for root in sorted(root_rows, key=lambda item: (item["employee_name"], item["employee_number"])):
        ws.cell(row=row_idx, column=1, value=root["employee_name"])
        ws.cell(row=row_idx, column=2, value=root["employee_number"])
        ws.cell(row=row_idx, column=3, value=root["department"])
        row_idx += 1

    widths = [20, 24, 14, 36]
    for col, width in enumerate(widths, start=1):
        ws.column_dimensions[get_column_letter(col)].width = width


def build_org_summary_slide_lines(summary_rows, stats):
    manager_count = sum(1 for row in summary_rows if row["is_manager"] == "כן")
    return [
        f"סה\"כ עובדים: {stats['employee_count']}",
        f"סה\"כ שורשים: {stats['root_count']}",
        f"סה\"כ מנהלים: {manager_count}",
        f"סה\"כ חריגים: {stats['exception_count']}",
    ]


def build_org_root_slide_lines(tree_rows, root_name):
    lines = []
    for row in tree_rows:
        if row["root_name"] != root_name:
            continue
        meta = []
        if row["department"]:
            meta.append(row["department"])
        if row["employee_number"]:
            meta.append(f"מס' שכר {row['employee_number']}")
        line = row["display_name"].strip()
        if meta:
            line += " | " + " | ".join(meta)
        lines.append((row["depth"], line))
    return lines


def build_org_exceptions_slide_lines(exception_rows):
    if not exception_rows:
        return ["אין חריגים"]
    return [f"{row['category']}: {row['employee_name']} | {row['detail']}".strip(" |") for row in exception_rows]


def format_display_name(name):
    parts = str(name or "").split()
    if len(parts) <= 1:
        return str(name or "").strip()
    return " ".join(reversed(parts))


def build_org_root_nodes(summary_rows, root_name):
    nodes = []
    for row in summary_rows:
        row_root_key = (row.get("root_employee_number", ""), row.get("root_id_number", ""), row["root_name"])
        if row_root_key != root_name:
            continue
        node_key = (row["employee_number"], row["id_number"], row["employee_name"])
        nodes.append(
            {
                "node_key": node_key,
                "employee_name": row["employee_name"],
                "employee_number": row["employee_number"],
                "id_number": row["id_number"],
                "department": row["department"],
                "direct_manager": row["direct_manager"],
                "depth": row["depth"],
            }
        )
    keys_by_name = defaultdict(list)
    for node in nodes:
        keys_by_name[node["employee_name"]].append(node["node_key"])
    children_map = defaultdict(list)
    for node in nodes:
        manager_keys = keys_by_name.get(node["direct_manager"], [])
        node["parent_key"] = manager_keys[0] if len(manager_keys) == 1 else None
        node["manager_resolution"] = "unique" if len(manager_keys) == 1 else ("ambiguous" if len(manager_keys) > 1 else "missing")
        children_map[node["parent_key"]].append(node)
    for parent_key, children in children_map.items():
        children.sort(key=lambda item: (item["employee_name"], item["department"]))
    top_nodes = [node for node in nodes if node.get("parent_key") is None]
    metrics = {}
    for top_node in top_nodes:
        compute_org_subtree_metrics(top_node, children_map, metrics)
    for node in nodes:
        node["direct_reports_count"] = metrics.get(node["node_key"], {}).get("direct_reports_count", 0)
        node["subtree_employee_count"] = metrics.get(node["node_key"], {}).get("subtree_employee_count", 1)
    return nodes, children_map


def compute_org_subtree_metrics(node, children_map, metrics):
    children = children_map.get(node["node_key"], [])
    subtree_total = 1
    for child in children:
        subtree_total += compute_org_subtree_metrics(child, children_map, metrics)
    metrics[node["node_key"]] = {
        "direct_reports_count": len(children),
        "subtree_employee_count": subtree_total,
    }
    return subtree_total


def compute_org_subtree_widths(node, children_map, widths):
    children = children_map.get(node["node_key"], [])
    if not children:
        widths[node["node_key"]] = 1.0
        return 1.0
    total_width = 0.0
    for child in children:
        total_width += compute_org_subtree_widths(child, children_map, widths)
    widths[node["node_key"]] = max(total_width, 1.0)
    return widths[node["node_key"]]


def layout_org_chart(node, children_map, widths, positions, left_units=0.0):
    node_width = widths[node["node_key"]]
    children = children_map.get(node["node_key"], [])
    if not children:
        positions[node["node_key"]] = left_units + (node_width / 2.0)
        return
    current_left = left_units
    for child in children:
        child_width = widths[child["node_key"]]
        layout_org_chart(child, children_map, widths, positions, current_left)
        current_left += child_width
    first_child = children[0]["node_key"]
    last_child = children[-1]["node_key"]
    positions[node["node_key"]] = (positions[first_child] + positions[last_child]) / 2.0


def build_org_chart_subset(top_nodes, children_map, level_limit=None):
    subset_nodes = []
    subset_keys = set()

    def visit(node, chart_depth):
        node_copy = dict(node)
        node_copy["chart_depth"] = chart_depth
        subset_nodes.append(node_copy)
        subset_keys.add(node["node_key"])
        if level_limit is not None and chart_depth >= level_limit:
            return
        for child in children_map.get(node["node_key"], []):
            visit(child, chart_depth + 1)

    for top_node in top_nodes:
        visit(top_node, 0)

    subset_children_map = defaultdict(list)
    for node in subset_nodes:
        parent_key = node.get("parent_key")
        if parent_key not in subset_keys:
            parent_key = None
        subset_children_map[parent_key].append(node)
    for parent_key, children in subset_children_map.items():
        children.sort(key=lambda item: (item["employee_name"], item["department"]))
    return subset_nodes, subset_children_map


def build_org_detail_slide_nodes(root_node, children_map):
    detail_nodes = []

    def visit(node):
        if node["direct_reports_count"] > 0 and (
            node["direct_reports_count"] >= 5 or node["subtree_employee_count"] >= 12
        ):
            detail_nodes.append(node)
        for child in children_map.get(node["node_key"], []):
            visit(child)

    for child in children_map.get(root_node["node_key"], []):
        visit(child)
    return detail_nodes


def build_department_summary_rows(summary_rows):
    department_rows = []
    grouped = defaultdict(list)
    for row in summary_rows:
        grouped[row["department"] or "ללא מחלקה"].append(row)
    for department, rows in sorted(grouped.items(), key=lambda item: (-len(item[1]), item[0])):
        managers = [row for row in rows if row["is_manager"] == "כן"]
        manager_names = [format_display_name(row["employee_name"]) for row in sorted(managers, key=lambda item: item["employee_name"])[:4]]
        department_rows.append(
            {
                "department": department,
                "employee_count": len(rows),
                "manager_count": len(managers),
                "manager_names": manager_names,
            }
        )
    return department_rows


def build_manager_summary_rows(summary_rows, root_nodes):
    manager_rows = []
    for root_node in root_nodes:
        nodes, children_map = build_org_root_nodes(summary_rows, root_node["root_key"])
        node_map = {node["node_key"]: node for node in nodes}
        for node in sorted(nodes, key=lambda item: (-item["subtree_employee_count"], item["employee_name"])):
            if node["direct_reports_count"] <= 0:
                continue
            direct_reports = children_map.get(node["node_key"], [])
            manager_rows.append(
                {
                    "manager_name": format_display_name(node["employee_name"]),
                    "department": node["department"] or "ללא מחלקה",
                    "direct_reports_count": node["direct_reports_count"],
                    "subtree_employee_count": node["subtree_employee_count"],
                    "employee_names": [format_display_name(child["employee_name"]) for child in direct_reports],
                }
            )
    return manager_rows


def add_pptx_bullets(text_frame, lines, font_size=18):
    text_frame.clear()
    if not lines:
        lines = [""]
    for index, line in enumerate(lines):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.text = str(line.get("text") if isinstance(line, dict) else line)
        paragraph.level = int(line.get("level", 0)) if isinstance(line, dict) else 0
        paragraph.alignment = PP_ALIGN.RIGHT
        if paragraph.runs:
            run = paragraph.runs[0]
            run.font.size = Pt(line.get("size", font_size) if isinstance(line, dict) else font_size)
            run.font.bold = bool(line.get("bold", False)) if isinstance(line, dict) else False


def add_pptx_slide(prs, title, body_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    title_frame = slide.shapes.title.text_frame
    for paragraph in title_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.RIGHT
        if paragraph.runs:
            paragraph.runs[0].font.size = Pt(26)
            paragraph.runs[0].font.bold = True
    body = slide.placeholders[1]
    add_pptx_bullets(body.text_frame, body_lines, font_size=18)
    return slide


def add_org_chart_slide(prs, title, top_nodes, children_map, note_text, level_limit=2, compact=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.45), Inches(0.2), Inches(12.2), Inches(0.55))
    title_frame = title_box.text_frame
    title_frame.clear()
    title_paragraph = title_frame.paragraphs[0]
    title_paragraph.text = title
    title_paragraph.alignment = PP_ALIGN.RIGHT
    if title_paragraph.runs:
        title_run = title_paragraph.runs[0]
        title_run.font.size = Pt(24)
        title_run.font.bold = True

    note_box = slide.shapes.add_textbox(Inches(0.45), Inches(0.78), Inches(12.2), Inches(0.3))
    note_frame = note_box.text_frame
    note_frame.clear()
    note_paragraph = note_frame.paragraphs[0]
    note_paragraph.text = note_text
    note_paragraph.alignment = PP_ALIGN.RIGHT
    if note_paragraph.runs:
        note_run = note_paragraph.runs[0]
        note_run.font.size = Pt(11)
        note_run.font.color.rgb = RGBColor(71, 85, 105)

    nodes, children_map = build_org_chart_subset(top_nodes, children_map, level_limit=level_limit)
    if not nodes:
        return slide, {}

    box_nodes = [node for node in nodes if node["direct_reports_count"] > 0]
    if not box_nodes:
        return slide, {}
    box_node_keys = {node["node_key"] for node in box_nodes}
    top_node_keys = {node["node_key"] for node in box_nodes if node.get("chart_depth", 0) == 0}
    box_children_map = defaultdict(list)
    for parent_key, children in children_map.items():
        visible_box_children = [child for child in children if child["node_key"] in box_node_keys]
        if visible_box_children or parent_key in box_node_keys or parent_key is None:
            box_children_map[parent_key].extend(visible_box_children)
            box_children_map[parent_key].sort(key=lambda item: (item["employee_name"], item["department"]))

    widths = {}
    total_units = 0.0
    left_units = 0.0
    positions = {}
    for top_node in sorted([node for node in box_nodes if node.get("chart_depth", 0) == 0], key=lambda item: (item["employee_name"], item["department"])):
        subtree_width = compute_org_subtree_widths(top_node, box_children_map, widths)
        layout_org_chart(top_node, box_children_map, widths, positions, left_units)
        left_units += subtree_width
        total_units += subtree_width

    slide_width = prs.slide_width
    slide_height = prs.slide_height
    left_margin = Inches(0.3)
    right_margin = Inches(0.3)
    top_margin = Inches(1.2)
    bottom_margin = Inches(0.35)
    chart_width = slide_width - left_margin - right_margin
    chart_height = slide_height - top_margin - bottom_margin
    max_depth = max((row.get("chart_depth", 0) for row in box_nodes), default=0)
    level_count = max_depth + 1
    vertical_gap = chart_height / max(level_count, 1)
    total_units = max(total_units, 1.0)

    if compact:
        base_box_height = min(int(Inches(0.62)), max(int(Inches(0.34)), int(vertical_gap * 0.52)))
        manager_box_height = min(int(Inches(0.84)), max(int(Inches(0.48)), int(vertical_gap * 0.72)))
        min_box_width = int(Inches(1.05))
        max_box_width = int(Inches(1.6))
        name_size_regular = 9.5
        name_size_manager = 9
        dept_size_regular = 7.5
        dept_size_manager = 7
        metrics_size = 6.5
    else:
        base_box_height = min(int(Inches(0.72)), max(int(Inches(0.42)), int(vertical_gap * 0.58)))
        manager_box_height = min(int(Inches(0.96)), max(int(Inches(0.56)), int(vertical_gap * 0.82)))
        min_box_width = int(Inches(1.35))
        max_box_width = int(Inches(1.95))
        name_size_regular = 11
        name_size_manager = 10.5
        dept_size_regular = 9
        dept_size_manager = 8.5
        metrics_size = 7.5
    shape_bounds = {}
    node_shapes = {}

    sorted_nodes = sorted(box_nodes, key=lambda item: (item.get("chart_depth", 0), positions.get(item["node_key"], 0.0), item["employee_name"]))
    for node in sorted_nodes:
        node_width_units = widths.get(node["node_key"], 1.0)
        center_unit = positions.get(node["node_key"], 0.5)
        depth = node.get("chart_depth", 0)
        children = box_children_map.get(node["node_key"], [])

        center_x = int(left_margin + (center_unit / total_units) * chart_width)
        center_y = int(top_margin + depth * vertical_gap + vertical_gap * 0.45)
        width_by_subtree = int((node_width_units / total_units) * chart_width * 0.72)
        box_width = max(min_box_width, min(max_box_width, width_by_subtree))
        box_height = manager_box_height if node["direct_reports_count"] > 0 else base_box_height
        x = max(int(left_margin), min(int(slide_width - right_margin - box_width), center_x - box_width // 2))
        y = max(int(top_margin), min(int(slide_height - bottom_margin - box_height), center_y - box_height // 2))

        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, box_width, box_height)
        shape.line.color.rgb = RGBColor(100, 116, 139)
        shape.line.width = Pt(1)
        shape.fill.solid()
        if depth == 0:
            shape.fill.fore_color.rgb = RGBColor(191, 219, 254)
        elif node["node_key"] in top_node_keys:
            shape.fill.fore_color.rgb = RGBColor(254, 240, 138)
        elif children:
            shape.fill.fore_color.rgb = RGBColor(220, 252, 231)
        else:
            shape.fill.fore_color.rgb = RGBColor(241, 245, 249)

        text_frame = shape.text_frame
        text_frame.clear()
        text_frame.word_wrap = True
        text_frame.margin_top = int(Inches(0.05))
        text_frame.margin_bottom = int(Inches(0.04))
        text_frame.margin_left = int(Inches(0.05))
        text_frame.margin_right = int(Inches(0.05))
        name_paragraph = text_frame.paragraphs[0]
        name_paragraph.text = format_display_name(node["employee_name"])
        name_paragraph.alignment = PP_ALIGN.CENTER
        if name_paragraph.runs:
            name_run = name_paragraph.runs[0]
            name_run.font.bold = True
            name_run.font.size = Pt(12 if depth == 0 else (name_size_manager if node["direct_reports_count"] > 0 else name_size_regular))
            name_run.font.color.rgb = RGBColor(30, 41, 59)

        dept_paragraph = text_frame.add_paragraph()
        dept_paragraph.text = node["department"] or ""
        dept_paragraph.alignment = PP_ALIGN.CENTER
        if dept_paragraph.runs:
            dept_run = dept_paragraph.runs[0]
            dept_run.font.size = Pt(dept_size_manager if node["direct_reports_count"] > 0 else dept_size_regular)
            dept_run.font.color.rgb = RGBColor(71, 85, 105)

        if node["direct_reports_count"] > 0:
            metrics_paragraph = text_frame.add_paragraph()
            metrics_paragraph.text = f"כפיפים ישירים: {node['direct_reports_count']} | סה\"כ בענף: {node['subtree_employee_count']}"
            metrics_paragraph.alignment = PP_ALIGN.CENTER
            if metrics_paragraph.runs:
                metrics_run = metrics_paragraph.runs[0]
                metrics_run.font.size = Pt(metrics_size)
                metrics_run.font.bold = True
                metrics_run.font.color.rgb = RGBColor(51, 65, 85)

        shape_bounds[node["node_key"]] = {"x": x, "y": y, "w": box_width, "h": box_height}
        node_shapes[node["node_key"]] = shape

    for node in sorted_nodes:
        parent_bounds = shape_bounds.get(node["node_key"])
        if not parent_bounds:
            continue
        for child in box_children_map.get(node["node_key"], []):
            child_bounds = shape_bounds.get(child["node_key"])
            if not child_bounds:
                continue
            connector = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                parent_bounds["x"] + parent_bounds["w"] // 2,
                parent_bounds["y"] + parent_bounds["h"],
                child_bounds["x"] + child_bounds["w"] // 2,
                child_bounds["y"],
            )
            connector.line.color.rgb = RGBColor(148, 163, 184)
            connector.line.width = Pt(1.2)
    return slide, node_shapes


def add_department_summary_slide(prs, title, department_rows, page_index=1, page_count=1):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.45), Inches(0.2), Inches(12.2), Inches(0.55))
    title_frame = title_box.text_frame
    title_frame.clear()
    title_paragraph = title_frame.paragraphs[0]
    title_paragraph.text = title if page_count == 1 else f"{title} ({page_index}/{page_count})"
    title_paragraph.alignment = PP_ALIGN.RIGHT
    if title_paragraph.runs:
        title_run = title_paragraph.runs[0]
        title_run.font.size = Pt(24)
        title_run.font.bold = True

    cols = 2
    rows = 3
    card_width = Inches(6.15)
    card_height = Inches(1.55)
    start_x = Inches(0.45)
    start_y = Inches(1.1)
    gap_x = Inches(0.18)
    gap_y = Inches(0.16)
    for index, department in enumerate(department_rows):
        row = index // cols
        col = index % cols
        x = start_x + col * (card_width + gap_x)
        y = start_y + row * (card_height + gap_y)
        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, card_width, card_height)
        card.line.color.rgb = RGBColor(15, 118, 110)
        card.line.width = Pt(1)
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(236, 253, 245)
        frame = card.text_frame
        frame.clear()
        frame.margin_top = int(Inches(0.05))
        frame.margin_bottom = int(Inches(0.04))
        frame.margin_left = int(Inches(0.06))
        frame.margin_right = int(Inches(0.06))

        p1 = frame.paragraphs[0]
        p1.text = department["department"]
        p1.alignment = PP_ALIGN.CENTER
        if p1.runs:
            p1.runs[0].font.bold = True
            p1.runs[0].font.size = Pt(13)
            p1.runs[0].font.color.rgb = RGBColor(15, 23, 42)
        p2 = frame.add_paragraph()
        p2.text = f"עובדים: {department['employee_count']} | מנהלים: {department['manager_count']}"
        p2.alignment = PP_ALIGN.CENTER
        if p2.runs:
            p2.runs[0].font.size = Pt(10)
            p2.runs[0].font.color.rgb = RGBColor(51, 65, 85)
        p3 = frame.add_paragraph()
        p3.text = "מנהלים מובילים: " + (", ".join(department["manager_names"]) if department["manager_names"] else "ללא")
        p3.alignment = PP_ALIGN.CENTER
        if p3.runs:
            p3.runs[0].font.size = Pt(9)
            p3.runs[0].font.color.rgb = RGBColor(71, 85, 105)
    return slide


def add_manager_summary_slide(prs, title, manager_rows, page_index=1, page_count=1):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.45), Inches(0.2), Inches(12.2), Inches(0.55))
    title_frame = title_box.text_frame
    title_frame.clear()
    title_paragraph = title_frame.paragraphs[0]
    title_paragraph.text = title if page_count == 1 else f"{title} ({page_index}/{page_count})"
    title_paragraph.alignment = PP_ALIGN.RIGHT
    if title_paragraph.runs:
        title_run = title_paragraph.runs[0]
        title_run.font.size = Pt(24)
        title_run.font.bold = True

    cols = 2
    card_width = Inches(6.1)
    card_height = Inches(2.0)
    start_x = Inches(0.45)
    start_y = Inches(1.05)
    gap_x = Inches(0.25)
    gap_y = Inches(0.18)
    for index, manager in enumerate(manager_rows):
        row = index // cols
        col = index % cols
        x = start_x + col * (card_width + gap_x)
        y = start_y + row * (card_height + gap_y)
        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, card_width, card_height)
        card.line.color.rgb = RGBColor(100, 116, 139)
        card.line.width = Pt(1)
        card.fill.solid()
        if manager["subtree_employee_count"] >= 12:
            card.fill.fore_color.rgb = RGBColor(254, 240, 138)
        elif manager["direct_reports_count"] >= 5:
            card.fill.fore_color.rgb = RGBColor(220, 252, 231)
        else:
            card.fill.fore_color.rgb = RGBColor(239, 246, 255)
        frame = card.text_frame
        frame.clear()
        frame.margin_top = int(Inches(0.05))
        frame.margin_bottom = int(Inches(0.04))
        frame.margin_left = int(Inches(0.06))
        frame.margin_right = int(Inches(0.06))

        p1 = frame.paragraphs[0]
        p1.text = manager["manager_name"]
        p1.alignment = PP_ALIGN.CENTER
        if p1.runs:
            p1.runs[0].font.bold = True
            p1.runs[0].font.size = Pt(13)
            p1.runs[0].font.color.rgb = RGBColor(30, 41, 59)
        p2 = frame.add_paragraph()
        p2.text = f"{manager['department']} | כפיפים ישירים: {manager['direct_reports_count']} | סה\"כ בענף: {manager['subtree_employee_count']}"
        p2.alignment = PP_ALIGN.CENTER
        if p2.runs:
            p2.runs[0].font.size = Pt(9.5)
            p2.runs[0].font.color.rgb = RGBColor(51, 65, 85)

        employee_names = manager["employee_names"]
        if employee_names:
            p3 = frame.add_paragraph()
            p3.text = "עובדים:"
            p3.alignment = PP_ALIGN.CENTER
            if p3.runs:
                p3.runs[0].font.size = Pt(9)
                p3.runs[0].font.bold = True
                p3.runs[0].font.color.rgb = RGBColor(30, 41, 59)

            names_per_line = 1 if len(employee_names) <= 4 else (2 if len(employee_names) <= 8 else 3)
            name_lines = []
            for start in range(0, len(employee_names), names_per_line):
                group = employee_names[start:start + names_per_line]
                name_lines.append(" • ".join(group))

            if len(employee_names) <= 4:
                name_font_size = 10.5
            elif len(employee_names) <= 8:
                name_font_size = 9.5
            else:
                name_font_size = 8.5

            for line in name_lines:
                p_name = frame.add_paragraph()
                p_name.text = line
                p_name.alignment = PP_ALIGN.CENTER
                if p_name.runs:
                    p_name.runs[0].font.size = Pt(name_font_size)
                    p_name.runs[0].font.color.rgb = RGBColor(51, 65, 85)
    return slide


def write_org_hierarchy_pptx(output_path, summary_rows, tree_rows, exception_rows, stats):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    manager_count = sum(1 for row in summary_rows if row["is_manager"] == "כן")
    department_rows = build_department_summary_rows(summary_rows)
    root_nodes = [
        {
            "root_key": (row["employee_number"], row["id_number"], row["employee_name"]),
            "title": format_display_name(row["employee_name"]),
        }
        for row in summary_rows
        if row["depth"] == 0
    ]
    manager_rows = build_manager_summary_rows(summary_rows, root_nodes)
    overview_top_nodes = []
    overview_children_map = defaultdict(list)
    for root_node in root_nodes:
        nodes, children_map = build_org_root_nodes(summary_rows, root_node["root_key"])
        actual_root = next((row for row in nodes if row["node_key"] == root_node["root_key"] and row["depth"] == 0), nodes[0] if nodes else None)
        if actual_root is None:
            continue
        top_nodes = [actual_root] + [
            node for node in nodes if node["node_key"] != actual_root["node_key"] and node.get("parent_key") is None
        ]
        overview_top_nodes.extend(top_nodes)
        for parent_key, children in children_map.items():
            overview_children_map[parent_key].extend(children)
    add_pptx_slide(
        prs,
        "סיכום מבנה ארגוני",
        [
            {"text": line, "size": 20}
            for line in (
                build_org_summary_slide_lines(summary_rows, stats)
                + [f"סה\"כ מחלקות: {len(department_rows)}"]
            )
        ],
    )
    add_org_chart_slide(
        prs,
        "מפת מנהלים ומחלקות",
        overview_top_nodes,
        overview_children_map,
        "תרשים ניהולי של מבנה הדיווח בין מנהלים. עובדים ללא צוות אינם מוצגים בתרשים זה.",
        level_limit=3,
        compact=True,
    )
    add_org_chart_slide(
        prs,
        "מבט-על ארגוני",
        overview_top_nodes,
        overview_children_map,
        (
            f"סקירת כלל הארגון לפי מחלקות ומנהלים | עובדים: {stats['employee_count']} | "
            f"מנהלים: {manager_count} | מחלקות: {len(department_rows)}"
        ),
        level_limit=None,
        compact=True,
    )
    department_page_size = 6
    department_pages = max(1, (len(department_rows) + department_page_size - 1) // department_page_size)
    for page_index in range(department_pages):
        start = page_index * department_page_size
        end = start + department_page_size
        add_department_summary_slide(
            prs,
            "סיכום מחלקות",
            department_rows[start:end],
            page_index=page_index + 1,
            page_count=department_pages,
        )
    manager_page_size = 4
    manager_pages = max(1, (len(manager_rows) + manager_page_size - 1) // manager_page_size)
    for page_index in range(manager_pages):
        start = page_index * manager_page_size
        end = start + manager_page_size
        add_manager_summary_slide(
            prs,
            "סיכום מנהלים וצוותים",
            manager_rows[start:end],
            page_index=page_index + 1,
            page_count=manager_pages,
        )
    add_pptx_slide(
        prs,
        "חריגים",
        [{"text": line, "size": 18} for line in build_org_exceptions_slide_lines(exception_rows)],
    )
    prs.save(output_path)


def run_org_hierarchy_report(input_path, output_path, extension, options=None):
    if extension != "csv":
        raise ValueError("Organizational hierarchy report currently supports CSV input only")
    options = options or {}
    output_type = options.get("output_type", "both").strip() or "both"
    warnings = build_org_hierarchy_mapping_warnings(options)
    summary_rows, tree_rows, exception_rows, stats = parse_org_hierarchy_csv(input_path, options)
    root_rows = [row for row in summary_rows if row["depth"] == 0]
    output_file = Path(output_path)
    excel_output_path = output_file.with_name(output_file.stem + ".xlsx")
    pptx_output_path = output_file.with_name(output_file.stem + ".pptx")
    wb = Workbook()
    write_org_hierarchy_summary(wb.active, summary_rows, stats, root_rows)
    write_org_hierarchy_tree(wb.create_sheet(), tree_rows)
    write_org_hierarchy_exceptions(wb.create_sheet(), exception_rows, root_rows)
    if output_type == "excel":
        wb.save(output_file)
        return {"warnings": warnings}
    if output_type == "powerpoint":
        write_org_hierarchy_pptx(str(output_file), summary_rows, tree_rows, exception_rows, stats)
        return {"warnings": warnings}

    wb.save(excel_output_path)
    write_org_hierarchy_pptx(str(pptx_output_path), summary_rows, tree_rows, exception_rows, stats)
    with ZipFile(output_file, "w", compression=ZIP_DEFLATED) as bundle:
        bundle.write(excel_output_path, arcname=excel_output_path.name)
        bundle.write(pptx_output_path, arcname=pptx_output_path.name)
    for temp_path in (excel_output_path, pptx_output_path):
        try:
            temp_path.unlink()
        except OSError:
            pass
    return {"warnings": warnings}


def parse_matan_missing_report(input_path, mapping):
    wb = xlrd.open_workbook(input_path)
    ws = wb.sheet_by_index(0)
    header_row = detect_matan_missing_header_row(ws)
    rows = []
    for row_index in range(header_row + 1, ws.nrows):
        employee_number = stringify_excel_value(extract_matan_missing_mapping_value(ws, mapping.get("employee_number_source"), row_index))
        employee_name = stringify_excel_value(extract_matan_missing_mapping_value(ws, mapping.get("employee_name_source"), row_index))
        id_number = stringify_excel_value(extract_matan_missing_mapping_value(ws, mapping.get("id_number_source"), row_index))
        badge_number = stringify_excel_value(extract_matan_missing_mapping_value(ws, mapping.get("badge_number_source"), row_index))
        passport_number = stringify_excel_value(extract_matan_missing_mapping_value(ws, mapping.get("passport_number_source"), row_index))
        if not any([employee_number, employee_name, id_number, badge_number, passport_number]):
            continue
        row = {
            "employee_number": employee_number,
            "id_number": id_number,
            "badge_number": badge_number,
            "passport_number": passport_number,
            "month": stringify_excel_value(extract_matan_missing_mapping_value(ws, mapping.get("month_source"), row_index)),
            "employee_name": employee_name,
            "standard_hours": parse_hours_or_zero(extract_matan_missing_mapping_value(ws, mapping.get("standard_hours_source"), row_index)),
            "missing_hours": parse_hours_or_zero(extract_matan_missing_mapping_value(ws, mapping.get("missing_hours_source"), row_index)),
            "attendance_hours": parse_hours_or_zero(extract_matan_missing_mapping_value(ws, mapping.get("attendance_hours_source"), row_index)),
            "vacation_hours": parse_hours_or_zero(extract_matan_missing_mapping_value(ws, mapping.get("vacation_hours_source"), row_index)),
            "sick_hours": parse_hours_or_zero(extract_matan_missing_mapping_value(ws, mapping.get("sick_hours_source"), row_index)),
            "reserve_hours": parse_hours_or_zero(extract_matan_missing_mapping_value(ws, mapping.get("reserve_hours_source"), row_index)),
            "pregnancy_hours": parse_hours_or_zero(extract_matan_missing_mapping_value(ws, mapping.get("pregnancy_hours_source"), row_index)),
            "special_child_hours": parse_hours_or_zero(extract_matan_missing_mapping_value(ws, mapping.get("special_child_hours_source"), row_index)),
            "absence_hours": parse_hours_or_zero(extract_matan_missing_mapping_value(ws, mapping.get("absence_hours_source"), row_index)),
        }
        rows.append(row)
    return rows


def apply_matan_missing_filters(rows, options):
    min_missing = parse_float_or_none(options.get("min_missing_hours", ""))
    max_missing = parse_float_or_none(options.get("max_missing_hours", ""))

    filtered = []
    for row in rows:
        missing_hours = row["missing_hours"] or 0.0
        if min_missing is not None and missing_hours < min_missing:
            continue
        if max_missing is not None and missing_hours > max_missing:
            continue
        filtered.append(row)
    return filtered


def get_matan_missing_selected_optional_fields(mapping):
    return [
        field_name
        for field_name in (
            "attendance_hours_source",
            "vacation_hours_source",
            "sick_hours_source",
            "reserve_hours_source",
            "pregnancy_hours_source",
            "special_child_hours_source",
            "absence_hours_source",
        )
        if mapping.get(field_name)
    ]


def write_matan_missing_summary(ws, filtered_rows, filters_used, mapping):
    ws.title = safe_sheet_title("סיכום שעות חסר", "סיכום שעות חסר")
    ws.sheet_view.rightToLeft = True
    ws.sheet_view.showGridLines = False
    ws.freeze_panes = "A10"

    total_missing = sum(row["missing_hours"] for row in filtered_rows)
    total_attendance = sum(row["attendance_hours"] for row in filtered_rows)
    total_standard = sum(row["standard_hours"] for row in filtered_rows)
    avg_missing = (total_missing / len(filtered_rows)) if filtered_rows else 0.0
    over_4 = sum(1 for row in filtered_rows if (row["missing_hours"] or 0.0) > 4.0)
    over_8 = sum(1 for row in filtered_rows if (row["missing_hours"] or 0.0) > 8.0)

    ws["A1"] = "דוח חוסר מול תקן"
    ws["A1"].font = Font(bold=True, size=18, color="0F172A")
    ws["A1"].fill = PatternFill(fill_type="solid", fgColor="BFDBFE")
    ws["A2"] = "סינון עובדים לפי שעות חוסר מול שעות תקן מתוך הדוח המרוכז"
    ws["A2"].font = Font(italic=True, size=11, color="475569")

    metrics = [
        ("עובדים בתוצאה", len(filtered_rows)),
        ("סה\"כ שעות חוסר", format_hours(total_missing)),
        ("ממוצע שעות חוסר", format_hours(avg_missing)),
        ("סה\"כ שעות נוכחות", format_hours(total_attendance)),
        ("סה\"כ שעות תקן", format_hours(total_standard)),
        ("עובדים מעל 4 שעות חוסר", over_4),
        ("עובדים מעל 8 שעות חוסר", over_8),
    ]
    for idx, (label, value) in enumerate(metrics, start=3):
        ws.cell(row=idx, column=1, value=label).font = Font(bold=True)
        ws.cell(row=idx, column=2, value=value)

    ws["D3"] = "פילטרים שהופעלו"
    ws["D3"].font = Font(bold=True)
    for idx, (label, value) in enumerate(filters_used.items(), start=4):
        ws.cell(row=idx, column=4, value=label).font = Font(bold=True)
        ws.cell(row=idx, column=5, value=value or "ללא")

    header_row = 10
    headers = ["שם עובד"]
    if mapping.get("employee_number_source"):
        headers.append("מספר עובד")
    if mapping.get("id_number_source"):
        headers.append("תעודת זהות")
    if mapping.get("badge_number_source"):
        headers.append("מספר תג")
    if mapping.get("passport_number_source"):
        headers.append("דרכון")
    if mapping.get("month_source"):
        headers.append("חודש")
    headers.extend(["שעות חוסר", "שעות תקן"])
    optional_columns = []
    optional_map = {
        "attendance_hours_source": "ש.נוכחות",
        "vacation_hours_source": "חופשה",
        "sick_hours_source": "מחלה",
        "reserve_hours_source": "מילואים",
        "pregnancy_hours_source": "שעות הריון",
        "special_child_hours_source": "שעות ילד מיוחד",
        "absence_hours_source": "היעדרות",
    }
    for field_name in get_matan_missing_selected_optional_fields(mapping):
        label = optional_map.get(field_name)
        if label:
            optional_columns.append((field_name, label))
            headers.append(label)
    for col, header in enumerate(headers, start=1):
        cell = ws.cell(row=header_row, column=col, value=header)
        cell.font = Font(bold=True, color="FFFFFF")
        cell.fill = PatternFill(fill_type="solid", fgColor="1E3A8A")
        cell.alignment = Alignment(horizontal="right")

    for row_idx, row in enumerate(filtered_rows, start=header_row + 1):
        values = [row["employee_name"]]
        if mapping.get("employee_number_source"):
            values.append(row["employee_number"])
        if mapping.get("id_number_source"):
            values.append(row["id_number"])
        if mapping.get("badge_number_source"):
            values.append(row["badge_number"])
        if mapping.get("passport_number_source"):
            values.append(row["passport_number"])
        if mapping.get("month_source"):
            values.append(row["month"])
        values.extend([format_hours(row["missing_hours"]), format_hours(row["standard_hours"])])
        for field_name, _ in optional_columns:
            row_key = field_name.replace("_source", "")
            values.append(format_hours(row.get(row_key, 0.0)))
        for col, value in enumerate(values, start=1):
            cell = ws.cell(row=row_idx, column=col, value=value)
            cell.alignment = Alignment(horizontal="right")
            if row_idx % 2 == 0:
                cell.fill = PatternFill(fill_type="solid", fgColor="F8FAFC")
    widths = [24, 16, 18, 16, 16, 10, 14, 14] + [12] * len(optional_columns)
    for col, width in enumerate(widths, start=1):
        ws.column_dimensions[get_column_letter(col)].width = width


def write_matan_missing_filtered(ws, filtered_rows, mapping):
    ws.title = safe_sheet_title("עובדים מסוננים", "עובדים מסוננים")
    ws.sheet_view.rightToLeft = True
    ws.sheet_view.showGridLines = False
    ws.freeze_panes = "A2"
    headers = ["שם עובד"]
    if mapping.get("employee_number_source"):
        headers.append("מספר עובד")
    if mapping.get("id_number_source"):
        headers.append("תעודת זהות")
    if mapping.get("badge_number_source"):
        headers.append("מספר תג")
    if mapping.get("passport_number_source"):
        headers.append("דרכון")
    if mapping.get("month_source"):
        headers.append("חודש")
    headers.extend(["שעות תקן", "שעות חוסר"])
    optional_columns = []
    optional_map = {
        "attendance_hours_source": "ש.נוכחות",
        "vacation_hours_source": "חופשה",
        "sick_hours_source": "מחלה",
        "reserve_hours_source": "מילואים",
        "pregnancy_hours_source": "שעות הריון",
        "special_child_hours_source": "שעות ילד מיוחד",
        "absence_hours_source": "היעדרות",
    }
    for field_name in get_matan_missing_selected_optional_fields(mapping):
        label = optional_map.get(field_name)
        if label:
            optional_columns.append((field_name, label))
            headers.append(label)
    for col, header in enumerate(headers, start=1):
        cell = ws.cell(row=1, column=col, value=header)
        cell.font = Font(bold=True, color="FFFFFF")
        cell.fill = PatternFill(fill_type="solid", fgColor="0F766E")
        cell.alignment = Alignment(horizontal="right")
    for row_idx, row in enumerate(filtered_rows, start=2):
        values = [row["employee_name"]]
        if mapping.get("employee_number_source"):
            values.append(row["employee_number"])
        if mapping.get("id_number_source"):
            values.append(row["id_number"])
        if mapping.get("badge_number_source"):
            values.append(row["badge_number"])
        if mapping.get("passport_number_source"):
            values.append(row["passport_number"])
        if mapping.get("month_source"):
            values.append(row["month"])
        values.extend([format_hours(row["standard_hours"]), format_hours(row["missing_hours"])])
        for field_name, _ in optional_columns:
            row_key = field_name.replace("_source", "")
            values.append(format_hours(row.get(row_key, 0.0)))
        for col, value in enumerate(values, start=1):
            cell = ws.cell(row=row_idx, column=col, value=value)
            cell.alignment = Alignment(horizontal="right")
    widths = [24, 16, 18, 16, 16, 10, 14, 14] + [12] * len(optional_columns)
    for col, width in enumerate(widths, start=1):
        ws.column_dimensions[get_column_letter(col)].width = width


def write_matan_missing_unmatched(ws, unmatched_rows):
    ws.title = safe_sheet_title("Unmatched Employees", "Unmatched Employees")
    ws.sheet_view.rightToLeft = True
    ws.sheet_view.showGridLines = False
    headers = ["Employee Number", "Employee Name", "Note"]
    for col, header in enumerate(headers, start=1):
        cell = ws.cell(row=1, column=col, value=header)
        cell.font = Font(bold=True, color="FFFFFF")
        cell.fill = PatternFill(fill_type="solid", fgColor="B91C1C")
    for row_idx, row in enumerate(unmatched_rows, start=2):
        ws.cell(row=row_idx, column=1, value=row["employee_number"])
        ws.cell(row=row_idx, column=2, value=row["employee_name"])
        ws.cell(row=row_idx, column=3, value="No matching record was found in the organization CSV")
    widths = [16, 24, 42]
    for col, width in enumerate(widths, start=1):
        ws.column_dimensions[get_column_letter(col)].width = width


def yes_no(value):
    return "כן" if value else "לא"


def parse_excel_date(workbook, value):
    if value in ("", None):
        return None
    if isinstance(value, (int, float)):
        try:
            return xlrd.xldate_as_datetime(float(value), workbook.datemode).date()
        except (ValueError, OverflowError):
            return None
    return None


def _detect_corrections_employee_meta(sheet):
    """Auto-detect employee name, ID, department, payroll number, tag from the sheet header section."""
    meta = {"employee_name": sheet.name, "department": "", "payroll_number": "", "id_number": "", "tag_number": ""}
    name_kws = ("שםעובד", "שםמלא", "שם")
    dept_kws = ("מחלקה", "מדור", "יחידה", "אגף")
    payroll_kws = ("מספרשכר", "מסשכר", "מספרעובד", "מסעובד", "מסעובד")
    id_kws = ("ת.ז", "תז", "תעודתזהות", "דרכון", "פספורט", "passport")
    tag_kws = ("מספרתג", "מסתג", "תג", "badge", "כרטיס")

    for r in range(min(12, sheet.nrows)):
        for c in range(sheet.ncols - 1):
            raw_label = str(sheet.cell_value(r, c)).strip()
            if not raw_label:
                continue
            tok = normalize_token(raw_label)
            next_val = ""
            for off in range(1, min(10, sheet.ncols - c)):
                v = str(sheet.cell_value(r, c + off)).strip()
                if v:
                    try:
                        fv = float(v)
                        if fv == int(fv):
                            v = str(int(fv))
                    except (ValueError, TypeError):
                        pass
                    next_val = v
                    break
            if not next_val:
                continue
            if (not meta["employee_name"] or meta["employee_name"] == sheet.name) and any(kw in tok for kw in name_kws):
                meta["employee_name"] = next_val
            elif not meta["department"] and any(kw in tok for kw in dept_kws):
                meta["department"] = next_val
            elif not meta["payroll_number"] and any(kw in tok for kw in payroll_kws):
                meta["payroll_number"] = next_val
            elif not meta["id_number"] and any(kw in tok for kw in id_kws):
                meta["id_number"] = next_val
            elif not meta["tag_number"] and any(kw in tok for kw in tag_kws):
                meta["tag_number"] = next_val
    return meta


def _detect_corrections_daily_structure(workbook, sheet):
    """Auto-detect header row, entry/exit columns, and date column for daily correction data."""
    entry_kws = ("כניסה",)
    exit_kws = ("יציאה",)

    best_row, best_score = 10, 0
    for r in range(min(18, sheet.nrows)):
        score = sum(
            1 for c in range(sheet.ncols)
            if any(k in normalize_token(str(sheet.cell_value(r, c))) for k in entry_kws + exit_kws)
        )
        if score > best_score:
            best_score, best_row = score, r

    data_start = best_row + 1
    header_tokens = [normalize_token(str(sheet.cell_value(best_row, c))) for c in range(sheet.ncols)]

    # Find columns that actually contain * corrections in data rows
    star_entry, star_exit = -1, -1
    for r in range(data_start, min(data_start + 80, sheet.nrows)):
        for c in range(sheet.ncols):
            if "*" in str(sheet.cell_value(r, c)):
                tok = header_tokens[c] if c < len(header_tokens) else ""
                if star_entry < 0 and any(k in tok for k in entry_kws):
                    star_entry = c
                elif star_exit < 0 and any(k in tok for k in exit_kws):
                    star_exit = c
        if star_entry >= 0 and star_exit >= 0:
            break

    # Fallback: use header column positions if no * found yet
    if star_entry < 0:
        cands = [c for c, t in enumerate(header_tokens) if any(k in t for k in entry_kws)]
        if cands:
            star_entry = cands[-1]
    if star_exit < 0:
        cands = [c for c, t in enumerate(header_tokens) if any(k in t for k in exit_kws)]
        if cands:
            star_exit = cands[-1]

    # Date column: first column with an Excel date serial in the first data row
    date_col = 0
    for c in range(sheet.ncols):
        v = sheet.cell_value(data_start, c) if data_start < sheet.nrows else None
        if isinstance(v, float) and 35000 < v < 65000:
            date_col = c
            break

    return {"header_row": best_row, "data_start": data_start, "entry_col": star_entry, "exit_col": star_exit, "date_col": date_col}


def parse_matan_manual_corrections(input_path, mapping=None):
    mapping = mapping or {}
    workbook = xlrd.open_workbook(input_path)
    employee_rows = []
    daily_rows = []

    for sheet in workbook.sheets():
        meta = _detect_corrections_employee_meta(sheet)
        struct = _detect_corrections_daily_structure(workbook, sheet)

        employee_name = meta["employee_name"]
        department = meta["department"]
        payroll_number = meta["payroll_number"]
        id_number = meta["id_number"]
        tag_number = meta["tag_number"]

        entry_col = struct["entry_col"]
        exit_col = struct["exit_col"]
        date_col = struct["date_col"]
        data_start = struct["data_start"]

        # Override with user-confirmed column mapping if provided
        if mapping.get("entry_col_source", "").startswith("col:"):
            try:
                entry_col = int(mapping["entry_col_source"].split(":", 1)[1])
            except ValueError:
                pass
        if mapping.get("exit_col_source", "").startswith("col:"):
            try:
                exit_col = int(mapping["exit_col_source"].split(":", 1)[1])
            except ValueError:
                pass
        if mapping.get("date_col_source", "").startswith("col:"):
            try:
                date_col = int(mapping["date_col_source"].split(":", 1)[1])
            except ValueError:
                pass

        raw_corrections = 0
        entry_corrections = 0
        exit_corrections = 0
        days_with_corrections = 0
        capped_corrections = 0
        work_days = 0
        month_days = 0
        detected_month = ""
        detected_year = 0

        for row_index in range(data_start, sheet.nrows):
            entry_value = str(get_sheet_cell(sheet, row_index, entry_col, "")).strip() if entry_col >= 0 else ""
            exit_value = str(get_sheet_cell(sheet, row_index, exit_col, "")).strip() if exit_col >= 0 else ""
            date_raw = get_sheet_cell(sheet, row_index, date_col, "")
            day_date = parse_excel_date(workbook, date_raw)

            has_time = bool(entry_value or exit_value)
            if not has_time and not day_date:
                continue

            if day_date and not month_days:
                month_days = calendar.monthrange(day_date.year, day_date.month)[1]
                detected_month = f"{day_date.month:02d}/{day_date.year}"
                detected_year = day_date.year
            if has_time:
                work_days += 1

            entry_corrected = "*" in entry_value
            exit_corrected = "*" in exit_value
            raw_daily = int(entry_corrected) + int(exit_corrected)
            capped_daily = min(raw_daily, 2)

            if raw_daily > 0:
                days_with_corrections += 1
                raw_corrections += raw_daily
                entry_corrections += int(entry_corrected)
                exit_corrections += int(exit_corrected)
                capped_corrections += capped_daily
                daily_rows.append({
                    "employee_name": employee_name,
                    "payroll_number": payroll_number,
                    "id_number": id_number,
                    "tag_number": tag_number,
                    "department": department,
                    "date": day_date.isoformat() if day_date else "",
                    "entry_value": entry_value,
                    "exit_value": exit_value,
                    "entry_corrected": entry_corrected,
                    "exit_corrected": exit_corrected,
                    "raw_daily_corrections": raw_daily,
                    "capped_daily_corrections": capped_daily,
                })

        employee_rows.append({
            "employee_name": employee_name,
            "payroll_number": payroll_number,
            "id_number": id_number,
            "tag_number": tag_number,
            "department": department,
            "detected_month": detected_month,
            "month_days": month_days,
            "work_days": work_days,
            "raw_correction_count": raw_corrections,
            "entry_correction_count": entry_corrections,
            "exit_correction_count": exit_corrections,
            "days_with_corrections": days_with_corrections,
            "capped_correction_count": capped_corrections,
            "avg_per_work_day": round(raw_corrections / work_days, 2) if work_days else 0.0,
        })

    return employee_rows, daily_rows


def apply_matan_manual_corrections_filters(rows, options):
    min_corrections = parse_float_or_none(options.get("min_corrections", ""))
    max_corrections = parse_float_or_none(options.get("max_corrections", ""))
    filtered = []
    for row in rows:
        correction_count = row["raw_correction_count"]
        if min_corrections is not None and correction_count < min_corrections:
            continue
        if max_corrections is not None and correction_count > max_corrections:
            continue
        filtered.append(row)
    return filtered


def subtract_months(base_date, months):
    if months <= 0:
        return base_date
    year = base_date.year
    month = base_date.month - months
    while month <= 0:
        month += 12
        year -= 1
    day = min(base_date.day, calendar.monthrange(year, month)[1])
    return date(year, month, day)


def looks_like_repeated_inactive_header(row_map):
    return normalize_token(row_map.get("employee_name_source", "")) == "שםעובד" or normalize_token(row_map.get("date_source", "")) == "תאריך"


def clean_daily_activity_value(value):
    text = stringify_excel_value(value)
    token = normalize_token(text)
    if token in {"כניסה", "יציאה", "אירוע", "סהכ", "שעות", "שםעובד", "תאריך"}:
        return ""
    return text


def parse_inactive_workers_report(input_path, extension, mapping, options):
    workbook_kind, workbook = open_excel_workbook(input_path, extension)
    sheet = iter_excel_sheets(workbook_kind, workbook)[0]
    rows, _cols = get_excel_dims(sheet, workbook_kind)
    header_row = detect_inactive_workers_header_row(sheet, workbook_kind)

    employees = {}
    all_dates = set()
    for row_index in range(header_row + 1, rows):
        date_value = parse_excel_date_generic(workbook_kind, workbook, extract_inactive_workers_mapping_value(sheet, workbook_kind, mapping.get("date_source"), row_index))
        row_map = {
            "employee_name_source": stringify_excel_value(extract_inactive_workers_mapping_value(sheet, workbook_kind, mapping.get("employee_name_source"), row_index)),
            "employee_number_source": stringify_excel_value(extract_inactive_workers_mapping_value(sheet, workbook_kind, mapping.get("employee_number_source"), row_index)),
            "badge_number_source": stringify_excel_value(extract_inactive_workers_mapping_value(sheet, workbook_kind, mapping.get("badge_number_source"), row_index)),
            "id_number_source": stringify_excel_value(extract_inactive_workers_mapping_value(sheet, workbook_kind, mapping.get("id_number_source"), row_index)),
            "passport_number_source": stringify_excel_value(extract_inactive_workers_mapping_value(sheet, workbook_kind, mapping.get("passport_number_source"), row_index)),
        }
        if looks_like_repeated_inactive_header({"employee_name_source": row_map["employee_name_source"], "date_source": stringify_excel_value(extract_inactive_workers_mapping_value(sheet, workbook_kind, mapping.get("date_source"), row_index))}):
            continue

        employee_name = row_map["employee_name_source"]
        employee_number = row_map["employee_number_source"]
        badge_number = row_map["badge_number_source"]
        id_number = row_map["id_number_source"]
        passport_number = row_map["passport_number_source"]
        if not any([employee_name, employee_number, badge_number, id_number, passport_number]):
            continue
        if not date_value:
            continue

        all_dates.add(date_value)
        entry_time = clean_daily_activity_value(extract_inactive_workers_mapping_value(sheet, workbook_kind, mapping.get("entry_time_source"), row_index))
        exit_time = clean_daily_activity_value(extract_inactive_workers_mapping_value(sheet, workbook_kind, mapping.get("exit_time_source"), row_index))
        total_hours_text = clean_daily_activity_value(extract_inactive_workers_mapping_value(sheet, workbook_kind, mapping.get("total_hours_source"), row_index))
        event_text = clean_daily_activity_value(extract_inactive_workers_mapping_value(sheet, workbook_kind, mapping.get("event_source"), row_index))
        department = stringify_excel_value(extract_inactive_workers_mapping_value(sheet, workbook_kind, mapping.get("department_source"), row_index))
        has_attendance_pair = bool(entry_time and exit_time)
        total_hours_value = parse_hours_or_zero(total_hours_text)
        has_total_hours = bool(total_hours_text and total_hours_value > 0)
        has_event = bool(event_text)
        has_activity = has_attendance_pair or has_total_hours or has_event

        employee_key = (
            employee_number,
            badge_number,
            id_number,
            passport_number,
            employee_name,
        )
        if employee_key not in employees:
            employees[employee_key] = {
                "employee_name": employee_name,
                "employee_number": employee_number,
                "badge_number": badge_number,
                "id_number": id_number,
                "passport_number": passport_number,
                "department": department,
                "last_active_date": None,
                "last_seen_date": date_value,
                "activity_days": 0,
                "event_only_days": 0,
                "system_notes": set(),
            }
        employee = employees[employee_key]
        employee["last_seen_date"] = max(employee["last_seen_date"], date_value)
        if department and not employee["department"]:
            employee["department"] = department
        if not employee_number and not badge_number and not id_number and not passport_number:
            employee["system_notes"].add("חסר מזהה עובד נוסף")
        if not employee_name:
            employee["system_notes"].add("חסר שם עובד")
        if has_activity:
            employee["activity_days"] += 1
            if has_event and not (has_attendance_pair or has_total_hours):
                employee["event_only_days"] += 1
            if employee["last_active_date"] is None or date_value > employee["last_active_date"]:
                employee["last_active_date"] = date_value

    reference_date = max(all_dates) if all_dates else date.today()
    unit = str(options.get("inactive_period_unit", "days") or "days").strip().lower()
    threshold_value = parse_int_or_none(options.get("inactive_period_value", "")) or 30
    threshold_value = max(1, threshold_value)
    cutoff_date = reference_date - timedelta(days=threshold_value) if unit == "days" else subtract_months(reference_date, threshold_value)

    inactive_rows = []
    for employee in employees.values():
        last_active_date = employee["last_active_date"]
        is_inactive = last_active_date is None or last_active_date < cutoff_date
        if not is_inactive:
            continue
        inactive_rows.append(
            {
                "employee_name": employee["employee_name"],
                "employee_number": employee["employee_number"],
                "badge_number": employee["badge_number"],
                "id_number": employee["id_number"],
                "passport_number": employee["passport_number"],
                "department": employee["department"],
                "last_active_date": last_active_date,
                "last_active_display": last_active_date.strftime("%d/%m/%Y") if last_active_date else "לא קיים מידע",
                "last_seen_date": employee["last_seen_date"],
                "activity_days": employee["activity_days"],
                "event_only_days": employee["event_only_days"],
                "status_reason": "לא זוהתה פעילות בכלל" if last_active_date is None else "לא זוהתה פעילות בטווח שנבדק",
                "system_notes": " | ".join(sorted(employee["system_notes"])),
            }
        )

    inactive_rows.sort(key=lambda row: (row["last_active_date"] is not None, row["last_active_date"] or date.min, row["employee_name"]))
    return inactive_rows, {
        "reference_date": reference_date,
        "cutoff_date": cutoff_date,
        "threshold_unit": unit,
        "threshold_value": threshold_value,
        "employee_count": len(employees),
        "inactive_count": len(inactive_rows),
        "span_days": ((reference_date - min(all_dates)).days + 1) if all_dates else 0,
    }


def write_inactive_workers_summary(ws, inactive_rows, meta, mapping):
    ws.title = safe_sheet_title("עובדים לא פעילים", "עובדים לא פעילים")
    ws.sheet_view.rightToLeft = True
    ws.sheet_view.showGridLines = False
    ws.freeze_panes = "A10"

    ws.merge_cells("A1:H1")
    ws["A1"] = "דוח עובדים לא פעילים"
    ws["A1"].font = Font(bold=True, size=16, color="0F172A")
    ws["A1"].fill = PatternFill(fill_type="solid", fgColor="BFD9FF")
    ws["A1"].alignment = Alignment(horizontal="center")
    ws.merge_cells("A2:H2")
    ws["A2"] = "איתור עובדים ללא פעילות בטווח שנבדק לפי דוח יומי"
    ws["A2"].font = Font(italic=True, size=10, color="475569")
    ws["A2"].alignment = Alignment(horizontal="center")

    unit_label = "ימים אחרונים" if meta.get("threshold_unit") == "days" else "חודשים אחרונים"
    metrics = [
        ("עובדים שנבדקו", meta.get("employee_count", 0), "DBEAFE"),
        ("עובדים לא פעילים", meta.get("inactive_count", 0), "FEE2E2"),
        ("תאריך ייחוס אחרון בקובץ", meta.get("reference_date").strftime("%d/%m/%Y") if meta.get("reference_date") else "—", "DCFCE7"),
        ("טווח בדיקה", f"{meta.get('threshold_value', 0)} {unit_label}", "FEF3C7"),
        ("תאריך חיתוך", meta.get("cutoff_date").strftime("%d/%m/%Y") if meta.get("cutoff_date") else "—", "E9D5FF"),
        ("מספר ימים שנכללו בקובץ", meta.get("span_days", 0), "FDE68A"),
    ]
    metric_blocks = [
        (4, "A4:B4", metrics[0]),
        (4, "D4:E4", metrics[1]),
        (4, "G4:H4", metrics[2]),
        (6, "A6:B6", metrics[3]),
        (6, "D6:E6", metrics[4]),
        (6, "G6:H6", metrics[5]),
    ]
    for start_row, cell_range, (label, value, fill_color) in metric_blocks:
        start_cell = cell_range.split(":")[0]
        ws.merge_cells(cell_range)
        ws[start_cell] = f"{label}\n{value}"
        ws[start_cell].font = Font(bold=True, color="0F172A", size=11)
        ws[start_cell].alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
        ws[start_cell].fill = PatternFill(fill_type="solid", fgColor=fill_color)
        ws.row_dimensions[start_row].height = 34

    header_row = 9
    headers = ["שם עובד"]
    if mapping.get("employee_number_source"):
        headers.append("מספר עובד")
    if mapping.get("badge_number_source"):
        headers.append("מספר תג")
    if mapping.get("id_number_source"):
        headers.append("תעודת זהות")
    if mapping.get("passport_number_source"):
        headers.append("דרכון")
    if mapping.get("department_source"):
        headers.append("מחלקה")
    headers.extend(["תאריך אחרון שזוהתה פעילות", "סיבת סימון", "ימי פעילות שזוהו"])
    if mapping.get("event_source"):
        headers.append("ימי אירוע ללא נוכחות")
    headers.append("הערות מערכת")

    for col, header in enumerate(headers, start=1):
        cell = ws.cell(row=header_row, column=col, value=header)
        cell.font = Font(bold=True, color="FFFFFF")
        cell.fill = PatternFill(fill_type="solid", fgColor="1E3A8A")
        cell.alignment = Alignment(horizontal="right", vertical="center")

    for row_idx, row in enumerate(inactive_rows, start=header_row + 1):
        values = [row["employee_name"]]
        if mapping.get("employee_number_source"):
            values.append(row["employee_number"])
        if mapping.get("badge_number_source"):
            values.append(row["badge_number"])
        if mapping.get("id_number_source"):
            values.append(row["id_number"])
        if mapping.get("passport_number_source"):
            values.append(row["passport_number"])
        if mapping.get("department_source"):
            values.append(row["department"])
        values.extend([row["last_active_display"], row["status_reason"], row["activity_days"]])
        if mapping.get("event_source"):
            values.append(row["event_only_days"])
        values.append(row.get("system_notes", ""))
        for col, value in enumerate(values, start=1):
            cell = ws.cell(row=row_idx, column=col, value=value)
            cell.alignment = Alignment(horizontal="right", vertical="top", wrap_text=True)
            if row_idx % 2 == 0:
                cell.fill = PatternFill(fill_type="solid", fgColor="F8FAFC")
        ws.row_dimensions[row_idx].height = 28

    widths = [26, 14, 14, 16, 16, 20, 22, 22, 14, 18, 28]
    for col, width in enumerate(widths[:len(headers)], start=1):
        ws.column_dimensions[get_column_letter(col)].width = width


def write_inactive_workers_by_department(ws, inactive_rows, mapping):
    ws.title = safe_sheet_title("לפי מחלקה", "לפי מחלקה")
    ws.sheet_view.rightToLeft = True
    ws.sheet_view.showGridLines = False

    ws.merge_cells("A1:F1")
    ws["A1"] = "עובדים לא פעילים לפי מחלקה"
    ws["A1"].font = Font(bold=True, size=17, color="0F172A")
    ws["A1"].fill = PatternFill(fill_type="solid", fgColor="C7D2FE")
    ws["A1"].alignment = Alignment(horizontal="center")

    grouped = {}
    for row in inactive_rows:
        department_name = row.get("department") or "ללא מחלקה"
        grouped.setdefault(department_name, []).append(row)

    current_row = 3
    for department_name in sorted(grouped):
        department_rows = grouped[department_name]
        ws.merge_cells(start_row=current_row, start_column=1, end_row=current_row, end_column=6)
        ws.cell(row=current_row, column=1, value=f"מחלקה: {department_name} | עובדים לא פעילים: {len(department_rows)}")
        ws.cell(row=current_row, column=1).font = Font(bold=True, color="0F172A")
        ws.cell(row=current_row, column=1).fill = PatternFill(fill_type="solid", fgColor="E0F2FE")
        ws.cell(row=current_row, column=1).alignment = Alignment(horizontal="right")
        current_row += 1

        headers = ["שם עובד", "מספר עובד", "מספר תג", "תאריך אחרון שזוהתה פעילות", "ימי פעילות שזוהו", "סיבת סימון"]
        for col, header in enumerate(headers, start=1):
            cell = ws.cell(row=current_row, column=col, value=header)
            cell.font = Font(bold=True, color="FFFFFF")
            cell.fill = PatternFill(fill_type="solid", fgColor="2563EB")
            cell.alignment = Alignment(horizontal="right")
        current_row += 1

        for row in department_rows:
            values = [
                row["employee_name"],
                row.get("employee_number", ""),
                row.get("badge_number", ""),
                row["last_active_display"],
                row["activity_days"],
                row["status_reason"],
            ]
            for col, value in enumerate(values, start=1):
                cell = ws.cell(row=current_row, column=col, value=value)
                cell.alignment = Alignment(horizontal="right", vertical="top", wrap_text=True)
                if current_row % 2 == 0:
                    cell.fill = PatternFill(fill_type="solid", fgColor="F8FAFC")
            current_row += 1
        current_row += 1

    for col, width in enumerate([26, 14, 14, 22, 16, 26], start=1):
        ws.column_dimensions[get_column_letter(col)].width = width


def write_matan_corrections_summary(ws, employee_rows, filters_used):
    ws.title = safe_sheet_title("סיכום תיקונים", "Corrections Summary")
    ws.sheet_view.rightToLeft = True
    ws.sheet_view.showGridLines = False
    ws.freeze_panes = "A12"

    ws["A1"] = "דוח תיקונים ידניים"
    ws["A1"].font = Font(bold=True, size=18, color="0F172A")
    ws["A1"].fill = PatternFill(fill_type="solid", fgColor="BFDBFE")

    # Detect month/days from the most common value across all employees
    month_label = ""
    month_days_global = 0
    if employee_rows:
        from collections import Counter
        month_counts = Counter(r.get("detected_month", "") for r in employee_rows if r.get("detected_month"))
        if month_counts:
            month_label = month_counts.most_common(1)[0][0]
        month_days_vals = [r.get("month_days", 0) for r in employee_rows if r.get("month_days")]
        if month_days_vals:
            month_days_global = max(month_days_vals)

    total_entry = sum(row["entry_correction_count"] for row in employee_rows)
    total_exit = sum(row["exit_correction_count"] for row in employee_rows)
    metrics = [
        ("עובדים בתוצאה", len(employee_rows), "DBEAFE"),
        ("חודש שזוהה", month_label or "לא זוהה", "E0F2FE"),
        ("ימי חודש", month_days_global or "לא זוהה", "E0F2FE"),
        ("סה\"כ תיקונים", sum(row["raw_correction_count"] for row in employee_rows), "FEE2E2"),
        ("תיקוני כניסה", total_entry, "FEF3C7"),
        ("תיקוני יציאה", total_exit, "FEF3C7"),
        ("סה\"כ ימים עם תיקונים", sum(row["days_with_corrections"] for row in employee_rows), "DCFCE7"),
    ]
    for idx, (label, value, fill_color) in enumerate(metrics, start=3):
        label_cell = ws.cell(row=idx, column=1, value=label)
        value_cell = ws.cell(row=idx, column=2, value=value)
        label_cell.font = Font(bold=True, color="334155")
        value_cell.font = Font(bold=True, color="0F172A")
        label_cell.fill = PatternFill(fill_type="solid", fgColor=fill_color)
        value_cell.fill = PatternFill(fill_type="solid", fgColor=fill_color)

    ws["D3"] = "סינונים"
    ws["D3"].font = Font(bold=True)
    for idx, (label, value) in enumerate(filters_used.items(), start=4):
        ws.cell(row=idx, column=4, value=label).font = Font(bold=True)
        ws.cell(row=idx, column=5, value=value or "ללא")

    header_row = 11
    headers = [
        "שם עובד",
        "מספר שכר",
        "תעודת זהות",
        "מספר תג",
        "מחלקה",
        "סה\"כ תיקונים",
        "תיקוני כניסה",
        "תיקוני יציאה",
        "ימים עם תיקונים",
        "ימי חודש",
        "ימי עבודה",
        "ממוצע תיקונים ליום עבודה",
    ]
    for col, header in enumerate(headers, start=1):
        cell = ws.cell(row=header_row, column=col, value=header)
        cell.font = Font(bold=True, color="FFFFFF")
        cell.fill = PatternFill(fill_type="solid", fgColor="1E3A8A")

    sorted_rows = sorted(employee_rows, key=lambda row: (-row["raw_correction_count"], row["employee_name"]))
    for row_idx, row in enumerate(sorted_rows, start=header_row + 1):
        values = [
            row["employee_name"],
            row["payroll_number"],
            row["id_number"],
            row.get("tag_number", ""),
            row["department"],
            row["raw_correction_count"],
            row["entry_correction_count"],
            row["exit_correction_count"],
            row["days_with_corrections"],
            row.get("month_days", ""),
            row.get("work_days", ""),
            row.get("avg_per_work_day", 0.0),
        ]
        for col, value in enumerate(values, start=1):
            ws.cell(row=row_idx, column=col, value=value)
            if row_idx % 2 == 0:
                ws.cell(row=row_idx, column=col).fill = PatternFill(fill_type="solid", fgColor="F8FAFC")
        ws.cell(row=row_idx, column=6).fill = PatternFill(fill_type="solid", fgColor="FEE2E2")
        ws.cell(row=row_idx, column=7).fill = PatternFill(fill_type="solid", fgColor="FEF3C7")
        ws.cell(row=row_idx, column=8).fill = PatternFill(fill_type="solid", fgColor="FEF3C7")
        ws.cell(row=row_idx, column=6).font = Font(bold=True, color="991B1B")
        ws.cell(row=row_idx, column=12).number_format = "0.00"

    widths = [24, 16, 16, 14, 22, 16, 16, 16, 18, 14, 14, 26]
    for col, width in enumerate(widths, start=1):
        ws.column_dimensions[get_column_letter(col)].width = width


def write_matan_corrections_by_department(ws, employee_rows):
    ws.title = safe_sheet_title("לפי מחלקה", "By Department")
    ws.sheet_view.rightToLeft = True
    ws.sheet_view.showGridLines = False
    ws.freeze_panes = "A2"

    dept_data = defaultdict(lambda: {"total": 0, "entry": 0, "exit": 0, "count": 0})
    for row in employee_rows:
        dept = row["department"] or "ללא מחלקה"
        dept_data[dept]["total"] += row["raw_correction_count"]
        dept_data[dept]["entry"] += row["entry_correction_count"]
        dept_data[dept]["exit"] += row["exit_correction_count"]
        dept_data[dept]["count"] += 1

    headers = ["מחלקה", "מספר עובדים", "סה\"כ תיקונים", "תיקוני כניסה", "תיקוני יציאה", "ממוצע תיקונים לעובד"]
    for col, header in enumerate(headers, start=1):
        cell = ws.cell(row=1, column=col, value=header)
        cell.font = Font(bold=True, color="FFFFFF")
        cell.fill = PatternFill(fill_type="solid", fgColor="7C3AED")

    sorted_depts = sorted(dept_data.items(), key=lambda x: -x[1]["total"])
    for row_idx, (dept_name, data) in enumerate(sorted_depts, start=2):
        count = data["count"]
        avg = round(data["total"] / count, 2) if count else 0.0
        values = [dept_name, count, data["total"], data["entry"], data["exit"], avg]
        for col, value in enumerate(values, start=1):
            cell = ws.cell(row=row_idx, column=col, value=value)
            if row_idx % 2 == 0:
                cell.fill = PatternFill(fill_type="solid", fgColor="F5F3FF")
        ws.cell(row=row_idx, column=3).font = Font(bold=True, color="5B21B6")
        ws.cell(row=row_idx, column=6).number_format = "0.00"

    widths = [28, 16, 18, 18, 18, 22]
    for col, width in enumerate(widths, start=1):
        ws.column_dimensions[get_column_letter(col)].width = width


def write_matan_corrections_daily(ws, daily_rows, allowed_names):
    ws.title = safe_sheet_title("פירוט יומי", "Daily Corrections")
    ws.sheet_view.rightToLeft = True
    ws.sheet_view.showGridLines = False
    ws.freeze_panes = "A2"

    headers = [
        "עובד",
        "מחלקה",
        "תאריך",
        "שעת כניסה",
        "שעת יציאה",
        "כניסה תוקנה",
        "יציאה תוקנה",
        "תיקונים יומיים",
    ]
    for col, header in enumerate(headers, start=1):
        cell = ws.cell(row=1, column=col, value=header)
        cell.font = Font(bold=True, color="FFFFFF")
        cell.fill = PatternFill(fill_type="solid", fgColor="0F766E")

    row_idx = 2
    for row in daily_rows:
        if row["employee_name"] not in allowed_names:
            continue
        values = [
            row["employee_name"],
            row.get("department", ""),
            row["date"],
            row.get("entry_value", ""),
            row.get("exit_value", ""),
            yes_no(row["entry_corrected"]),
            yes_no(row["exit_corrected"]),
            row["raw_daily_corrections"],
        ]
        for col, value in enumerate(values, start=1):
            ws.cell(row=row_idx, column=col, value=value)
            if row_idx % 2 == 0:
                ws.cell(row=row_idx, column=col).fill = PatternFill(fill_type="solid", fgColor="ECFDF5")
        ws.cell(row=row_idx, column=4).fill = PatternFill(fill_type="solid", fgColor="FEF2F2" if row["entry_corrected"] else "F8FAFC")
        ws.cell(row=row_idx, column=5).fill = PatternFill(fill_type="solid", fgColor="FEF2F2" if row["exit_corrected"] else "F8FAFC")
        row_idx += 1

    widths = [24, 22, 14, 16, 16, 16, 16, 16]
    for col, width in enumerate(widths, start=1):
        ws.column_dimensions[get_column_letter(col)].width = width


def run_matan_manual_corrections(input_path, output_path, extension, options=None):
    if extension != "xls":
        raise ValueError("Matan manual-corrections tool currently supports XLS export only")
    options = options or {}
    mapping = {k: v for k, v in options.items() if k.endswith("_source")}
    employee_rows, daily_rows = parse_matan_manual_corrections(input_path, mapping)
    filtered_rows = apply_matan_manual_corrections_filters(employee_rows, options)
    allowed_names = {row["employee_name"] for row in filtered_rows}
    wb = Workbook()
    write_matan_corrections_summary(
        wb.active,
        filtered_rows,
        {
            "מינימום תיקונים": options.get("min_corrections", ""),
            "מקסימום תיקונים": options.get("max_corrections", ""),
        },
    )
    write_matan_corrections_by_department(wb.create_sheet(), filtered_rows)
    write_matan_corrections_daily(wb.create_sheet(), daily_rows, allowed_names)
    wb.save(output_path)


def has_rimon_work_activity(entry_value, exit_value, total_hours):
    return any([str(entry_value).strip(), str(exit_value).strip(), parse_hours_value(total_hours)])


def is_rimon_error_text(error_text):
    text = str(error_text).strip()
    if not text:
        return False
    return text not in {"יום חסר"}


def extract_rimon_mapping_value(sheet, workbook_kind, mapping_value, row_index=None):
    source = str(mapping_value or "").strip()
    if not source:
        return ""
    if source == "meta:employee_name":
        return find_rimon_meta_value(sheet, workbook_kind, ["שם לתצוגה", "שם עובד"])
    if source == "meta:payroll_number":
        return find_rimon_meta_value(sheet, workbook_kind, ["מספר שכר", "מספר עובד", "מספר בשכר", "תג עובד", "שכר"])
    if source == "meta:department":
        return find_rimon_meta_value(sheet, workbook_kind, ["מחלקה"])
    if source == "meta:id_number":
        return find_rimon_meta_value(sheet, workbook_kind, ["תעודת זהות", "דרכון"])
    if source.startswith("meta_cell:"):
        try:
            _, row_text, col_text = source.split(":", 2)
            meta_row = int(row_text)
            meta_col = int(col_text)
        except ValueError:
            return ""
        return stringify_excel_value(get_excel_cell(sheet, workbook_kind, meta_row, meta_col, ""))
    if source.startswith("col:") and row_index is not None:
        try:
            col_index = int(source.split(":", 1)[1])
        except ValueError:
            return ""
        return get_excel_cell(sheet, workbook_kind, row_index, col_index, "")
    return ""


def get_flamingo_sheet_name(sheet, workbook_kind):
    return sheet.title if workbook_kind == "xlsx" else sheet.name


def get_flamingo_sheet_dims(sheet, workbook_kind):
    return get_excel_dims(sheet, workbook_kind)


def get_flamingo_sheet_cell(sheet, workbook_kind, row_index, col_index, default=""):
    return get_excel_cell(sheet, workbook_kind, row_index, col_index, default)


def find_sheet_label_row(sheet, workbook_kind, label_text):
    normalized_label = normalize_token(label_text)
    rows, cols = get_flamingo_sheet_dims(sheet, workbook_kind)
    for row_index in range(rows):
        for col_index in range(cols):
            token = normalize_token(get_flamingo_sheet_cell(sheet, workbook_kind, row_index, col_index))
            if token == normalized_label:
                return row_index
    return -1


def find_sheet_label_position(sheet, workbook_kind, label_text):
    """Find the (row, col) position of a label in the sheet. Returns (-1, -1) if not found."""
    normalized_label = normalize_token(label_text)
    rows, cols = get_flamingo_sheet_dims(sheet, workbook_kind)
    for row_index in range(rows):
        for col_index in range(cols):
            token = normalize_token(get_flamingo_sheet_cell(sheet, workbook_kind, row_index, col_index))
            if token == normalized_label:
                return row_index, col_index
    return -1, -1


def sheet_has_label(sheet, workbook_kind, label_text):
    normalized_label = normalize_token(label_text)
    rows, cols = get_flamingo_sheet_dims(sheet, workbook_kind)
    for row_index in range(rows):
        for col_index in range(cols):
            token = normalize_token(get_flamingo_sheet_cell(sheet, workbook_kind, row_index, col_index))
            if token == normalized_label:
                return True
    return False


def flamingo_sheet_has_daily(sheet, workbook_kind):
    return sheet_has_label(sheet, workbook_kind, "תאריך") and sheet_has_label(sheet, workbook_kind, "כניסה")


def flamingo_sheet_has_summary(sheet, workbook_kind):
    return (
        find_sheet_label_row(sheet, workbook_kind, "נתונים כללים") >= 0
        or (
            not flamingo_sheet_has_daily(sheet, workbook_kind)
            and (
                sheet_has_label(sheet, workbook_kind, "נוכחות")
                or sheet_has_label(sheet, workbook_kind, "שעות לתשלום")
                or sheet_has_label(sheet, workbook_kind, "שעות משולמות")
            )
        )
    )


def iter_flamingo_worker_blocks(workbook_kind, workbook):
    sheets = iter_excel_sheets(workbook_kind, workbook)
    index = 0
    while index < len(sheets):
        current_sheet = sheets[index]
        if not flamingo_sheet_has_daily(current_sheet, workbook_kind):
            index += 1
            continue
        summary_sheet = None
        if flamingo_sheet_has_summary(current_sheet, workbook_kind):
            summary_sheet = current_sheet
            index += 1
        elif index + 1 < len(sheets):
            next_sheet = sheets[index + 1]
            if flamingo_sheet_has_summary(next_sheet, workbook_kind) and not flamingo_sheet_has_daily(next_sheet, workbook_kind):
                summary_sheet = next_sheet
                index += 2
            else:
                index += 1
        else:
            index += 1
        yield current_sheet, summary_sheet


def find_first_non_empty_in_row(sheet, row_index, start_col=0):
    if row_index >= sheet.nrows:
        return ""
    for col_index in range(start_col, sheet.ncols):
        value = sheet.cell_value(row_index, col_index)
        if value not in ("", None):
            return value
    return ""


def parse_flamingo_source(source):
    text = str(source or "").strip()
    if not text or ":" not in text:
        return "", ""
    source_type, source_label = text.split(":", 1)
    return source_type, source_label


def find_value_by_label_nearby(sheet, workbook_kind, label_text, max_col_distance=8, min_row=0, max_row=None):
    normalized_label = normalize_token(label_text)
    rows, cols = get_flamingo_sheet_dims(sheet, workbook_kind)
    last_row = rows if max_row is None else min(max_row, rows)
    for row_index in range(min_row, last_row):
        for col_index in range(cols):
            token = normalize_token(get_flamingo_sheet_cell(sheet, workbook_kind, row_index, col_index))
            if token != normalized_label:
                continue
            for next_col in range(col_index + 1, min(cols, col_index + max_col_distance + 1)):
                candidate = get_flamingo_sheet_cell(sheet, workbook_kind, row_index, next_col)
                if candidate not in ("", None):
                    candidate_token = normalize_token(candidate)
                    if candidate_token in FLAMINGO_META_LABEL_TOKENS or candidate_token in {"תאריך", "יום", "כניסה", "יציאה", "אירוע", "סהכ", "סה\"כ", "תקן", "חוסר"}:
                        continue
                    return candidate
    return ""


def find_flamingo_summary_value_by_label(detail_sheet, summary_sheet, workbook_kind, label_text):
    normalized_label = normalize_token(label_text)
    search_sheets = []
    if summary_sheet is not None:
        search_sheets.append(summary_sheet)
    if summary_sheet is None or summary_sheet is detail_sheet:
        search_sheets.append(detail_sheet)
    for sheet in search_sheets:
        rows, cols = get_flamingo_sheet_dims(sheet, workbook_kind)
        summary_start_row = find_sheet_label_row(sheet, workbook_kind, "נתונים כללים")
        start_row = summary_start_row if summary_start_row >= 0 else 0
        for row_index in range(start_row, rows):
            row_values = [get_flamingo_sheet_cell(sheet, workbook_kind, row_index, c) for c in range(cols)]
            for col_index, raw in enumerate(row_values):
                if normalize_token(raw) != normalized_label:
                    continue
                for next_col in range(col_index + 1, len(row_values)):
                    candidate = row_values[next_col]
                    parsed_hours = try_parse_hours_value(candidate)
                    if parsed_hours is not None:
                        return candidate
                    if isinstance(candidate, (int, float)) and candidate not in (0, 0.0):
                        return candidate
                    if str(candidate).strip():
                        try:
                            float(str(candidate).strip().replace(",", "."))
                            return candidate
                        except ValueError:
                            continue
    return ""


def extract_flamingo_mapping_value(detail_sheet, summary_sheet, workbook_kind, source):
    source_type, source_label = parse_flamingo_source(source)
    if source_type == "meta":
        return find_value_by_label_nearby(detail_sheet, workbook_kind, source_label, max_col_distance=10, min_row=0, max_row=18)
    if source_type == "summary":
        return find_flamingo_summary_value_by_label(detail_sheet, summary_sheet, workbook_kind, source_label)
    return ""


def default_rimon_mapping():
    return {
        "employee_name_source": "meta:employee_name",
        "payroll_number_source": "meta:payroll_number",
        "date_source": "col:0",
        "day_name_source": "col:6",
        "entry_time_source": "col:8",
        "exit_time_source": "col:12",
        "total_hours_source": "col:20",
        "standard_hours_source": "col:25",
        "missing_hours_source": "col:30",
        "event_source": "col:17",
        "error_text_source": "col:51",
        "department_source": "meta:department",
        "id_number_source": "meta:id_number",
    }


def default_attendance_alerts_mapping():
    m = default_rimon_mapping()
    m.update({
        "break_hours_source": "col:28",
        "regular_hours_source": "col:34",
        "overtime_100_source": "col:38",
        "overtime_125_source": "col:41",
        "overtime_150_source": "col:46",
        "overtime_200_source": "col:54",
    })
    return m


def build_rimon_mapping_warnings(mapping):
    warnings = []
    if not mapping.get("error_text_source"):
        warnings.append("לא נקלט שדה שגיאות. הדוח יורץ ללא זיהוי ימי שגיאה.")
    if not mapping.get("event_source"):
        warnings.append("לא נקלט שדה אירוע. הדוח יורץ, אך לקבלת תוצאה מדויקת יש לייצר מחדש עם שדה האירוע.")
    if not mapping.get("date_source"):
        warnings.append("לא נקלט שדה תאריך. הדוח יורץ, אך לקבלת דוח תקין ומדויק יש לייצר מחדש עם שדה התאריך.")
    if not mapping.get("payroll_number_source"):
        warnings.append("לא נקלט שדה מספר עובד. הדוח יורץ, אך לקבלת דוח תקין ומדויק יש לייצר מחדש עם שדה מספר העובד.")
    if not mapping.get("employee_name_source"):
        warnings.append("לא נקלט שדה שם עובד. הדוח יורץ, אך יהיה פחות ברור לקריאה.")
    return warnings


def parse_rimon_home_office_report(input_path, extension, mapping):
    workbook_kind, workbook = open_excel_workbook(input_path, extension)
    employee_rows = []
    daily_rows = []
    detected_company_name = ""
    detected_months = []
    detected_dates = set()

    for sheet in iter_excel_sheets(workbook_kind, workbook):
        rows, _ = get_excel_dims(sheet, workbook_kind)
        header_row = detect_rimon_header_row(sheet, workbook_kind)
        employee_name = extract_rimon_mapping_value(sheet, workbook_kind, mapping.get("employee_name_source"))
        department = extract_rimon_mapping_value(sheet, workbook_kind, mapping.get("department_source"))
        payroll_number = extract_rimon_mapping_value(sheet, workbook_kind, mapping.get("payroll_number_source"))
        id_number = extract_rimon_mapping_value(sheet, workbook_kind, mapping.get("id_number_source"))
        if not employee_name:
            employee_name = getattr(sheet, "name", "עובד")
        if not stringify_excel_value(extract_rimon_mapping_value(sheet, workbook_kind, mapping.get("employee_name_source"))) and not stringify_excel_value(extract_rimon_mapping_value(sheet, workbook_kind, mapping.get("payroll_number_source"))):
            continue
        if not detected_company_name:
            detected_company_name = stringify_excel_value(get_excel_cell(sheet, workbook_kind, 0, 0, "")) or stringify_excel_value(get_excel_cell(sheet, workbook_kind, 1, 42, ""))

        grouped_dates = {}
        current_date = None

        for row_index in range(header_row + 1, rows):
            row_date = parse_excel_date_generic(
                workbook_kind,
                workbook,
                extract_rimon_mapping_value(sheet, workbook_kind, mapping.get("date_source"), row_index),
            )
            if row_date:
                current_date = row_date
            if current_date is None:
                continue

            event_value = extract_rimon_mapping_value(sheet, workbook_kind, mapping.get("event_source"), row_index)
            error_text = extract_rimon_mapping_value(sheet, workbook_kind, mapping.get("error_text_source"), row_index)
            day_name = stringify_excel_value(extract_rimon_mapping_value(sheet, workbook_kind, mapping.get("day_name_source"), row_index))
            entry_time = stringify_excel_value(extract_rimon_mapping_value(sheet, workbook_kind, mapping.get("entry_time_source"), row_index))
            exit_time = stringify_excel_value(extract_rimon_mapping_value(sheet, workbook_kind, mapping.get("exit_time_source"), row_index))
            total_hours = stringify_excel_value(extract_rimon_mapping_value(sheet, workbook_kind, mapping.get("total_hours_source"), row_index))
            standard_hours = stringify_excel_value(extract_rimon_mapping_value(sheet, workbook_kind, mapping.get("standard_hours_source"), row_index))
            missing_hours = stringify_excel_value(extract_rimon_mapping_value(sheet, workbook_kind, mapping.get("missing_hours_source"), row_index))
            if not any([row_date, event_value, error_text, entry_time, exit_time, total_hours, standard_hours, missing_hours]):
                continue

            day_key = current_date.isoformat()
            if day_key not in grouped_dates:
                grouped_dates[day_key] = {
                    "employee_name": employee_name,
                    "payroll_number": payroll_number,
                    "id_number": id_number,
                    "department": department,
                    "date": day_key,
                    "home_office": False,
                    "missing_absence": False,
                    "error": False,
                    "day_name": day_name,
                    "entry_time": "",
                    "exit_time": "",
                    "total_hours": "",
                    "standard_hours": "",
                    "missing_hours": "",
                    "events": [],
                    "errors": [],
                }

            grouped = grouped_dates[day_key]
            normalized_event = str(event_value or "").strip()
            grouped["home_office"] = grouped["home_office"] or normalized_event == "עבודה מהבית"
            if day_name and not grouped["day_name"]:
                grouped["day_name"] = day_name
            if entry_time and not grouped["entry_time"]:
                grouped["entry_time"] = entry_time
            if exit_time and not grouped["exit_time"]:
                grouped["exit_time"] = exit_time
            if total_hours and not grouped["total_hours"]:
                grouped["total_hours"] = total_hours
            if standard_hours and not grouped["standard_hours"]:
                grouped["standard_hours"] = standard_hours
            if missing_hours and not grouped["missing_hours"]:
                grouped["missing_hours"] = missing_hours
            if event_value and event_value not in grouped["events"]:
                grouped["events"].append(event_value)
            if error_text and error_text not in grouped["errors"]:
                grouped["errors"].append(error_text)
            grouped["error"] = grouped["error"] or bool(error_text)

        office_days = 0
        home_office_days = 0
        missing_absence_days = 0
        left_days = 0
        error_days = 0
        standard_hours_total = 0.0
        missing_hours_total = 0.0

        for day_key in sorted(grouped_dates):
            grouped = grouped_dates[day_key]
            has_home_event = grouped["home_office"]
            normalized_events = [str(event).strip() for event in grouped["events"] if str(event).strip()]
            has_leave_event = any("פיטור" in event for event in normalized_events)
            has_absence_event = any(
                keyword in event
                for event in normalized_events
                for keyword in ("חופשה", "מחלה", "היעדר", "חג", "אבל", "מילואים")
            )
            has_other_work_event = any(
                event != "עבודה מהבית"
                and "פיטור" not in event
                and not any(keyword in event for keyword in ("חופשה", "מחלה", "היעדר", "חג", "אבל", "מילואים"))
                for event in normalized_events
            )
            has_entry = bool(grouped["entry_time"])
            has_exit = bool(grouped["exit_time"])
            has_complete_attendance = has_entry and has_exit
            has_partial_attendance = (has_entry and not has_exit) or (has_exit and not has_entry)
            office_work = has_other_work_event or (has_complete_attendance and not has_home_event)
            left_employee = has_leave_event and not office_work and not has_home_event and not has_complete_attendance
            grouped["missing_absence"] = has_absence_event and not left_employee and not office_work and not has_home_event and not has_complete_attendance
            missing_day_error = not has_complete_attendance and not has_partial_attendance and not normalized_events and not grouped["missing_absence"] and not left_employee
            if has_partial_attendance and "חסר דיווח" not in grouped["errors"]:
                grouped["errors"].append("חסר דיווח")
            if missing_day_error and "יום חסר" not in grouped["errors"]:
                grouped["errors"].append("יום חסר")
            grouped["error"] = grouped["error"] or missing_day_error or has_partial_attendance

            if has_home_event and not office_work:
                home_office_days += 1
            if office_work:
                office_days += 1
            if grouped["missing_absence"]:
                missing_absence_days += 1
            if left_employee:
                left_days += 1
            if grouped["error"]:
                error_days += 1
            standard_hours_total += parse_hours_or_zero(grouped["standard_hours"])
            missing_hours_total += parse_hours_or_zero(grouped["missing_hours"])
            detected_months.append(datetime.fromisoformat(grouped["date"]).month)
            detected_dates.add(grouped["date"])

            daily_rows.append(
                {
                    "employee_name": grouped["employee_name"],
                    "date": grouped["date"],
                    "day_name": grouped["day_name"],
                    "entry_time": grouped["entry_time"],
                    "exit_time": grouped["exit_time"],
                    "home_office": grouped["home_office"],
                    "office_work": office_work,
                    "missing_absence": grouped["missing_absence"],
                    "error": grouped["error"],
                    "event": " | ".join(grouped["events"]),
                    "total_hours": grouped["total_hours"],
                    "standard_hours": grouped["standard_hours"],
                    "missing_hours": grouped["missing_hours"],
                    "error_text": " | ".join(grouped["errors"]),
                }
            )

        employee_rows.append(
            {
                "employee_name": employee_name,
                "payroll_number": payroll_number,
                "id_number": id_number,
                "department": department,
                "office_work_days": office_days,
                "home_office_days": home_office_days,
                "missing_absence_days": missing_absence_days,
                "left_days": left_days,
                "error_days": error_days,
                "total_grouped_dates": len(grouped_dates),
                "standard_hours_total": standard_hours_total,
                "missing_hours_total": missing_hours_total,
            }
        )

    if workbook_kind == "xlsx":
        workbook.close()
    report_month = ""
    if detected_months:
        month_names_he = {
            1: "ינואר",
            2: "פברואר",
            3: "מרץ",
            4: "אפריל",
            5: "מאי",
            6: "יוני",
            7: "יולי",
            8: "אוגוסט",
            9: "ספטמבר",
            10: "אוקטובר",
            11: "נובמבר",
            12: "דצמבר",
        }
        month_number = max(set(detected_months), key=detected_months.count)
        report_month = month_names_he.get(month_number, "")
    report_meta = {
        "company_name": detected_company_name,
        "report_month": report_month,
        "identified_day_count": len(detected_dates),
    }
    return employee_rows, daily_rows, report_meta


def write_rimon_home_office_summary(ws, employee_rows, report_meta):
    ws.title = safe_sheet_title("סיכום כולל", "Overall Summary")
    ws.sheet_view.rightToLeft = True
    ws.sheet_view.showGridLines = False

    headers = [
        "שם עובד",
        "מספר שכר",
        "תעודת זהות",
        "מחלקה",
        "ימי עבודה מהמשרד",
        "ימי עבודה מהבית",
        "ימי היעדרות",
        "ימי שגיאה",
        "ימי עזיבה",
        "סה\"כ ימי עבודה שזוהו",
        "סה\"כ שעות תקן",
        "סה\"כ שעות חוסר",
    ]

    ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=len(headers))
    ws["A1"] = "דוח סיכום עבודה מהבית מול עבודה מהמשרד"
    ws["A1"].font = Font(bold=True, size=18, color="0F172A")
    ws["A1"].fill = PatternFill(fill_type="solid", fgColor="BFDBFE")
    ws["A1"].alignment = Alignment(horizontal="center")

    company_name = report_meta.get("company_name", "")
    if company_name:
        ws.merge_cells(start_row=2, start_column=1, end_row=2, end_column=len(headers))
        ws["A2"] = company_name
        ws["A2"].font = Font(bold=True, size=13, color="334155")
        ws["A2"].alignment = Alignment(horizontal="center")

    total_days_identified = report_meta.get("identified_day_count", 0)
    total_home_days = sum(row["home_office_days"] for row in employee_rows)
    total_office_days = sum(row["office_work_days"] for row in employee_rows)
    total_absence_days = sum(row["missing_absence_days"] for row in employee_rows)
    total_error_days = sum(row["error_days"] for row in employee_rows)
    total_standard_hours = sum(row.get("standard_hours_total", 0.0) or 0.0 for row in employee_rows)
    total_missing_hours = sum(row.get("missing_hours_total", 0.0) or 0.0 for row in employee_rows)
    metrics = [
        ("סה\"כ עובדים שנקלטו", len(employee_rows), "E0F2FE"),
        ("סה\"כ ימים לחודש שזוהו", total_days_identified, "E0F2FE"),
        ("חודש הדוח", report_meta.get("report_month", ""), "E0F2FE"),
        ("סה\"כ ימי עבודה", total_office_days + total_home_days, "DBEAFE"),
        ("סה\"כ ימי עבודה מהבית", total_home_days, "DDD6FE"),
        ("סה\"כ ימי עבודה מהמשרד", total_office_days, "DCFCE7"),
        ("סה\"כ שגיאות", total_error_days, "FEE2E2"),
        ("סה\"כ היעדרויות", total_absence_days, "FEF3C7"),
        ("סה\"כ שעות תקן", format_hours(total_standard_hours), "E0F2FE"),
        ("סה\"כ שעות חוסר", format_hours(total_missing_hours), "FEF3C7"),
    ]
    for idx, (label, value, fill_color) in enumerate(metrics, start=4):
        label_cell = ws.cell(row=idx, column=1, value=label)
        value_cell = ws.cell(row=idx, column=2, value=value)
        label_cell.font = Font(bold=True, color="334155")
        value_cell.font = Font(bold=True, color="0F172A")
        label_cell.fill = PatternFill(fill_type="solid", fgColor=fill_color)
        value_cell.fill = PatternFill(fill_type="solid", fgColor=fill_color)
        label_cell.alignment = Alignment(horizontal="right")
        value_cell.alignment = Alignment(horizontal="right")

    header_row = len(metrics) + 5
    ws.freeze_panes = "A" + str(header_row + 1)
    for col, header in enumerate(headers, start=1):
        cell = ws.cell(row=header_row, column=col, value=header)
        cell.font = Font(bold=True, color="FFFFFF")
        cell.fill = PatternFill(fill_type="solid", fgColor="1E3A8A")
        cell.alignment = Alignment(horizontal="right")

    sorted_rows = sorted(employee_rows, key=lambda row: (row["employee_name"], row["payroll_number"]))
    for row_idx, row in enumerate(sorted_rows, start=header_row + 1):
        values = [
            row["employee_name"],
            row["payroll_number"],
            row["id_number"],
            row["department"],
            row["office_work_days"],
            row["home_office_days"],
            row["missing_absence_days"],
            row["error_days"],
            row.get("left_days", 0),
            row["office_work_days"] + row["home_office_days"],
            format_hours(row.get("standard_hours_total", 0.0)),
            format_hours(row.get("missing_hours_total", 0.0)),
        ]
        for col, value in enumerate(values, start=1):
            cell = ws.cell(row=row_idx, column=col, value=value)
            cell.alignment = Alignment(horizontal="right")
            if row_idx % 2 == 0:
                cell.fill = PatternFill(fill_type="solid", fgColor="F8FAFC")

    widths = [24, 16, 16, 24, 18, 18, 18, 14, 14, 28, 16, 16]
    for col, width in enumerate(widths, start=1):
        ws.column_dimensions[get_column_letter(col)].width = width


def write_rimon_home_office_daily(ws, daily_rows):
    ws.title = safe_sheet_title("פירוט יומי", "Daily Breakdown")
    ws.sheet_view.rightToLeft = True
    ws.sheet_view.showGridLines = False
    ws.freeze_panes = "A2"

    headers = [
        "עובד",
        "תאריך",
        "סוג יום",
        "שעת כניסה",
        "שעת יציאה",
        "סה\"כ שעות",
        "עבודה מהבית",
        "עבודה מהמשרד",
        "היעדרות",
        "שגיאה",
        "אירוע",
        "שעות תקן",
        "שעות חוסר",
        "פירוט שגיאה",
    ]
    for col, header in enumerate(headers, start=1):
        cell = ws.cell(row=1, column=col, value=header)
        cell.font = Font(bold=True, color="FFFFFF")
        cell.fill = PatternFill(fill_type="solid", fgColor="0F766E")
        cell.alignment = Alignment(horizontal="right")

    sorted_rows = sorted(daily_rows, key=lambda row: (row["employee_name"], row["date"]))
    for row_idx, row in enumerate(sorted_rows, start=2):
        values = [
            row["employee_name"],
            row["date"],
            row["day_name"],
            row["entry_time"],
            row["exit_time"],
            row["total_hours"],
            yes_no(row["home_office"]),
            yes_no(row["office_work"]),
            yes_no(row["missing_absence"]),
            yes_no(row["error"]),
            row["event"],
            row["standard_hours"],
            row["missing_hours"],
            row["error_text"],
        ]
        for col, value in enumerate(values, start=1):
            cell = ws.cell(row=row_idx, column=col, value=value)
            cell.alignment = Alignment(horizontal="right")
            if row_idx % 2 == 0:
                cell.fill = PatternFill(fill_type="solid", fgColor="ECFDF5")

    widths = [24, 14, 12, 12, 12, 12, 14, 14, 12, 10, 20, 14, 14, 22]
    for col, width in enumerate(widths, start=1):
        ws.column_dimensions[get_column_letter(col)].width = width


def run_rimon_home_office_summary(input_path, output_path, extension, options=None):
    if extension not in {"xls", "xlsx"}:
        raise ValueError("Rimon home-office summary currently supports XLS and XLSX uploads only")
    options = options or {}
    mapping = default_rimon_mapping()
    mapping.update({key: value for key, value in options.items() if key.endswith("_source")})
    employee_rows, daily_rows, report_meta = parse_rimon_home_office_report(input_path, extension, mapping)
    wb = Workbook()
    write_rimon_home_office_summary(wb.active, employee_rows, report_meta)
    write_rimon_home_office_daily(wb.create_sheet(), daily_rows)
    wb.save(output_path)
    return {"warnings": build_rimon_mapping_warnings(mapping)}


# ---- Attendance Alerts ----

ATTENDANCE_ALERT_TYPES = [
    {"id": "day_over_8h", "label": "מעל 8 שעות עבודה ביום", "color": "FFFDE7", "severity": 2, "category": "long_day"},
    {"id": "day_over_9h", "label": "מעל 9 שעות עבודה ביום", "color": "FFF3E0", "severity": 4, "category": "long_day"},
    {"id": "exit_after_20", "label": "עבד אחרי השעה 20:00", "color": "FFF3E0", "severity": 3, "category": "late_exit"},
    {"id": "exit_after_21", "label": "עבד אחרי השעה 21:00", "color": "FFEBEE", "severity": 5, "category": "late_exit"},
    {"id": "exit_after_22", "label": "עבד אחרי השעה 22:00", "color": "FFEBEE", "severity": 6, "category": "late_exit"},
    {"id": "week_over_40h", "label": "עבד מעל 40 שעות שבועיות", "color": "EDE7F6", "severity": 1, "category": "weekly"},
]
ATTENDANCE_ALERT_BY_ID = {a["id"]: a for a in ATTENDANCE_ALERT_TYPES}

ATTENDANCE_ALERT_COLORS = {
    "FFEBEE": PatternFill(fill_type="solid", fgColor="FFEBEE"),
    "FFF3E0": PatternFill(fill_type="solid", fgColor="FFF3E0"),
    "FFFDE7": PatternFill(fill_type="solid", fgColor="FFFDE7"),
    "EDE7F6": PatternFill(fill_type="solid", fgColor="EDE7F6"),
}


def clean_time_text(raw):
    if not raw:
        return ""
    text = str(raw).strip()
    while text and text[0] in ("*", "?"):
        text = text[1:]
    return text.strip()


def parse_exit_hour(text):
    text = clean_time_text(text)
    if not text:
        return None
    if ":" in text:
        parts = text.split(":", 1)
        try:
            return int(parts[0]) + int(parts[1]) / 60.0
        except (ValueError, IndexError):
            return None
    try:
        val = float(text)
        if 0 < val < 1:
            return val * 24
        return val
    except (ValueError, TypeError):
        return None


def compute_day_total_hours(text):
    cleaned = clean_time_text(text)
    return try_parse_hours_value(cleaned)


def detect_day_alerts(exit_text, total_text, enabled_alerts=None):
    alerts = []
    exit_hour = parse_exit_hour(exit_text)
    total_hours = compute_day_total_hours(total_text)
    enabled = enabled_alerts or {a["id"] for a in ATTENDANCE_ALERT_TYPES}
    if exit_hour is not None:
        if exit_hour >= 22.0 and "exit_after_22" in enabled:
            alerts.append({"id": "exit_after_22", "label": "עבד אחרי 22:00", "color": "FFEBEE", "severity": 6})
        elif exit_hour >= 21.0 and "exit_after_21" in enabled:
            alerts.append({"id": "exit_after_21", "label": "עבד אחרי 21:00", "color": "FFEBEE", "severity": 5})
        elif exit_hour >= 20.0 and "exit_after_20" in enabled:
            alerts.append({"id": "exit_after_20", "label": "עבד אחרי 20:00", "color": "FFF3E0", "severity": 3})
    if total_hours is not None:
        if total_hours > 9 and "day_over_9h" in enabled:
            alerts.append({"id": "day_over_9h", "label": "מעל 9 שעות ביום", "color": "FFF3E0", "severity": 4})
        elif total_hours > 8 and "day_over_8h" in enabled:
            alerts.append({"id": "day_over_8h", "label": "מעל 8 שעות ביום", "color": "FFFDE7", "severity": 2})
    return alerts


def get_israeli_week_sunday(date_obj):
    if date_obj is None:
        return None
    weekday = date_obj.weekday()
    offset = (weekday + 1) % 7
    return date_obj - timedelta(days=offset)


def get_israeli_week_range(date_obj):
    sunday = get_israeli_week_sunday(date_obj)
    if sunday is None:
        return None, None
    saturday = sunday + timedelta(days=6)
    return sunday, saturday


def parse_attendance_alerts_report(input_path, extension, mapping, enabled_alerts=None):
    workbook_kind, workbook = open_excel_workbook(input_path, extension)
    sheets = iter_excel_sheets(workbook_kind, workbook)
    employee_summaries = []
    employee_daily_data = {}
    all_dates = []
    enabled = enabled_alerts or {a["id"] for a in ATTENDANCE_ALERT_TYPES}

    for sheet in sheets:
        header_row = detect_rimon_header_row(sheet, workbook_kind)
        emp_name = str(extract_rimon_mapping_value(sheet, workbook_kind, mapping.get("employee_name_source", "")) or "").strip()
        emp_number = str(extract_rimon_mapping_value(sheet, workbook_kind, mapping.get("payroll_number_source", "")) or "").strip()
        department = str(extract_rimon_mapping_value(sheet, workbook_kind, mapping.get("department_source", "")) or "").strip()
        id_number = str(extract_rimon_mapping_value(sheet, workbook_kind, mapping.get("id_number_source", "")) or "").strip()
        if not emp_name:
            continue
        rows, cols = get_excel_dims(sheet, workbook_kind)
        raw_days = []
        for row_idx in range(header_row + 1, rows):
            date_raw = extract_rimon_mapping_value(sheet, workbook_kind, mapping.get("date_source", ""), row_idx)
            date_val = parse_excel_date_generic(workbook_kind, workbook, date_raw)
            entry_raw = extract_rimon_mapping_value(sheet, workbook_kind, mapping.get("entry_time_source", ""), row_idx)
            exit_raw = extract_rimon_mapping_value(sheet, workbook_kind, mapping.get("exit_time_source", ""), row_idx)
            total_raw = extract_rimon_mapping_value(sheet, workbook_kind, mapping.get("total_hours_source", ""), row_idx)
            standard_raw = extract_rimon_mapping_value(sheet, workbook_kind, mapping.get("standard_hours_source", ""), row_idx)
            break_raw = extract_rimon_mapping_value(sheet, workbook_kind, mapping.get("break_hours_source", ""), row_idx)
            missing_raw = extract_rimon_mapping_value(sheet, workbook_kind, mapping.get("missing_hours_source", ""), row_idx)
            regular_raw = extract_rimon_mapping_value(sheet, workbook_kind, mapping.get("regular_hours_source", ""), row_idx)
            ot100_raw = extract_rimon_mapping_value(sheet, workbook_kind, mapping.get("overtime_100_source", ""), row_idx)
            ot125_raw = extract_rimon_mapping_value(sheet, workbook_kind, mapping.get("overtime_125_source", ""), row_idx)
            ot150_raw = extract_rimon_mapping_value(sheet, workbook_kind, mapping.get("overtime_150_source", ""), row_idx)
            ot200_raw = extract_rimon_mapping_value(sheet, workbook_kind, mapping.get("overtime_200_source", ""), row_idx)
            event_raw = extract_rimon_mapping_value(sheet, workbook_kind, mapping.get("event_source", ""), row_idx)
            error_raw = extract_rimon_mapping_value(sheet, workbook_kind, mapping.get("error_text_source", ""), row_idx)
            day_name_raw = extract_rimon_mapping_value(sheet, workbook_kind, mapping.get("day_name_source", ""), row_idx)

            entry_text = clean_time_text(stringify_excel_value(entry_raw))
            exit_text = clean_time_text(stringify_excel_value(exit_raw))
            total_text = clean_time_text(stringify_excel_value(total_raw))
            standard_text = clean_time_text(stringify_excel_value(standard_raw))
            break_text = clean_time_text(stringify_excel_value(break_raw))
            missing_text = clean_time_text(stringify_excel_value(missing_raw))
            regular_text = clean_time_text(stringify_excel_value(regular_raw))
            ot100_text = clean_time_text(stringify_excel_value(ot100_raw))
            ot125_text = clean_time_text(stringify_excel_value(ot125_raw))
            ot150_text = clean_time_text(stringify_excel_value(ot150_raw))
            ot200_text = clean_time_text(stringify_excel_value(ot200_raw))
            event_text = stringify_excel_value(event_raw).strip()
            error_text = stringify_excel_value(error_raw).strip()
            day_name_text = stringify_excel_value(day_name_raw).strip()

            if not date_val and not entry_text and not exit_text and not total_text and not event_text:
                continue
            raw_days.append({
                "date": date_val, "day_name": day_name_text,
                "entry": entry_text, "exit": exit_text, "total": total_text,
                "standard": standard_text, "break": break_text, "missing": missing_text,
                "regular": regular_text, "ot100": ot100_text, "ot125": ot125_text,
                "ot150": ot150_text, "ot200": ot200_text,
                "event": event_text, "error": error_text,
            })

        # Merge split-shift continuation rows (same date)
        merged_days = []
        for day in raw_days:
            if merged_days and day["date"] and merged_days[-1]["date"] == day["date"]:
                prev = merged_days[-1]
                if day["entry"] and not prev["entry"]:
                    prev["entry"] = day["entry"]
                if day["exit"]:
                    prev["exit"] = day["exit"]
                if day["total"]:
                    prev["total"] = day["total"]
                for field in ("standard", "break", "missing", "regular", "ot100", "ot125", "ot150", "ot200"):
                    if day[field] and not prev[field]:
                        prev[field] = day[field]
                if day["event"] and day["event"] != prev.get("event", ""):
                    prev["event"] = (prev.get("event", "") + " / " + day["event"]).strip(" / ")
                if day["error"] and day["error"] != prev.get("error", ""):
                    prev["error"] = (prev.get("error", "") + " / " + day["error"]).strip(" / ")
            else:
                merged_days.append(dict(day))

        # Detect per-day alerts
        total_alerts = 0
        daily_data = []
        for day in merged_days:
            alerts = detect_day_alerts(day["exit"], day["total"], enabled)
            total_alerts += len(alerts)
            max_severity = max((a["severity"] for a in alerts), default=0)
            row_color = ""
            if max_severity > 0:
                row_color = next((a["color"] for a in alerts if a["severity"] == max_severity), "")
            day["alerts"] = alerts
            day["row_color"] = row_color
            day["alert_texts"] = " | ".join(a["label"] for a in sorted(alerts, key=lambda x: -x["severity"]))
            day["weekly_alert"] = False
            day["weekly_range_label"] = ""
            if day["date"]:
                all_dates.append(day["date"])
            daily_data.append(day)

        # Weekly 40h check — only complete weeks (Sun-Sat fully within data range)
        if "week_over_40h" in enabled and daily_data:
            emp_dates = sorted(d["date"] for d in daily_data if d["date"])
            if emp_dates:
                first_date = emp_dates[0]
                last_date = emp_dates[-1]
                weekly_totals = {}
                weekly_date_ranges = {}
                for day in daily_data:
                    if not day["date"]:
                        continue
                    sunday = get_israeli_week_sunday(day["date"])
                    saturday = sunday + timedelta(days=6)
                    # Only count complete weeks: both Sunday and Saturday must be within data range
                    if sunday < first_date or saturday > last_date:
                        continue
                    wk_key = (sunday, saturday)
                    total_h = compute_day_total_hours(day["total"])
                    if total_h is not None:
                        weekly_totals[wk_key] = weekly_totals.get(wk_key, 0) + total_h
                    weekly_date_ranges[wk_key] = (sunday, saturday)

                for day in daily_data:
                    if not day["date"]:
                        continue
                    sunday = get_israeli_week_sunday(day["date"])
                    saturday = sunday + timedelta(days=6)
                    wk_key = (sunday, saturday)
                    if wk_key in weekly_totals and weekly_totals[wk_key] > 40:
                        day["weekly_alert"] = True
                        range_label = sunday.strftime("%d/%m") + "-" + saturday.strftime("%d/%m")
                        day["weekly_range_label"] = range_label
                        weekly_label = "מעל 40ש שבועיות (" + range_label + ")"
                        weekly_alert = {"id": "week_over_40h", "label": weekly_label, "color": "EDE7F6", "severity": 1}
                        day["alerts"].append(weekly_alert)
                        if not day["alert_texts"]:
                            day["alert_texts"] = weekly_label
                        else:
                            day["alert_texts"] += " | " + weekly_label
                        total_alerts += 1

        alert_counts = {}
        for day in daily_data:
            for a in day["alerts"]:
                alert_counts[a["id"]] = alert_counts.get(a["id"], 0) + 1

        employee_summaries.append({
            "name": emp_name, "number": emp_number, "department": department,
            "id_number": id_number, "total_alerts": total_alerts,
            "alert_counts": alert_counts, "days_count": len(merged_days),
        })
        employee_daily_data[emp_name] = daily_data

    if workbook_kind == "xlsx":
        workbook.close()

    report_meta = {
        "total_employees": len(employee_summaries),
        "total_alerts": sum(e["total_alerts"] for e in employee_summaries),
        "employees_with_alerts": sum(1 for e in employee_summaries if e["total_alerts"] > 0),
        "date_range": (min(all_dates).strftime("%d/%m/%Y") + " - " + max(all_dates).strftime("%d/%m/%Y")) if all_dates else "",
    }
    return employee_summaries, employee_daily_data, report_meta


def write_attendance_alerts_summary(ws, summaries, employee_daily_data, meta, enabled_alerts=None):
    ws.title = safe_sheet_title("סיכום התראות", "Alert Summary")
    ws.sheet_view.rightToLeft = True
    ws.sheet_view.showGridLines = False
    title_font = Font(name="Calibri", size=16, bold=True, color="1E3A8A")
    section_font = Font(name="Calibri", size=14, bold=True, color="1E3A8A")
    header_font = Font(name="Calibri", size=11, bold=True, color="FFFFFF")
    header_fill = PatternFill(fill_type="solid", fgColor="1E3A8A")
    enabled = enabled_alerts or {a["id"] for a in ATTENDANCE_ALERT_TYPES}

    # --- Title ---
    ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=7)
    ws["A1"] = "דוח התראות נוכחות"
    ws["A1"].font = title_font
    ws["A1"].alignment = Alignment(horizontal="center")
    ws["A1"].fill = PatternFill(fill_type="solid", fgColor="BFDBFE")
    ws.cell(row=2, column=1, value=meta.get("date_range", "")).font = Font(name="Calibri", size=11, color="64748B")

    # --- Metrics ---
    metrics = [
        ("סה\"כ עובדים", str(meta["total_employees"]), "E0F2FE"),
        ("סה\"כ התראות", str(meta["total_alerts"]), "FEE2E2"),
        ("עובדים עם התראות", str(meta["employees_with_alerts"]), "FEF3C7"),
    ]
    for idx, (label, value, fill_color) in enumerate(metrics):
        label_cell = ws.cell(row=4, column=idx * 2 + 1, value=label)
        value_cell = ws.cell(row=4, column=idx * 2 + 2, value=value)
        label_cell.font = Font(bold=True, color="334155")
        value_cell.font = Font(bold=True, size=14, color="1E3A8A")
        label_cell.fill = PatternFill(fill_type="solid", fgColor=fill_color)
        value_cell.fill = PatternFill(fill_type="solid", fgColor=fill_color)
        label_cell.alignment = Alignment(horizontal="right")
        value_cell.alignment = Alignment(horizontal="center")

    # --- Employee summary table ---
    row_num = 6
    emp_headers = ["שם עובד", "מספר עובד", "מחלקה", "סה\"כ התראות", "יציאה מאוחרת", "יום ארוך", "שבוע מעל 40ש"]
    for col_idx, h in enumerate(emp_headers, 1):
        cell = ws.cell(row=row_num, column=col_idx, value=h)
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = Alignment(horizontal="center")
    sorted_summaries = sorted(summaries, key=lambda x: -x["total_alerts"])
    for emp in sorted_summaries:
        row_num += 1
        exit_alerts = emp["alert_counts"].get("exit_after_20", 0) + emp["alert_counts"].get("exit_after_21", 0) + emp["alert_counts"].get("exit_after_22", 0)
        day_alerts = emp["alert_counts"].get("day_over_8h", 0) + emp["alert_counts"].get("day_over_9h", 0)
        week_alerts = emp["alert_counts"].get("week_over_40h", 0)
        values = [emp["name"], emp["number"], emp["department"], emp["total_alerts"], exit_alerts, day_alerts, week_alerts]
        for col_idx, val in enumerate(values, 1):
            cell = ws.cell(row=row_num, column=col_idx, value=val)
            cell.alignment = Alignment(horizontal="right")
            if row_num % 2 == 0:
                cell.fill = PatternFill(fill_type="solid", fgColor="F8FAFC")
            if emp["total_alerts"] > 0 and col_idx == 4:
                cell.font = Font(bold=True, color="DC2626")

    # --- Department summary table ---
    has_departments = any(e["department"] for e in summaries)
    if has_departments:
        row_num += 2
        ws.cell(row=row_num, column=1, value="סיכום לפי מחלקות").font = section_font
        row_num += 1
        dept_headers = ["מחלקה", "עובדים", "סה\"כ התראות", "יציאה מאוחרת", "יום ארוך", "שבוע מעל 40ש"]
        for col_idx, h in enumerate(dept_headers, 1):
            cell = ws.cell(row=row_num, column=col_idx, value=h)
            cell.font = header_font
            cell.fill = PatternFill(fill_type="solid", fgColor="334155")
            cell.alignment = Alignment(horizontal="center")
        dept_data = {}
        for emp in summaries:
            dept = emp["department"] or "(ללא מחלקה)"
            if dept not in dept_data:
                dept_data[dept] = {"count": 0, "total": 0, "exit": 0, "day": 0, "week": 0}
            dept_data[dept]["count"] += 1
            dept_data[dept]["total"] += emp["total_alerts"]
            dept_data[dept]["exit"] += emp["alert_counts"].get("exit_after_20", 0) + emp["alert_counts"].get("exit_after_21", 0) + emp["alert_counts"].get("exit_after_22", 0)
            dept_data[dept]["day"] += emp["alert_counts"].get("day_over_8h", 0) + emp["alert_counts"].get("day_over_9h", 0)
            dept_data[dept]["week"] += emp["alert_counts"].get("week_over_40h", 0)
        for dept_name in sorted(dept_data.keys(), key=lambda d: -dept_data[d]["total"]):
            row_num += 1
            d = dept_data[dept_name]
            for col_idx, val in enumerate([dept_name, d["count"], d["total"], d["exit"], d["day"], d["week"]], 1):
                cell = ws.cell(row=row_num, column=col_idx, value=val)
                cell.alignment = Alignment(horizontal="right")
                if row_num % 2 == 0:
                    cell.fill = PatternFill(fill_type="solid", fgColor="F8FAFC")

    # --- Per alert type breakdown ---
    alert_type_order = [a for a in ATTENDANCE_ALERT_TYPES if a["id"] in enabled]
    for alert_type in alert_type_order:
        instances = []
        for emp in sorted_summaries:
            daily = employee_daily_data.get(emp["name"], [])
            for day in daily:
                for alert in day.get("alerts", []):
                    if alert["id"] == alert_type["id"]:
                        instances.append({"name": emp["name"], "number": emp["number"], "department": emp["department"], "date": day.get("date"), "detail": day.get("exit", "") if "exit" in alert["id"] else day.get("total", "")})
        if not instances:
            continue
        row_num += 2
        ws.cell(row=row_num, column=1, value=alert_type["label"] + " (" + str(len(instances)) + ")").font = section_font
        row_num += 1
        alert_color = alert_type["color"]
        alert_fill = ATTENDANCE_ALERT_COLORS.get(alert_color)
        sub_headers = ["שם עובד", "מספר עובד", "מחלקה", "תאריך", "פירוט"]
        for col_idx, h in enumerate(sub_headers, 1):
            cell = ws.cell(row=row_num, column=col_idx, value=h)
            cell.font = header_font
            cell.fill = PatternFill(fill_type="solid", fgColor="475569")
            cell.alignment = Alignment(horizontal="center")
        for inst in sorted(instances, key=lambda x: (x["name"], x["date"] or date.min)):
            row_num += 1
            date_str = inst["date"].strftime("%d/%m/%Y") if inst["date"] else ""
            detail = ""
            if "exit" in alert_type["id"]:
                detail = "יציאה: " + inst["detail"] if inst["detail"] else ""
            elif "day_over" in alert_type["id"]:
                detail = 'סה"כ: ' + inst["detail"] if inst["detail"] else ""
            elif "week" in alert_type["id"]:
                detail = "שבוע מעל 40 שעות"
            for col_idx, val in enumerate([inst["name"], inst["number"], inst["department"], date_str, detail], 1):
                cell = ws.cell(row=row_num, column=col_idx, value=val)
                cell.alignment = Alignment(horizontal="right")
                if alert_fill:
                    cell.fill = alert_fill

    widths = [22, 14, 22, 14, 22, 14, 18]
    for col_idx, w in enumerate(widths, 1):
        ws.column_dimensions[get_column_letter(col_idx)].width = w


def write_attendance_alerts_employee_tab(ws, name, daily_data, emp_summary):
    ws.title = safe_sheet_title(name, "Employee")
    ws.sheet_view.rightToLeft = True
    ws.sheet_view.showGridLines = False

    header_font = Font(name="Calibri", size=11, bold=True, color="FFFFFF")
    header_fill = PatternFill(fill_type="solid", fgColor="1E3A8A")
    purple_border = Border(left=Side(style="medium", color="7C3AED"))

    # --- Employee summary header ---
    ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=8)
    ws["A1"] = name
    ws["A1"].font = Font(name="Calibri", size=14, bold=True, color="1E3A8A")
    ws["A1"].alignment = Alignment(horizontal="center")
    ws["A1"].fill = PatternFill(fill_type="solid", fgColor="BFDBFE")

    total_alerts = emp_summary.get("total_alerts", 0)
    alert_counts = emp_summary.get("alert_counts", {})
    exit_alerts = alert_counts.get("exit_after_20", 0) + alert_counts.get("exit_after_21", 0) + alert_counts.get("exit_after_22", 0)
    day_alerts = alert_counts.get("day_over_8h", 0) + alert_counts.get("day_over_9h", 0)
    week_alerts = alert_counts.get("week_over_40h", 0)

    summary_items = [
        ("סה\"כ התראות", str(total_alerts), "FEE2E2"),
        ("יציאה מאוחרת", str(exit_alerts), "FFF3E0"),
        ("יום ארוך", str(day_alerts), "FFFDE7"),
        ("שבוע מעל 40ש", str(week_alerts), "EDE7F6"),
    ]
    for idx, (label, value, color) in enumerate(summary_items):
        lc = ws.cell(row=2, column=idx * 2 + 1, value=label)
        vc = ws.cell(row=2, column=idx * 2 + 2, value=value)
        lc.font = Font(bold=True, size=10, color="334155")
        vc.font = Font(bold=True, size=13, color="1E3A8A")
        lc.fill = PatternFill(fill_type="solid", fgColor=color)
        vc.fill = PatternFill(fill_type="solid", fgColor=color)
        lc.alignment = Alignment(horizontal="right")
        vc.alignment = Alignment(horizontal="center")

    # Alert dates mini-summary
    alert_dates = sorted(set(d["date"] for d in daily_data if d.get("alerts") and d.get("date")))
    dates_text = ", ".join(d.strftime("%d/%m") for d in alert_dates[:15])
    if len(alert_dates) > 15:
        dates_text += "..."
    if dates_text:
        ws.cell(row=3, column=1, value="תאריכי התראות: " + dates_text).font = Font(size=10, color="64748B")

    # --- Daily data table ---
    data_start_row = 5
    headers = [
        "תאריך", "יום", "כניסה", "יציאה", "סה\"כ שעות", "תקן", "הפסקה", "חוסר",
        "רגילות", "100%", "125%", "150%", "200%", "אירוע", "שגיאות", "התראות",
    ]
    for col_idx, h in enumerate(headers, 1):
        cell = ws.cell(row=data_start_row, column=col_idx, value=h)
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = Alignment(horizontal="center")
    ws.freeze_panes = "A" + str(data_start_row + 1)

    for row_offset, day in enumerate(daily_data):
        row_idx = data_start_row + 1 + row_offset
        date_str = day["date"].strftime("%d/%m/%Y") if day.get("date") else ""
        values = [
            date_str, day.get("day_name", ""), day.get("entry", ""), day.get("exit", ""),
            day.get("total", ""), day.get("standard", ""), day.get("break", ""), day.get("missing", ""),
            day.get("regular", ""), day.get("ot100", ""), day.get("ot125", ""),
            day.get("ot150", ""), day.get("ot200", ""),
            day.get("event", ""), day.get("error", ""), day.get("alert_texts", ""),
        ]
        row_color = day.get("row_color", "")
        fill = ATTENDANCE_ALERT_COLORS.get(row_color) if row_color else None
        for col_idx, val in enumerate(values, 1):
            cell = ws.cell(row=row_idx, column=col_idx, value=val)
            cell.alignment = Alignment(horizontal="right")
            if fill:
                cell.fill = fill
            elif row_idx % 2 == 0:
                cell.fill = PatternFill(fill_type="solid", fgColor="F8FAFC")
            if day.get("weekly_alert"):
                cell.border = purple_border

    widths = [12, 10, 8, 8, 10, 8, 8, 8, 8, 8, 8, 8, 8, 18, 18, 28]
    for col_idx, w in enumerate(widths, 1):
        ws.column_dimensions[get_column_letter(col_idx)].width = w


def run_attendance_alerts(input_path, output_path, extension, options=None):
    if extension not in {"xls", "xlsx"}:
        raise ValueError("כלי התראות הנוכחות תומך בקבצי XLS ו-XLSX בלבד")
    options = options or {}
    mapping = default_attendance_alerts_mapping()
    mapping.update({key: value for key, value in options.items() if key.endswith("_source")})
    enabled_alerts = set()
    for a in ATTENDANCE_ALERT_TYPES:
        if options.get("alert_" + a["id"]) != "0":
            enabled_alerts.add(a["id"])
    if not enabled_alerts:
        enabled_alerts = {a["id"] for a in ATTENDANCE_ALERT_TYPES}
    summaries, daily_data, meta = parse_attendance_alerts_report(input_path, extension, mapping, enabled_alerts)
    wb = Workbook()
    write_attendance_alerts_summary(wb.active, summaries, daily_data, meta, enabled_alerts)
    for emp in summaries:
        emp_daily = daily_data.get(emp["name"], [])
        if emp_daily:
            write_attendance_alerts_employee_tab(wb.create_sheet(), emp["name"], emp_daily, emp)
    wb.save(output_path)
    return {"warnings": []}


def run_matan_missing_filter(input_path, output_path, extension, options=None):
    if extension != "xls":
        raise ValueError("Matan missing-hours tool currently supports XLS export only")
    options = options or {}
    mapping = default_matan_missing_mapping()
    mapping.update({key: value for key, value in options.items() if key.endswith("_source")})
    rows = parse_matan_missing_report(input_path, mapping)
    filtered_rows = apply_matan_missing_filters(rows, options)
    wb = Workbook()
    write_matan_missing_summary(
        wb.active,
        filtered_rows,
        {
            "מינימום שעות חוסר": options.get("min_missing_hours", ""),
            "מקסימום שעות חוסר": options.get("max_missing_hours", ""),
        },
        mapping,
    )
    write_matan_missing_filtered(wb.create_sheet(), filtered_rows, mapping)
    wb.save(output_path)
    return {"warnings": build_matan_missing_mapping_warnings(mapping)}


def run_inactive_workers_report(input_path, output_path, extension, options=None):
    if extension not in {"xls", "xlsx"}:
        raise ValueError("Inactive workers report currently supports XLS and XLSX uploads only")
    options = options or {}
    mapping = default_inactive_workers_mapping()
    mapping.update({key: value for key, value in options.items() if key.endswith("_source")})
    inactive_rows, meta = parse_inactive_workers_report(input_path, extension, mapping, options)
    wb = Workbook()
    write_inactive_workers_summary(wb.active, inactive_rows, meta, mapping)
    write_inactive_workers_by_department(wb.create_sheet(), inactive_rows, mapping)
    wb.save(output_path)
    return {"warnings": build_inactive_workers_mapping_warnings(mapping, options)}


def run_flamingo_payroll(input_path, output_path, extension, options=None):
    if extension not in {"xls", "xlsx"}:
        raise ValueError("Flamingo payroll currently supports XLS and XLSX uploads only")
    options = options or {}
    workbook_kind, workbook = open_excel_workbook(input_path, extension)
    mapping = default_flamingo_mapping()
    mapping.update({key: value for key, value in options.items() if key.endswith("_source")})
    manual_hourly_rate_text = str(options.get("manual_hourly_rate", "") or "").strip()

    worker_rows = []
    for detail_sheet, summary_sheet in iter_flamingo_worker_blocks(workbook_kind, workbook):
        worker_rows.append(extract_flamingo_worker_pair(detail_sheet, summary_sheet, workbook_kind, mapping, manual_hourly_rate_text))

    output_wb = Workbook()
    summary_ws = output_wb.active
    write_flamingo_summary_sheet(summary_ws, worker_rows)
    write_flamingo_attention_sheet(output_wb.create_sheet(), worker_rows)
    write_flamingo_department_sheet(output_wb.create_sheet(), worker_rows)
    write_flamingo_top_earners_sheet(output_wb.create_sheet(), worker_rows)
    output_wb.save(output_path)
    return {"warnings": build_flamingo_mapping_warnings(mapping, manual_hourly_rate_text)}


# ── Department Payroll ──────────────────────────────────────────────

def default_dept_payroll_mapping():
    return {
        "worker_name_source": "meta:שם לתצוגה",
        "dept_number_source": "meta:מספר מחלקה",
        "department_source": "meta:מחלקה",
        "payable_hours_source": "summary:שעות לתשלום",
        "id_number_source": "meta:תעודת זהות",
        "phone_source": "",
        "passport_source": "meta:דרכון",
        "event_source": "meta:אירוע",
        "entry_time_source": "meta:כניסה",
        "exit_time_source": "meta:יציאה",
        "date_source": "meta:תאריך",
        "total_hours_source": "meta:סה\"כ",
        "hourly_rate_source": "meta:תעריף שעתי",
        "housing_source": "__auto__",
    }


def find_value_by_label_nearby_or_below(sheet, workbook_kind, label_text, max_col_distance=8, min_row=0, max_row=None):
    """Like find_value_by_label_nearby but also checks the row below the label."""
    normalized_label = normalize_token(label_text)
    rows, cols = get_flamingo_sheet_dims(sheet, workbook_kind)
    last_row = rows if max_row is None else min(max_row, rows)
    for row_index in range(min_row, last_row):
        for col_index in range(cols):
            token = normalize_token(get_flamingo_sheet_cell(sheet, workbook_kind, row_index, col_index))
            if token != normalized_label:
                continue
            # Try to the right first (standard)
            for next_col in range(col_index + 1, min(cols, col_index + max_col_distance + 1)):
                candidate = get_flamingo_sheet_cell(sheet, workbook_kind, row_index, next_col)
                if candidate not in ("", None):
                    candidate_token = normalize_token(candidate)
                    if candidate_token in FLAMINGO_META_LABEL_TOKENS:
                        continue
                    return candidate
            # Fallback: check the row below at the same column
            if row_index + 1 < rows:
                below = get_flamingo_sheet_cell(sheet, workbook_kind, row_index + 1, col_index)
                if below not in ("", None):
                    return below
    return ""


def extract_dept_payroll_worker(detail_sheet, summary_sheet, workbook_kind, mapping):
    """Extract enriched worker data including daily rows, housing, hourly rate, and client breakdown."""
    worker_name = stringify_excel_value(extract_flamingo_mapping_value(detail_sheet, summary_sheet, workbook_kind, mapping.get("worker_name_source")))
    worker_name = worker_name or get_flamingo_sheet_name(detail_sheet, workbook_kind)
    # For dept_number, try standard extraction first, then enhanced below-label search
    dept_number = stringify_excel_value(extract_flamingo_mapping_value(detail_sheet, summary_sheet, workbook_kind, mapping.get("dept_number_source")))
    if not dept_number:
        source = mapping.get("dept_number_source", "")
        if source.startswith("meta:"):
            label = source.split(":", 1)[1]
            dept_number = stringify_excel_value(find_value_by_label_nearby_or_below(detail_sheet, workbook_kind, label))
    department = stringify_excel_value(extract_flamingo_mapping_value(detail_sheet, summary_sheet, workbook_kind, mapping.get("department_source")))
    id_number = stringify_excel_value(extract_flamingo_mapping_value(detail_sheet, summary_sheet, workbook_kind, mapping.get("id_number_source")))
    phone = stringify_excel_value(extract_flamingo_mapping_value(detail_sheet, summary_sheet, workbook_kind, mapping.get("phone_source")))
    passport = stringify_excel_value(extract_flamingo_mapping_value(detail_sheet, summary_sheet, workbook_kind, mapping.get("passport_source")))

    payable_hours_raw = extract_flamingo_mapping_value(detail_sheet, summary_sheet, workbook_kind, mapping.get("payable_hours_source"))
    notes = []
    status = "OK"
    try:
        payable_hours = parse_hours_value(payable_hours_raw)
    except ValueError:
        payable_hours = None
        status = "שעות לתשלום לא תקינות"
        notes.append(f"ערך שעות התשלום אינו תקין: {payable_hours_raw}")
    if payable_hours is None and status == "OK":
        status = "חסרות שעות לתשלום"
        notes.append("לא זוהו שעות לתשלום עבור העובד.")

    # ── Extract housing charge ──
    housing_charge = 0
    rows, cols = get_flamingo_sheet_dims(detail_sheet, workbook_kind)
    housing_source = mapping.get("housing_source", "__auto__")
    if housing_source and housing_source not in ("__auto__", ""):
        # User selected a specific field from the report
        raw = extract_flamingo_mapping_value(detail_sheet, summary_sheet, workbook_kind, housing_source)
        if raw not in ("", None):
            try:
                housing_charge = float(str(raw).replace(",", ""))
            except (ValueError, TypeError):
                # Try extracting number from text like "חיוב דירה 800"
                m = re.search(r'(\d+(?:\.\d+)?)', str(raw))
                if m:
                    housing_charge = float(m.group(1))
    if not housing_charge and housing_source in ("__auto__", ""):
        # Auto mode: scan notes area for "חיוב דירה NNN"
        notes_label_row = find_sheet_label_row(detail_sheet, workbook_kind, "הערות")
        if notes_label_row >= 0:
            for r in range(max(0, notes_label_row - 1), min(rows, notes_label_row + 5)):
                for c in range(cols):
                    cell_val = stringify_excel_value(get_flamingo_sheet_cell(detail_sheet, workbook_kind, r, c))
                    if not cell_val:
                        continue
                    housing_match = re.search(r'חיוב\s*דירה\s*(\d+(?:\.\d+)?)', cell_val)
                    if housing_match:
                        housing_charge = float(housing_match.group(1))
                        break
                if housing_charge:
                    break
        # Fallback: scan meta area for housing
        if not housing_charge:
            for r in range(min(rows, 12)):
                for c in range(cols):
                    cell_val = stringify_excel_value(get_flamingo_sheet_cell(detail_sheet, workbook_kind, r, c))
                    if cell_val and re.search(r'חיוב\s*דירה\s*(\d+(?:\.\d+)?)', cell_val):
                        housing_charge = float(re.search(r'חיוב\s*דירה\s*(\d+(?:\.\d+)?)', cell_val).group(1))
                        break
                if housing_charge:
                    break

    # ── Find daily header row and column positions ──
    daily_header_row = find_sheet_label_row(detail_sheet, workbook_kind, "תאריך")
    col_positions = {}
    daily_col_fields = {
        "date_source": "date",
        "entry_time_source": "entry_time",
        "exit_time_source": "exit_time",
        "event_source": "client_name",
        "total_hours_source": "total_hours",
        "hourly_rate_source": "hourly_rate",
    }
    if daily_header_row >= 0:
        for c in range(cols):
            header_val = normalize_token(get_flamingo_sheet_cell(detail_sheet, workbook_kind, daily_header_row, c))
            if not header_val:
                continue
            for field_name, col_key in daily_col_fields.items():
                source = mapping.get(field_name, "")
                if source.startswith("meta:"):
                    source_label = normalize_token(source.split(":", 1)[1])
                    if source_label and source_label == header_val:
                        col_positions[col_key] = c
                        break
            # Also detect by known tokens as fallback
            if "date" not in col_positions and header_val == "תאריך":
                col_positions.setdefault("date", c)
            if "entry_time" not in col_positions and header_val == "כניסה":
                col_positions.setdefault("entry_time", c)
            if "exit_time" not in col_positions and header_val == "יציאה":
                col_positions.setdefault("exit_time", c)
            if "client_name" not in col_positions and header_val == "אירוע":
                col_positions.setdefault("client_name", c)
            if "total_hours" not in col_positions and (header_val in ("סהכ", "סה\"כ", "סהכשעות")):
                col_positions.setdefault("total_hours", c)
            if "hourly_rate" not in col_positions and header_val in ("תעריףשעתי", "תעריף"):
                col_positions.setdefault("hourly_rate", c)

    # Also try to find יום column for day name
    if daily_header_row >= 0:
        for c in range(cols):
            header_val = normalize_token(get_flamingo_sheet_cell(detail_sheet, workbook_kind, daily_header_row, c))
            if header_val == "יום":
                col_positions["day_name"] = c
                break

    # ── Extract hourly rate from daily rows ──
    hourly_rate = 0
    if "hourly_rate" in col_positions and daily_header_row >= 0:
        for r in range(daily_header_row + 1, min(rows, daily_header_row + 35)):
            val = get_flamingo_sheet_cell(detail_sheet, workbook_kind, r, col_positions["hourly_rate"])
            if val not in ("", None):
                try:
                    hourly_rate = float(str(val).replace(",", ""))
                    break
                except (ValueError, TypeError):
                    pass

    # ── Helper: format date/time cell values ──
    def format_date_cell(raw):
        """Convert raw cell value to a readable date string (dd/mm/yyyy)."""
        if raw in ("", None):
            return ""
        if isinstance(raw, datetime):
            return raw.strftime("%d/%m/%Y")
        if isinstance(raw, date):
            return raw.strftime("%d/%m/%Y")
        if isinstance(raw, (int, float)):
            try:
                if workbook_kind == "xls":
                    dt = xlrd.xldate_as_datetime(float(raw), detail_sheet.book.datemode)
                else:
                    # openpyxl serial date
                    from openpyxl.utils.datetime import from_excel
                    dt = from_excel(raw)
                return dt.strftime("%d/%m/%Y")
            except (ValueError, OverflowError, TypeError):
                pass
        return str(raw).strip()

    def format_time_cell(raw):
        """Convert raw cell value to HH:MM time string."""
        if raw in ("", None):
            return ""
        if isinstance(raw, datetime):
            return raw.strftime("%H:%M")
        if isinstance(raw, time):
            return raw.strftime("%H:%M")
        if isinstance(raw, (int, float)):
            # Excel stores times as fraction of day
            try:
                total_seconds = round(float(raw) * 86400)
                hours = total_seconds // 3600
                minutes = (total_seconds % 3600) // 60
                if 0 <= hours <= 23:
                    return f"{hours:02d}:{minutes:02d}"
            except (ValueError, OverflowError):
                pass
        s = str(raw).strip()
        if s.startswith("*"):
            s = s.lstrip("*")
        return s

    # ── Extract daily rows ──
    daily_rows = []
    if daily_header_row >= 0:
        for r in range(daily_header_row + 1, min(rows, daily_header_row + 35)):
            date_raw = get_flamingo_sheet_cell(detail_sheet, workbook_kind, r, col_positions.get("date", 0)) if "date" in col_positions else ""
            date_val = format_date_cell(date_raw)
            if not date_val:
                # Check if there's any data in event/hours columns
                event_val = stringify_excel_value(get_flamingo_sheet_cell(detail_sheet, workbook_kind, r, col_positions.get("client_name", 0))) if "client_name" in col_positions else ""
                hours_val = get_flamingo_sheet_cell(detail_sheet, workbook_kind, r, col_positions.get("total_hours", 0)) if "total_hours" in col_positions else ""
                if not event_val and not hours_val:
                    continue
            day_name = stringify_excel_value(get_flamingo_sheet_cell(detail_sheet, workbook_kind, r, col_positions.get("day_name", 0))) if "day_name" in col_positions else ""
            entry_time_raw = get_flamingo_sheet_cell(detail_sheet, workbook_kind, r, col_positions.get("entry_time", 0)) if "entry_time" in col_positions else ""
            entry_time = format_time_cell(entry_time_raw)
            exit_time_raw = get_flamingo_sheet_cell(detail_sheet, workbook_kind, r, col_positions.get("exit_time", 0)) if "exit_time" in col_positions else ""
            exit_time = format_time_cell(exit_time_raw)
            client_name = stringify_excel_value(get_flamingo_sheet_cell(detail_sheet, workbook_kind, r, col_positions.get("client_name", 0))) if "client_name" in col_positions else ""
            total_hours_raw = get_flamingo_sheet_cell(detail_sheet, workbook_kind, r, col_positions.get("total_hours", 0)) if "total_hours" in col_positions else ""
            total_hours = 0
            if total_hours_raw not in ("", None):
                try:
                    total_hours = parse_hours_value(total_hours_raw)
                except (ValueError, TypeError):
                    total_hours = 0
            # Include any row that has meaningful data
            # A real daily row must have a date — rows without dates are summary/total rows
            has_any_data = client_name or total_hours or entry_time or exit_time
            if has_any_data and date_val:
                # Track what's missing for red-highlighting in output
                missing = []
                if not entry_time:
                    missing.append("חסרה כניסה")
                if not exit_time:
                    missing.append("חסרה יציאה")
                if not client_name:
                    missing.append("חסר אירוע/לקוח")
                if not total_hours:
                    missing.append("חסר סה\"כ שעות")
                daily_rows.append({
                    "date": date_val,
                    "day_name": day_name,
                    "entry_time": entry_time,
                    "exit_time": exit_time,
                    "client_name": client_name,
                    "total_hours": total_hours,
                    "missing": missing,
                })

    # ── Extract תנועות מיוחדות section ──
    special_movements = {}
    tnuot_row, tnuot_col = find_sheet_label_position(detail_sheet, workbook_kind, "תנועות מיוחדות")
    if tnuot_row >= 0:
        # Scan only in the column range where the header was found (± 10 cols)
        scan_col_start = max(0, tnuot_col - 2)
        scan_col_end = min(cols, tnuot_col + 12)
        for r in range(tnuot_row + 1, min(rows, tnuot_row + 20)):
            for c in range(scan_col_start, scan_col_end):
                label_val = stringify_excel_value(get_flamingo_sheet_cell(detail_sheet, workbook_kind, r, c))
                if not label_val:
                    continue
                # Look for a numeric value to the right
                for nc in range(c + 1, min(cols, c + 10)):
                    val = get_flamingo_sheet_cell(detail_sheet, workbook_kind, r, nc)
                    if val in ("", None):
                        continue
                    try:
                        hours_val = parse_hours_value(val)
                        if hours_val and hours_val > 0:
                            special_movements[label_val] = hours_val
                        break
                    except (ValueError, TypeError):
                        break
                break  # only process the first non-empty label per row

    return {
        "worker_name": worker_name,
        "dept_number": dept_number,
        "department": department,
        "id_number": id_number,
        "phone": phone,
        "passport": passport,
        "payable_hours": payable_hours,
        "housing_charge": housing_charge,
        "hourly_rate": hourly_rate,
        "daily_rows": daily_rows,
        "special_movements": special_movements,
        "status": status,
        "notes": " | ".join(notes),
    }


def parse_dept_settings_json(dept_settings_json):
    """Parse client settings JSON. Key is customer_name (matches תנועה מיוחדת in report)."""
    if isinstance(dept_settings_json, str):
        dept_settings_json = json.loads(dept_settings_json)
    settings_by_client = {}
    for entry in dept_settings_json:
        key = str(entry.get("customer_name", "")).strip()
        if not key:
            continue
        try:
            charge_rate = float(str(entry.get("charge_rate", 0) or 0).replace(",", ""))
        except (ValueError, TypeError):
            charge_rate = 0
        settings_by_client[key] = {
            "customer_name": key,
            "region_manager": str(entry.get("region_manager", "") or "").strip(),
            "address": str(entry.get("address", "") or "").strip(),
            "charge_rate": charge_rate,
            "contact_person": str(entry.get("contact_person", "") or "").strip(),
            "contact_phone": str(entry.get("contact_phone", "") or "").strip(),
        }
    return settings_by_client


def parse_clients_excel(clients_path, extension):
    """Parse a clients Excel file with client info and worker assignments.
    Returns (clients_by_name, worker_to_client) where:
    - clients_by_name: {normalized_client_name: {name, charge_rate, address, manager, contact, phone}}
    - worker_to_client: {normalized_worker_name: normalized_client_name}
    """
    workbook_kind, workbook = open_excel_workbook(clients_path, extension)
    sheets = iter_excel_sheets(workbook_kind, workbook)
    if not sheets:
        return {}, {}

    sheet = sheets[0]
    rows_count, cols_count = get_excel_dims(sheet, workbook_kind)

    # Find header row by looking for known headers
    header_row = -1
    col_map = {}
    target_headers = {
        "שםהעסק": "business_name",
        "תעריףגביה": "charge_rate",
        "מנהלאזור": "area_manager",
        "כתובתהעסק": "address",
        "אישקשר": "contact",
        "טלפון": "phone",
        "שםהעובד": "worker_name",
        "שכרשעתי": "hourly_rate",
        "מספרעובדים": "worker_count",
        "3%מיסים": "taxes",
        "ברוטו": "gross",
        "280ביטוח": "insurance",
        "חיובדירה": "housing",
        "חיובמפרעות": "advances",
        "חיובויזה": "visa",
        "נטולתשלום": "net_pay",
        "מספרפספורט": "passport",
    }
    for r in range(min(rows_count, 10)):
        found_count = 0
        for c in range(cols_count):
            val = normalize_token(get_excel_cell(sheet, workbook_kind, r, c))
            if val in target_headers:
                col_map[target_headers[val]] = c
                found_count += 1
        if found_count >= 2:
            header_row = r
            break

    if header_row < 0:
        return {}, {}

    clients_by_name = {}
    worker_to_client = {}
    current_client = None

    for r in range(header_row + 1, rows_count):
        biz_name = stringify_excel_value(get_excel_cell(sheet, workbook_kind, r, col_map.get("business_name", 0))) if "business_name" in col_map else ""
        worker_name = stringify_excel_value(get_excel_cell(sheet, workbook_kind, r, col_map.get("worker_name", 0))) if "worker_name" in col_map else ""

        if biz_name:
            current_client = biz_name
            norm_client = normalize_token(biz_name)
            if norm_client not in clients_by_name:
                charge_rate = 0
                if "charge_rate" in col_map:
                    raw = get_excel_cell(sheet, workbook_kind, r, col_map["charge_rate"])
                    try:
                        charge_rate = float(str(raw or 0).replace(",", ""))
                    except (ValueError, TypeError):
                        charge_rate = 0
                clients_by_name[norm_client] = {
                    "name": biz_name,
                    "charge_rate": charge_rate,
                    "address": stringify_excel_value(get_excel_cell(sheet, workbook_kind, r, col_map.get("address", 0))) if "address" in col_map else "",
                    "area_manager": stringify_excel_value(get_excel_cell(sheet, workbook_kind, r, col_map.get("area_manager", 0))) if "area_manager" in col_map else "",
                    "contact": stringify_excel_value(get_excel_cell(sheet, workbook_kind, r, col_map.get("contact", 0))) if "contact" in col_map else "",
                    "phone": stringify_excel_value(get_excel_cell(sheet, workbook_kind, r, col_map.get("phone", 0))) if "phone" in col_map else "",
                }

        if worker_name and current_client:
            worker_to_client[normalize_token(worker_name)] = normalize_token(current_client)

    return clients_by_name, worker_to_client


def write_dept_payroll_output(output_path, worker_rows, dept_settings):
    INSURANCE = 280
    wb = Workbook()
    ws = wb.active
    ws.sheet_view.rightToLeft = True
    ws.sheet_view.showGridLines = False
    ws.title = safe_sheet_title("סיכום לפי מחלקות", "סיכום")
    ws.freeze_panes = "A2"

    header_font = Font(bold=True, size=12, color="FFFFFF")
    header_fill = PatternFill(start_color="1E3A8A", end_color="1E3A8A", fill_type="solid")
    dept_font = Font(bold=True, size=11, color="1E3A8A")
    dept_fill = PatternFill(start_color="DBEAFE", end_color="DBEAFE", fill_type="solid")
    total_font = Font(bold=True, size=10, color="15803D")
    total_fill = PatternFill(start_color="DCFCE7", end_color="DCFCE7", fill_type="solid")
    attn_fill = PatternFill(start_color="FEE2E2", end_color="FEE2E2", fill_type="solid")
    num_fmt = '#,##0.00'
    thin_border = Border(
        left=Side(style="thin", color="E2E8F0"),
        right=Side(style="thin", color="E2E8F0"),
        top=Side(style="thin", color="E2E8F0"),
        bottom=Side(style="thin", color="E2E8F0"),
    )

    worker_headers = [
        "שם עובד", "מספר נייד", "מספר פספורט", "ת.ז / דרכון",
        "שעות", "שכר שעתי", "ברוטו", "3% מיסים", "280 ביטוח",
        "חיוב דירה", "נטו לתשלום", "תעריף גביה", "חיוב ללקוח",
    ]
    dept_info_headers = [
        "מנהל אזור", "כתובת העסק", "שם לקוח", "איש קשר", "טלפון איש קשר",
    ]

    # Group workers by department number
    dept_groups = {}
    unmatched = []
    for w in worker_rows:
        dnum = str(w.get("dept_number", "") or "").strip()
        if dnum:
            dept_groups.setdefault(dnum, []).append(w)
        else:
            unmatched.append(w)

    row = 1
    grand_total_hours = 0
    grand_total_gross = 0
    grand_total_net = 0
    grand_total_charge = 0
    grand_worker_count = 0

    for dept_num in sorted(dept_groups.keys()):
        workers = dept_groups[dept_num]

        # Department header row
        dept_label = f"מחלקה: {dept_num}"
        ws.merge_cells(start_row=row, start_column=1, end_row=row, end_column=len(worker_headers))
        cell = ws.cell(row=row, column=1, value=dept_label)
        cell.font = dept_font
        cell.fill = dept_fill
        cell.alignment = Alignment(horizontal="right")
        row += 1

        # Column headers for workers
        for col_idx, header in enumerate(worker_headers, 1):
            cell = ws.cell(row=row, column=col_idx, value=header)
            cell.font = header_font
            cell.fill = header_fill
            cell.alignment = Alignment(horizontal="center")
            cell.border = thin_border
        row += 1

        dept_total_hours = 0
        dept_total_gross = 0
        dept_total_net = 0
        dept_total_charge = 0

        for w in workers:
            hours = w.get("payable_hours") or 0
            if w["status"] != "OK":
                hours = 0
            hourly_rate = w.get("hourly_rate") or 0
            housing = w.get("housing_charge") or 0
            gross = round(hours * hourly_rate, 2) if hourly_rate else 0
            taxes = round(gross * 0.03, 2)
            net = round(gross - taxes - INSURANCE - housing, 2) if gross else 0
            charge_total = 0

            values = [
                w.get("worker_name", ""),
                w.get("phone", ""),
                w.get("passport", ""),
                w.get("id_number", ""),
                hours,
                hourly_rate,
                gross,
                taxes,
                INSURANCE if gross else 0,
                housing if gross else 0,
                net,
                0,
                charge_total,
            ]
            for col_idx, val in enumerate(values, 1):
                cell = ws.cell(row=row, column=col_idx, value=val)
                cell.border = thin_border
                cell.alignment = Alignment(horizontal="center")
                if isinstance(val, float):
                    cell.number_format = num_fmt
            if w["status"] != "OK":
                for col_idx in range(1, len(worker_headers) + 1):
                    ws.cell(row=row, column=col_idx).fill = attn_fill
            row += 1
            dept_total_hours += hours
            dept_total_gross += gross
            dept_total_net += net
            dept_total_charge += charge_total

        # Department totals row
        ws.cell(row=row, column=1, value=f"סה\"כ מחלקה {dept_num} ({len(workers)} עובדים)")
        ws.cell(row=row, column=1).font = total_font
        ws.cell(row=row, column=1).fill = total_fill
        ws.cell(row=row, column=5, value=round(dept_total_hours, 2)).number_format = num_fmt
        ws.cell(row=row, column=7, value=round(dept_total_gross, 2)).number_format = num_fmt
        ws.cell(row=row, column=11, value=round(dept_total_net, 2)).number_format = num_fmt
        ws.cell(row=row, column=13, value=round(dept_total_charge, 2)).number_format = num_fmt
        for col_idx in range(1, len(worker_headers) + 1):
            ws.cell(row=row, column=col_idx).fill = total_fill
            ws.cell(row=row, column=col_idx).font = total_font
            ws.cell(row=row, column=col_idx).border = thin_border
        row += 2  # blank row between departments

        grand_total_hours += dept_total_hours
        grand_total_gross += dept_total_gross
        grand_total_net += dept_total_net
        grand_total_charge += dept_total_charge
        grand_worker_count += len(workers)

    # Grand total row
    ws.merge_cells(start_row=row, start_column=1, end_row=row, end_column=4)
    ws.cell(row=row, column=1, value=f"סה\"כ כללי ({grand_worker_count} עובדים)")
    ws.cell(row=row, column=1).font = Font(bold=True, size=12, color="1E3A8A")
    ws.cell(row=row, column=5, value=round(grand_total_hours, 2)).number_format = num_fmt
    ws.cell(row=row, column=7, value=round(grand_total_gross, 2)).number_format = num_fmt
    ws.cell(row=row, column=11, value=round(grand_total_net, 2)).number_format = num_fmt
    ws.cell(row=row, column=13, value=round(grand_total_charge, 2)).number_format = num_fmt
    for col_idx in range(1, len(worker_headers) + 1):
        cell = ws.cell(row=row, column=col_idx)
        cell.fill = PatternFill(start_color="E0E7FF", end_color="E0E7FF", fill_type="solid")
        cell.font = Font(bold=True, size=11, color="1E3A8A")
        cell.border = thin_border

    # Auto-width columns
    for col_idx in range(1, len(worker_headers) + 1):
        ws.column_dimensions[get_column_letter(col_idx)].width = 16

    # ── Attention sheet ──
    attn_ws = wb.create_sheet(safe_sheet_title("שימו לב", "שימו לב"))
    attn_ws.sheet_view.rightToLeft = True
    attn_ws.sheet_view.showGridLines = False
    attn_headers = ["שם עובד", "מספר מחלקה", "שם מחלקה", "בעיה", "פירוט"]
    for col_idx, header in enumerate(attn_headers, 1):
        cell = attn_ws.cell(row=1, column=col_idx, value=header)
        cell.font = header_font
        cell.fill = PatternFill(start_color="DC2626", end_color="DC2626", fill_type="solid")
        cell.alignment = Alignment(horizontal="center")
    attn_row = 2
    for w in unmatched:
        attn_ws.cell(row=attn_row, column=1, value=w.get("worker_name", ""))
        attn_ws.cell(row=attn_row, column=2, value=w.get("dept_number", ""))
        attn_ws.cell(row=attn_row, column=3, value=w.get("department", ""))
        attn_ws.cell(row=attn_row, column=4, value="לא נמצאה מחלקה מתאימה")
        attn_ws.cell(row=attn_row, column=5, value=w.get("notes", ""))
        attn_row += 1
    for w in worker_rows:
        if w["status"] != "OK" and w not in unmatched:
            attn_ws.cell(row=attn_row, column=1, value=w.get("worker_name", ""))
            attn_ws.cell(row=attn_row, column=2, value=w.get("dept_number", ""))
            attn_ws.cell(row=attn_row, column=3, value=w.get("department", ""))
            attn_ws.cell(row=attn_row, column=4, value=w.get("status", ""))
            attn_ws.cell(row=attn_row, column=5, value=w.get("notes", ""))
            attn_row += 1
    for dept_num, ds in dept_settings.items():
        if dept_num not in dept_groups:
            attn_ws.cell(row=attn_row, column=1, value="")
            attn_ws.cell(row=attn_row, column=2, value=dept_num)
            attn_ws.cell(row=attn_row, column=3, value=ds.get("customer_name", ""))
            attn_ws.cell(row=attn_row, column=4, value="מחלקה ללא עובדים")
            attn_row += 1
    for col_idx in range(1, len(attn_headers) + 1):
        attn_ws.column_dimensions[get_column_letter(col_idx)].width = 22

    wb.save(output_path)


def write_dept_payroll_output_v2(output_path, worker_rows, dept_settings, clients_info):
    """Write 3-tab output: charge per client, pay to employee, company summary."""
    INSURANCE = 280
    clients_by_name, worker_to_client = clients_info or ({}, {})

    wb = Workbook()

    # ── Professional color palette ──
    header_font = Font(bold=True, size=11, color="FFFFFF")
    header_fill = PatternFill(start_color="1E3A8A", end_color="1E3A8A", fill_type="solid")
    section_font = Font(bold=True, size=11, color="1E3A8A")
    section_fill = PatternFill(start_color="DBEAFE", end_color="DBEAFE", fill_type="solid")
    # Client/worker subtotal — amber/gold
    subtotal_font = Font(bold=True, size=10, color="92400E")
    subtotal_fill = PatternFill(start_color="FEF3C7", end_color="FEF3C7", fill_type="solid")
    # Daily total row per worker — teal
    daily_total_font = Font(bold=True, size=10, color="115E59")
    daily_total_fill = PatternFill(start_color="CCFBF1", end_color="CCFBF1", fill_type="solid")
    # Grand total — bold indigo
    grand_font = Font(bold=True, size=12, color="FFFFFF")
    grand_fill = PatternFill(start_color="312E81", end_color="312E81", fill_type="solid")
    # Error/missing — red
    attn_fill = PatternFill(start_color="FEE2E2", end_color="FEE2E2", fill_type="solid")
    attn_font = Font(color="991B1B")
    # Alternating row fills
    alt_fill_a = PatternFill(start_color="FFFFFF", end_color="FFFFFF", fill_type="solid")
    alt_fill_b = PatternFill(start_color="F8FAFC", end_color="F8FAFC", fill_type="solid")
    # Calculation block — slate
    calc_header_font = Font(bold=True, size=10, color="FFFFFF")
    calc_header_fill = PatternFill(start_color="475569", end_color="475569", fill_type="solid")
    calc_value_fill = PatternFill(start_color="F1F5F9", end_color="F1F5F9", fill_type="solid")
    calc_value_font = Font(bold=True, size=11, color="0F172A")
    # Info text — subtle
    info_font = Font(size=10, color="64748B")

    num_fmt = '#,##0.00'
    thin_border = Border(
        left=Side(style="thin", color="CBD5E1"),
        right=Side(style="thin", color="CBD5E1"),
        top=Side(style="thin", color="CBD5E1"),
        bottom=Side(style="thin", color="CBD5E1"),
    )

    def style_header_row(ws, row, headers):
        for col_idx, h in enumerate(headers, 1):
            cell = ws.cell(row=row, column=col_idx, value=h)
            cell.font = header_font
            cell.fill = header_fill
            cell.alignment = Alignment(horizontal="center", vertical="center")
            cell.border = thin_border
        ws.row_dimensions[row].height = 28

    def write_cell(ws, row, col, value, font=None, fill=None, fmt=None, align=None):
        cell = ws.cell(row=row, column=col, value=value)
        cell.border = thin_border
        cell.alignment = Alignment(horizontal=align or "center", vertical="center")
        if font:
            cell.font = font
        if fill:
            cell.fill = fill
        if fmt and isinstance(value, (int, float)):
            cell.number_format = fmt
        return cell

    # ── Build client→workers mapping from daily rows ──
    # Each worker's daily rows contain client_name (from אירוע column)
    client_workers = {}  # {client_display_name: [{worker, daily_rows_for_client, subtotal_hours}]}
    worker_calculations = []  # enriched worker data with calculations

    for w in worker_rows:
        hours = w.get("payable_hours") or 0
        hourly_rate = w.get("hourly_rate") or 0
        housing = w.get("housing_charge") or 0

        gross = round(hours * hourly_rate, 2) if hourly_rate and hours else 0
        taxes = round(gross * 0.03, 2)
        net = round(gross - taxes - INSURANCE - housing, 2) if gross else 0

        # Group daily rows by client
        client_daily = {}
        unassigned_daily = []
        for dr in w.get("daily_rows", []):
            cn = dr.get("client_name", "").strip()
            if cn:
                client_daily.setdefault(cn, []).append(dr)
            else:
                unassigned_daily.append(dr)

        # Build per-client subtotals from תנועות מיוחדות or daily sums
        # Filter out non-client entries (vacation, sick, holidays)
        non_client_tokens = {"חופשה", "מחלה", "חג", "מילואים", "אבל", "היעדרות"}
        client_hours = {}
        attendance_events = {}
        for cn, hrs in w.get("special_movements", {}).items():
            cn_token = normalize_token(cn)
            is_non_client = any(normalize_token(kw) == cn_token or normalize_token(kw) in cn_token for kw in non_client_tokens)
            if is_non_client:
                attendance_events[cn] = hrs
            else:
                client_hours[cn] = hrs
        # Supplement with daily row sums for clients not in special_movements
        for cn, drs in client_daily.items():
            cn_token = normalize_token(cn)
            is_non_client = any(normalize_token(kw) == cn_token or normalize_token(kw) in cn_token for kw in non_client_tokens)
            if is_non_client:
                if cn not in attendance_events:
                    attendance_events[cn] = round(sum(dr.get("total_hours", 0) for dr in drs), 2)
            elif cn not in client_hours:
                client_hours[cn] = round(sum(dr.get("total_hours", 0) for dr in drs), 2)

        unassigned_hours = round(sum(dr.get("total_hours", 0) for dr in unassigned_daily), 2)

        w_calc = {
            **w,
            "hourly_rate": hourly_rate,
            "housing_charge": housing,
            "gross": gross,
            "taxes": taxes,
            "insurance": INSURANCE if gross else 0,
            "net": net,
            "client_daily": client_daily,
            "unassigned_daily": unassigned_daily,
            "unassigned_hours": unassigned_hours,
            "client_hours": client_hours,
            "attendance_events": attendance_events,
        }
        worker_calculations.append(w_calc)

        # Register worker under each client they worked for
        for cn in client_hours:
            if cn not in client_workers:
                client_workers[cn] = []
            client_workers[cn].append(w_calc)

    # ── TAB 1: חיוב ללקוח (Charge per Client) ──
    ws1 = wb.active
    ws1.sheet_view.rightToLeft = True
    ws1.sheet_view.showGridLines = False
    ws1.title = safe_sheet_title("חיוב ללקוח", "חיוב")
    ws1.freeze_panes = "A2"

    row = 1
    grand_charge_total = 0

    for client_name in sorted(client_workers.keys()):
        workers_for_client = client_workers[client_name]
        norm_cn = normalize_token(client_name)
        client_info = clients_by_name.get(norm_cn, {})
        charge_rate = client_info.get("charge_rate", 0)

        # If no charge_rate from clients file, try dept_settings (keyed by customer_name)
        if not charge_rate:
            ds = dept_settings.get(client_name, {})
            if not ds:
                # Try normalized match
                for ds_key, ds_val in dept_settings.items():
                    if normalize_token(ds_key) == norm_cn:
                        ds = ds_val
                        break
            if ds.get("charge_rate"):
                charge_rate = ds["charge_rate"]
            # Also pull client info from dept_settings if not from clients file
            if not client_info and ds:
                client_info = {
                    "name": ds.get("customer_name", client_name),
                    "area_manager": ds.get("region_manager", ""),
                    "address": ds.get("address", ""),
                    "contact": ds.get("contact_person", ""),
                    "phone": ds.get("contact_phone", ""),
                    "charge_rate": ds.get("charge_rate", 0),
                }

        # Client header
        ws1.merge_cells(start_row=row, start_column=1, end_row=row, end_column=8)
        cell = ws1.cell(row=row, column=1, value=f"לקוח: {client_info.get('name', client_name)}")
        cell.font = section_font
        cell.fill = section_fill
        cell.alignment = Alignment(horizontal="right")
        row += 1

        # Client info line
        info_parts = []
        if client_info.get("area_manager"):
            info_parts.append(f"מנהל אזור: {client_info['area_manager']}")
        if client_info.get("address"):
            info_parts.append(f"כתובת: {client_info['address']}")
        if client_info.get("contact"):
            info_parts.append(f"איש קשר: {client_info['contact']}")
        if client_info.get("phone"):
            info_parts.append(f"טלפון: {client_info['phone']}")
        if charge_rate:
            info_parts.append(f"תעריף גביה: {charge_rate}")
        if info_parts:
            ws1.merge_cells(start_row=row, start_column=1, end_row=row, end_column=8)
            info_cell = ws1.cell(row=row, column=1, value=" | ".join(info_parts))
            info_cell.font = Font(size=10, color="475569")
            info_cell.alignment = Alignment(horizontal="right")
            row += 1

        # Worker headers for this client
        client_headers = ["שם עובד", "תאריך", "כניסה", "יציאה", "שעות", "סה\"כ שעות לעובד", "תעריף גביה", "חיוב"]
        style_header_row(ws1, row, client_headers)
        row += 1

        client_total_hours = 0
        client_total_charge = 0

        for wc in workers_for_client:
            worker_client_hours = wc["client_hours"].get(client_name, 0)
            worker_daily_for_client = wc.get("client_daily", {}).get(client_name, [])
            worker_charge = round(worker_client_hours * charge_rate, 2) if charge_rate else 0

            if worker_daily_for_client:
                first = True
                for idx, dr in enumerate(worker_daily_for_client):
                    row_fill = alt_fill_b if idx % 2 else alt_fill_a
                    write_cell(ws1, row, 1, wc["worker_name"] if first else "", fill=row_fill, align="right")
                    write_cell(ws1, row, 2, dr.get("date", ""), fill=row_fill)
                    write_cell(ws1, row, 3, dr.get("entry_time", ""), fill=row_fill)
                    write_cell(ws1, row, 4, dr.get("exit_time", ""), fill=row_fill)
                    write_cell(ws1, row, 5, dr.get("total_hours", 0), fmt=num_fmt, fill=row_fill)
                    if first:
                        write_cell(ws1, row, 6, worker_client_hours, fmt=num_fmt, fill=row_fill)
                        write_cell(ws1, row, 7, charge_rate, fmt=num_fmt, fill=row_fill)
                        write_cell(ws1, row, 8, worker_charge, fmt=num_fmt, fill=row_fill)
                    else:
                        for c in (6, 7, 8):
                            write_cell(ws1, row, c, "", fill=row_fill)
                    row += 1
                    first = False
            else:
                write_cell(ws1, row, 1, wc["worker_name"], align="right")
                for c in (2, 3, 4, 5):
                    write_cell(ws1, row, c, "")
                write_cell(ws1, row, 6, worker_client_hours, fmt=num_fmt)
                write_cell(ws1, row, 7, charge_rate, fmt=num_fmt)
                write_cell(ws1, row, 8, worker_charge, fmt=num_fmt)
                row += 1

            client_total_hours += worker_client_hours
            client_total_charge += worker_charge

        # Client subtotal row — amber
        write_cell(ws1, row, 1, f"סה\"כ {client_info.get('name', client_name)} ({len(workers_for_client)} עובדים)", font=subtotal_font, fill=subtotal_fill, align="right")
        for col_idx in range(2, 6):
            write_cell(ws1, row, col_idx, "", fill=subtotal_fill)
        write_cell(ws1, row, 6, round(client_total_hours, 2), font=subtotal_font, fill=subtotal_fill, fmt=num_fmt)
        write_cell(ws1, row, 7, "", fill=subtotal_fill)
        write_cell(ws1, row, 8, round(client_total_charge, 2), font=subtotal_font, fill=subtotal_fill, fmt=num_fmt)
        row += 2

        grand_charge_total += client_total_charge

    # ── Unassigned hours section (workers with hours not tied to any client) ──
    unassigned_error_fill = PatternFill(start_color="FEE2E2", end_color="FEE2E2", fill_type="solid")
    unassigned_error_font = Font(color="991B1B")
    unassigned_header_font = Font(bold=True, size=11, color="991B1B")
    unassigned_header_fill = PatternFill(start_color="FECACA", end_color="FECACA", fill_type="solid")

    workers_with_unassigned = [wc for wc in worker_calculations if wc.get("unassigned_daily")]
    if workers_with_unassigned:
        ws1.merge_cells(start_row=row, start_column=1, end_row=row, end_column=8)
        cell = ws1.cell(row=row, column=1, value="שעות ללא שיוך לקוח")
        cell.font = unassigned_header_font
        cell.fill = unassigned_header_fill
        cell.alignment = Alignment(horizontal="right")
        row += 1

        unassigned_headers = ["שם עובד", "תאריך", "כניסה", "יציאה", "שעות", "הערה", "", ""]
        style_header_row(ws1, row, unassigned_headers)
        row += 1

        total_unassigned = 0
        for wc in workers_with_unassigned:
            first = True
            for dr in wc["unassigned_daily"]:
                c1 = write_cell(ws1, row, 1, wc["worker_name"] if first else "")
                c2 = write_cell(ws1, row, 2, dr.get("date", ""))
                c3 = write_cell(ws1, row, 3, dr.get("entry_time", ""))
                c4 = write_cell(ws1, row, 4, dr.get("exit_time", ""))
                c5 = write_cell(ws1, row, 5, dr.get("total_hours", 0), fmt=num_fmt)
                missing = dr.get("missing", [])
                note = " | ".join(missing) if missing else "חסר אירוע/לקוח"
                c6 = write_cell(ws1, row, 6, note)
                for c in [c1, c2, c3, c4, c5, c6]:
                    c.fill = unassigned_error_fill
                    c.font = unassigned_error_font
                row += 1
                first = False
            total_unassigned += wc.get("unassigned_hours", 0)

        write_cell(ws1, row, 1, f"סה\"כ שעות ללא שיוך ({len(workers_with_unassigned)} עובדים)", font=Font(bold=True, color="991B1B"), fill=unassigned_error_fill)
        for col_idx in range(2, 9):
            write_cell(ws1, row, col_idx, "", fill=unassigned_error_fill)
        write_cell(ws1, row, 5, round(total_unassigned, 2), font=Font(bold=True, color="991B1B"), fill=unassigned_error_fill, fmt=num_fmt)
        row += 2

    # ── Workers with NO data at all ──
    workers_no_data = [wc for wc in worker_calculations if not wc.get("daily_rows") and not wc.get("special_movements")]
    if workers_no_data:
        ws1.merge_cells(start_row=row, start_column=1, end_row=row, end_column=8)
        cell = ws1.cell(row=row, column=1, value="עובדים ללא נתוני נוכחות")
        cell.font = unassigned_header_font
        cell.fill = unassigned_header_fill
        cell.alignment = Alignment(horizontal="right")
        row += 1
        for wc in workers_no_data:
            c1 = write_cell(ws1, row, 1, wc["worker_name"])
            c2 = write_cell(ws1, row, 2, "לא נמצאו נתוני עבודה בדוח")
            for c in [c1, c2]:
                c.fill = unassigned_error_fill
                c.font = unassigned_error_font
            row += 1
        row += 1

    # Grand total — bold dark indigo
    if client_workers:
        write_cell(ws1, row, 1, "סה\"כ חיוב כל הלקוחות", font=grand_font, fill=grand_fill, align="right")
        for col_idx in range(2, 8):
            write_cell(ws1, row, col_idx, "", fill=grand_fill)
        write_cell(ws1, row, 8, round(grand_charge_total, 2), font=grand_font, fill=grand_fill, fmt=num_fmt)
        ws1.merge_cells(start_row=row, start_column=1, end_row=row, end_column=7)
        ws1.row_dimensions[row].height = 32

    # Column widths
    tab1_widths = {1: 20, 2: 14, 3: 10, 4: 10, 5: 12, 6: 18, 7: 14, 8: 14}
    for col_idx, w in tab1_widths.items():
        ws1.column_dimensions[get_column_letter(col_idx)].width = w

    # ── TAB 2: תשלום לעובד (Pay to Employee) ──
    ws2 = wb.create_sheet(safe_sheet_title("תשלום לעובד", "תשלום"))
    ws2.sheet_view.rightToLeft = True
    ws2.sheet_view.showGridLines = False
    ws2.freeze_panes = "A2"

    row = 1
    grand_total_net = 0
    grand_total_gross = 0

    for wc in worker_calculations:
        # Worker header
        ws2.merge_cells(start_row=row, start_column=1, end_row=row, end_column=7)
        worker_label = f"עובד: {wc['worker_name']}"
        if wc.get("id_number"):
            worker_label += f" | ת.ז: {wc['id_number']}"
        if wc.get("passport"):
            worker_label += f" | דרכון: {wc['passport']}"
        if wc.get("phone"):
            worker_label += f" | נייד: {wc['phone']}"
        if wc.get("department"):
            worker_label += f" | מחלקה: {wc['department']}"
        cell = ws2.cell(row=row, column=1, value=worker_label)
        cell.font = Font(bold=True, size=12, color="1E3A8A")
        cell.fill = section_fill
        cell.alignment = Alignment(horizontal="right", vertical="center")
        ws2.row_dimensions[row].height = 30
        row += 1

        # Daily breakdown headers
        daily_headers = ["תאריך", "יום", "כניסה", "יציאה", "אירוע/לקוח", "שעות", "הערות"]
        style_header_row(ws2, row, daily_headers)
        row += 1

        daily_rows_list = wc.get("daily_rows", [])
        daily_hours_sum = 0
        if not daily_rows_list:
            c1 = write_cell(ws2, row, 1, "לא נמצאו נתוני נוכחות לעובד זה", font=Font(bold=True, color="991B1B"), fill=attn_fill, align="right")
            for col_idx in range(2, 8):
                write_cell(ws2, row, col_idx, "", fill=attn_fill)
            row += 1
        else:
            for idx, dr in enumerate(daily_rows_list):
                missing = dr.get("missing", [])
                row_has_error = bool(missing)
                if row_has_error:
                    rf, rfont = attn_fill, attn_font
                else:
                    rf = alt_fill_b if idx % 2 else alt_fill_a
                    rfont = None
                write_cell(ws2, row, 1, dr.get("date", ""), fill=rf, font=rfont)
                write_cell(ws2, row, 2, dr.get("day_name", ""), fill=rf, font=rfont)
                write_cell(ws2, row, 3, dr.get("entry_time", "") or ("—" if row_has_error else ""), fill=rf, font=rfont)
                write_cell(ws2, row, 4, dr.get("exit_time", "") or ("—" if row_has_error else ""), fill=rf, font=rfont)
                write_cell(ws2, row, 5, dr.get("client_name", "") or ("—" if row_has_error else ""), fill=rf, font=rfont, align="right")
                hrs = dr.get("total_hours", 0)
                write_cell(ws2, row, 6, hrs, fmt=num_fmt, fill=rf, font=rfont)
                daily_hours_sum += hrs if isinstance(hrs, (int, float)) else 0
                if missing:
                    write_cell(ws2, row, 7, " | ".join(missing), fill=rf, font=Font(size=9, color="991B1B"))
                else:
                    write_cell(ws2, row, 7, "", fill=rf)
                row += 1

            # ── Daily total row — teal ──
            write_cell(ws2, row, 1, "סה\"כ נוכחות", font=daily_total_font, fill=daily_total_fill, align="right")
            for c in range(2, 6):
                write_cell(ws2, row, c, "", fill=daily_total_fill)
            write_cell(ws2, row, 6, round(daily_hours_sum, 2), font=daily_total_font, fill=daily_total_fill, fmt=num_fmt)
            write_cell(ws2, row, 7, f"{len(daily_rows_list)} ימים", font=Font(size=9, color="115E59"), fill=daily_total_fill)
            row += 1

        # Per-client subtotals from תנועות מיוחדות
        if wc.get("client_hours"):
            row += 1
            write_cell(ws2, row, 1, "פירוט לפי לקוח:", font=Font(bold=True, size=10, color="1E3A8A"), align="right")
            row += 1
            for cn_idx, (cn, hrs) in enumerate(wc["client_hours"].items()):
                rf = alt_fill_b if cn_idx % 2 else alt_fill_a
                write_cell(ws2, row, 1, cn, fill=rf, align="right")
                for c in range(2, 6):
                    write_cell(ws2, row, c, "", fill=rf)
                write_cell(ws2, row, 6, hrs, fmt=num_fmt, fill=rf)
                write_cell(ws2, row, 7, "", fill=rf)
                row += 1

        # Attendance events
        if wc.get("attendance_events"):
            row += 1
            write_cell(ws2, row, 1, "נוכחות אחר:", font=Font(bold=True, size=10, color="6B21A8"), align="right")
            row += 1
            purple_font = Font(size=10, color="6B21A8")
            purple_fill = PatternFill(start_color="FAF5FF", end_color="FAF5FF", fill_type="solid")
            for ev_name, ev_hrs in wc["attendance_events"].items():
                write_cell(ws2, row, 1, ev_name, font=purple_font, fill=purple_fill, align="right")
                for c in range(2, 6):
                    write_cell(ws2, row, c, "", fill=purple_fill)
                write_cell(ws2, row, 6, ev_hrs, fmt=num_fmt, font=purple_font, fill=purple_fill)
                write_cell(ws2, row, 7, "", fill=purple_fill)
                row += 1

        # Unassigned hours
        if wc.get("unassigned_hours"):
            write_cell(ws2, row, 1, "שעות ללא שיוך לקוח", font=Font(bold=True, size=10, color="991B1B"), fill=attn_fill, align="right")
            for c in range(2, 6):
                write_cell(ws2, row, c, "", fill=attn_fill)
            write_cell(ws2, row, 6, wc["unassigned_hours"], fmt=num_fmt, fill=attn_fill, font=Font(bold=True, color="991B1B"))
            write_cell(ws2, row, 7, "", fill=attn_fill)
            row += 1

        # ── Payment calculation block — slate header ──
        row += 1
        calc_headers = ["שעות לתשלום", "תעריף שעתי", "ברוטו", "3% מיסים", "ביטוח (280)", "חיוב דירה", "נטו לתשלום"]
        for col_idx, h in enumerate(calc_headers, 1):
            cell = ws2.cell(row=row, column=col_idx, value=h)
            cell.font = calc_header_font
            cell.fill = calc_header_fill
            cell.alignment = Alignment(horizontal="center", vertical="center")
            cell.border = thin_border
        ws2.row_dimensions[row].height = 26
        row += 1

        hours = wc.get("payable_hours") or 0
        net_val = wc.get("net", 0)
        for col_idx, val in enumerate([hours, wc.get("hourly_rate", 0), wc.get("gross", 0), wc.get("taxes", 0), wc.get("insurance", 0), wc.get("housing_charge", 0), net_val], 1):
            cf = calc_value_fill
            cfo = calc_value_font
            if wc["status"] != "OK":
                cf = attn_fill
                cfo = Font(bold=True, size=11, color="991B1B")
            write_cell(ws2, row, col_idx, val, fmt=num_fmt, fill=cf, font=cfo)
        ws2.row_dimensions[row].height = 28
        row += 3

        grand_total_net += wc.get("net", 0)
        grand_total_gross += wc.get("gross", 0)

    # Grand total for employees — dark indigo
    write_cell(ws2, row, 1, f"סה\"כ תשלום לכל העובדים ({len(worker_calculations)} עובדים)", font=grand_font, fill=grand_fill, align="right")
    write_cell(ws2, row, 2, "", fill=grand_fill)
    write_cell(ws2, row, 3, round(grand_total_gross, 2), font=grand_font, fill=grand_fill, fmt=num_fmt)
    for col_idx in range(4, 7):
        write_cell(ws2, row, col_idx, "", fill=grand_fill)
    write_cell(ws2, row, 7, round(grand_total_net, 2), font=grand_font, fill=grand_fill, fmt=num_fmt)
    ws2.merge_cells(start_row=row, start_column=1, end_row=row, end_column=2)
    ws2.row_dimensions[row].height = 32

    # Column widths
    tab2_widths = {1: 16, 2: 12, 3: 10, 4: 10, 5: 18, 6: 12, 7: 26}
    for col_idx, w in tab2_widths.items():
        ws2.column_dimensions[get_column_letter(col_idx)].width = w

    # ── TAB 3: סיכום חברה (Company Summary) — Professional dashboard ──
    ws3 = wb.create_sheet(safe_sheet_title("סיכום חברה", "סיכום"))
    ws3.sheet_view.rightToLeft = True
    ws3.sheet_view.showGridLines = False

    row = 1
    # Title bar — full-width dark indigo
    ws3.merge_cells(start_row=row, start_column=1, end_row=row, end_column=5)
    cell = ws3.cell(row=row, column=1, value="סיכום חברה")
    cell.font = Font(bold=True, size=16, color="FFFFFF")
    cell.fill = grand_fill
    cell.alignment = Alignment(horizontal="center", vertical="center")
    for c in range(2, 6):
        ws3.cell(row=row, column=c).fill = grand_fill
    ws3.row_dimensions[row].height = 38
    row += 2

    # ── KPI summary cards ──
    total_unassigned_hours = round(sum(wc.get("unassigned_hours", 0) for wc in worker_calculations), 2)
    total_assigned_hours = round(sum(sum(wc.get("client_hours", {}).values()) for wc in worker_calculations), 2)
    expected_profit = round(grand_charge_total - grand_total_net, 2)

    kpi_items = [
        ("סה\"כ חיוב ללקוחות", round(grand_charge_total, 2), "15803D"),
        ("סה\"כ שולם לעובדים (נטו)", round(grand_total_net, 2), "1E3A8A"),
        ("סה\"כ ברוטו עובדים", round(grand_total_gross, 2), "475569"),
        ("רווח צפוי (חיוב − נטו)", expected_profit, "15803D" if expected_profit >= 0 else "DC2626"),
        ("סה\"כ שעות משויכות ללקוחות", total_assigned_hours, "115E59"),
        ("סה\"כ שעות ללא שיוך לקוח", total_unassigned_hours, "991B1B" if total_unassigned_hours > 0 else "475569"),
    ]
    # KPI header
    kpi_header_fill = PatternFill(start_color="1E3A8A", end_color="1E3A8A", fill_type="solid")
    ws3.merge_cells(start_row=row, start_column=1, end_row=row, end_column=2)
    cell = ws3.cell(row=row, column=1, value="מדדים עיקריים")
    cell.font = Font(bold=True, size=12, color="FFFFFF")
    cell.fill = kpi_header_fill
    cell.alignment = Alignment(horizontal="center", vertical="center")
    ws3.cell(row=row, column=2).fill = kpi_header_fill
    ws3.row_dimensions[row].height = 30
    row += 1

    for idx, (label, val, color) in enumerate(kpi_items):
        rf = alt_fill_b if idx % 2 else alt_fill_a
        write_cell(ws3, row, 1, label, font=Font(bold=True, size=11, color="334155"), fill=rf, align="right")
        write_cell(ws3, row, 2, val, fmt=num_fmt, font=Font(bold=True, size=12, color=color), fill=rf)
        ws3.row_dimensions[row].height = 26
        row += 1

    # Grand total bar
    write_cell(ws3, row, 1, "מספר עובדים", font=grand_font, fill=grand_fill, align="right")
    write_cell(ws3, row, 2, len(worker_calculations), font=grand_font, fill=grand_fill)
    ws3.row_dimensions[row].height = 30
    row += 2

    # ── Per-client breakdown ──
    ws3.merge_cells(start_row=row, start_column=1, end_row=row, end_column=5)
    cell = ws3.cell(row=row, column=1, value="פירוט לפי לקוח")
    cell.font = Font(bold=True, size=13, color="FFFFFF")
    cell.fill = header_fill
    cell.alignment = Alignment(horizontal="center", vertical="center")
    for c in range(2, 6):
        ws3.cell(row=row, column=c).fill = header_fill
    ws3.row_dimensions[row].height = 30
    row += 1
    client_breakdown_headers = ["לקוח", "שעות", "תעריף גביה", "סה\"כ חיוב", "עובדים"]
    style_header_row(ws3, row, client_breakdown_headers)
    row += 1

    client_grand_hours = 0
    client_grand_charge = 0
    for cn_idx, cn in enumerate(sorted(client_workers.keys())):
        workers_for_cn = client_workers[cn]
        norm_cn = normalize_token(cn)
        ci = clients_by_name.get(norm_cn, {})
        cr = ci.get("charge_rate", 0)
        if not cr:
            ds = dept_settings.get(cn, {})
            if not ds:
                for ds_key, ds_val in dept_settings.items():
                    if normalize_token(ds_key) == norm_cn:
                        ds = ds_val
                        break
            if ds.get("charge_rate"):
                cr = ds["charge_rate"]
        total_hrs = round(sum(wc["client_hours"].get(cn, 0) for wc in workers_for_cn), 2)
        total_ch = round(total_hrs * cr, 2) if cr else 0
        client_grand_hours += total_hrs
        client_grand_charge += total_ch
        rf = alt_fill_b if cn_idx % 2 else alt_fill_a
        write_cell(ws3, row, 1, ci.get("name", cn), fill=rf, align="right")
        write_cell(ws3, row, 2, total_hrs, fmt=num_fmt, fill=rf)
        write_cell(ws3, row, 3, cr, fmt=num_fmt, fill=rf)
        write_cell(ws3, row, 4, total_ch, fmt=num_fmt, fill=rf, font=Font(bold=True, size=10, color="15803D") if total_ch else None)
        write_cell(ws3, row, 5, len(workers_for_cn), fill=rf)
        if not cr:
            for c in range(1, 6):
                ws3.cell(row=row, column=c).font = Font(color="991B1B", italic=True)
        row += 1

    # Client grand total — amber
    write_cell(ws3, row, 1, "סה\"כ לקוחות", font=subtotal_font, fill=subtotal_fill, align="right")
    write_cell(ws3, row, 2, round(client_grand_hours, 2), font=subtotal_font, fill=subtotal_fill, fmt=num_fmt)
    write_cell(ws3, row, 3, "", fill=subtotal_fill)
    write_cell(ws3, row, 4, round(client_grand_charge, 2), font=subtotal_font, fill=subtotal_fill, fmt=num_fmt)
    write_cell(ws3, row, 5, "", fill=subtotal_fill)
    ws3.row_dimensions[row].height = 28
    row += 2

    # ── Per-worker breakdown ──
    ws3.merge_cells(start_row=row, start_column=1, end_row=row, end_column=5)
    cell = ws3.cell(row=row, column=1, value="פירוט לפי עובד")
    cell.font = Font(bold=True, size=13, color="FFFFFF")
    cell.fill = header_fill
    cell.alignment = Alignment(horizontal="center", vertical="center")
    for c in range(2, 6):
        ws3.cell(row=row, column=c).fill = header_fill
    ws3.row_dimensions[row].height = 30
    row += 1
    worker_breakdown_headers = ["שם עובד", "שעות נוכחות", "ברוטו", "נטו", "מצב"]
    style_header_row(ws3, row, worker_breakdown_headers)
    row += 1

    for wc_idx, wc in enumerate(worker_calculations):
        rf = alt_fill_b if wc_idx % 2 else alt_fill_a
        status_ok = wc.get("status") == "OK"
        status_txt = "תקין" if status_ok else wc.get("status", "—")
        status_color = "15803D" if status_ok else "991B1B"
        write_cell(ws3, row, 1, wc["worker_name"], fill=rf, font=Font(bold=True, size=10), align="right")
        write_cell(ws3, row, 2, wc.get("payable_hours") or 0, fmt=num_fmt, fill=rf)
        write_cell(ws3, row, 3, wc.get("gross", 0), fmt=num_fmt, fill=rf)
        write_cell(ws3, row, 4, wc.get("net", 0), fmt=num_fmt, fill=rf, font=Font(bold=True, size=10, color="1E3A8A"))
        write_cell(ws3, row, 5, status_txt, fill=rf, font=Font(bold=True, size=10, color=status_color))
        row += 1

    # Worker grand total — dark indigo
    write_cell(ws3, row, 1, f"סה\"כ {len(worker_calculations)} עובדים", font=grand_font, fill=grand_fill, align="right")
    write_cell(ws3, row, 2, round(sum(wc.get("payable_hours") or 0 for wc in worker_calculations), 2), font=grand_font, fill=grand_fill, fmt=num_fmt)
    write_cell(ws3, row, 3, round(grand_total_gross, 2), font=grand_font, fill=grand_fill, fmt=num_fmt)
    write_cell(ws3, row, 4, round(grand_total_net, 2), font=grand_font, fill=grand_fill, fmt=num_fmt)
    write_cell(ws3, row, 5, "", fill=grand_fill)
    ws3.row_dimensions[row].height = 32

    # Column widths
    tab3_widths = {1: 24, 2: 16, 3: 16, 4: 18, 5: 14}
    for col_idx, w in tab3_widths.items():
        ws3.column_dimensions[get_column_letter(col_idx)].width = w

    # ── TAB 4: שימו לב (Alerts) — Professional heatmap dashboard ──
    attn_ws = wb.create_sheet(safe_sheet_title("שימו לב", "שימו לב"))
    attn_ws.sheet_view.rightToLeft = True
    attn_ws.sheet_view.showGridLines = False

    # Title bar
    attn_ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=4)
    title_cell = attn_ws.cell(row=1, column=1, value="שימו לב — התראות ובעיות")
    title_cell.font = Font(bold=True, size=14, color="FFFFFF")
    title_cell.fill = grand_fill
    title_cell.alignment = Alignment(horizontal="center", vertical="center")
    for c in range(2, 5):
        attn_ws.cell(row=1, column=c).fill = grand_fill
    attn_ws.row_dimensions[1].height = 36

    # Headers
    attn_headers = ["סוג", "שם", "בעיה", "פירוט"]
    for col_idx, header in enumerate(attn_headers, 1):
        cell = attn_ws.cell(row=2, column=col_idx, value=header)
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = Alignment(horizontal="center", vertical="center")
        cell.border = thin_border
    attn_ws.row_dimensions[2].height = 28

    # Priority levels with icons: 1=critical (red), 2=warning (orange), 3=info (blue)
    critical_fill = PatternFill(start_color="FEE2E2", end_color="FEE2E2", fill_type="solid")
    critical_font = Font(bold=True, size=10, color="991B1B")
    warning_fill = PatternFill(start_color="FEF3C7", end_color="FEF3C7", fill_type="solid")
    warning_font = Font(size=10, color="92400E")
    info_fill_alert = PatternFill(start_color="EFF6FF", end_color="EFF6FF", fill_type="solid")
    info_font_alert = Font(size=10, color="1E40AF")

    alerts = []  # (priority, type_label, name, problem, detail)

    # 1. Clients missing required fields (charge_rate is critical)
    for cn in sorted(client_workers.keys()):
        norm_cn = normalize_token(cn)
        ci = clients_by_name.get(norm_cn, {})
        ds = dept_settings.get(cn, {})
        if not ds:
            for ds_key, ds_val in dept_settings.items():
                if normalize_token(ds_key) == norm_cn:
                    ds = ds_val
                    break
        cr = ci.get("charge_rate", 0) or ds.get("charge_rate", 0)
        if not cr:
            alerts.append((1, "לקוח", cn, "חסר תעריף גביה", "לא ניתן לחשב חיוב ללקוח ללא תעריף. הוסף את הלקוח בטבלת הגדרת פרטי לקוחות."))
        # Optional missing fields
        has_info = bool(ci) or bool(ds)
        if has_info:
            missing_optional = []
            if not (ci.get("area_manager") or ds.get("region_manager")):
                missing_optional.append("מנהל אזור")
            if not (ci.get("address") or ds.get("address")):
                missing_optional.append("כתובת")
            if not (ci.get("contact") or ds.get("contact_person")):
                missing_optional.append("איש קשר")
            if not (ci.get("phone") or ds.get("contact_phone")):
                missing_optional.append("טלפון")
            if missing_optional:
                alerts.append((3, "לקוח", cn, "פרטים חסרים", "חסרים: " + ", ".join(missing_optional)))
        else:
            alerts.append((2, "לקוח", cn, "לקוח לא מוגדר", "הלקוח לא נמצא בהגדרות. הוסף אותו בטבלת הגדרת פרטי לקוחות."))

    # 2. Worker-level issues (critical)
    for w in worker_rows:
        if w["status"] != "OK":
            alerts.append((1, "עובד", w.get("worker_name", ""), w["status"], w.get("notes", "")))

    # 3. Workers without client assignment (warning)
    for w in worker_rows:
        if w["status"] == "OK" and not w.get("daily_rows") and not w.get("special_movements"):
            alerts.append((2, "עובד", w.get("worker_name", ""), "לא זוהו לקוחות", "לא נמצאו ימי עבודה עם אירוע/לקוח"))

    # 4. Workers without hourly rate (warning)
    for wc in worker_calculations:
        if not wc.get("hourly_rate"):
            alerts.append((2, "עובד", wc["worker_name"], "חסר תעריף שעתי", "לא ניתן לחשב ברוטו ונטו ללא תעריף שעתי."))

    # 5. Workers with unassigned hours (warning)
    for wc in worker_calculations:
        if wc.get("unassigned_hours"):
            alerts.append((2, "עובד", wc["worker_name"], f"שעות ללא שיוך לקוח ({wc['unassigned_hours']})", "ימי עבודה ללא אירוע/לקוח — השעות לא ייחשבו בחיוב ללקוח."))

    # Sort by priority (1=critical first, 3=info last)
    alerts.sort(key=lambda a: a[0])

    attn_row = 3
    for priority, type_label, name, problem, detail in alerts:
        if priority == 1:
            row_fill, row_font = critical_fill, critical_font
            type_prefix = "קריטי"
        elif priority == 2:
            row_fill, row_font = warning_fill, warning_font
            type_prefix = "אזהרה"
        else:
            row_fill, row_font = info_fill_alert, info_font_alert
            type_prefix = "מידע"
        for col_idx, val in enumerate([f"{type_prefix} — {type_label}", name, problem, detail], 1):
            cell = attn_ws.cell(row=attn_row, column=col_idx, value=val)
            cell.fill = row_fill
            cell.font = row_font
            cell.alignment = Alignment(horizontal="right", vertical="center", wrap_text=True)
            cell.border = thin_border
        attn_ws.row_dimensions[attn_row].height = 24
        attn_row += 1

    if attn_row == 3:
        # No alerts — write a green "all good" message
        green_fill = PatternFill(start_color="DCFCE7", end_color="DCFCE7", fill_type="solid")
        attn_ws.merge_cells(start_row=3, start_column=1, end_row=3, end_column=4)
        ok_cell = attn_ws.cell(row=3, column=1, value="הכל תקין — לא נמצאו בעיות")
        ok_cell.font = Font(bold=True, size=13, color="166534")
        ok_cell.fill = green_fill
        ok_cell.alignment = Alignment(horizontal="center", vertical="center")
        for c in range(2, 5):
            attn_ws.cell(row=3, column=c).fill = green_fill
        attn_ws.row_dimensions[3].height = 36
    else:
        # Summary bar at the bottom
        attn_row += 1
        critical_count = sum(1 for a in alerts if a[0] == 1)
        warning_count = sum(1 for a in alerts if a[0] == 2)
        info_count = sum(1 for a in alerts if a[0] == 3)
        summary_txt = f"סה\"כ: {len(alerts)} התראות"
        if critical_count:
            summary_txt += f" | {critical_count} קריטיות"
        if warning_count:
            summary_txt += f" | {warning_count} אזהרות"
        if info_count:
            summary_txt += f" | {info_count} מידע"
        attn_ws.merge_cells(start_row=attn_row, start_column=1, end_row=attn_row, end_column=4)
        sum_cell = attn_ws.cell(row=attn_row, column=1, value=summary_txt)
        sum_cell.font = Font(bold=True, size=11, color="FFFFFF")
        sum_cell.fill = calc_header_fill
        sum_cell.alignment = Alignment(horizontal="center", vertical="center")
        for c in range(2, 5):
            attn_ws.cell(row=attn_row, column=c).fill = calc_header_fill
        attn_ws.row_dimensions[attn_row].height = 30

    for col_idx, width in enumerate([18, 22, 24, 42], 1):
        attn_ws.column_dimensions[get_column_letter(col_idx)].width = width

    wb.save(output_path)


def run_dept_payroll(input_path, output_path, extension, options=None):
    if extension not in {"xls", "xlsx"}:
        raise ValueError("הכלי תומך בקבצי XLS ו-XLSX בלבד")
    options = options or {}
    workbook_kind, workbook = open_excel_workbook(input_path, extension)
    mapping = default_dept_payroll_mapping()
    mapping.update({key: value for key, value in options.items() if key.endswith("_source")})

    worker_rows = []
    for detail_sheet, summary_sheet in iter_flamingo_worker_blocks(workbook_kind, workbook):
        worker_rows.append(extract_dept_payroll_worker(detail_sheet, summary_sheet, workbook_kind, mapping))

    # Parse department settings from options
    dept_settings_json = options.get("dept_settings_json", "[]")
    dept_settings = parse_dept_settings_json(dept_settings_json)

    # Apply manual overrides for hourly rate and housing
    # Empty dropdown = manual mode (user entered a fixed value)
    # Non-empty dropdown = report field (value extracted per worker)
    hourly_rate_source = mapping.get("hourly_rate_source", "")
    if not hourly_rate_source:
        # Manual mode — apply fixed rate to all workers
        try:
            manual_hourly_rate = float(str(options.get("manual_hourly_rate", 0) or 0).replace(",", ""))
        except (ValueError, TypeError):
            manual_hourly_rate = 0
        if manual_hourly_rate:
            for w in worker_rows:
                w["hourly_rate"] = manual_hourly_rate

    housing_source = mapping.get("housing_source", "__auto__")
    if not housing_source:
        # Manual mode — apply fixed housing to all workers
        try:
            manual_housing = float(str(options.get("manual_housing", 0) or 0).replace(",", ""))
        except (ValueError, TypeError):
            manual_housing = 0
        if manual_housing:
            for w in worker_rows:
                w["housing_charge"] = manual_housing

    # Parse clients Excel if provided
    clients_info = ({}, {})
    clients_file_path = options.get("clients_file_path", "")
    if clients_file_path and os.path.exists(clients_file_path):
        clients_ext = get_extension(clients_file_path)
        if clients_ext in ("xls", "xlsx"):
            try:
                clients_info = parse_clients_excel(clients_file_path, clients_ext)
            except Exception:
                pass  # silently fallback if clients file is invalid

    # Use v2 output if we have daily rows / client data
    has_daily_data = any(w.get("daily_rows") or w.get("special_movements") for w in worker_rows)
    if has_daily_data or clients_info[0]:
        write_dept_payroll_output_v2(output_path, worker_rows, dept_settings, clients_info)
    else:
        write_dept_payroll_output(output_path, worker_rows, dept_settings)

    warnings = []
    problem_count = sum(1 for w in worker_rows if w["status"] != "OK")
    if problem_count:
        warnings.append(f"{problem_count} עובדים עם בעיות בנתונים — ראו גיליון 'שימו לב'.")
    no_clients = sum(1 for w in worker_rows if w["status"] == "OK" and not w.get("daily_rows") and not w.get("special_movements"))
    if no_clients and has_daily_data:
        warnings.append(f"{no_clients} עובדים ללא שיוך לקוח — ראו גיליון 'שימו לב'.")
    return {"warnings": warnings}


app = Flask(__name__)
app.secret_key = os.environ.get("SECRET_KEY", "zman_emet_secret_2024")
app.config["MAX_CONTENT_LENGTH"] = MAX_UPLOAD_SIZE

BASE_DIR = Path(__file__).resolve().parent
DB = str(BASE_DIR / "platform.db")
UPLOAD_FOLDER = BASE_DIR / "uploads"
OUTPUT_FOLDER = BASE_DIR / "outputs"
UPLOAD_FOLDER.mkdir(exist_ok=True)
OUTPUT_FOLDER.mkdir(exist_ok=True)


class DatabaseConnection:
    def __init__(self):
        self.conn = None
        self.is_postgres = bool(DATABASE_URL)

    def __enter__(self):
        if self.is_postgres:
            if psycopg is None:
                raise RuntimeError("PostgreSQL driver is not installed")
            self.conn = psycopg.connect(DATABASE_URL, row_factory=dict_row)
        else:
            self.conn = sqlite3.connect(DB)
            self.conn.row_factory = sqlite3.Row
        return self

    def __exit__(self, exc_type, exc, tb):
        if self.conn is None:
            return False
        try:
            if exc_type is not None:
                self.conn.rollback()
        finally:
            self.conn.close()
        return False

    def execute(self, sql, params=()):
        if self.conn is None:
            raise RuntimeError("Database connection is not open")
        if self.is_postgres:
            sql = sql.replace("?", "%s").replace("AUTOINCREMENT", "").replace("INTEGER PRIMARY KEY", "SERIAL PRIMARY KEY")
            sql = sql.replace("INSERT OR IGNORE INTO permissions(user_id,script_id) VALUES (%s,%s)", "INSERT INTO permissions(user_id,script_id) VALUES (%s,%s) ON CONFLICT DO NOTHING")
        return self.conn.execute(sql, params)

    def commit(self):
        if self.conn is not None:
            self.conn.commit()


def get_db():
    return DatabaseConnection()


def is_integrity_error(exc):
    if isinstance(exc, sqlite3.IntegrityError):
        return True
    if psycopg is not None and isinstance(exc, psycopg.IntegrityError):
        return True
    return False


FLOW_LANGUAGES = {"en", "he"}

FLOW_TEXTS = {
    "en": {
        "topbar_greeting": "Hello, ",
        "logout": "Logout",
        "login_page_title": "Login",
        "login_error": "Incorrect username or password",
        "login_subtitle": "Attendance and payroll tools portal",
        "login_username": "Username",
        "login_password": "Password",
        "login_submit": "Log In",
        "dashboard_page_title": "Your Tools",
        "dashboard_empty": "No tools are available for your account yet",
        "dashboard_greeting": "Welcome, ",
        "dashboard_intro": "Your available tools:",
        "run_access_denied": "You do not have access to this tool",
        "run_extra_file_type_error": "The extra file type is not supported",
        "run_missing_extra_file_error": "A required extra file is missing",
        "run_unexpected_error_prefix": "Unexpected processing error: ",
        "back_arrow": "&#8592;",
        "scripts": {
            "nikuy": {
                "name": "Attendance Report Cleanup",
                "desc": "Clean a detailed monthly attendance report by removing asterisks and question marks",
                "success_title": "File is ready",
                "success_action": "Download cleaned file",
                "retry_action": "Process another file",
                "submit_label": "Run",
                "back_label": "Back to tools",
                "empty_error": "No file selected",
                "unsupported_error": "Unsupported file type",
                "invalid_error": "The uploaded file is not a valid Excel file",
                "empty_file_error": "The uploaded file is empty",
                "too_large_error": "The uploaded file is too large",
                "processing_error": "The uploaded file could not be processed",
                "processing_title": "File is being processed",
                "processing_note": "Preparing the cleaned report may take a few minutes. Please keep this page open.",
                "file_picker_label": "Choose file",
            },
            "flamingo_payroll": {
                "name": "Fixed Hourly-Rate Payroll Summary",
                "desc": "Generate a clear monthly payroll summary based on a detailed report and a fixed hourly rate",
                "success_title": "Payroll file is ready",
                "success_action": "Download payroll summary",
                "retry_action": "Process another payroll file",
                "submit_label": "Create payroll summary",
                "back_label": "Back to tools",
                "empty_error": "No file selected",
                "unsupported_error": "Please upload the original Flamingo XLS export",
                "invalid_error": "The uploaded file is not a valid Excel file",
                "empty_file_error": "The uploaded file is empty",
                "too_large_error": "The uploaded file is too large",
                "processing_error": "Could not generate a payroll summary from this file",
                "processing_title": "Payroll summary is being prepared",
                "processing_note": "The system is calculating payable hours and salary for all employees. This may take a few minutes.",
                "file_picker_label": "Choose Flamingo file",
            },
            "matan_missing": {
                "name": "Missing Hours Report",
                "desc": "Filter employees by their monthly missing-hours totals with a clear summary ready for follow-up",
                "success_title": "Missing-hours report is ready",
                "success_action": "Download report",
                "retry_action": "Process another file",
                "submit_label": "Create report",
                "back_label": "Back to tools",
                "empty_error": "No file selected",
                "unsupported_error": "Please upload the original XLS missing-hours report",
                "invalid_error": "The uploaded file is not a valid Excel file",
                "empty_file_error": "The uploaded file is empty",
                "too_large_error": "The uploaded file is too large",
                "processing_error": "Could not generate the missing-hours report from this file",
                "processing_title": "Report is being prepared",
                "processing_note": "The system is filtering the missing-hours report. This may take a few minutes.",
                "file_picker_label": "Choose missing-hours report",
                "filter_fields": [
                    {"label": "Minimum missing hours", "placeholder": "For example 4"},
                    {"label": "Maximum missing hours", "placeholder": "For example 8"},
                ],
            },
        },
    },
    "he": {
        "topbar_greeting": "שלום, ",
        "logout": "התנתקות",
        "login_page_title": "כניסה",
        "login_error": "שם המשתמש או הסיסמה שגויים",
        "login_subtitle": "פורטל כלי נוכחות ושכר",
        "login_username": "שם משתמש",
        "login_password": "סיסמה",
        "login_submit": "כניסה",
        "dashboard_page_title": "הכלים שלך",
        "dashboard_empty": "עדיין אין כלים זמינים לחשבון שלך",
        "dashboard_greeting": "ברוך/ה הבא/ה, ",
        "dashboard_intro": "",
        "run_access_denied": "אין לך גישה לכלי הזה",
        "run_extra_file_type_error": "סוג הקובץ הנוסף אינו נתמך",
        "run_missing_extra_file_error": "חסר קובץ נוסף נדרש",
        "run_unexpected_error_prefix": "שגיאת עיבוד לא צפויה: ",
        "back_arrow": "&#8594;",
        "scripts": {
            "nikuy": {
                "name": "ניקוי דוח נוכחות",
                "desc": "ניקוי סימני כוכביות וסימני שאלה מדוח נוכחות מפורט חודשי",
                "help_label": "דרישות לקובץ",
                "help_title": "מה צריך להעלות?",
                "help_intro": "יש להעלות דוח נוכחות מפורט חודשי. המערכת תנקה אוטומטית סימני * וגם סימני ? מתוך הדוח.",
                "help_items": ["אפשר להמשיך לעבוד בפלטפורמה בזמן שהדוח בעיבוד", "כשהדוח מוכן הוא יופיע בדוחות המוכנים להורדה", "העיבוד עשוי להימשך מספר דקות וזה תקין גם בדוחות גדולים"],
                "help_note": "דוחות בעיבוד ברקע נשמרים עד להורדה או עד 3 ימים, המוקדם מביניהם.",
                "success_title": "הקובץ מוכן",
                "success_action": "הורדת הקובץ הנקי",
                "retry_action": "עיבוד קובץ נוסף",
                "submit_label": "הפעל",
                "back_label": "חזרה לכלים",
                "empty_error": "לא נבחר קובץ",
                "unsupported_error": "סוג הקובץ אינו נתמך",
                "invalid_error": "הקובץ שהועלה אינו קובץ אקסל תקין",
                "empty_file_error": "הקובץ שהועלה ריק",
                "too_large_error": "הקובץ שהועלה גדול מדי",
                "processing_error": "לא ניתן היה לעבד את הקובץ שהועלה",
                "processing_title": "הקובץ נמצא בעיבוד",
                "processing_note": "הכנת הדוח הנקי עשויה להימשך כמה דקות. אפשר להמשיך לעבוד בפלטפורמה, וכשהדוח יהיה מוכן הוא יופיע בדוחות המוכנים להורדה.",
                "file_picker_label": "בחירת קובץ",
            },
            "flamingo_payroll": {
                "name": "סיכום שכר לפי תעריף שעתי קבוע",
                "desc": "סיכום שכר חכם מתוך דוח מפורט חודשי, עם אישור שדות לפני חישוב",
                "help_label": "דרישות לקובץ",
                "help_title": "מה צריך להעלות?",
                "help_intro": "יש להעלות דוח מפורט חודשי הכולל אזור סיכום שממנו אפשר למשוך את שעות התשלום.",
                "help_items": ["המערכת מזהה שדות שכר מרכזיים מתוך הדוח", "מבקשת אישור שדות לפני הרצת החישוב", "ומפיקה סיכום שכר ברור לפי העובדים שנקלטו"],
                "help_note": "אפשר לעבוד גם עם דוחות מסודרים וגם עם דוחות פחות מסודרים, כל עוד הלקוח מאשר את השדות הנכונים.",
                "rules_label": "איך הסקריפט מחשב",
                "rules_title": "מה חשוב לאשר לפני חישוב השכר?",
                "rules_intro": "הסקריפט מחשב את השכר לפי שני שדות קריטיים שחייבים להיות נכונים:",
                "rules_items": ["תעריף שעתי - אפשר למשוך מהדוח או להזין ידנית", "שעות לתשלום בפועל - הלקוח בוחר את השדה שממנו יחושב השכר", "אם מוזן תעריף ידני, כל העובדים בדוח יחושבו לפי אותו תעריף", "שדות כמו נוכחות, תקן וחוסר משמשים לבקרה ולהצלבה"],
                "rules_note": "לפני ההרצה המערכת תציע זיהוי אוטומטי, אבל הלקוח הוא זה שמאשר את השדות הקריטיים.",
                "success_title": "קובץ השכר מוכן",
                "success_action": "הורדת סיכום השכר",
                "retry_action": "עיבוד קובץ שכר נוסף",
                "submit_label": "יצירת סיכום שכר",
                "back_label": "חזרה לכלים",
                "empty_error": "לא נבחר קובץ",
                "unsupported_error": "יש להעלות דוח מפורט חודשי בפורמט Excel",
                "invalid_error": "הקובץ שהועלה אינו קובץ אקסל תקין",
                "empty_file_error": "הקובץ שהועלה ריק",
                "too_large_error": "הקובץ שהועלה גדול מדי",
                "processing_error": "לא ניתן היה ליצור סיכום שכר מהקובץ הזה",
                "processing_title": "סיכום השכר בהכנה",
                "processing_note": "המערכת מחשבת את שעות התשלום והשכר לכל העובדים. זה עשוי להימשך כמה דקות.",
                "file_picker_label": "בחירת קובץ דוח מפורט חודשי",
            },
            "matan_missing": {
                "name": "דוח חוסר מול תקן",
                "desc": "אפשרות לסינון עובדים לפי כמות שעות החוסר החודשיות שלהם כולל סיכום ברור ונוח לטיפול",
                "help_label": "דרישות לקובץ",
                "help_title": "מה צריך להעלות?",
                "help_intro": "יש להעלות דוח מרוכז של חוסר מול תקן. לפני ההרצה המערכת תזהה את השדות ותבקש אישור.",
                "help_items": ["המערכת מזהה שדות חובה ושדות משלימים מתוך הדוח", "מבקשת אישור שדות לפני יצירת הדוח", "ומציגה רק את העובדים שעומדים בתנאי הסינון שהוגדרו"],
                "help_note": "אם הלקוח בוחר שדות משלימים כמו חופשה, מחלה, היעדרות או ש.נוכחות, הם יופיעו גם בפלט.",
                "rules_label": "איך הסקריפט מחשב",
                "rules_title": "מה חשוב לאשר לפני יצירת הדוח?",
                "rules_intro": "הסקריפט לא מחשב מחדש את החוסר. הוא משתמש בערכים שכבר קיימים בדוח ומסנן לפיהם.",
                "rules_items": ["השדות הקריטיים הם חוסר וש.תקן, ולכן חשוב לאשר שהם ממופים נכון", "יש לבחור שם עובד ולפחות מזהה נוסף אחד: מספר עובד, תעודת זהות, מספר תג או דרכון", "שדות כמו ש.נוכחות, חופשה, מחלה והיעדרות הם שדות משלימים בלבד", "היעדרות אינה חוסר. הדוח משתמש בשדה החוסר כפי שהוא מופיע בקובץ המקור"],
                "rules_note": "לפני ההרצה המערכת תציע זיהוי אוטומטי, אבל הלקוח הוא זה שמאשר את השדות הסופיים.",
                "success_title": "דוח שעות החסר מוכן",
                "success_action": "הורדת הדוח",
                "retry_action": "עיבוד קובץ נוסף",
                "submit_label": "יצירת דוח",
                "back_label": "חזרה לכלים",
                "empty_error": "לא נבחר קובץ",
                "unsupported_error": "יש להעלות את דוח שעות החסר המקורי מסוג XLS",
                "invalid_error": "הקובץ שהועלה אינו קובץ אקסל תקין",
                "empty_file_error": "הקובץ שהועלה ריק",
                "too_large_error": "הקובץ שהועלה גדול מדי",
                "processing_error": "לא ניתן היה ליצור את דוח שעות החסר מהקובץ הזה",
                "processing_title": "הדוח בהכנה",
                "processing_note": "המערכת מסננת את דוח שעות החסר. זה עשוי להימשך כמה דקות.",
                "file_picker_label": "בחירת דוח חוסר מול תקן",
                "filter_fields": [
                    {"label": "מינימום שעות חסר", "placeholder": "לדוגמה 4"},
                    {"label": "מקסימום שעות חסר", "placeholder": "לדוגמה 8"},
                ],
            },
        },
    },
}


def get_flow_lang():
    requested = request.args.get("lang")
    if requested in FLOW_LANGUAGES:
        session["flow_lang"] = requested
    saved = session.get("flow_lang", "he")
    return saved if saved in FLOW_LANGUAGES else "he"


def get_flow_dir(lang):
    return "rtl" if lang == "he" else "ltr"


def get_flow_text(lang):
    return FLOW_TEXTS.get(lang, FLOW_TEXTS["en"])


def get_localized_script(script, lang):
    localized = dict(script)
    script_texts = get_flow_text(lang)["scripts"].get(script["id"], {})
    localized.update({k: v for k, v in script_texts.items() if k not in {"filter_fields", "extra_uploads"}})
    if script.get("filter_fields"):
        localized_fields = []
        field_overrides = script_texts.get("filter_fields", [])
        for index, field in enumerate(script.get("filter_fields", [])):
            merged = dict(field)
            if index < len(field_overrides):
                merged.update(field_overrides[index])
            localized_fields.append(merged)
        localized["filter_fields"] = localized_fields
    if script.get("extra_uploads"):
        localized_uploads = []
        upload_overrides = script_texts.get("extra_uploads", [])
        for index, upload in enumerate(script.get("extra_uploads", [])):
            merged = dict(upload)
            if index < len(upload_overrides):
                merged.update(upload_overrides[index])
            localized_uploads.append(merged)
        localized["extra_uploads"] = localized_uploads
    return localized


def build_lang_switch(lang):
    current_path = request.path
    buttons = []
    for code, label in (("en", "EN"), ("he", "HE")):
        cls = "lang-btn active" if code == lang else "lang-btn"
        buttons.append('<a href="' + current_path + '?lang=' + code + '" class="' + cls + '">' + label + "</a>")
    return '<div class="lang-switch">' + "".join(buttons) + "</div>"

SCRIPTS = {
    "nikuy": {
        "id": "nikuy",
        "name": "ניקוי דוח נוכחות",
        "desc": "ניקוי סימני כוכביות וסימני שאלה מדוח נוכחות מפורט חודשי",
        "accept": ".xls,.xlsx",
        "icon": "🧹",
    }
}

SCRIPT_REGISTRY = {
    "nikuy": {
        **SCRIPTS["nikuy"],
        "processor": run_attendance_cleanup,
        "output_suffix": "cleaned",
        "background_queue": True,
        "success_title": "File is ready",
        "success_action": "Download cleaned file",
        "retry_action": "Process another file",
        "submit_label": "Run",
        "back_label": "Back to tools",
        "empty_error": "No file selected",
        "unsupported_error": "Unsupported file type",
        "invalid_error": "The uploaded file is not a valid Excel file",
        "empty_file_error": "The uploaded file is empty",
        "too_large_error": "The uploaded file is too large",
        "processing_error": "The uploaded file could not be processed",
        "processing_title": "File is being processed",
        "processing_note": "Preparing the cleaned report may take a few minutes. Please keep this page open.",
        "file_picker_label": "Choose file",
    }
}

SCRIPTS["flamingo_payroll"] = {
    "id": "flamingo_payroll",
    "name": "סיכום שכר לפי תעריף שעתי קבוע",
    "desc": "הפקת סיכום שכר לפי דוח מפורט חודשי בצורה ברורה ומוכנה לבדיקה",
    "help_label": "מה הסקריפט עושה",
    "help_title": "מה צריך להעלות?",
    "help_intro": "יש להעלות דוח מפורט חודשי הכולל אזור סיכום שממנו אפשר למשוך את שעות התשלום.",
    "help_items": ["המערכת מזהה שדות שכר מרכזיים מתוך הדוח", "מבקשת אישור שדות לפני הרצת החישוב", "ומפיקה סיכום שכר ברור לפי העובדים שנקלטו"],
    "help_note": "אפשר לעבוד גם עם דוחות מסודרים וגם עם דוחות פחות מסודרים, כל עוד הלקוח מאשר את השדות הנכונים.",
    "rules_label": "איך הסקריפט מחשב",
    "rules_title": "מה חשוב לאשר לפני חישוב השכר?",
    "rules_intro": "הסקריפט מחשב את השכר לפי שני שדות קריטיים שחייבים להיות נכונים:",
    "rules_items": [
        "תעריף שעתי - אפשר למשוך מהדוח או להזין ידנית",
        "שעות לתשלום בפועל - הלקוח בוחר את השדה שממנו יחושב השכר",
        "אם מוזן תעריף ידני, כל העובדים בדוח יחושבו לפי אותו תעריף",
        "שדות כמו נוכחות, תקן וחוסר משמשים לבקרה ולהצלבה",
    ],
    "rules_note": "לפני ההרצה המערכת תציע זיהוי אוטומטי, אבל הלקוח הוא זה שמאשר את השדות הקריטיים.",
    "accept": ".xls,.xlsx",
    "icon": "$",
}

SCRIPT_REGISTRY["flamingo_payroll"] = {
    **SCRIPTS["flamingo_payroll"],
    "processor": run_flamingo_payroll,
    "output_suffix": "payment_report",
    "requires_mapping_confirmation": True,
    "success_title": "Payroll file is ready",
    "success_action": "Download payroll summary",
    "retry_action": "Process another payroll file",
    "submit_label": "Create payroll summary",
    "back_label": "Back to tools",
    "empty_error": "No file selected",
    "unsupported_error": "Please upload a monthly detailed Excel report",
    "invalid_error": "The uploaded file is not a valid Excel file",
    "empty_file_error": "The uploaded file is empty",
    "too_large_error": "The uploaded file is too large",
    "processing_error": "Could not generate a payroll summary from this file",
    "processing_title": "Payroll summary is being prepared",
    "processing_note": "The system is calculating payable hours and salary for all employees. This may take a few minutes.",
    "file_picker_label": "Choose monthly detailed report",
}

SCRIPTS["matan_missing"] = {
    "id": "matan_missing",
    "name": "דוח חוסר מול תקן",
    "desc": "אפשרות לסינון עובדים לפי כמות שעות החוסר החודשיות שלהם כולל סיכום ברור ונוח לטיפול",
    "help_label": "דרישות לקובץ",
    "help_title": "מה צריך להעלות?",
    "help_intro": "יש להעלות דוח מרוכז של חוסר מול תקן. לפני ההרצה המערכת תזהה את השדות ותבקש אישור.",
    "help_items": ["המערכת מזהה שדות חובה ושדות משלימים מתוך הדוח", "מבקשת אישור שדות לפני יצירת הדוח", "ומציגה רק את העובדים שעומדים בתנאי הסינון שהוגדרו"],
    "help_note": "אם הלקוח בוחר שדות משלימים כמו חופשה, מחלה, היעדרות או ש.נוכחות, הם יופיעו גם בפלט.",
    "rules_label": "איך הסקריפט מחשב",
    "rules_title": "מה חשוב לאשר לפני יצירת הדוח?",
    "rules_intro": "הסקריפט לא מחשב מחדש את החוסר. הוא משתמש בערכים שכבר קיימים בדוח ומסנן לפיהם.",
    "rules_items": [
        "השדות הקריטיים הם חוסר וש.תקן, ולכן חשוב לאשר שהם ממופים נכון",
        "יש לבחור שם עובד ולפחות מזהה נוסף אחד: מספר עובד, תעודת זהות, מספר תג או דרכון",
        "שדות כמו ש.נוכחות, חופשה, מחלה והיעדרות הם שדות משלימים בלבד",
        "היעדרות אינה חוסר. הדוח משתמש בשדה החוסר כפי שהוא מופיע בקובץ המקור",
    ],
    "rules_note": "לפני ההרצה המערכת תציע זיהוי אוטומטי, אבל הלקוח הוא זה שמאשר את השדות הסופיים.",
    "accept": ".xls",
    "icon": "📊",
}

SCRIPTS["inactive_workers"] = {
    "id": "inactive_workers",
    "name": "איתור עובדים לא פעילים",
    "desc": "איתור עובדים שלא זוהתה אצלם פעילות בטווח הימים או החודשים האחרונים מתוך דוח יומי",
    "help_label": "דרישות לקובץ",
    "help_title": "מה צריך להעלות?",
    "help_intro": "יש להעלות דוח יומי מתאריך עד תאריך, רצוי של לפחות 3 חודשים.",
    "help_items": ["המערכת מזהה את שדות הפעילות מתוך הדוח", "מבקשת אישור שדות לפני יצירת הדוח", "ובודקת אם לעובד הייתה פעילות בטווח הימים או החודשים האחרונים"],
    "help_note": "פעילות יכולה להיחשב לפי כניסה ויציאה יחד, או לפי שדה סה\"כ שעות. אפשר גם לבחור אירוע כשדה משלים.",
    "rules_label": "איך הסקריפט מחשב",
    "rules_title": "מה נחשב פעילות של עובד?",
    "rules_intro": "הסקריפט בודק אם לעובד הייתה פעילות בטווח שנבחר לפי תאריך הייחוס האחרון שקיים בקובץ.",
    "rules_items": [
        "שדות חובה: שם עובד, תאריך ולפחות מזהה עובד אחד נוסף",
        "כדי לזהות פעילות יש לבחור או כניסה ויציאה יחד, או שדה סה\"כ שעות",
        "אירוע הוא שדה אופציונלי, ואם הוא נבחר הוא יכול להיחשב גם הוא כפעילות",
        "אם לא זוהתה לעובד פעילות בכלל, הוא יסומן כלא פעיל ויוצג שלא קיים מידע על פעילות קודמת",
        "אם זוהתה פעילות בעבר אך לא בטווח שנבדק, יוצג התאריך האחרון שבו זוהתה פעילות",
    ],
    "rules_note": "יום חסר או מצב בלי כניסה ובלי יציאה אינם נחשבים פעילות.",
    "accept": ".xls,.xlsx",
    "icon": "🕵️",
}

SCRIPTS["matan_manual_corrections"] = {
    "id": "matan_manual_corrections",
    "name": "דוח תיקונים ידניים",
    "desc": "איתור וסיכום של תיקוני נוכחות ידניים מתוך הדוח, כולל ספירה לפי כניסות, יציאות ומחלקות",
    "help_label": "דרישות לקובץ",
    "help_title": "מה צריך להעלות?",
    "help_intro": "יש להעלות דוח מפורט חודשי הכולל תיקוני כניסה ויציאה ידניים.",
    "help_items": [
        "המערכת מזהה אוטומטית את עמודות הכניסה והיציאה מתוך הדוח",
        "מזהה תיקונים לפי סימן כוכבית (*) לפני ערך הזמן",
        "מזהה שם עובד, ת.ז., מספר תג, מספר שכר ומחלקה מפרטי הגיליון",
        "מפיקה סיכום לכל עובד ולשונית סיכום נפרדת לפי מחלקות",
    ],
    "help_note": "הכלי תומך בפורמטים שונים של דוח מפורט חודשי — אין צורך לוודא עמודות ידנית",
    "rules_label": "איך הכלי מזהה תיקונים",
    "rules_title": "מה נחשב תיקון ידני?",
    "rules_intro": "הכלי מחפש כוכבית (*) לפני ערך הזמן בשדות הכניסה והיציאה.",
    "rules_items": [
        "תיקון כניסה = ערך כניסה שמתחיל בסימן * כגון *08:30",
        "תיקון יציאה = ערך יציאה שמתחיל בסימן * כגון *17:00",
        "לכל עובד נספרים תיקוני כניסה ותיקוני יציאה בנפרד",
        "לשונית מחלקות מציגה סיכום תיקונים לפי מחלקה וממוצע לעובד",
        "הפירוט היומי כולל את ערכי הזמן המתוקנים בפועל",
    ],
    "rules_note": "הספירה מבוצעת על כל עובד בנפרד לאורך כל ימי החודש שבדוח",
    "accept": ".xls",
    "icon": "📝",
}

SCRIPTS["rimon_home_office_summary"] = {
    "id": "rimon_home_office_summary",
    "name": "סיכום עבודה מהבית והמשרד",
    "desc": "סיכום חכם של עבודה מהבית, עבודה מהמשרד, היעדרויות, עזיבות ושגיאות מתוך דוח הנוכחות",
    "help_label": "דרישות לקובץ",
    "help_title": "מה צריך להעלות?",
    "help_intro": "יש להעלות דוח מפורט חודשי.",
    "help_items": ["המערכת מזהה אוטומטית ימי עבודה מהבית וימי עבודה מהמשרד", "מזהה היעדרויות, עזיבות ושגיאות בדיווח", "ומסכמת גם שעות תקן ושעות חוסר לכל עובד"],
    "help_note": "הפלט מחזיר סיכום כולל ופירוט יומי ברור, נוח לבדיקה ולטיפול",
    "rules_label": "איך הסקריפט מחשב",
    "rules_title": "איך הסקריפט מחשב את הימים?",
    "rules_intro": "הסקריפט מסכם כל יום לפי כללי ההכרעה הבאים:",
    "rules_items": [
        "אם יש גם עבודה מהבית וגם עבודה מהמשרד באותו יום, עבודה מהמשרד גוברת",
        "אם יש גם נוכחות וגם היעדרות באותו יום, נוכחות גוברת",
        "אם יש שעת כניסה ושעת יציאה, היום נחשב לנוכחות",
        "אם יש רק שעת כניסה או רק שעת יציאה, היום מסומן כשגיאה מסוג חסר דיווח",
        "אם אין נוכחות ואין אירוע, היום מסומן כשגיאה מסוג יום חסר",
        "אם מזוהה אירוע פיטורין, היום נספר בקטגוריית עובדים שעזבו ולא כהיעדרות",
        "הפלט כולל גם סיכום של שעות תקן ושעות חוסר לכל עובד",
    ],
    "rules_note": "כך אפשר להבין בדיוק איך כל יום מסווג ואיך הסיכומים מחושבים בפלט הסופי.",
    "accept": ".xls,.xlsx",
    "icon": "🏠",
}

SCRIPTS["attendance_alerts"] = {
    "id": "attendance_alerts",
    "name": "התראות נוכחות",
    "desc": "זיהוי חריגות שעות ויציאות מאוחרות מתוך דוח נוכחות חודשי — עם סימון צבעוני אוטומטי",
    "help_label": "דרישות לקובץ",
    "help_title": "מה צריך להעלות?",
    "help_intro": "יש להעלות דוח מפורט חודשי (לשונית לכל עובד).",
    "help_items": [
        "המערכת מזהה אוטומטית יציאות מאוחרות (אחרי 20:00 ו-21:00)",
        "מזהה ימים ארוכים (מעל 8, 9 ו-12 שעות)",
        "מזהה שבועות עם מעל 40 שעות עבודה",
        "כל חריגה מסומנת בצבע לפי חומרה",
    ],
    "help_note": "הפלט כולל לשונית סיכום כוללת ולשונית לכל עובד עם שורות צבעוניות.",
    "rules_label": "כללי ההתראות",
    "rules_title": "מתי נוצרת התראה?",
    "rules_intro": "המערכת בודקת כל יום עבודה לפי הכללים הבאים:",
    "rules_items": [
        "יציאה אחרי 20:00 — כתום (חומרה בינונית)",
        "יציאה אחרי 21:00 — אדום (חומרה גבוהה, גובר על הכתום)",
        "סה\"כ שעות ביום מעל 8 — צהוב בהיר",
        "סה\"כ שעות ביום מעל 9 — כתום",
        "סה\"כ שעות ביום מעל 12 — אדום (חומרה קריטית)",
        "סה\"כ שעות בשבוע (ראשון-שבת) מעל 40 — סגול (גבול שמאלי)",
    ],
    "rules_note": "צבע השורה נקבע לפי ההתראה בעלת החומרה הגבוהה ביותר באותו יום.",
    "accept": ".xls,.xlsx",
    "icon": "🚨",
}

SCRIPTS["dept_payroll"] = {
    "id": "dept_payroll",
    "name": "חיוב לקוחות ותשלום עובדים",
    "desc": "מערכת חיוב לקוחות ותשלום עובדים לחברות כוח אדם — מפיקה דוח חיוב ללקוח, תשלום לעובד וסיכום חברה",
    "help_label": "מה הכלי עושה",
    "help_title": "איך זה עובד?",
    "help_intro": "יש להעלות דוח נוכחות חודשי מפורט. אופציונלית — קובץ לקוחות עם תעריפי גביה ופרטי קשר.",
    "help_items": [
        "המערכת מזהה אוטומטית לקוחות מעמודת האירוע ותעריף שעתי מהדוח",
        "חיוב דירה מזוהה אוטומטית מהערות העובד",
        "ניתן להעלות קובץ לקוחות עם תעריפי גביה, כתובות ואנשי קשר",
        "הפלט כולל 3 גיליונות: חיוב ללקוח, תשלום לעובד, סיכום חברה",
    ],
    "help_note": "ניתן לשמור הגדרות כתבנית לשימוש חוזר בכל חודש.",
    "rules_label": "איך המערכת מחשבת",
    "rules_title": "חישוב חיוב ותשלום",
    "rules_intro": "כל עובד משויך ללקוחות לפי עמודת האירוע. החישוב:",
    "rules_items": [
        "ברוטו = שעות לתשלום × תעריף שעתי (מהדוח או מהגדרות)",
        "3% מיסים = ברוטו × 0.03",
        "ביטוח = 280 ₪ קבוע",
        "נטו = ברוטו − מיסים − ביטוח − חיוב דירה",
        "חיוב ללקוח = שעות בלקוח × תעריף גביה",
    ],
    "rules_note": "עובדים ללא שיוך לקוח יופיעו בגיליון 'שימו לב' לטיפול.",
    "accept": ".xls,.xlsx",
    "icon": "🏢",
}

SCRIPTS["org_hierarchy_report"] = {
    "id": "org_hierarchy_report",
    "name": "תרשים מבנה ארגוני",
    "desc": "הפקת תרשים מבנה ארגוני ודוחות סיכום לפי מנהלים, מחלקות ומבנה הדיווח בארגון, כולל פלט אקסל ו-PowerPoint",
    "help_label": "מה הסקריפט עושה",
    "help_title": "מה צריך להעלות?",
    "help_intro": "יש להעלות קובץ מבנה ארגוני בפורמט CSV.",
    "help_items": ["המערכת מזהה שדות מרכזיים מתוך הקובץ", "מבקשת אישור שדות לפני יצירת הדוח", "ומפיקה דוח מבנה ארגוני עם פלט אקסל, PowerPoint או שניהם יחד"],
    "help_note": "אפשר לעבוד גם עם קבצים מסודרים וגם עם קבצים פחות מסודרים, כל עוד מאשרים את השדות הנכונים.",
    "rules_label": "איך הסקריפט בונה את ההיררכיה",
    "rules_title": "מה חשוב לאשר לפני יצירת הדוח?",
    "rules_intro": "הסקריפט בונה את ההיררכיה לפי שלושה שדות קריטיים שחייבים להיות נכונים:",
    "rules_items": [
        "שם עובד - זהו הצומת שמופיע בהיררכיה",
        "מנהל ישיר - זהו הקשר שקובע למי כל עובד מדווח",
        "מחלקה - משמשת לסיכומים, לקיבוץ ולהבנת המבנה הארגוני",
        "מומלץ לבחור גם מזהה נוסף לעובד: מספר עובד, תעודת זהות או דרכון",
        "שדה מנהל הוא אופציונלי בלבד. אם עובד מוגדר כמנהל ישיר של אחרים, הוא מזוהה כמנהל גם בלי סימון מפורש",
    ],
    "rules_note": "לפני ההרצה המערכת תציע זיהוי אוטומטי, אבל הלקוח הוא זה שמאשר את המיפוי הסופי.",
    "accept": ".csv",
    "icon": "🌳",
}

SCRIPT_REGISTRY["matan_missing"] = {
    **SCRIPTS["matan_missing"],
    "processor": run_matan_missing_filter,
    "output_suffix": "missing_vs_standard_report",
    "requires_mapping_confirmation": True,
    "success_title": "Missing-hours report is ready",
    "success_action": "Download report",
    "retry_action": "Process another file",
    "submit_label": "Create report",
    "back_label": "Back to tools",
    "empty_error": "No file selected",
    "unsupported_error": "Please upload the original XLS missing-hours report",
    "invalid_error": "The uploaded file is not a valid Excel file",
    "empty_file_error": "The uploaded file is empty",
    "too_large_error": "The uploaded file is too large",
    "processing_error": "Could not generate the missing-hours report from this file",
    "processing_title": "Report is being prepared",
    "processing_note": "The system is filtering the missing-hours report. This may take a few minutes.",
    "file_picker_label": "Choose missing-hours report",
    "filter_fields": [
        {"name": "min_missing_hours", "label": "Minimum missing hours", "placeholder": "For example 4"},
        {"name": "max_missing_hours", "label": "Maximum missing hours", "placeholder": "For example 8"}
    ],
}

SCRIPT_REGISTRY["inactive_workers"] = {
    **SCRIPTS["inactive_workers"],
    "processor": run_inactive_workers_report,
    "output_suffix": "inactive_workers_report",
    "requires_mapping_confirmation": True,
    "success_title": "דוח העובדים הלא פעילים מוכן",
    "success_action": "הורדת הדוח",
    "retry_action": "עיבוד קובץ נוסף",
    "submit_label": "יצירת דוח עובדים לא פעילים",
    "back_label": "חזרה לכלים",
    "empty_error": "לא נבחר קובץ",
    "unsupported_error": "יש להעלות דוח יומי מקורי מסוג XLS או XLSX",
    "invalid_error": "הקובץ שהועלה אינו קובץ אקסל תקין",
    "empty_file_error": "הקובץ שהועלה ריק",
    "too_large_error": "הקובץ שהועלה גדול מדי",
    "processing_error": "לא ניתן היה להפיק את דוח העובדים הלא פעילים מהקובץ הזה",
    "processing_title": "דוח העובדים הלא פעילים בהכנה",
    "processing_note": "המערכת בודקת עובדים ללא פעילות בטווח שנבחר. הפעולה יכולה להימשך כמה דקות.",
    "file_picker_label": "בחירת דוח יומי",
}

SCRIPT_REGISTRY["matan_manual_corrections"] = {
    **SCRIPTS["matan_manual_corrections"],
    "processor": run_matan_manual_corrections,
    "output_suffix": "matan_manual_corrections",
    "requires_mapping_confirmation": True,
    "success_title": "דוח התיקונים מוכן",
    "success_action": "הורדת הדוח",
    "retry_action": "עיבוד קובץ נוסף",
    "submit_label": "יצירת דוח תיקונים",
    "back_label": "חזרה לכלים",
    "empty_error": "לא נבחר קובץ",
    "unsupported_error": "יש להעלות את דוח הנוכחות החודשי המפורט המקורי מסוג XLS",
    "invalid_error": "הקובץ שהועלה אינו קובץ אקסל תקין",
    "empty_file_error": "הקובץ שהועלה ריק",
    "too_large_error": "הקובץ שהועלה גדול מדי",
    "processing_error": "לא ניתן היה להפיק את דוח התיקונים מהקובץ הזה",
    "processing_title": "דוח התיקונים בהכנה",
    "processing_note": "המערכת סופרת תיקוני נוכחות ידניים לכל עובד. הפעולה עשויה להימשך כמה דקות.",
    "file_picker_label": "בחירת דוח נוכחות חודשי",
}

SCRIPT_REGISTRY["rimon_home_office_summary"] = {
    **SCRIPTS["rimon_home_office_summary"],
    "processor": run_rimon_home_office_summary,
    "output_suffix": "home_office_report",
    "success_title": "דוח הסיכום מוכן",
    "success_action": "הורדת הדוח",
    "retry_action": "עיבוד קובץ נוסף",
    "submit_label": "יצירת דוח סיכום",
    "back_label": "חזרה לכלים",
    "empty_error": "לא נבחר קובץ",
    "unsupported_error": "יש להעלות דוח מפורט חודשי מסוג XLS או XLSX",
    "invalid_error": "הקובץ שהועלה אינו קובץ אקסל תקין",
    "empty_file_error": "הקובץ שהועלה ריק",
    "too_large_error": "הקובץ שהועלה גדול מדי",
    "processing_error": "לא ניתן היה להפיק את דוח הסיכום מהקובץ הזה",
    "processing_title": "דוח הסיכום בהכנה",
    "processing_note": "המערכת מקבצת תאריכים וסופרת ימי משרד, עבודה מהבית, היעדרות ושגיאות. הפעולה עשויה להימשך כמה דקות.",
    "file_picker_label": "בחירת דוח מפורט חודשי",
    "requires_mapping_confirmation": True,
}

SCRIPT_REGISTRY["attendance_alerts"] = {
    **SCRIPTS["attendance_alerts"],
    "processor": run_attendance_alerts,
    "output_suffix": "attendance_alerts_report",
    "requires_mapping_confirmation": True,
    "success_title": "דוח ההתראות מוכן",
    "success_action": "הורדת הדוח",
    "retry_action": "עיבוד קובץ נוסף",
    "submit_label": "יצירת דוח התראות",
    "back_label": "חזרה לכלים",
    "empty_error": "לא נבחר קובץ",
    "unsupported_error": "יש להעלות דוח מפורט חודשי מסוג XLS או XLSX",
    "invalid_error": "הקובץ שהועלה אינו קובץ אקסל תקין",
    "empty_file_error": "הקובץ שהועלה ריק",
    "too_large_error": "הקובץ שהועלה גדול מדי",
    "processing_error": "לא ניתן היה להפיק את דוח ההתראות מהקובץ הזה",
    "processing_title": "דוח ההתראות בהכנה",
    "processing_note": "המערכת בודקת חריגות שעות ויציאות מאוחרות. הפעולה עשויה להימשך כמה דקות.",
    "file_picker_label": "בחירת דוח מפורט חודשי",
}

SCRIPT_REGISTRY["dept_payroll"] = {
    **SCRIPTS["dept_payroll"],
    "processor": run_dept_payroll,
    "output_suffix": "dept_payroll_report",
    "requires_mapping_confirmation": True,
    "success_title": "דוח השכר לפי מחלקות מוכן",
    "success_action": "הורדת הדוח",
    "retry_action": "עיבוד קובץ נוסף",
    "submit_label": "הפקת דוח שכר לפי מחלקות",
    "back_label": "חזרה לכלים",
    "empty_error": "לא נבחר קובץ",
    "unsupported_error": "יש להעלות דוח מפורט חודשי מסוג XLS או XLSX",
    "invalid_error": "הקובץ שהועלה אינו קובץ אקסל תקין",
    "empty_file_error": "הקובץ שהועלה ריק",
    "too_large_error": "הקובץ שהועלה גדול מדי",
    "processing_error": "לא ניתן היה להפיק את דוח השכר מהקובץ הזה",
    "processing_title": "דוח השכר לפי מחלקות בהכנה",
    "processing_note": "המערכת מחשבת שכר לכל עובד לפי מחלקה. הפעולה עשויה להימשך כמה דקות.",
    "file_picker_label": "בחירת דוח נוכחות מפורט חודשי",
}

SCRIPT_REGISTRY["org_hierarchy_report"] = {
    **SCRIPTS["org_hierarchy_report"],
    "processor": run_org_hierarchy_report,
    "output_suffix": "org_hierarchy_report",
    "requires_mapping_confirmation": True,
    "output_extension": "zip",
    "output_option_name": "output_type",
    "output_extension_map": {"excel": "xlsx", "powerpoint": "pptx", "both": "zip"},
    "success_title": "דוח המבנה הארגוני מוכן",
    "success_action": "הורדת הקובץ",
    "retry_action": "עיבוד קובץ נוסף",
    "submit_label": "יצירת דוח מבנה ארגוני",
    "back_label": "חזרה לכלים",
    "empty_error": "לא נבחר קובץ",
    "unsupported_error": "יש להעלות את קובץ המבנה הארגוני המקורי מסוג CSV",
    "invalid_error": "הקובץ שהועלה אינו קובץ תקין",
    "empty_file_error": "הקובץ שהועלה ריק",
    "too_large_error": "הקובץ שהועלה גדול מדי",
    "processing_error": "לא ניתן היה להפיק את דוח המבנה הארגוני מהקובץ הזה",
    "processing_title": "דוח המבנה הארגוני בהכנה",
    "processing_note": "המערכת בונה דוח אקסל, מצגת או קובץ ZIP לפי הבחירה שלך. הפעולה עשויה להימשך כמה דקות.",
    "file_picker_label": "בחירת קובץ מבנה ארגוני",
    "filter_fields": [
        {
            "name": "output_type",
            "label": "סוג פלט",
            "type": "select",
            "default": "powerpoint",
            "options": [
                {"value": "excel", "label": "אקסל בלבד"},
                {"value": "powerpoint", "label": "PowerPoint בלבד"},
                {"value": "both", "label": "XL+PowerPoint"},
            ],
        },
    ],
}

SCRIPTS = SCRIPT_REGISTRY


def get_script(script_id):
    return SCRIPT_REGISTRY.get(script_id)


def build_output_filename(script, uid, options=None):
    suffix = script.get("output_suffix", "output")
    extension = script.get("output_extension", "xlsx")
    option_name = script.get("output_option_name")
    extension_map = script.get("output_extension_map", {})
    if option_name and options:
        selected = str(options.get(option_name, "")).strip()
        if selected in extension_map:
            extension = extension_map[selected]
    return f"{uid}_{suffix}.{extension}"


def execute_script(script, input_path, output_path, extension, options=None):
    processor = script.get("processor")
    if processor is None:
        raise ValueError("Script processor is not configured")
    return processor(input_path, output_path, extension, options)


CSS = """
* { box-sizing: border-box; margin: 0; padding: 0; }
body { font-family: 'Segoe UI', Arial, sans-serif; background: #f0f4ff; min-height: 100vh; direction: inherit; position: relative; }
body::before { content: ''; position: fixed; inset: 0; background-image: radial-gradient(rgba(37,99,235,.03) 1px, transparent 1px); background-size: 32px 32px; pointer-events: none; z-index: 0; }
.topbar { background: linear-gradient(135deg, #0f1b3d 0%, #1e3a8a 100%); color: white; padding: 0 2rem; height: 62px; display: flex; align-items: center; justify-content: space-between; position: sticky; top: 0; z-index: 50; box-shadow: 0 2px 20px rgba(15,23,42,.15); }
.topbar h1 { font-size: 18px; font-weight: 800; letter-spacing: -0.3px; }
.topbar a { color: #93c5fd; font-size: 13px; text-decoration: none; transition: color .2s; }
.topbar a:hover { color: #ffffff; }
.wrap { max-width: 900px; margin: 2rem auto; padding: 0 1rem; position: relative; z-index: 1; animation: panelFadeIn .5s ease-out; }
.login-wrap { max-width: 420px; margin: 0 auto; padding: 2rem 1rem; position: relative; z-index: 1; animation: panelFadeIn .6s ease-out; min-height: 100vh; display: flex; flex-direction: column; align-items: stretch; justify-content: center; }
.login-bg { min-height: 100vh; background: linear-gradient(160deg, #0f1b3d 0%, #1e3a8a 35%, #2563eb 70%, #3b82f6 100%); position: relative; overflow: hidden; }
.login-bg::before { content: ''; position: absolute; inset: 0; background-image: radial-gradient(rgba(255,255,255,.04) 1px, transparent 1px); background-size: 40px 40px; pointer-events: none; }
.login-bg .login-wrap { padding-top: 0; padding-bottom: 0; }
.login-orb { position: absolute; border-radius: 50%; filter: blur(80px); opacity: .12; pointer-events: none; }
.login-orb-1 { width: 400px; height: 400px; background: #60a5fa; top: -100px; right: -100px; animation: loginOrbFloat 10s ease-in-out infinite; }
.login-orb-2 { width: 300px; height: 300px; background: #a78bfa; bottom: -80px; left: -80px; animation: loginOrbFloat 12s ease-in-out infinite reverse; }
@keyframes loginOrbFloat { 0%,100% { transform: translate(0,0) scale(1); } 50% { transform: translate(30px,-20px) scale(1.05); } }
.card { background: white; border: 1px solid #e2e8f0; border-radius: 18px; box-shadow: 0 4px 24px rgba(37,99,235,.07); padding: 2rem; margin-bottom: 1.5rem; transition: box-shadow .3s cubic-bezier(.4,0,.2,1), border-color .3s, transform .3s cubic-bezier(.4,0,.2,1); }
.card:hover { box-shadow: 0 8px 32px rgba(37,99,235,.1); border-color: #bfdbfe; }
.card h2 { font-size: 17px; font-weight: 800; color: #1e3a8a; margin-bottom: 1rem; padding-bottom: .75rem; border-bottom: 1.5px solid #e0e7ff; }
label.field-label { font-size: 13px; font-weight: 600; color: #374151; margin-bottom: 5px; display: block; }
input[type=text], input[type=password], input[type=email] { padding: 11px 14px; border: 1.5px solid #e2e8f0; border-radius: 12px; font-size: 14px; font-family: inherit; outline: none; width: 100%; margin-bottom: .75rem; transition: border-color .2s, box-shadow .2s; background: #fafcff; }
input[type=text]:focus, input[type=password]:focus, input[type=email]:focus { border-color: #2563eb; box-shadow: 0 0 0 3px rgba(37,99,235,.1); background: #ffffff; }
select { padding: 11px 14px; border: 1.5px solid #e2e8f0; border-radius: 12px; font-size: 14px; font-family: inherit; outline: none; width: 100%; background: #fafcff; transition: border-color .2s, box-shadow .2s; }
select:focus { border-color: #2563eb; box-shadow: 0 0 0 3px rgba(37,99,235,.1); }
textarea { padding: 11px 14px; border: 1.5px solid #e2e8f0; border-radius: 12px; font-size: 14px; font-family: inherit; outline: none; width: 100%; background: #fafcff; transition: border-color .2s, box-shadow .2s; resize: vertical; }
textarea:focus { border-color: #2563eb; box-shadow: 0 0 0 3px rgba(37,99,235,.1); background: #ffffff; }
.btn { padding: 11px 22px; border: none; border-radius: 12px; font-size: 14px; font-weight: 700; cursor: pointer; font-family: inherit; transition: all .25s cubic-bezier(.4,0,.2,1); display: inline-flex; align-items: center; justify-content: center; gap: 6px; }
.btn:active { transform: scale(.97); }
.btn-blue { background: linear-gradient(135deg, #2563eb 0%, #1d4ed8 100%); color: white; box-shadow: 0 2px 8px rgba(37,99,235,.25); }
.btn-blue:hover { background: linear-gradient(135deg, #1d4ed8 0%, #1e40af 100%); box-shadow: 0 4px 16px rgba(37,99,235,.3); transform: translateY(-1px); }
.btn-red { background: #fee2e2; color: #dc2626; border: 1px solid #fecaca; }
.btn-red:hover { background: #fecaca; transform: translateY(-1px); }
.btn-gray { background: #f1f5f9; color: #475569; border: 1px solid #e2e8f0; }
.btn-gray:hover { background: #e2e8f0; transform: translateY(-1px); }
.flash { background: #f0fdf4; border: 1px solid #86efac; color: #15803d; border-radius: 12px; padding: 12px 16px; font-size: 13px; margin-bottom: 1rem; }
.flash-err { background: #fef2f2; border: 1px solid #fecaca; color: #dc2626; border-radius: 12px; padding: 12px 16px; font-size: 13px; margin-bottom: 1rem; }
.flash-stack { position: fixed; top: 92px; right: max(14px, calc((100vw - 900px)/2 - 110px)); z-index: 120; display: flex; flex-direction: column; gap: 10px; width: min(340px, calc(100vw - 28px)); }
.flash-toast { background: #f0fdf4; border: 1px solid #86efac; color: #15803d; border-radius: 14px; padding: 12px 16px; font-size: 13px; line-height: 1.7; box-shadow: 0 14px 34px rgba(15,23,42,.14); opacity: 0; transform: translateY(-8px); animation: flashToastIn .22s ease-out forwards; backdrop-filter: blur(8px); }
.flash-toast.dismiss { animation: flashToastOut .22s ease-in forwards; }
table { width: 100%; border-collapse: separate; border-spacing: 0; font-size: 13px; }
th { text-align: start; padding: 12px 14px; background: linear-gradient(180deg, #f8fafc 0%, #f1f5f9 100%); color: #475569; font-weight: 700; font-size: 12px; text-transform: uppercase; letter-spacing: .5px; border-bottom: 1.5px solid #e2e8f0; }
td { padding: 13px 14px; border-bottom: 1px solid #f1f5f9; vertical-align: middle; transition: background .15s; }
tr:hover td { background: #f8fafc; }
.badge { display: inline-block; padding: 4px 12px; border-radius: 99px; font-size: 11px; font-weight: 700; background: #f1f5f9; color: #64748b; }
.form-row { display: flex; gap: 12px; flex-wrap: wrap; align-items: flex-end; }
.form-group { flex: 1; min-width: 130px; }
.drop-zone { border: 2px dashed #bfdbfe; border-radius: 16px; padding: 2.5rem 2rem; text-align: center; cursor: pointer; background: linear-gradient(180deg, #fafcff 0%, #eff6ff 100%); margin-bottom: 1rem; transition: all .3s cubic-bezier(.4,0,.2,1); }
.drop-zone:hover { border-color: #2563eb; background: linear-gradient(180deg, #eff6ff 0%, #dbeafe 100%); transform: translateY(-2px); box-shadow: 0 8px 24px rgba(37,99,235,.1); }
.success-box { padding: 1.5rem; background: linear-gradient(180deg, #f0fdf4 0%, #ecfdf5 100%); border: 1.5px solid #86efac; border-radius: 16px; text-align: center; margin-top: 1rem; }
.dl-btn { display: inline-block; padding: 13px 32px; background: linear-gradient(135deg, #16a34a 0%, #15803d 100%); color: white; border-radius: 12px; font-size: 15px; font-weight: 700; text-decoration: none; box-shadow: 0 2px 8px rgba(22,163,74,.3); transition: all .25s cubic-bezier(.4,0,.2,1); }
.dl-btn:hover { transform: translateY(-2px); box-shadow: 0 6px 20px rgba(22,163,74,.35); }
.processing-box { display: none; margin-top: 1rem; padding: 1.25rem 1.5rem; background: linear-gradient(180deg, #eff6ff 0%, #dbeafe 100%); border: 1.5px solid #bfdbfe; border-radius: 16px; }
.processing-box.show { display: block; animation: panelFadeIn .3s ease-out; }
.progress-track { width: 100%; height: 10px; border-radius: 999px; background: #dbeafe; overflow: hidden; margin: .9rem 0 .7rem; }
.progress-bar { width: 45%; height: 100%; border-radius: 999px; background: linear-gradient(90deg, #2563eb 0%, #60a5fa 50%, #93c5fd 100%); animation: loadingSlide 1.6s ease-in-out infinite; }
.processing-note { font-size: 14px; color: #1d4ed8; font-weight: 700; }
.processing-subnote { font-size: 12px; color: #64748b; line-height: 1.6; }
.lang-switch { display: inline-flex; align-items: center; gap: 6px; background: rgba(255,255,255,.14); border: 1px solid rgba(255,255,255,.2); border-radius: 999px; padding: 4px; }
.lang-switch.standalone { background: #ffffff; border-color: #dbeafe; box-shadow: 0 4px 16px rgba(37,99,235,.08); margin-bottom: 1rem; }
.lang-btn { display: inline-block; min-width: 38px; padding: 6px 10px; border-radius: 999px; font-size: 12px; font-weight: 700; text-align: center; text-decoration: none; color: #cbd5e1; transition: all .2s; }
.lang-switch.standalone .lang-btn { color: #64748b; }
.lang-btn.active { background: #ffffff; color: #1e3a8a; }
.lang-switch.standalone .lang-btn.active { background: #2563eb; color: #ffffff; }
@keyframes loadingSlide {
  0% { transform: translateX(0); }
  50% { transform: translateX(120%); }
  100% { transform: translateX(0); }
}
@keyframes panelFadeIn {
  from { opacity: 0; transform: translateY(16px); }
  to { opacity: 1; transform: translateY(0); }
}
@keyframes flashToastIn {
  from { opacity: 0; transform: translateY(-8px); }
  to { opacity: 1; transform: translateY(0); }
}
@keyframes flashToastOut {
  from { opacity: 1; transform: translateY(0); }
  to { opacity: 0; transform: translateY(-8px); }
}
.modal-bg { display: none; position: fixed; inset: 0; background: rgba(15,23,42,.5); backdrop-filter: blur(4px); z-index: 100; align-items: center; justify-content: center; }
.modal-box { background: white; border-radius: 20px; padding: 2rem; width: 340px; box-shadow: 0 24px 48px rgba(15,23,42,.2); }
.admin-user-grid { display:grid; grid-template-columns:repeat(auto-fit,minmax(320px,1fr)); gap:16px; }
.admin-user-card { background:linear-gradient(180deg,#ffffff 0%,#f8fbff 100%); border:1px solid #dbeafe; border-radius:18px; padding:18px; box-shadow:0 8px 28px rgba(37,99,235,.08); }
.admin-user-head { display:flex; align-items:flex-start; justify-content:space-between; gap:12px; margin-bottom:14px; }
.admin-user-title { font-size:17px; font-weight:800; color:#0f172a; line-height:1.4; }
.admin-user-sub { font-size:12px; color:#64748b; line-height:1.7; }
.admin-user-status { display:inline-flex; align-items:center; padding:7px 12px; border-radius:999px; font-size:12px; font-weight:800; white-space:nowrap; }
.admin-user-meta { display:grid; grid-template-columns:repeat(2,minmax(0,1fr)); gap:10px; margin-bottom:14px; }
.admin-user-meta-box { background:#ffffff; border:1px solid #e2e8f0; border-radius:14px; padding:12px; }
.admin-user-meta-box .k { font-size:11px; color:#64748b; margin-bottom:4px; }
.admin-user-meta-box .v { font-size:14px; font-weight:700; color:#0f172a; line-height:1.6; word-break:break-word; }
.admin-user-section { margin-top:14px; padding-top:14px; border-top:1px solid #e2e8f0; }
.admin-user-section-title { font-size:13px; font-weight:800; color:#1e3a8a; margin-bottom:10px; }
.admin-user-perms { display:flex; flex-wrap:wrap; gap:8px; }
.admin-user-perms label { display:inline-flex; align-items:center; gap:6px; padding:8px 10px; border-radius:999px; background:#ffffff; border:1px solid #dbeafe; font-size:12px; color:#334155; }
.admin-user-actions { display:flex; gap:8px; flex-wrap:wrap; }
.admin-user-summary { display:grid; grid-template-columns:repeat(auto-fit,minmax(140px,1fr)); gap:10px; margin-bottom:16px; }
.admin-user-summary-box { background:#f8fafc; border:1px solid #e2e8f0; border-radius:14px; padding:12px; }
.admin-user-summary-box .k { font-size:12px; color:#64748b; margin-bottom:4px; }
.admin-user-summary-box .v { font-size:22px; font-weight:800; color:#0f172a; }
.admin-float-nav { position:fixed; top:92px; right:max(8px, calc((100vw - 900px)/2 - 88px)); z-index:20; display:flex; flex-direction:column; gap:8px; }
.admin-float-nav a { box-shadow:0 8px 24px rgba(15,23,42,.12); background:#ffffff; border:1px solid #dbeafe; min-width:72px; padding:9px 10px; font-size:12px; }
.support-request-list { display:flex; flex-direction:column; gap:14px; }
.support-request-card { background:linear-gradient(180deg,#ffffff 0%,#f8fbff 100%); border:1px solid #dbeafe; border-radius:18px; padding:16px; }
.support-request-card-head { display:flex; align-items:flex-start; justify-content:space-between; gap:12px; margin-bottom:12px; flex-wrap:wrap; }
.support-request-card-meta { display:grid; grid-template-columns:repeat(auto-fit,minmax(150px,1fr)); gap:10px; margin-bottom:12px; }
.support-request-card-box { background:#ffffff; border:1px solid #e2e8f0; border-radius:14px; padding:12px; }
.support-request-card-box .k { font-size:11px; color:#64748b; margin-bottom:4px; }
.support-request-card-box .v { font-size:14px; font-weight:700; color:#0f172a; line-height:1.7; word-break:break-word; }
.support-request-message { background:#ffffff; border:1px solid #e2e8f0; border-radius:14px; padding:14px; font-size:13px; color:#334155; line-height:1.9; white-space:pre-wrap; margin-bottom:12px; }
.admin-collapsible-summary { list-style:none; cursor:pointer; padding:18px 20px; display:flex; align-items:center; justify-content:space-between; gap:12px; }
.admin-collapsible-sub { font-size:13px; color:#64748b; line-height:1.7; }
.admin-support-summary { display:grid; grid-template-columns:repeat(auto-fit,minmax(180px,1fr)); gap:10px; margin-bottom:14px; }
.admin-support-summary-box { background:#f8fafc; border:1px solid #e2e8f0; border-radius:14px; padding:12px; }
.admin-support-summary-box .k { font-size:12px; color:#64748b; margin-bottom:4px; }
.admin-support-summary-box .v { font-size:20px; font-weight:800; color:#0f172a; }
@media (max-width: 1280px) { .admin-float-nav { right:10px; } .admin-float-nav a { min-width:64px; font-size:11px; padding:8px 9px; } .flash-stack { right: 10px; } }
@media (max-width: 1100px) { .admin-float-nav { position:static; margin-bottom:1rem; flex-direction:row; flex-wrap:wrap; } .admin-float-nav a { min-width:unset; } .flash-stack { top: 74px; right: 12px; } }
"""


def get_table_columns(db, table_name):
    if db.is_postgres:
        rows = db.execute(
            "SELECT column_name FROM information_schema.columns WHERE table_schema = current_schema() AND table_name = %s",
            (table_name,),
        ).fetchall()
        return {row["column_name"] for row in rows}
    rows = db.execute(f"PRAGMA table_info({table_name})").fetchall()
    return {row["name"] for row in rows}


def init_db():
    with get_db() as db:
        db.execute(
            """CREATE TABLE IF NOT EXISTS users (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            username TEXT UNIQUE NOT NULL,
            password TEXT NOT NULL,
            full_name TEXT,
            is_admin INTEGER DEFAULT 0,
            active INTEGER DEFAULT 1)"""
        )
        db.execute(
            """CREATE TABLE IF NOT EXISTS permissions (
            user_id INTEGER, script_id TEXT,
            PRIMARY KEY (user_id, script_id))"""
        )
        db.execute(
            """CREATE TABLE IF NOT EXISTS activity_logs (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            user_id INTEGER,
            username TEXT,
            full_name TEXT,
            event_type TEXT NOT NULL,
            action_label TEXT NOT NULL,
            script_id TEXT,
            script_name TEXT,
            details TEXT,
            created_at TEXT NOT NULL)"""
        )
        db.execute("CREATE INDEX IF NOT EXISTS idx_activity_logs_created_at ON activity_logs(created_at)")
        db.execute(
            """CREATE TABLE IF NOT EXISTS mapping_templates (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            user_id INTEGER NOT NULL,
            script_id TEXT NOT NULL,
            name TEXT NOT NULL,
            mapping_json TEXT NOT NULL,
            created_at TEXT NOT NULL,
            updated_at TEXT NOT NULL)"""
        )
        db.execute("CREATE INDEX IF NOT EXISTS idx_mapping_templates_user_script ON mapping_templates(user_id, script_id)")
        db.execute(
            """CREATE TABLE IF NOT EXISTS support_requests (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            user_id INTEGER NOT NULL,
            username TEXT,
            full_name TEXT,
            company_name TEXT,
            email TEXT,
            phone TEXT,
            request_type TEXT NOT NULL,
            script_id TEXT,
            script_name TEXT,
            message TEXT NOT NULL,
            status TEXT NOT NULL,
            created_at TEXT NOT NULL)"""
        )
        db.execute("CREATE INDEX IF NOT EXISTS idx_support_requests_created_at ON support_requests(created_at)")
        db.execute(
            """CREATE TABLE IF NOT EXISTS report_jobs (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            user_id INTEGER NOT NULL,
            username TEXT,
            full_name TEXT,
            company_name TEXT,
            script_id TEXT NOT NULL,
            script_name TEXT NOT NULL,
            original_filename TEXT,
            input_path TEXT,
            input_ext TEXT,
            output_filename TEXT,
            status TEXT NOT NULL,
            status_note TEXT,
            created_at TEXT NOT NULL,
            started_at TEXT,
            completed_at TEXT,
            expires_at TEXT,
            downloaded_at TEXT)"""
        )
        db.execute("CREATE INDEX IF NOT EXISTS idx_report_jobs_user_created_at ON report_jobs(user_id, created_at)")
        db.execute("CREATE INDEX IF NOT EXISTS idx_report_jobs_status ON report_jobs(status)")
        # ── Marketplace tables ──────────────────────────────────────
        db.execute(
            """CREATE TABLE IF NOT EXISTS marketplace_tools (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            creator_id INTEGER NOT NULL,
            name TEXT NOT NULL,
            description TEXT,
            icon TEXT DEFAULT '🔧',
            category TEXT DEFAULT 'general',
            definition_json TEXT NOT NULL,
            status TEXT NOT NULL DEFAULT 'draft',
            is_public INTEGER DEFAULT 0,
            usage_count INTEGER DEFAULT 0,
            created_at TEXT NOT NULL,
            updated_at TEXT NOT NULL,
            approved_at TEXT,
            approved_by INTEGER)"""
        )
        db.execute("CREATE INDEX IF NOT EXISTS idx_marketplace_tools_creator ON marketplace_tools(creator_id)")
        db.execute("CREATE INDEX IF NOT EXISTS idx_marketplace_tools_status ON marketplace_tools(status)")
        db.execute(
            """CREATE TABLE IF NOT EXISTS tool_ratings (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            tool_id INTEGER NOT NULL,
            user_id INTEGER NOT NULL,
            rating INTEGER NOT NULL DEFAULT 1,
            created_at TEXT NOT NULL)"""
        )
        db.execute("CREATE UNIQUE INDEX IF NOT EXISTS idx_tool_ratings_unique ON tool_ratings(tool_id, user_id)")
        db.execute(
            """CREATE TABLE IF NOT EXISTS tool_comments (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            tool_id INTEGER NOT NULL,
            user_id INTEGER NOT NULL,
            username TEXT,
            full_name TEXT,
            body TEXT NOT NULL,
            created_at TEXT NOT NULL)"""
        )
        db.execute("CREATE INDEX IF NOT EXISTS idx_tool_comments_tool ON tool_comments(tool_id)")
        db.execute(
            """CREATE TABLE IF NOT EXISTS tool_installs (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            tool_id INTEGER NOT NULL,
            user_id INTEGER NOT NULL,
            created_at TEXT NOT NULL)"""
        )
        db.execute("CREATE UNIQUE INDEX IF NOT EXISTS idx_tool_installs_unique ON tool_installs(tool_id, user_id)")
        db.execute(
            """CREATE TABLE IF NOT EXISTS tool_chat_sessions (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            user_id INTEGER NOT NULL,
            messages_json TEXT NOT NULL DEFAULT '[]',
            tool_id INTEGER,
            status TEXT NOT NULL DEFAULT 'active',
            created_at TEXT NOT NULL,
            updated_at TEXT NOT NULL)"""
        )
        db.execute("CREATE INDEX IF NOT EXISTS idx_tool_chat_sessions_user ON tool_chat_sessions(user_id)")
        # ── Tool requests (developer briefs from AI chat) ─────────
        db.execute(
            """CREATE TABLE IF NOT EXISTS tool_requests (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            user_id INTEGER NOT NULL,
            username TEXT,
            full_name TEXT,
            company_name TEXT,
            tool_name TEXT NOT NULL,
            request_type TEXT NOT NULL DEFAULT 'new',
            existing_tool_name TEXT,
            brief_text TEXT NOT NULL,
            chat_session_id INTEGER,
            status TEXT NOT NULL DEFAULT 'pending',
            admin_notes TEXT,
            created_at TEXT NOT NULL,
            updated_at TEXT NOT NULL)"""
        )
        db.execute("CREATE INDEX IF NOT EXISTS idx_tool_requests_user ON tool_requests(user_id)")
        db.execute("CREATE INDEX IF NOT EXISTS idx_tool_requests_status ON tool_requests(status)")
        # ── Chat token usage (cost tracking) ──────────────────────
        db.execute(
            """CREATE TABLE IF NOT EXISTS chat_token_usage (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            user_id INTEGER NOT NULL,
            session_id INTEGER,
            input_tokens INTEGER NOT NULL DEFAULT 0,
            output_tokens INTEGER NOT NULL DEFAULT 0,
            estimated_cost_usd REAL NOT NULL DEFAULT 0,
            created_at TEXT NOT NULL)"""
        )
        db.execute("CREATE INDEX IF NOT EXISTS idx_chat_token_usage_user ON chat_token_usage(user_id)")
        # ── End marketplace tables ────────────────────────────────

        existing_columns = get_table_columns(db, "users")
        desired_columns = {
            "company_name": "TEXT",
            "company_id": "TEXT",
            "email": "TEXT",
            "phone": "TEXT",
            "join_date": "TEXT",
            "trial_start_date": "TEXT",
            "service_valid_until": "TEXT",
            "billing_mode": "TEXT DEFAULT 'monthly'",
            "trial_days": "INTEGER DEFAULT 30",
        }
        for column_name, column_sql in desired_columns.items():
            if column_name not in existing_columns:
                db.execute(f"ALTER TABLE users ADD COLUMN {column_name} {column_sql}")
        if not db.execute("SELECT id FROM users WHERE username='admin'").fetchone():
            db.execute(
                "INSERT INTO users(username,password,full_name,is_admin) VALUES (?,?,?,1)",
                ("admin", generate_password_hash("admin123"), "מנהל מערכת"),
            )
        db.commit()


init_db()


def add_flash(msg):
    session.setdefault("msgs", []).append(msg)


def pop_flashes():
    msgs = session.pop("msgs", [])
    if not msgs:
        return ""
    return '<div class="flash-stack" id="flashStack">' + "".join('<div class="flash-toast">' + m + "</div>" for m in msgs) + "</div>"


def generate_temp_password(length=10):
    alphabet = string.ascii_letters + string.digits
    return "".join(secrets.choice(alphabet) for _ in range(length))


def esc(value):
    return html.escape(str(value or ""))


def format_ui_datetime(value):
    text = str(value or "").strip()
    if not text:
        return ""
    for fmt in ("%Y-%m-%d %H:%M:%S", "%Y-%m-%dT%H:%M:%S"):
        try:
            return datetime.strptime(text, fmt).strftime("%d/%m/%Y %H:%M")
        except ValueError:
            continue
    return text


def parse_datetime_value(value):
    text = str(value or "").strip()
    if not text:
        return None
    for fmt in ("%Y-%m-%d %H:%M:%S", "%Y-%m-%dT%H:%M:%S"):
        try:
            return datetime.strptime(text, fmt)
        except ValueError:
            continue
    return None


def parse_iso_date(value):
    text = str(value or "").strip()
    if not text:
        return None
    for fmt in ("%Y-%m-%d", "%d/%m/%Y"):
        try:
            return datetime.strptime(text, fmt).date()
        except ValueError:
            continue
    return None


def format_ui_date(value, lang="he"):
    parsed = parse_iso_date(value)
    if not parsed:
        return "לא הוגדר" if lang == "he" else "Not set"
    return parsed.strftime("%d/%m/%Y")


def billing_mode_label(value, lang="he"):
    normalized = str(value or "monthly").strip().lower()
    labels = {
        "monthly": "חודשי" if lang == "he" else "Monthly",
        "yearly_prepaid": "שנתי מראש" if lang == "he" else "Yearly prepaid",
    }
    return labels.get(normalized, normalized or ("לא הוגדר" if lang == "he" else "Not set"))


def log_user_activity(event_type, action_label, script_id="", script_name="", details=""):
    user_id = session.get("user_id")
    if not user_id or session.get("is_admin"):
        return
    try:
        with get_db() as db:
            db.execute(
                """INSERT INTO activity_logs(
                user_id, username, full_name, event_type, action_label, script_id, script_name, details, created_at
                ) VALUES (?,?,?,?,?,?,?,?,?)""",
                (
                    user_id,
                    session.get("username", ""),
                    session.get("name", ""),
                    str(event_type or "")[:80],
                    str(action_label or "")[:120],
                    str(script_id or "")[:80],
                    str(script_name or "")[:120],
                    str(details or "")[:240],
                    datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                ),
            )
            db.commit()
    except Exception:
        pass


def create_support_request(user_row, request_type, message, script_id="", script_name=""):
    now_text = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    with get_db() as db:
        db.execute(
            """INSERT INTO support_requests(
            user_id, username, full_name, company_name, email, phone, request_type, script_id, script_name, message, status, created_at
            ) VALUES (?,?,?,?,?,?,?,?,?,?,?,?)""",
            (
                user_row["id"],
                user_row["username"] or "",
                user_row["full_name"] or "",
                user_row["company_name"] or "",
                user_row["email"] or "",
                user_row["phone"] or "",
                request_type,
                script_id or "",
                script_name or "",
                message,
                "pending",
                now_text,
            ),
        )
        db.commit()


def now_text():
    return datetime.now().strftime("%Y-%m-%d %H:%M:%S")


def report_job_status_label(status):
    labels = {
        "pending": "ממתין",
        "processing": "בעיבוד",
        "ready": "מוכן להורדה",
        "failed": "נכשל",
        "downloaded": "הורד",
        "expired": "פג תוקף",
    }
    return labels.get(str(status or "").strip().lower(), str(status or ""))


def cleanup_report_file(path_text):
    path = Path(str(path_text or "").strip())
    if not path:
        return
    try:
        if path.exists():
            path.unlink()
    except OSError:
        pass


def expire_report_jobs():
    current = now_text()
    with get_db() as db:
        expired_jobs = db.execute(
            "SELECT * FROM report_jobs WHERE status='ready' AND expires_at IS NOT NULL AND expires_at<=?",
            (current,),
        ).fetchall()
        for job in expired_jobs:
            cleanup_report_file(job["output_filename"])
            db.execute(
                "UPDATE report_jobs SET status='expired', status_note=? WHERE id=?",
                ("הדוח נשמר למשך 3 ימים ולאחר מכן פג תוקפו.", job["id"]),
            )
        if expired_jobs:
            db.commit()


def create_report_job(user_row, script_id, script_name, original_filename, input_path, input_ext, output_filename):
    with get_db() as db:
        created_at = now_text()
        params = (
            user_row["id"],
            user_row["username"] or "",
            user_row["full_name"] or "",
            user_row["company_name"] or "",
            script_id,
            script_name,
            original_filename or "",
            input_path,
            input_ext,
            output_filename,
            "pending",
            "הדוח ממתין לתחילת עיבוד.",
            created_at,
        )
        if db.is_postgres:
            job_row = db.execute(
                """INSERT INTO report_jobs(
                user_id, username, full_name, company_name, script_id, script_name, original_filename,
                input_path, input_ext, output_filename, status, status_note, created_at
                ) VALUES (?,?,?,?,?,?,?,?,?,?,?,?,?) RETURNING id""",
                params,
            ).fetchone()
        else:
            db.execute(
                """INSERT INTO report_jobs(
                user_id, username, full_name, company_name, script_id, script_name, original_filename,
                input_path, input_ext, output_filename, status, status_note, created_at
                ) VALUES (?,?,?,?,?,?,?,?,?,?,?,?,?)""",
                params,
            )
            job_row = db.execute("SELECT last_insert_rowid() AS id").fetchone()
        db.commit()
        return int(job_row["id"])


def get_user_report_jobs(user_id):
    expire_report_jobs()
    with get_db() as db:
        return db.execute(
            "SELECT * FROM report_jobs WHERE user_id=? ORDER BY created_at DESC, id DESC",
            (user_id,),
        ).fetchall()


def start_cleanup_report_job(job_id):
    def worker():
        with get_db() as db:
            job = db.execute("SELECT * FROM report_jobs WHERE id=?", (job_id,)).fetchone()
            if not job or job["status"] not in {"pending", "processing"}:
                return
            db.execute(
                "UPDATE report_jobs SET status='processing', status_note=?, started_at=? WHERE id=?",
                ("הקובץ נמצא בעיבוד. אפשר להמשיך לעבוד בפלטפורמה ובסיומו הוא יופיע בדוחות המוכנים להורדה.", now_text(), job_id),
            )
            db.commit()
        try:
            execute_script(
                SCRIPT_REGISTRY["nikuy"],
                job["input_path"],
                str(OUTPUT_FOLDER / job["output_filename"]),
                job["input_ext"],
                {},
            )
            with get_db() as db:
                db.execute(
                    "UPDATE report_jobs SET status='ready', status_note=?, completed_at=?, expires_at=? WHERE id=?",
                    (
                        "הדוח מוכן להורדה. הוא יישמר עד להורדה או עד 3 ימים, המוקדם מביניהם.",
                        now_text(),
                        (datetime.now() + timedelta(days=3)).strftime("%Y-%m-%d %H:%M:%S"),
                        job_id,
                    ),
                )
                db.commit()
        except Exception:
            with get_db() as db:
                db.execute(
                    "UPDATE report_jobs SET status='failed', status_note=?, completed_at=? WHERE id=?",
                    ("הדוח לא הושלם. מומלץ לנסות שוב עם הקובץ המקורי.", now_text(), job_id),
                )
                db.commit()
        finally:
            cleanup_report_file(job["input_path"])

    threading.Thread(target=worker, daemon=True).start()


def support_status_meta(status_value):
    normalized = str(status_value or "pending").strip().lower()
    mapping = {
        "pending": {
            "label": "ממתין להתייחסות מנהל מערכת",
            "short_label": "ממתין",
            "bg": "#fff7ed",
            "fg": "#c2410c",
        },
        "accepted": {
            "label": "התקבל",
            "short_label": "התקבל",
            "bg": "#eff6ff",
            "fg": "#1d4ed8",
        },
        "resolved": {
            "label": "טופל",
            "short_label": "טופל",
            "bg": "#ecfdf5",
            "fg": "#047857",
        },
    }
    return mapping.get(normalized, mapping["pending"])


def resolve_script_from_output_name(filename):
    logical_name = filename.split("_", 1)[-1] if "_" in filename else filename
    for script in SCRIPT_REGISTRY.values():
        suffix = script.get("output_suffix", "")
        if suffix and logical_name.startswith(suffix + "."):
            return script
    return None


RIMON_MAPPING_FIELDS = [
    {"name": "employee_name_source", "label": "שם עובד", "required": True},
    {"name": "payroll_number_source", "label": "מספר עובד", "required": True},
    {"name": "date_source", "label": "תאריך", "required": True},
    {"name": "day_name_source", "label": "סוג יום", "required": False},
    {"name": "entry_time_source", "label": "שעת כניסה", "required": False},
    {"name": "exit_time_source", "label": "שעת יציאה", "required": False},
    {"name": "total_hours_source", "label": "סה\"כ שעות", "required": False},
    {"name": "standard_hours_source", "label": "שעות תקן", "required": False},
    {"name": "missing_hours_source", "label": "שעות חוסר", "required": False},
    {"name": "event_source", "label": "אירוע", "required": True},
    {"name": "error_text_source", "label": "שדה שגיאה", "required": False},
    {"name": "department_source", "label": "מחלקה", "required": False},
    {"name": "id_number_source", "label": "תעודת זהות", "required": False},
]

ATTENDANCE_ALERTS_MAPPING_FIELDS = [
    {"name": "employee_name_source", "label": "שם עובד", "required": True},
    {"name": "payroll_number_source", "label": "מספר עובד", "required": True},
    {"name": "date_source", "label": "תאריך", "required": True},
    {"name": "day_name_source", "label": "סוג יום", "required": False},
    {"name": "entry_time_source", "label": "שעת כניסה", "required": False},
    {"name": "exit_time_source", "label": "שעת יציאה", "required": False},
    {"name": "total_hours_source", "label": "סה\"כ שעות", "required": False},
    {"name": "standard_hours_source", "label": "שעות תקן", "required": False},
    {"name": "break_hours_source", "label": "הפסקה", "required": False},
    {"name": "missing_hours_source", "label": "שעות חוסר", "required": False},
    {"name": "regular_hours_source", "label": "שעות רגילות", "required": False},
    {"name": "overtime_100_source", "label": 'שעות נוספות 100%', "required": False},
    {"name": "overtime_125_source", "label": 'שעות נוספות 125%', "required": False},
    {"name": "overtime_150_source", "label": 'שעות נוספות 150%', "required": False},
    {"name": "overtime_200_source", "label": 'שעות נוספות 200%', "required": False},
    {"name": "event_source", "label": "אירוע", "required": True},
    {"name": "error_text_source", "label": "שדה שגיאה", "required": False},
    {"name": "department_source", "label": "מחלקה", "required": False},
    {"name": "id_number_source", "label": "תעודת זהות", "required": False},
]

FLAMINGO_MAPPING_FIELDS = [
    {"name": "worker_name_source", "label": "שם עובד", "required": True},
    {"name": "worker_number_source", "label": "מספר עובד", "required": False},
    {"name": "id_number_source", "label": "תעודת זהות / דרכון", "required": False},
    {"name": "department_source", "label": "מחלקה", "required": False},
    {"name": "hourly_rate_source", "label": "תעריף שעתי", "required": False, "critical": True},
    {"name": "payable_hours_source", "label": "שעות לתשלום בפועל", "required": True, "critical": True},
    {"name": "attendance_hours_source", "label": "נוכחות", "required": False},
    {"name": "standard_hours_source", "label": "תקן", "required": False},
    {"name": "missing_hours_source", "label": "חוסר", "required": False},
    {"name": "start_date_source", "label": "תחילת עבודה", "required": False},
]

MATAN_MISSING_MAPPING_FIELDS = [
    {"name": "employee_name_source", "label": "שם עובד", "required": True},
    {"name": "employee_number_source", "label": "מספר עובד", "required": False},
    {"name": "id_number_source", "label": "תעודת זהות", "required": False},
    {"name": "badge_number_source", "label": "מספר תג", "required": False},
    {"name": "passport_number_source", "label": "דרכון", "required": False},
    {"name": "month_source", "label": "חודש", "required": False},
    {"name": "standard_hours_source", "label": "ש.תקן", "required": True, "critical": True},
    {"name": "missing_hours_source", "label": "חוסר", "required": True, "critical": True},
    {"name": "attendance_hours_source", "label": "ש.נוכחות", "required": False},
    {"name": "vacation_hours_source", "label": "חופשה", "required": False},
    {"name": "sick_hours_source", "label": "מחלה", "required": False},
    {"name": "reserve_hours_source", "label": "מילואים", "required": False},
    {"name": "pregnancy_hours_source", "label": "שעות הריון", "required": False},
    {"name": "special_child_hours_source", "label": "שעות ילד מיוחד", "required": False},
    {"name": "absence_hours_source", "label": "היעדרות", "required": False},
]

INACTIVE_WORKERS_MAPPING_FIELDS = [
    {"name": "employee_name_source", "label": "שם עובד", "required": True},
    {"name": "employee_number_source", "label": "מספר עובד", "required": False},
    {"name": "badge_number_source", "label": "מספר תג", "required": False},
    {"name": "id_number_source", "label": "תעודת זהות", "required": False},
    {"name": "passport_number_source", "label": "דרכון", "required": False},
    {"name": "date_source", "label": "תאריך", "required": True},
    {"name": "entry_time_source", "label": "כניסה", "required": False, "critical": True},
    {"name": "exit_time_source", "label": "יציאה", "required": False, "critical": True},
    {"name": "total_hours_source", "label": "סה\"כ שעות", "required": False, "critical": True},
    {"name": "event_source", "label": "אירוע", "required": False},
    {"name": "department_source", "label": "מחלקה", "required": False},
]

ORG_HIERARCHY_MAPPING_FIELDS = [
    {"name": "employee_name_source", "label": "שם עובד", "required": True, "critical": True},
    {"name": "direct_manager_source", "label": "מנהל ישיר", "required": True, "critical": True},
    {"name": "department_source", "label": "מחלקה", "required": True, "critical": True},
    {"name": "employee_number_source", "label": "מספר עובד", "required": False},
    {"name": "id_number_source", "label": "תעודת זהות", "required": False},
    {"name": "passport_number_source", "label": "דרכון", "required": False},
    {"name": "manager_flag_source", "label": "סימון מנהל", "required": False},
    {"name": "email_source", "label": "אימייל", "required": False},
    {"name": "secondary_email_source", "label": "אימייל נוסף", "required": False},
    {"name": "app_access_source", "label": "הרשאה לאפליקציה", "required": False},
    {"name": "employment_percent_source", "label": "אחוז משרה", "required": False},
    {"name": "agreement_number_source", "label": "מס' הסכם", "required": False},
    {"name": "agreement_name_source", "label": "שם הסכם", "required": False},
]

MATAN_CORRECTIONS_MAPPING_FIELDS = [
    {"name": "entry_col_source", "label": "עמודת כניסה", "required": True, "critical": True},
    {"name": "exit_col_source", "label": "עמודת יציאה", "required": True, "critical": True},
    {"name": "date_col_source", "label": "עמודת תאריך", "required": False},
]

DEPT_PAYROLL_MAPPING_FIELDS = [
    {"name": "worker_name_source", "label": "שם עובד", "required": True},
    {"name": "dept_number_source", "label": "מספר מחלקה", "required": False},
    {"name": "department_source", "label": "שם מחלקה", "required": False},
    {"name": "payable_hours_source", "label": "שעות לתשלום בפועל", "required": True, "critical": True},
    {"name": "id_number_source", "label": "תעודת זהות / דרכון", "required": False},
    {"name": "phone_source", "label": "מספר נייד", "required": False},
    {"name": "passport_source", "label": "מספר פספורט", "required": False},
    {"name": "event_source", "label": "אירוע / לקוח", "required": True, "critical": True},
    {"name": "entry_time_source", "label": "כניסה", "required": False},
    {"name": "exit_time_source", "label": "יציאה", "required": False},
    {"name": "date_source", "label": "תאריך", "required": False},
    {"name": "total_hours_source", "label": 'סה"כ שעות', "required": True, "critical": True},
    {"name": "hourly_rate_source", "label": "תעריף שעתי", "required": False, "has_manual": True},
    {"name": "housing_source", "label": "חיוב דירה", "required": False, "has_manual": True},
]

DEPT_PAYROLL_SUGGESTION_KEYWORDS = {
    "worker_name_source": ["שםלתצוגה", "שםעובד", "עובד", "name"],
    "dept_number_source": ["מספרמחלקה", "מחלקה", "department"],
    "department_source": ["מחלקה", "department"],
    "payable_hours_source": ["שלתשלום", "שעותלתשלום", "שעותמשולמות", "סהכ", "רגילות", "נוכחות"],
    "id_number_source": ["תעודתזהות", "דרכון", "זהות", "id"],
    "phone_source": ["נייד", "טלפון", "phone", "mobile"],
    "passport_source": ["פספורט", "דרכון", "passport"],
    "event_source": ["אירוע", "event", "לקוח", "סטטוס"],
    "entry_time_source": ["כניסה", "entry", "checkin"],
    "exit_time_source": ["יציאה", "exit", "checkout"],
    "date_source": ["תאריך", "date"],
    "total_hours_source": ["סהכ", "סה\"כ", "total", "hours", "שעות"],
    "hourly_rate_source": ["תעריףשעתי", "תעריף", "rate", "hourly"],
    "housing_source": ["חיובדירה", "דירה", "housing", "הערות"],
}


RIMON_SUGGESTION_KEYWORDS = {
    "employee_name_source": ["שםלתצוגה", "שםעובד", "עובד", "employee", "name"],
    "payroll_number_source": ["מספרשכר", "מספרעובד", "שכר", "עובד", "payroll", "employeeid"],
    "date_source": ["תאריך", "date"],
    "day_name_source": ["יום", "day"],
    "entry_time_source": ["כניסה", "entry", "checkin"],
    "exit_time_source": ["יציאה", "exit", "checkout"],
    "total_hours_source": ["סהכ", "סה\"כ", "total", "hours"],
    "standard_hours_source": ["תקן", "ש.תקן", "standard"],
    "missing_hours_source": ["חוסר", "missing"],
    "event_source": ["אירוע", "event", "סטטוס"],
    "error_text_source": ["שגיאה", "שגיאות", "error", "errors"],
    "department_source": ["מחלקה", "department"],
    "id_number_source": ["תעודתזהות", "זהות", "דרכון", "id", "identity"],
}

ATTENDANCE_ALERTS_SUGGESTION_KEYWORDS = {
    "employee_name_source": ["שםלתצוגה", "שםעובד", "עובד", "employee", "name"],
    "payroll_number_source": ["מספרשכר", "מספרעובד", "שכר", "עובד", "payroll", "employeeid"],
    "date_source": ["תאריך", "date"],
    "day_name_source": ["יום", "day"],
    "entry_time_source": ["כניסה", "entry", "checkin"],
    "exit_time_source": ["יציאה", "exit", "checkout"],
    "total_hours_source": ["סהכ", "סה\"כ", "total", "hours"],
    "standard_hours_source": ["תקן", "ש.תקן", "standard"],
    "break_hours_source": ["הפסקה", "break"],
    "missing_hours_source": ["חוסר", "missing"],
    "regular_hours_source": ["רגילות", "regular"],
    "overtime_100_source": ["100", "נוספות100"],
    "overtime_125_source": ["125", "נוספות125"],
    "overtime_150_source": ["150", "נוספות150"],
    "overtime_200_source": ["200", "נוספות200"],
    "event_source": ["אירוע", "event", "סטטוס"],
    "error_text_source": ["שגיאה", "שגיאות", "error", "errors"],
    "department_source": ["מחלקה", "department"],
    "id_number_source": ["תעודתזהות", "זהות", "דרכון", "id", "identity"],
}

FLAMINGO_SUGGESTION_KEYWORDS = {
    "worker_name_source": ["שםלתצוגה", "שםעובד", "עובד", "name"],
    "worker_number_source": ["מספרבשכר", "מספרעובד", "מפעלבשכר", "employee"],
    "id_number_source": ["תעודתזהות", "דרכון", "זהות", "id"],
    "department_source": ["מחלקה", "department"],
    "hourly_rate_source": ["תעריף", "שעה", "rate"],
    "payable_hours_source": ["שעותלתשלום", "שעותמשולמות", "רגילות", "נוכחות"],
    "attendance_hours_source": ["נוכחות"],
    "standard_hours_source": ["תקן"],
    "missing_hours_source": ["חוסר"],
    "start_date_source": ["תחילתעבודה"],
}

MATAN_MISSING_SUGGESTION_KEYWORDS = {
    "employee_name_source": ["שםעובד", "שם", "עובד"],
    "employee_number_source": ["מספרעובד", "מספר", "עובד"],
    "id_number_source": ["תעודתזהות", "זהות"],
    "badge_number_source": ["מספרתג", "תג"],
    "passport_number_source": ["דרכון", "passport"],
    "month_source": ["חודש"],
    "standard_hours_source": ["ש.תקן", "תקן", "שעותתקן"],
    "missing_hours_source": ["חוסר", "שעותחוסר"],
    "attendance_hours_source": ["ש.נוכחות", "נוכחות"],
    "vacation_hours_source": ["חופשה"],
    "sick_hours_source": ["מחלה"],
    "reserve_hours_source": ["מילואים"],
    "pregnancy_hours_source": ["הריון"],
    "special_child_hours_source": ["ילדמיחד", "ילדמיוחד"],
    "absence_hours_source": ["היעדרות"],
}

INACTIVE_WORKERS_SUGGESTION_KEYWORDS = {
    "employee_name_source": ["שםעובד", "שם", "עובד"],
    "employee_number_source": ["מספרעובד", "מספר", "עובד"],
    "badge_number_source": ["תג", "מספרתג"],
    "id_number_source": ["תעודתזהות", "זהות"],
    "passport_number_source": ["דרכון", "passport"],
    "date_source": ["תאריך", "date"],
    "entry_time_source": ["כניסה", "checkin", "entry"],
    "exit_time_source": ["יציאה", "checkout", "exit"],
    "total_hours_source": ["סהכ", "סה\"כ", "שעות", "total"],
    "event_source": ["אירוע", "event"],
    "department_source": ["מחלקה", "department"],
}

ORG_HIERARCHY_SUGGESTION_KEYWORDS = {
    "employee_name_source": ["שםעובד", "עובד", "name"],
    "direct_manager_source": ["מנהלישיר", "directmanager", "manager"],
    "department_source": ["מחלקה", "department"],
    "employee_number_source": ["מספרעובד", "מספרשכר", "שכר", "employeeid", "payroll"],
    "id_number_source": ["תז", "תעותזהות", "תעודתזהות", "id"],
    "passport_number_source": ["דרכון", "passport"],
    "manager_flag_source": ["מנהל", "ismanager", "managerflag"],
    "email_source": ["אימייל", "אימיל", "email", "mail"],
    "secondary_email_source": ["אימיילנוסף", "מיילנוסף", "secondaryemail", "additionalemail"],
    "app_access_source": ["הרשאהלאפליקציה", "אפליקציה", "access", "permission"],
    "employment_percent_source": ["אחוזמשרה", "משרה", "percent", "fte"],
    "agreement_number_source": ["מסהסכם", "מספרהסכם", "agreementnumber"],
    "agreement_name_source": ["שםהסכם", "agreementname"],
}

RIMON_META_LABEL_TOKENS = {
    normalize_token(label)
    for label in [
        "שם לתצוגה",
        "שם עובד",
        "מחלקה",
        "תג עובד",
        "מספר שכר",
        "מספר עובד",
        "מספר בשכר",
        "תעודת זהות",
        "דרכון",
        "תחילת עבודה",
    ]
}

FLAMINGO_META_LABEL_TOKENS = {
    normalize_token(label)
    for label in [
        "שם לתצוגה",
        "שם עובד",
        "מחלקה",
        "מספר בשכר",
        "מספר עובד",
        "מס' מפעל בשכר",
        "תעודת זהות",
        "דרכון",
        "תחילת עבודה",
        "הערות",
    ]
}


def open_excel_workbook(input_path, extension):
    ext = str(extension or get_extension(input_path)).lower()
    if ext == "xlsx":
        return "xlsx", load_workbook(input_path, data_only=True, read_only=True)
    return "xls", xlrd.open_workbook(input_path)


def iter_excel_sheets(workbook_kind, workbook):
    if workbook_kind == "xlsx":
        return list(workbook.worksheets)
    return workbook.sheets()


def get_excel_dims(sheet, workbook_kind):
    if workbook_kind == "xlsx":
        return int(sheet.max_row or 0), int(sheet.max_column or 0)
    return sheet.nrows, sheet.ncols


def get_excel_cell(sheet, workbook_kind, row_index, col_index, default=""):
    rows, cols = get_excel_dims(sheet, workbook_kind)
    if row_index < 0 or col_index < 0 or row_index >= rows or col_index >= cols:
        return default
    if workbook_kind == "xlsx":
        value = sheet.cell(row=row_index + 1, column=col_index + 1).value
        return default if value is None else value
    return get_sheet_cell(sheet, row_index, col_index, default)


def stringify_excel_value(value):
    if value is None:
        return ""
    if isinstance(value, datetime):
        return value.strftime("%Y-%m-%d %H:%M:%S")
    if isinstance(value, date):
        return value.strftime("%Y-%m-%d")
    return str(value).strip()


def parse_excel_date_generic(workbook_kind, workbook, value):
    if value in ("", None):
        return None
    if isinstance(value, datetime):
        return value.date()
    if isinstance(value, date):
        return value
    if workbook_kind == "xls":
        return parse_excel_date(workbook, value)
    return None


def looks_like_excel_date_sample(sample_text):
    text = str(sample_text or "").strip()
    if not text:
        return False
    try:
        number = float(text)
        return 20000 <= number <= 60000
    except ValueError:
        pass
    for fmt in ("%Y-%m-%d", "%d/%m/%Y", "%d.%m.%Y"):
        try:
            datetime.strptime(text, fmt)
            return True
        except ValueError:
            continue
    return False


def looks_like_time_sample(sample_text):
    text = str(sample_text or "").strip().replace("*", "").replace("?", "")
    return bool(re.fullmatch(r"\d{1,2}:\d{2}", text))


def looks_like_day_name_sample(sample_text):
    text = str(sample_text or "").strip()
    return text in {"א", "ב", "ג", "ד", "ה", "ו", "שבת", "יום שישי", "ראשון", "שני", "שלישי", "רביעי", "חמישי", "שישי"}


def is_rimon_option_relevant_for_field(field_name, option):
    header_token = normalize_token(option.get("header", ""))
    exact_token = normalize_token(option.get("exact_header", ""))
    sample_text = option.get("sample", "")
    resolved_token = exact_token or header_token

    if field_name == "date_source":
        return looks_like_excel_date_sample(sample_text) or "תאריך" in resolved_token or "date" in resolved_token
    if field_name == "day_name_source":
        return looks_like_day_name_sample(sample_text) or resolved_token == "יום" or "day" in resolved_token
    if field_name == "entry_time_source":
        return looks_like_time_sample(sample_text) and ("כניסה" in resolved_token or "entry" in resolved_token or "checkin" in resolved_token)
    if field_name == "exit_time_source":
        return looks_like_time_sample(sample_text) and ("יציאה" in resolved_token or "exit" in resolved_token or "checkout" in resolved_token)
    if field_name == "total_hours_source":
        return looks_like_time_sample(sample_text) and ("סהכ" in resolved_token or "שעות" in resolved_token or "total" in resolved_token or "hours" in resolved_token)
    if field_name == "event_source":
        sample_token = normalize_token(str(sample_text))
        return "אירוע" in resolved_token or "event" in resolved_token or "סטטוס" in resolved_token or any(
            keyword in sample_token for keyword in ("עבודהמהבית", "חופשה", "מחלה", "מילואים", "היעדר")
        )
    if field_name == "error_text_source":
        sample_token = normalize_token(str(sample_text))
        return "שגיאה" in resolved_token or "error" in resolved_token or any(
            keyword in sample_token for keyword in ("יוםחסר", "חסרדיווח", "שגיאה")
        )
    if field_name == "employee_name_source":
        return any(token in resolved_token for token in ("שם", "employee", "name", "עובד"))
    if field_name == "payroll_number_source":
        return any(token in resolved_token for token in ("מספר", "תג", "payroll", "employeeid", "עובד", "שכר"))
    if field_name == "department_source":
        return any(token in resolved_token for token in ("מחלקה", "department"))
    if field_name == "id_number_source":
        return any(token in resolved_token for token in ("זהות", "דרכון", "id", "identity"))
    return False


def dedupe_rimon_field_options(options):
    grouped = {}
    order = []
    for option in options:
        key = normalize_token(option.get("exact_header", "") or option.get("header", ""))
        if not key:
            key = option.get("value", "")
        if key not in grouped:
            grouped[key] = []
            order.append(key)
        grouped[key].append(option)

    deduped = []
    for key in order:
        candidates = grouped[key]
        with_sample = [option for option in candidates if option.get("sample")]
        if with_sample:
            deduped.extend(with_sample[:1])
            continue
        exact = [option for option in candidates if option.get("exact_header")]
        if exact:
            deduped.extend(exact[:1])
            continue
        deduped.extend(candidates[:1])
    return deduped


def filter_rimon_table_options_for_display(options):
    exact_sample_signatures = set()
    for option in options:
        if option.get("source_kind") != "table_exact":
            continue
        header_token = normalize_token(option.get("exact_header", "") or option.get("header", ""))
        sample_token = normalize_token(option.get("sample", ""))
        if header_token and sample_token:
            exact_sample_signatures.add((header_token, sample_token))

    filtered = []
    seen_nearby_signatures = set()
    for option in options:
        if option.get("source_kind") != "table_nearby":
            filtered.append(option)
            continue
        header_token = normalize_token(option.get("exact_header", "") or option.get("header", ""))
        sample_token = normalize_token(option.get("sample", ""))
        if not sample_token:
            continue
        signature = (header_token, sample_token)
        if signature in exact_sample_signatures or signature in seen_nearby_signatures:
            continue
        seen_nearby_signatures.add(signature)
        filtered.append(option)
    return filtered


def filter_rimon_options_for_field(field_name, options):
    if field_name == "date_source":
        return [option for option in options if looks_like_excel_date_sample(option.get("sample", ""))]
    return options


def default_matan_missing_mapping():
    return {
        "employee_name_source": "col:2",
        "employee_number_source": "col:0",
        "id_number_source": "",
        "badge_number_source": "",
        "passport_number_source": "",
        "month_source": "col:1",
        "standard_hours_source": "col:3",
        "missing_hours_source": "col:4",
        "attendance_hours_source": "col:6",
        "vacation_hours_source": "col:7",
        "sick_hours_source": "col:8",
        "reserve_hours_source": "col:9",
        "pregnancy_hours_source": "col:10",
        "special_child_hours_source": "col:13",
        "absence_hours_source": "col:14",
    }


def build_matan_missing_mapping_warnings(mapping):
    warnings = []
    if not mapping.get("standard_hours_source"):
        warnings.append("לא נבחר שדה שעות תקן. בלי השדה הזה הדוח לא יהיה אמין.")
    if not mapping.get("missing_hours_source"):
        warnings.append("לא נבחר שדה חוסר. בלי השדה הזה לא ניתן לסנן את העובדים נכון.")
    identifier_sources = [
        mapping.get("employee_number_source"),
        mapping.get("id_number_source"),
        mapping.get("badge_number_source"),
        mapping.get("passport_number_source"),
    ]
    if not any(identifier_sources):
        warnings.append("לא נבחר מזהה נוסף לעובד. מומלץ לבחור מספר עובד, תעודת זהות, מספר תג או דרכון.")
    return warnings


def default_inactive_workers_mapping():
    return {
        "employee_name_source": "col:2",
        "employee_number_source": "",
        "badge_number_source": "col:4",
        "id_number_source": "",
        "passport_number_source": "",
        "date_source": "col:0",
        "entry_time_source": "col:7",
        "exit_time_source": "col:8",
        "total_hours_source": "col:15",
        "event_source": "col:9",
        "department_source": "col:30",
    }


def build_inactive_workers_mapping_warnings(mapping, options):
    warnings = []
    if not mapping.get("date_source"):
        warnings.append("לא נבחר שדה תאריך. בלי השדה הזה לא ניתן לבדוק אי-פעילות.")
    if not mapping.get("employee_name_source"):
        warnings.append("לא נבחר שדה שם עובד. בלי השדה הזה הדוח לא יהיה קריא.")
    identifier_sources = [
        mapping.get("employee_number_source"),
        mapping.get("badge_number_source"),
        mapping.get("id_number_source"),
        mapping.get("passport_number_source"),
    ]
    if not any(identifier_sources):
        warnings.append("לא נבחר מזהה נוסף לעובד. מומלץ לבחור מספר עובד, מספר תג, תעודת זהות או דרכון.")
    has_attendance_pair = bool(mapping.get("entry_time_source") and mapping.get("exit_time_source"))
    has_total_hours = bool(mapping.get("total_hours_source"))
    if not has_attendance_pair and not has_total_hours:
        warnings.append("יש לבחור כניסה ויציאה יחד, או לחלופין שדה סה\"כ שעות.")
    threshold_value = parse_int_or_none(options.get("inactive_period_value", ""))
    if threshold_value is None or threshold_value <= 0:
        warnings.append("ערך בדיקת אי-הפעילות אינו תקין. מומלץ להזין מספר חיובי של ימים או חודשים.")
    return warnings


def detect_inactive_workers_header_row(sheet, workbook_kind):
    rows, cols = get_excel_dims(sheet, workbook_kind)
    best_row = 0
    best_score = -1
    for row_index in range(min(rows, 40)):
        row_tokens = [normalize_token(get_excel_cell(sheet, workbook_kind, row_index, col_index, "")) for col_index in range(cols)]
        score = 0
        if "תאריך" in row_tokens:
            score += 3
        if "שםעובד" in row_tokens:
            score += 3
        if "כניסה" in row_tokens:
            score += 2
        if "יציאה" in row_tokens:
            score += 2
        if any(token in row_tokens for token in ("סהכ", "סהכ\"")):
            score += 1
        if "אירוע" in row_tokens:
            score += 1
        if score > best_score:
            best_score = score
            best_row = row_index
    return best_row


def extract_inactive_workers_mapping_value(sheet, workbook_kind, source, row_index):
    text = str(source or "").strip()
    if not text.startswith("col:"):
        return ""
    try:
        col_index = int(text.split(":", 1)[1])
    except ValueError:
        return ""
    return get_excel_cell(sheet, workbook_kind, row_index, col_index, "")


def build_inactive_workers_mapping_options(input_path, extension):
    workbook_kind, workbook = open_excel_workbook(input_path, extension)
    sheet = iter_excel_sheets(workbook_kind, workbook)[0]
    rows, cols = get_excel_dims(sheet, workbook_kind)
    header_row = detect_inactive_workers_header_row(sheet, workbook_kind)

    options = [{"value": "", "label": "לא נבחר", "source_kind": "empty"}]
    for col_index in range(cols):
        header = stringify_excel_value(get_excel_cell(sheet, workbook_kind, header_row, col_index, ""))
        if not header:
            continue
        sample = ""
        for row_index in range(header_row + 1, rows):
            candidate = stringify_excel_value(get_excel_cell(sheet, workbook_kind, row_index, col_index, ""))
            if not candidate:
                continue
            token = normalize_token(candidate)
            if token in {"תאריך", "יום", "שםעובד", "תג", "הסכם", "כניסה", "יציאה", "אירוע", "נוכחות", "כניסהמ", "יציאהמ", "סהכ", "תקן", "חוסר", "מחלקה", "שגיאות"}:
                continue
            sample = candidate
            break
        options.append(
            {
                "value": f"col:{col_index}",
                "label": f"עמודה {get_column_letter(col_index + 1)} - {header}" + (f" (לדוגמה: {sample})" if sample else ""),
                "source_kind": "table_exact",
                "match_token": normalize_token(header),
                "header": header,
                "sample": sample,
            }
        )

    preferred_tokens = {
        "employee_name_source": ["שםעובד"],
        "employee_number_source": ["מספרעובד"],
        "badge_number_source": ["תג", "מספרתג"],
        "id_number_source": ["תעודתזהות"],
        "passport_number_source": ["דרכון"],
        "date_source": ["תאריך"],
        "entry_time_source": ["כניסה"],
        "exit_time_source": ["יציאה"],
        "total_hours_source": ["סהכ", "שעות"],
        "event_source": ["אירוע"],
        "department_source": ["מחלקה"],
    }

    options_by_field = {}
    suggestions = {}
    for field in INACTIVE_WORKERS_MAPPING_FIELDS:
        field_name = field["name"]
        field_options = [options[0]]
        keywords = INACTIVE_WORKERS_SUGGESTION_KEYWORDS.get(field_name, [])
        for option in options[1:]:
            token = option.get("match_token", "")
            if any(keyword in token for keyword in keywords):
                field_options.append(option)
        for option in options[1:]:
            if option["value"] not in {item["value"] for item in field_options}:
                field_options.append(option)
        options_by_field[field_name] = field_options

        suggested = ""
        for preferred in preferred_tokens.get(field_name, []):
            for option in field_options[1:]:
                token = option.get("match_token", "")
                if preferred == token or preferred in token:
                    suggested = option["value"]
                    break
            if suggested:
                break
        if not suggested:
            for option in field_options[1:]:
                token = option.get("match_token", "")
                if any(keyword in token for keyword in keywords):
                    suggested = option["value"]
                    break
        suggestions[field_name] = suggested

    return {
        "header_row": header_row,
        "options_by_field": options_by_field,
        "suggestions": suggestions,
        "suggested_template_name": "תבנית עובדים לא פעילים",
    }


def extract_matan_missing_mapping_value(sheet, source, row_index):
    text = str(source or "").strip()
    if not text.startswith("col:"):
        return ""
    try:
        col_index = int(text.split(":", 1)[1])
    except ValueError:
        return ""
    return get_sheet_cell(sheet, row_index, col_index, "")


def build_matan_missing_mapping_options(input_path, extension):
    if extension != "xls":
        raise ValueError("Matan missing-hours tool currently supports XLS export only")
    wb = xlrd.open_workbook(input_path)
    ws = wb.sheet_by_index(0)
    header_row = detect_matan_missing_header_row(ws)
    options = [{"value": "", "label": "לא נבחר", "source_kind": "empty"}]
    for col_index in range(ws.ncols):
        header = stringify_excel_value(get_sheet_cell(ws, header_row, col_index, ""))
        if not header:
            continue
        sample = ""
        for row_index in range(header_row + 1, ws.nrows):
            candidate = stringify_excel_value(get_sheet_cell(ws, row_index, col_index, ""))
            if candidate:
                sample = candidate
                break
        options.append(
            {
                "value": f"col:{col_index}",
                "label": f"עמודה {get_column_letter(col_index + 1)} - {header}" + (f" (לדוגמה: {sample})" if sample else ""),
                "source_kind": "table_exact",
                "match_token": normalize_token(header),
                "header": header,
                "sample": sample,
            }
        )

    options_by_field = {}
    suggestions = {}
    preferred_tokens = {
        "employee_name_source": ["שםעובד"],
        "employee_number_source": ["מספרעובד"],
        "id_number_source": ["תעודתזהות"],
        "badge_number_source": ["מספרתג", "תג"],
        "passport_number_source": ["דרכון"],
        "month_source": ["חודש"],
        "standard_hours_source": ["ש.תקן", "תקן"],
        "missing_hours_source": ["חוסר"],
        "attendance_hours_source": ["ש.נוכחות", "נוכחות"],
        "vacation_hours_source": ["חופשה"],
        "sick_hours_source": ["מחלה"],
        "reserve_hours_source": ["מילואים"],
        "pregnancy_hours_source": ["הריון"],
        "special_child_hours_source": ["ילדמיחד", "ילדמיוחד"],
        "absence_hours_source": ["היעדרות"],
    }
    for field in MATAN_MISSING_MAPPING_FIELDS:
        field_name = field["name"]
        field_options = [options[0]]
        keywords = MATAN_MISSING_SUGGESTION_KEYWORDS.get(field_name, [])
        for option in options[1:]:
            token = option.get("match_token", "")
            if any(keyword in token for keyword in keywords):
                field_options.append(option)
        for option in options[1:]:
            if option["value"] not in {item["value"] for item in field_options}:
                field_options.append(option)
        options_by_field[field_name] = field_options

        suggested = ""
        for preferred in preferred_tokens.get(field_name, []):
            for option in field_options[1:]:
                token = option.get("match_token", "")
                if preferred == token or preferred in token:
                    suggested = option["value"]
                    break
            if suggested:
                break
        if not suggested:
            for option in field_options[1:]:
                token = option.get("match_token", "")
                if any(keyword in token for keyword in keywords):
                    suggested = option["value"]
                    break
        suggestions[field_name] = suggested

    return {
        "header_row": header_row,
        "options_by_field": options_by_field,
        "suggestions": suggestions,
        "suggested_template_name": "תבנית שעות חסר",
    }


def default_org_hierarchy_mapping():
    return {field["name"]: "" for field in ORG_HIERARCHY_MAPPING_FIELDS}


def build_org_hierarchy_mapping_warnings(mapping):
    warnings = []
    if not mapping.get("employee_name_source"):
        warnings.append("לא נבחר שדה שם עובד. בלי השדה הזה לא ניתן לבנות את מבנה הדיווח.")
    if not mapping.get("direct_manager_source"):
        warnings.append("לא נבחר שדה מנהל ישיר. בלי השדה הזה לא ניתן לבנות את ההיררכיה נכון.")
    if not mapping.get("department_source"):
        warnings.append("לא נבחר שדה מחלקה. הפלט עדיין יופק, אבל יהיה פחות שימושי לסיכומי מחלקות.")
    identifier_sources = [
        mapping.get("employee_number_source"),
        mapping.get("id_number_source"),
        mapping.get("passport_number_source"),
    ]
    if not any(identifier_sources):
        warnings.append("לא נבחר מזהה נוסף לעובד. מומלץ לבחור מספר עובד, תעודת זהות או דרכון.")
    return warnings


def build_org_hierarchy_mapping_options(input_path, extension):
    if extension != "csv":
        raise ValueError("Organizational hierarchy tool currently supports CSV input only")
    with open(input_path, "r", encoding="utf-8-sig", newline="") as handle:
        rows = list(csv.reader(handle))
    headers = rows[0] if rows else []
    samples_by_index = {}
    for col_index, _header in enumerate(headers):
        sample = ""
        for row in rows[1:]:
            if col_index < len(row):
                candidate = str(row[col_index] or "").strip()
                if candidate:
                    sample = candidate
                    break
        samples_by_index[col_index] = sample

    base_options = [{"value": "", "label": "לא נבחר", "source_kind": "empty"}]
    for col_index, header in enumerate(headers):
        header_text = str(header or "").strip()
        if not header_text:
            continue
        sample = samples_by_index.get(col_index, "")
        base_options.append(
            {
                "value": f"header:{header_text}",
                "label": f"עמודה {get_column_letter(col_index + 1)} - {header_text}" + (f" (לדוגמה: {sample})" if sample else ""),
                "source_kind": "table_exact",
                "match_token": normalize_token(header_text),
                "header": header_text,
                "sample": sample,
            }
        )

    preferred_tokens = {
        "employee_name_source": ["שםעובד"],
        "direct_manager_source": ["מנהלישיר"],
        "department_source": ["מחלקה"],
        "employee_number_source": ["מספרשכר", "שכר", "מספרעובד"],
        "id_number_source": ["ת.ז", "תז", "תעודתזהות"],
        "passport_number_source": ["דרכון"],
        "manager_flag_source": ["מנהל"],
        "email_source": ["אימייל"],
        "secondary_email_source": ["אימיילנוסף"],
        "app_access_source": ["הרשאהלאפליקציה"],
        "employment_percent_source": ["אחוזמשרה"],
        "agreement_number_source": ["מסהסכם", "מספרהסכם"],
        "agreement_name_source": ["שםהסכם"],
    }

    options_by_field = {}
    suggestions = {}
    for field in ORG_HIERARCHY_MAPPING_FIELDS:
        field_name = field["name"]
        field_options = [base_options[0]]
        keywords = ORG_HIERARCHY_SUGGESTION_KEYWORDS.get(field_name, [])
        for option in base_options[1:]:
            token = option.get("match_token", "")
            if any(keyword in token for keyword in keywords):
                field_options.append(option)
        for option in base_options[1:]:
            if option["value"] not in {item["value"] for item in field_options}:
                field_options.append(option)
        options_by_field[field_name] = field_options

        suggested = ""
        for preferred in preferred_tokens.get(field_name, []):
            for option in field_options[1:]:
                token = option.get("match_token", "")
                if preferred == token:
                    suggested = option["value"]
                    break
            if suggested:
                break
        if not suggested:
            for preferred in preferred_tokens.get(field_name, []):
                for option in field_options[1:]:
                    token = option.get("match_token", "")
                    if preferred in token:
                        suggested = option["value"]
                        break
                if suggested:
                    break
        if not suggested:
            for option in field_options[1:]:
                token = option.get("match_token", "")
                if any(keyword in token for keyword in keywords):
                    suggested = option["value"]
                    break
        suggestions[field_name] = suggested

    return {
        "options_by_field": options_by_field,
        "suggestions": suggestions,
        "suggested_template_name": "תבנית ארגונית",
    }


def default_flamingo_mapping():
    return {
        "worker_name_source": "meta:שם לתצוגה",
        "worker_number_source": "meta:מספר בשכר",
        "id_number_source": "meta:תעודת זהות",
        "department_source": "meta:מחלקה",
        "hourly_rate_source": "",
        "payable_hours_source": "summary:שעות לתשלום",
        "attendance_hours_source": "summary:נוכחות",
        "standard_hours_source": "summary:תקן",
        "missing_hours_source": "summary:חוסר",
        "start_date_source": "meta:תחילת עבודה",
    }


def build_flamingo_mapping_warnings(mapping, manual_hourly_rate_text):
    warnings = []
    if str(manual_hourly_rate_text or "").strip():
        warnings.append(f"כל העובדים בדוח חושבו לפי התעריף השעתי: {str(manual_hourly_rate_text).strip()}")
    elif not mapping.get("hourly_rate_source"):
        warnings.append("לא נבחר שדה תעריף שעתי ולא הוזן תעריף ידני.")
    if not mapping.get("payable_hours_source"):
        warnings.append("לא נבחר שדה שעות לתשלום בפועל. ללא השדה הזה חישוב השכר לא יהיה תקין.")
    if not mapping.get("worker_number_source") and not mapping.get("id_number_source"):
        warnings.append("לא נבחר מספר עובד או תעודת זהות. הזיהוי יתבסס על שם עובד בלבד ויכול להיות פחות חזק.")
    return warnings


def collect_flamingo_meta_candidates(detail_sheet, workbook_kind):
    candidates = []
    seen = set()
    rows, cols = get_flamingo_sheet_dims(detail_sheet, workbook_kind)
    date_header_row = find_sheet_label_row(detail_sheet, workbook_kind, "תאריך")
    max_meta_row = min(rows, 20)
    if date_header_row >= 0:
        max_meta_row = min(max_meta_row, date_header_row)
    for row_index in range(max_meta_row):
        row_values = [get_flamingo_sheet_cell(detail_sheet, workbook_kind, row_index, c) for c in range(cols)]
        for col_index, raw_label in enumerate(row_values):
            label_text = str(raw_label or "").strip()
            if not label_text:
                continue
            normalized = normalize_token(label_text)
            if normalized in {"שםלתצוגה", "שםעובד", "מחלקה", "מספרבשכר", "מספרעובד", "מסמפעלבשכר", "תעודתזהות", "דרכון", "תחילתעבודה", "הערות"}:
                for next_col in range(col_index + 1, min(cols, col_index + 11)):
                    candidate = row_values[next_col]
                    if candidate in ("", None):
                        continue
                    candidate_token = normalize_token(candidate)
                    if candidate_token in FLAMINGO_META_LABEL_TOKENS or candidate_token in {"תאריך", "יום", "כניסה", "יציאה", "אירוע", "סהכ", "סה\"כ", "תקן", "חוסר"}:
                        continue
                    source = f"meta:{label_text}"
                    if source in seen:
                        break
                    candidates.append(
                        {
                            "value": source,
                            "label": f"שדה עליון: {label_text} (לדוגמה: {stringify_excel_value(candidate)})",
                            "source_kind": "meta",
                            "match_token": normalized,
                        }
                    )
                    seen.add(source)
                    break
    return candidates


def collect_flamingo_summary_candidates(detail_sheet, summary_sheet, workbook_kind):
    search_sheets = []
    if summary_sheet is not None:
        search_sheets.append(summary_sheet)
    if summary_sheet is None or summary_sheet is detail_sheet:
        search_sheets.append(detail_sheet)
    candidates = []
    seen = set()
    relevant_keywords = (
        "נוכחות",
        "תקן",
        "חוסר",
        "שעותלתשלום",
        "שעותמשולמות",
        "רגילות",
        "100",
        "125",
        "150",
        "175",
        "200",
    )
    for sheet in search_sheets:
        rows, cols = get_flamingo_sheet_dims(sheet, workbook_kind)
        summary_start_row = find_sheet_label_row(sheet, workbook_kind, "נתונים כללים")
        start_row = summary_start_row if summary_start_row >= 0 else 0
        for row_index in range(start_row, rows):
            row_values = [get_flamingo_sheet_cell(sheet, workbook_kind, row_index, c) for c in range(cols)]
            for col_index, raw_label in enumerate(row_values):
                label_text = str(raw_label or "").strip()
                if not label_text:
                    continue
                normalized = normalize_token(label_text)
                if not normalized or not any(keyword in normalized for keyword in relevant_keywords):
                    continue
                for next_col in range(col_index + 1, len(row_values)):
                    candidate = row_values[next_col]
                    if candidate in ("", None):
                        continue
                    parsed_hours = try_parse_hours_value(candidate)
                    if parsed_hours is None and not isinstance(candidate, (int, float)):
                        candidate_text = str(candidate).strip()
                        if not candidate_text:
                            continue
                        try:
                            float(candidate_text.replace(",", "."))
                        except ValueError:
                            continue
                    source = f"summary:{label_text}"
                    if source in seen:
                        break
                    candidates.append(
                        {
                            "value": source,
                            "label": f"שדה סיכום: {label_text} (לדוגמה: {stringify_excel_value(candidate)})",
                            "source_kind": "table_exact",
                            "match_token": normalized,
                        }
                    )
                    seen.add(source)
                    break
    return candidates


# ── Department Payroll Mapping ──────────────────────────────────────

def collect_dept_daily_header_candidates(detail_sheet, workbook_kind):
    """Scan the daily data header row (the one with תאריך, כניסה, etc.)
    for columns like מספר מחלקה and סה"כ that are useful for dept_payroll."""
    candidates = []
    seen = set()
    rows, cols = get_flamingo_sheet_dims(detail_sheet, workbook_kind)
    header_row = find_sheet_label_row(detail_sheet, workbook_kind, "תאריך")
    if header_row < 0:
        return candidates
    # Relevant column headers to look for in the daily header row
    relevant_tokens = {
        "מספרמחלקה", "סהכ", "סה\"כ", "סהכשעות",
        "רגילות", "נוכחות", "תקן", "חוסר",
        "אירוע", "כניסה", "יציאה", "תאריך",
        "תעריףשעתי", "תעריף",
    }
    for col_index in range(cols):
        label_text = str(get_flamingo_sheet_cell(detail_sheet, workbook_kind, header_row, col_index) or "").strip()
        if not label_text:
            continue
        normalized = normalize_token(label_text)
        if not any(kw in normalized for kw in relevant_tokens):
            continue
        # Get a sample value from the first data row below
        sample = ""
        for data_row in range(header_row + 1, min(rows, header_row + 6)):
            val = get_flamingo_sheet_cell(detail_sheet, workbook_kind, data_row, col_index)
            if val not in ("", None):
                sample = stringify_excel_value(val)
                break
        source = f"meta:{label_text}"
        if source in seen:
            continue
        candidates.append({
            "value": source,
            "label": f"עמודה יומית: {label_text}" + (f" (לדוגמה: {sample})" if sample else ""),
            "source_kind": "meta",
            "match_token": normalized,
        })
        seen.add(source)
    return candidates


def collect_dept_summary_candidates(detail_sheet, summary_sheet, workbook_kind):
    """Like collect_flamingo_summary_candidates but with broader keyword matching
    to also catch 'ש.לתשלום ואירועים' and similar variants."""
    search_sheets = []
    if summary_sheet is not None:
        search_sheets.append(summary_sheet)
    if summary_sheet is None or summary_sheet is detail_sheet:
        search_sheets.append(detail_sheet)
    candidates = []
    seen = set()
    relevant_keywords = (
        "נוכחות", "תקן", "חוסר",
        "שעותלתשלום", "שעותמשולמות", "שלתשלום", "לתשלום",
        "רגילות", "100", "125", "150", "175", "200",
    )
    for sheet in search_sheets:
        rows, cols = get_flamingo_sheet_dims(sheet, workbook_kind)
        summary_start_row = find_sheet_label_row(sheet, workbook_kind, "נתונים כללים")
        start_row = summary_start_row if summary_start_row >= 0 else 0
        for row_index in range(start_row, rows):
            row_values = [get_flamingo_sheet_cell(sheet, workbook_kind, row_index, c) for c in range(cols)]
            for col_index, raw_label in enumerate(row_values):
                label_text = str(raw_label or "").strip()
                if not label_text:
                    continue
                normalized = normalize_token(label_text)
                # Strip dots for broader matching (ש.לתשלום → שלתשלום)
                normalized_nodot = normalized.replace(".", "")
                if not normalized or not any(keyword in normalized or keyword in normalized_nodot for keyword in relevant_keywords):
                    continue
                found_value = False
                for next_col in range(col_index + 1, len(row_values)):
                    candidate = row_values[next_col]
                    if candidate in ("", None):
                        continue
                    parsed_hours = try_parse_hours_value(candidate)
                    if parsed_hours is None and not isinstance(candidate, (int, float)):
                        candidate_text = str(candidate).strip()
                        if not candidate_text:
                            continue
                        try:
                            float(candidate_text.replace(",", "."))
                        except ValueError:
                            continue
                    source = f"summary:{label_text}"
                    if source in seen:
                        found_value = True
                        break
                    candidates.append({
                        "value": source,
                        "label": f"שדה סיכום: {label_text} (לדוגמה: {stringify_excel_value(candidate)})",
                        "source_kind": "table_exact",
                        "match_token": normalized,
                    })
                    seen.add(source)
                    found_value = True
                    break
                # For key payable-hours labels, add even without a sample value
                if not found_value:
                    payable_keywords = ("לתשלום", "שלתשלום", "שעותלתשלום", "נוכחות")
                    if any(kw in normalized_nodot for kw in payable_keywords):
                        source = f"summary:{label_text}"
                        if source not in seen:
                            candidates.append({
                                "value": source,
                                "label": f"שדה סיכום: {label_text}",
                                "source_kind": "table_exact",
                                "match_token": normalized_nodot,
                            })
                            seen.add(source)
    return candidates


def build_dept_payroll_mapping_options(input_path, extension):
    workbook_kind, workbook = open_excel_workbook(input_path, extension)
    worker_blocks = list(iter_flamingo_worker_blocks(workbook_kind, workbook))
    if not worker_blocks:
        raise ValueError("לא זוהו גיליונות עובדים בדוח הנוכחות")
    # Extract all unique client names from תנועות מיוחדות summary table across all workers
    # Use the same extraction logic as extract_dept_payroll_worker for consistency
    NON_CLIENT_KEYWORDS = {"חופשה", "מחלה", "חג", "מילואים", "אבל", "היעדרות", "סהכ", "סה\"כ"}
    detected_clients = set()
    for detail_sheet_block, summary_sheet_block in worker_blocks:
        tnuot_row, tnuot_col = find_sheet_label_position(detail_sheet_block, workbook_kind, "תנועות מיוחדות")
        if tnuot_row < 0:
            continue
        rows_count, cols_count = get_flamingo_sheet_dims(detail_sheet_block, workbook_kind)
        # Scan only in the column range where the header was found (± a few cols)
        scan_col_start = max(0, tnuot_col - 2)
        scan_col_end = min(cols_count, tnuot_col + 12)
        for r in range(tnuot_row + 1, min(rows_count, tnuot_row + 20)):
            for c in range(scan_col_start, scan_col_end):
                label_val = stringify_excel_value(get_flamingo_sheet_cell(detail_sheet_block, workbook_kind, r, c))
                if not label_val or not label_val.strip():
                    continue
                label_clean = label_val.strip()
                # Skip non-client entries (vacation, sick, holidays, totals)
                label_token = normalize_token(label_clean)
                if any(normalize_token(kw) == label_token or normalize_token(kw) in label_token for kw in NON_CLIENT_KEYWORDS):
                    break
                # Look for a numeric value to the right (confirming it's a real entry with hours)
                for nc in range(c + 1, min(cols_count, c + 10)):
                    val = get_flamingo_sheet_cell(detail_sheet_block, workbook_kind, r, nc)
                    if val in ("", None):
                        continue
                    try:
                        hours_val = parse_hours_value(val)
                        if hours_val is not None and hours_val > 0:
                            detected_clients.add(label_clean)
                        break
                    except (ValueError, TypeError):
                        break
                break  # only process the first non-empty label per row
    detected_clients = sorted(detected_clients)
    detail_sheet, summary_sheet = worker_blocks[0]
    meta_options = collect_flamingo_meta_candidates(detail_sheet, workbook_kind)
    header_options = collect_dept_daily_header_candidates(detail_sheet, workbook_kind)
    summary_options = collect_dept_summary_candidates(detail_sheet, summary_sheet, workbook_kind)
    base_options = meta_options + header_options + summary_options
    options_by_field = {}
    suggestions = {}
    for field in DEPT_PAYROLL_MAPPING_FIELDS:
        field_name = field["name"]
        options = [{"value": "", "label": "ללא — הזנה ידנית" if field.get("has_manual") else "לא נבחר", "source_kind": "empty"}]
        if field_name == "housing_source":
            options.append({"value": "__auto__", "label": "זיהוי אוטומטי מהערות הדוח", "source_kind": "meta"})
        daily_header_fields = {"event_source", "entry_time_source", "exit_time_source", "date_source", "total_hours_source", "hourly_rate_source", "housing_source"}
        for option in base_options:
            token = option.get("match_token", "")
            keywords = DEPT_PAYROLL_SUGGESTION_KEYWORDS.get(field_name, [])
            if field_name == "payable_hours_source":
                if any(keyword in token for keyword in keywords) or option.get("source_kind") == "table_exact":
                    options.append(option)
            elif field_name in daily_header_fields:
                if any(keyword in token for keyword in keywords):
                    options.append(option)
                elif option.get("source_kind") == "meta":
                    options.append(option)
            elif field_name == "dept_number_source":
                options.append(option)
            elif option.get("source_kind") == "meta":
                options.append(option)
        deduped = []
        seen_values = set()
        for option in options:
            if option["value"] in seen_values:
                continue
            deduped.append(option)
            seen_values.add(option["value"])
        options_by_field[field_name] = deduped

        suggested = ""
        keywords = DEPT_PAYROLL_SUGGESTION_KEYWORDS.get(field_name, [])
        if field_name == "payable_hours_source":
            for preferred in ["נוכחות", "שלתשלום", "שעותלתשלום", "שעותמשולמות", "סהכ", "רגילות"]:
                for option in deduped:
                    token = option.get("match_token", "")
                    if preferred in token:
                        suggested = option["value"]
                        break
                if suggested:
                    break
        elif field_name == "dept_number_source":
            for preferred in ["מספרמחלקה"]:
                for option in deduped:
                    token = option.get("match_token", "")
                    if preferred in token:
                        suggested = option["value"]
                        break
                if suggested:
                    break
            if not suggested:
                for option in deduped:
                    token = option.get("match_token", "")
                    if "מחלקה" in token:
                        suggested = option["value"]
                        break
        elif field_name == "housing_source":
            suggested = "__auto__"
        else:
            for option in deduped:
                token = option.get("match_token", "")
                if any(keyword in token for keyword in keywords):
                    suggested = option["value"]
                    break
        suggestions[field_name] = suggested

    return {
        "options_by_field": options_by_field,
        "suggestions": suggestions,
        "suggested_template_name": "תבנית לקוחות",
        "detected_clients": detected_clients,
    }


def build_dept_payroll_mapping_form(script_id, temp_upload_path, temp_upload_ext, inspection, current_mapping, templates, template_name_value, dept_settings_json_value, manual_values=None):
    template_options = '<option value="">ללא תבנית שמורה</option>'
    for template in templates:
        template_options += '<option value="' + str(template["id"]) + '">' + esc(template["name"]) + '</option>'

    mapping_labels = {field["name"]: field["label"] for field in DEPT_PAYROLL_MAPPING_FIELDS}
    template_payload = {
        str(template["id"]): {key: str(value or "") for key, value in template["mapping"].items()}
        for template in templates
    }

    manual_values = manual_values or {}
    manual_hourly_rate_value = str(manual_values.get("manual_hourly_rate", "") or "")
    manual_housing_value = str(manual_values.get("manual_housing", "") or "")

    mapping_fields_html = ""
    for field in DEPT_PAYROLL_MAPPING_FIELDS:
        field_name = field["name"]
        current_value = str(current_mapping.get(field_name, "") or "")
        options = inspection["options_by_field"].get(field_name, [])
        blank_options = [option for option in options if not option.get("value")]
        meta_options = [option for option in options if option.get("source_kind") == "meta"]
        summary_options = [option for option in options if option.get("source_kind") == "table_exact"]

        def render_option(option):
            selected = ' selected' if option["value"] == current_value else ""
            return '<option value="' + esc(option["value"]) + '" data-base-label="' + esc(option["label"]) + '" data-source-kind="' + esc(option.get("source_kind", "empty")) + '"' + selected + '>' + esc(option["label"]) + '</option>'

        select_options = ""
        for option in blank_options:
            select_options += render_option(option)
        if meta_options:
            select_options += '<optgroup label="שדות עליונים">'
            for option in meta_options:
                select_options += render_option(option)
            select_options += '</optgroup>'
        if summary_options:
            select_options += '<optgroup label="שדות סיכום">'
            for option in summary_options:
                select_options += render_option(option)
            select_options += '</optgroup>'

        required_badge = ' <span style="color:#dc2626">*</span>' if field["required"] else ' <span style="color:#94a3b8">(אופציונלי)</span>'
        wrapper_style = ""
        if field.get("critical"):
            wrapper_style = 'background:#fff7ed;border:1px solid #fdba74;border-radius:12px;padding:10px 10px 12px'
        # Manual input for hourly_rate_source and housing_source
        # Show manual input when dropdown is "ללא" (empty) — means user wants to enter manually
        # Hide when a report field is selected — value comes from report
        manual_input_html = ""
        if field_name == "hourly_rate_source":
            show_manual = not current_value  # empty = manual mode
            manual_input_html = (
                '<div data-manual-wrap="hourly_rate" style="' + ('display:block' if show_manual else 'display:none') + ';margin-top:8px">'
                + '<label class="field-label" style="margin-bottom:4px;font-size:12px">תעריף שעתי (ידני)</label>'
                + '<input type="text" name="manual_hourly_rate" value="' + esc(manual_hourly_rate_value) + '" placeholder="לדוגמה 45.5" style="margin-bottom:0;padding:7px 10px;font-size:13px">'
                + '<div style="font-size:11px;color:#9a3412;line-height:1.5;margin-top:4px">תעריף אחיד לכל העובדים. לחישוב לפי הדוח — בחר שדה מהרשימה למעלה.</div>'
                + '</div>'
            )
        elif field_name == "housing_source":
            show_manual = not current_value  # empty = manual mode
            manual_input_html = (
                '<div data-manual-wrap="housing" style="' + ('display:block' if show_manual else 'display:none') + ';margin-top:8px">'
                + '<label class="field-label" style="margin-bottom:4px;font-size:12px">חיוב דירה (ידני)</label>'
                + '<input type="text" name="manual_housing" value="' + esc(manual_housing_value) + '" placeholder="לדוגמה 800" style="margin-bottom:0;padding:7px 10px;font-size:13px">'
                + '<div style="font-size:11px;color:#9a3412;line-height:1.5;margin-top:4px">סכום אחיד לכל העובדים. לזיהוי אוטומטי — בחר &laquo;זיהוי אוטומטי מהערות&raquo; למעלה.</div>'
                + '</div>'
            )
        mapping_fields_html += (
            '<div style="' + wrapper_style + '"><label class="field-label">' + field["label"] + required_badge + '</label>'
            + ('<div style="font-size:12px;color:#9a3412;line-height:1.6;margin:-4px 0 8px">שדה קריטי לחישוב. יש לוודא שזהו המקור הנכון.</div>' if field.get("critical") else '')
            + '<select name="' + field_name + '" data-mapping-field="1" data-field-label="' + esc(field["label"]) + '" style="padding:9px 12px;border:1.5px solid #e2e8f0;border-radius:8px;font-size:13px;font-family:inherit;outline:none;width:100%;margin-bottom:0;background:white;transition:background-color .15s ease,border-color .15s ease,box-shadow .15s ease">'
            + select_options
            + '</select>'
            + manual_input_html
            + '</div>'
        )

    # Department settings section — manual entry table
    dept_data_escaped = esc(dept_settings_json_value or "[]")

    dept_settings_section = (
        '<div style="background:#f0fdf4;border:1px solid #86efac;border-radius:14px;padding:1rem;margin-top:1rem">'
        + '<div style="font-size:15px;font-weight:700;color:#166534;margin-bottom:10px">הגדרת פרטי לקוחות</div>'
        + '<div style="font-size:13px;color:#475569;margin-bottom:12px">הגדר לכל לקוח (לפי שם התנועה המיוחדת בדוח) את תעריף הגביה ופרטי הקשר. שדות עם <span style="color:#dc2626">*</span> הם חובה.</div>'
        + '<div id="deptSettingsTable">'
        + '<div style="overflow-x:auto"><table style="width:100%;border-collapse:collapse;font-size:13px" id="deptTable">'
        + '<thead><tr style="background:#166534;color:white">'
        + '<th style="padding:8px 4px;width:30px;text-align:center">#</th>'
        + '<th style="padding:8px 6px;white-space:nowrap">שם לקוח (תנועה מיוחדת) <span style="color:#fca5a5">*</span></th>'
        + '<th style="padding:8px 6px;white-space:nowrap">מנהל אזור</th>'
        + '<th style="padding:8px 6px;white-space:nowrap">כתובת</th>'
        + '<th style="padding:8px 6px;white-space:nowrap">תעריף גביה <span style="color:#fca5a5">*</span></th>'
        + '<th style="padding:8px 6px;white-space:nowrap">איש קשר</th>'
        + '<th style="padding:8px 6px;white-space:nowrap">טלפון</th>'
        + '<th style="padding:8px 6px;white-space:nowrap">העתקה</th>'
        + '<th style="padding:8px 4px"></th>'
        + '</tr></thead>'
        + '<tbody id="deptTableBody"></tbody>'
        + '</table></div>'
        + '<div style="display:flex;gap:8px;margin-top:10px;flex-wrap:wrap">'
        + '<button type="button" id="addDeptRow" class="btn btn-blue" style="padding:6px 16px;font-size:13px">+ הוסף לקוח</button>'
        + '<label for="deptFileInput" class="btn btn-gray" style="padding:6px 16px;font-size:13px;cursor:pointer;display:inline-flex;align-items:center">טען מקובץ אקסל</label>'
        + '<input type="file" id="deptFileInput" accept=".xls,.xlsx,.csv" style="display:none">'
        + '</div>'
        + '</div>'
        + '<input type="hidden" name="dept_settings_json" id="deptSettingsJsonInput" value="' + dept_data_escaped + '">'
        + '</div>'
    )

    return (
        '<form method="POST" id="deptPayrollMappingForm" enctype="multipart/form-data">'
        + '<input type="hidden" name="flow_mode" value="confirm_mapping">'
        + '<input type="hidden" name="temp_upload_path" value="' + esc(temp_upload_path) + '">'
        + '<input type="hidden" name="temp_upload_ext" value="' + esc(temp_upload_ext) + '">'
        + '<div style="background:#fafcff;border:1px solid #dbeafe;border-radius:14px;padding:1rem;margin-bottom:1rem">'
        + '<div style="font-size:15px;font-weight:700;color:#1e3a8a;margin-bottom:10px">אישור שדות מדוח הנוכחות</div>'
        + '<div style="display:grid;grid-template-columns:260px minmax(0,1fr);gap:14px;align-items:start">'
        # Templates sidebar
        + '<div style="background:#ffffff;border:1px solid #e2e8f0;border-radius:12px;padding:12px">'
        + '<div style="font-size:14px;font-weight:700;color:#0f172a;margin-bottom:10px">תבניות שמורות</div>'
        + '<label class="field-label">בחירת תבנית</label>'
        + '<div style="display:grid;grid-template-columns:minmax(0,1fr) auto;gap:8px;align-items:center;margin-bottom:12px">'
        + '<select id="selectedDeptTemplateId" name="selected_template_id" style="padding:9px 12px;border:1.5px solid #e2e8f0;border-radius:8px;font-size:13px;font-family:inherit;outline:none;width:100%;background:white">' + template_options + '</select>'
        + '<button type="submit" name="mapping_action" value="delete_template" class="btn btn-gray" style="min-width:104px;padding-inline:14px;white-space:nowrap">מחיקה</button>'
        + '</div>'
        + '<input type="hidden" name="save_template" id="deptSaveTemplateFlag" value="0">'
        + '<input type="hidden" name="template_name" id="deptTemplateNameInput" value="">'
        + '<div style="font-size:12px;color:#64748b;line-height:1.7;margin-top:10px">התבנית שומרת את מיפוי השדות וגם את הגדרת פרטי הלקוחות, כך שלא צריך להזין אותם מחדש בכל חודש.</div>'
        + '</div>'
        # Mapping fields
        + '<div>'
        + '<div style="display:flex;gap:8px;flex-wrap:wrap;margin-bottom:12px">'
        + '<span style="display:inline-flex;align-items:center;gap:6px;padding:6px 10px;border-radius:999px;background:#fff7ed;border:1px solid #fdba74;font-size:12px;color:#9a3412">שדה קריטי</span>'
        + '<span style="display:inline-flex;align-items:center;gap:6px;padding:6px 10px;border-radius:999px;background:#eff6ff;border:1px solid #bfdbfe;font-size:12px;color:#1d4ed8">שדה עליון</span>'
        + '<span style="display:inline-flex;align-items:center;gap:6px;padding:6px 10px;border-radius:999px;background:#ecfdf5;border:1px solid #86efac;font-size:12px;color:#166534">שדה סיכום</span>'
        + '</div>'
        + '<div style="display:grid;grid-template-columns:repeat(auto-fit,minmax(220px,1fr));gap:10px;margin-bottom:12px">' + mapping_fields_html + '</div>'
        + '<div style="font-size:12px;color:#64748b;line-height:1.7;margin-bottom:12px">שדות קריטיים: אירוע/לקוח וסה"כ שעות. בלי אלה לא ניתן לבנות דוח חיוב ותשלום.</div>'
        + '</div></div>'
        + '</div>'
        # Clients file upload section
        + '<div style="background:#fef3c7;border:1px solid #fbbf24;border-radius:14px;padding:1rem;margin-top:1rem">'
        + '<div style="display:flex;align-items:center;justify-content:space-between;margin-bottom:10px">'
        + '<div style="font-size:15px;font-weight:700;color:#92400e">קובץ לקוחות (אופציונלי)</div>'
        + '<a href="/sample-clients-file" class="btn btn-gray" style="text-decoration:none;display:inline-flex;align-items:center;gap:6px;padding:6px 14px;font-size:12px;font-weight:600;border-radius:8px;white-space:nowrap" target="_blank">'
        + '<svg width="14" height="14" viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round" stroke-linejoin="round"><path d="M21 15v4a2 2 0 01-2 2H5a2 2 0 01-2-2v-4"/><polyline points="7 10 12 15 17 10"/><line x1="12" y1="15" x2="12" y2="3"/></svg>'
        + 'הורד קובץ לדוגמה</a></div>'
        + '<div style="font-size:13px;color:#78350f;margin-bottom:12px">ניתן להעלות קובץ אקסל עם נתוני לקוחות — תעריפי גביה, כתובות, אנשי קשר. הורד את הקובץ לדוגמה, מלא את הפרטים ושמור.</div>'
        + '<input type="file" name="clients_file" accept=".xls,.xlsx" style="padding:8px;border:1.5px solid #e2e8f0;border-radius:8px;font-size:13px;font-family:inherit;width:100%;background:white">'
        + '</div>'
        # Department settings section
        + dept_settings_section
        # Submit buttons
        + '<div style="display:flex;gap:10px;flex-wrap:wrap;justify-content:center;margin-top:1rem">'
        + '<button type="submit" id="deptPayrollConfirmButton" name="mapping_action" value="confirm" class="btn btn-blue" style="min-width:220px">אשר הכל והפעל חישוב</button>'
        + '<a href="/run/' + script_id + '" class="btn btn-gray" style="text-decoration:none;display:inline-flex;align-items:center;justify-content:center;min-width:150px">העלאת קובץ חדש</a>'
        + '</div>'
        # Save template prompt modal
        + '<div id="deptSaveTemplateModal" style="display:none;position:fixed;inset:0;background:rgba(15,23,42,.45);backdrop-filter:blur(3px);z-index:90;align-items:center;justify-content:center;padding:20px">'
        + '<div style="width:100%;max-width:380px;background:#ffffff;border:1px solid #dbeafe;border-radius:18px;box-shadow:0 20px 50px rgba(15,23,42,.18);padding:28px 24px;text-align:center">'
        + '<div style="font-size:28px;margin-bottom:10px">💾</div>'
        + '<div style="font-size:17px;font-weight:800;color:#1e3a8a;margin-bottom:8px">שמירת תבנית</div>'
        + '<div style="font-size:13px;line-height:1.7;color:#475569;margin-bottom:18px">האם לשמור את ההגדרות הנוכחיות כתבנית?<br>בפעם הבאה תוכל לטעון אותן בלחיצה אחת.</div>'
        + '<div style="display:flex;gap:10px;justify-content:center">'
        + '<button type="button" id="deptSaveYes" class="btn btn-blue" style="min-width:120px;font-size:14px">כן, שמור</button>'
        + '<button type="button" id="deptSaveNo" class="btn btn-gray" style="min-width:120px;font-size:14px">לא, המשך</button>'
        + '</div></div></div>'
        # Processing overlay
        + '<div id="deptPayrollProcessingOverlay" style="display:none;position:fixed;inset:0;background:rgba(248,250,252,.78);backdrop-filter:blur(2px);z-index:80;align-items:center;justify-content:center;padding:20px">'
        + '<div style="width:100%;max-width:320px;background:#ffffff;border:1px solid #dbeafe;border-radius:18px;box-shadow:0 20px 50px rgba(15,23,42,.14);padding:24px 20px;text-align:center">'
        + '<div style="width:42px;height:42px;border-radius:999px;border:3px solid #bfdbfe;border-top-color:#2563eb;margin:0 auto 14px;animation:deptSpin .9s linear infinite"></div>'
        + '<div style="font-size:16px;font-weight:800;color:#1e3a8a;margin-bottom:6px">חישוב השכר בהכנה</div>'
        + '<div style="font-size:13px;line-height:1.7;color:#475569">המערכת מחשבת שכר לכל עובד לפי מחלקה. פעולה זו עשויה להימשך מעט זמן.</div>'
        + '</div></div>'
        # JavaScript
        + '<script>'
        + '(function(){'
        + 'var templateSelect=document.getElementById("selectedDeptTemplateId");'
        + 'var fieldSelects=Array.prototype.slice.call(document.querySelectorAll(\'select[data-mapping-field="1"]\'));'
        + 'var form=document.getElementById("deptPayrollMappingForm");'
        + 'var confirmButton=document.getElementById("deptPayrollConfirmButton");'
        + 'var overlay=document.getElementById("deptPayrollProcessingOverlay");'
        + 'var templateMappings=' + json.dumps(template_payload, ensure_ascii=False) + ';'
        + 'var fieldLabels=' + json.dumps(mapping_labels, ensure_ascii=False) + ';'
        + 'var deptJsonInput=document.getElementById("deptSettingsJsonInput");'
        + 'var deptTableBody=document.getElementById("deptTableBody");'
        + 'var selectStyles={critical:{bg:"#fff7ed",border:"#fb923c",shadow:"rgba(249,115,22,.14)"},meta:{bg:"#eff6ff",border:"#60a5fa",shadow:"rgba(59,130,246,.12)"},table_exact:{bg:"#ecfdf5",border:"#4ade80",shadow:"rgba(34,197,94,.14)"},empty:{bg:"#ffffff",border:"#e2e8f0",shadow:"rgba(148,163,184,.08)"}};'
        + 'function toggleManualInputs(sel){if(sel.name==="hourly_rate_source"){var w=document.querySelector("[data-manual-wrap=\'hourly_rate\']");if(w)w.style.display=(!sel.value)?"block":"none";}if(sel.name==="housing_source"){var w2=document.querySelector("[data-manual-wrap=\'housing\']");if(w2)w2.style.display=(!sel.value)?"block":"none";}}'
        + 'function applySelectVisual(sel){var opt=sel.options[sel.selectedIndex];var kind=(opt&&opt.getAttribute("data-source-kind"))||"empty";var style=selectStyles[kind]||selectStyles.empty;sel.style.backgroundColor=style.bg;sel.style.borderColor=style.border;sel.style.boxShadow="0 0 0 3px "+style.shadow;toggleManualInputs(sel);}'
        + 'function refreshOptionLabels(){var assignments={};fieldSelects.forEach(function(sel){if(sel.value){assignments[sel.value]=sel.name;}});fieldSelects.forEach(function(sel){Array.prototype.forEach.call(sel.options,function(opt){var base=opt.getAttribute("data-base-label")||opt.text;var assigned=assignments[opt.value];var suffix="";if(opt.value && assigned && assigned!==sel.name && opt.value!=="__auto__"){suffix=" [נבחר עבור "+(fieldLabels[assigned]||assigned)+"]";}opt.text=base+suffix;});applySelectVisual(sel);});}'
        + 'function clearDuplicateSelections(changedSelect){if(!changedSelect.value || changedSelect.value==="__auto__"){refreshOptionLabels();return;}fieldSelects.forEach(function(sel){if(sel!==changedSelect && sel.value===changedSelect.value){sel.value="";}});refreshOptionLabels();}'
        # Department table functions
        + 'var deptFields=["customer_name","region_manager","address","charge_rate","contact_person","contact_phone"];'
        + 'var copyableFields=["region_manager","address","charge_rate","contact_person","contact_phone"];'
        + 'var detectedClients=' + json.dumps(inspection.get("detected_clients", []), ensure_ascii=False) + ';'
        # renumberRows: update row numbers and copy-dropdown options
        + 'function renumberRows(){var trs=deptTableBody.querySelectorAll("tr");trs.forEach(function(tr,i){var numCell=tr.querySelector("[data-row-num]");if(numCell)numCell.textContent=(i+1);});}'
        # getRowData: read all field values from a <tr>
        + 'function getRowData(tr){var d={};tr.querySelectorAll("input[data-dept-field]").forEach(function(inp){d[inp.getAttribute("data-dept-field")]=inp.value;});return d;}'
        # buildCopyMenu: create the copy dropdown for a row
        + 'function buildCopyMenu(tr){var td=document.createElement("td");td.style.padding="4px 3px";td.style.whiteSpace="nowrap";'
        + 'var sel=document.createElement("select");sel.style.cssText="padding:4px 6px;border:1.5px solid #e2e8f0;border-radius:6px;font-size:11px;font-family:inherit;background:white;cursor:pointer;min-width:90px";'
        + 'sel.innerHTML=\'<option value="">העתק מ...</option>\';'
        + 'sel.addEventListener("change",function(){if(!sel.value)return;var parts=sel.value.split(":");var mode=parts[0];var srcIdx=parseInt(parts[1]);'
        + 'var allTrs=Array.prototype.slice.call(deptTableBody.querySelectorAll("tr"));if(srcIdx<0||srcIdx>=allTrs.length){sel.value="";return;}'
        + 'var srcData=getRowData(allTrs[srcIdx]);var targetInputs=tr.querySelectorAll("input[data-dept-field]");'
        + 'if(mode==="all"){targetInputs.forEach(function(inp){var f=inp.getAttribute("data-dept-field");if(f!=="customer_name"&&srcData[f]){inp.value=srcData[f];}});}'
        + 'else{targetInputs.forEach(function(inp){if(inp.getAttribute("data-dept-field")===mode&&srcData[mode]){inp.value=srcData[mode];}});}'
        + 'sel.value="";syncDeptJson();});'
        + 'td.appendChild(sel);return td;}'
        # refreshCopyMenus: rebuild all copy dropdown options based on current rows
        + 'function refreshCopyMenus(){var trs=Array.prototype.slice.call(deptTableBody.querySelectorAll("tr"));'
        + 'trs.forEach(function(tr,myIdx){var sel=tr.querySelector("select");if(!sel)return;'
        + 'var prev=sel.value;sel.innerHTML=\'<option value="">העתק מ...</option>\';'
        + 'trs.forEach(function(otherTr,otherIdx){if(otherIdx===myIdx)return;'
        + 'var od=getRowData(otherTr);var label=od.customer_name||(\"שורה \"+(otherIdx+1));'
        + 'var hasData=copyableFields.some(function(f){return od[f];});if(!hasData)return;'
        + 'sel.innerHTML+=\'<option value="all:\'+otherIdx+\'">כל הפרטים מ\'+label+\'</option>\';'
        + 'copyableFields.forEach(function(f){if(!od[f])return;'
        + 'var fLabel={region_manager:"מנהל אזור",address:"כתובת",charge_rate:"תעריף גביה",contact_person:"איש קשר",contact_phone:"טלפון"}[f]||f;'
        + 'sel.innerHTML+=\'<option value="\'+f+\':\'+otherIdx+\'">\'+fLabel+\' מ\'+label+\'</option>\';});'
        + '});sel.value=prev;});}'
        # createDeptRow: build a single table row
        + 'function createDeptRow(data){data=data||{};var tr=document.createElement("tr");tr.style.borderBottom="1px solid #e2e8f0";'
        # Row number cell
        + 'var tdNum=document.createElement("td");tdNum.style.cssText="padding:4px 3px;text-align:center;font-weight:700;color:#64748b;font-size:13px";tdNum.setAttribute("data-row-num","1");tr.appendChild(tdNum);'
        # Field cells
        + 'deptFields.forEach(function(f){var td=document.createElement("td");td.style.padding="4px 3px";'
        + 'if(f==="customer_name"){var inp=document.createElement("input");inp.type="text";inp.setAttribute("data-dept-field",f);inp.setAttribute("list","detectedClientsList");inp.value=data[f]||"";inp.placeholder="בחר מהרשימה...";inp.style.cssText="width:100%;padding:6px 8px;border:1.5px solid #86efac;border-radius:6px;font-size:12px;font-family:inherit;outline:none;min-width:120px;background:#f0fdf4";inp.addEventListener("input",function(){syncDeptJson();refreshCopyMenus();});td.appendChild(inp);}'
        + 'else{var inp=document.createElement("input");inp.type="text";inp.setAttribute("data-dept-field",f);inp.value=data[f]||"";inp.style.cssText="width:100%;padding:6px 8px;border:1.5px solid #e2e8f0;border-radius:6px;font-size:12px;font-family:inherit;outline:none;min-width:80px";inp.addEventListener("input",function(){syncDeptJson();refreshCopyMenus();});td.appendChild(inp);}'
        + 'tr.appendChild(td);});'
        # Copy menu cell
        + 'var copyTd=buildCopyMenu(tr);tr.appendChild(copyTd);'
        # Delete button cell
        + 'var tdDel=document.createElement("td");tdDel.style.padding="4px 3px";var btn=document.createElement("button");btn.type="button";btn.textContent="✕";btn.style.cssText="background:#fee2e2;color:#dc2626;border:none;border-radius:6px;padding:4px 10px;cursor:pointer;font-size:13px";btn.addEventListener("click",function(){tr.remove();syncDeptJson();renumberRows();refreshCopyMenus();});tdDel.appendChild(btn);tr.appendChild(tdDel);'
        + 'deptTableBody.appendChild(tr);renumberRows();return tr;}'
        + 'function syncDeptJson(){var rows=deptTableBody.querySelectorAll("tr");var data=[];rows.forEach(function(tr){var entry={};var inputs=tr.querySelectorAll("input[data-dept-field]");inputs.forEach(function(inp){entry[inp.getAttribute("data-dept-field")]=inp.value;});if(entry.customer_name){data.push(entry);}});deptJsonInput.value=JSON.stringify(data);}'
        + 'function loadDeptRows(arr){deptTableBody.innerHTML="";(arr||[]).forEach(function(d){createDeptRow(d);});if(!arr||arr.length===0){createDeptRow();}syncDeptJson();refreshCopyMenus();}'
        # Build datalist for detected clients
        + '(function(){var dl=document.createElement("datalist");dl.id="detectedClientsList";detectedClients.forEach(function(c){var opt=document.createElement("option");opt.value=c;dl.appendChild(opt);});document.body.appendChild(dl);})();'
        # Load initial data — if empty, seed rows from detected clients
        + 'try{var initData=JSON.parse(deptJsonInput.value||"[]");if(initData.length>0){loadDeptRows(initData);}else if(detectedClients.length>0){var seeded=detectedClients.map(function(c){return{customer_name:c};});loadDeptRows(seeded);}else{createDeptRow();}}catch(e){createDeptRow();}'
        + 'document.getElementById("addDeptRow").addEventListener("click",function(){createDeptRow();syncDeptJson();refreshCopyMenus();});'
        # File upload for dept settings
        + 'document.getElementById("deptFileInput").addEventListener("change",function(e){'
        + 'var file=e.target.files[0];if(!file)return;'
        + 'var reader=new FileReader();reader.onload=function(ev){'
        + 'try{'
        + 'var text=ev.target.result;var lines=text.split("\\n").filter(function(l){return l.trim();});'
        + 'if(lines.length<2){alert("הקובץ ריק או לא תקין");return;}'
        + 'var headers=lines[0].split(",");var data=[];'
        + 'for(var i=1;i<lines.length;i++){var vals=lines[i].split(",");var entry={};'
        + 'for(var j=0;j<deptFields.length&&j<vals.length;j++){entry[deptFields[j]]=(vals[j]||"").replace(/^"|"$/g,"").trim();}'
        + 'if(entry.customer_name)data.push(entry);}'
        + 'loadDeptRows(data);'
        + '}catch(err){alert("שגיאה בקריאת הקובץ: "+err.message);}'
        + '};'
        + 'if(file.name.endsWith(".csv")){reader.readAsText(file,"UTF-8");}else{alert("לטעינה מקובץ יש להשתמש בפורמט CSV. ניתן לשמור אקסל כ-CSV.");}'
        + '});'
        # Template apply/load
        + 'function applyTemplate(templateId){var mapping=templateMappings[templateId]||{};if(!templateId){refreshOptionLabels();return;}fieldSelects.forEach(function(sel){sel.value=mapping[sel.name]||"";});var seen={};fieldSelects.forEach(function(sel){if(sel.value && seen[sel.value]){sel.value="";}else if(sel.value){seen[sel.value]=true;}});refreshOptionLabels();if(mapping.dept_settings_json){try{loadDeptRows(JSON.parse(mapping.dept_settings_json));}catch(e){}}}'
        + 'fieldSelects.forEach(function(sel){sel.addEventListener("change",function(){clearDuplicateSelections(sel);});});'
        + 'if(templateSelect){templateSelect.addEventListener("change",function(){applyTemplate(this.value);});}'
        # Submit handler — show save-template prompt before processing
        + 'var saveModal=document.getElementById("deptSaveTemplateModal");'
        + 'var saveFlag=document.getElementById("deptSaveTemplateFlag");'
        + 'var saveNameInput=document.getElementById("deptTemplateNameInput");'
        + 'var pendingSubmit=false;'
        + 'function doActualSubmit(){syncDeptJson();if(confirmButton){confirmButton.disabled=true;confirmButton.textContent="החישוב התחיל...";}if(overlay){overlay.style.display="flex";}document.body.style.overflow="hidden";pendingSubmit=true;var h=document.createElement("input");h.type="hidden";h.name="mapping_action";h.value="confirm";form.appendChild(h);form.submit();}'
        + 'document.getElementById("deptSaveYes").addEventListener("click",function(){saveFlag.value="1";var now=new Date();var pad=function(n){return n<10?"0"+n:n;};saveNameInput.value="תבנית "+pad(now.getDate())+"/"+pad(now.getMonth()+1)+"/"+now.getFullYear()+" "+pad(now.getHours())+":"+pad(now.getMinutes());saveModal.style.display="none";doActualSubmit();});'
        + 'document.getElementById("deptSaveNo").addEventListener("click",function(){saveFlag.value="0";saveNameInput.value="";saveModal.style.display="none";doActualSubmit();});'
        + 'if(form){form.addEventListener("submit",function(event){var submitter=event.submitter||document.activeElement;if(!submitter||submitter.value!=="confirm"){return;}if(pendingSubmit){return;}event.preventDefault();syncDeptJson();saveModal.style.display="flex";});}'
        + 'refreshOptionLabels();'
        + '})();'
        + '</script>'
        + '<style>@keyframes deptSpin{from{transform:rotate(0deg);}to{transform:rotate(360deg);}}</style>'
        + '</form>'
    )


def build_flamingo_mapping_options(input_path, extension):
    workbook_kind, workbook = open_excel_workbook(input_path, extension)
    worker_blocks = list(iter_flamingo_worker_blocks(workbook_kind, workbook))
    if not worker_blocks:
        raise ValueError("Could not identify worker sheets in this payroll report")
    detail_sheet, summary_sheet = worker_blocks[0]
    meta_options = collect_flamingo_meta_candidates(detail_sheet, workbook_kind)
    summary_options = collect_flamingo_summary_candidates(detail_sheet, summary_sheet, workbook_kind)
    base_options = meta_options + summary_options
    options_by_field = {}
    suggestions = {}
    for field in FLAMINGO_MAPPING_FIELDS:
        field_name = field["name"]
        options = [{"value": "", "label": "לא נבחר", "source_kind": "empty"}]
        if field_name == "hourly_rate_source":
            options.append({"value": "__manual__", "label": "הזנה ידנית של תעריף שעתי", "source_kind": "critical"})
        for option in base_options:
            token = option.get("match_token", "")
            keywords = FLAMINGO_SUGGESTION_KEYWORDS.get(field_name, [])
            if field_name == "payable_hours_source":
                if any(keyword in token for keyword in keywords):
                    options.append(option)
            elif field_name == "hourly_rate_source":
                if option.get("source_kind") == "meta":
                    options.append(option)
            elif field_name in {"attendance_hours_source", "standard_hours_source", "missing_hours_source"}:
                if any(keyword in token for keyword in keywords):
                    options.append(option)
            elif option.get("source_kind") == "meta":
                options.append(option)
        deduped = []
        seen_values = set()
        for option in options:
            if option["value"] in seen_values:
                continue
            deduped.append(option)
            seen_values.add(option["value"])
        options_by_field[field_name] = deduped

        suggested = ""
        if field_name == "payable_hours_source":
            preferred_order = ["שעותלתשלום", "שעותמשולמות", "נוכחות", "רגילות"]
            for preferred in preferred_order:
                for option in deduped:
                    token = option.get("match_token", "")
                    if preferred in token:
                        suggested = option["value"]
                        break
                if suggested:
                    break
        else:
            keywords = FLAMINGO_SUGGESTION_KEYWORDS.get(field_name, [])
            if field_name == "hourly_rate_source":
                for preferred in ["תעריף", "rate", "שעה"]:
                    for option in deduped:
                        token = option.get("match_token", "")
                        if preferred in token:
                            suggested = option["value"]
                            break
                    if suggested:
                        break
            else:
                for option in deduped:
                    token = option.get("match_token", "")
                    if any(keyword in token for keyword in keywords):
                        suggested = option["value"]
                        break
        suggestions[field_name] = suggested

    return {
        "options_by_field": options_by_field,
        "suggestions": suggestions,
        "suggested_template_name": "תבנית שכר",
    }


def detect_rimon_header_row(sheet, workbook_kind):
    rows, cols = get_excel_dims(sheet, workbook_kind)
    best_row = 11 if rows > 11 else 0
    best_score = -1
    for row_index in range(min(rows, 25)):
        score = 0
        for col_index in range(cols):
            token = normalize_token(stringify_excel_value(get_excel_cell(sheet, workbook_kind, row_index, col_index, "")))
            if token in {"תאריך", "אירוע", "כניסה", "יציאה", "סהכ", "שגיאות", "שגיאה"}:
                score += 3
            elif token:
                score += 1
        if score > best_score:
            best_score = score
            best_row = row_index
    return best_row


def find_rimon_first_sample(sheet, workbook_kind, col_index, start_row):
    rows, _ = get_excel_dims(sheet, workbook_kind)
    for row_index in range(start_row + 1, min(rows, start_row + 10)):
        text = stringify_excel_value(get_excel_cell(sheet, workbook_kind, row_index, col_index, ""))
        if text:
            return text
    return ""


def resolve_rimon_column_header(sheet, workbook_kind, header_row, col_index):
    header_text = stringify_excel_value(get_excel_cell(sheet, workbook_kind, header_row, col_index, ""))
    if header_text:
        return header_text, False
    for offset in (1, -1, 2, -2):
        candidate = stringify_excel_value(get_excel_cell(sheet, workbook_kind, header_row, col_index + offset, ""))
        if candidate:
            return candidate, True
    return "", False


def find_rimon_meta_value(sheet, workbook_kind, labels, fallback_cells=()):
    rows, cols = get_excel_dims(sheet, workbook_kind)
    normalized_labels = {normalize_token(label) for label in labels}
    for row_index in range(min(rows, 12)):
        for col_index in range(cols):
            token = normalize_token(stringify_excel_value(get_excel_cell(sheet, workbook_kind, row_index, col_index, "")))
            if token in normalized_labels:
                for next_col in range(col_index + 1, min(cols, col_index + 6)):
                    candidate = stringify_excel_value(get_excel_cell(sheet, workbook_kind, row_index, next_col, ""))
                    if candidate and normalize_token(candidate) in RIMON_META_LABEL_TOKENS:
                        break
                    if candidate:
                        return candidate
    for row_index, col_index in fallback_cells:
        candidate = stringify_excel_value(get_excel_cell(sheet, workbook_kind, row_index, col_index, ""))
        if candidate:
            return candidate
    return ""


def find_rimon_meta_candidates(sheet, workbook_kind, labels, fallback_cells=()):
    rows, cols = get_excel_dims(sheet, workbook_kind)
    normalized_labels = {normalize_token(label) for label in labels}
    candidates = []
    seen_sources = set()

    for row_index in range(min(rows, 12)):
        for col_index in range(cols):
            raw_label = stringify_excel_value(get_excel_cell(sheet, workbook_kind, row_index, col_index, ""))
            token = normalize_token(raw_label)
            if token not in normalized_labels:
                continue
            for next_col in range(col_index + 1, min(cols, col_index + 8)):
                candidate = stringify_excel_value(get_excel_cell(sheet, workbook_kind, row_index, next_col, ""))
                candidate_token = normalize_token(candidate)
                if candidate and candidate_token in RIMON_META_LABEL_TOKENS:
                    break
                if candidate:
                    source = f"meta_cell:{row_index}:{next_col}"
                    if source not in seen_sources:
                        candidates.append(
                            {
                                "value": source,
                                "label": f"שדה עליון: {raw_label} (לדוגמה: {candidate})",
                                "source_kind": "meta",
                            }
                        )
                        seen_sources.add(source)
                    break

    return candidates


def build_rimon_mapping_options(input_path, extension):
    workbook_kind, workbook = open_excel_workbook(input_path, extension)
    sheets = iter_excel_sheets(workbook_kind, workbook)
    first_sheet = sheets[0]
    header_row = detect_rimon_header_row(first_sheet, workbook_kind)
    _, cols = get_excel_dims(first_sheet, workbook_kind)
    table_options = []
    seen_values = set()
    for col_index in range(cols):
        exact_header = stringify_excel_value(get_excel_cell(first_sheet, workbook_kind, header_row, col_index, ""))
        header_text, from_nearby = resolve_rimon_column_header(first_sheet, workbook_kind, header_row, col_index)
        sample = find_rimon_first_sample(first_sheet, workbook_kind, col_index, header_row)
        if not header_text and not sample:
            continue
        column_letter = get_column_letter(col_index + 1)
        resolved_header = header_text or f"עמודה {column_letter}"
        label = f"עמודה {column_letter} - {resolved_header}"
        if from_nearby:
            label += " (הכותרת נלקחה מהעמודה ליד)"
        if sample:
            label += f" (לדוגמה: {sample})"
        option_value = f"col:{col_index}"
        if option_value in seen_values:
            continue
        seen_values.add(option_value)
        table_options.append(
            {
                "value": option_value,
                "label": label,
                "header": resolved_header,
                "exact_header": exact_header,
                "from_nearby": from_nearby,
                "sample": sample,
                "source_kind": "table_nearby" if from_nearby else "table_exact",
            }
        )

    meta_options = {
        "employee_name_source": [],
        "payroll_number_source": [],
        "department_source": [],
        "id_number_source": [],
    }
    alias_meta_values = {
        "employee_name_source": find_rimon_meta_value(first_sheet, workbook_kind, ["שם לתצוגה", "שם עובד"]),
        "payroll_number_source": find_rimon_meta_value(first_sheet, workbook_kind, ["מספר שכר", "מספר עובד", "מספר בשכר", "תג עובד"]),
        "department_source": find_rimon_meta_value(first_sheet, workbook_kind, ["מחלקה"]),
        "id_number_source": find_rimon_meta_value(first_sheet, workbook_kind, ["תעודת זהות", "דרכון"]),
    }
    alias_meta_labels = {
        "employee_name_source": "שדה עליון: שם עובד",
        "payroll_number_source": "שדה עליון: מספר עובד",
        "department_source": "שדה עליון: מחלקה",
        "id_number_source": "שדה עליון: תעודת זהות",
    }
    candidate_meta_labels = {
        "employee_name_source": (["שם לתצוגה", "שם עובד"], []),
        "payroll_number_source": (["מספר שכר", "מספר עובד", "מספר בשכר", "תג עובד"], []),
        "department_source": (["מחלקה"], []),
        "id_number_source": (["תעודת זהות", "דרכון"], []),
    }
    for field_name, field_value in alias_meta_values.items():
        if field_value:
            meta_options[field_name].append(
                {
                    "value": "meta:" + field_name.replace("_source", ""),
                    "label": alias_meta_labels[field_name] + f" (לדוגמה: {field_value})",
                    "source_kind": "meta",
                }
            )
        labels, fallback_cells = candidate_meta_labels[field_name]
        meta_options[field_name].extend(find_rimon_meta_candidates(first_sheet, workbook_kind, labels))

    visible_table_options = filter_rimon_table_options_for_display(table_options)
    options_by_field = {}
    suggestions = {}
    for field in RIMON_MAPPING_FIELDS:
        field_name = field["name"]
        options = [{"value": "", "label": "לא נבחר", "source_kind": "empty"}]
        options.extend(meta_options.get(field_name, []))
        options.extend(filter_rimon_options_for_field(field_name, visible_table_options))
        options_by_field[field_name] = options

        suggested = ""
        if meta_options.get(field_name):
            suggested = meta_options[field_name][0]["value"]
        if not suggested:
            keywords = RIMON_SUGGESTION_KEYWORDS.get(field_name, [])
            ranked_options = sorted(visible_table_options, key=lambda option: (option.get("from_nearby", False), option["value"]))
            if field_name == "date_source":
                for option in ranked_options:
                    header_token = normalize_token(option["header"])
                    exact_token = normalize_token(option.get("exact_header", ""))
                    if (
                        (any(keyword in exact_token for keyword in keywords) or any(keyword in header_token for keyword in keywords))
                        and looks_like_excel_date_sample(option.get("sample", ""))
                    ):
                        suggested = option["value"]
                        break
            elif field_name == "day_name_source":
                for option in ranked_options:
                    header_token = normalize_token(option["header"])
                    exact_token = normalize_token(option.get("exact_header", ""))
                    if (
                        (header_token == "יום" or exact_token == "יום" or any(keyword in exact_token for keyword in keywords))
                        and looks_like_day_name_sample(option.get("sample", ""))
                    ):
                        suggested = option["value"]
                        break
            elif field_name in {"entry_time_source", "exit_time_source", "total_hours_source"}:
                for option in ranked_options:
                    header_token = normalize_token(option["header"])
                    exact_token = normalize_token(option.get("exact_header", ""))
                    if (
                        (any(keyword in exact_token for keyword in keywords) or (not option.get("from_nearby") and any(keyword in header_token for keyword in keywords)))
                        and looks_like_time_sample(option.get("sample", ""))
                    ):
                        suggested = option["value"]
                        break
            if not suggested:
                for option in ranked_options:
                    header_token = normalize_token(option["header"])
                    exact_token = normalize_token(option.get("exact_header", ""))
                    if any(keyword in exact_token for keyword in keywords):
                        suggested = option["value"]
                        break
                    if not option.get("from_nearby") and any(keyword in header_token for keyword in keywords):
                        suggested = option["value"]
                        break
            if not suggested:
                for option in ranked_options:
                    token = normalize_token(option["header"])
                    if any(keyword in token for keyword in keywords):
                        suggested = option["value"]
                        break
        suggestions[field_name] = suggested

    if workbook_kind == "xlsx":
        workbook.close()
    suggested_name_parts = []
    for source in (suggestions.get("employee_name_source", ""), suggestions.get("event_source", ""), suggestions.get("date_source", "")):
        for options in options_by_field.values():
            for option in options:
                if option["value"] == source and option["label"] != "לא נבחר":
                    raw = option["label"].split(" - ", 1)[-1].split(" (", 1)[0]
                    if raw and raw not in suggested_name_parts:
                        suggested_name_parts.append(raw)
                    break
    suggested_template_name = " / ".join(suggested_name_parts[:2]) or "תבנית רימון"
    return {
        "header_row": header_row,
        "options_by_field": options_by_field,
        "suggestions": suggestions,
        "suggested_template_name": suggested_template_name,
    }


def get_mapping_templates(user_id, script_id):
    with get_db() as db:
        rows = db.execute(
            "SELECT * FROM mapping_templates WHERE user_id=? AND script_id=? ORDER BY updated_at DESC, id DESC",
            (user_id, script_id),
        ).fetchall()
    templates = []
    for row in rows:
        try:
            mapping = json.loads(row["mapping_json"] or "{}")
        except Exception:
            mapping = {}
        templates.append({"id": row["id"], "name": row["name"], "mapping": mapping})
    return templates


def get_next_mapping_template_name(templates):
    max_index = 0
    for template in templates:
        match = re.fullmatch(r"\s*תבנית:\s*(\d+)\s*", str(template.get("name", "")))
        if match:
            max_index = max(max_index, int(match.group(1)))
    return f"תבנית: {max_index + 1}"


def save_mapping_template(user_id, script_id, template_name, mapping):
    now_text = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    mapping_json = json.dumps(mapping, ensure_ascii=False)
    with get_db() as db:
        db.execute(
            "INSERT INTO mapping_templates(user_id, script_id, name, mapping_json, created_at, updated_at) VALUES (?,?,?,?,?,?)",
            (user_id, script_id, template_name, mapping_json, now_text, now_text),
        )
        db.commit()


def delete_mapping_template(user_id, script_id, template_id):
    selected_id = str(template_id or "").strip()
    if not selected_id:
        return False
    with get_db() as db:
        db.execute(
            "DELETE FROM mapping_templates WHERE id=? AND user_id=? AND script_id=?",
            (int(selected_id), user_id, script_id),
        )
        db.commit()
    return True


def apply_selected_template(default_mapping, templates, template_id):
    selected_id = str(template_id or "").strip()
    if not selected_id:
        return dict(default_mapping), None
    for template in templates:
        if str(template["id"]) == selected_id:
            merged = dict(default_mapping)
            merged.update(template["mapping"])
            return merged, template
    return dict(default_mapping), None


def build_matan_corrections_mapping_options(input_path, extension):
    workbook = xlrd.open_workbook(input_path)
    first_sheet = workbook.sheets()[0]
    struct = _detect_corrections_daily_structure(workbook, first_sheet)
    header_row = struct["header_row"]
    data_start = struct["data_start"]

    table_options = []
    seen_values = set()
    for col_index in range(first_sheet.ncols):
        header_text = str(first_sheet.cell_value(header_row, col_index)).strip() if header_row < first_sheet.nrows else ""
        sample = ""
        for r in range(data_start, min(data_start + 20, first_sheet.nrows)):
            v = str(first_sheet.cell_value(r, col_index)).strip()
            if v:
                sample = v
                break
        if not header_text and not sample:
            continue
        column_letter = get_column_letter(col_index + 1)
        label = f"עמודה {column_letter}"
        if header_text:
            label += f" - {header_text}"
        if sample:
            label += f" (לדוגמה: {sample})"
        option_value = f"col:{col_index}"
        if option_value in seen_values:
            continue
        seen_values.add(option_value)
        table_options.append({
            "value": option_value,
            "label": label,
            "header": header_text,
            "sample": sample,
            "source_kind": "table_exact",
        })

    suggestions = {}
    if struct["entry_col"] >= 0:
        suggestions["entry_col_source"] = f"col:{struct['entry_col']}"
    if struct["exit_col"] >= 0:
        suggestions["exit_col_source"] = f"col:{struct['exit_col']}"
    if struct["date_col"] >= 0:
        suggestions["date_col_source"] = f"col:{struct['date_col']}"

    blank_option = [{"value": "", "label": "לא נבחר", "source_kind": "empty"}]
    options_by_field = {field["name"]: blank_option + list(table_options) for field in MATAN_CORRECTIONS_MAPPING_FIELDS}
    return {"options_by_field": options_by_field, "suggestions": suggestions}


def build_matan_corrections_mapping_form(script_id, temp_upload_path, temp_upload_ext, inspection, current_mapping, templates, template_name_value, current_filters):
    template_options = '<option value="">ללא תבנית שמורה</option>'
    for template in templates:
        template_options += '<option value="' + str(template["id"]) + '">' + esc(template["name"]) + '</option>'

    mapping_labels = {field["name"]: field["label"] for field in MATAN_CORRECTIONS_MAPPING_FIELDS}
    template_payload = {
        str(template["id"]): {key: str(value or "") for key, value in template["mapping"].items()}
        for template in templates
    }

    mapping_fields_html = ""
    for field in MATAN_CORRECTIONS_MAPPING_FIELDS:
        field_name = field["name"]
        current_value = str(current_mapping.get(field_name, "") or "")
        options = inspection["options_by_field"].get(field_name, [])
        blank_options = [o for o in options if not o.get("value")]
        table_exact_options = [o for o in options if o.get("source_kind") == "table_exact"]

        select_options = ""

        def render_opt(option):
            selected = " selected" if option["value"] == current_value else ""
            return (
                '<option value="' + esc(option["value"]) + '" data-base-label="' + esc(option["label"]) + '" data-source-kind="' + esc(option.get("source_kind", "empty")) + '"' + selected + ">"
                + esc(option["label"])
                + "</option>"
            )

        for option in blank_options:
            select_options += render_opt(option)
        if table_exact_options:
            select_options += '<optgroup label="עמודות מהקובץ">'
            for option in table_exact_options:
                select_options += render_opt(option)
            select_options += '</optgroup>'

        is_critical = field.get("critical", False)
        required_badge = ' <span style="color:#dc2626">*</span>' if field["required"] else ' <span style="color:#94a3b8">(אופציונלי)</span>'
        border_style = "border:1.5px solid #fca5a5;" if is_critical else "border:1.5px solid #e2e8f0;"
        mapping_fields_html += (
            '<div><label class="field-label">' + field["label"] + required_badge + '</label>'
            + '<select name="' + field_name + '" data-mapping-field="1" data-field-label="' + esc(field["label"]) + '" style="padding:9px 12px;' + border_style + 'border-radius:8px;font-size:13px;font-family:inherit;outline:none;width:100%;margin-bottom:0;background:white;transition:background-color .15s ease,border-color .15s ease,box-shadow .15s ease">'
            + select_options
            + '</select></div>'
        )

    min_corrections_val = esc(str(current_filters.get("min_corrections", "") or ""))
    max_corrections_val = esc(str(current_filters.get("max_corrections", "") or ""))
    filters_html = (
        '<div style="margin-top:12px;padding-top:12px;border-top:1px solid #e2e8f0">'
        + '<div style="font-size:13px;font-weight:600;color:#334155;margin-bottom:8px">סינון תוצאות (אופציונלי)</div>'
        + '<div style="display:grid;grid-template-columns:1fr 1fr;gap:10px">'
        + '<div><label class="field-label">מינימום תיקונים</label><input type="text" name="min_corrections" value="' + min_corrections_val + '" placeholder="לדוגמה 4" style="margin-bottom:0"></div>'
        + '<div><label class="field-label">מקסימום תיקונים</label><input type="text" name="max_corrections" value="' + max_corrections_val + '" placeholder="לדוגמה 12" style="margin-bottom:0"></div>'
        + '</div></div>'
    )

    return (
        '<form method="POST" id="mappingConfirmForm">'
        + '<input type="hidden" name="flow_mode" value="confirm_mapping">'
        + '<input type="hidden" name="temp_upload_path" value="' + esc(temp_upload_path) + '">'
        + '<input type="hidden" name="temp_upload_ext" value="' + esc(temp_upload_ext) + '">'
        + '<div style="background:#fafcff;border:1px solid #dbeafe;border-radius:14px;padding:1rem;margin-bottom:1rem">'
        + '<div style="font-size:15px;font-weight:700;color:#1e3a8a;margin-bottom:10px">אישור שדות לפני עיבוד</div>'
        + '<div style="font-size:12px;color:#475569;background:#f0fdf4;border:1px solid #86efac;border-radius:8px;padding:8px 12px;margin-bottom:12px">שם עובד, מחלקה, ת.ז ומספר תג מזוהים אוטומטית מכותרות כל גיליון. יש לאשר כאן רק את עמודות הכניסה, היציאה והתאריך בטבלה היומית.</div>'
        + '<div style="display:grid;grid-template-columns:260px minmax(0,1fr);gap:14px;align-items:start">'
        + '<div style="background:#ffffff;border:1px solid #e2e8f0;border-radius:12px;padding:12px">'
        + '<div style="font-size:14px;font-weight:700;color:#0f172a;margin-bottom:10px">תבניות שמורות</div>'
        + '<label class="field-label">בחירת תבנית</label>'
        + '<div style="display:grid;grid-template-columns:minmax(0,1fr) auto;gap:8px;align-items:center;margin-bottom:12px">'
        + '<select id="selectedTemplateId" name="selected_template_id" style="padding:9px 12px;border:1.5px solid #e2e8f0;border-radius:8px;font-size:13px;font-family:inherit;outline:none;width:100%;background:white">' + template_options + '</select>'
        + '<button type="submit" name="mapping_action" value="delete_template" class="btn btn-gray" style="min-width:104px;padding-inline:14px;white-space:nowrap">מחיקה</button>'
        + '</div>'
        + '<label style="display:flex;align-items:center;gap:6px;font-size:13px;color:#334155;margin-bottom:10px"><input type="checkbox" name="save_template" value="1"> שמור כתבנית חדשה</label>'
        + '<label class="field-label">שם תבנית חדשה</label>'
        + '<input type="text" name="template_name" value="' + esc(template_name_value) + '" placeholder="שם תבנית" style="margin-bottom:0">'
        + '</div>'
        + '<div>'
        + '<div style="display:flex;gap:8px;flex-wrap:wrap;margin-bottom:12px">'
        + '<span style="display:inline-flex;align-items:center;gap:6px;padding:6px 10px;border-radius:999px;background:#ecfdf5;border:1px solid #86efac;font-size:12px;color:#166534">עמודה מהטבלה</span>'
        + '</div>'
        + '<div style="display:grid;grid-template-columns:repeat(auto-fit,minmax(220px,1fr));gap:10px;margin-bottom:12px">' + mapping_fields_html + '</div>'
        + '<div style="font-size:12px;color:#64748b;line-height:1.7;margin-bottom:12px">שדות חובה: עמודת כניסה ועמודת יציאה. המערכת מזהה תיקון לפי סימן * בערך הזמן.</div>'
        + filters_html
        + '<div style="display:flex;gap:10px;flex-wrap:wrap;justify-content:center;margin-top:14px"><button type="submit" id="mappingConfirmButton" name="mapping_action" value="confirm" class="btn btn-blue" style="min-width:220px">אשר הכל והפעל עיבוד</button><a href="/run/' + script_id + '" class="btn btn-gray" style="text-decoration:none;display:inline-flex;align-items:center;justify-content:center;min-width:150px">העלאת קובץ חדש</a></div>'
        + '</div></div>'
        + '</div>'
        + '<div id="mappingProcessingOverlay" style="display:none;position:fixed;inset:0;background:rgba(248,250,252,.78);backdrop-filter:blur(2px);z-index:80;align-items:center;justify-content:center;padding:20px">'
        + '<div style="width:100%;max-width:320px;background:#ffffff;border:1px solid #dbeafe;border-radius:18px;box-shadow:0 20px 50px rgba(15,23,42,.14);padding:24px 20px;text-align:center">'
        + '<div style="width:42px;height:42px;border-radius:999px;border:3px solid #bfdbfe;border-top-color:#2563eb;margin:0 auto 14px;animation:mappingSpin .9s linear infinite"></div>'
        + '<div style="font-size:16px;font-weight:800;color:#1e3a8a;margin-bottom:6px">הדוח בהכנה</div>'
        + '<div style="font-size:13px;line-height:1.7;color:#475569">המערכת מאשרת את השדות ומעבדת את הקובץ. בדוחות גדולים הפעולה יכולה להימשך מעט זמן.</div>'
        + '</div></div>'
        + '<script>'
        + '(function(){'
        + 'var templateSelect=document.getElementById("selectedTemplateId");'
        + 'var fieldSelects=Array.prototype.slice.call(document.querySelectorAll(\'select[data-mapping-field="1"]\'));'
        + 'var form=document.getElementById("mappingConfirmForm");'
        + 'var confirmButton=document.getElementById("mappingConfirmButton");'
        + 'var overlay=document.getElementById("mappingProcessingOverlay");'
        + 'var templateMappings=' + json.dumps(template_payload, ensure_ascii=False) + ';'
        + 'var fieldLabels=' + json.dumps(mapping_labels, ensure_ascii=False) + ';'
        + 'var selectStyles={table_exact:{bg:"#ecfdf5",border:"#4ade80",shadow:"rgba(34,197,94,.14)"},empty:{bg:"#ffffff",border:"#e2e8f0",shadow:"rgba(148,163,184,.08)"}};'
        + 'function applySelectVisual(sel){var opt=sel.options[sel.selectedIndex];var kind=(opt&&opt.getAttribute("data-source-kind"))||"empty";var style=selectStyles[kind]||selectStyles.empty;sel.style.backgroundColor=style.bg;sel.style.borderColor=style.border;sel.style.boxShadow="0 0 0 3px "+style.shadow;}'
        + 'function refreshOptionLabels(){var assignments={};fieldSelects.forEach(function(sel){if(sel.value){assignments[sel.value]=sel.name;}});fieldSelects.forEach(function(sel){Array.prototype.forEach.call(sel.options,function(opt){var base=opt.getAttribute("data-base-label")||opt.text;var assigned=assignments[opt.value];var suffix="";if(opt.value && assigned && assigned!==sel.name){suffix=" [נבחר עבור "+(fieldLabels[assigned]||assigned)+"]";}opt.text=base+suffix;});applySelectVisual(sel);});}'
        + 'function clearDuplicateSelections(changedSelect){if(!changedSelect.value){refreshOptionLabels();return;}fieldSelects.forEach(function(sel){if(sel!==changedSelect && sel.value===changedSelect.value){sel.value="";}});refreshOptionLabels();}'
        + 'function applyTemplate(templateId){var mapping=templateMappings[templateId]||{};if(!templateId){refreshOptionLabels();return;}fieldSelects.forEach(function(sel){sel.value=mapping[sel.name]||"";});var seen={};fieldSelects.forEach(function(sel){if(sel.value && seen[sel.value]){sel.value="";}else if(sel.value){seen[sel.value]=true;}});refreshOptionLabels();}'
        + 'fieldSelects.forEach(function(sel){sel.addEventListener("change",function(){clearDuplicateSelections(sel);});});'
        + 'if(templateSelect){templateSelect.addEventListener("change",function(){applyTemplate(this.value);});}'
        + 'if(form){form.addEventListener("submit",function(event){var submitter=event.submitter||document.activeElement;if(!submitter||submitter.value!=="confirm"){return;}if(confirmButton){confirmButton.disabled=true;confirmButton.textContent="העיבוד התחיל...";}if(overlay){overlay.style.display="flex";}document.body.style.overflow="hidden";});}'
        + 'refreshOptionLabels();'
        + '})();'
        + '</script>'
        + '<style>@keyframes mappingSpin{from{transform:rotate(0deg);}to{transform:rotate(360deg);}}</style>'
        + '</form>'
    )


def build_rimon_mapping_form(script_id, temp_upload_path, temp_upload_ext, inspection, current_mapping, templates, template_name_value):
    template_options = '<option value="">ללא תבנית שמורה</option>'
    for template in templates:
        template_options += '<option value="' + str(template["id"]) + '">' + esc(template["name"]) + '</option>'

    mapping_labels = {field["name"]: field["label"] for field in RIMON_MAPPING_FIELDS}
    template_payload = {
        str(template["id"]): {key: str(value or "") for key, value in template["mapping"].items()}
        for template in templates
    }

    mapping_fields_html = ""
    for field in RIMON_MAPPING_FIELDS:
        field_name = field["name"]
        current_value = str(current_mapping.get(field_name, "") or "")
        options = inspection["options_by_field"].get(field_name, [])
        blank_options = [option for option in options if not option.get("value")]
        meta_options = [option for option in options if option.get("source_kind") == "meta"]
        table_exact_options = [option for option in options if option.get("source_kind") == "table_exact"]
        table_nearby_options = [option for option in options if option.get("source_kind") == "table_nearby"]

        select_options = ""

        def render_option(option):
            selected = ' selected' if option["value"] == current_value else ""
            return (
                '<option value="' + esc(option["value"]) + '" data-base-label="' + esc(option["label"]) + '" data-source-kind="' + esc(option.get("source_kind", "empty")) + '"' + selected + ">"
                + esc(option["label"])
                + "</option>"
            )

        for option in blank_options:
            select_options += render_option(option)
        if meta_options:
            select_options += '<optgroup label="שדות עליונים">'
            for option in meta_options:
                select_options += render_option(option)
            select_options += '</optgroup>'
        if table_exact_options:
            select_options += '<optgroup label="שדות מהטבלה">'
            for option in table_exact_options:
                select_options += render_option(option)
            select_options += '</optgroup>'
        if table_nearby_options:
            select_options += '<optgroup label="שדות מהטבלה (כותרת זוהתה מעמודה סמוכה)">'
            for option in table_nearby_options:
                select_options += render_option(option)
            select_options += '</optgroup>'

        required_badge = ' <span style="color:#dc2626">*</span>' if field["required"] else ' <span style="color:#94a3b8">(אופציונלי)</span>'
        mapping_fields_html += (
            '<div><label class="field-label">' + field["label"] + required_badge + '</label>'
            + '<select name="' + field_name + '" data-mapping-field="1" data-field-label="' + esc(field["label"]) + '" style="padding:9px 12px;border:1.5px solid #e2e8f0;border-radius:8px;font-size:13px;font-family:inherit;outline:none;width:100%;margin-bottom:0;background:white;transition:background-color .15s ease,border-color .15s ease,box-shadow .15s ease">'
            + select_options
            + '</select></div>'
        )

    return (
        '<form method="POST" id="mappingConfirmForm">'
        + '<input type="hidden" name="flow_mode" value="confirm_mapping">'
        + '<input type="hidden" name="temp_upload_path" value="' + esc(temp_upload_path) + '">'
        + '<input type="hidden" name="temp_upload_ext" value="' + esc(temp_upload_ext) + '">'
        + '<div style="background:#fafcff;border:1px solid #dbeafe;border-radius:14px;padding:1rem;margin-bottom:1rem">'
        + '<div style="font-size:15px;font-weight:700;color:#1e3a8a;margin-bottom:10px">אישור שדות לפני עיבוד</div>'
        + '<div style="display:grid;grid-template-columns:260px minmax(0,1fr);gap:14px;align-items:start">'
        + '<div style="background:#ffffff;border:1px solid #e2e8f0;border-radius:12px;padding:12px">'
        + '<div style="font-size:14px;font-weight:700;color:#0f172a;margin-bottom:10px">תבניות שמורות</div>'
        + '<label class="field-label">בחירת תבנית</label>'
        + '<div style="display:grid;grid-template-columns:minmax(0,1fr) auto;gap:8px;align-items:center;margin-bottom:12px">'
        + '<select id="selectedTemplateId" name="selected_template_id" style="padding:9px 12px;border:1.5px solid #e2e8f0;border-radius:8px;font-size:13px;font-family:inherit;outline:none;width:100%;background:white">' + template_options + '</select>'
        + '<button type="submit" name="mapping_action" value="delete_template" class="btn btn-gray" style="min-width:104px;padding-inline:14px;white-space:nowrap">מחיקה</button>'
        + '</div>'
        + '<label style="display:flex;align-items:center;gap:6px;font-size:13px;color:#334155;margin-bottom:10px"><input type="checkbox" name="save_template" value="1"> שמור כתבנית חדשה</label>'
        + '<label class="field-label">שם תבנית חדשה</label>'
        + '<input type="text" name="template_name" value="' + esc(template_name_value) + '" placeholder="שם תבנית" style="margin-bottom:0">'
        + '<div style="font-size:12px;color:#64748b;line-height:1.7;margin-top:10px">בחירת תבנית תעדכן את כל השדות בהתאם. שמירה תיצור תבנית חדשה בלבד ולא תדרוס תבנית קיימת.</div>'
        + '</div>'
        + '<div>'
        + '<div style="display:flex;gap:8px;flex-wrap:wrap;margin-bottom:12px">'
        + '<span style="display:inline-flex;align-items:center;gap:6px;padding:6px 10px;border-radius:999px;background:#eff6ff;border:1px solid #bfdbfe;font-size:12px;color:#1d4ed8">שדה עליון</span>'
        + '<span style="display:inline-flex;align-items:center;gap:6px;padding:6px 10px;border-radius:999px;background:#ecfdf5;border:1px solid #86efac;font-size:12px;color:#166534">שדה מהטבלה</span>'
        + '<span style="display:inline-flex;align-items:center;gap:6px;padding:6px 10px;border-radius:999px;background:#fffbeb;border:1px solid #fcd34d;font-size:12px;color:#92400e">שדה מהטבלה שזוהה לפי כותרת סמוכה</span>'
        + '</div>'
        + '<div style="display:grid;grid-template-columns:repeat(auto-fit,minmax(220px,1fr));gap:10px;margin-bottom:12px">' + mapping_fields_html + '</div>'
        + '<div style="font-size:12px;color:#64748b;line-height:1.7;margin-bottom:12px">שדות חובה: שם עובד, מספר עובד, תאריך ואירוע. אם אותו שדה נבחר בקטגוריה אחרת, הבחירה הקודמת תנוקה אוטומטית.</div>'
        + '<div style="display:flex;gap:10px;flex-wrap:wrap;justify-content:center"><button type="submit" id="mappingConfirmButton" name="mapping_action" value="confirm" class="btn btn-blue" style="min-width:220px">אשר הכל והפעל עיבוד</button><a href="/run/' + script_id + '" class="btn btn-gray" style="text-decoration:none;display:inline-flex;align-items:center;justify-content:center;min-width:150px">העלאת קובץ חדש</a></div>'
        + '</div></div>'
        + '</div>'
        + '<div id="mappingProcessingOverlay" style="display:none;position:fixed;inset:0;background:rgba(248,250,252,.78);backdrop-filter:blur(2px);z-index:80;align-items:center;justify-content:center;padding:20px">'
        + '<div style="width:100%;max-width:320px;background:#ffffff;border:1px solid #dbeafe;border-radius:18px;box-shadow:0 20px 50px rgba(15,23,42,.14);padding:24px 20px;text-align:center">'
        + '<div style="width:42px;height:42px;border-radius:999px;border:3px solid #bfdbfe;border-top-color:#2563eb;margin:0 auto 14px;animation:mappingSpin .9s linear infinite"></div>'
        + '<div style="font-size:16px;font-weight:800;color:#1e3a8a;margin-bottom:6px">הדוח בהכנה</div>'
        + '<div style="font-size:13px;line-height:1.7;color:#475569">המערכת מאשרת את השדות ומעבדת את הקובץ. בדוחות גדולים הפעולה יכולה להימשך מעט זמן.</div>'
        + '</div></div>'
        + '<script>'
        + '(function(){'
        + 'var templateSelect=document.getElementById("selectedTemplateId");'
        + 'var fieldSelects=Array.prototype.slice.call(document.querySelectorAll(\'select[data-mapping-field="1"]\'));'
        + 'var form=document.getElementById("mappingConfirmForm");'
        + 'var confirmButton=document.getElementById("mappingConfirmButton");'
        + 'var overlay=document.getElementById("mappingProcessingOverlay");'
        + 'var templateMappings=' + json.dumps(template_payload, ensure_ascii=False) + ';'
        + 'var fieldLabels=' + json.dumps(mapping_labels, ensure_ascii=False) + ';'
        + 'var selectStyles={meta:{bg:"#eff6ff",border:"#60a5fa",shadow:"rgba(59,130,246,.12)"},table_exact:{bg:"#ecfdf5",border:"#4ade80",shadow:"rgba(34,197,94,.14)"},table_nearby:{bg:"#fffbeb",border:"#f59e0b",shadow:"rgba(245,158,11,.16)"},empty:{bg:"#ffffff",border:"#e2e8f0",shadow:"rgba(148,163,184,.08)"}};'
        + 'function applySelectVisual(sel){var opt=sel.options[sel.selectedIndex];var kind=(opt&&opt.getAttribute("data-source-kind"))||"empty";var style=selectStyles[kind]||selectStyles.empty;sel.style.backgroundColor=style.bg;sel.style.borderColor=style.border;sel.style.boxShadow="0 0 0 3px "+style.shadow;}'
        + 'function refreshOptionLabels(){'
        + 'var assignments={};'
        + 'fieldSelects.forEach(function(sel){if(sel.value){assignments[sel.value]=sel.name;}});'
        + 'fieldSelects.forEach(function(sel){Array.prototype.forEach.call(sel.options,function(opt){var base=opt.getAttribute("data-base-label")||opt.text;var assigned=assignments[opt.value];var suffix="";if(opt.value && assigned && assigned!==sel.name){suffix=" [נבחר עבור "+(fieldLabels[assigned]||assigned)+"]";}opt.text=base+suffix;});applySelectVisual(sel);});'
        + '}'
        + 'function clearDuplicateSelections(changedSelect){if(!changedSelect.value){refreshOptionLabels();return;}fieldSelects.forEach(function(sel){if(sel!==changedSelect && sel.value===changedSelect.value){sel.value="";}});refreshOptionLabels();}'
        + 'function applyTemplate(templateId){var mapping=templateMappings[templateId]||{};if(!templateId){refreshOptionLabels();return;}fieldSelects.forEach(function(sel){sel.value=mapping[sel.name]||"";});var seen={};fieldSelects.forEach(function(sel){if(sel.value && seen[sel.value]){sel.value="";}else if(sel.value){seen[sel.value]=true;}});refreshOptionLabels();}'
        + 'fieldSelects.forEach(function(sel){sel.addEventListener("change",function(){clearDuplicateSelections(sel);});});'
        + 'if(templateSelect){templateSelect.addEventListener("change",function(){applyTemplate(this.value);});}'
        + 'if(form){form.addEventListener("submit",function(event){var submitter=event.submitter||document.activeElement;if(!submitter||submitter.value!=="confirm"){return;}if(confirmButton){confirmButton.disabled=true;confirmButton.textContent="העיבוד התחיל...";}if(overlay){overlay.style.display="flex";}document.body.style.overflow="hidden";});}'
        + 'refreshOptionLabels();'
        + '})();'
        + '</script>'
        + '<style>@keyframes mappingSpin{from{transform:rotate(0deg);}to{transform:rotate(360deg);}}</style>'
        + '</form>'
    )


def build_attendance_alerts_mapping_options(input_path, extension):
    base = build_rimon_mapping_options(input_path, extension)
    options_by_field = dict(base["options_by_field"])
    suggestions = dict(base["suggestions"])
    visible_table_options = []
    for key, opts in options_by_field.items():
        for opt in opts:
            if opt.get("source_kind") in ("table_exact", "table_nearby") and opt not in visible_table_options:
                visible_table_options.append(opt)
    extra_fields = [f for f in ATTENDANCE_ALERTS_MAPPING_FIELDS if f["name"] not in options_by_field]
    for field in extra_fields:
        field_name = field["name"]
        options = [{"value": "", "label": "לא נבחר", "source_kind": "empty"}]
        options.extend(visible_table_options)
        options_by_field[field_name] = options
        suggested = ""
        keywords = ATTENDANCE_ALERTS_SUGGESTION_KEYWORDS.get(field_name, [])
        ranked = sorted(visible_table_options, key=lambda o: (o.get("from_nearby", False), o["value"]))
        for option in ranked:
            header_token = normalize_token(option["header"])
            exact_token = normalize_token(option.get("exact_header", ""))
            if any(kw in exact_token for kw in keywords):
                suggested = option["value"]
                break
            if not option.get("from_nearby") and any(kw in header_token for kw in keywords):
                suggested = option["value"]
                break
        if not suggested:
            for option in ranked:
                token = normalize_token(option["header"])
                if any(kw in token for kw in keywords):
                    suggested = option["value"]
                    break
        suggestions[field_name] = suggested
    base["options_by_field"] = options_by_field
    base["suggestions"] = suggestions
    base["suggested_template_name"] = base.get("suggested_template_name", "תבנית התראות")
    return base


def build_attendance_alerts_mapping_form(script_id, temp_upload_path, temp_upload_ext, inspection, current_mapping, templates, template_name_value):
    template_options = '<option value="">ללא תבנית שמורה</option>'
    for template in templates:
        template_options += '<option value="' + str(template["id"]) + '">' + esc(template["name"]) + '</option>'
    mapping_labels = {field["name"]: field["label"] for field in ATTENDANCE_ALERTS_MAPPING_FIELDS}
    template_payload = {
        str(template["id"]): {key: str(value or "") for key, value in template["mapping"].items()}
        for template in templates
    }
    mapping_fields_html = ""
    for field in ATTENDANCE_ALERTS_MAPPING_FIELDS:
        field_name = field["name"]
        current_value = str(current_mapping.get(field_name, "") or "")
        options = inspection["options_by_field"].get(field_name, [])
        blank_options = [o for o in options if not o.get("value")]
        meta_options = [o for o in options if o.get("source_kind") == "meta"]
        table_exact_options = [o for o in options if o.get("source_kind") == "table_exact"]
        table_nearby_options = [o for o in options if o.get("source_kind") == "table_nearby"]
        select_options = ""

        def render_option(option):
            selected = ' selected' if option["value"] == current_value else ""
            return (
                '<option value="' + esc(option["value"]) + '" data-base-label="' + esc(option["label"]) + '" data-source-kind="' + esc(option.get("source_kind", "empty")) + '"' + selected + ">"
                + esc(option["label"])
                + "</option>"
            )

        for option in blank_options:
            select_options += render_option(option)
        if meta_options:
            select_options += '<optgroup label="שדות עליונים">'
            for option in meta_options:
                select_options += render_option(option)
            select_options += '</optgroup>'
        if table_exact_options:
            select_options += '<optgroup label="שדות מהטבלה">'
            for option in table_exact_options:
                select_options += render_option(option)
            select_options += '</optgroup>'
        if table_nearby_options:
            select_options += '<optgroup label="שדות מהטבלה (כותרת זוהתה מעמודה סמוכה)">'
            for option in table_nearby_options:
                select_options += render_option(option)
            select_options += '</optgroup>'

        required_badge = ' <span style="color:#dc2626">*</span>' if field["required"] else ' <span style="color:#94a3b8">(אופציונלי)</span>'
        mapping_fields_html += (
            '<div><label class="field-label">' + field["label"] + required_badge + '</label>'
            + '<select name="' + field_name + '" data-mapping-field="1" data-field-label="' + esc(field["label"]) + '" style="padding:9px 12px;border:1.5px solid #e2e8f0;border-radius:8px;font-size:13px;font-family:inherit;outline:none;width:100%;margin-bottom:0;background:white;transition:background-color .15s ease,border-color .15s ease,box-shadow .15s ease">'
            + select_options
            + '</select></div>'
        )

    return (
        '<form method="POST" id="mappingConfirmForm">'
        + '<input type="hidden" name="flow_mode" value="confirm_mapping">'
        + '<input type="hidden" name="temp_upload_path" value="' + esc(temp_upload_path) + '">'
        + '<input type="hidden" name="temp_upload_ext" value="' + esc(temp_upload_ext) + '">'
        + '<div style="background:#fafcff;border:1px solid #dbeafe;border-radius:14px;padding:1rem;margin-bottom:1rem">'
        + '<div style="font-size:15px;font-weight:700;color:#1e3a8a;margin-bottom:10px">אישור שדות לפני עיבוד</div>'
        + '<div style="display:grid;grid-template-columns:260px minmax(0,1fr);gap:14px;align-items:start">'
        + '<div style="background:#ffffff;border:1px solid #e2e8f0;border-radius:12px;padding:12px">'
        + '<div style="font-size:14px;font-weight:700;color:#0f172a;margin-bottom:10px">תבניות שמורות</div>'
        + '<label class="field-label">בחירת תבנית</label>'
        + '<div style="display:grid;grid-template-columns:minmax(0,1fr) auto;gap:8px;align-items:center;margin-bottom:12px">'
        + '<select id="selectedTemplateId" name="selected_template_id" style="padding:9px 12px;border:1.5px solid #e2e8f0;border-radius:8px;font-size:13px;font-family:inherit;outline:none;width:100%;background:white">' + template_options + '</select>'
        + '<button type="submit" name="mapping_action" value="delete_template" class="btn btn-gray" style="min-width:104px;padding-inline:14px;white-space:nowrap">מחיקה</button>'
        + '</div>'
        + '<label style="display:flex;align-items:center;gap:6px;font-size:13px;color:#334155;margin-bottom:10px"><input type="checkbox" name="save_template" value="1"> שמור כתבנית חדשה</label>'
        + '<label class="field-label">שם תבנית חדשה</label>'
        + '<input type="text" name="template_name" value="' + esc(template_name_value) + '" placeholder="שם תבנית" style="margin-bottom:0">'
        + '<div style="font-size:12px;color:#64748b;line-height:1.7;margin-top:10px">בחירת תבנית תעדכן את כל השדות בהתאם. שמירה תיצור תבנית חדשה בלבד ולא תדרוס תבנית קיימת.</div>'
        + '</div>'
        + '<div>'
        + '<div style="display:flex;gap:8px;flex-wrap:wrap;margin-bottom:12px">'
        + '<span style="display:inline-flex;align-items:center;gap:6px;padding:6px 10px;border-radius:999px;background:#eff6ff;border:1px solid #bfdbfe;font-size:12px;color:#1d4ed8">שדה עליון</span>'
        + '<span style="display:inline-flex;align-items:center;gap:6px;padding:6px 10px;border-radius:999px;background:#ecfdf5;border:1px solid #86efac;font-size:12px;color:#166534">שדה מהטבלה</span>'
        + '<span style="display:inline-flex;align-items:center;gap:6px;padding:6px 10px;border-radius:999px;background:#fffbeb;border:1px solid #fcd34d;font-size:12px;color:#92400e">שדה מהטבלה שזוהה לפי כותרת סמוכה</span>'
        + '</div>'
        + '<div style="display:grid;grid-template-columns:repeat(auto-fit,minmax(220px,1fr));gap:10px;margin-bottom:12px">' + mapping_fields_html + '</div>'
        + '<div style="font-size:12px;color:#64748b;line-height:1.7;margin-bottom:12px">שדות חובה: שם עובד, מספר עובד, תאריך ואירוע. שדות נוספים (הפסקה, שעות רגילות, שעות נוספות) ישפרו את הדוח.</div>'
        + '<div style="background:#f8fafc;border:1px solid #e2e8f0;border-radius:12px;padding:12px;margin-bottom:14px">'
        + '<div style="font-size:14px;font-weight:700;color:#0f172a;margin-bottom:10px">בחירת התראות לדוח</div>'
        + '<div style="display:grid;grid-template-columns:repeat(auto-fit,minmax(220px,1fr));gap:8px">'
        + ''.join(
            '<label style="display:flex;align-items:center;gap:8px;font-size:13px;color:#334155;padding:6px 10px;border-radius:8px;background:white;border:1px solid #e2e8f0;cursor:pointer">'
            + '<input type="checkbox" name="alert_' + a["id"] + '" value="1" checked>'
            + '<span style="width:12px;height:12px;border-radius:3px;background:#' + a["color"] + ';border:1px solid #cbd5e1;flex-shrink:0"></span>'
            + a["label"]
            + '</label>'
            for a in ATTENDANCE_ALERT_TYPES
        )
        + '</div>'
        + '<div style="font-size:11px;color:#94a3b8;margin-top:6px">בטל סימון של התראות שלא רלוונטיות. רק התראות מסומנות יופיעו בדוח.</div>'
        + '</div>'
        + '<div style="display:flex;gap:10px;flex-wrap:wrap;justify-content:center"><button type="submit" id="mappingConfirmButton" name="mapping_action" value="confirm" class="btn btn-blue" style="min-width:220px">אשר הכל והפעל עיבוד</button><a href="/run/' + script_id + '" class="btn btn-gray" style="text-decoration:none;display:inline-flex;align-items:center;justify-content:center;min-width:150px">העלאת קובץ חדש</a></div>'
        + '</div></div>'
        + '</div>'
        + '<div id="mappingProcessingOverlay" style="display:none;position:fixed;inset:0;background:rgba(248,250,252,.78);backdrop-filter:blur(2px);z-index:80;align-items:center;justify-content:center;padding:20px">'
        + '<div style="width:100%;max-width:320px;background:#ffffff;border:1px solid #dbeafe;border-radius:18px;box-shadow:0 20px 50px rgba(15,23,42,.14);padding:24px 20px;text-align:center">'
        + '<div style="width:42px;height:42px;border-radius:999px;border:3px solid #bfdbfe;border-top-color:#2563eb;margin:0 auto 14px;animation:mappingSpin .9s linear infinite"></div>'
        + '<div style="font-size:16px;font-weight:800;color:#1e3a8a;margin-bottom:6px">הדוח בהכנה</div>'
        + '<div style="font-size:13px;line-height:1.7;color:#475569">המערכת מאשרת את השדות ומעבדת את הקובץ. בדוחות גדולים הפעולה יכולה להימשך מעט זמן.</div>'
        + '</div></div>'
        + '<script>'
        + '(function(){'
        + 'var templateSelect=document.getElementById("selectedTemplateId");'
        + 'var fieldSelects=Array.prototype.slice.call(document.querySelectorAll(\'select[data-mapping-field="1"]\'));'
        + 'var form=document.getElementById("mappingConfirmForm");'
        + 'var confirmButton=document.getElementById("mappingConfirmButton");'
        + 'var overlay=document.getElementById("mappingProcessingOverlay");'
        + 'var templateMappings=' + json.dumps(template_payload, ensure_ascii=False) + ';'
        + 'var fieldLabels=' + json.dumps(mapping_labels, ensure_ascii=False) + ';'
        + 'var selectStyles={meta:{bg:"#eff6ff",border:"#60a5fa",shadow:"rgba(59,130,246,.12)"},table_exact:{bg:"#ecfdf5",border:"#4ade80",shadow:"rgba(34,197,94,.14)"},table_nearby:{bg:"#fffbeb",border:"#f59e0b",shadow:"rgba(245,158,11,.16)"},empty:{bg:"#ffffff",border:"#e2e8f0",shadow:"rgba(148,163,184,.08)"}};'
        + 'function applySelectVisual(sel){var opt=sel.options[sel.selectedIndex];var kind=(opt&&opt.getAttribute("data-source-kind"))||"empty";var style=selectStyles[kind]||selectStyles.empty;sel.style.backgroundColor=style.bg;sel.style.borderColor=style.border;sel.style.boxShadow="0 0 0 3px "+style.shadow;}'
        + 'function refreshOptionLabels(){'
        + 'var assignments={};'
        + 'fieldSelects.forEach(function(sel){if(sel.value){assignments[sel.value]=sel.name;}});'
        + 'fieldSelects.forEach(function(sel){Array.prototype.forEach.call(sel.options,function(opt){var base=opt.getAttribute("data-base-label")||opt.text;var assigned=assignments[opt.value];var suffix="";if(opt.value && assigned && assigned!==sel.name){suffix=" [נבחר עבור "+(fieldLabels[assigned]||assigned)+"]";}opt.text=base+suffix;});applySelectVisual(sel);});'
        + '}'
        + 'function clearDuplicateSelections(changedSelect){if(!changedSelect.value){refreshOptionLabels();return;}fieldSelects.forEach(function(sel){if(sel!==changedSelect && sel.value===changedSelect.value){sel.value="";}});refreshOptionLabels();}'
        + 'function applyTemplate(templateId){var mapping=templateMappings[templateId]||{};if(!templateId){refreshOptionLabels();return;}fieldSelects.forEach(function(sel){sel.value=mapping[sel.name]||"";});var seen={};fieldSelects.forEach(function(sel){if(sel.value && seen[sel.value]){sel.value="";}else if(sel.value){seen[sel.value]=true;}});refreshOptionLabels();}'
        + 'fieldSelects.forEach(function(sel){sel.addEventListener("change",function(){clearDuplicateSelections(sel);});});'
        + 'if(templateSelect){templateSelect.addEventListener("change",function(){applyTemplate(this.value);});}'
        + 'if(form){form.addEventListener("submit",function(event){var submitter=event.submitter||document.activeElement;if(!submitter||submitter.value!=="confirm"){return;}if(confirmButton){confirmButton.disabled=true;confirmButton.textContent="העיבוד התחיל...";}if(overlay){overlay.style.display="flex";}document.body.style.overflow="hidden";});}'
        + 'refreshOptionLabels();'
        + '})();'
        + '</script>'
        + '<style>@keyframes mappingSpin{from{transform:rotate(0deg);}to{transform:rotate(360deg);}}</style>'
        + '</form>'
    )


def build_matan_missing_mapping_form(script_id, temp_upload_path, temp_upload_ext, inspection, current_mapping, templates, template_name_value, current_filters):
    template_options = '<option value="">ללא תבנית שמורה</option>'
    for template in templates:
        template_options += '<option value="' + str(template["id"]) + '">' + esc(template["name"]) + '</option>'

    mapping_labels = {field["name"]: field["label"] for field in MATAN_MISSING_MAPPING_FIELDS}
    template_payload = {
        str(template["id"]): {key: str(value or "") for key, value in template["mapping"].items()}
        for template in templates
    }

    fields_html = ""
    for field in MATAN_MISSING_MAPPING_FIELDS:
        field_name = field["name"]
        current_value = str(current_mapping.get(field_name, "") or "")
        options = inspection["options_by_field"].get(field_name, [])

        select_options = ""
        for option in options:
            selected = ' selected' if option["value"] == current_value else ""
            select_options += (
                '<option value="' + esc(option["value"]) + '" data-base-label="' + esc(option["label"]) + '" data-source-kind="' + esc(option.get("source_kind", "empty")) + '"' + selected + ">"
                + esc(option["label"])
                + "</option>"
            )

        required_badge = ' <span style="color:#dc2626">*</span>' if field["required"] else ' <span style="color:#94a3b8">(אופציונלי)</span>'
        wrapper_style = ""
        if field.get("critical"):
            wrapper_style = 'background:#fff7ed;border:1px solid #fdba74;border-radius:12px;padding:10px 10px 12px'
        fields_html += (
            '<div style="' + wrapper_style + '"><label class="field-label">' + field["label"] + required_badge + '</label>'
            + ('<div style="font-size:12px;color:#9a3412;line-height:1.6;margin:-4px 0 8px">שדה קריטי לסינון שעות החוסר. יש לוודא שזהו השדה הנכון.</div>' if field.get("critical") else '')
            + '<select name="' + field_name + '" data-mapping-field="1" data-field-label="' + esc(field["label"]) + '" style="padding:9px 12px;border:1.5px solid #e2e8f0;border-radius:8px;font-size:13px;font-family:inherit;outline:none;width:100%;margin-bottom:0;background:white;transition:background-color .15s ease,border-color .15s ease,box-shadow .15s ease">'
            + select_options
            + '</select></div>'
        )

    return (
        '<form method="POST" id="matanMappingConfirmForm">'
        + '<input type="hidden" name="flow_mode" value="confirm_mapping">'
        + '<input type="hidden" name="temp_upload_path" value="' + esc(temp_upload_path) + '">'
        + '<input type="hidden" name="temp_upload_ext" value="' + esc(temp_upload_ext) + '">'
        + '<div style="background:#fafcff;border:1px solid #dbeafe;border-radius:14px;padding:1rem;margin-bottom:1rem">'
        + '<div style="font-size:15px;font-weight:700;color:#1e3a8a;margin-bottom:10px">אישור שדות לפני עיבוד</div>'
        + '<div style="display:grid;grid-template-columns:260px minmax(0,1fr);gap:14px;align-items:start">'
        + '<div style="background:#ffffff;border:1px solid #e2e8f0;border-radius:12px;padding:12px">'
        + '<div style="font-size:14px;font-weight:700;color:#0f172a;margin-bottom:10px">תבניות שמורות</div>'
        + '<label class="field-label">בחירת תבנית</label>'
        + '<div style="display:grid;grid-template-columns:minmax(0,1fr) auto;gap:8px;align-items:center;margin-bottom:12px">'
        + '<select id="selectedMatanTemplateId" name="selected_template_id" style="padding:9px 12px;border:1.5px solid #e2e8f0;border-radius:8px;font-size:13px;font-family:inherit;outline:none;width:100%;background:white">' + template_options + '</select>'
        + '<button type="submit" name="mapping_action" value="delete_template" class="btn btn-gray" style="min-width:104px;padding-inline:14px;white-space:nowrap">מחיקה</button>'
        + '</div>'
        + '<label style="display:flex;align-items:center;gap:6px;font-size:13px;color:#334155;margin-bottom:10px"><input type="checkbox" name="save_template" value="1"> שמור כתבנית חדשה</label>'
        + '<label class="field-label">שם תבנית חדשה</label>'
        + '<input type="text" name="template_name" value="' + esc(template_name_value) + '" placeholder="שם תבנית" style="margin-bottom:14px">'
        + '<div style="font-size:14px;font-weight:700;color:#0f172a;margin-bottom:10px">תנאי סינון</div>'
        + '<label class="field-label">מינימום שעות חוסר</label>'
        + '<input type="text" name="min_missing_hours" value="' + esc(current_filters.get("min_missing_hours", "")) + '" placeholder="לדוגמה 4" style="margin-bottom:10px">'
        + '<label class="field-label">מקסימום שעות חוסר</label>'
        + '<input type="text" name="max_missing_hours" value="' + esc(current_filters.get("max_missing_hours", "")) + '" placeholder="לדוגמה 8" style="margin-bottom:0">'
        + '<div style="font-size:12px;color:#64748b;line-height:1.7;margin-top:10px">בחירת תבנית תעדכן את השדות בלבד. תנאי הסינון נשארים לפי מה שהוזן במסך הנוכחי.</div>'
        + '</div>'
        + '<div>'
        + '<div style="display:flex;gap:8px;flex-wrap:wrap;margin-bottom:12px">'
        + '<span style="display:inline-flex;align-items:center;gap:6px;padding:6px 10px;border-radius:999px;background:#fff7ed;border:1px solid #fdba74;font-size:12px;color:#9a3412">שדה קריטי לסינון</span>'
        + '<span style="display:inline-flex;align-items:center;gap:6px;padding:6px 10px;border-radius:999px;background:#ecfdf5;border:1px solid #86efac;font-size:12px;color:#166534">שדה מהדוח</span>'
        + '</div>'
        + '<div style="display:grid;grid-template-columns:repeat(auto-fit,minmax(220px,1fr));gap:10px;margin-bottom:12px">' + fields_html + '</div>'
        + '<div style="font-size:12px;color:#64748b;line-height:1.7;margin-bottom:12px">שדות חובה: שם עובד, שעות תקן וחוסר. בנוסף יש לבחור לפחות מזהה אחד נוסף: מספר עובד, תעודת זהות, מספר תג או דרכון.</div>'
        + '<div style="display:flex;gap:10px;flex-wrap:wrap;justify-content:center"><button type="submit" id="matanMappingConfirmButton" name="mapping_action" value="confirm" class="btn btn-blue" style="min-width:220px">אשר הכל והפעל עיבוד</button><a href="/run/' + script_id + '" class="btn btn-gray" style="text-decoration:none;display:inline-flex;align-items:center;justify-content:center;min-width:150px">העלאת קובץ חדש</a></div>'
        + '</div></div>'
        + '</div>'
        + '<div id="matanProcessingOverlay" style="display:none;position:fixed;inset:0;background:rgba(248,250,252,.78);backdrop-filter:blur(2px);z-index:80;align-items:center;justify-content:center;padding:20px">'
        + '<div style="width:100%;max-width:320px;background:#ffffff;border:1px solid #dbeafe;border-radius:18px;box-shadow:0 20px 50px rgba(15,23,42,.14);padding:24px 20px;text-align:center">'
        + '<div style="width:42px;height:42px;border-radius:999px;border:3px solid #bfdbfe;border-top-color:#2563eb;margin:0 auto 14px;animation:mappingSpin .9s linear infinite"></div>'
        + '<div style="font-size:16px;font-weight:800;color:#1e3a8a;margin-bottom:6px">הדוח בהכנה</div>'
        + '<div style="font-size:13px;line-height:1.7;color:#475569">המערכת מאשרת את השדות, מסננת את העובדים ובונה את הדוח. הפעולה יכולה להימשך מעט זמן.</div>'
        + '</div></div>'
        + '<script>'
        + '(function(){'
        + 'var templateSelect=document.getElementById("selectedMatanTemplateId");'
        + 'var fieldSelects=Array.prototype.slice.call(document.querySelectorAll(\'select[data-mapping-field="1"]\'));'
        + 'var form=document.getElementById("matanMappingConfirmForm");'
        + 'var confirmButton=document.getElementById("matanMappingConfirmButton");'
        + 'var overlay=document.getElementById("matanProcessingOverlay");'
        + 'var templateMappings=' + json.dumps(template_payload, ensure_ascii=False) + ';'
        + 'var fieldLabels=' + json.dumps(mapping_labels, ensure_ascii=False) + ';'
        + 'var selectStyles={critical:{bg:"#fff7ed",border:"#fb923c",shadow:"rgba(249,115,22,.14)"},table_exact:{bg:"#ecfdf5",border:"#4ade80",shadow:"rgba(34,197,94,.14)"},empty:{bg:"#ffffff",border:"#e2e8f0",shadow:"rgba(148,163,184,.08)"}};'
        + 'function applySelectVisual(sel){var fieldName=sel.name;var isCritical=(fieldName==="standard_hours_source"||fieldName==="missing_hours_source");var kind=isCritical?"critical":(((sel.options[sel.selectedIndex]||{}).getAttribute&&sel.options[sel.selectedIndex].getAttribute("data-source-kind"))||"table_exact");var style=selectStyles[kind]||selectStyles.empty;sel.style.backgroundColor=style.bg;sel.style.borderColor=style.border;sel.style.boxShadow="0 0 0 3px "+style.shadow;}'
        + 'function refreshOptionLabels(){var assignments={};fieldSelects.forEach(function(sel){if(sel.value){assignments[sel.value]=sel.name;}});fieldSelects.forEach(function(sel){Array.prototype.forEach.call(sel.options,function(opt){var base=opt.getAttribute("data-base-label")||opt.text;var assigned=assignments[opt.value];var suffix="";if(opt.value && assigned && assigned!==sel.name){suffix=" [נבחר עבור "+(fieldLabels[assigned]||assigned)+"]";}opt.text=base+suffix;});applySelectVisual(sel);});}'
        + 'function clearDuplicateSelections(changedSelect){if(!changedSelect.value){refreshOptionLabels();return;}fieldSelects.forEach(function(sel){if(sel!==changedSelect && sel.value===changedSelect.value){sel.value="";}});refreshOptionLabels();}'
        + 'function applyTemplate(templateId){var mapping=templateMappings[templateId]||{};if(!templateId){refreshOptionLabels();return;}fieldSelects.forEach(function(sel){sel.value=mapping[sel.name]||"";});var seen={};fieldSelects.forEach(function(sel){if(sel.value && seen[sel.value]){sel.value="";}else if(sel.value){seen[sel.value]=true;}});refreshOptionLabels();}'
        + 'fieldSelects.forEach(function(sel){sel.addEventListener("change",function(){clearDuplicateSelections(sel);});});'
        + 'if(templateSelect){templateSelect.addEventListener("change",function(){applyTemplate(this.value);});}'
        + 'if(form){form.addEventListener("submit",function(event){var submitter=event.submitter||document.activeElement;if(!submitter||submitter.value!=="confirm"){return;}if(confirmButton){confirmButton.disabled=true;confirmButton.textContent="העיבוד התחיל...";}if(overlay){overlay.style.display="flex";}document.body.style.overflow="hidden";});}'
        + 'refreshOptionLabels();'
        + '})();'
        + '</script>'
        + '<style>@keyframes mappingSpin{from{transform:rotate(0deg);}to{transform:rotate(360deg);}}</style>'
        + '</form>'
    )


def build_inactive_workers_mapping_form(script_id, temp_upload_path, temp_upload_ext, inspection, current_mapping, templates, template_name_value, current_filters):
    template_options = '<option value="">ללא תבנית שמורה</option>'
    for template in templates:
        template_options += '<option value="' + str(template["id"]) + '">' + esc(template["name"]) + '</option>'

    mapping_labels = {field["name"]: field["label"] for field in INACTIVE_WORKERS_MAPPING_FIELDS}
    template_payload = {
        str(template["id"]): {key: str(value or "") for key, value in template["mapping"].items()}
        for template in templates
    }

    fields_html = ""
    for field in INACTIVE_WORKERS_MAPPING_FIELDS:
        field_name = field["name"]
        current_value = str(current_mapping.get(field_name, "") or "")
        options = inspection["options_by_field"].get(field_name, [])
        select_options = ""
        for option in options:
            selected = ' selected' if option["value"] == current_value else ""
            select_options += (
                '<option value="' + esc(option["value"]) + '" data-base-label="' + esc(option["label"]) + '" data-source-kind="' + esc(option.get("source_kind", "empty")) + '"' + selected + ">"
                + esc(option["label"])
                + "</option>"
            )
        required_badge = ' <span style="color:#dc2626">*</span>' if field["required"] else ' <span style="color:#94a3b8">(אופציונלי)</span>'
        wrapper_style = ""
        if field.get("critical"):
            wrapper_style = 'background:#fff7ed;border:1px solid #fdba74;border-radius:12px;padding:10px 10px 12px'
        fields_html += (
            '<div style="' + wrapper_style + '"><label class="field-label">' + field["label"] + required_badge + '</label>'
            + ('<div style="font-size:12px;color:#9a3412;line-height:1.6;margin:-4px 0 8px">שדה קריטי לזיהוי פעילות. יש לוודא שהוא ממופה נכון.</div>' if field.get("critical") else '')
            + '<select name="' + field_name + '" data-mapping-field="1" data-field-label="' + esc(field["label"]) + '" style="padding:9px 12px;border:1.5px solid #e2e8f0;border-radius:8px;font-size:13px;font-family:inherit;outline:none;width:100%;margin-bottom:0;background:white;transition:background-color .15s ease,border-color .15s ease,box-shadow .15s ease">'
            + select_options
            + '</select></div>'
        )

    unit_value = current_filters.get("inactive_period_unit", "days")
    threshold_value = current_filters.get("inactive_period_value", "")

    return (
        '<form method="POST" id="inactiveWorkersMappingForm">'
        + '<input type="hidden" name="flow_mode" value="confirm_mapping">'
        + '<input type="hidden" name="temp_upload_path" value="' + esc(temp_upload_path) + '">'
        + '<input type="hidden" name="temp_upload_ext" value="' + esc(temp_upload_ext) + '">'
        + '<div style="background:#fafcff;border:1px solid #dbeafe;border-radius:14px;padding:1rem;margin-bottom:1rem">'
        + '<div style="font-size:15px;font-weight:700;color:#1e3a8a;margin-bottom:10px">אישור שדות לפני עיבוד</div>'
        + '<div style="display:grid;grid-template-columns:260px minmax(0,1fr);gap:14px;align-items:start">'
        + '<div style="background:#ffffff;border:1px solid #e2e8f0;border-radius:12px;padding:12px">'
        + '<div style="font-size:14px;font-weight:700;color:#0f172a;margin-bottom:10px">תבניות שמורות</div>'
        + '<label class="field-label">בחירת תבנית</label>'
        + '<div style="display:grid;grid-template-columns:minmax(0,1fr) auto;gap:8px;align-items:center;margin-bottom:12px">'
        + '<select id="selectedInactiveTemplateId" name="selected_template_id" style="padding:9px 12px;border:1.5px solid #e2e8f0;border-radius:8px;font-size:13px;font-family:inherit;outline:none;width:100%;background:white">' + template_options + '</select>'
        + '<button type="submit" name="mapping_action" value="delete_template" class="btn btn-gray" style="min-width:104px;padding-inline:14px;white-space:nowrap">מחיקה</button>'
        + '</div>'
        + '<label style="display:flex;align-items:center;gap:6px;font-size:13px;color:#334155;margin-bottom:10px"><input type="checkbox" name="save_template" value="1"> שמור כתבנית חדשה</label>'
        + '<label class="field-label">שם תבנית חדשה</label>'
        + '<input type="text" name="template_name" value="' + esc(template_name_value) + '" placeholder="שם תבנית" style="margin-bottom:14px">'
        + '<div style="margin-top:14px;background:linear-gradient(180deg,#f8fbff 0%,#eef6ff 100%);border:1px solid #bfdbfe;border-radius:14px;padding:12px 12px 14px;box-shadow:inset 0 1px 0 rgba(255,255,255,.7)">'
        + '<div style="font-size:13px;font-weight:800;color:#1d4ed8;margin-bottom:4px">סינון הסקריפט</div>'
        + '<div style="font-size:14px;font-weight:700;color:#0f172a;margin-bottom:10px">טווח הבדיקה</div>'
        + '<label class="field-label">סוג בדיקה</label>'
        + '<select name="inactive_period_unit" style="padding:9px 12px;border:1.5px solid #bfdbfe;border-radius:10px;font-size:13px;font-family:inherit;outline:none;width:100%;margin-bottom:10px;background:white">'
        + '<option value="days"' + (' selected' if unit_value == "days" else '') + '>ימים אחרונים</option>'
        + '<option value="months"' + (' selected' if unit_value == "months" else '') + '>חודשים אחרונים</option>'
        + '</select>'
        + '<label class="field-label">ערך הבדיקה</label>'
        + '<input type="text" name="inactive_period_value" value="' + esc(threshold_value) + '" placeholder="לדוגמה 30 או 3" style="margin-bottom:0;border-color:#bfdbfe;background:#ffffff">'
        + '<div style="font-size:12px;color:#475569;line-height:1.7;margin-top:10px">המערכת תבדוק האם לעובד הייתה פעילות בטווח שנבחר, לפי תאריך הייחוס האחרון שקיים בקובץ.</div>'
        + '</div>'
        + '</div>'
        + '<div>'
        + '<div style="display:flex;gap:8px;flex-wrap:wrap;margin-bottom:12px">'
        + '<span style="display:inline-flex;align-items:center;gap:6px;padding:6px 10px;border-radius:999px;background:#fff7ed;border:1px solid #fdba74;font-size:12px;color:#9a3412">שדה קריטי לזיהוי פעילות</span>'
        + '<span style="display:inline-flex;align-items:center;gap:6px;padding:6px 10px;border-radius:999px;background:#ecfdf5;border:1px solid #86efac;font-size:12px;color:#166534">שדה מהדוח</span>'
        + '</div>'
        + '<div style="display:grid;grid-template-columns:repeat(auto-fit,minmax(220px,1fr));gap:10px;margin-bottom:12px">' + fields_html + '</div>'
        + '<div style="font-size:12px;color:#64748b;line-height:1.7;margin-bottom:12px">שדות חובה: שם עובד ותאריך. בנוסף יש לבחור לפחות מזהה אחד נוסף, וגם לבחור או כניסה ויציאה יחד, או לחלופין שדה סה&quot;כ שעות.</div>'
        + '<div style="display:flex;gap:10px;flex-wrap:wrap;justify-content:center"><button type="submit" id="inactiveWorkersConfirmButton" name="mapping_action" value="confirm" class="btn btn-blue" style="min-width:220px">אשר הכל והפעל עיבוד</button><a href="/run/' + script_id + '" class="btn btn-gray" style="text-decoration:none;display:inline-flex;align-items:center;justify-content:center;min-width:150px">העלאת קובץ חדש</a></div>'
        + '</div></div>'
        + '</div>'
        + '<div id="inactiveWorkersProcessingOverlay" style="display:none;position:fixed;inset:0;background:rgba(248,250,252,.78);backdrop-filter:blur(2px);z-index:80;align-items:center;justify-content:center;padding:20px">'
        + '<div style="width:100%;max-width:320px;background:#ffffff;border:1px solid #dbeafe;border-radius:18px;box-shadow:0 20px 50px rgba(15,23,42,.14);padding:24px 20px;text-align:center">'
        + '<div style="width:42px;height:42px;border-radius:999px;border:3px solid #bfdbfe;border-top-color:#2563eb;margin:0 auto 14px;animation:mappingSpin .9s linear infinite"></div>'
        + '<div style="font-size:16px;font-weight:800;color:#1e3a8a;margin-bottom:6px">הדוח בהכנה</div>'
        + '<div style="font-size:13px;line-height:1.7;color:#475569">המערכת מאשרת את השדות ובודקת עובדים ללא פעילות בטווח שנבחר. בקבצים גדולים הפעולה יכולה להימשך מעט זמן.</div>'
        + '</div></div>'
        + '<script>'
        + '(function(){'
        + 'var templateSelect=document.getElementById("selectedInactiveTemplateId");'
        + 'var fieldSelects=Array.prototype.slice.call(document.querySelectorAll(\'select[data-mapping-field="1"]\'));'
        + 'var form=document.getElementById("inactiveWorkersMappingForm");'
        + 'var confirmButton=document.getElementById("inactiveWorkersConfirmButton");'
        + 'var overlay=document.getElementById("inactiveWorkersProcessingOverlay");'
        + 'var templateMappings=' + json.dumps(template_payload, ensure_ascii=False) + ';'
        + 'var fieldLabels=' + json.dumps(mapping_labels, ensure_ascii=False) + ';'
        + 'var selectStyles={critical:{bg:"#fff7ed",border:"#fb923c",shadow:"rgba(249,115,22,.14)"},table_exact:{bg:"#ecfdf5",border:"#4ade80",shadow:"rgba(34,197,94,.14)"},empty:{bg:"#ffffff",border:"#e2e8f0",shadow:"rgba(148,163,184,.08)"}};'
        + 'function applySelectVisual(sel){var isCritical=(sel.name==="entry_time_source"||sel.name==="exit_time_source"||sel.name==="total_hours_source");var kind=isCritical?"critical":"table_exact";var style=selectStyles[kind]||selectStyles.empty;sel.style.backgroundColor=style.bg;sel.style.borderColor=style.border;sel.style.boxShadow="0 0 0 3px "+style.shadow;}'
        + 'function refreshOptionLabels(){var assignments={};fieldSelects.forEach(function(sel){if(sel.value){assignments[sel.value]=sel.name;}});fieldSelects.forEach(function(sel){Array.prototype.forEach.call(sel.options,function(opt){var base=opt.getAttribute("data-base-label")||opt.text;var assigned=assignments[opt.value];var suffix="";if(opt.value&&assigned&&assigned!==sel.name){suffix=" [נבחר עבור "+(fieldLabels[assigned]||assigned)+"]";}opt.text=base+suffix;});applySelectVisual(sel);});}'
        + 'function clearDuplicateSelections(changedSelect){if(!changedSelect.value){refreshOptionLabels();return;}fieldSelects.forEach(function(sel){if(sel!==changedSelect&&sel.value===changedSelect.value){sel.value="";}});refreshOptionLabels();}'
        + 'function applyTemplate(templateId){var mapping=templateMappings[templateId]||{};if(!templateId){refreshOptionLabels();return;}fieldSelects.forEach(function(sel){sel.value=mapping[sel.name]||"";});var seen={};fieldSelects.forEach(function(sel){if(sel.value&&seen[sel.value]){sel.value="";}else if(sel.value){seen[sel.value]=true;}});refreshOptionLabels();}'
        + 'fieldSelects.forEach(function(sel){sel.addEventListener("change",function(){clearDuplicateSelections(sel);});});'
        + 'if(templateSelect){templateSelect.addEventListener("change",function(){applyTemplate(this.value);});}'
        + 'if(form){form.addEventListener("submit",function(event){var submitter=event.submitter||document.activeElement;if(!submitter||submitter.value!=="confirm"){return;}if(confirmButton){confirmButton.disabled=true;confirmButton.textContent="העיבוד התחיל...";}if(overlay){overlay.style.display="flex";}document.body.style.overflow="hidden";});}'
        + 'refreshOptionLabels();'
        + '})();'
        + '</script>'
        + '<style>@keyframes mappingSpin{from{transform:rotate(0deg);}to{transform:rotate(360deg);}}</style>'
        + '</form>'
    )


def build_org_hierarchy_mapping_form(script_id, temp_upload_path, temp_upload_ext, inspection, current_mapping, templates, template_name_value, current_output_type):
    template_options = '<option value="">ללא תבנית שמורה</option>'
    for template in templates:
        template_options += '<option value="' + str(template["id"]) + '">' + esc(template["name"]) + '</option>'

    mapping_labels = {field["name"]: field["label"] for field in ORG_HIERARCHY_MAPPING_FIELDS}
    template_payload = {
        str(template["id"]): {key: str(value or "") for key, value in template["mapping"].items()}
        for template in templates
    }

    fields_html = ""
    for field in ORG_HIERARCHY_MAPPING_FIELDS:
        field_name = field["name"]
        current_value = str(current_mapping.get(field_name, "") or "")
        options = inspection["options_by_field"].get(field_name, [])
        select_options = ""
        for option in options:
            selected = ' selected' if option["value"] == current_value else ""
            select_options += (
                '<option value="' + esc(option["value"]) + '" data-base-label="' + esc(option["label"]) + '" data-source-kind="' + esc(option.get("source_kind", "empty")) + '"' + selected + ">"
                + esc(option["label"])
                + "</option>"
            )

        required_badge = ' <span style="color:#dc2626">*</span>' if field["required"] else ' <span style="color:#94a3b8">(אופציונלי)</span>'
        wrapper_style = ""
        if field.get("critical"):
            wrapper_style = 'background:#fff7ed;border:1px solid #fdba74;border-radius:12px;padding:10px 10px 12px'
        fields_html += (
            '<div style="' + wrapper_style + '"><label class="field-label">' + field["label"] + required_badge + '</label>'
            + ('<div style="font-size:12px;color:#9a3412;line-height:1.6;margin:-4px 0 8px">שדה קריטי לבניית ההיררכיה. יש לוודא שזהו המקור הנכון.</div>' if field.get("critical") else '')
            + '<select name="' + field_name + '" data-mapping-field="1" data-field-label="' + esc(field["label"]) + '" style="padding:9px 12px;border:1.5px solid #e2e8f0;border-radius:8px;font-size:13px;font-family:inherit;outline:none;width:100%;margin-bottom:0;background:white;transition:background-color .15s ease,border-color .15s ease,box-shadow .15s ease">'
            + select_options
            + '</select></div>'
        )

    return (
        '<form method="POST" id="orgHierarchyMappingForm">'
        + '<input type="hidden" name="flow_mode" value="confirm_mapping">'
        + '<input type="hidden" name="temp_upload_path" value="' + esc(temp_upload_path) + '">'
        + '<input type="hidden" name="temp_upload_ext" value="' + esc(temp_upload_ext) + '">'
        + '<div style="background:#fafcff;border:1px solid #dbeafe;border-radius:14px;padding:1rem;margin-bottom:1rem">'
        + '<div style="font-size:15px;font-weight:700;color:#1e3a8a;margin-bottom:10px">אישור שדות לפני עיבוד</div>'
        + '<div style="display:grid;grid-template-columns:260px minmax(0,1fr);gap:14px;align-items:start">'
        + '<div style="background:#ffffff;border:1px solid #e2e8f0;border-radius:12px;padding:12px">'
        + '<div style="font-size:14px;font-weight:700;color:#0f172a;margin-bottom:10px">תבניות שמורות</div>'
        + '<label class="field-label">בחירת תבנית</label>'
        + '<div style="display:grid;grid-template-columns:minmax(0,1fr) auto;gap:8px;align-items:center;margin-bottom:12px">'
        + '<select id="selectedTemplateId" name="selected_template_id" style="padding:9px 12px;border:1.5px solid #e2e8f0;border-radius:8px;font-size:13px;font-family:inherit;outline:none;width:100%;background:white">' + template_options + '</select>'
        + '<button type="submit" name="mapping_action" value="delete_template" class="btn btn-gray" style="min-width:104px;padding-inline:14px;white-space:nowrap">מחיקה</button>'
        + '</div>'
        + '<label style="display:flex;align-items:center;gap:6px;font-size:13px;color:#334155;margin-bottom:10px"><input type="checkbox" name="save_template" value="1"> שמור כתבנית חדשה</label>'
        + '<label class="field-label">שם תבנית חדשה</label>'
        + '<input type="text" name="template_name" value="' + esc(template_name_value) + '" placeholder="שם תבנית" style="margin-bottom:0">'
        + '<div style="font-size:12px;color:#64748b;line-height:1.7;margin-top:10px">בחירת תבנית תעדכן את כל השדות בהתאם. שמירה תיצור תבנית חדשה בלבד ולא תדרוס תבנית קיימת.</div>'
        + '</div>'
        + '<div>'
        + '<div style="margin-bottom:12px"><label class="field-label">סוג פלט</label><select name="output_type" style="padding:9px 12px;border:1.5px solid #e2e8f0;border-radius:8px;font-size:13px;font-family:inherit;outline:none;width:100%;background:white">'
        + '<option value="excel"' + (' selected' if current_output_type == 'excel' else '') + '>אקסל בלבד</option>'
        + '<option value="powerpoint"' + (' selected' if current_output_type == 'powerpoint' else '') + '>PowerPoint בלבד</option>'
        + '<option value="both"' + (' selected' if current_output_type == 'both' else '') + '>XL+PowerPoint</option>'
        + '</select></div>'
        + '<div style="display:flex;gap:8px;flex-wrap:wrap;margin-bottom:12px">'
        + '<span style="display:inline-flex;align-items:center;gap:6px;padding:6px 10px;border-radius:999px;background:#fff7ed;border:1px solid #fdba74;font-size:12px;color:#9a3412">שדה קריטי</span>'
        + '<span style="display:inline-flex;align-items:center;gap:6px;padding:6px 10px;border-radius:999px;background:#ecfdf5;border:1px solid #86efac;font-size:12px;color:#166534">שדה מהקובץ</span>'
        + '</div>'
        + '<div style="display:grid;grid-template-columns:repeat(auto-fit,minmax(220px,1fr));gap:10px;margin-bottom:12px">' + fields_html + '</div>'
        + '<div style="font-size:12px;color:#64748b;line-height:1.7;margin-bottom:12px">שדות חובה: שם עובד, מנהל ישיר ומחלקה. בנוסף יש לבחור לפחות מזהה אחד נוסף: מספר עובד, תעודת זהות או דרכון.</div>'
        + '<div style="display:flex;gap:10px;flex-wrap:wrap;justify-content:center"><button type="submit" id="orgHierarchyConfirmButton" name="mapping_action" value="confirm" class="btn btn-blue" style="min-width:220px">אשר הכל והפעל עיבוד</button><a href="/run/' + script_id + '" class="btn btn-gray" style="text-decoration:none;display:inline-flex;align-items:center;justify-content:center;min-width:150px">העלאת קובץ חדש</a></div>'
        + '</div></div>'
        + '</div>'
        + '<div id="orgHierarchyProcessingOverlay" style="display:none;position:fixed;inset:0;background:rgba(248,250,252,.78);backdrop-filter:blur(2px);z-index:80;align-items:center;justify-content:center;padding:20px">'
        + '<div style="width:100%;max-width:320px;background:#ffffff;border:1px solid #dbeafe;border-radius:18px;box-shadow:0 20px 50px rgba(15,23,42,.14);padding:24px 20px;text-align:center">'
        + '<div style="width:42px;height:42px;border-radius:999px;border:3px solid #bfdbfe;border-top-color:#2563eb;margin:0 auto 14px;animation:mappingSpin .9s linear infinite"></div>'
        + '<div style="font-size:16px;font-weight:800;color:#1e3a8a;margin-bottom:6px">הדוח בהכנה</div>'
        + '<div style="font-size:13px;line-height:1.7;color:#475569">המערכת מאשרת את השדות ובונה את דוח המבנה הארגוני. בקבצים גדולים הפעולה יכולה להימשך מעט זמן.</div>'
        + '</div></div>'
        + '<script>'
        + '(function(){'
        + 'var templateSelect=document.getElementById("selectedTemplateId");'
        + 'var fieldSelects=Array.prototype.slice.call(document.querySelectorAll(\'select[data-mapping-field="1"]\'));'
        + 'var form=document.getElementById("orgHierarchyMappingForm");'
        + 'var confirmButton=document.getElementById("orgHierarchyConfirmButton");'
        + 'var overlay=document.getElementById("orgHierarchyProcessingOverlay");'
        + 'var templateMappings=' + json.dumps(template_payload, ensure_ascii=False) + ';'
        + 'var fieldLabels=' + json.dumps(mapping_labels, ensure_ascii=False) + ';'
        + 'var selectStyles={critical:{bg:"#fff7ed",border:"#fb923c",shadow:"rgba(249,115,22,.14)"},table_exact:{bg:"#ecfdf5",border:"#4ade80",shadow:"rgba(34,197,94,.14)"},empty:{bg:"#ffffff",border:"#e2e8f0",shadow:"rgba(148,163,184,.08)"}};'
        + 'function applySelectVisual(sel){var fieldName=sel.name;var isCritical=(fieldName==="employee_name_source"||fieldName==="direct_manager_source"||fieldName==="department_source");var kind=isCritical?"critical":(((sel.options[sel.selectedIndex]||{}).getAttribute&&sel.options[sel.selectedIndex].getAttribute("data-source-kind"))||"table_exact");var style=selectStyles[kind]||selectStyles.empty;sel.style.backgroundColor=style.bg;sel.style.borderColor=style.border;sel.style.boxShadow="0 0 0 3px "+style.shadow;}'
        + 'function refreshOptionLabels(){var assignments={};fieldSelects.forEach(function(sel){if(sel.value){assignments[sel.value]=sel.name;}});fieldSelects.forEach(function(sel){Array.prototype.forEach.call(sel.options,function(opt){var base=opt.getAttribute("data-base-label")||opt.text;var assigned=assignments[opt.value];var suffix="";if(opt.value&&assigned&&assigned!==sel.name){suffix=" [נבחר עבור "+(fieldLabels[assigned]||assigned)+"]";}opt.text=base+suffix;});applySelectVisual(sel);});}'
        + 'function clearDuplicateSelections(changedSelect){if(!changedSelect.value){refreshOptionLabels();return;}fieldSelects.forEach(function(sel){if(sel!==changedSelect&&sel.value===changedSelect.value){sel.value="";}});refreshOptionLabels();}'
        + 'function applyTemplate(templateId){var mapping=templateMappings[templateId]||{};if(!templateId){refreshOptionLabels();return;}fieldSelects.forEach(function(sel){sel.value=mapping[sel.name]||"";});var seen={};fieldSelects.forEach(function(sel){if(sel.value&&seen[sel.value]){sel.value="";}else if(sel.value){seen[sel.value]=true;}});refreshOptionLabels();}'
        + 'fieldSelects.forEach(function(sel){sel.addEventListener("change",function(){clearDuplicateSelections(sel);});});'
        + 'if(templateSelect){templateSelect.addEventListener("change",function(){applyTemplate(this.value);});}'
        + 'if(form){form.addEventListener("submit",function(event){var submitter=event.submitter||document.activeElement;if(!submitter||submitter.value!=="confirm"){return;}if(confirmButton){confirmButton.disabled=true;confirmButton.textContent="העיבוד התחיל...";}if(overlay){overlay.style.display="flex";}document.body.style.overflow="hidden";});}'
        + 'refreshOptionLabels();'
        + '})();'
        + '</script>'
        + '<style>@keyframes mappingSpin{from{transform:rotate(0deg);}to{transform:rotate(360deg);}}</style>'
        + '</form>'
    )


def build_flamingo_mapping_form(script_id, temp_upload_path, temp_upload_ext, inspection, current_mapping, templates, template_name_value, manual_hourly_rate_value):
    template_options = '<option value="">ללא תבנית שמורה</option>'
    for template in templates:
        template_options += '<option value="' + str(template["id"]) + '">' + esc(template["name"]) + '</option>'

    mapping_labels = {field["name"]: field["label"] for field in FLAMINGO_MAPPING_FIELDS}
    template_payload = {
        str(template["id"]): {key: str(value or "") for key, value in template["mapping"].items()}
        for template in templates
    }

    mapping_fields_html = ""
    for field in FLAMINGO_MAPPING_FIELDS:
        field_name = field["name"]
        current_value = str(current_mapping.get(field_name, "") or "")
        options = inspection["options_by_field"].get(field_name, [])
        blank_options = [option for option in options if not option.get("value")]
        critical_options = [option for option in options if option.get("source_kind") == "critical"]
        meta_options = [option for option in options if option.get("source_kind") == "meta"]
        summary_options = [option for option in options if option.get("source_kind") == "table_exact"]

        def render_option(option):
            selected = ' selected' if option["value"] == current_value else ""
            return (
                '<option value="' + esc(option["value"]) + '" data-base-label="' + esc(option["label"]) + '" data-source-kind="' + esc(option.get("source_kind", "empty")) + '"' + selected + ">"
                + esc(option["label"])
                + "</option>"
            )

        select_options = ""
        for option in blank_options:
            select_options += render_option(option)
        if critical_options:
            select_options += '<optgroup label="בחירה ידנית">'
            for option in critical_options:
                select_options += render_option(option)
            select_options += '</optgroup>'
        if meta_options:
            select_options += '<optgroup label="שדות עליונים">'
            for option in meta_options:
                select_options += render_option(option)
            select_options += '</optgroup>'
        if summary_options:
            select_options += '<optgroup label="שדות סיכום">'
            for option in summary_options:
                select_options += render_option(option)
            select_options += '</optgroup>'

        required_badge = ' <span style="color:#dc2626">*</span>' if field["required"] else ' <span style="color:#94a3b8">(אופציונלי)</span>'
        wrapper_style = ""
        if field.get("critical"):
            wrapper_style = 'background:#fff7ed;border:1px solid #fdba74;border-radius:12px;padding:10px 10px 12px'
        mapping_fields_html += (
            '<div style="' + wrapper_style + '"><label class="field-label">' + field["label"] + required_badge + '</label>'
            + ('<div style="font-size:12px;color:#9a3412;line-height:1.6;margin:-4px 0 8px">שדה קריטי לחישוב השכר. יש לוודא שזהו המקור הנכון.</div>' if field.get("critical") else '')
            + '<select name="' + field_name + '" data-mapping-field="1" data-field-label="' + esc(field["label"]) + '" style="padding:9px 12px;border:1.5px solid #e2e8f0;border-radius:8px;font-size:13px;font-family:inherit;outline:none;width:100%;margin-bottom:0;background:white;transition:background-color .15s ease,border-color .15s ease,box-shadow .15s ease">'
            + select_options
            + '</select>'
            + (
                '<div data-manual-rate-wrap="1" style="' + ('display:block' if current_value == "__manual__" else 'display:none') + ';margin-top:10px">'
                + '<label class="field-label" style="margin-bottom:6px">תעריף שעתי ידני</label>'
                + '<input type="text" name="manual_hourly_rate" value="' + esc(manual_hourly_rate_value) + '" placeholder="לדוגמה 45.5" style="margin-bottom:0">'
                + '<div style="font-size:12px;color:#9a3412;line-height:1.6;margin-top:8px">אם תבחר תעריף ידני, כל העובדים בדוח יחושבו לפי אותו תעריף.</div>'
                + '</div>'
                if field_name == "hourly_rate_source" else ""
            )
            + '</div>'
        )

    return (
        '<form method="POST" id="flamingoMappingConfirmForm">'
        + '<input type="hidden" name="flow_mode" value="confirm_mapping">'
        + '<input type="hidden" name="temp_upload_path" value="' + esc(temp_upload_path) + '">'
        + '<input type="hidden" name="temp_upload_ext" value="' + esc(temp_upload_ext) + '">'
        + '<div style="background:#fafcff;border:1px solid #dbeafe;border-radius:14px;padding:1rem;margin-bottom:1rem">'
        + '<div style="font-size:15px;font-weight:700;color:#1e3a8a;margin-bottom:10px">אישור שדות לפני חישוב שכר</div>'
        + '<div style="display:grid;grid-template-columns:260px minmax(0,1fr);gap:14px;align-items:start">'
        + '<div style="background:#ffffff;border:1px solid #e2e8f0;border-radius:12px;padding:12px">'
        + '<div style="font-size:14px;font-weight:700;color:#0f172a;margin-bottom:10px">תבניות שמורות</div>'
        + '<label class="field-label">בחירת תבנית</label>'
        + '<div style="display:grid;grid-template-columns:minmax(0,1fr) auto;gap:8px;align-items:center;margin-bottom:12px">'
        + '<select id="selectedFlamingoTemplateId" name="selected_template_id" style="padding:9px 12px;border:1.5px solid #e2e8f0;border-radius:8px;font-size:13px;font-family:inherit;outline:none;width:100%;background:white">' + template_options + '</select>'
        + '<button type="submit" name="mapping_action" value="delete_template" class="btn btn-gray" style="min-width:104px;padding-inline:14px;white-space:nowrap">מחיקה</button>'
        + '</div>'
        + '<label style="display:flex;align-items:center;gap:6px;font-size:13px;color:#334155;margin-bottom:10px"><input type="checkbox" name="save_template" value="1"> שמור כתבנית חדשה</label>'
        + '<label class="field-label">שם תבנית חדשה</label>'
        + '<input type="text" name="template_name" value="' + esc(template_name_value) + '" placeholder="שם תבנית" style="margin-bottom:0">'
        + '<div style="font-size:12px;color:#64748b;line-height:1.7;margin-top:10px">בחירת תבנית תעדכן את כל השדות בהתאם. שמירה תיצור תבנית חדשה בלבד ולא תדרוס תבנית קיימת.</div>'
        + '</div>'
        + '<div>'
        + '<div style="display:flex;gap:8px;flex-wrap:wrap;margin-bottom:12px">'
        + '<span style="display:inline-flex;align-items:center;gap:6px;padding:6px 10px;border-radius:999px;background:#fff7ed;border:1px solid #fdba74;font-size:12px;color:#9a3412">שדה קריטי לחישוב</span>'
        + '<span style="display:inline-flex;align-items:center;gap:6px;padding:6px 10px;border-radius:999px;background:#eff6ff;border:1px solid #bfdbfe;font-size:12px;color:#1d4ed8">שדה עליון</span>'
        + '<span style="display:inline-flex;align-items:center;gap:6px;padding:6px 10px;border-radius:999px;background:#ecfdf5;border:1px solid #86efac;font-size:12px;color:#166534">שדה סיכום</span>'
        + '</div>'
        + '<div style="display:grid;grid-template-columns:repeat(auto-fit,minmax(220px,1fr));gap:10px;margin-bottom:12px">' + mapping_fields_html + '</div>'
        + '<div style="font-size:12px;color:#64748b;line-height:1.7;margin-bottom:12px">שדות קריטיים: תעריף שעתי ושעות לתשלום בפועל. בלי השניים האלה חישוב השכר לא יהיה אמין.</div>'
        + '<div style="display:flex;gap:10px;flex-wrap:wrap;justify-content:center"><button type="submit" id="flamingoMappingConfirmButton" name="mapping_action" value="confirm" class="btn btn-blue" style="min-width:220px">אשר הכל והפעל חישוב</button><a href="/run/' + script_id + '" class="btn btn-gray" style="text-decoration:none;display:inline-flex;align-items:center;justify-content:center;min-width:150px">העלאת קובץ חדש</a></div>'
        + '</div></div>'
        + '</div>'
        + '<div id="flamingoProcessingOverlay" style="display:none;position:fixed;inset:0;background:rgba(248,250,252,.78);backdrop-filter:blur(2px);z-index:80;align-items:center;justify-content:center;padding:20px">'
        + '<div style="width:100%;max-width:320px;background:#ffffff;border:1px solid #dbeafe;border-radius:18px;box-shadow:0 20px 50px rgba(15,23,42,.14);padding:24px 20px;text-align:center">'
        + '<div style="width:42px;height:42px;border-radius:999px;border:3px solid #bfdbfe;border-top-color:#2563eb;margin:0 auto 14px;animation:mappingSpin .9s linear infinite"></div>'
        + '<div style="font-size:16px;font-weight:800;color:#1e3a8a;margin-bottom:6px">חישוב השכר בהכנה</div>'
        + '<div style="font-size:13px;line-height:1.7;color:#475569">המערכת מאשרת את השדות ומחשבת את השכר. בדוחות גדולים הפעולה יכולה להימשך מעט זמן.</div>'
        + '</div></div>'
        + '<script>'
        + '(function(){'
        + 'var templateSelect=document.getElementById("selectedFlamingoTemplateId");'
        + 'var fieldSelects=Array.prototype.slice.call(document.querySelectorAll(\'select[data-mapping-field="1"]\'));'
        + 'var form=document.getElementById("flamingoMappingConfirmForm");'
        + 'var confirmButton=document.getElementById("flamingoMappingConfirmButton");'
        + 'var overlay=document.getElementById("flamingoProcessingOverlay");'
        + 'var templateMappings=' + json.dumps(template_payload, ensure_ascii=False) + ';'
        + 'var fieldLabels=' + json.dumps(mapping_labels, ensure_ascii=False) + ';'
        + 'var selectStyles={critical:{bg:"#fff7ed",border:"#fb923c",shadow:"rgba(249,115,22,.14)"},meta:{bg:"#eff6ff",border:"#60a5fa",shadow:"rgba(59,130,246,.12)"},table_exact:{bg:"#ecfdf5",border:"#4ade80",shadow:"rgba(34,197,94,.14)"},empty:{bg:"#ffffff",border:"#e2e8f0",shadow:"rgba(148,163,184,.08)"}};'
        + 'function applySelectVisual(sel){var opt=sel.options[sel.selectedIndex];var kind=(opt&&opt.getAttribute("data-source-kind"))||"empty";var style=selectStyles[kind]||selectStyles.empty;sel.style.backgroundColor=style.bg;sel.style.borderColor=style.border;sel.style.boxShadow="0 0 0 3px "+style.shadow;if(sel.name==="hourly_rate_source"){var wrap=document.querySelector("[data-manual-rate-wrap=\'1\']");if(wrap){wrap.style.display=(sel.value==="__manual__")?"block":"none";}}}'
        + 'function refreshOptionLabels(){var assignments={};fieldSelects.forEach(function(sel){if(sel.value && sel.value!=="__manual__"){assignments[sel.value]=sel.name;}});fieldSelects.forEach(function(sel){Array.prototype.forEach.call(sel.options,function(opt){var base=opt.getAttribute("data-base-label")||opt.text;var assigned=assignments[opt.value];var suffix="";if(opt.value && opt.value!=="__manual__" && assigned && assigned!==sel.name){suffix=" [נבחר עבור "+(fieldLabels[assigned]||assigned)+"]";}opt.text=base+suffix;});applySelectVisual(sel);});}'
        + 'function clearDuplicateSelections(changedSelect){if(!changedSelect.value || changedSelect.value==="__manual__"){refreshOptionLabels();return;}fieldSelects.forEach(function(sel){if(sel!==changedSelect && sel.value===changedSelect.value){sel.value="";}});refreshOptionLabels();}'
        + 'function applyTemplate(templateId){var mapping=templateMappings[templateId]||{};if(!templateId){refreshOptionLabels();return;}fieldSelects.forEach(function(sel){sel.value=mapping[sel.name]||"";});var manualRate=document.querySelector(\'input[name="manual_hourly_rate"]\');if(manualRate){manualRate.value=mapping.manual_hourly_rate||manualRate.value||"";}var seen={};fieldSelects.forEach(function(sel){if(sel.value && sel.value!=="__manual__" && seen[sel.value]){sel.value="";}else if(sel.value && sel.value!=="__manual__"){seen[sel.value]=true;}});refreshOptionLabels();}'
        + 'fieldSelects.forEach(function(sel){sel.addEventListener("change",function(){clearDuplicateSelections(sel);});});'
        + 'if(templateSelect){templateSelect.addEventListener("change",function(){applyTemplate(this.value);});}'
        + 'if(form){form.addEventListener("submit",function(event){var submitter=event.submitter||document.activeElement;if(!submitter||submitter.value!=="confirm"){return;}if(confirmButton){confirmButton.disabled=true;confirmButton.textContent="החישוב התחיל...";}if(overlay){overlay.style.display="flex";}document.body.style.overflow="hidden";});}'
        + 'refreshOptionLabels();'
        + '})();'
        + '</script>'
        + '<style>@keyframes mappingSpin{from{transform:rotate(0deg);}to{transform:rotate(360deg);}}</style>'
        + '</form>'
    )


def get_account_status(user_row):
    today = date.today()
    trial_start = parse_iso_date(user_row["trial_start_date"])
    valid_until = parse_iso_date(user_row["service_valid_until"])

    if valid_until:
        if valid_until >= today:
            return {
                "status_key": "active",
                "status_label_he": "לקוח פעיל",
                "status_label_en": "Active customer",
                "renewal_date": valid_until,
                "days_remaining": None,
            }
        return {
            "status_key": "expired",
            "status_label_he": "שירות שפג תוקפו",
            "status_label_en": "Expired service",
            "renewal_date": valid_until,
            "days_remaining": None,
        }

    if trial_start:
        trial_days = user_row["trial_days"] or 30
        days_remaining = max(0, trial_days - (today - trial_start).days)
        if days_remaining == 0:
            return {
                "status_key": "expired",
                "status_label_he": "לא בשירות",
                "status_label_en": "Not in service",
                "renewal_date": None,
                "days_remaining": None,
            }
        return {
            "status_key": "trial",
            "status_label_he": "ניסיון ל-" + str(trial_days) + " יום",
            "status_label_en": str(trial_days) + "-day trial",
            "renewal_date": None,
            "days_remaining": days_remaining,
        }

    return {
        "status_key": "unknown",
        "status_label_he": "לא בשירות",
        "status_label_en": "Not in service",
        "renewal_date": None,
        "days_remaining": None,
    }


def render(title, body, nav=True, lang="en", topbar_greeting="Hello, ", logout_label="Logout", show_lang_switch=False):
    direction = get_flow_dir(lang)
    topbar = ""
    if nav:
        name = session.get("name", "")
        lang_switch = build_lang_switch(lang) if show_lang_switch else ""
        topbar = (
            '<div class="topbar">'
            '<a href="/dashboard" style="text-decoration:none;color:inherit;display:flex;align-items:center;gap:8px"><h1>&#9201; Scriptly</h1></a>'
            '<div style="display:flex;gap:16px;align-items:center">'
            + lang_switch
            + '<span style="font-size:13px;color:#93c5fd;font-weight:500">' + topbar_greeting + name + "</span>"
            + '<a href="/logout" style="padding:6px 14px;border:1px solid rgba(255,255,255,.2);border-radius:999px;font-size:12px;font-weight:600;transition:all .2s">' + logout_label + '</a>'
            "</div></div>"
        )
    wrap_cls = "wrap" if nav else "login-wrap"
    body_open = "<body>"
    body_close = ""
    if not nav:
        body_open = '<body><div class="login-bg"><div class="login-orb login-orb-1"></div><div class="login-orb login-orb-2"></div>'
        body_close = '</div>'
    return (
        '<!DOCTYPE html><html dir="' + direction + '" lang="' + lang + '">'
        "<head><meta charset=\"UTF-8\">"
        '<meta name="viewport" content="width=device-width,initial-scale=1">'
        "<title>" + title + " | Scriptly</title>"
        "<style>" + CSS + "</style></head>"
        + body_open
        + topbar
        + '<div class="' + wrap_cls + '">'
        + pop_flashes()
        + body
        + "</div>"
        + body_close
        + '<script>'
        + '(function(){'
        + 'function resetTransientUi(){'
        + 'document.querySelectorAll(".modal-bg").forEach(function(el){el.style.display="none";});'
        + 'document.body.style.removeProperty("overflow");'
        + '}'
        + 'function initFlashToasts(){'
        + 'var stack=document.getElementById("flashStack");'
        + 'if(!stack){return;}'
        + 'var toasts=stack.querySelectorAll(".flash-toast");'
        + 'toasts.forEach(function(toast,index){'
        + 'window.setTimeout(function(){'
        + 'toast.classList.add("dismiss");'
        + 'window.setTimeout(function(){'
        + 'if(toast&&toast.parentNode){toast.parentNode.removeChild(toast);}'
        + 'if(stack && !stack.children.length && stack.parentNode){stack.parentNode.removeChild(stack);}'
        + '},220);'
        + '},5000+(index*150));'
        + '});'
        + '}'
        + 'if(document.readyState==="loading"){document.addEventListener("DOMContentLoaded",resetTransientUi);}else{resetTransientUi();}'
        + 'if(document.readyState==="loading"){document.addEventListener("DOMContentLoaded",initFlashToasts);}else{initFlashToasts();}'
        + 'window.addEventListener("pageshow",resetTransientUi);'
        + '})();'
        + "</script></body></html>"
    )


def login_required(f):
    from functools import wraps

    @wraps(f)
    def decorated(*args, **kwargs):
        if "user_id" not in session:
            return redirect("/login")
        return f(*args, **kwargs)

    return decorated


def admin_required(f):
    from functools import wraps

    @wraps(f)
    def decorated(*args, **kwargs):
        if not session.get("is_admin"):
            return redirect("/dashboard")
        return f(*args, **kwargs)

    return decorated


@app.route("/")
def homepage():
    if "user_id" in session:
        return redirect("/admin" if session.get("is_admin") else "/dashboard")
    return (
        '<!DOCTYPE html>'
        '<html dir="rtl" lang="he">'
        '<head>'
        '<meta charset="UTF-8">'
        '<meta name="viewport" content="width=device-width,initial-scale=1">'
        '<title>Scriptly — כלי AI חכמים לניהול נוכחות ושכר</title>'
        '<meta name="description" content="פלטפורמת AI שמייצרת כלים חכמים לניהול נוכחות ושכר. תארו מה אתם צריכים בעברית — והכלי מוכן. התחילו בחינם.">'
        '<meta name="keywords" content="AI, אוטומציה, דוחות נוכחות, שכר, HR, בינה מלאכותית, כלים חכמים, חינם, ללא כרטיס אשראי">'
        '<meta property="og:title" content="Scriptly — כלי AI חכמים לניהול נוכחות ושכר">'
        '<meta property="og:description" content="תארו מה אתם צריכים בעברית — ה-AI בונה את הכלי. התחילו בחינם.">'
        '<meta property="og:type" content="website">'
        '<meta property="og:url" content="https://script-ly.com">'
        '<style>'
        '*, *::before, *::after { box-sizing: border-box; margin: 0; padding: 0; }'
        'html { scroll-behavior: smooth; }'
        "body { font-family: 'Segoe UI', -apple-system, BlinkMacSystemFont, Arial, sans-serif; background: #ffffff; color: #1e293b; direction: rtl; overflow-x: hidden; }"
        '.hp-nav { position: fixed; top: 0; right: 0; left: 0; z-index: 100; height: 64px; display: flex; align-items: center; justify-content: space-between; padding: 0 2.5rem; transition: background .4s, box-shadow .4s; }'
        '.hp-nav.scrolled { background: rgba(255,255,255,.92); backdrop-filter: blur(16px); -webkit-backdrop-filter: blur(16px); box-shadow: 0 1px 20px rgba(0,0,0,.06); }'
        '.hp-nav-brand { color: white; font-size: 20px; font-weight: 800; text-decoration: none; display: flex; align-items: center; gap: 10px; letter-spacing: -0.3px; transition: color .4s; }'
        '.hp-nav.scrolled .hp-nav-brand { color: #1e3a8a; }'
        '.hp-nav-links { display: flex; gap: 8px; align-items: center; }'
        '.hp-nav-links a { font-size: 14px; text-decoration: none; padding: 8px 18px; border-radius: 10px; transition: all .25s; font-weight: 500; }'
        '.hp-nav-links .nav-login { color: rgba(255,255,255,.85); }'
        '.hp-nav-links .nav-login:hover { color: white; background: rgba(255,255,255,.12); }'
        '.hp-nav.scrolled .nav-login { color: #475569; }'
        '.hp-nav.scrolled .nav-login:hover { color: #1e3a8a; background: #f1f5f9; }'
        '.hp-nav-links .btn-cta { background: rgba(255,255,255,.18); color: white; font-weight: 600; border: 1.5px solid rgba(255,255,255,.3); }'
        '.hp-nav-links .btn-cta:hover { background: rgba(255,255,255,.28); border-color: rgba(255,255,255,.5); transform: translateY(-1px); }'
        '.hp-nav.scrolled .btn-cta { background: #2563eb; color: white; border-color: #2563eb; }'
        '.hp-nav.scrolled .btn-cta:hover { background: #1d4ed8; border-color: #1d4ed8; transform: translateY(-1px); box-shadow: 0 4px 16px rgba(37,99,235,.3); }'
        '.hp-hero { min-height: 100vh; display: flex; align-items: center; justify-content: center; text-align: center; padding: 100px 2rem 4rem; position: relative; overflow: hidden; background: linear-gradient(160deg, #0f1b3d 0%, #1e3a8a 35%, #2563eb 70%, #3b82f6 100%); }'
        '.hp-hero-bg { position: absolute; inset: 0; overflow: hidden; pointer-events: none; }'
        '.hp-hero-orb { position: absolute; border-radius: 50%; filter: blur(80px); opacity: 0.12; }'
        '.hp-hero-orb.orb1 { width: 600px; height: 600px; background: #60a5fa; top: -200px; left: -100px; animation: orbMove1 14s ease-in-out infinite; }'
        '.hp-hero-orb.orb2 { width: 450px; height: 450px; background: #818cf8; bottom: -150px; right: -100px; animation: orbMove2 12s ease-in-out 3s infinite; }'
        '.hp-hero-orb.orb3 { width: 300px; height: 300px; background: #34d399; top: 50%; left: 60%; animation: orbMove3 10s ease-in-out 1s infinite; }'
        '.hp-hero-grid { position: absolute; inset: 0; background-image: radial-gradient(rgba(255,255,255,.04) 1px, transparent 1px); background-size: 40px 40px; }'
        '.hp-hero-content { position: relative; z-index: 2; max-width: 800px; }'
        '.hp-hero-badge { display: inline-flex; align-items: center; gap: 8px; background: rgba(255,255,255,.1); border: 1px solid rgba(255,255,255,.15); border-radius: 100px; padding: 8px 20px 8px 16px; font-size: 13px; color: #bfdbfe; margin-bottom: 2rem; animation: fadeInUp .7s ease-out both; backdrop-filter: blur(4px); }'
        '.hp-hero-badge-dot { width: 8px; height: 8px; background: #34d399; border-radius: 50%; animation: dotPulse 2s ease-in-out infinite; }'
        '.hp-hero h1 { font-size: clamp(32px, 5.5vw, 56px); font-weight: 800; color: white; margin-bottom: 1.25rem; line-height: 1.25; letter-spacing: -0.5px; animation: fadeInUp .7s ease-out .15s both; }'
        '.hp-hero h1 span { background: linear-gradient(135deg, #93c5fd, #67e8f9); -webkit-background-clip: text; -webkit-text-fill-color: transparent; background-clip: text; }'
        '.hp-hero-sub { font-size: clamp(16px, 2.5vw, 20px); color: #94a3b8; margin-bottom: 2.5rem; line-height: 1.75; animation: fadeInUp .7s ease-out .3s both; max-width: 600px; margin-left: auto; margin-right: auto; }'
        '.hp-hero-btns { display: flex; gap: 16px; justify-content: center; flex-wrap: wrap; animation: fadeInUp .7s ease-out .45s both; }'
        '.hp-btn-primary { display: inline-flex; align-items: center; gap: 10px; padding: 16px 36px; border-radius: 14px; font-size: 17px; font-weight: 700; text-decoration: none; background: white; color: #1e3a8a; box-shadow: 0 4px 24px rgba(0,0,0,.15), 0 0 0 0 rgba(255,255,255,.3); transition: all .3s cubic-bezier(.4,0,.2,1); }'
        '.hp-btn-primary:hover { transform: translateY(-3px); box-shadow: 0 12px 40px rgba(0,0,0,.2), 0 0 0 3px rgba(255,255,255,.1); }'
        '.hp-btn-primary svg { width: 20px; height: 20px; transition: transform .3s; }'
        '.hp-btn-primary:hover svg { transform: translateX(-4px); }'
        '.hp-btn-secondary { display: inline-flex; align-items: center; gap: 8px; padding: 16px 32px; border-radius: 14px; font-size: 16px; font-weight: 600; text-decoration: none; color: rgba(255,255,255,.9); border: 1.5px solid rgba(255,255,255,.2); background: rgba(255,255,255,.05); transition: all .3s; backdrop-filter: blur(4px); }'
        '.hp-btn-secondary:hover { border-color: rgba(255,255,255,.4); background: rgba(255,255,255,.1); transform: translateY(-2px); }'
        '.hp-hero-stats { display: flex; gap: 3rem; justify-content: center; margin-top: 3.5rem; animation: fadeInUp .7s ease-out .6s both; }'
        '.hp-hero-stat { text-align: center; }'
        '.hp-hero-stat-num { font-size: 28px; font-weight: 800; color: white; display: block; }'
        '.hp-hero-stat-label { font-size: 13px; color: #94a3b8; margin-top: 2px; }'
        '.hp-section { padding: 6rem 2rem; }'
        '.hp-section-inner { max-width: 1100px; margin: 0 auto; }'
        '.hp-section-label { display: inline-block; font-size: 12px; font-weight: 700; text-transform: uppercase; letter-spacing: 2px; color: #2563eb; background: #eff6ff; border-radius: 100px; padding: 6px 16px; margin-bottom: 1rem; }'
        '.hp-section-title { font-size: clamp(24px, 4vw, 38px); font-weight: 800; color: #0f172a; margin-bottom: .75rem; line-height: 1.3; letter-spacing: -0.3px; }'
        '.hp-section-sub { font-size: 16px; color: #64748b; margin-bottom: 3.5rem; line-height: 1.7; max-width: 550px; }'
        '.hp-section-center .hp-section-sub { margin-left: auto; margin-right: auto; text-align: center; }'
        '.hp-section-center .hp-section-title { text-align: center; }'
        '.hp-section-center .hp-section-label { text-align: center; }'
        '.hp-features { display: grid; grid-template-columns: repeat(3, 1fr); gap: 1.5rem; }'
        '.hp-feature { background: #ffffff; border: 1px solid #e2e8f0; border-radius: 20px; padding: 2.25rem; position: relative; overflow: hidden; opacity: 0; transform: translateY(30px); transition: opacity .6s ease-out, transform .6s ease-out, border-color .3s, box-shadow .3s; }'
        '.hp-feature.visible { opacity: 1; transform: translateY(0); }'
        '.hp-feature:hover { border-color: #bfdbfe; box-shadow: 0 8px 32px rgba(37,99,235,.08); }'
        '.hp-feature-icon-wrap { width: 52px; height: 52px; border-radius: 14px; display: flex; align-items: center; justify-content: center; font-size: 26px; margin-bottom: 1.25rem; }'
        '.hp-feature h3 { font-size: 17px; font-weight: 700; color: #0f172a; margin-bottom: .5rem; line-height: 1.4; }'
        '.hp-feature p { font-size: 14px; color: #64748b; line-height: 1.7; }'
        '.feat-ai .hp-feature-icon-wrap { background: linear-gradient(135deg, #ede9fe, #ddd6fe); }'
        '.feat-tools .hp-feature-icon-wrap { background: linear-gradient(135deg, #dbeafe, #bfdbfe); }'
        '.feat-market .hp-feature-icon-wrap { background: linear-gradient(135deg, #d1fae5, #a7f3d0); }'
        '.feat-file .hp-feature-icon-wrap { background: linear-gradient(135deg, #fef3c7, #fde68a); }'
        '.feat-il .hp-feature-icon-wrap { background: linear-gradient(135deg, #fce7f3, #fbcfe8); }'
        '.feat-free .hp-feature-icon-wrap { background: linear-gradient(135deg, #ccfbf1, #99f6e4); }'
        '.hp-showcase { background: linear-gradient(135deg, #f8fafc 0%, #eff6ff 100%); }'
        '.hp-showcase-grid { display: grid; grid-template-columns: 1fr 1fr; gap: 4rem; align-items: center; }'
        '.hp-showcase-visual { background: #0f172a; border-radius: 20px; padding: 2rem; box-shadow: 0 24px 64px rgba(0,0,0,.12); position: relative; overflow: hidden; }'
        '.hp-showcase-visual::before { content: ""; position: absolute; top: 0; right: 0; left: 0; height: 40px; background: #1e293b; border-radius: 20px 20px 0 0; }'
        '.hp-chat-dots { position: absolute; top: 12px; right: 16px; display: flex; gap: 6px; z-index: 1; }'
        '.hp-chat-dots span { width: 10px; height: 10px; border-radius: 50%; }'
        '.hp-chat-dots span:nth-child(1) { background: #ef4444; }'
        '.hp-chat-dots span:nth-child(2) { background: #eab308; }'
        '.hp-chat-dots span:nth-child(3) { background: #22c55e; }'
        '.hp-chat-msgs { padding-top: 1.5rem; display: flex; flex-direction: column; gap: 1rem; }'
        '.hp-chat-msg { padding: 12px 16px; border-radius: 14px; font-size: 14px; line-height: 1.6; max-width: 85%; opacity: 0; animation: chatMsgIn .5s ease-out forwards; }'
        '.hp-chat-msg.user { background: #2563eb; color: white; align-self: flex-start; border-bottom-right-radius: 4px; }'
        '.hp-chat-msg.ai { background: #1e293b; color: #e2e8f0; align-self: flex-end; border-bottom-left-radius: 4px; border: 1px solid #334155; }'
        '.hp-chat-msg.ai .tag { display: inline-block; background: #334155; color: #67e8f9; font-size: 11px; padding: 2px 8px; border-radius: 6px; margin-top: 6px; font-weight: 600; }'
        '.hp-chat-msg:nth-child(1) { animation-delay: .3s; }'
        '.hp-chat-msg:nth-child(2) { animation-delay: 1s; }'
        '.hp-chat-msg:nth-child(3) { animation-delay: 1.8s; }'
        '.hp-showcase-text h3 { font-size: clamp(20px, 3vw, 28px); font-weight: 800; color: #0f172a; margin-bottom: 1rem; line-height: 1.35; }'
        '.hp-showcase-text p { font-size: 15px; color: #64748b; line-height: 1.8; margin-bottom: 1.5rem; }'
        '.hp-showcase-checks { list-style: none; display: flex; flex-direction: column; gap: 12px; }'
        '.hp-showcase-checks li { display: flex; align-items: flex-start; gap: 10px; font-size: 14px; color: #334155; line-height: 1.6; }'
        '.hp-showcase-checks .check { width: 22px; height: 22px; background: #d1fae5; border-radius: 50%; display: flex; align-items: center; justify-content: center; flex-shrink: 0; color: #047857; font-size: 12px; font-weight: 700; margin-top: 2px; }'
        '.hp-how { background: #ffffff; }'
        '.hp-steps { display: grid; grid-template-columns: repeat(3, 1fr); gap: 2.5rem; position: relative; }'
        '.hp-steps::before { content: ""; position: absolute; top: 36px; right: 80px; left: 80px; height: 2px; background: linear-gradient(90deg, #e2e8f0, #bfdbfe, #e2e8f0); z-index: 0; }'
        '.hp-step { text-align: center; position: relative; z-index: 1; opacity: 0; transform: translateY(30px); transition: opacity .6s ease-out, transform .6s ease-out; }'
        '.hp-step.visible { opacity: 1; transform: translateY(0); }'
        '.hp-step-num { width: 56px; height: 56px; background: linear-gradient(135deg, #2563eb, #1d4ed8); color: white; border-radius: 50%; display: inline-flex; align-items: center; justify-content: center; font-size: 22px; font-weight: 800; margin-bottom: 1.25rem; box-shadow: 0 4px 16px rgba(37,99,235,.25), 0 0 0 6px #eff6ff; position: relative; }'
        '.hp-step h3 { font-size: 17px; font-weight: 700; color: #0f172a; margin-bottom: .5rem; }'
        '.hp-step p { font-size: 14px; color: #64748b; line-height: 1.7; max-width: 260px; margin: 0 auto; }'
        '.hp-scripts { background: linear-gradient(135deg, #f8fafc 0%, #f1f5f9 100%); }'
        '.hp-script-cards { display: grid; grid-template-columns: repeat(auto-fit, minmax(280px, 1fr)); gap: 1.25rem; }'
        '.hp-script-card { background: white; border: 1px solid #e2e8f0; border-radius: 16px; padding: 1.75rem; display: flex; align-items: flex-start; gap: 1rem; opacity: 0; transform: translateY(20px); transition: all .5s ease-out; }'
        '.hp-script-card.visible { opacity: 1; transform: translateY(0); }'
        '.hp-script-card:hover { border-color: #bfdbfe; box-shadow: 0 4px 20px rgba(37,99,235,.06); }'
        '.hp-script-icon { width: 44px; height: 44px; border-radius: 12px; background: linear-gradient(135deg, #eff6ff, #dbeafe); display: flex; align-items: center; justify-content: center; font-size: 22px; flex-shrink: 0; }'
        '.hp-script-info h4 { font-size: 15px; font-weight: 700; color: #0f172a; margin-bottom: 4px; }'
        '.hp-script-info p { font-size: 13px; color: #64748b; line-height: 1.6; }'
        '.hp-trust-section { background: #ffffff; }'
        '.hp-trust-grid { display: grid; grid-template-columns: repeat(4, 1fr); gap: 2rem; }'
        '.hp-trust-item { text-align: center; padding: 2rem 1.5rem; border-radius: 16px; border: 1px solid #f1f5f9; opacity: 0; transform: translateY(20px); transition: all .5s ease-out; }'
        '.hp-trust-item.visible { opacity: 1; transform: translateY(0); }'
        '.hp-trust-item:hover { border-color: #e2e8f0; background: #f8fafc; }'
        '.hp-trust-icon { font-size: 32px; margin-bottom: .75rem; }'
        '.hp-trust-item h4 { font-size: 15px; font-weight: 700; color: #0f172a; margin-bottom: .25rem; }'
        '.hp-trust-item p { font-size: 13px; color: #64748b; line-height: 1.6; }'
        '.hp-cta-section { background: linear-gradient(160deg, #0f1b3d 0%, #1e3a8a 50%, #2563eb 100%); padding: 6rem 2rem; position: relative; overflow: hidden; }'
        '.hp-cta-inner { max-width: 700px; margin: 0 auto; text-align: center; position: relative; z-index: 2; }'
        '.hp-cta-inner h2 { font-size: clamp(24px, 4vw, 38px); font-weight: 800; color: white; margin-bottom: 1rem; line-height: 1.3; }'
        '.hp-cta-inner p { font-size: 17px; color: #94a3b8; margin-bottom: 2.5rem; line-height: 1.7; }'
        '.hp-cta-btn { display: inline-flex; align-items: center; gap: 10px; padding: 18px 44px; border-radius: 14px; font-size: 18px; font-weight: 700; text-decoration: none; background: white; color: #1e3a8a; box-shadow: 0 4px 24px rgba(0,0,0,.15); transition: all .3s cubic-bezier(.4,0,.2,1); }'
        '.hp-cta-btn:hover { transform: translateY(-3px); box-shadow: 0 12px 40px rgba(0,0,0,.2); }'
        '.hp-cta-badges { display: flex; gap: 2rem; justify-content: center; flex-wrap: wrap; margin-top: 2rem; }'
        '.hp-cta-badge { font-size: 14px; font-weight: 500; color: #94a3b8; display: flex; align-items: center; gap: 6px; }'
        '.hp-cta-badge svg { color: #34d399; width: 18px; height: 18px; }'
        '.hp-footer { background: #0f172a; color: #94a3b8; padding: 3rem 2rem; }'
        '.hp-footer-inner { max-width: 1100px; margin: 0 auto; display: flex; align-items: center; justify-content: space-between; flex-wrap: wrap; gap: 1rem; }'
        '.hp-footer-brand { font-size: 16px; font-weight: 700; color: #e2e8f0; display: flex; align-items: center; gap: 8px; }'
        '.hp-footer-links { display: flex; gap: 1.5rem; }'
        '.hp-footer-links a { color: #94a3b8; text-decoration: none; font-size: 14px; transition: color .2s; }'
        '.hp-footer-links a:hover { color: #e2e8f0; }'
        '.hp-footer-copy { font-size: 13px; color: #475569; width: 100%; text-align: center; margin-top: 1.5rem; padding-top: 1.5rem; border-top: 1px solid #1e293b; }'
        '@keyframes fadeInUp { from { opacity: 0; transform: translateY(30px); } to { opacity: 1; transform: translateY(0); } }'
        '@keyframes orbMove1 { 0%, 100% { transform: translate(0, 0) scale(1); } 50% { transform: translate(40px, -30px) scale(1.05); } }'
        '@keyframes orbMove2 { 0%, 100% { transform: translate(0, 0) scale(1); } 50% { transform: translate(-30px, 20px) scale(1.08); } }'
        '@keyframes orbMove3 { 0%, 100% { transform: translate(0, 0) scale(1); } 50% { transform: translate(20px, 30px) scale(.95); } }'
        '@keyframes dotPulse { 0%, 100% { opacity: 1; transform: scale(1); } 50% { opacity: .5; transform: scale(.8); } }'
        '@keyframes chatMsgIn { from { opacity: 0; transform: translateY(12px); } to { opacity: 1; transform: translateY(0); } }'
        '@media (max-width: 900px) {'
        '  .hp-features { grid-template-columns: repeat(2, 1fr); }'
        '  .hp-showcase-grid { grid-template-columns: 1fr; gap: 2.5rem; }'
        '  .hp-trust-grid { grid-template-columns: repeat(2, 1fr); }'
        '  .hp-steps::before { display: none; }'
        '}'
        '@media (max-width: 600px) {'
        '  .hp-nav { padding: 0 1rem; height: 56px; }'
        '  .hp-hero { padding: 80px 1.25rem 3rem; min-height: auto; }'
        '  .hp-hero-stats { gap: 1.5rem; }'
        '  .hp-section { padding: 4rem 1.25rem; }'
        '  .hp-features { grid-template-columns: 1fr; }'
        '  .hp-steps { grid-template-columns: 1fr; gap: 2rem; }'
        '  .hp-trust-grid { grid-template-columns: 1fr; }'
        '  .hp-hero-btns { flex-direction: column; align-items: center; }'
        '  .hp-btn-primary, .hp-btn-secondary { width: 100%; justify-content: center; }'
        '  .hp-footer-inner { flex-direction: column; text-align: center; }'
        '  .hp-footer-links { justify-content: center; }'
        '  .hp-script-cards { grid-template-columns: 1fr; }'
        '}'
        '</style>'
        '</head>'
        '<body>'
        '<nav class="hp-nav" id="hpNav">'
        '<a href="/" class="hp-nav-brand">'
        '<svg viewBox="0 0 28 28" fill="none" width="28" height="28"><circle cx="14" cy="14" r="13" stroke="currentColor" stroke-width="2"/><path d="M9 14h10M14 9v10" stroke="currentColor" stroke-width="2" stroke-linecap="round"/></svg>'
        'Scriptly</a>'
        '<div class="hp-nav-links">'
        '<a href="/login" class="nav-login">← כניסה</a>'
        '<a href="/register" class="btn-cta">התחל בחינם</a>'
        '</div>'
        '</nav>'
        '<section class="hp-hero">'
        '<div class="hp-hero-bg">'
        '<div class="hp-hero-grid"></div>'
        '<div class="hp-hero-orb orb1"></div>'
        '<div class="hp-hero-orb orb2"></div>'
        '<div class="hp-hero-orb orb3"></div>'
        '</div>'
        '<div class="hp-hero-content">'
        '<div class="hp-hero-badge">'
        '<span class="hp-hero-badge-dot"></span>'
        'פלטפורמת AI חדשה לשוק ה-HR הישראלי'
        '</div>'
        '<h1>בנו כלים חכמים<br>עם <span>בינה מלאכותית</span></h1>'
        '<p class="hp-hero-sub">תארו מה אתם צריכים בעברית פשוטה &mdash; וה-AI של Scriptly בונה את הכלי בשבילכם.<br>ניהול נוכחות, שכר ודוחות מעולם לא היה כל כך פשוט.</p>'
        '<div class="hp-hero-btns">'
        '<a href="/register" class="hp-btn-primary">'
        '<svg viewBox="0 0 20 20" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round"><path d="M15 10H5M9 6l-4 4 4 4"/></svg>'
        'התחילו בחינם'
        '</a>'
        '<a href="#features" class="hp-btn-secondary">גלו את היכולות</a>'
        '</div>'
        '<div class="hp-hero-stats">'
        '<div class="hp-hero-stat"><span class="hp-hero-stat-num">30</span><span class="hp-hero-stat-label">ימי ניסיון חינם</span></div>'
        '<div class="hp-hero-stat"><span class="hp-hero-stat-num">7+</span><span class="hp-hero-stat-label">כלים מובנים</span></div>'
        '<div class="hp-hero-stat"><span class="hp-hero-stat-num">AI</span><span class="hp-hero-stat-label">בונה כלים בשבילכם</span></div>'
        '</div>'
        '</div>'
        '</section>'
        '<section class="hp-section" id="features">'
        '<div class="hp-section-inner hp-section-center">'
        '<span class="hp-section-label">&#10024; יכולות הפלטפורמה</span>'
        '<h2 class="hp-section-title">הכל במקום אחד &mdash; חכם ומהיר</h2>'
        '<p class="hp-section-sub">Scriptly משלבת בינה מלאכותית עם כלים מוכנים לשימוש, כדי שתוכלו להתמקד במה שבאמת חשוב.</p>'
        '<div class="hp-features">'
        '<div class="hp-feature feat-ai"><div class="hp-feature-icon-wrap">&#129302;</div><h3>בונה כלים עם AI</h3><p>תארו בעברית את הכלי שאתם צריכים &mdash; ה-AI מייצר אותו עבורכם. ללא קוד, ללא מתכנת.</p></div>'
        '<div class="hp-feature feat-tools"><div class="hp-feature-icon-wrap">&#128736;</div><h3>כלים מוכנים לשימוש</h3><p>ניקוי נוכחות, סיכום שכר, דוח חוסר, עובדים לא פעילים, מבנה ארגוני &mdash; מוכן להפעלה.</p></div>'
        '<div class="hp-feature feat-market"><div class="hp-feature-icon-wrap">&#127968;</div><h3>שוק כלים קהילתי</h3><p>גלו כלים שנוצרו על ידי משתמשים אחרים. התקינו, דרגו, והשתמשו בכלים מוכחים.</p></div>'
        '<div class="hp-feature feat-file"><div class="hp-feature-icon-wrap">&#128196;</div><h3>עיבוד קבצים מיידי</h3><p>העלו קובץ Excel או CSV, קבלו דוח מעובד בחזרה. תוצאות מדויקות תוך שניות.</p></div>'
        '<div class="hp-feature feat-il"><div class="hp-feature-icon-wrap">&#127470;&#127473;</div><h3>בנוי ל-HR ישראלי</h3><p>עברית מלאה, הבנה של מערכות נוכחות ושכר ישראליות, מושגים מדיני עבודה מקומיים.</p></div>'
        '<div class="hp-feature feat-free"><div class="hp-feature-icon-wrap">&#127873;</div><h3>30 יום חינם</h3><p>התנסו בכל היכולות &mdash; ללא כרטיס אשראי, ללא התחייבות, ללא הגבלה.</p></div>'
        '</div></div></section>'
        '<section class="hp-section hp-showcase"><div class="hp-section-inner"><div class="hp-showcase-grid">'
        '<div class="hp-showcase-text">'
        '<span class="hp-section-label">&#129302; בונה הכלים</span>'
        '<h3>ספרו ל-AI מה אתם צריכים &mdash;<br>והכלי מוכן</h3>'
        '<p>לא צריך לדעת לתכנת. לא צריך לחכות לפיתוח. פשוט כתבו בעברית מה הכלי צריך לעשות, וה-AI של Scriptly בונה אותו בזמן אמת.</p>'
        '<ul class="hp-showcase-checks">'
        '<li><span class="check">&#10003;</span>תיאור בעברית פשוטה &mdash; בלי קוד</li>'
        '<li><span class="check">&#10003;</span>תוצאה מיידית &mdash; כלי מוכן לשימוש</li>'
        '<li><span class="check">&#10003;</span>אפשר לפרסם בשוק הכלים הקהילתי</li>'
        '<li><span class="check">&#10003;</span>עיבוד חכם: סינון, קיבוץ, חישובים, פיבוט ועוד</li>'
        '</ul></div>'
        '<div class="hp-showcase-visual"><div class="hp-chat-dots"><span></span><span></span><span></span></div><div class="hp-chat-msgs">'
        '<div class="hp-chat-msg user">אני צריך כלי שמסנן עובדים שעבדו פחות מ-100 שעות בחודש ומסכם לפי מחלקה</div>'
        '<div class="hp-chat-msg ai">הבנתי! בניתי כלי שמבצע:<br>1. סינון לפי סף שעות<br>2. קיבוץ לפי מחלקה<br>3. סיכום שעות ועובדים<br><span class="tag">&#10003; הכלי מוכן להרצה</span></div>'
        '<div class="hp-chat-msg user">מעולה, אפשר להוסיף גם ממוצע שעות?</div>'
        '</div></div></div></div></section>'
        '<section class="hp-section hp-how"><div class="hp-section-inner hp-section-center">'
        '<span class="hp-section-label">&#128640; איך זה עובד</span>'
        '<h2 class="hp-section-title">שלושה צעדים &mdash; מקובץ גולמי לדוח מוכן</h2>'
        '<p class="hp-section-sub">תהליך פשוט ומהיר &mdash; בלי הגדרות מסובכות</p>'
        '<div class="hp-steps">'
        '<div class="hp-step"><div class="hp-step-num">1</div><h3>בחרו כלי או בנו חדש</h3><p>בחרו כלי מוכן מהספרייה, התקינו מהשוק הקהילתי, או תארו ל-AI מה אתם צריכים.</p></div>'
        '<div class="hp-step"><div class="hp-step-num">2</div><h3>העלו קובץ ואשרו</h3><p>העלו את קובץ ה-Excel, המערכת מזהה את השדות &mdash; אשרו בלחיצה ותנו לה לעבוד.</p></div>'
        '<div class="hp-step"><div class="hp-step-num">3</div><h3>הורידו את התוצאה</h3><p>קבלו דוח מעוצב ומוכן &mdash; Excel, PowerPoint או ZIP &mdash; מוכן להדפסה או להנהלה.</p></div>'
        '</div></div></section>'
        '<section class="hp-section hp-scripts"><div class="hp-section-inner hp-section-center">'
        '<span class="hp-section-label">&#128736; כלים מובנים</span>'
        '<h2 class="hp-section-title">כלים מוכנים שחוסכים שעות עבודה</h2>'
        '<p class="hp-section-sub">הפעילו כלים מוכנים לשימוש מיידי &mdash; בלי הגדרות, בלי המתנה</p>'
        '<div class="hp-script-cards">'
        '<div class="hp-script-card"><div class="hp-script-icon">&#128203;</div><div class="hp-script-info"><h4>ניקוי דוח נוכחות</h4><p>ניקוי ועיבוד אוטומטי של דוחות נוכחות גולמיים</p></div></div>'
        '<div class="hp-script-card"><div class="hp-script-icon">&#128176;</div><div class="hp-script-info"><h4>סיכום שכר לפי תעריף</h4><p>חישוב שכר אוטומטי לפי שעות ותעריפים</p></div></div>'
        '<div class="hp-script-card"><div class="hp-script-icon">&#128200;</div><div class="hp-script-info"><h4>דוח חוסר מול תקן</h4><p>השוואת שעות בפועל מול תקן עם סינון חכם</p></div></div>'
        '<div class="hp-script-card"><div class="hp-script-icon">&#128101;</div><div class="hp-script-info"><h4>עובדים לא פעילים</h4><p>איתור עובדים ללא נוכחות בתקופה נבחרת</p></div></div>'
        '<div class="hp-script-card"><div class="hp-script-icon">&#128221;</div><div class="hp-script-info"><h4>תיקונים ידניים</h4><p>דוח תיקונים ידניים עם סיכום מפורט</p></div></div>'
        '<div class="hp-script-card"><div class="hp-script-icon">&#127970;</div><div class="hp-script-info"><h4>מבנה ארגוני</h4><p>יצירת מבנה ארגוני ב-Excel, PowerPoint ו-ZIP</p></div></div>'
        '</div></div></section>'
        '<section class="hp-section hp-trust-section"><div class="hp-section-inner hp-section-center">'
        '<span class="hp-section-label">&#128170; למה Scriptly</span>'
        '<h2 class="hp-section-title">פלטפורמה שנבנתה עבורכם</h2>'
        '<p class="hp-section-sub">כל מה שאנשי HR ושכר בישראל צריכים &mdash; במקום אחד</p>'
        '<div class="hp-trust-grid">'
        '<div class="hp-trust-item"><div class="hp-trust-icon">&#127470;&#127473;</div><h4>עברית מלאה</h4><p>ממשק, דוחות וחוויה מלאה בעברית &mdash; כולל ה-AI</p></div>'
        '<div class="hp-trust-item"><div class="hp-trust-icon">&#128274;</div><h4>אבטחה ופרטיות</h4><p>הקבצים שלכם מאובטחים ולא נשמרים לאחר העיבוד</p></div>'
        '<div class="hp-trust-item"><div class="hp-trust-icon">&#9889;</div><h4>מהיר ומיידי</h4><p>ללא התקנה, ללא הגדרות &mdash; פשוט נכנסים ומתחילים</p></div>'
        '<div class="hp-trust-item"><div class="hp-trust-icon">&#128640;</div><h4>מתעדכן כל הזמן</h4><p>כלים חדשים ושיפורים מתווספים באופן שוטף</p></div>'
        '</div></div></section>'
        '<section class="hp-cta-section">'
        '<div class="hp-hero-orb orb1" style="width:400px;height:400px;top:-100px;right:-80px;position:absolute"></div>'
        '<div class="hp-hero-orb orb2" style="width:300px;height:300px;bottom:-80px;left:-60px;position:absolute"></div>'
        '<div class="hp-cta-inner">'
        '<h2>מוכנים להתחיל?<br>ניסיון חינם למשך 30 יום</h2>'
        '<p>הצטרפו ל-Scriptly וגלו איך AI חוסך לכם שעות עבודה כל חודש.<br>ללא כרטיס אשראי. ללא התחייבות.</p>'
        '<a href="/register" class="hp-cta-btn">התחילו בחינם &larr;</a>'
        '<div class="hp-cta-badges">'
        '<span class="hp-cta-badge"><svg viewBox="0 0 20 20" fill="currentColor"><path d="M16.707 5.293a1 1 0 010 1.414l-8 8a1 1 0 01-1.414 0l-4-4a1 1 0 111.414-1.414L8 12.586l7.293-7.293a1 1 0 011.414 0z"/></svg>ללא כרטיס אשראי</span>'
        '<span class="hp-cta-badge"><svg viewBox="0 0 20 20" fill="currentColor"><path d="M16.707 5.293a1 1 0 010 1.414l-8 8a1 1 0 01-1.414 0l-4-4a1 1 0 111.414-1.414L8 12.586l7.293-7.293a1 1 0 011.414 0z"/></svg>30 יום חינם</span>'
        '<span class="hp-cta-badge"><svg viewBox="0 0 20 20" fill="currentColor"><path d="M16.707 5.293a1 1 0 010 1.414l-8 8a1 1 0 01-1.414 0l-4-4a1 1 0 111.414-1.414L8 12.586l7.293-7.293a1 1 0 011.414 0z"/></svg>ללא התחייבות</span>'
        '</div></div></section>'
        '<footer class="hp-footer"><div class="hp-footer-inner">'
        '<div class="hp-footer-brand">'
        '<svg viewBox="0 0 28 28" fill="none" width="22" height="22"><circle cx="14" cy="14" r="13" stroke="currentColor" stroke-width="2"/><path d="M9 14h10M14 9v10" stroke="currentColor" stroke-width="2" stroke-linecap="round"/></svg>'
        'Scriptly</div>'
        '<div class="hp-footer-links">'
        '<a href="/login">כניסה</a>'
        '<a href="/register">הרשמה</a>'
        '<a href="/support">תמיכה</a>'
        '</div>'
        '<div class="hp-footer-copy">&#169; 2026 Scriptly &mdash; כלי AI חכמים לניהול נוכחות ושכר</div>'
        '</div></footer>'
        '<script>'
        '(function(){'
        'var nav=document.getElementById("hpNav");'
        'function checkScroll(){if(window.scrollY>60){nav.classList.add("scrolled")}else{nav.classList.remove("scrolled")}}'
        'window.addEventListener("scroll",checkScroll,{passive:true});'
        'checkScroll();'
        'var els=document.querySelectorAll(".hp-feature,.hp-step,.hp-script-card,.hp-trust-item");'
        'if(!("IntersectionObserver" in window)){els.forEach(function(e){e.classList.add("visible")});return;}'
        'var obs=new IntersectionObserver(function(entries){'
        'entries.forEach(function(entry){'
        'if(entry.isIntersecting){entry.target.classList.add("visible");obs.unobserve(entry.target);}'
        '});'
        '},{threshold:0.1,rootMargin:"0px 0px -40px 0px"});'
        'els.forEach(function(el,i){'
        'var siblings=el.parentElement.querySelectorAll("."+el.className.split(" ")[0]);'
        'var idx=Array.prototype.indexOf.call(siblings,el);'
        'el.style.transitionDelay=(idx*0.1)+"s";'
        'obs.observe(el);'
        '});'
        'document.querySelectorAll(' + "'" + 'a[href^="#"]' + "'" + ').forEach(function(a){'
        'a.addEventListener("click",function(e){'
        'var target=document.querySelector(this.getAttribute("href"));'
        'if(target){e.preventDefault();target.scrollIntoView({behavior:"smooth",block:"start"});}'
        '});'
        '});'
        'var chatMsgs=document.querySelectorAll(".hp-chat-msg");'
        'if(chatMsgs.length&&"IntersectionObserver" in window){'
        'var chatObs=new IntersectionObserver(function(entries){'
        'entries.forEach(function(entry){'
        'if(entry.isIntersecting){'
        'chatMsgs.forEach(function(msg){msg.style.animation="none";msg.offsetHeight;msg.style.animation="";});'
        'chatObs.unobserve(entry.target);'
        '}'
        '});'
        '},{threshold:0.3});'
        'chatObs.observe(document.querySelector(".hp-showcase-visual"));'
        '}'
        '})();'
        '</script>'
        '</body></html>'
    )


@app.route("/register", methods=["GET", "POST"])
def register():
    if "user_id" in session:
        return redirect("/admin" if session.get("is_admin") else "/dashboard")

    lang = get_flow_lang()
    error = ""

    if request.method == "POST":
        full_name = request.form.get("full_name", "").strip()
        company_name = request.form.get("company_name", "").strip()
        company_id = request.form.get("company_id", "").strip()
        email = request.form.get("email", "").strip()
        phone = request.form.get("phone", "").strip()
        username = request.form.get("username", "").strip()
        password = request.form.get("password", "").strip()

        if not full_name or not email or not username or not password:
            error = '<div class="flash-err">יש למלא את כל שדות החובה</div>' if lang == "he" else '<div class="flash-err">Please fill all required fields</div>'
        elif len(password) < 4:
            error = '<div class="flash-err">הסיסמה חייבת להכיל לפחות 4 תווים</div>' if lang == "he" else '<div class="flash-err">Password must be at least 4 characters</div>'
        else:
            try:
                join_date = date.today().isoformat()
                with get_db() as db:
                    db.execute(
                        """INSERT INTO users(
                        username, password, full_name, company_name, company_id, email, phone,
                        join_date, trial_start_date, trial_days, active, is_admin, billing_mode
                        ) VALUES (?,?,?,?,?,?,?,?,?,?,?,?,?)""",
                        (
                            username,
                            generate_password_hash(password),
                            full_name,
                            company_name,
                            company_id,
                            email,
                            phone,
                            join_date,
                            join_date,
                            30,
                            1,
                            0,
                            "monthly",
                        ),
                    )
                    db.commit()
                add_flash("נרשמת בהצלחה! כעת ניתן להתחבר" if lang == "he" else "Registration successful! You can now log in")
                return redirect("/login")
            except Exception:
                error = '<div class="flash-err">שם המשתמש כבר קיים במערכת</div>' if lang == "he" else '<div class="flash-err">Username already exists</div>'

    lbl = {
        "he": {"title": "הרשמה", "full_name": "שם מלא *", "company_name": "שם חברה", "company_id": "ח.פ / מזהה חברה", "email": "אימייל *", "phone": "טלפון", "username": "שם משתמש *", "password": "סיסמה *", "submit": "הרשמה", "login_link": "יש לך חשבון?", "login_text": "התחברות"},
        "en": {"title": "Register", "full_name": "Full Name *", "company_name": "Company Name", "company_id": "Company ID", "email": "Email *", "phone": "Phone", "username": "Username *", "password": "Password *", "submit": "Register", "login_link": "Already have an account?", "login_text": "Log in"},
    }
    t = lbl.get(lang, lbl["he"])

    svg_logo = '<svg viewBox="0 0 28 28" fill="none" width="36" height="36" style="color:#2563eb"><circle cx="14" cy="14" r="13" stroke="currentColor" stroke-width="2"/><path d="M9 14h10M14 9v10" stroke="currentColor" stroke-width="2" stroke-linecap="round"/></svg>'
    body = (
        '<div style="background:rgba(255,255,255,.95);backdrop-filter:blur(20px);-webkit-backdrop-filter:blur(20px);padding:2.5rem;border-radius:24px;box-shadow:0 20px 60px rgba(0,0,0,.2),0 0 0 1px rgba(255,255,255,.1);border:1px solid rgba(255,255,255,.3)">'
        '<div style="text-align:center;margin-bottom:2rem">'
        '<div style="display:inline-flex;align-items:center;justify-content:center;margin-bottom:8px">' + svg_logo + '</div>'
        '<h1 style="font-size:24px;font-weight:800;color:#0f172a;letter-spacing:-0.5px">Scriptly</h1>'
        '<p style="font-size:13px;color:#64748b;margin-top:6px;line-height:1.6">' + t["title"] + '</p>'
        '</div>'
        + error
        + '<form method="POST">'
        + '<label class="field-label">' + t["full_name"] + '</label>'
        '<input type="text" name="full_name" required value="' + esc(request.form.get("full_name", "")) + '" style="background:rgba(248,250,255,.8);border-color:#dbeafe">'
        + '<label class="field-label">' + t["company_name"] + '</label>'
        '<input type="text" name="company_name" value="' + esc(request.form.get("company_name", "")) + '" style="background:rgba(248,250,255,.8);border-color:#dbeafe">'
        + '<label class="field-label">' + t["company_id"] + '</label>'
        '<input type="text" name="company_id" value="' + esc(request.form.get("company_id", "")) + '" style="background:rgba(248,250,255,.8);border-color:#dbeafe">'
        + '<label class="field-label">' + t["email"] + '</label>'
        '<input type="text" name="email" required value="' + esc(request.form.get("email", "")) + '" style="background:rgba(248,250,255,.8);border-color:#dbeafe">'
        + '<label class="field-label">' + t["phone"] + '</label>'
        '<input type="text" name="phone" value="' + esc(request.form.get("phone", "")) + '" style="background:rgba(248,250,255,.8);border-color:#dbeafe">'
        + '<label class="field-label">' + t["username"] + '</label>'
        '<input type="text" name="username" required value="' + esc(request.form.get("username", "")) + '" style="background:rgba(248,250,255,.8);border-color:#dbeafe">'
        + '<label class="field-label">' + t["password"] + '</label>'
        '<input type="password" name="password" required style="background:rgba(248,250,255,.8);border-color:#dbeafe">'
        + '<button type="submit" class="btn btn-blue" style="width:100%;padding:14px;font-size:15px;margin-top:.75rem;border-radius:14px">' + t["submit"] + '</button>'
        '</form>'
        '<p style="text-align:center;margin-top:1.25rem;font-size:13px;color:#64748b">'
        + t["login_link"] + ' '
        '<a href="/login" style="color:#2563eb;font-weight:700;text-decoration:none">' + t["login_text"] + '</a>'
        '</p>'
        '<p style="text-align:center;margin-top:1rem"><a href="/" style="font-size:13px;color:#64748b;text-decoration:none;font-weight:600">'
        + ("&#8592; חזרה לדף הבית" if lang == "he" else "&#8592; Back to homepage")
        + '</a></p>'
        '<p style="text-align:center;margin-top:.75rem;font-size:11px;color:#94a3b8">&#169; Scriptly</p>'
        '</div>'
    )
    return render(t["title"], body, nav=False, lang=lang)


@app.route("/login", methods=["GET", "POST"])
def login():
    if "user_id" in session:
        return redirect("/admin" if session.get("is_admin") else "/dashboard")

    lang = get_flow_lang()
    text = get_flow_text(lang)
    error = ""
    if request.method == "POST":
        username = request.form["username"].strip()
        password = request.form["password"]
        with get_db() as db:
            user = db.execute("SELECT * FROM users WHERE username=? AND active=1", (username,)).fetchone()
        if user and check_password_hash(user["password"], password):
            session.update(
                {
                    "user_id": user["id"],
                    "username": user["username"],
                    "name": user["full_name"],
                    "is_admin": bool(user["is_admin"]),
                }
            )
            return redirect("/admin" if user["is_admin"] else "/dashboard")
        error = '<div class="flash-err">' + text["login_error"] + "</div>"

    svg_logo = '<svg viewBox="0 0 28 28" fill="none" width="36" height="36" style="color:#2563eb"><circle cx="14" cy="14" r="13" stroke="currentColor" stroke-width="2"/><path d="M9 14h10M14 9v10" stroke="currentColor" stroke-width="2" stroke-linecap="round"/></svg>'
    body = (
        '<div style="text-align:center;margin-bottom:1.5rem">'
        + build_lang_switch(lang).replace('class="lang-switch"', 'class="lang-switch standalone"').replace('border-color: #dbeafe', 'border-color: rgba(255,255,255,.2)').replace('background: #ffffff', 'background: rgba(255,255,255,.12)').replace('color: #64748b', 'color: rgba(255,255,255,.6)').replace('background: #2563eb; color: #ffffff', 'background: rgba(255,255,255,.2); color: #ffffff')
        + '</div>'
        + '<div style="background:rgba(255,255,255,.95);backdrop-filter:blur(20px);-webkit-backdrop-filter:blur(20px);padding:2.5rem;border-radius:24px;box-shadow:0 20px 60px rgba(0,0,0,.2),0 0 0 1px rgba(255,255,255,.1);border:1px solid rgba(255,255,255,.3)">'
        '<div style="text-align:center;margin-bottom:2rem">'
        '<div style="display:inline-flex;align-items:center;justify-content:center;margin-bottom:8px">' + svg_logo + '</div>'
        '<h1 style="font-size:24px;font-weight:800;color:#0f172a;letter-spacing:-0.5px">Scriptly</h1>'
        + '<p style="font-size:13px;color:#64748b;margin-top:6px;line-height:1.6">' + text["login_subtitle"] + '</p>'
        "</div>"
        + error
        + '<form method="POST">'
        + '<label class="field-label">' + text["login_username"] + '</label>'
        '<input type="text" name="username" required autofocus style="background:rgba(248,250,255,.8);border-color:#dbeafe">'
        + '<label class="field-label">' + text["login_password"] + '</label>'
        '<input type="password" name="password" required style="background:rgba(248,250,255,.8);border-color:#dbeafe">'
        + '<button type="submit" class="btn btn-blue" style="width:100%;padding:14px;font-size:15px;margin-top:.75rem;border-radius:14px">' + text["login_submit"] + '</button>'
        "</form>"
        '<p style="text-align:center;margin-top:1.25rem;font-size:13px;color:#64748b">'
        + ("אין לך חשבון? " if lang == "he" else "Don't have an account? ")
        + '<a href="/register" style="color:#2563eb;font-weight:700;text-decoration:none">'
        + ("הרשמה בחינם" if lang == "he" else "Register free")
        + "</a></p>"
        '<p style="text-align:center;margin-top:1rem"><a href="/" style="font-size:13px;color:#64748b;text-decoration:none;font-weight:600">'
        + ("&#8592; חזרה לדף הבית" if lang == "he" else "&#8592; Back to homepage")
        + '</a></p>'
        '<p style="text-align:center;margin-top:.75rem;font-size:11px;color:#94a3b8">&#169; Scriptly</p>'
        "</div>"
    )
    return render(text["login_page_title"], body, nav=False, lang=lang)


@app.route("/logout")
def logout():
    session.clear()
    return redirect("/login")


@app.route("/dashboard")
@login_required
def dashboard():
    if session.get("is_admin"):
        return redirect("/admin")

    lang = get_flow_lang()
    text = get_flow_text(lang)
    expire_report_jobs()
    with get_db() as db:
        user = db.execute("SELECT * FROM users WHERE id=?", (session["user_id"],)).fetchone()
        perms = db.execute("SELECT script_id FROM permissions WHERE user_id=?", (session["user_id"],)).fetchall()
        report_jobs = db.execute("SELECT * FROM report_jobs WHERE user_id=? AND status<>'downloaded' ORDER BY created_at DESC, id DESC LIMIT 8", (session["user_id"],)).fetchall()
        installed_tools = db.execute(
            """SELECT t.id, t.name, t.description, t.icon FROM marketplace_tools t
            JOIN tool_installs i ON i.tool_id = t.id
            WHERE i.user_id=? AND t.status='approved' AND t.is_public=1
            ORDER BY t.name""",
            (session["user_id"],),
        ).fetchall()
    if user is None:
        session.clear()
        return redirect("/login")

    allowed = [get_localized_script(SCRIPTS[p["script_id"]], lang) for p in perms if p["script_id"] in SCRIPTS]
    status = get_account_status(user)
    status_label = status["status_label_he"] if lang == "he" else status["status_label_en"]
    not_set = "לא הוגדר" if lang == "he" else "Not set"
    status_colors = {
        "trial": ("#fff7ed", "#c2410c"),
        "active": ("#ecfdf5", "#047857"),
        "expired": ("#fef2f2", "#b91c1c"),
        "unknown": ("#f8fafc", "#475569"),
    }
    status_bg, status_fg = status_colors.get(status["status_key"], ("#f8fafc", "#475569"))
    if status["status_key"] == "trial":
        service_note = f"נותרו {status['days_remaining']} ימי ניסיון" if lang == "he" else f"{status['days_remaining']} trial days remaining"
    elif status["renewal_date"]:
        service_note = ("בתוקף עד " if lang == "he" else "Valid until ") + format_ui_date(status["renewal_date"], lang)
    else:
        service_note = "ללא פרטי שירות נוספים" if lang == "he" else "No additional service details"
    collapsed_company_name = user["company_name"] or user["full_name"] or user["username"] or not_set
    service_compact = status_label
    if status["status_key"] == "trial" and status["days_remaining"] is not None:
        service_compact += (" • נותרו " + str(status["days_remaining"]) + " ימים" if lang == "he" else f" • {status['days_remaining']} days left")

    cards = ""
    for script in allowed:
        cards += (
            '<a href="/run/' + script["id"] + '" style="background:white;border:1px solid #e2e8f0;border-radius:18px;box-shadow:0 2px 12px rgba(37,99,235,.06);padding:1.5rem;text-decoration:none;display:block;transition:all .3s cubic-bezier(.4,0,.2,1)"'
            ' onmouseenter="this.style.transform=\'translateY(-3px)\';this.style.boxShadow=\'0 8px 28px rgba(37,99,235,.12)\';this.style.borderColor=\'#bfdbfe\'"'
            ' onmouseleave="this.style.transform=\'none\';this.style.boxShadow=\'0 2px 12px rgba(37,99,235,.06)\';this.style.borderColor=\'#e2e8f0\'">'
            '<div style="width:48px;height:48px;border-radius:14px;background:linear-gradient(135deg,#eff6ff,#dbeafe);display:flex;align-items:center;justify-content:center;font-size:24px;margin-bottom:.75rem">' + script["icon"] + "</div>"
            '<div style="font-size:15px;font-weight:800;color:#0f172a;margin-bottom:4px">' + script["name"] + "</div>"
            '<div style="font-size:12px;color:#64748b;line-height:1.6">' + script["desc"] + "</div>"
            "</a>"
        )
    # Add installed marketplace tools
    for mt in installed_tools:
        cards += (
            '<a href="/marketplace/tool/' + str(mt["id"]) + '/run" style="background:linear-gradient(135deg,#ffffff 0%,#f0fdf4 100%);border:1px solid #86efac;border-radius:18px;box-shadow:0 2px 12px rgba(22,163,74,.06);padding:1.5rem;text-decoration:none;display:block;transition:all .3s cubic-bezier(.4,0,.2,1)"'
            ' onmouseenter="this.style.transform=\'translateY(-3px)\';this.style.boxShadow=\'0 8px 28px rgba(22,163,74,.12)\'"'
            ' onmouseleave="this.style.transform=\'none\';this.style.boxShadow=\'0 2px 12px rgba(22,163,74,.06)\'">'
            '<div style="width:48px;height:48px;border-radius:14px;background:linear-gradient(135deg,#ecfdf5,#d1fae5);display:flex;align-items:center;justify-content:center;font-size:24px;margin-bottom:.75rem">' + esc(mt["icon"] or "🔧") + "</div>"
            '<div style="font-size:15px;font-weight:800;color:#047857;margin-bottom:4px">' + esc(mt["name"]) + '</div>'
            '<div style="font-size:12px;color:#64748b;line-height:1.6">' + esc((mt["description"] or "")[:80]) + '</div>'
            '<div style="font-size:10px;color:#94a3b8;margin-top:6px">🛒 מתוך שוק הכלים</div>'
            "</a>"
        )

    if not allowed and not installed_tools:
        cards = (
            '<div style="text-align:center;padding:3rem;color:#94a3b8">'
            '<div style="font-size:48px;margin-bottom:1rem">&#128274;</div>'
            "<div>" + text["dashboard_empty"] + "</div>"
            "</div>"
        )

    info_items = [
        ("שם חברה" if lang == "he" else "Company", user["company_name"] or user["full_name"] or user["username"]),
        ("שם משתמש" if lang == "he" else "Username", user["username"]),
        ("ח.פ / מזהה חברה" if lang == "he" else "Company ID", user["company_id"]),
        ("תאריך הצטרפות" if lang == "he" else "Join date", format_ui_date(user["join_date"], lang)),
    ]
    info_grid = "".join(
        '<div style="background:linear-gradient(180deg,#f8fafc 0%,#f1f5f9 100%);border:1px solid #e2e8f0;border-radius:14px;padding:14px;transition:border-color .2s">'
        '<div style="font-size:11px;color:#64748b;margin-bottom:6px;text-transform:uppercase;letter-spacing:.5px;font-weight:600">' + esc(label) + '</div>'
        '<div style="font-size:15px;font-weight:700;color:#0f172a">' + esc(value or not_set) + "</div>"
        "</div>"
        for label, value in info_items
    )
    status_styles = {
        "pending": ("#eff6ff", "#1d4ed8"),
        "processing": ("#fff7ed", "#c2410c"),
        "ready": ("#ecfdf5", "#047857"),
        "failed": ("#fef2f2", "#b91c1c"),
        "expired": ("#f8fafc", "#475569"),
        "downloaded": ("#f8fafc", "#475569"),
    }
    report_job_rows = ""
    for job in report_jobs:
        pill_bg, pill_fg = status_styles.get(job["status"], ("#f8fafc", "#475569"))
        action_html = (
            '<a href="/report-jobs/' + str(job["id"]) + '/download" class="btn btn-blue" style="text-decoration:none;display:inline-flex;align-items:center;justify-content:center;min-width:140px">הורדת הדוח</a>'
            if job["status"] == "ready"
            else ""
        )
        report_job_rows += (
            '<div style="border:1px solid #e2e8f0;border-radius:16px;padding:16px 18px;background:#ffffff;transition:border-color .2s,box-shadow .2s" onmouseenter="this.style.borderColor=\'#bfdbfe\';this.style.boxShadow=\'0 4px 16px rgba(37,99,235,.06)\'" onmouseleave="this.style.borderColor=\'#e2e8f0\';this.style.boxShadow=\'none\'">'
            '<div style="display:flex;align-items:flex-start;justify-content:space-between;gap:14px;flex-wrap:wrap">'
            '<div>'
            '<div style="font-size:15px;font-weight:800;color:#0f172a;margin-bottom:6px">' + esc(job["script_name"]) + '</div>'
            '<div style="font-size:13px;color:#64748b;margin-bottom:6px">' + esc(job["original_filename"]) + '</div>'
            '<div style="display:inline-flex;align-items:center;padding:6px 10px;border-radius:999px;background:' + pill_bg + ';color:' + pill_fg + ';font-size:12px;font-weight:800">' + esc(report_job_status_label(job["status"])) + '</div>'
            + ('<div style="font-size:12px;color:#475569;margin-top:8px;line-height:1.7">' + esc(job["status_note"]) + '</div>' if job["status_note"] else '')
            + '</div>'
            '<div style="display:flex;flex-direction:column;align-items:flex-end;gap:8px">'
            '<div style="font-size:12px;color:#64748b">' + esc(format_ui_datetime(job["created_at"])) + '</div>'
            + action_html
            + '</div>'
            '</div></div>'
        )
    report_jobs_html = (
        '<div class="card" style="margin-top:1rem;background:linear-gradient(135deg,#ffffff 0%,#f8fbff 100%);border:1px solid #dbeafe;border-radius:18px">'
        '<div style="display:flex;align-items:center;justify-content:space-between;gap:14px;flex-wrap:wrap;margin-bottom:14px">'
        '<div><div style="font-size:18px;font-weight:800;color:#1e3a8a;margin-bottom:6px">📋 דוחות בעיבוד והורדה</div>'
        '<div style="font-size:14px;color:#475569;line-height:1.7">דוחות שנשלחו לעיבוד ברקע יופיעו כאן. הדוח נשמר עד להורדה או עד 3 ימים, המוקדם מביניהם.</div></div>'
        '</div>'
        + (report_job_rows if report_job_rows else '<div style="text-align:center;padding:1rem 0;color:#94a3b8">עדיין אין דוחות בעיבוד או דוחות מוכנים להורדה</div>')
        + '</div>'
    )
    body = (
        '<h2 style="font-size:clamp(22px,3.5vw,28px);font-weight:800;color:#0f172a;margin-bottom:.4rem;letter-spacing:-0.3px">' + text["dashboard_greeting"]
        + esc(session["name"])
        + ' &#128075;</h2>'
        + ('<p style="font-size:15px;color:#64748b;margin-bottom:1.5rem;line-height:1.7">' + text["dashboard_intro"] + "</p>" if text["dashboard_intro"] else "")
        + '<div style="display:grid;grid-template-columns:1.15fr .85fr;gap:1rem;margin-bottom:1rem">'
        + '<details class="card" style="margin:0;padding:0;overflow:hidden;border:1px solid #e2e8f0" id="accountDetailsCard">'
        + '<summary style="list-style:none;cursor:pointer;padding:20px 22px;display:flex;align-items:center;justify-content:space-between;gap:12px;transition:background .2s"'
        + ' onmouseenter="this.style.background=\'#f8fafc\'" onmouseleave="this.style.background=\'transparent\'">'
        + '<div><div style="font-size:17px;font-weight:800;color:#0f172a;margin-bottom:6px">' + ("פרטי חשבון ולקוח" if lang == "he" else "Account and company details") + '</div>'
        + '<div style="font-size:14px;color:#475569;font-weight:600">' + esc(collapsed_company_name) + '</div></div>'
        + '<span style="font-size:16px;color:#94a3b8;transition:transform .2s">▼</span>'
        + '</summary>'
        + '<div style="padding:0 20px 20px"><div style="display:grid;grid-template-columns:repeat(auto-fit,minmax(180px,1fr));gap:12px">'
        + info_grid
        + '</div></div></details>'
        + '<details class="card" style="margin:0;padding:0;overflow:hidden;background:linear-gradient(180deg,#ffffff 0%,#f8fbff 100%);border:1px solid #dbeafe" id="serviceStatusCard">'
        + '<summary style="list-style:none;cursor:pointer;padding:20px 22px;display:flex;align-items:center;justify-content:space-between;gap:12px;transition:background .2s"'
        + ' onmouseenter="this.style.background=\'#f8fafc\'" onmouseleave="this.style.background=\'transparent\'">'
        + '<div><div style="font-size:17px;font-weight:800;color:#0f172a;margin-bottom:6px">' + ("סטטוס שירות" if lang == "he" else "Service status") + '</div>'
        + '<div style="display:inline-flex;align-items:center;padding:8px 12px;border-radius:999px;background:' + status_bg + ';color:' + status_fg + ';font-size:14px;font-weight:800;margin-bottom:8px">' + esc(status_label) + '</div>'
        + '<div style="font-size:13px;color:#475569">' + esc(service_compact) + '</div></div>'
        + '<span style="font-size:16px;color:#94a3b8;transition:transform .2s">▼</span>'
        + '</summary>'
        + '<div style="padding:0 22px 22px">'
        + '<div style="font-size:14px;color:#334155;margin-bottom:10px;line-height:1.7">' + esc(service_note) + '</div>'
        + '<div style="font-size:12px;color:#64748b;margin-bottom:4px">' + ("מסלול חיוב" if lang == "he" else "Billing mode") + '</div>'
        + '<div style="font-size:15px;font-weight:700;color:#0f172a">' + esc(billing_mode_label(user["billing_mode"], lang)) + '</div>'
        + '</div></details>'
        + '</div>'
        + '<div class="card" style="margin:0;border:1px solid #e2e8f0"><div style="font-size:18px;font-weight:800;color:#0f172a;margin-bottom:4px">'
        + ("הכלים הזמינים לך" if lang == "he" else "Your available tools")
        + '</div><div style="font-size:13px;color:#64748b;margin-bottom:16px">'
        + ("לחץ על כלי כדי להתחיל לעבוד" if lang == "he" else "Click a tool to get started")
        + '</div><div style="display:grid;grid-template-columns:repeat(auto-fill,minmax(210px,1fr));gap:1rem">'
        + cards
        + '</div></div>'
        + '<div class="card" style="margin-top:1rem;background:linear-gradient(135deg,#f0fdf4 0%,#ecfdf5 100%);border:1px solid #86efac;border-radius:18px">'
        + '<div style="display:flex;align-items:center;justify-content:space-between;gap:16px;flex-wrap:wrap">'
        + '<div><div style="font-size:18px;font-weight:800;color:#047857;margin-bottom:6px">' + ("🛒 שוק הכלים" if lang == "he" else "🛒 Tool Marketplace") + '</div>'
        + '<div style="font-size:14px;color:#475569;line-height:1.7">' + ("גלה כלים שנבנו על ידי משתמשים אחרים או צור כלי משלך עם AI" if lang == "he" else "Discover tools built by other users or create your own with AI") + '</div></div>'
        + '<div style="display:flex;gap:10px;flex-wrap:wrap">'
        + '<a href="/tools/my-tools" class="btn btn-gray" style="text-decoration:none;min-width:120px">' + ("&#128188; הכלים שלי" if lang == "he" else "&#128188; My tools") + '</a>'
        + '<a href="/marketplace" class="btn" style="text-decoration:none;min-width:120px;background:linear-gradient(135deg,#047857,#059669);color:white;box-shadow:0 2px 8px rgba(4,120,87,.25)">' + ("גלה כלים" if lang == "he" else "Browse tools") + '</a>'
        + '<a href="/tools/create" class="btn btn-blue" style="text-decoration:none;min-width:120px">' + ("&#129302; צור כלי" if lang == "he" else "&#129302; Create tool") + '</a>'
        + '</div></div></div>'
        + report_jobs_html
        + '<div class="card" style="margin-top:1rem;background:linear-gradient(135deg,#ffffff 0%,#f8fbff 100%);border:1px solid #dbeafe;border-radius:18px">'
        + '<div style="display:flex;align-items:center;justify-content:space-between;gap:16px;flex-wrap:wrap">'
        + '<div><div style="font-size:18px;font-weight:800;color:#1e3a8a;margin-bottom:6px">' + ("💬 שירות לקוחות" if lang == "he" else "💬 Customer support") + '</div>'
        + '<div style="font-size:14px;color:#475569;line-height:1.7">' + ("לדיווח על תקלה או לכל שאלה" if lang == "he" else "Report an issue or ask a question") + '</div></div>'
        + '<a href="/support" class="btn btn-blue" style="text-decoration:none;min-width:180px">' + ("פתיחת פנייה" if lang == "he" else "Open request") + '</a>'
        + '</div></div>'
        + '<div class="card" style="margin-top:1rem;background:linear-gradient(135deg,#eff6ff 0%,#f8fafc 100%);border:1px solid #bfdbfe;border-radius:18px">'
        + '<div style="font-size:18px;font-weight:800;color:#1e3a8a;margin-bottom:8px">' + ("🔒 אבטחת מידע ופרטיות" if lang == "he" else "🔒 Security and privacy") + '</div>'
        + '<div style="font-size:14px;line-height:1.8;color:#334155">'
        + (
            "הדוחות שאתם מעלים משמשים לעיבוד בלבד ואינם נשמרים כחלק ממאגר קבוע. ברוב הכלים קבצי העבודה ותוצרי העיבוד נמחקים אוטומטית מיד לאחר ההורדה. בכלי ניקוי דוח נוכחות, אם הדוח נשלח לעיבוד ברקע, הוא נשמר עד להורדה או עד 3 ימים, המוקדם מביניהם, ולאחר מכן נמחק אוטומטית."
            if lang == "he"
            else "Uploaded reports are used only for processing and are not kept as part of a permanent data store. In most tools, working files and outputs are deleted immediately after download. For the attendance cleanup tool, when a report is sent to background processing it is kept until download or for up to 3 days, whichever comes first."
        )
        + '</div></div>'
        + '<details id="serviceTerms" style="margin-top:1rem;background:#ffffff;border:1px solid #e2e8f0;border-radius:18px;box-shadow:0 2px 16px rgba(0,0,0,.04);overflow:hidden">'
        + '<summary style="list-style:none;cursor:pointer;padding:18px 22px;font-size:15px;font-weight:800;color:#0f172a;display:flex;align-items:center;justify-content:space-between;background:linear-gradient(180deg,#ffffff 0%,#f8fafc 100%);direction:rtl;text-align:right;transition:background .2s"'
        + ' onmouseenter="this.style.background=\'#f8fafc\'" onmouseleave="this.style.background=\'linear-gradient(180deg,#ffffff 0%,#f8fafc 100%)\'">'
        + '<span>ℹ️ אודות השירות</span><span style="font-size:16px;color:#94a3b8">▼</span></summary>'
        + '<div style="padding:0 22px 22px;font-size:14px;line-height:1.9;color:#334155;direction:rtl;text-align:right">'
        + 'הפלטפורמה מרכזת כלים ודוחות שפותחו מתוך צרכים אמיתיים שעלו מהשטח.<br>'
        + 'לקוחות מנויים נהנים מגישה לכלל הכלים הזמינים במערכת.'
        + '</div></details>'
        + '<script>'
        + 'function trackUserActivity(eventType, actionLabel, scriptId, scriptName, details){try{var data=new FormData();data.append("event_type",eventType||"");data.append("action_label",actionLabel||"");data.append("script_id",scriptId||"");data.append("script_name",scriptName||"");data.append("details",details||"");if(navigator.sendBeacon){navigator.sendBeacon("/activity",data);}else{fetch("/activity",{method:"POST",body:data,credentials:"same-origin",keepalive:true});}}catch(e){}}'
        + 'var serviceTerms=document.getElementById("serviceTerms");if(serviceTerms){serviceTerms.addEventListener("toggle",function(){if(this.open){trackUserActivity("open_service_terms","פתח אודות השירות","","","אודות השירות");}});}'
        + '</script>'
    )
    return render(
        text["dashboard_page_title"],
        body,
        lang=lang,
        topbar_greeting=text["topbar_greeting"],
        logout_label=text["logout"],
        show_lang_switch=True,
    )


@app.route("/support", methods=["GET", "POST"])
@login_required
def support():
    if session.get("is_admin"):
        return redirect("/admin")

    lang = get_flow_lang()
    text = get_flow_text(lang)
    with get_db() as db:
        user = db.execute("SELECT * FROM users WHERE id=?", (session["user_id"],)).fetchone()
        perms = db.execute("SELECT script_id FROM permissions WHERE user_id=?", (session["user_id"],)).fetchall()
    if user is None:
        session.clear()
        return redirect("/login")

    status = get_account_status(user)
    allowed_scripts = [get_localized_script(SCRIPTS[p["script_id"]], lang) for p in perms if p["script_id"] in SCRIPTS]
    support_type = request.form.get("support_type", "new_tool").strip() or "new_tool"
    selected_script_id = request.form.get("support_script_id", "").strip()
    message_text = request.form.get("support_message", "").strip()
    error = ""
    success = ""

    if request.method == "POST":
        if support_type not in {"new_tool", "existing_tool"}:
            support_type = "new_tool"
        if not message_text:
            error = '<div class="flash-err">יש לכתוב את תוכן הפנייה לפני השליחה.</div>'
        elif support_type == "existing_tool" and status["status_key"] != "active":
            error = '<div class="flash-err">רק לקוחות בשירות פעיל מקבלים תמיכה על כלים קיימים. יש לפנות למנהל הפלטפורמה להסדרת השירות.</div>'
        else:
            selected_script = next((script for script in allowed_scripts if script["id"] == selected_script_id), None)
            if support_type == "existing_tool" and selected_script is None:
                error = '<div class="flash-err">יש לבחור כלי קיים מהרשימה לפני שליחת הפנייה.</div>'
            else:
                create_support_request(
                    user,
                    support_type,
                    message_text,
                    selected_script["id"] if selected_script else "",
                    selected_script["name"] if selected_script else "",
                )
                log_user_activity(
                    "open_support_request",
                    "שלח פנייה לשירות לקוחות",
                    selected_script["id"] if selected_script else "",
                    selected_script["name"] if selected_script else "",
                    "בקשה לכלי חדש" if support_type == "new_tool" else "תמיכה בכלי קיים",
                )
                if support_type == "new_tool":
                    success = '<div class="flash">תודה על פנייתך. נתפנה לעיון בבקשה בקרוב וניצור קשר דרך הנייד או האימייל הכתובים לנו במערכת.</div>'
                else:
                    success = '<div class="flash">הפנייה נקלטה בהצלחה. נתפנה לטפל בה בהקדם וניצור קשר לפי הפרטים שכתובים לנו.</div>'
                support_type = "new_tool"
                selected_script_id = ""
                message_text = ""

    with get_db() as db:
        existing_requests = db.execute(
            "SELECT * FROM support_requests WHERE user_id=? ORDER BY created_at DESC, id DESC",
            (session["user_id"],),
        ).fetchall()

    script_options = '<option value="">בחירת כלי</option>'
    for script in allowed_scripts:
        script_options += '<option value="' + esc(script["id"]) + '"' + (' selected' if selected_script_id == script["id"] else '') + '>' + esc(script["name"]) + '</option>'

    inactive_support_note = ""
    if status["status_key"] != "active":
        inactive_support_note = '<div id="existingSupportWarning" style="background:#fef2f2;border:1px solid #fecaca;color:#b91c1c;border-radius:12px;padding:12px 14px;font-size:13px;line-height:1.7;margin-bottom:12px">רק לקוחות בשירות פעיל מקבלים תמיכה על כלים קיימים. יש לפנות למנהל הפלטפורמה להסדרת השירות.</div>'

    request_rows = ""
    for entry in existing_requests:
        meta = support_status_meta(entry["status"])
        request_type_label = "בקשה לכלי חדש" if entry["request_type"] == "new_tool" else "תמיכה בכלי קיים"
        request_rows += (
            '<details style="border:1px solid #e2e8f0;border-radius:16px;padding:16px;background:linear-gradient(180deg,#ffffff 0%,#f8fafc 100%);margin-bottom:10px;transition:border-color .2s">'
            '<summary style="list-style:none;cursor:pointer;display:flex;align-items:flex-start;justify-content:space-between;gap:12px;flex-wrap:wrap">'
            '<div>'
            '<div style="font-size:14px;font-weight:800;color:#0f172a">' + esc(request_type_label) + '</div>'
            '<div style="font-size:12px;color:#64748b;margin-top:4px">' + esc(format_ui_datetime(entry["created_at"])) + (' • ' + esc(entry["script_name"]) if entry["script_name"] else '') + '</div>'
            '</div>'
            '<div style="display:flex;align-items:center;gap:10px;flex-wrap:wrap"><span style="display:inline-flex;align-items:center;padding:7px 12px;border-radius:999px;background:' + meta["bg"] + ';color:' + meta["fg"] + ';font-size:12px;font-weight:800">' + esc(meta["label"]) + '</span><span style="font-size:16px;color:#94a3b8">▼</span></div>'
            '</summary>'
            '<div style="font-size:13px;color:#334155;line-height:1.8;white-space:pre-wrap;margin-top:14px;padding-top:12px;border-top:1px solid #e2e8f0">' + esc(entry["message"] or "") + '</div>'
            '</details>'
        )
    requests_html = (
        '<div style="margin-top:1.25rem;padding-top:1.25rem;border-top:1px solid #e2e8f0">'
        '<div style="font-size:16px;font-weight:800;color:#1e3a8a;margin-bottom:10px">הפניות שלך</div>'
        + request_rows
        + '</div>'
    ) if existing_requests else ""

    body = (
        '<a href="/dashboard" style="color:#2563eb;font-size:13px;font-weight:600;text-decoration:none;display:inline-flex;align-items:center;gap:4px;margin-bottom:1rem">' + text["back_arrow"] + ' ' + ("חזרה לכלים" if lang == "he" else "Back to tools") + '</a>'
        + '<div class="card" style="border-radius:18px">'
        + '<div style="font-size:20px;font-weight:800;color:#0f172a;margin-bottom:8px">💬 שירות לקוחות</div>'
        + '<div style="font-size:14px;color:#64748b;line-height:1.7;margin-bottom:1.25rem">אפשר לפתוח פנייה לבקשה לכלי חדש או לקבלת תמיכה בכלי קיים.</div>'
        + success
        + error
        + '<form method="POST" id="supportRequestForm">'
        + '<div style="display:grid;grid-template-columns:repeat(auto-fit,minmax(220px,1fr));gap:12px;margin-bottom:14px">'
        + '<label style="display:flex;align-items:flex-start;gap:10px;background:#f8fafc;border:1px solid #dbeafe;border-radius:14px;padding:14px;cursor:pointer"><input type="radio" name="support_type" value="new_tool"' + (' checked' if support_type == 'new_tool' else '') + '><span><span style="display:block;font-size:15px;font-weight:700;color:#0f172a;margin-bottom:4px">בקשה לכלי חדש</span><span style="display:block;font-size:12px;color:#64748b;line-height:1.6">כתיבת בקשה חופשית לכלי או דוח חדש שתרצו שנבחן.</span></span></label>'
        + '<label style="display:flex;align-items:flex-start;gap:10px;background:#f8fafc;border:1px solid #dbeafe;border-radius:14px;padding:14px;cursor:pointer"><input type="radio" name="support_type" value="existing_tool"' + (' checked' if support_type == 'existing_tool' else '') + '><span><span style="display:block;font-size:15px;font-weight:700;color:#0f172a;margin-bottom:4px">תמיכה בכלי קיים</span><span style="display:block;font-size:12px;color:#64748b;line-height:1.6">דיווח על תקלה או בעיה באחד הכלים שכבר פתוחים לכם במערכת.</span></span></label>'
        + '</div>'
        + '<div id="existingToolFields" style="display:' + ('block' if support_type == 'existing_tool' else 'none') + ';margin-bottom:12px">'
        + inactive_support_note
        + '<label class="field-label">בחירת כלי קיים</label>'
        + '<select name="support_script_id">' + script_options + '</select>'
        + '</div>'
        + '<label class="field-label">תיאור הפנייה</label>'
        + '<textarea name="support_message" rows="7" style="margin-bottom:14px" placeholder="אפשר לכתוב כאן חופשי את הבקשה או את תיאור התקלה">' + esc(message_text) + '</textarea>'
        + '<div style="display:flex;gap:10px;justify-content:flex-end;flex-wrap:wrap"><a href="/dashboard" class="btn btn-gray" style="text-decoration:none;display:inline-flex;align-items:center;justify-content:center">חזרה</a><button type="submit" class="btn btn-blue">שליחת פנייה</button></div>'
        + '</form>'
        + requests_html
        + '<script>'
        + '(function(){var radios=Array.prototype.slice.call(document.querySelectorAll(\'input[name="support_type"]\'));var block=document.getElementById("existingToolFields");function refresh(){var selected=(document.querySelector(\'input[name="support_type"]:checked\')||{}).value||"new_tool";if(block){block.style.display=selected==="existing_tool"?"block":"none";}}radios.forEach(function(r){r.addEventListener("change",refresh);});refresh();})();'
        + '</script>'
        + '</div>'
    )
    return render(
        "שירות לקוחות" if lang == "he" else "Customer support",
        body,
        lang=lang,
        topbar_greeting=text["topbar_greeting"],
        logout_label=text["logout"],
        show_lang_switch=True,
    )


@app.route("/run/<script_id>", methods=["GET", "POST"])
@login_required
def run_script(script_id):
    if session.get("is_admin"):
        return redirect("/admin")

    lang = get_flow_lang()
    text = get_flow_text(lang)
    with get_db() as db:
        perm = db.execute(
            "SELECT 1 FROM permissions WHERE user_id=? AND script_id=?",
            (session["user_id"], script_id),
        ).fetchone()

    base_script = get_script(script_id)
    scr = get_localized_script(base_script, lang) if base_script else None

    if not perm or scr is None:
        add_flash(text["run_access_denied"])
        return redirect("/dashboard")
    if request.method == "GET":
        log_user_activity("open_script", "פתח כלי", script_id, scr["name"], "")
    result = None
    error = ""
    info_message = ""
    processing_warnings = []
    mapping_confirmation_html = ""
    mapping_templates = []
    selected_template = None

    if request.method == "POST":
        flow_mode = request.form.get("flow_mode", "").strip()
        if scr.get("requires_mapping_confirmation") and flow_mode != "confirm_mapping":
            file_obj = request.files.get("file")
            validation_error, ext = validate_upload(file_obj)
            if validation_error == "missing":
                error = '<div class="flash-err">' + scr["empty_error"] + '</div>'
            elif validation_error == "unsupported":
                error = '<div class="flash-err">' + scr["unsupported_error"] + '</div>'
            elif validation_error == "invalid_excel":
                error = '<div class="flash-err">' + scr["invalid_error"] + '</div>'
            elif validation_error == "empty":
                error = '<div class="flash-err">' + scr["empty_file_error"] + '</div>'
            elif validation_error == "too_large":
                error = '<div class="flash-err">' + scr["too_large_error"] + '</div>'
            else:
                uid = str(uuid.uuid4())[:8]
                inp = str(UPLOAD_FOLDER / f"{uid}_mapping.{ext}")
                file_obj.save(inp)
                if script_id == "flamingo_payroll":
                    inspection = build_flamingo_mapping_options(inp, ext)
                    selected_mapping = dict(default_flamingo_mapping())
                    selected_mapping.update(inspection["suggestions"])
                elif script_id == "matan_missing":
                    inspection = build_matan_missing_mapping_options(inp, ext)
                    selected_mapping = dict(default_matan_missing_mapping())
                    selected_mapping.update(inspection["suggestions"])
                elif script_id == "inactive_workers":
                    inspection = build_inactive_workers_mapping_options(inp, ext)
                    selected_mapping = dict(default_inactive_workers_mapping())
                    selected_mapping.update(inspection["suggestions"])
                elif script_id == "org_hierarchy_report":
                    inspection = build_org_hierarchy_mapping_options(inp, ext)
                    selected_mapping = dict(default_org_hierarchy_mapping())
                    selected_mapping.update(inspection["suggestions"])
                elif script_id == "dept_payroll":
                    inspection = build_dept_payroll_mapping_options(inp, ext)
                    selected_mapping = dict(default_dept_payroll_mapping())
                    selected_mapping.update(inspection["suggestions"])
                elif script_id == "matan_manual_corrections":
                    inspection = build_matan_corrections_mapping_options(inp, ext)
                    selected_mapping = dict(inspection["suggestions"])
                elif script_id == "attendance_alerts":
                    inspection = build_attendance_alerts_mapping_options(inp, ext)
                    selected_mapping = dict(default_attendance_alerts_mapping())
                    selected_mapping.update(inspection["suggestions"])
                else:
                    inspection = build_rimon_mapping_options(inp, ext)
                    selected_mapping = dict(inspection["suggestions"])
                mapping_templates = get_mapping_templates(session["user_id"], script_id)
                info_message = '<div class="flash" style="background:#eff6ff;border-color:#bfdbfe;color:#1d4ed8">המערכת זיהתה שדות אפשריים. נא לאשר או לתקן לפני הרצת הדוח.</div>'
                if script_id == "flamingo_payroll":
                    mapping_confirmation_html = build_flamingo_mapping_form(
                        script_id,
                        inp,
                        ext,
                        inspection,
                        selected_mapping,
                        mapping_templates,
                        get_next_mapping_template_name(mapping_templates),
                        "",
                    )
                elif script_id == "matan_missing":
                    mapping_confirmation_html = build_matan_missing_mapping_form(
                        script_id,
                        inp,
                        ext,
                        inspection,
                        selected_mapping,
                        mapping_templates,
                        get_next_mapping_template_name(mapping_templates),
                        {
                            "min_missing_hours": request.form.get("min_missing_hours", "").strip(),
                            "max_missing_hours": request.form.get("max_missing_hours", "").strip(),
                        },
                    )
                elif script_id == "inactive_workers":
                    mapping_confirmation_html = build_inactive_workers_mapping_form(
                        script_id,
                        inp,
                        ext,
                        inspection,
                        selected_mapping,
                        mapping_templates,
                        get_next_mapping_template_name(mapping_templates),
                        {
                            "inactive_period_unit": request.form.get("inactive_period_unit", "").strip() or "days",
                            "inactive_period_value": request.form.get("inactive_period_value", "").strip(),
                        },
                    )
                elif script_id == "org_hierarchy_report":
                    mapping_confirmation_html = build_org_hierarchy_mapping_form(
                        script_id,
                        inp,
                        ext,
                        inspection,
                        selected_mapping,
                        mapping_templates,
                        get_next_mapping_template_name(mapping_templates),
                        request.form.get("output_type", "").strip() or "powerpoint",
                    )
                elif script_id == "dept_payroll":
                    mapping_confirmation_html = build_dept_payroll_mapping_form(
                        script_id,
                        inp,
                        ext,
                        inspection,
                        selected_mapping,
                        mapping_templates,
                        get_next_mapping_template_name(mapping_templates),
                        "[]",
                        {},
                    )
                elif script_id == "matan_manual_corrections":
                    mapping_confirmation_html = build_matan_corrections_mapping_form(
                        script_id,
                        inp,
                        ext,
                        inspection,
                        selected_mapping,
                        mapping_templates,
                        get_next_mapping_template_name(mapping_templates),
                        {
                            "min_corrections": request.form.get("min_corrections", "").strip(),
                            "max_corrections": request.form.get("max_corrections", "").strip(),
                        },
                    )
                elif script_id == "attendance_alerts":
                    mapping_confirmation_html = build_attendance_alerts_mapping_form(
                        script_id,
                        inp,
                        ext,
                        inspection,
                        selected_mapping,
                        mapping_templates,
                        get_next_mapping_template_name(mapping_templates),
                    )
                else:
                    mapping_confirmation_html = build_rimon_mapping_form(
                        script_id,
                        inp,
                        ext,
                        inspection,
                        selected_mapping,
                        mapping_templates,
                        get_next_mapping_template_name(mapping_templates),
                    )
        elif scr.get("requires_mapping_confirmation") and flow_mode == "confirm_mapping":
            inp = request.form.get("temp_upload_path", "").strip()
            ext = request.form.get("temp_upload_ext", "").strip().lower()
            mapping = {}
            mapping_fields = FLAMINGO_MAPPING_FIELDS if script_id == "flamingo_payroll" else DEPT_PAYROLL_MAPPING_FIELDS if script_id == "dept_payroll" else MATAN_MISSING_MAPPING_FIELDS if script_id == "matan_missing" else INACTIVE_WORKERS_MAPPING_FIELDS if script_id == "inactive_workers" else ORG_HIERARCHY_MAPPING_FIELDS if script_id == "org_hierarchy_report" else MATAN_CORRECTIONS_MAPPING_FIELDS if script_id == "matan_manual_corrections" else ATTENDANCE_ALERTS_MAPPING_FIELDS if script_id == "attendance_alerts" else RIMON_MAPPING_FIELDS
            for field in mapping_fields:
                mapping[field["name"]] = request.form.get(field["name"], "").strip()
            manual_hourly_rate = request.form.get("manual_hourly_rate", "").strip() if script_id == "flamingo_payroll" else ""
            matan_filters = {
                "min_missing_hours": request.form.get("min_missing_hours", "").strip(),
                "max_missing_hours": request.form.get("max_missing_hours", "").strip(),
            }
            inactive_filters = {
                "inactive_period_unit": request.form.get("inactive_period_unit", "").strip() or "days",
                "inactive_period_value": request.form.get("inactive_period_value", "").strip(),
            }
            org_options = {
                "output_type": request.form.get("output_type", "").strip() or "powerpoint",
            }
            corrections_filters = {
                "min_corrections": request.form.get("min_corrections", "").strip(),
                "max_corrections": request.form.get("max_corrections", "").strip(),
            }
            dept_settings_json = request.form.get("dept_settings_json", "[]").strip() if script_id == "dept_payroll" else "[]"
            dept_manual_values = {
                "manual_hourly_rate": request.form.get("manual_hourly_rate", "").strip(),
                "manual_housing": request.form.get("manual_housing", "").strip(),
            } if script_id == "dept_payroll" else {}
            mapping_templates = get_mapping_templates(session["user_id"], script_id)
            mapping_action = request.form.get("mapping_action", "confirm").strip() or "confirm"
            selected_template_id = request.form.get("selected_template_id", "").strip()
            if mapping_action == "delete_template":
                if not inp or not os.path.exists(inp):
                    error = '<div class="flash-err">הקובץ הזמני לא נמצא. יש להעלות את הדוח מחדש.</div>'
                else:
                    if selected_template_id:
                        delete_mapping_template(session["user_id"], script_id, selected_template_id)
                        info_message = '<div class="flash" style="background:#eff6ff;border-color:#bfdbfe;color:#1d4ed8">התבנית נמחקה.</div>'
                    else:
                        info_message = '<div class="flash-err">לא נבחרה תבנית למחיקה.</div>'
                    inspection = build_flamingo_mapping_options(inp, ext) if script_id == "flamingo_payroll" else build_dept_payroll_mapping_options(inp, ext) if script_id == "dept_payroll" else build_matan_missing_mapping_options(inp, ext) if script_id == "matan_missing" else build_inactive_workers_mapping_options(inp, ext) if script_id == "inactive_workers" else build_org_hierarchy_mapping_options(inp, ext) if script_id == "org_hierarchy_report" else build_matan_corrections_mapping_options(inp, ext) if script_id == "matan_manual_corrections" else build_attendance_alerts_mapping_options(inp, ext) if script_id == "attendance_alerts" else build_rimon_mapping_options(inp, ext)
                    mapping_templates = get_mapping_templates(session["user_id"], script_id)
                    if script_id == "flamingo_payroll":
                        mapping_confirmation_html = build_flamingo_mapping_form(
                            script_id,
                            inp,
                            ext,
                            inspection,
                            mapping,
                            mapping_templates,
                            get_next_mapping_template_name(mapping_templates),
                            manual_hourly_rate,
                        )
                    elif script_id == "dept_payroll":
                        mapping_confirmation_html = build_dept_payroll_mapping_form(
                            script_id,
                            inp,
                            ext,
                            inspection,
                            mapping,
                            mapping_templates,
                            get_next_mapping_template_name(mapping_templates),
                            dept_settings_json,
                            dept_manual_values,
                        )
                    elif script_id == "matan_missing":
                        mapping_confirmation_html = build_matan_missing_mapping_form(
                            script_id,
                            inp,
                            ext,
                            inspection,
                            mapping,
                            mapping_templates,
                            get_next_mapping_template_name(mapping_templates),
                            matan_filters,
                        )
                    elif script_id == "inactive_workers":
                        mapping_confirmation_html = build_inactive_workers_mapping_form(
                            script_id,
                            inp,
                            ext,
                            inspection,
                            mapping,
                            mapping_templates,
                            get_next_mapping_template_name(mapping_templates),
                            inactive_filters,
                        )
                    elif script_id == "org_hierarchy_report":
                        mapping_confirmation_html = build_org_hierarchy_mapping_form(
                            script_id,
                            inp,
                            ext,
                            inspection,
                            mapping,
                            mapping_templates,
                            get_next_mapping_template_name(mapping_templates),
                            org_options["output_type"],
                        )
                    elif script_id == "matan_manual_corrections":
                        mapping_confirmation_html = build_matan_corrections_mapping_form(
                            script_id,
                            inp,
                            ext,
                            inspection,
                            mapping,
                            mapping_templates,
                            get_next_mapping_template_name(mapping_templates),
                            corrections_filters,
                        )
                    elif script_id == "attendance_alerts":
                        mapping_confirmation_html = build_attendance_alerts_mapping_form(
                            script_id,
                            inp,
                            ext,
                            inspection,
                            mapping,
                            mapping_templates,
                            get_next_mapping_template_name(mapping_templates),
                        )
                    else:
                        mapping_confirmation_html = build_rimon_mapping_form(
                            script_id,
                            inp,
                            ext,
                            inspection,
                            mapping,
                            mapping_templates,
                            get_next_mapping_template_name(mapping_templates),
                        )
            elif mapping_action == "apply_template":
                inspection = build_flamingo_mapping_options(inp, ext) if script_id == "flamingo_payroll" else build_dept_payroll_mapping_options(inp, ext) if script_id == "dept_payroll" else build_matan_missing_mapping_options(inp, ext) if script_id == "matan_missing" else build_inactive_workers_mapping_options(inp, ext) if script_id == "inactive_workers" else build_org_hierarchy_mapping_options(inp, ext) if script_id == "org_hierarchy_report" else build_matan_corrections_mapping_options(inp, ext) if script_id == "matan_manual_corrections" else build_attendance_alerts_mapping_options(inp, ext) if script_id == "attendance_alerts" else build_rimon_mapping_options(inp, ext)
                selected_mapping, selected_template = apply_selected_template(
                    dict(default_flamingo_mapping()) if script_id == "flamingo_payroll" else dict(default_dept_payroll_mapping()) if script_id == "dept_payroll" else dict(default_matan_missing_mapping()) if script_id == "matan_missing" else dict(default_inactive_workers_mapping()) if script_id == "inactive_workers" else dict(default_org_hierarchy_mapping()) if script_id == "org_hierarchy_report" else dict(default_attendance_alerts_mapping()) if script_id == "attendance_alerts" else dict(inspection["suggestions"]),
                    mapping_templates,
                    selected_template_id,
                )
                selected_mapping.update({key: value for key, value in mapping.items() if value})
                if script_id == "flamingo_payroll" and selected_template:
                    manual_hourly_rate = str(selected_template["mapping"].get("manual_hourly_rate", manual_hourly_rate) or "")
                if script_id == "dept_payroll" and selected_template:
                    dept_settings_json = str(selected_template["mapping"].get("dept_settings_json", dept_settings_json) or "[]")
                    dept_manual_values = {
                        "manual_hourly_rate": str(selected_template["mapping"].get("manual_hourly_rate", "") or ""),
                        "manual_housing": str(selected_template["mapping"].get("manual_housing", "") or ""),
                    }
                info_message = '<div class="flash" style="background:#eff6ff;border-color:#bfdbfe;color:#1d4ed8">התבנית נטענה. אפשר לבדוק את השדות ואז להריץ.</div>'
                if script_id == "flamingo_payroll":
                    mapping_confirmation_html = build_flamingo_mapping_form(
                        script_id,
                        inp,
                        ext,
                        inspection,
                        selected_mapping,
                        mapping_templates,
                        request.form.get("template_name", "").strip() or get_next_mapping_template_name(mapping_templates),
                        manual_hourly_rate,
                    )
                elif script_id == "dept_payroll":
                    mapping_confirmation_html = build_dept_payroll_mapping_form(
                        script_id,
                        inp,
                        ext,
                        inspection,
                        selected_mapping,
                        mapping_templates,
                        request.form.get("template_name", "").strip() or get_next_mapping_template_name(mapping_templates),
                        dept_settings_json,
                        dept_manual_values,
                    )
                elif script_id == "matan_missing":
                    mapping_confirmation_html = build_matan_missing_mapping_form(
                        script_id,
                        inp,
                        ext,
                        inspection,
                        selected_mapping,
                        mapping_templates,
                        request.form.get("template_name", "").strip() or get_next_mapping_template_name(mapping_templates),
                        matan_filters,
                    )
                elif script_id == "inactive_workers":
                    mapping_confirmation_html = build_inactive_workers_mapping_form(
                        script_id,
                        inp,
                        ext,
                        inspection,
                        selected_mapping,
                        mapping_templates,
                        request.form.get("template_name", "").strip() or get_next_mapping_template_name(mapping_templates),
                        inactive_filters,
                    )
                elif script_id == "org_hierarchy_report":
                    mapping_confirmation_html = build_org_hierarchy_mapping_form(
                        script_id,
                        inp,
                        ext,
                        inspection,
                        selected_mapping,
                        mapping_templates,
                        request.form.get("template_name", "").strip() or get_next_mapping_template_name(mapping_templates),
                        org_options["output_type"],
                    )
                elif script_id == "matan_manual_corrections":
                    mapping_confirmation_html = build_matan_corrections_mapping_form(
                        script_id,
                        inp,
                        ext,
                        inspection,
                        selected_mapping,
                        mapping_templates,
                        request.form.get("template_name", "").strip() or get_next_mapping_template_name(mapping_templates),
                        corrections_filters,
                    )
                elif script_id == "attendance_alerts":
                    mapping_confirmation_html = build_attendance_alerts_mapping_form(
                        script_id,
                        inp,
                        ext,
                        inspection,
                        selected_mapping,
                        mapping_templates,
                        request.form.get("template_name", "").strip() or get_next_mapping_template_name(mapping_templates),
                    )
                else:
                    mapping_confirmation_html = build_rimon_mapping_form(
                        script_id,
                        inp,
                        ext,
                        inspection,
                        selected_mapping,
                        mapping_templates,
                        request.form.get("template_name", "").strip() or get_next_mapping_template_name(mapping_templates),
                    )
            elif not inp or not os.path.exists(inp):
                error = '<div class="flash-err">הקובץ הזמני לא נמצא. יש להעלות את הדוח מחדש.</div>'
            else:
                if script_id == "flamingo_payroll":
                    if not mapping.get("payable_hours_source"):
                        inspection = build_flamingo_mapping_options(inp, ext)
                        error = '<div class="flash-err">יש לבחור שדה שעות לתשלום בפועל לפני חישוב השכר.</div>'
                        mapping_confirmation_html = build_flamingo_mapping_form(
                            script_id,
                            inp,
                            ext,
                            inspection,
                            mapping,
                            mapping_templates,
                            request.form.get("template_name", "").strip() or get_next_mapping_template_name(mapping_templates),
                            manual_hourly_rate,
                        )
                        return render(
                            scr["name"],
                            '<a href="/dashboard" style="color:#2563eb;font-size:13px;font-weight:600;text-decoration:none;display:inline-flex;align-items:center;gap:4px;margin-bottom:1rem">' + text["back_arrow"] + ' ' + scr["back_label"] + '</a>'
                            + '<div class="card"><div style="font-size:40px;margin-bottom:.5rem">' + scr["icon"] + '</div><div style="font-size:20px;font-weight:700;color:#1e3a8a;margin-bottom:4px">' + scr["name"] + '</div>'
                            + error + mapping_confirmation_html + '</div>',
                            lang=lang,
                            topbar_greeting=text["topbar_greeting"],
                            logout_label=text["logout"],
                            show_lang_switch=True,
                        )
                    if mapping.get("hourly_rate_source") == "__manual__" and not manual_hourly_rate:
                        inspection = build_flamingo_mapping_options(inp, ext)
                        error = '<div class="flash-err">נבחר תעריף שעתי ידני, אבל לא הוזן ערך לתעריף.</div>'
                        mapping_confirmation_html = build_flamingo_mapping_form(
                            script_id,
                            inp,
                            ext,
                            inspection,
                            mapping,
                            mapping_templates,
                            request.form.get("template_name", "").strip() or get_next_mapping_template_name(mapping_templates),
                            manual_hourly_rate,
                        )
                        return render(
                            scr["name"],
                            '<a href="/dashboard" style="color:#2563eb;font-size:13px;font-weight:600;text-decoration:none;display:inline-flex;align-items:center;gap:4px;margin-bottom:1rem">' + text["back_arrow"] + ' ' + scr["back_label"] + '</a>'
                            + '<div class="card"><div style="font-size:40px;margin-bottom:.5rem">' + scr["icon"] + '</div><div style="font-size:20px;font-weight:700;color:#1e3a8a;margin-bottom:4px">' + scr["name"] + '</div>'
                            + error + mapping_confirmation_html + '</div>',
                            lang=lang,
                            topbar_greeting=text["topbar_greeting"],
                            logout_label=text["logout"],
                            show_lang_switch=True,
                        )
                elif script_id == "dept_payroll":
                    if not mapping.get("payable_hours_source"):
                        inspection = build_dept_payroll_mapping_options(inp, ext)
                        error = '<div class="flash-err">יש לבחור שדה שעות לתשלום בפועל.</div>'
                        mapping_confirmation_html = build_dept_payroll_mapping_form(
                            script_id, inp, ext, inspection, mapping, mapping_templates,
                            request.form.get("template_name", "").strip() or get_next_mapping_template_name(mapping_templates),
                            dept_settings_json, dept_manual_values,
                        )
                        return render(
                            scr["name"],
                            '<a href="/dashboard" style="color:#2563eb;font-size:13px;font-weight:600;text-decoration:none;display:inline-flex;align-items:center;gap:4px;margin-bottom:1rem">' + text["back_arrow"] + ' ' + scr["back_label"] + '</a>'
                            + '<div class="card"><div style="font-size:40px;margin-bottom:.5rem">' + scr["icon"] + '</div><div style="font-size:20px;font-weight:700;color:#1e3a8a;margin-bottom:4px">' + scr["name"] + '</div>'
                            + error + mapping_confirmation_html + '</div>',
                            lang=lang,
                            topbar_greeting=text["topbar_greeting"],
                            logout_label=text["logout"],
                            show_lang_switch=True,
                        )
                    if not mapping.get("event_source") or not mapping.get("total_hours_source"):
                        inspection = build_dept_payroll_mapping_options(inp, ext)
                        error = '<div class="flash-err">יש לבחור שדה אירוע/לקוח ושדה סה"כ שעות.</div>'
                        mapping_confirmation_html = build_dept_payroll_mapping_form(
                            script_id, inp, ext, inspection, mapping, mapping_templates,
                            request.form.get("template_name", "").strip() or get_next_mapping_template_name(mapping_templates),
                            dept_settings_json, dept_manual_values,
                        )
                        return render(
                            scr["name"],
                            '<a href="/dashboard" style="color:#2563eb;font-size:13px;font-weight:600;text-decoration:none;display:inline-flex;align-items:center;gap:4px;margin-bottom:1rem">' + text["back_arrow"] + ' ' + scr["back_label"] + '</a>'
                            + '<div class="card"><div style="font-size:40px;margin-bottom:.5rem">' + scr["icon"] + '</div><div style="font-size:20px;font-weight:700;color:#1e3a8a;margin-bottom:4px">' + scr["name"] + '</div>'
                            + error + mapping_confirmation_html + '</div>',
                            lang=lang,
                            topbar_greeting=text["topbar_greeting"],
                            logout_label=text["logout"],
                            show_lang_switch=True,
                        )
                    dept_manual_missing = []
                    if not mapping.get("hourly_rate_source") and not dept_manual_values.get("manual_hourly_rate"):
                        dept_manual_missing.append("תעריף שעתי")
                    if not mapping.get("housing_source") and not dept_manual_values.get("manual_housing"):
                        dept_manual_missing.append("חיוב דירה")
                    if dept_manual_missing:
                        inspection = build_dept_payroll_mapping_options(inp, ext)
                        error = '<div class="flash-err">נבחרה הזנה ידנית אבל לא הוזן ערך עבור: ' + ', '.join(dept_manual_missing) + '.</div>'
                        mapping_confirmation_html = build_dept_payroll_mapping_form(
                            script_id, inp, ext, inspection, mapping, mapping_templates,
                            request.form.get("template_name", "").strip() or get_next_mapping_template_name(mapping_templates),
                            dept_settings_json, dept_manual_values,
                        )
                        return render(
                            scr["name"],
                            '<a href="/dashboard" style="color:#2563eb;font-size:13px;font-weight:600;text-decoration:none;display:inline-flex;align-items:center;gap:4px;margin-bottom:1rem">' + text["back_arrow"] + ' ' + scr["back_label"] + '</a>'
                            + '<div class="card"><div style="font-size:40px;margin-bottom:.5rem">' + scr["icon"] + '</div><div style="font-size:20px;font-weight:700;color:#1e3a8a;margin-bottom:4px">' + scr["name"] + '</div>'
                            + error + mapping_confirmation_html + '</div>',
                            lang=lang,
                            topbar_greeting=text["topbar_greeting"],
                            logout_label=text["logout"],
                            show_lang_switch=True,
                        )
                elif script_id == "matan_missing":
                    identifier_values = [
                        mapping.get("employee_number_source"),
                        mapping.get("id_number_source"),
                        mapping.get("badge_number_source"),
                        mapping.get("passport_number_source"),
                    ]
                    if not mapping.get("employee_name_source"):
                        inspection = build_matan_missing_mapping_options(inp, ext)
                        error = '<div class="flash-err">יש לבחור שדה שם עובד לפני יצירת הדוח.</div>'
                        mapping_confirmation_html = build_matan_missing_mapping_form(
                            script_id,
                            inp,
                            ext,
                            inspection,
                            mapping,
                            mapping_templates,
                            request.form.get("template_name", "").strip() or get_next_mapping_template_name(mapping_templates),
                            matan_filters,
                        )
                        return render(
                            scr["name"],
                            '<a href="/dashboard" style="color:#2563eb;font-size:13px;font-weight:600;text-decoration:none;display:inline-flex;align-items:center;gap:4px;margin-bottom:1rem">' + text["back_arrow"] + ' ' + scr["back_label"] + '</a>'
                            + '<div class="card"><div style="font-size:40px;margin-bottom:.5rem">' + scr["icon"] + '</div><div style="font-size:20px;font-weight:700;color:#1e3a8a;margin-bottom:4px">' + scr["name"] + '</div>'
                            + error + mapping_confirmation_html + '</div>',
                            lang=lang,
                            topbar_greeting=text["topbar_greeting"],
                            logout_label=text["logout"],
                            show_lang_switch=True,
                        )
                    if not mapping.get("standard_hours_source") or not mapping.get("missing_hours_source"):
                        inspection = build_matan_missing_mapping_options(inp, ext)
                        error = '<div class="flash-err">יש לבחור שדה שעות תקן ושדה חוסר לפני יצירת הדוח.</div>'
                        mapping_confirmation_html = build_matan_missing_mapping_form(
                            script_id,
                            inp,
                            ext,
                            inspection,
                            mapping,
                            mapping_templates,
                            request.form.get("template_name", "").strip() or get_next_mapping_template_name(mapping_templates),
                            matan_filters,
                        )
                        return render(
                            scr["name"],
                            '<a href="/dashboard" style="color:#2563eb;font-size:13px;font-weight:600;text-decoration:none;display:inline-flex;align-items:center;gap:4px;margin-bottom:1rem">' + text["back_arrow"] + ' ' + scr["back_label"] + '</a>'
                            + '<div class="card"><div style="font-size:40px;margin-bottom:.5rem">' + scr["icon"] + '</div><div style="font-size:20px;font-weight:700;color:#1e3a8a;margin-bottom:4px">' + scr["name"] + '</div>'
                            + error + mapping_confirmation_html + '</div>',
                            lang=lang,
                            topbar_greeting=text["topbar_greeting"],
                            logout_label=text["logout"],
                            show_lang_switch=True,
                        )
                    if not any(identifier_values):
                        inspection = build_matan_missing_mapping_options(inp, ext)
                        error = '<div class="flash-err">יש לבחור לפחות מזהה אחד נוסף: מספר עובד, תעודת זהות, מספר תג או דרכון.</div>'
                        mapping_confirmation_html = build_matan_missing_mapping_form(
                            script_id,
                            inp,
                            ext,
                            inspection,
                            mapping,
                            mapping_templates,
                            request.form.get("template_name", "").strip() or get_next_mapping_template_name(mapping_templates),
                            matan_filters,
                        )
                        return render(
                            scr["name"],
                            '<a href="/dashboard" style="color:#2563eb;font-size:13px;font-weight:600;text-decoration:none;display:inline-flex;align-items:center;gap:4px;margin-bottom:1rem">' + text["back_arrow"] + ' ' + scr["back_label"] + '</a>'
                            + '<div class="card"><div style="font-size:40px;margin-bottom:.5rem">' + scr["icon"] + '</div><div style="font-size:20px;font-weight:700;color:#1e3a8a;margin-bottom:4px">' + scr["name"] + '</div>'
                            + error + mapping_confirmation_html + '</div>',
                            lang=lang,
                            topbar_greeting=text["topbar_greeting"],
                            logout_label=text["logout"],
                            show_lang_switch=True,
                        )
                elif script_id == "inactive_workers":
                    identifier_values = [
                        mapping.get("employee_number_source"),
                        mapping.get("badge_number_source"),
                        mapping.get("id_number_source"),
                        mapping.get("passport_number_source"),
                    ]
                    has_entry_exit = bool(mapping.get("entry_time_source") and mapping.get("exit_time_source"))
                    has_total_hours = bool(mapping.get("total_hours_source"))
                    period_value = inactive_filters.get("inactive_period_value", "")
                    if not mapping.get("employee_name_source") or not mapping.get("date_source"):
                        inspection = build_inactive_workers_mapping_options(inp, ext)
                        error = '<div class="flash-err">יש לבחור שדה שם עובד ושדה תאריך לפני יצירת הדוח.</div>'
                        mapping_confirmation_html = build_inactive_workers_mapping_form(
                            script_id,
                            inp,
                            ext,
                            inspection,
                            mapping,
                            mapping_templates,
                            request.form.get("template_name", "").strip() or get_next_mapping_template_name(mapping_templates),
                            inactive_filters,
                        )
                        return render(
                            scr["name"],
                            '<a href="/dashboard" style="color:#2563eb;font-size:13px;font-weight:600;text-decoration:none;display:inline-flex;align-items:center;gap:4px;margin-bottom:1rem">' + text["back_arrow"] + ' ' + scr["back_label"] + '</a>'
                            + '<div class="card"><div style="font-size:40px;margin-bottom:.5rem">' + scr["icon"] + '</div><div style="font-size:20px;font-weight:700;color:#1e3a8a;margin-bottom:4px">' + scr["name"] + '</div>'
                            + error + mapping_confirmation_html + '</div>',
                            lang=lang,
                            topbar_greeting=text["topbar_greeting"],
                            logout_label=text["logout"],
                            show_lang_switch=True,
                        )
                    if not any(identifier_values):
                        inspection = build_inactive_workers_mapping_options(inp, ext)
                        error = '<div class="flash-err">יש לבחור לפחות מזהה אחד נוסף: מספר עובד, מספר תג, תעודת זהות או דרכון.</div>'
                        mapping_confirmation_html = build_inactive_workers_mapping_form(
                            script_id,
                            inp,
                            ext,
                            inspection,
                            mapping,
                            mapping_templates,
                            request.form.get("template_name", "").strip() or get_next_mapping_template_name(mapping_templates),
                            inactive_filters,
                        )
                        return render(
                            scr["name"],
                            '<a href="/dashboard" style="color:#2563eb;font-size:13px;font-weight:600;text-decoration:none;display:inline-flex;align-items:center;gap:4px;margin-bottom:1rem">' + text["back_arrow"] + ' ' + scr["back_label"] + '</a>'
                            + '<div class="card"><div style="font-size:40px;margin-bottom:.5rem">' + scr["icon"] + '</div><div style="font-size:20px;font-weight:700;color:#1e3a8a;margin-bottom:4px">' + scr["name"] + '</div>'
                            + error + mapping_confirmation_html + '</div>',
                            lang=lang,
                            topbar_greeting=text["topbar_greeting"],
                            logout_label=text["logout"],
                            show_lang_switch=True,
                        )
                    if not has_entry_exit and not has_total_hours:
                        inspection = build_inactive_workers_mapping_options(inp, ext)
                        error = '<div class="flash-err">יש לבחור או שדה כניסה ושדה יציאה יחד, או לחלופין שדה סה&quot;כ שעות.</div>'
                        mapping_confirmation_html = build_inactive_workers_mapping_form(
                            script_id,
                            inp,
                            ext,
                            inspection,
                            mapping,
                            mapping_templates,
                            request.form.get("template_name", "").strip() or get_next_mapping_template_name(mapping_templates),
                            inactive_filters,
                        )
                        return render(
                            scr["name"],
                            '<a href="/dashboard" style="color:#2563eb;font-size:13px;font-weight:600;text-decoration:none;display:inline-flex;align-items:center;gap:4px;margin-bottom:1rem">' + text["back_arrow"] + ' ' + scr["back_label"] + '</a>'
                            + '<div class="card"><div style="font-size:40px;margin-bottom:.5rem">' + scr["icon"] + '</div><div style="font-size:20px;font-weight:700;color:#1e3a8a;margin-bottom:4px">' + scr["name"] + '</div>'
                            + error + mapping_confirmation_html + '</div>',
                            lang=lang,
                            topbar_greeting=text["topbar_greeting"],
                            logout_label=text["logout"],
                            show_lang_switch=True,
                        )
                    try:
                        period_int = int(period_value)
                    except (TypeError, ValueError):
                        period_int = 0
                    if period_int <= 0:
                        inspection = build_inactive_workers_mapping_options(inp, ext)
                        error = '<div class="flash-err">יש להזין ערך חיובי לטווח הבדיקה בימים או בחודשים.</div>'
                        mapping_confirmation_html = build_inactive_workers_mapping_form(
                            script_id,
                            inp,
                            ext,
                            inspection,
                            mapping,
                            mapping_templates,
                            request.form.get("template_name", "").strip() or get_next_mapping_template_name(mapping_templates),
                            inactive_filters,
                        )
                        return render(
                            scr["name"],
                            '<a href="/dashboard" style="color:#2563eb;font-size:13px;font-weight:600;text-decoration:none;display:inline-flex;align-items:center;gap:4px;margin-bottom:1rem">' + text["back_arrow"] + ' ' + scr["back_label"] + '</a>'
                            + '<div class="card"><div style="font-size:40px;margin-bottom:.5rem">' + scr["icon"] + '</div><div style="font-size:20px;font-weight:700;color:#1e3a8a;margin-bottom:4px">' + scr["name"] + '</div>'
                            + error + mapping_confirmation_html + '</div>',
                            lang=lang,
                            topbar_greeting=text["topbar_greeting"],
                            logout_label=text["logout"],
                            show_lang_switch=True,
                        )
                elif script_id == "org_hierarchy_report":
                    identifier_values = [
                        mapping.get("employee_number_source"),
                        mapping.get("id_number_source"),
                        mapping.get("passport_number_source"),
                    ]
                    if not mapping.get("employee_name_source"):
                        inspection = build_org_hierarchy_mapping_options(inp, ext)
                        error = '<div class="flash-err">יש לבחור שדה שם עובד לפני יצירת הדוח.</div>'
                        mapping_confirmation_html = build_org_hierarchy_mapping_form(
                            script_id,
                            inp,
                            ext,
                            inspection,
                            mapping,
                            mapping_templates,
                            request.form.get("template_name", "").strip() or get_next_mapping_template_name(mapping_templates),
                            org_options["output_type"],
                        )
                        return render(
                            scr["name"],
                            '<a href="/dashboard" style="color:#2563eb;font-size:13px;font-weight:600;text-decoration:none;display:inline-flex;align-items:center;gap:4px;margin-bottom:1rem">' + text["back_arrow"] + ' ' + scr["back_label"] + '</a>'
                            + '<div class="card"><div style="font-size:40px;margin-bottom:.5rem">' + scr["icon"] + '</div><div style="font-size:20px;font-weight:700;color:#1e3a8a;margin-bottom:4px">' + scr["name"] + '</div>'
                            + error + mapping_confirmation_html + '</div>',
                            lang=lang,
                            topbar_greeting=text["topbar_greeting"],
                            logout_label=text["logout"],
                            show_lang_switch=True,
                        )
                    if not mapping.get("direct_manager_source") or not mapping.get("department_source"):
                        inspection = build_org_hierarchy_mapping_options(inp, ext)
                        error = '<div class="flash-err">יש לבחור שדה מנהל ישיר ושדה מחלקה לפני יצירת הדוח.</div>'
                        mapping_confirmation_html = build_org_hierarchy_mapping_form(
                            script_id,
                            inp,
                            ext,
                            inspection,
                            mapping,
                            mapping_templates,
                            request.form.get("template_name", "").strip() or get_next_mapping_template_name(mapping_templates),
                            org_options["output_type"],
                        )
                        return render(
                            scr["name"],
                            '<a href="/dashboard" style="color:#2563eb;font-size:13px;font-weight:600;text-decoration:none;display:inline-flex;align-items:center;gap:4px;margin-bottom:1rem">' + text["back_arrow"] + ' ' + scr["back_label"] + '</a>'
                            + '<div class="card"><div style="font-size:40px;margin-bottom:.5rem">' + scr["icon"] + '</div><div style="font-size:20px;font-weight:700;color:#1e3a8a;margin-bottom:4px">' + scr["name"] + '</div>'
                            + error + mapping_confirmation_html + '</div>',
                            lang=lang,
                            topbar_greeting=text["topbar_greeting"],
                            logout_label=text["logout"],
                            show_lang_switch=True,
                        )
                    if not any(identifier_values):
                        inspection = build_org_hierarchy_mapping_options(inp, ext)
                        error = '<div class="flash-err">יש לבחור לפחות מזהה אחד נוסף: מספר עובד, תעודת זהות או דרכון.</div>'
                        mapping_confirmation_html = build_org_hierarchy_mapping_form(
                            script_id,
                            inp,
                            ext,
                            inspection,
                            mapping,
                            mapping_templates,
                            request.form.get("template_name", "").strip() or get_next_mapping_template_name(mapping_templates),
                            org_options["output_type"],
                        )
                        return render(
                            scr["name"],
                            '<a href="/dashboard" style="color:#2563eb;font-size:13px;font-weight:600;text-decoration:none;display:inline-flex;align-items:center;gap:4px;margin-bottom:1rem">' + text["back_arrow"] + ' ' + scr["back_label"] + '</a>'
                            + '<div class="card"><div style="font-size:40px;margin-bottom:.5rem">' + scr["icon"] + '</div><div style="font-size:20px;font-weight:700;color:#1e3a8a;margin-bottom:4px">' + scr["name"] + '</div>'
                            + error + mapping_confirmation_html + '</div>',
                            lang=lang,
                            topbar_greeting=text["topbar_greeting"],
                            logout_label=text["logout"],
                            show_lang_switch=True,
                        )
                elif script_id == "matan_manual_corrections":
                    if not mapping.get("entry_col_source") or not mapping.get("exit_col_source"):
                        inspection = build_matan_corrections_mapping_options(inp, ext)
                        error = '<div class="flash-err">יש לבחור עמודת כניסה ועמודת יציאה לפני יצירת הדוח.</div>'
                        mapping_confirmation_html = build_matan_corrections_mapping_form(
                            script_id,
                            inp,
                            ext,
                            inspection,
                            mapping,
                            mapping_templates,
                            request.form.get("template_name", "").strip() or get_next_mapping_template_name(mapping_templates),
                            corrections_filters,
                        )
                        return render(
                            scr["name"],
                            '<a href="/dashboard" style="color:#2563eb;font-size:13px;font-weight:600;text-decoration:none;display:inline-flex;align-items:center;gap:4px;margin-bottom:1rem">' + text["back_arrow"] + ' ' + scr["back_label"] + '</a>'
                            + '<div class="card"><div style="font-size:40px;margin-bottom:.5rem">' + scr["icon"] + '</div><div style="font-size:20px;font-weight:700;color:#1e3a8a;margin-bottom:4px">' + scr["name"] + '</div>'
                            + error + mapping_confirmation_html + '</div>',
                            lang=lang,
                            topbar_greeting=text["topbar_greeting"],
                            logout_label=text["logout"],
                            show_lang_switch=True,
                        )
                uid = str(uuid.uuid4())[:8]
                options = {key: value for key, value in mapping.items()}
                if script_id == "flamingo_payroll":
                    options["manual_hourly_rate"] = manual_hourly_rate
                elif script_id == "dept_payroll":
                    options["dept_settings_json"] = dept_settings_json
                    if dept_manual_values.get("manual_hourly_rate"):
                        options["manual_hourly_rate"] = dept_manual_values["manual_hourly_rate"]
                    if dept_manual_values.get("manual_housing"):
                        options["manual_housing"] = dept_manual_values["manual_housing"]
                    # Handle optional clients file upload
                    clients_file = request.files.get("clients_file")
                    if clients_file and clients_file.filename:
                        clients_ext = get_extension(clients_file.filename)
                        if clients_ext in ("xls", "xlsx"):
                            clients_temp_path = str(UPLOAD_FOLDER / f"{uid}_clients.{clients_ext}")
                            clients_file.save(clients_temp_path)
                            options["clients_file_path"] = clients_temp_path
                elif script_id == "matan_missing":
                    options.update(matan_filters)
                elif script_id == "inactive_workers":
                    options.update(inactive_filters)
                elif script_id == "org_hierarchy_report":
                    options.update(org_options)
                elif script_id == "matan_manual_corrections":
                    options.update(corrections_filters)
                elif script_id == "attendance_alerts":
                    for a in ATTENDANCE_ALERT_TYPES:
                        key = "alert_" + a["id"]
                        options[key] = "1" if request.form.get(key) else "0"
                result_name = build_output_filename(scr, uid, options)
                out = str(OUTPUT_FOLDER / result_name)
                try:
                    execution_result = execute_script(scr, inp, out, ext, options) or {}
                    result = result_name
                    processing_warnings = execution_result.get("warnings", [])
                    log_user_activity("generate_report", "הפיק דוח", script_id, scr["name"], result_name)
                    if request.form.get("save_template") == "1":
                        template_name = request.form.get("template_name", "").strip() or get_next_mapping_template_name(mapping_templates)
                        template_mapping = dict(mapping)
                        if script_id == "flamingo_payroll" and manual_hourly_rate:
                            template_mapping["manual_hourly_rate"] = manual_hourly_rate
                        if script_id == "dept_payroll" and dept_settings_json:
                            template_mapping["dept_settings_json"] = dept_settings_json
                        if script_id == "dept_payroll" and dept_manual_values.get("manual_hourly_rate"):
                            template_mapping["manual_hourly_rate"] = dept_manual_values["manual_hourly_rate"]
                        if script_id == "dept_payroll" and dept_manual_values.get("manual_housing"):
                            template_mapping["manual_housing"] = dept_manual_values["manual_housing"]
                        save_mapping_template(session["user_id"], script_id, template_name, template_mapping)
                except (xlrd.biffh.XLRDError, BadZipFile, OSError, ValueError):
                    error = '<div class="flash-err">' + scr["processing_error"] + '</div>'
                except Exception as e:
                    error = '<div class="flash-err">' + text["run_unexpected_error_prefix"] + str(e) + "</div>"
                finally:
                    try:
                        os.remove(inp)
                    except OSError:
                        pass
                    # Clean up clients file if uploaded
                    clients_temp = options.get("clients_file_path", "")
                    if clients_temp:
                        try:
                            os.remove(clients_temp)
                        except OSError:
                            pass
        else:
            file_obj = request.files.get("file")
            validation_error, ext = validate_upload(file_obj)
            if validation_error == "missing":
                error = '<div class="flash-err">' + scr["empty_error"] + '</div>'
            elif validation_error == "unsupported":
                error = '<div class="flash-err">' + scr["unsupported_error"] + '</div>'
            elif validation_error == "invalid_excel":
                error = '<div class="flash-err">' + scr["invalid_error"] + '</div>'
            elif validation_error == "empty":
                error = '<div class="flash-err">' + scr["empty_file_error"] + '</div>'
            elif validation_error == "too_large":
                error = '<div class="flash-err">' + scr["too_large_error"] + '</div>'
            else:
                uid = str(uuid.uuid4())[:8]
                inp = str(UPLOAD_FOLDER / f"{uid}.{ext}")
                options = {}
                extra_paths = []
                for field in scr.get("filter_fields", []):
                    options[field["name"]] = request.form.get(field["name"], "").strip()
                result_name = build_output_filename(scr, uid, options)
                out = str(OUTPUT_FOLDER / result_name)
                for upload in scr.get("extra_uploads", []):
                    extra_file = request.files.get(upload["name"])
                    if extra_file and extra_file.filename:
                        extra_ext = get_extension(extra_file.filename)
                        expected = upload.get("accept", "").lstrip(".").lower()
                        if expected and extra_ext != expected:
                            error = '<div class="flash-err">' + text["run_extra_file_type_error"] + "</div>"
                            break
                        extra_path = str(UPLOAD_FOLDER / f"{uid}_{upload['name']}.{extra_ext or 'dat'}")
                        extra_file.save(extra_path)
                        options[f"{upload['name']}_path"] = extra_path
                        extra_paths.append(extra_path)
                    elif upload.get("required"):
                        error = '<div class="flash-err">' + text["run_missing_extra_file_error"] + "</div>"
                        break
                if error:
                    for path in extra_paths:
                        try:
                            os.remove(path)
                        except OSError:
                            pass
                else:
                    file_obj.save(inp)
                    if scr.get("background_queue"):
                        with get_db() as db:
                            user_row = db.execute("SELECT * FROM users WHERE id=?", (session["user_id"],)).fetchone()
                        if user_row is None:
                            try:
                                os.remove(inp)
                            except OSError:
                                pass
                            error = '<div class="flash-err">המשתמש לא נמצא. יש להתחבר מחדש.</div>'
                        else:
                            job_id = create_report_job(user_row, script_id, scr["name"], file_obj.filename, inp, ext, result_name)
                            start_cleanup_report_job(job_id)
                            log_user_activity("queue_report", "שלח דוח לעיבוד ברקע", script_id, scr["name"], result_name)
                            add_flash("הדוח נקלט לעיבוד ברקע. אפשר להמשיך לעבוד בפלטפורמה, וכשהוא יהיה מוכן הוא יופיע בדוחות המוכנים להורדה.")
                            return redirect("/dashboard")
                    else:
                        try:
                            execution_result = execute_script(scr, inp, out, ext, options) or {}
                            result = result_name
                            processing_warnings = execution_result.get("warnings", [])
                            log_user_activity("generate_report", "הפיק דוח", script_id, scr["name"], result_name)
                        except (xlrd.biffh.XLRDError, BadZipFile, OSError, ValueError):
                            error = '<div class="flash-err">' + scr["processing_error"] + '</div>'
                        except Exception as e:
                            error = '<div class="flash-err">' + text["run_unexpected_error_prefix"] + str(e) + "</div>"
                        finally:
                            try:
                                os.remove(inp)
                            except OSError:
                                pass
                            for path in extra_paths:
                                try:
                                    os.remove(path)
                                except OSError:
                                    pass

    if result:
        warning_html = ""
        if processing_warnings:
            warning_html = '<div style="text-align:right;background:#fff7ed;border:1px solid #fdba74;color:#9a3412;border-radius:12px;padding:12px 14px;margin-bottom:14px;line-height:1.8">' + "<br>".join(esc(item) for item in processing_warnings) + "</div>"
        content = (
            warning_html
            + '<div class="success-box">'
            '<div style="font-size:32px;margin-bottom:6px">&#9989;</div>'
            '<div style="font-size:16px;font-weight:700;color:#15803d;margin-bottom:10px">' + scr["success_title"] + '</div>'
            '<a href="/download/' + result + '" class="dl-btn">&#8681; ' + scr["success_action"] + '</a>'
            '<br><br><a href="/run/' + script_id + '" style="font-size:13px;color:#2563eb">' + scr["retry_action"] + '</a>'
            '</div>'
        )
    else:
        extra_uploads_html = ""
        for upload in scr.get("extra_uploads", []):
            extra_uploads_html += (
                '<div style="margin-top:1rem;padding-top:1rem;border-top:1px solid #e2e8f0">'
                + '<div style="font-size:14px;font-weight:600;color:#334155;margin-bottom:8px">' + upload["label"] + '</div>'
                + '<input type="file" name="' + upload["name"] + '" accept="' + upload.get("accept", "") + '" style="width:100%;max-width:420px;margin:0 auto 8px;display:block;font-family:inherit">'
                + '<div style="font-size:12px;color:#94a3b8">' + upload.get("help", "") + '</div>'
                + '</div>'
            )
        filter_fields_html = ""
        if scr.get("filter_fields"):
            filter_fields_html += '<div style="display:grid;grid-template-columns:repeat(auto-fit,minmax(180px,1fr));gap:10px;margin-bottom:1rem">'
            for field in scr.get("filter_fields", []):
                if field.get("type") == "select":
                    select_options = ""
                    default_value = field.get("default", "")
                    for option in field.get("options", []):
                        selected = ' selected' if option.get("value") == default_value else ""
                        select_options += '<option value="' + option.get("value", "") + '"' + selected + '>' + option.get("label", "") + '</option>'
                    filter_fields_html += (
                        '<div><label class="field-label">' + field["label"] + '</label>'
                        + '<select name="' + field["name"] + '" style="padding:9px 12px;border:1.5px solid #e2e8f0;border-radius:8px;font-size:13px;font-family:inherit;outline:none;width:100%;margin-bottom:0;background:white">'
                        + select_options
                        + '</select></div>'
                    )
                else:
                    filter_fields_html += (
                        '<div><label class="field-label">' + field["label"] + '</label>'
                        + '<input type="text" name="' + field["name"] + '" placeholder="' + field.get("placeholder", "") + '" style="margin-bottom:0"></div>'
                    )
            filter_fields_html += '</div>'
        if mapping_confirmation_html:
            content = error + info_message + mapping_confirmation_html
        else:
            content = (
                error
                + info_message
                + '<form method="POST" enctype="multipart/form-data" id="uploadForm">'
                + filter_fields_html
                + '<div style="background:#fafcff;border:2px dashed #c7d7f5;border-radius:14px;padding:1.5rem;margin-bottom:1rem;text-align:center">'
                + '<div style="font-size:32px;margin-bottom:8px">&#128194;</div>'
                + '<div style="font-size:15px;font-weight:600;color:#1e40af;margin-bottom:12px">' + scr["file_picker_label"] + '</div>'
                + '<input type="file" name="file" id="fi" accept="' + scr["accept"] + '" style="width:100%;max-width:420px;margin:0 auto 10px;display:block;font-family:inherit">'
                + '<div style="font-size:12px;color:#94a3b8" id="lbl">' + scr["accept"] + '</div>'
                + extra_uploads_html
                + '</div>'
                + '<button type="submit" class="btn btn-blue" id="gb" style="width:100%;padding:13px;font-size:15px;font-weight:700">' + scr["icon"] + ' ' + scr["submit_label"] + '</button>'
                + '<div class="processing-box" id="processingBox">'
                + '<div class="processing-note">' + scr["processing_title"] + '</div>'
                + '<div class="progress-track"><div class="progress-bar"></div></div>'
                + '<div class="processing-subnote">' + scr["processing_note"] + '</div>'
                + '</div>'
                + '</form>'
            )

    help_trigger_html = ""
    help_modal_html = ""
    rules_trigger_html = ""
    rules_modal_html = ""
    activity_script_id = json.dumps(script_id)
    activity_script_name = json.dumps(scr["name"])
    if lang == "he" and scr.get("help_title"):
        help_items_html = "".join(
            '<li style="margin-bottom:6px">' + esc(item) + '</li>'
            for item in scr.get("help_items", [])
        )
        help_trigger_html = (
            '<button type="button" onclick="openHelpModal()" style="display:inline-flex;align-items:center;gap:6px;border:1px solid #bfdbfe;background:#eff6ff;color:#1d4ed8;border-radius:999px;padding:6px 10px;font-size:12px;font-weight:700;cursor:pointer;font-family:inherit">'
            '<span style="display:inline-flex;align-items:center;justify-content:center;width:18px;height:18px;border-radius:999px;background:#dbeafe">?</span>'
            + esc(scr.get("help_label", "מידע נוסף"))
            + '</button>'
        )
        help_modal_html = (
            '<div class="modal-bg" id="helpModal" onclick="closeHelpModal(event)">'
            '<div class="modal-box" style="width:100%;max-width:560px;padding:1.5rem 1.5rem 1.25rem;border-radius:18px">'
            '<div style="font-size:18px;font-weight:800;color:#0f172a;margin-bottom:.75rem">' + esc(scr["help_title"]) + '</div>'
            '<div style="font-size:14px;line-height:1.8;color:#334155">'
            + esc(scr.get("help_intro", ""))
            + '<ul style="margin:.6rem 0 .75rem;padding-inline-start:1.2rem">'
            + help_items_html
            + '</ul>'
            + esc(scr.get("help_note", ""))
            + '</div>'
            + '<div style="margin-top:1rem"><button type="button" class="btn btn-blue" style="width:100%" onclick="closeHelpModal()">סגור</button></div>'
            + '</div></div>'
        )
    if lang == "he" and scr.get("rules_title"):
        rules_items_html = "".join(
            '<li style="margin-bottom:6px">' + esc(item) + '</li>'
            for item in scr.get("rules_items", [])
        )
        rules_trigger_html = (
            '<button type="button" onclick="openRulesModal()" style="display:inline-flex;align-items:center;gap:6px;border:1px solid #c7d2fe;background:#eef2ff;color:#4338ca;border-radius:999px;padding:6px 10px;font-size:12px;font-weight:700;cursor:pointer;font-family:inherit">'
            '<span style="display:inline-flex;align-items:center;justify-content:center;width:18px;height:18px;border-radius:999px;background:#e0e7ff">i</span>'
            + esc(scr.get("rules_label", "איך זה עובד"))
            + '</button>'
        )
        rules_modal_html = (
            '<div class="modal-bg" id="rulesModal" onclick="closeRulesModal(event)">'
            '<div class="modal-box" style="width:100%;max-width:560px;padding:1.5rem 1.5rem 1.25rem;border-radius:18px">'
            '<div style="font-size:18px;font-weight:800;color:#0f172a;margin-bottom:.75rem">' + esc(scr["rules_title"]) + '</div>'
            '<div style="font-size:14px;line-height:1.8;color:#334155">'
            + esc(scr.get("rules_intro", ""))
            + '<ul style="margin:.6rem 0 .75rem;padding-inline-start:1.2rem">'
            + rules_items_html
            + '</ul>'
            + esc(scr.get("rules_note", ""))
            + '</div>'
            + '<div style="margin-top:1rem"><button type="button" class="btn btn-blue" style="width:100%" onclick="closeRulesModal()">סגור</button></div>'
            + '</div></div>'
        )

    body = (
        '<a href="/dashboard" style="color:#2563eb;font-size:13px;font-weight:600;text-decoration:none;display:inline-flex;align-items:center;gap:4px;margin-bottom:1rem">' + text["back_arrow"] + ' ' + scr["back_label"] + '</a>'
        + '<div class="card" style="border-radius:18px">'
        + '<div style="width:56px;height:56px;border-radius:16px;background:linear-gradient(135deg,#eff6ff,#dbeafe);display:flex;align-items:center;justify-content:center;font-size:28px;margin-bottom:.75rem">' + scr["icon"] + '</div>'
        + '<div style="font-size:22px;font-weight:800;color:#0f172a;margin-bottom:6px;letter-spacing:-0.3px">' + scr["name"] + '</div>'
        + '<div style="display:flex;align-items:center;justify-content:center;gap:10px;flex-wrap:wrap;margin-bottom:1.75rem">'
        + '<div style="font-size:14px;color:#64748b;line-height:1.6">' + scr["desc"] + '</div>'
        + help_trigger_html
        + rules_trigger_html
        + '</div>'
        + content
        + '</div>'
        + help_modal_html
        + rules_modal_html
        + '<script>'
        + 'var fileInput=document.getElementById("fi");'
        + 'var label=document.getElementById("lbl");'
        + 'var button=document.getElementById("gb");'
        + 'var form=document.getElementById("uploadForm");'
        + 'function trackUserActivity(eventType, actionLabel, scriptId, scriptName, details){try{var data=new FormData();data.append("event_type",eventType||"");data.append("action_label",actionLabel||"");data.append("script_id",scriptId||"");data.append("script_name",scriptName||"");data.append("details",details||"");if(navigator.sendBeacon){navigator.sendBeacon("/activity",data);}else{fetch("/activity",{method:"POST",body:data,credentials:"same-origin",keepalive:true});}}catch(e){}}'
        + 'var activityScriptId=' + activity_script_id + ';'
        + 'var activityScriptName=' + activity_script_name + ';'
        + 'function openHelpModal(){var modal=document.getElementById("helpModal");if(modal){modal.style.display="flex";trackUserActivity("open_help_popup","פתח מידע והנחיות נוספות",activityScriptId,activityScriptName,"");}}'
        + 'function closeHelpModal(event){if(event && event.target && event.target.id!=="helpModal"){return;}var modal=document.getElementById("helpModal");if(modal){modal.style.display="none";}}'
        + 'function openRulesModal(){var modal=document.getElementById("rulesModal");if(modal){modal.style.display="flex";trackUserActivity("open_logic_popup","פתח הסבר על אופן החישוב",activityScriptId,activityScriptName,"");}}'
        + 'function closeRulesModal(event){if(event && event.target && event.target.id!=="rulesModal"){return;}var modal=document.getElementById("rulesModal");if(modal){modal.style.display="none";}}'
        + 'if(fileInput && label){fileInput.addEventListener("change", function(){if(this.files && this.files.length){label.textContent=this.files[0].name;}});}'
        + 'if(form){form.addEventListener("submit", function(event){if(!fileInput || !fileInput.files || !fileInput.files.length){event.preventDefault();return false;}button.disabled=true;button.textContent="' + scr["processing_title"] + '";var box=document.getElementById("processingBox");if(box){box.classList.add("show");}return true;});}'
        + '</script>'
    )
    return render(
        scr["name"],
        body,
        lang=lang,
        topbar_greeting=text["topbar_greeting"],
        logout_label=text["logout"],
        show_lang_switch=True,
    )

@app.route("/download/<filename>")
@login_required
def download(filename):
    path = OUTPUT_FOLDER / filename
    if not path.exists():
        add_flash("File not found")
        return redirect("/dashboard")
    download_name = filename.split("_", 1)[-1] if "_" in filename else filename
    script = resolve_script_from_output_name(filename)
    if script:
        log_user_activity("download_report", "הוריד דוח", script.get("id", ""), script.get("name", ""), download_name)
    return send_file(path, as_attachment=True, download_name=download_name)


@app.route("/sample-clients-file")
@login_required
def download_sample_clients_file():
    """Generate and download a professional sample clients Excel file."""
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    wb = Workbook()
    ws = wb.active
    ws.title = "נתוני לקוחות"
    ws.sheet_view.rightToLeft = True

    # ── Title row ──
    title_font = Font(bold=True, size=14, color="FFFFFF")
    title_fill = PatternFill(start_color="1E3A8A", end_color="1E3A8A", fill_type="solid")
    ws.merge_cells("A1:J1")
    title_cell = ws.cell(row=1, column=1, value="קובץ נתוני לקוחות — Scriptly")
    title_cell.font = title_font
    title_cell.fill = title_fill
    title_cell.alignment = Alignment(horizontal="center", vertical="center")
    ws.row_dimensions[1].height = 32

    # ── Subtitle row ──
    sub_font = Font(size=11, color="475569", italic=True)
    ws.merge_cells("A2:J2")
    sub_cell = ws.cell(row=2, column=1, value="מלא את פרטי הלקוחות בשורות מתחת לכותרות. שדות עם * הם חובה. ניתן להוסיף שורות לפי הצורך.")
    sub_cell.font = sub_font
    sub_cell.alignment = Alignment(horizontal="center", vertical="center")
    ws.row_dimensions[2].height = 24

    # ── Headers ──
    headers = [
        ("שם העסק *", 22, True),
        ("תעריף גביה *", 14, True),
        ("מנהל אזור", 16, False),
        ("כתובת העסק", 22, False),
        ("איש קשר", 16, False),
        ("טלפון", 14, False),
        ("שם העובד", 18, False),
        ("מספר פספורט", 14, False),
        ("שכר שעתי", 12, False),
        ("חיוב דירה", 12, False),
    ]
    header_font = Font(bold=True, size=11, color="FFFFFF")
    header_fill_req = PatternFill(start_color="166534", end_color="166534", fill_type="solid")
    header_fill_opt = PatternFill(start_color="1D4ED8", end_color="1D4ED8", fill_type="solid")
    thin_border = Border(
        left=Side(style="thin", color="CBD5E1"),
        right=Side(style="thin", color="CBD5E1"),
        top=Side(style="thin", color="CBD5E1"),
        bottom=Side(style="thin", color="CBD5E1"),
    )
    for col_idx, (label, width, required) in enumerate(headers, 1):
        cell = ws.cell(row=3, column=col_idx, value=label)
        cell.font = header_font
        cell.fill = header_fill_req if required else header_fill_opt
        cell.alignment = Alignment(horizontal="center", vertical="center")
        cell.border = thin_border
        ws.column_dimensions[get_column_letter(col_idx)].width = width
    ws.row_dimensions[3].height = 28

    # ── Example rows ──
    example_data = [
        ["אצל לקוח", 55, "יוסי כהן", "רחוב הרצל 10, תל אביב", "דנה לוי", "050-1234567", "ישראל ישראלי", "12345678", 42, 800],
        ["", "", "", "", "", "", "מוחמד אחמד", "87654321", 42, 800],
        ["בתפקיד", 48, "שרה מזרחי", "שדרות רוטשילד 5, ת\"א", "אבי רון", "052-9876543", "פרסה רוגלי", "11223344", 38, 0],
        ["דייל", 60, "", "אילת, מלון רויאל", "משה דוד", "054-5551234", "דוד דוד", "55667788", 45, 600],
    ]
    example_font = Font(size=11, color="334155")
    alt_fill_1 = PatternFill(start_color="F8FAFC", end_color="F8FAFC", fill_type="solid")
    alt_fill_2 = PatternFill(start_color="FFFFFF", end_color="FFFFFF", fill_type="solid")
    for row_idx, row_data in enumerate(example_data, 4):
        fill = alt_fill_1 if row_idx % 2 == 0 else alt_fill_2
        for col_idx, val in enumerate(row_data, 1):
            cell = ws.cell(row=row_idx, column=col_idx, value=val)
            cell.font = example_font
            cell.alignment = Alignment(horizontal="right" if isinstance(val, str) else "center", vertical="center")
            cell.border = thin_border
            cell.fill = fill

    # ── Legend section ──
    legend_row = len(example_data) + 5
    ws.merge_cells(start_row=legend_row, start_column=1, end_row=legend_row, end_column=6)
    legend_title = ws.cell(row=legend_row, column=1, value="הסברים")
    legend_title.font = Font(bold=True, size=12, color="1E3A8A")
    legend_title.alignment = Alignment(horizontal="right")
    legend_items = [
        ("שם העסק *", "שם הלקוח כפי שמופיע בעמודת התנועה המיוחדת בדוח הנוכחות. חובה."),
        ("תעריף גביה *", "הסכום לשעה שנגבה מהלקוח. חובה לחישוב חיוב."),
        ("מנהל אזור", "שם מנהל האזור האחראי על הלקוח."),
        ("כתובת העסק", "כתובת אתר העבודה של הלקוח."),
        ("איש קשר / טלפון", "פרטי קשר של הלקוח להצגה בדוח."),
        ("שם העובד", "שם העובד שמוצב אצל הלקוח. ניתן להוסיף מספר שורות לאותו לקוח."),
        ("שכר שעתי / חיוב דירה", "פרטי שכר והטבות ספציפיים לעובד (אופציונלי — ניתן להגדיר גם במסך הראשי)."),
    ]
    info_font = Font(size=10, color="475569")
    label_font = Font(size=10, color="1E3A8A", bold=True)
    for i, (label, desc) in enumerate(legend_items, 1):
        r = legend_row + i
        cell_label = ws.cell(row=r, column=1, value=label)
        cell_label.font = label_font
        cell_label.alignment = Alignment(horizontal="right")
        ws.merge_cells(start_row=r, start_column=2, end_row=r, end_column=6)
        cell_desc = ws.cell(row=r, column=2, value=desc)
        cell_desc.font = info_font
        cell_desc.alignment = Alignment(horizontal="right")

    # Save to temp file
    import tempfile
    tmp = tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False, dir=str(OUTPUT_FOLDER))
    wb.save(tmp.name)
    tmp.close()
    return send_file(tmp.name, as_attachment=True, download_name="קובץ_לקוחות_לדוגמה.xlsx")


@app.route("/report-jobs/<int:job_id>/download")
@login_required
def download_report_job(job_id):
    expire_report_jobs()
    with get_db() as db:
        job = db.execute("SELECT * FROM report_jobs WHERE id=?", (job_id,)).fetchone()
    if not job or job["user_id"] != session.get("user_id"):
        add_flash("הדוח לא נמצא")
        return redirect("/dashboard")
    if job["status"] != "ready":
        add_flash("הדוח עדיין לא מוכן להורדה")
        return redirect("/dashboard")
    path = OUTPUT_FOLDER / job["output_filename"]
    if not path.exists():
        with get_db() as db:
            db.execute("UPDATE report_jobs SET status='expired', status_note=? WHERE id=?", ("הדוח אינו זמין יותר להורדה.", job_id))
            db.commit()
        add_flash("הדוח אינו זמין יותר להורדה")
        return redirect("/dashboard")
    data = path.read_bytes()
    cleanup_report_file(path)
    with get_db() as db:
        db.execute(
            "UPDATE report_jobs SET status='downloaded', downloaded_at=?, status_note=? WHERE id=?",
            (now_text(), "הדוח הורד ונמחק אוטומטית מהמערכת.", job_id),
        )
        db.commit()
    log_user_activity("download_report", "הוריד דוח", job["script_id"], job["script_name"], job["output_filename"])
    return send_file(BytesIO(data), as_attachment=True, download_name=job["output_filename"].split("_", 1)[-1] if "_" in job["output_filename"] else job["output_filename"])


@app.route("/activity", methods=["POST"])
@login_required
def activity():
    if session.get("is_admin"):
        return "", 204
    event_type = request.form.get("event_type", "").strip()
    action_label = request.form.get("action_label", "").strip()
    if event_type and action_label:
        log_user_activity(
            event_type,
            action_label,
            request.form.get("script_id", "").strip(),
            request.form.get("script_name", "").strip(),
            request.form.get("details", "").strip(),
        )
    return "", 204


@app.route("/admin")
@login_required
@admin_required
def admin():
    activity_user_id = request.args.get("activity_user_id", "").strip()
    activity_range = request.args.get("activity_range", "all").strip() or "all"
    activity_from = request.args.get("activity_from", "").strip()
    activity_to = request.args.get("activity_to", "").strip()
    activity_event = request.args.get("activity_event", "all").strip() or "all"
    activity_limit_raw = request.args.get("activity_limit", "50").strip() or "50"
    activity_limit = 50
    if activity_limit_raw.isdigit():
        activity_limit = max(50, min(500, int(activity_limit_raw)))
    with get_db() as db:
        users = db.execute("SELECT * FROM users WHERE is_admin=0").fetchall()
        perms = db.execute("SELECT * FROM permissions").fetchall()
        all_activity_logs = db.execute("SELECT * FROM activity_logs ORDER BY created_at DESC, id DESC").fetchall()
        support_requests = db.execute("SELECT * FROM support_requests ORDER BY created_at DESC, id DESC").fetchall()
        # Tool requests (developer briefs)
        tool_requests_rows = db.execute("SELECT * FROM tool_requests ORDER BY created_at DESC, id DESC").fetchall()
        pending_tool_requests = sum(1 for r in tool_requests_rows if r["status"] == "pending")
        # Token usage per user (cost monitoring)
        token_usage_by_user = db.execute(
            """SELECT user_id, SUM(input_tokens) as total_input, SUM(output_tokens) as total_output,
               SUM(estimated_cost_usd) as total_cost, COUNT(*) as request_count
               FROM chat_token_usage GROUP BY user_id ORDER BY total_cost DESC"""
        ).fetchall()
        total_platform_cost = sum(row["total_cost"] or 0 for row in token_usage_by_user)

    user_perms = {}
    for perm in perms:
        user_perms.setdefault(perm["user_id"], set()).add(perm["script_id"])

    active_customers = 0
    trial_customers = 0
    inactive_customers = 0
    user_cards = ""
    for user in users:
        uid = user["id"]
        status = get_account_status(user)
        if status["status_key"] == "active":
            active_customers += 1
        elif status["status_key"] == "trial":
            trial_customers += 1
        else:
            inactive_customers += 1
        service_style = {
            "active": ("#ecfdf5", "#047857"),
            "trial": ("#fff7ed", "#c2410c"),
            "expired": ("#fef2f2", "#b91c1c"),
            "unknown": ("#f8fafc", "#475569"),
        }.get(status["status_key"], ("#f8fafc", "#475569"))
        checks = ""
        for sid, script in SCRIPTS.items():
            checked = "checked" if (uid in user_perms and sid in user_perms[uid]) else ""
            checks += (
                '<label>'
                '<input type="checkbox" name="scripts" value="' + sid + '" ' + checked + ">"
                + script["icon"]
                + " "
                + script["name"]
                + "</label>"
            )
        is_active = bool(user["active"])
        active_color = "#047857" if is_active else "#b91c1c"
        active_bg = "#ecfdf5" if is_active else "#fef2f2"
        active_label = "פעיל" if is_active else "מושבת"
        trial_extend_html = ""
        if status["status_key"] == "trial":
            trial_extend_html = (
                '<div class="admin-user-section"><div class="admin-user-section-title">הארכת ניסיון</div>'
                '<form method="POST" action="/admin/extend_trial/' + str(uid) + '" style="display:flex;gap:8px;align-items:center;flex-wrap:wrap">'
                '<select name="extend_days" style="padding:6px 10px;border:1.5px solid #e2e8f0;border-radius:8px;font-size:13px;font-family:inherit">'
                '<option value="30">30 יום</option><option value="60">60 יום</option><option value="90">90 יום</option></select>'
                '<button type="submit" class="btn btn-blue" style="font-size:12px;padding:6px 14px">הארכה</button>'
                '<span style="font-size:12px;color:#64748b">תקופת ניסיון נוכחית: ' + str(user["trial_days"] or 30) + ' יום</span>'
                '</form></div>'
            )
        current_account_type = "active" if user["service_valid_until"] else "trial"
        user_cards += (
            '<details class="admin-user-card">'
            '<summary class="admin-collapsible-summary">'
            '<div><div class="admin-user-title">' + esc(user["company_name"] or user["full_name"] or user["username"]) + '</div>'
            '<div class="admin-user-sub">@' + esc(user["username"]) + ' • ח.פ: ' + esc(user["company_id"] or "לא הוגדר") + '</div>'
            '<div class="admin-user-sub">' + esc(user["full_name"] or "לא הוגדר") + ' • ' + esc(user["email"] or "ללא אימייל") + '</div></div>'
            '<div style="display:flex;align-items:center;gap:10px;flex-wrap:wrap">'
            '<span style="display:inline-flex;align-items:center;padding:5px 10px;border-radius:999px;background:' + active_bg + ';color:' + active_color + ';font-size:11px;font-weight:800">' + active_label + '</span>'
            '<span class="admin-user-status" style="background:' + service_style[0] + ';color:' + service_style[1] + '">' + esc(status["status_label_he"]) + '</span>'
            '<span style="font-size:18px;color:#64748b">+</span></div>'
            '</summary>'
            # Active toggle
            '<div style="padding:12px 20px 0;display:flex;align-items:center;gap:10px;flex-wrap:wrap">'
            '<form method="POST" action="/admin/toggle_active/' + str(uid) + '" style="display:inline">'
            '<label style="display:inline-flex;align-items:center;gap:8px;cursor:pointer;font-size:13px;font-weight:700;color:#334155">'
            '<span>גישה למערכת:</span>'
            '<span style="position:relative;display:inline-block;width:44px;height:24px">'
            '<input type="checkbox" onchange="this.form.submit()" ' + ('checked ' if is_active else '') + 'style="opacity:0;width:0;height:0;position:absolute">'
            '<span style="position:absolute;top:0;left:0;right:0;bottom:0;background:' + ('#047857' if is_active else '#cbd5e1') + ';border-radius:24px;transition:background .2s"></span>'
            '<span style="position:absolute;top:2px;' + ('left:22px' if is_active else 'left:2px') + ';width:20px;height:20px;background:white;border-radius:50%;transition:left .2s;box-shadow:0 1px 3px rgba(0,0,0,.2)"></span>'
            '</span>'
            '<span style="font-size:12px;color:' + active_color + '">' + active_label + '</span>'
            '</label></form></div>'
            # Meta boxes
            '<div class="admin-user-meta">'
            '<div class="admin-user-meta-box"><div class="k">איש קשר</div><div class="v">' + esc(user["full_name"] or "לא הוגדר") + '</div></div>'
            '<div class="admin-user-meta-box"><div class="k">מסלול חיוב</div><div class="v">' + esc(billing_mode_label(user["billing_mode"], "he")) + '</div></div>'
            '<div class="admin-user-meta-box"><div class="k">אימייל</div><div class="v">' + esc(user["email"] or "ללא אימייל") + '</div></div>'
            '<div class="admin-user-meta-box"><div class="k">טלפון</div><div class="v">' + esc(user["phone"] or "ללא טלפון") + '</div></div>'
            '<div class="admin-user-meta-box"><div class="k">תאריך הצטרפות</div><div class="v">' + esc(format_ui_date(user["join_date"], "he")) + '</div></div>'
            '<div class="admin-user-meta-box"><div class="k">שירות</div><div class="v">'
            + esc(status["status_label_he"])
            + ('<br><span style="font-size:12px;font-weight:600;color:#64748b">נותרו ' + str(status["days_remaining"]) + ' ימים</span>' if status["days_remaining"] is not None else "")
            + ('<br><span style="font-size:12px;font-weight:600;color:#64748b">בתוקף עד ' + esc(format_ui_date(status["renewal_date"], "he")) + '</span>' if status["renewal_date"] else "")
            + '</div></div>'
            '</div>'
            # Edit form
            '<div class="admin-user-section"><div class="admin-user-section-title">עריכת פרטים</div>'
            '<form method="POST" action="/admin/edit_user/' + str(uid) + '">'
            '<div class="form-row">'
            '<div class="form-group"><label class="field-label">שם מלא</label><input type="text" name="full_name" value="' + esc(user["full_name"] or "") + '" style="margin-bottom:0"></div>'
            '<div class="form-group"><label class="field-label">שם חברה</label><input type="text" name="company_name" value="' + esc(user["company_name"] or "") + '" style="margin-bottom:0"></div>'
            '<div class="form-group"><label class="field-label">ח.פ / מזהה חברה</label><input type="text" name="company_id" value="' + esc(user["company_id"] or "") + '" style="margin-bottom:0"></div>'
            '</div><div class="form-row">'
            '<div class="form-group"><label class="field-label">אימייל</label><input type="text" name="email" value="' + esc(user["email"] or "") + '" style="margin-bottom:0"></div>'
            '<div class="form-group"><label class="field-label">טלפון</label><input type="text" name="phone" value="' + esc(user["phone"] or "") + '" style="margin-bottom:0"></div>'
            '<div class="form-group"><label class="field-label">מסלול חיוב</label><select name="billing_mode" style="margin-bottom:0"><option value="monthly"' + (' selected' if (user["billing_mode"] or "monthly") == "monthly" else '') + '>חודשי</option><option value="yearly_prepaid"' + (' selected' if user["billing_mode"] == "yearly_prepaid" else '') + '>שנתי מראש</option></select></div>'
            '</div><div class="form-row">'
            '<div class="form-group"><label class="field-label">סוג חשבון</label><select name="account_type" style="margin-bottom:0"><option value="trial"' + (' selected' if current_account_type == "trial" else '') + '>תקופת ניסיון</option><option value="active"' + (' selected' if current_account_type == "active" else '') + '>שירות פעיל</option></select></div>'
            '<div class="form-group"><label class="field-label">תחילת ניסיון</label><input type="text" name="trial_start_date" value="' + esc(user["trial_start_date"] or "") + '" placeholder="YYYY-MM-DD" style="margin-bottom:0"></div>'
            '<div class="form-group"><label class="field-label">בתוקף עד</label><input type="text" name="service_valid_until" value="' + esc(user["service_valid_until"] or "") + '" placeholder="YYYY-MM-DD" style="margin-bottom:0"></div>'
            '<button type="submit" class="btn btn-blue" style="height:40px;align-self:flex-end;font-size:12px;padding:6px 14px">שמירת שינויים</button>'
            '</div></form></div>'
            # Trial extend
            + trial_extend_html
            # Permissions
            + '<div class="admin-user-section"><div class="admin-user-section-title">כלים והרשאות</div>'
            '<form method="POST" action="/admin/permissions/' + str(uid) + '"><div class="admin-user-perms">'
            + checks
            + '</div><button type="submit" class="btn btn-gray" style="margin-top:10px;font-size:12px;padding:6px 14px">שמירת הרשאות</button></form></div>'
            '<div class="admin-user-section"><div class="admin-user-section-title">פעולות ניהול</div><div class="admin-user-actions">'
            + '<button type="button" class="btn btn-gray" style="font-size:12px;padding:6px 14px" onclick="openPass(' + str(uid) + ',' + json.dumps(user["full_name"] or user["company_name"] or user["username"] or "") + ')">שינוי סיסמה</button>'
            + '<form method="POST" action="/admin/resetpass/' + str(uid) + '" style="display:inline"><button type="submit" class="btn btn-gray" style="font-size:12px;padding:6px 14px">סיסמה זמנית</button></form>'
            + '<a href="/admin/delete/' + str(uid) + '" onclick="return confirm(\'האם למחוק את הלקוח הזה?\');" class="btn btn-red" style="text-decoration:none;font-size:12px;padding:6px 14px">מחיקה</a>'
            + '</div></div>'
            '</details>'
        )

    users_overview = (
        '<div class="admin-user-summary">'
        '<div class="admin-user-summary-box"><div class="k">סה"כ לקוחות</div><div class="v">' + str(len(users)) + '</div></div>'
        '<div class="admin-user-summary-box"><div class="k">בשירות פעיל</div><div class="v">' + str(active_customers) + '</div></div>'
        '<div class="admin-user-summary-box"><div class="k">בתקופת ניסיון</div><div class="v">' + str(trial_customers) + '</div></div>'
        '<div class="admin-user-summary-box"><div class="k">לא בשירות</div><div class="v">' + str(inactive_customers) + '</div></div>'
        '</div>'
    ) if users else ""
    table = (users_overview + '<div class="admin-user-grid">' + user_cards + '</div>') if users else '<p style="color:#94a3b8;text-align:center;padding:2rem">עדיין אין לקוחות במערכת</p>'

    user_lookup = {str(user["id"]): user for user in users}
    customer_options = ""
    for user in sorted(users, key=lambda item: ((item["company_name"] or item["full_name"] or item["username"] or "").lower(), str(item["id"]))):
        option_label = user["company_name"] or user["full_name"] or user["username"]
        if user["company_name"] and user["full_name"]:
            option_label += " — " + user["full_name"]
        option_label += " @" + (user["username"] or "")
        customer_options += '<option value="' + str(user["id"]) + '"' + (' selected' if activity_user_id == str(user["id"]) else '') + '>' + esc(option_label) + '</option>'

    filtered_activity_logs = []
    today = date.today()
    range_start = None
    range_end = None
    if activity_range == "last_30":
        range_start = datetime.combine(today - timedelta(days=29), datetime.min.time())
        range_end = datetime.combine(today, datetime.max.time())
    elif activity_range == "custom":
        parsed_from = parse_iso_date(activity_from)
        parsed_to = parse_iso_date(activity_to)
        if parsed_from:
            range_start = datetime.combine(parsed_from, datetime.min.time())
        if parsed_to:
            range_end = datetime.combine(parsed_to, datetime.max.time())

    for entry in all_activity_logs:
        if activity_user_id and str(entry["user_id"] or "") != activity_user_id:
            continue
        if activity_event != "all" and str(entry["event_type"] or "") != activity_event:
            continue
        created_at = parse_datetime_value(entry["created_at"])
        if range_start and created_at and created_at < range_start:
            continue
        if range_end and created_at and created_at > range_end:
            continue
        if range_start and not created_at:
            continue
        if range_end and not created_at:
            continue
        filtered_activity_logs.append(entry)

    total_activity = len(filtered_activity_logs)
    generated_reports = sum(1 for entry in filtered_activity_logs if entry["event_type"] == "generate_report")
    opened_tools = sum(1 for entry in filtered_activity_logs if entry["event_type"] == "open_script")
    help_opens = sum(1 for entry in filtered_activity_logs if entry["event_type"] == "open_help_popup")
    terms_opens = sum(1 for entry in filtered_activity_logs if entry["event_type"] == "open_service_terms")
    unique_users_count = len({entry["user_id"] for entry in filtered_activity_logs if entry["user_id"]})

    displayed_activity_logs = filtered_activity_logs[:activity_limit]

    def build_summary_card(label, value, event_type):
        return (
            '<button type="button" data-activity-event="' + esc(event_type) + '" style="background:#f8fafc;border:1px solid ' + ("#93c5fd" if activity_event == event_type or (event_type == "all" and activity_event == "all") else "#e2e8f0") + ';border-radius:12px;padding:12px;text-decoration:none;display:block;width:100%;text-align:right;font-family:inherit;cursor:pointer">'
            + '<div style="font-size:12px;color:#64748b;margin-bottom:6px">' + label + '</div>'
            + '<div style="font-size:20px;font-weight:800;color:#0f172a">' + str(value) + '</div>'
            + '</button>'
        )

    activity_rows = ""
    for entry in displayed_activity_logs:
        user_label = entry["full_name"] or entry["username"] or ("משתמש #" + str(entry["user_id"]))
        activity_rows += (
            "<tr>"
            '<td>' + esc(format_ui_datetime(entry["created_at"])) + "</td>"
            '<td><div style="font-weight:700;color:#0f172a">' + esc(user_label) + '</div><div style="font-size:12px;color:#64748b">@' + esc(entry["username"] or "") + "</div></td>"
            '<td>' + esc(entry["action_label"]) + "</td>"
            '<td>' + esc(entry["script_name"] or "ללא כלי") + "</td>"
            '<td>' + esc(entry["details"] or "—") + "</td>"
            "</tr>"
        )
    activity_table_inner = (
        "<table><thead><tr><th>מתי</th><th>משתמש</th><th>פעולה</th><th>כלי</th><th>פרטים</th></tr></thead><tbody>"
        + activity_rows
        + "</tbody></table>"
    ) if displayed_activity_logs else '<p style="color:#94a3b8;text-align:center;padding:2rem">אין לוגים שתואמים את הסינון הנוכחי</p>'
    activity_table = '<div id="activityTableWrap">' + activity_table_inner + '</div>'
    activity_summary = (
        '<div style="display:grid;grid-template-columns:repeat(auto-fit,minmax(160px,1fr));gap:10px;margin-bottom:1rem">'
        + build_summary_card("סה\"כ אירועים", total_activity, "all")
        + '<div id="uniqueUsersCard" style="background:#f8fafc;border:1px solid #e2e8f0;border-radius:12px;padding:12px;text-align:right">'
          '<div style="font-size:12px;color:#64748b;margin-bottom:6px">משתמשים שונים</div>'
          '<div id="uniqueUsersValue" style="font-size:20px;font-weight:800;color:#0f172a">' + str(unique_users_count) + '</div></div>'
        + build_summary_card("כלים שנפתחו", opened_tools, "open_script")
        + build_summary_card("דוחות שהופקו", generated_reports, "generate_report")
        + build_summary_card("חלונות מידע שנפתחו", help_opens, "open_help_popup")
        + build_summary_card("פירוט שירות ומחיר", terms_opens, "open_service_terms")
        + '</div>'
    )
    activity_filter_bar = (
        '<form method="GET" action="/admin" id="activityFilterForm" style="display:grid;grid-template-columns:repeat(auto-fit,minmax(160px,1fr));gap:10px;margin-bottom:1rem">'
        '<div><label class="field-label">לקוח</label><select id="activityUserId" name="activity_user_id" style="padding:9px 12px;border:1.5px solid #e2e8f0;border-radius:8px;font-size:13px;font-family:inherit;outline:none;width:100%;margin-bottom:0;background:white"><option value="">כל הלקוחות</option>' + customer_options + '</select></div>'
        '<div><label class="field-label">טווח תאריכים</label><select id="activityRange" name="activity_range" style="padding:9px 12px;border:1.5px solid #e2e8f0;border-radius:8px;font-size:13px;font-family:inherit;outline:none;width:100%;margin-bottom:0;background:white">'
        + '<option value="all"' + (' selected' if activity_range == "all" else '') + '>כל התקופה</option>'
        + '<option value="last_30"' + (' selected' if activity_range == "last_30" else '') + '>30 הימים האחרונים</option>'
        + '<option value="custom"' + (' selected' if activity_range == "custom" else '') + '>מתאריך עד תאריך</option>'
        + '</select></div>'
        '<div><label class="field-label">מתאריך</label><input type="text" id="activityFrom" name="activity_from" value="' + esc(activity_from) + '" placeholder="YYYY-MM-DD" style="margin-bottom:0"></div>'
        '<div><label class="field-label">עד תאריך</label><input type="text" id="activityTo" name="activity_to" value="' + esc(activity_to) + '" placeholder="YYYY-MM-DD" style="margin-bottom:0"></div>'
        '<div><label class="field-label">כמות שורות להצגה</label><select id="activityLimit" name="activity_limit" style="padding:9px 12px;border:1.5px solid #e2e8f0;border-radius:8px;font-size:13px;font-family:inherit;outline:none;width:100%;margin-bottom:0;background:white">'
        + '<option value="50"' + (' selected' if activity_limit == 50 else '') + '>50</option>'
        + '<option value="100"' + (' selected' if activity_limit == 100 else '') + '>100</option>'
        + '<option value="200"' + (' selected' if activity_limit == 200 else '') + '>200</option>'
        + '<option value="500"' + (' selected' if activity_limit == 500 else '') + '>500</option>'
        + '</select></div>'
        '<input type="hidden" id="activityEvent" name="activity_event" value="' + esc(activity_event) + '">'
        '<div style="display:flex;gap:8px;align-items:flex-end"><button type="submit" class="btn btn-blue" style="height:40px">סינון</button><button type="button" id="activityReset" class="btn btn-gray" style="height:40px">איפוס</button></div>'
        '</form>'
        '<div style="font-size:12px;color:#64748b;margin-bottom:1rem">הלוגים זמינים מרגע שההקלטה הופעלה. פעולות ישנות יותר לא ניתנות לשחזור רטרואקטיבית.</div>'
        + ('<div id="activityShowing" style="font-size:12px;color:#1d4ed8;margin-bottom:1rem">מוצגות ' + str(len(displayed_activity_logs)) + ' מתוך ' + str(total_activity) + ' רשומות תואמות.</div>' if total_activity > activity_limit else '<div id="activityShowing" style="font-size:12px;color:#1d4ed8;margin-bottom:1rem"></div>')
    )

    support_rows = ""
    pending_support = 0
    for entry in support_requests:
        meta = support_status_meta(entry["status"])
        request_type_label = "בקשה לכלי חדש" if entry["request_type"] == "new_tool" else "תמיכה בכלי קיים"
        customer_label = entry["company_name"] or entry["full_name"] or entry["username"] or ("משתמש #" + str(entry["user_id"]))
        if str(entry["status"] or "pending").strip().lower() not in {"accepted", "resolved"}:
            pending_support += 1
        contact_bits = []
        if entry["email"]:
            contact_bits.append("מייל: " + entry["email"])
        if entry["phone"]:
            contact_bits.append("טלפון: " + entry["phone"])
        contact_text = "<br>".join(esc(bit) for bit in contact_bits) if contact_bits else "—"
        support_rows += (
            '<details class="support-request-card">'
            '<summary class="admin-collapsible-summary" style="padding:0 0 12px 0">'
            '<div><div style="font-size:17px;font-weight:800;color:#0f172a;margin-bottom:4px">' + esc(customer_label) + '</div>'
            + ('<div style="font-size:12px;color:#64748b">@' + esc(entry["username"] or "") + '</div>' if entry["username"] else '')
            + '<div class="admin-collapsible-sub">' + esc(request_type_label) + ' • ' + esc(format_ui_datetime(entry["created_at"])) + (' • ' + esc(entry["script_name"]) if entry["script_name"] else '') + '</div>'
            + '</div>'
            '<div style="display:flex;align-items:center;gap:10px;flex-wrap:wrap"><span style="display:inline-flex;align-items:center;padding:7px 12px;border-radius:999px;background:' + meta["bg"] + ';color:' + meta["fg"] + ';font-size:12px;font-weight:800">' + esc(meta["label"]) + '</span><span style="font-size:18px;color:#64748b">+</span></div>'
            '</summary>'
            '<div class="support-request-card-meta">'
            '<div class="support-request-card-box"><div class="k">מועד פתיחה</div><div class="v">' + esc(format_ui_datetime(entry["created_at"])) + '</div></div>'
            '<div class="support-request-card-box"><div class="k">סוג פנייה</div><div class="v">' + esc(request_type_label) + '</div></div>'
            '<div class="support-request-card-box"><div class="k">כלי</div><div class="v">' + esc(entry["script_name"] or "—") + '</div></div>'
            '<div class="support-request-card-box"><div class="k">פרטי קשר</div><div class="v">' + contact_text + '</div></div>'
            '</div>'
            '<div class="support-request-message">' + esc(entry["message"] or "") + '</div>'
            '<div style="display:flex;gap:6px;flex-wrap:wrap">'
            + ('<form method="POST" action="/admin/support/' + str(entry["id"]) + '/status" style="display:inline"><input type="hidden" name="status" value="accepted"><button type="submit" class="btn btn-gray" style="font-size:12px;padding:5px 12px">התקבל</button></form>' if str(entry["status"] or "pending").strip().lower() != "accepted" else "")
            + ('<form method="POST" action="/admin/support/' + str(entry["id"]) + '/status" style="display:inline"><input type="hidden" name="status" value="resolved"><button type="submit" class="btn btn-blue" style="font-size:12px;padding:5px 12px">טופל</button></form>' if str(entry["status"] or "pending").strip().lower() != "resolved" else "")
            + '</div>'
            '</details>'
        )
    support_table = ('<div class="support-request-list">' + support_rows + '</div>') if support_requests else '<p style="color:#94a3b8;text-align:center;padding:2rem">עדיין אין פניות שירות מצד לקוחות</p>'

    # Build tool requests table
    tool_requests_html = ""
    if tool_requests_rows:
        req_status_meta = {
            "pending": ("ממתין לפיתוח", "#fff7ed", "#c2410c"),
            "in_progress": ("בפיתוח", "#eff6ff", "#1d4ed8"),
            "completed": ("הושלם", "#ecfdf5", "#047857"),
            "cancelled": ("בוטל", "#fef2f2", "#b91c1c"),
        }
        for req in tool_requests_rows:
            s_label, s_bg, s_fg = req_status_meta.get(req["status"], ("לא ידוע", "#f8fafc", "#475569"))
            req_date = format_ui_datetime(req["created_at"]) if req["created_at"] else ""
            req_type_label = "כלי חדש" if req["request_type"] == "new" else "שיפור כלי"
            tool_requests_html += (
                '<details class="support-request-card" style="margin-bottom:10px">'
                '<summary style="display:flex;justify-content:space-between;align-items:center;cursor:pointer;padding:12px 16px">'
                '<div style="flex:1">'
                '<div style="font-weight:700;color:#0f172a;font-size:14px">&#128736; ' + esc(req["tool_name"] or "") + '</div>'
                '<div style="font-size:12px;color:#64748b;margin-top:2px">'
                + esc(req["full_name"] or req["username"] or "") + ' &middot; ' + esc(req_type_label) + ' &middot; ' + esc(req_date) + '</div></div>'
                '<span style="font-size:11px;padding:3px 10px;border-radius:99px;font-weight:700;background:' + s_bg + ';color:' + s_fg + '">' + s_label + '</span>'
                '</summary>'
                '<div style="padding:0 16px 16px;border-top:1px solid #e2e8f0">'
                '<div style="margin-top:12px;background:#f8fafc;border:1px solid #e2e8f0;border-radius:10px;padding:14px;font-size:13px;line-height:1.8;white-space:pre-wrap;direction:rtl">'
                + esc(req["brief_text"] or "") + '</div>'
                '<div style="margin-top:10px;display:flex;gap:8px;flex-wrap:wrap">'
            )
            if req["status"] == "pending":
                tool_requests_html += (
                    '<form method="POST" action="/admin/tool-request/' + str(req["id"]) + '/status">'
                    '<input type="hidden" name="status" value="in_progress">'
                    '<button type="submit" class="btn btn-blue" style="font-size:12px;border-radius:8px;padding:6px 14px">&#9997; התחל פיתוח</button></form>'
                    '<form method="POST" action="/admin/tool-request/' + str(req["id"]) + '/status">'
                    '<input type="hidden" name="status" value="cancelled">'
                    '<button type="submit" class="btn btn-gray" style="font-size:12px;border-radius:8px;padding:6px 14px">&#10005; בטל</button></form>'
                )
            elif req["status"] == "in_progress":
                tool_requests_html += (
                    '<form method="POST" action="/admin/tool-request/' + str(req["id"]) + '/status">'
                    '<input type="hidden" name="status" value="completed">'
                    '<button type="submit" class="btn btn-blue" style="font-size:12px;border-radius:8px;padding:6px 14px;background:#047857">&#9989; סמן כהושלם</button></form>'
                )
            tool_requests_html += '</div></div></details>'
    else:
        tool_requests_html = '<p style="color:#94a3b8;text-align:center;padding:2rem">אין בקשות כלים כרגע</p>'

    # Build cost monitoring table
    # Map user_ids to names
    user_name_map = {u["id"]: (u["full_name"] or u["username"] or f"#{u['id']}") for u in users}
    user_name_map[1] = "מנהל מערכת"  # admin
    cost_table_html = ""
    if token_usage_by_user:
        cost_table_html += '<table><thead><tr><th>לקוח</th><th>בקשות</th><th>טוקנים (קלט)</th><th>טוקנים (פלט)</th><th>עלות ($)</th></tr></thead><tbody>'
        for row in token_usage_by_user:
            uname = user_name_map.get(row["user_id"], f"#{row['user_id']}")
            cost_table_html += (
                '<tr>'
                '<td style="font-weight:700">' + esc(uname) + '</td>'
                '<td>' + str(row["request_count"] or 0) + '</td>'
                '<td>' + f'{row["total_input"] or 0:,}' + '</td>'
                '<td>' + f'{row["total_output"] or 0:,}' + '</td>'
                '<td style="font-weight:700;color:#b91c1c">$' + f'{row["total_cost"] or 0:.4f}' + '</td>'
                '</tr>'
            )
        cost_table_html += '</tbody></table>'
        cost_table_html += '<div style="text-align:left;margin-top:12px;font-size:14px;font-weight:800;color:#b91c1c">סה"כ עלות פלטפורמה: $' + f'{total_platform_cost:.4f}' + '</div>'
    else:
        cost_table_html = '<p style="color:#94a3b8;text-align:center;padding:2rem">אין נתוני שימוש עדיין</p>'

    activity_logs_payload = json.dumps(
        [
            {
                "user_id": str(entry["user_id"] or ""),
                "username": entry["username"] or "",
                "full_name": entry["full_name"] or "",
                "event_type": entry["event_type"] or "",
                "action_label": entry["action_label"] or "",
                "script_name": entry["script_name"] or "",
                "details": entry["details"] or "",
                "created_at": entry["created_at"] or "",
                "display_when": format_ui_datetime(entry["created_at"]),
            }
            for entry in all_activity_logs
        ],
        ensure_ascii=False,
    )

    admin_side_nav = (
        '<div class="admin-float-nav">'
        '<a href="#adminAddUser" class="btn btn-gray" style="text-decoration:none;justify-content:center">הוספת לקוח</a>'
        '<a href="#adminUsers" class="btn btn-gray" style="text-decoration:none;justify-content:center">לקוחות</a>'
        '<a href="#adminSupport" class="btn btn-gray" style="text-decoration:none;justify-content:center">פניות שירות'
        + (' (' + str(pending_support) + ')' if pending_support else '')
        + '</a>'
        '<a href="#adminToolRequests" class="btn btn-gray" style="text-decoration:none;justify-content:center">בקשות כלים'
        + (' (' + str(pending_tool_requests) + ')' if pending_tool_requests else '')
        + '</a>'
        '<a href="#adminAICosts" class="btn btn-gray" style="text-decoration:none;justify-content:center">עלויות AI</a>'
        '<a href="#adminLogs" class="btn btn-gray" style="text-decoration:none;justify-content:center">לוגים</a>'
        '</div>'
    )

    body = (
        admin_side_nav
        +
        '<div class="card" id="adminAddUser"><h2>&#10133; הוספת לקוח חדש</h2><form method="POST" action="/admin/add_user"><div class="form-row">'
        '<div class="form-group"><label class="field-label">שם מלא</label><input type="text" name="full_name" placeholder="שם הלקוח" required style="margin-bottom:0"></div>'
        '<div class="form-group"><label class="field-label">שם חברה</label><input type="text" name="company_name" placeholder="שם החברה" style="margin-bottom:0"></div>'
        '<div class="form-group"><label class="field-label">ח.פ / מזהה חברה</label><input type="text" name="company_id" placeholder="ח.פ / מזהה חברה" style="margin-bottom:0"></div>'
        '<div class="form-group"><label class="field-label">שם משתמש</label><input type="text" name="username" placeholder="שם משתמש להתחברות" required style="margin-bottom:0"></div>'
        '<div class="form-group"><label class="field-label">סיסמה</label><input type="password" name="password" placeholder="סיסמה ראשונית" required style="margin-bottom:0"></div>'
        '</div><div class="form-row">'
        '<div class="form-group"><label class="field-label">אימייל</label><input type="text" name="email" placeholder="אימייל" style="margin-bottom:0"></div>'
        '<div class="form-group"><label class="field-label">טלפון</label><input type="text" name="phone" placeholder="טלפון" style="margin-bottom:0"></div>'
        '<div class="form-group"><label class="field-label">מסלול חיוב</label><select name="billing_mode" style="margin-bottom:0"><option value="monthly">חודשי</option><option value="yearly_prepaid">שנתי מראש</option></select></div>'
        '<div class="form-group"><label class="field-label">סוג חשבון</label><select name="account_type" style="margin-bottom:0"><option value="trial">תקופת ניסיון</option><option value="active">שירות פעיל</option></select></div>'
            '<div class="form-group"><label class="field-label">ימי ניסיון</label><select name="trial_days" style="margin-bottom:0"><option value="30">30 יום</option><option value="60">60 יום</option><option value="90">90 יום</option></select></div>'
        '<div class="form-group"><label class="field-label">בתוקף עד</label><input type="text" name="service_valid_until" placeholder="YYYY-MM-DD" style="margin-bottom:0"></div>'
        '<button type="submit" class="btn btn-blue" style="height:40px;align-self:flex-end">הוספה</button></div></form></div>'
        '<details class="card" id="adminUsers" style="padding:0;overflow:hidden" dir="rtl">'
        '<summary class="admin-collapsible-summary">'
        '<div><div style="font-size:22px;font-weight:800;color:#0f172a;margin-bottom:4px">&#128101; לקוחות במערכת</div><div class="admin-collapsible-sub">רשימת הלקוחות מקופלת כברירת מחדל. פתיחה לפי צורך לעבודה נוחה יותר.</div></div>'
        '<span style="font-size:18px;color:#64748b">+</span>'
        '</summary><div style="padding:0 20px 20px">'
        + table
        + '</div></details>'
        '<div class="card" style="background:linear-gradient(135deg,#f0fdf4 0%,#ecfdf5 100%);border:1px solid #86efac" dir="rtl">'
        '<div style="display:flex;align-items:center;justify-content:space-between;gap:14px;flex-wrap:wrap">'
        '<div><div style="font-size:18px;font-weight:800;color:#047857;margin-bottom:4px">🛒 ניהול שוק הכלים</div>'
        '<div style="font-size:13px;color:#475569">בדיקה ואישור כלים שנוצרו על ידי משתמשים</div></div>'
        '<a href="/admin/tools" class="btn btn-blue" style="text-decoration:none;display:inline-flex;align-items:center;justify-content:center;min-width:160px;background:#047857">ניהול כלים</a>'
        '</div></div>'
        '<details class="card" id="adminSupport" style="padding:0;overflow:hidden" dir="rtl">'
        '<summary class="admin-collapsible-summary">'
        '<div><div style="font-size:22px;font-weight:800;color:#0f172a;margin-bottom:4px">&#128172; פניות שירות לקוחות</div><div class="admin-collapsible-sub">'
        + ('יש ' + str(pending_support) + ' פניות שממתינות להתייחסות' if pending_support else 'אין כרגע פניות שממתינות להתייחסות')
        + '</div></div>'
        '<span style="font-size:18px;color:#64748b">+</span>'
        '</summary><div style="padding:0 20px 20px">'
        + support_table
        + '</div></details>'
        '<details class="card" id="adminToolRequests" style="padding:0;overflow:hidden" dir="rtl">'
        '<summary class="admin-collapsible-summary">'
        '<div><div style="font-size:22px;font-weight:800;color:#0f172a;margin-bottom:4px">&#128736; בקשות כלים</div><div class="admin-collapsible-sub">'
        + ('יש ' + str(pending_tool_requests) + ' בקשות שממתינות לפיתוח' if pending_tool_requests else 'אין כרגע בקשות ממתינות')
        + '</div></div>'
        '<span style="font-size:18px;color:#64748b">+</span>'
        '</summary><div style="padding:0 20px 20px">'
        + tool_requests_html
        + '</div></details>'
        '<details class="card" id="adminAICosts" style="padding:0;overflow:hidden" dir="rtl">'
        '<summary class="admin-collapsible-summary">'
        '<div><div style="font-size:22px;font-weight:800;color:#0f172a;margin-bottom:4px">&#128176; עלויות AI לפי לקוח</div><div class="admin-collapsible-sub">ניטור שימוש בטוקנים ועלויות API</div></div>'
        '<span style="font-size:18px;color:#64748b">+</span>'
        '</summary><div style="padding:0 20px 20px">'
        + cost_table_html
        + '</div></details>'
        '<details class="card" id="adminLogs" style="padding:0;overflow:hidden" dir="rtl">'
        '<summary class="admin-collapsible-summary">'
        '<div><div style="font-size:22px;font-weight:800;color:#0f172a;margin-bottom:4px">&#128221; לוג פעילות משתמשים</div><div style="font-size:13px;color:#64748b">פתיחה לפי צורך בלבד לצפייה ועבודה על הלוגים</div></div>'
        '<span style="font-size:18px;color:#64748b">+</span>'
        '</summary><div style="padding:0 20px 20px">'
        '<div id="activityPanel">'
        + activity_filter_bar
        + activity_summary
        + activity_table
        + '</div></div></details>'
        + '</div><div class="modal-bg" id="passModal"><div class="modal-box"><h3 style="font-size:15px;font-weight:700;margin-bottom:1rem;color:#1e3a8a">שינוי סיסמה &#8212; <span id="pname"></span></h3>'
        '<form method="POST" id="pform"><input type="password" name="new_password" placeholder="סיסמה חדשה" required>'
        '<div style="display:flex;gap:8px;margin-top:.5rem;justify-content:flex-end"><button type="button" class="btn btn-gray" onclick="closePass()">ביטול</button>'
        '<button type="submit" class="btn btn-blue">עדכון</button></div></form></div></div>'
        '<script>'
        'function openPass(id,name){document.getElementById("pname").textContent=name||"";document.getElementById("pform").action="/admin/setpass/"+id;document.getElementById("passModal").style.display="flex";}'
        'function closePass(){document.getElementById("passModal").style.display="none";}'
        '(function(){'
        'var allLogs=' + activity_logs_payload + ';'
        'var userSelect=document.getElementById("activityUserId");'
        'var rangeSelect=document.getElementById("activityRange");'
        'var fromInput=document.getElementById("activityFrom");'
        'var toInput=document.getElementById("activityTo");'
        'var limitSelect=document.getElementById("activityLimit");'
        'var eventInput=document.getElementById("activityEvent");'
        'var form=document.getElementById("activityFilterForm");'
        'var summaryButtons=document.querySelectorAll("[data-activity-event]");'
        'var panel=document.getElementById("activityPanel");'
        'function escHtml(value){return String(value||"").replace(/&/g,"&amp;").replace(/</g,"&lt;").replace(/>/g,"&gt;").replace(/"/g,"&quot;");}'
        'function parseDateOnly(text){if(!text){return null;}var m=String(text).match(/^(\\d{4})-(\\d{2})-(\\d{2})$/);if(!m){return null;}return new Date(Number(m[1]),Number(m[2])-1,Number(m[3]),0,0,0,0);}'
        'function parseCreatedAt(text){if(!text){return null;}var m=String(text).match(/^(\\d{4})-(\\d{2})-(\\d{2}) (\\d{2}):(\\d{2}):(\\d{2})$/);if(!m){return null;}return new Date(Number(m[1]),Number(m[2])-1,Number(m[3]),Number(m[4]),Number(m[5]),Number(m[6]),0);}'
        'function getFilteredLogs(){var userId=userSelect?userSelect.value:"";var range=rangeSelect?rangeSelect.value:"all";var fromDate=parseDateOnly(fromInput?fromInput.value:"");var toDate=parseDateOnly(toInput?toInput.value:"");var eventType=eventInput?eventInput.value:"all";var today=new Date();today.setHours(0,0,0,0);var last30Start=new Date(today.getTime());last30Start.setDate(today.getDate()-29);var last30End=new Date(today.getTime());last30End.setHours(23,59,59,999);return allLogs.filter(function(entry){if(userId&&entry.user_id!==userId){return false;}if(eventType&&eventType!=="all"&&entry.event_type!==eventType){return false;}var created=parseCreatedAt(entry.created_at);if(range==="last_30"){if(!created||created<last30Start||created>last30End){return false;}}else if(range==="custom"){if(fromDate&&(!created||created<fromDate)){return false;}if(toDate){var end=new Date(toDate.getTime());end.setHours(23,59,59,999);if(!created||created>end){return false;}}}return true;});}'
        'function countBy(logs,eventType){if(eventType==="all"){return logs.length;}var total=0;logs.forEach(function(entry){if(entry.event_type===eventType){total+=1;}});return total;}'
        'function countUniqueUsers(logs){var seen={};logs.forEach(function(entry){if(entry.user_id){seen[entry.user_id]=true;}});return Object.keys(seen).length;}'
        'function renderActivity(){if(!panel){return;}var filtered=getFilteredLogs();var limit=Math.max(50,Math.min(500,parseInt(limitSelect&&limitSelect.value||"50",10)||50));var displayed=filtered.slice(0,limit);var currentEvent=eventInput?eventInput.value:"all";summaryButtons.forEach(function(btn){var active=(btn.getAttribute("data-activity-event")||"all")===currentEvent;btn.style.borderColor=active?"#93c5fd":"#e2e8f0";var valueNode=btn.querySelectorAll("div")[1];if(valueNode){var eventKey=btn.getAttribute("data-activity-event")||"all";valueNode.textContent=String(countBy(filtered,eventKey));}});var rowsHtml="";if(displayed.length){displayed.forEach(function(entry){var userLabel=entry.full_name||entry.username||("משתמש #"+entry.user_id);rowsHtml+="<tr><td>"+escHtml(entry.display_when)+"</td><td><div style=\\"font-weight:700;color:#0f172a\\">"+escHtml(userLabel)+"</div><div style=\\"font-size:12px;color:#64748b\\">@"+escHtml(entry.username||"")+"</div></td><td>"+escHtml(entry.action_label||"")+"</td><td>"+escHtml(entry.script_name||"ללא כלי")+"</td><td>"+escHtml(entry.details||"—")+"</td></tr>";});rowsHtml="<table><thead><tr><th>מתי</th><th>משתמש</th><th>פעולה</th><th>כלי</th><th>פרטים</th></tr></thead><tbody>"+rowsHtml+"</tbody></table>";}else{rowsHtml=\'<p style="color:#94a3b8;text-align:center;padding:2rem">אין לוגים שתואמים את הסינון הנוכחי</p>\';}var showingNode=document.getElementById("activityShowing");if(showingNode){showingNode.textContent=filtered.length>limit?("מוצגות "+displayed.length+" מתוך "+filtered.length+" רשומות תואמות."): ""; }var existingTable=document.getElementById("activityTableWrap");if(existingTable){existingTable.innerHTML=rowsHtml;}else{var wrap=document.createElement("div");wrap.id="activityTableWrap";wrap.innerHTML=rowsHtml;panel.appendChild(wrap);}var uuNode=document.getElementById("uniqueUsersValue");if(uuNode){uuNode.textContent=String(countUniqueUsers(filtered));} }'
        'if(form){form.addEventListener("submit",function(ev){ev.preventDefault();renderActivity();});}'
        'summaryButtons.forEach(function(btn){btn.addEventListener("click",function(){if(eventInput){eventInput.value=this.getAttribute("data-activity-event")||"all";}renderActivity();});});'
        'var resetBtn=document.getElementById("activityReset");if(resetBtn){resetBtn.addEventListener("click",function(){if(userSelect){userSelect.value="";}if(rangeSelect){rangeSelect.value="all";}if(fromInput){fromInput.value="";}if(toInput){toInput.value="";}if(limitSelect){limitSelect.value="50";}if(eventInput){eventInput.value="all";}renderActivity();});}'
        '})();'
        '</script>'
    )
    return render("ניהול מערכת", body)


@app.route("/admin/add_user", methods=["POST"])
@login_required
@admin_required
def add_user():
    username = request.form["username"].strip()
    password = request.form["password"]
    full_name = request.form["full_name"].strip()
    company_name = request.form.get("company_name", "").strip()
    company_id = request.form.get("company_id", "").strip()
    email = request.form.get("email", "").strip()
    phone = request.form.get("phone", "").strip()
    billing_mode = request.form.get("billing_mode", "monthly").strip() or "monthly"
    account_type = request.form.get("account_type", "trial").strip() or "trial"
    service_valid_until = request.form.get("service_valid_until", "").strip()
    trial_days_raw = request.form.get("trial_days", "30").strip()
    trial_days = int(trial_days_raw) if trial_days_raw in ("30", "60", "90") else 30
    join_date = date.today().isoformat()
    trial_start_date = join_date if account_type == "trial" else ""
    service_until_value = service_valid_until if account_type == "active" else ""
    try:
        with get_db() as db:
            db.execute(
                """INSERT INTO users(
                username,password,full_name,company_name,company_id,email,phone,join_date,trial_start_date,service_valid_until,billing_mode,trial_days
                ) VALUES (?,?,?,?,?,?,?,?,?,?,?,?)""",
                (
                    username,
                    generate_password_hash(password),
                    full_name,
                    company_name,
                    company_id,
                    email,
                    phone,
                    join_date,
                    trial_start_date,
                    service_until_value,
                    billing_mode,
                    trial_days,
                ),
            )
            db.commit()
        add_flash("הלקוח " + full_name + " נוצר בהצלחה")
    except Exception as exc:
        if not is_integrity_error(exc):
            raise
        add_flash("שם המשתמש כבר קיים במערכת")
    return redirect("/admin")


@app.route("/admin/edit_user/<int:uid>", methods=["POST"])
@login_required
@admin_required
def edit_user(uid):
    full_name = request.form.get("full_name", "").strip()
    company_name = request.form.get("company_name", "").strip()
    company_id = request.form.get("company_id", "").strip()
    email = request.form.get("email", "").strip()
    phone = request.form.get("phone", "").strip()
    billing_mode = request.form.get("billing_mode", "monthly").strip() or "monthly"
    account_type = request.form.get("account_type", "").strip()
    service_valid_until = request.form.get("service_valid_until", "").strip()
    trial_start_date = request.form.get("trial_start_date", "").strip()
    with get_db() as db:
        user = db.execute("SELECT * FROM users WHERE id=? AND is_admin=0", (uid,)).fetchone()
        if not user:
            add_flash("לקוח לא נמצא")
            return redirect("/admin")
        trial_start_value = user["trial_start_date"] or ""
        service_until_value = user["service_valid_until"] or ""
        if account_type == "trial":
            trial_start_value = trial_start_date or user["trial_start_date"] or date.today().isoformat()
            service_until_value = ""
        elif account_type == "active":
            service_until_value = service_valid_until or user["service_valid_until"] or ""
            trial_start_value = ""
        db.execute(
            """UPDATE users SET full_name=?, company_name=?, company_id=?, email=?, phone=?,
            billing_mode=?, trial_start_date=?, service_valid_until=? WHERE id=?""",
            (full_name, company_name, company_id, email, phone, billing_mode,
             trial_start_value, service_until_value, uid),
        )
        db.commit()
    add_flash("פרטי הלקוח עודכנו בהצלחה")
    return redirect("/admin")


@app.route("/admin/toggle_active/<int:uid>", methods=["POST"])
@login_required
@admin_required
def toggle_active(uid):
    with get_db() as db:
        user = db.execute("SELECT id, active FROM users WHERE id=? AND is_admin=0", (uid,)).fetchone()
        if not user:
            add_flash("לקוח לא נמצא")
            return redirect("/admin")
        new_active = 0 if user["active"] else 1
        db.execute("UPDATE users SET active=? WHERE id=?", (new_active, uid))
        db.commit()
    add_flash("סטטוס הלקוח עודכן בהצלחה")
    return redirect("/admin")


@app.route("/admin/extend_trial/<int:uid>", methods=["POST"])
@login_required
@admin_required
def extend_trial(uid):
    extend_days = request.form.get("extend_days", "30").strip()
    if extend_days not in ("30", "60", "90"):
        extend_days = "30"
    extend_days = int(extend_days)
    with get_db() as db:
        user = db.execute("SELECT id, trial_days FROM users WHERE id=? AND is_admin=0", (uid,)).fetchone()
        if not user:
            add_flash("לקוח לא נמצא")
            return redirect("/admin")
        current_trial_days = user["trial_days"] or 30
        db.execute("UPDATE users SET trial_days=? WHERE id=?", (current_trial_days + extend_days, uid))
        db.commit()
    add_flash("תקופת הניסיון הוארכה ב-" + str(extend_days) + " ימים")
    return redirect("/admin")


@app.route("/admin/support/<int:request_id>/status", methods=["POST"])
@login_required
@admin_required
def update_support_request_status(request_id):
    new_status = request.form.get("status", "").strip().lower()
    if new_status not in {"accepted", "resolved"}:
        add_flash("סטטוס הפנייה אינו תקין")
        return redirect("/admin#adminSupport")
    with get_db() as db:
        db.execute("UPDATE support_requests SET status=? WHERE id=?", (new_status, request_id))
        db.commit()
    add_flash("סטטוס הפנייה עודכן בהצלחה")
    return redirect("/admin#adminSupport")


@app.route("/admin/tool-request/<int:request_id>/status", methods=["POST"])
@login_required
@admin_required
def update_tool_request_status(request_id):
    new_status = request.form.get("status", "").strip().lower()
    if new_status not in {"in_progress", "completed", "cancelled"}:
        add_flash("סטטוס הבקשה אינו תקין")
        return redirect("/admin#adminToolRequests")
    now = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    with get_db() as db:
        db.execute("UPDATE tool_requests SET status=?, updated_at=? WHERE id=?", (new_status, now, request_id))
        db.commit()
    add_flash("סטטוס הבקשה עודכן בהצלחה")
    return redirect("/admin#adminToolRequests")


@app.route("/admin/delete/<int:uid>")
@login_required
@admin_required
def delete_user(uid):
    with get_db() as db:
        db.execute("DELETE FROM users WHERE id=?", (uid,))
        db.execute("DELETE FROM permissions WHERE user_id=?", (uid,))
        db.commit()
    add_flash("הלקוח נמחק")
    return redirect("/admin")


@app.route("/admin/setpass/<int:uid>", methods=["POST"])
@login_required
@admin_required
def set_password(uid):
    with get_db() as db:
        db.execute(
            "UPDATE users SET password=? WHERE id=?",
            (generate_password_hash(request.form["new_password"]), uid),
        )
        db.commit()
    add_flash("הסיסמה עודכנה")
    return redirect("/admin")


@app.route("/admin/resetpass/<int:uid>", methods=["POST"])
@login_required
@admin_required
def reset_password(uid):
    temp_password = generate_temp_password()
    with get_db() as db:
        user = db.execute("SELECT full_name FROM users WHERE id=?", (uid,)).fetchone()
        db.execute(
            "UPDATE users SET password=? WHERE id=?",
            (generate_password_hash(temp_password), uid),
        )
        db.commit()
    name = user["full_name"] if user else str(uid)
    add_flash("סיסמה זמנית עבור " + name + ": " + temp_password)
    return redirect("/admin")


@app.route("/admin/permissions/<int:uid>", methods=["POST"])
@login_required
@admin_required
def set_permissions(uid):
    selected = request.form.getlist("scripts")
    with get_db() as db:
        db.execute("DELETE FROM permissions WHERE user_id=?", (uid,))
        for script_id in selected:
            if script_id in SCRIPTS:
                db.execute(
                    "INSERT OR IGNORE INTO permissions(user_id,script_id) VALUES (?,?)",
                    (uid, script_id),
                )
        db.commit()
    add_flash("ההרשאות עודכנו")
    return redirect("/admin")


# ════════════════════════════════════════════════════════════════════
# MARKETPLACE — Tool Execution Engine, AI Chat, Gallery & Routes
# ════════════════════════════════════════════════════════════════════

TOOL_ACTIONS = {
    "filter", "group_by", "sum", "count", "average", "min", "max",
    "sort", "rename_column", "select_columns", "add_column", "format_number",
    "concatenate", "split_column", "fill_missing", "remove_duplicates",
    "pivot", "unpivot", "date_extract", "math",
}

TOOL_CATEGORIES = {
    "general": "כללי",
    "payroll": "שכר",
    "attendance": "נוכחות",
    "hr": "משאבי אנוש",
    "finance": "כספים",
    "reports": "דוחות",
}


def validate_tool_definition(definition):
    """Validate a structured tool definition dict. Returns (ok, errors)."""
    errors = []
    if not isinstance(definition, dict):
        return False, ["Tool definition must be a JSON object"]
    if not definition.get("name"):
        errors.append("Tool must have a name")
    if not definition.get("input_type"):
        errors.append("Tool must specify input_type (xlsx, xls, csv)")
    steps = definition.get("steps", [])
    if not isinstance(steps, list) or len(steps) == 0:
        errors.append("Tool must have at least one processing step")
    for i, step in enumerate(steps):
        if not isinstance(step, dict):
            errors.append(f"Step {i+1} must be an object")
            continue
        action = step.get("action")
        if action not in TOOL_ACTIONS:
            errors.append(f"Step {i+1}: unknown action '{action}'. Allowed: {', '.join(sorted(TOOL_ACTIONS))}")
    output = definition.get("output_format", {})
    if not isinstance(output, dict):
        errors.append("output_format must be an object")
    return len(errors) == 0, errors


def execute_tool_definition(definition, input_path, output_path, extension):
    """Execute a structured JSON tool definition against an input file.

    Reads the input file into a pandas DataFrame, applies each step
    sequentially, then writes the result to output_path as xlsx.
    Returns dict with warnings list.
    """
    if pd is None:
        raise RuntimeError("pandas is required for marketplace tools. Install with: pip install pandas")

    # Read input
    if extension in ("xlsx", "xls"):
        df = pd.read_excel(input_path)
    elif extension == "csv":
        df = pd.read_csv(input_path)
    else:
        raise ValueError(f"Unsupported file type: {extension}")

    warnings = []
    steps = definition.get("steps", [])

    for i, step in enumerate(steps):
        action = step.get("action", "")
        try:
            df = _apply_tool_step(df, step, warnings)
        except Exception as exc:
            warnings.append(f"Step {i+1} ({action}) failed: {exc}")

    # Apply output formatting
    output_format = definition.get("output_format", {})
    title = output_format.get("title", definition.get("name", "Report"))
    columns = output_format.get("columns")
    if columns and isinstance(columns, list):
        available = [c for c in columns if c in df.columns]
        if available:
            df = df[available]

    # Write output
    wb = Workbook()
    ws = wb.active
    ws.title = title[:31] if title else "Report"

    # Write header
    for col_idx, col_name in enumerate(df.columns, 1):
        cell = ws.cell(row=1, column=col_idx, value=str(col_name))
        cell.font = Font(bold=True, color="FFFFFF", size=11)
        cell.fill = PatternFill(start_color="1E3A8A", end_color="1E3A8A", fill_type="solid")
        cell.alignment = Alignment(horizontal="center", vertical="center")

    # Write data
    for row_idx, row in enumerate(df.itertuples(index=False), 2):
        for col_idx, value in enumerate(row, 1):
            cell = ws.cell(row=row_idx, column=col_idx)
            if pd.isna(value):
                cell.value = ""
            elif isinstance(value, float):
                cell.value = round(value, 2)
                cell.number_format = '#,##0.00'
            else:
                cell.value = value
            if row_idx % 2 == 0:
                cell.fill = PatternFill(start_color="F0F4FF", end_color="F0F4FF", fill_type="solid")

    # Auto-width columns
    for col_idx in range(1, len(df.columns) + 1):
        letter = get_column_letter(col_idx)
        max_len = max(len(str(ws.cell(row=r, column=col_idx).value or "")) for r in range(1, min(ws.max_row + 1, 50)))
        ws.column_dimensions[letter].width = min(max(max_len + 2, 10), 50)

    ws.freeze_panes = "A2"
    wb.save(output_path)
    return {"warnings": warnings}


def _apply_tool_step(df, step, warnings):
    """Apply a single processing step to a DataFrame."""
    action = step.get("action", "")

    if action == "filter":
        field = step.get("field", "")
        operator = step.get("operator", "==")
        value = step.get("value")
        if field not in df.columns:
            warnings.append(f"Filter: column '{field}' not found, skipping")
            return df
        col = df[field]
        if operator == ">":
            df = df[pd.to_numeric(col, errors="coerce") > float(value)]
        elif operator == "<":
            df = df[pd.to_numeric(col, errors="coerce") < float(value)]
        elif operator == ">=":
            df = df[pd.to_numeric(col, errors="coerce") >= float(value)]
        elif operator == "<=":
            df = df[pd.to_numeric(col, errors="coerce") <= float(value)]
        elif operator == "==":
            df = df[col.astype(str) == str(value)]
        elif operator == "!=":
            df = df[col.astype(str) != str(value)]
        elif operator == "contains":
            df = df[col.astype(str).str.contains(str(value), na=False)]
        elif operator == "not_contains":
            df = df[~col.astype(str).str.contains(str(value), na=False)]
        elif operator == "is_empty":
            df = df[col.isna() | (col.astype(str).str.strip() == "")]
        elif operator == "not_empty":
            df = df[col.notna() & (col.astype(str).str.strip() != "")]

    elif action == "group_by":
        field = step.get("field", "")
        agg_field = step.get("agg_field", "")
        agg_func = step.get("agg_func", "sum")
        if field not in df.columns:
            warnings.append(f"Group by: column '{field}' not found")
            return df
        if agg_field and agg_field in df.columns:
            df[agg_field] = pd.to_numeric(df[agg_field], errors="coerce")
            agg_map = {"sum": "sum", "count": "count", "average": "mean", "min": "min", "max": "max"}
            df = df.groupby(field, as_index=False).agg({agg_field: agg_map.get(agg_func, "sum")})
        else:
            df = df.groupby(field, as_index=False).size()
            df.columns = [field, "count"]

    elif action == "sum":
        field = step.get("field", "")
        alias = step.get("alias", f"sum_{field}")
        if field in df.columns:
            df[field] = pd.to_numeric(df[field], errors="coerce")
            total = df[field].sum()
            summary_row = {col: "" for col in df.columns}
            summary_row[field] = total
            first_col = df.columns[0]
            summary_row[first_col] = alias
            df = pd.concat([df, pd.DataFrame([summary_row])], ignore_index=True)

    elif action == "count":
        field = step.get("field", "")
        alias = step.get("alias", f"count_{field}")
        if field in df.columns:
            df[alias] = df.groupby(field)[field].transform("count")

    elif action == "average":
        field = step.get("field", "")
        alias = step.get("alias", f"avg_{field}")
        if field in df.columns:
            df[field] = pd.to_numeric(df[field], errors="coerce")
            df[alias] = df[field].mean()

    elif action == "min":
        field = step.get("field", "")
        alias = step.get("alias", f"min_{field}")
        if field in df.columns:
            df[field] = pd.to_numeric(df[field], errors="coerce")
            df[alias] = df[field].min()

    elif action == "max":
        field = step.get("field", "")
        alias = step.get("alias", f"max_{field}")
        if field in df.columns:
            df[field] = pd.to_numeric(df[field], errors="coerce")
            df[alias] = df[field].max()

    elif action == "sort":
        field = step.get("field", "")
        order = step.get("order", "asc")
        if field in df.columns:
            ascending = order != "desc"
            try:
                df = df.sort_values(field, ascending=ascending, na_position="last")
            except TypeError:
                df[field] = pd.to_numeric(df[field], errors="coerce")
                df = df.sort_values(field, ascending=ascending, na_position="last")

    elif action == "rename_column":
        old_name = step.get("old_name", "")
        new_name = step.get("new_name", "")
        if old_name in df.columns and new_name:
            df = df.rename(columns={old_name: new_name})

    elif action == "select_columns":
        columns = step.get("columns", [])
        available = [c for c in columns if c in df.columns]
        if available:
            df = df[available]

    elif action == "add_column":
        name = step.get("name", "")
        value = step.get("value", "")
        if name:
            df[name] = value

    elif action == "format_number":
        field = step.get("field", "")
        decimals = step.get("decimals", 2)
        if field in df.columns:
            df[field] = pd.to_numeric(df[field], errors="coerce").round(int(decimals))

    elif action == "concatenate":
        fields = step.get("fields", [])
        separator = step.get("separator", " ")
        new_name = step.get("new_name", "combined")
        available = [f for f in fields if f in df.columns]
        if len(available) >= 2:
            df[new_name] = df[available].astype(str).agg(separator.join, axis=1)

    elif action == "fill_missing":
        field = step.get("field", "")
        fill_value = step.get("value", "")
        if field in df.columns:
            df[field] = df[field].fillna(fill_value)

    elif action == "remove_duplicates":
        fields = step.get("fields")
        if fields and isinstance(fields, list):
            available = [f for f in fields if f in df.columns]
            if available:
                df = df.drop_duplicates(subset=available)
        else:
            df = df.drop_duplicates()

    elif action == "math":
        field_a = step.get("field_a", "")
        field_b = step.get("field_b", "")
        operator = step.get("operator", "+")
        result_name = step.get("result_name", "result")
        if field_a in df.columns:
            a = pd.to_numeric(df[field_a], errors="coerce")
            if field_b in df.columns:
                b = pd.to_numeric(df[field_b], errors="coerce")
            else:
                try:
                    b = float(field_b) if field_b else 0
                except ValueError:
                    b = 0
            if operator == "+":
                df[result_name] = a + b
            elif operator == "-":
                df[result_name] = a - b
            elif operator == "*":
                df[result_name] = a * b
            elif operator == "/":
                df[result_name] = a / b.replace(0, float("nan")) if isinstance(b, pd.Series) else a / (b if b != 0 else float("nan"))

    elif action == "date_extract":
        field = step.get("field", "")
        part = step.get("part", "month")
        new_name = step.get("new_name", f"{field}_{part}")
        if field in df.columns:
            dt = pd.to_datetime(df[field], errors="coerce")
            if part == "year":
                df[new_name] = dt.dt.year
            elif part == "month":
                df[new_name] = dt.dt.month
            elif part == "day":
                df[new_name] = dt.dt.day
            elif part == "weekday":
                df[new_name] = dt.dt.day_name()

    elif action == "pivot":
        index = step.get("index", "")
        columns = step.get("columns", "")
        values = step.get("values", "")
        if index in df.columns and columns in df.columns and values in df.columns:
            df[values] = pd.to_numeric(df[values], errors="coerce")
            df = df.pivot_table(index=index, columns=columns, values=values, aggfunc="sum", fill_value=0).reset_index()
            df.columns = [str(c) if not isinstance(c, tuple) else "_".join(str(x) for x in c) for c in df.columns]

    return df


# ── AI Chat System Prompt ─────────────────────────────────────────

TOOL_CREATION_SYSTEM_PROMPT_BASE = """אתה עוזר אפיון כלים עבור פלטפורמת Scriptly — מערכת לעיבוד דוחות HR ושכר.

## מי המשתמש שלך
- אנשי HR ושכר שעובדים עם Excel, CSV ודוחות נוכחות.
- הם לא טכניים — הם לא מכירים JSON, קוד, או מונחים טכניים.
- הם מדברים עברית.
- הם רוצים תוצאות — לא להבין איך זה עובד מאחורי הקלעים.

## איך אתה מתנהג
- **דבר בעברית פשוטה** — אל תשתמש במונחים טכניים כמו JSON, API, function, schema.
- **היה חם ומעודד** — תגרום למשתמש להרגיש בטוח.
- **הנחה אותו צעד אחרי צעד** — אל תשאל 5 שאלות בבת אחת.
- **אל תציג קוד** — אל תציג Python, JavaScript, או כל שפת תכנות.

יש לך שלושה מצבים. אתה עובר ביניהם בצורה חלקה באותה שיחה.

═══════════════════════════════════
מצב א — ניתוח חופשי ("נסה עכשיו")
═══════════════════════════════════
כשלקוח מעלה קובץ ושואל שאלה חופשית:
- נתח את הקובץ (כותרות + 3-5 שורות בלבד, לא נתונים אמיתיים)
- ענה על שאלתו
- אם התשובה שימושית, הצע: "רוצה שאבנה לך כלי שיעשה את זה אוטומטית בכל פעם שתעלה קובץ כזה?"
- אם הלקוח אומר כן — עבור למצב ב

═══════════════════════════════════
מצב ב — אפיון כלי חדש ("בנה לי כלי")
═══════════════════════════════════
המטרה: לאסוף מספיק מידע כדי שמפתח יוכל לבנות את הכלי ללא שאלות נוספות.

שלב 1 — הבנת המטרה:
- "מה אתה רוצה לקבל בסוף? תאר לי את התוצאה האידיאלית."
- "באיזה מצבים תשתמש בכלי הזה?"

שלב 2 — הבנת הקלט:
- בקש קובץ לדוגמה (אמיתי או מדומה)
- רשום את כל העמודות שזיהית
- אשר עם הלקוח: "זיהיתי את העמודות הבאות: [רשימה]. האם זה נכון?"
- שאל: "האם יש עמודות נוספות שלא תמיד מופיעות?"

שלב 3 — הבנת הלוגיקה:
- "מה הכלי צריך לעשות עם המידע? פרט צעד אחר צעד."
- "האם יש חישובים מיוחדים? לדוגמה — שעות נוספות, אחוזים, תנאים?"
- "האם יש מקרים חריגים שהכלי צריך לטפל בהם אחרת?"

שלב 4 — משתנים:
- "האם יש ערכים שמשתנים מפעם לפעם? לדוגמה — תאריך, סף מספרי, שם מחלקה?"
- "מי קובע אותם — אתה בכל הרצה, או שהם קבועים?"

שלב 5 — הפלט:
- "איך אתה רוצה לקבל את התוצאה? Excel? PDF? טבלה על המסך?"
- "אילו עמודות צריכות להופיע בתוצאה?"
- "האם יש עיצוב מיוחד — צבעים, מיון, כותרות?"

שלב 6 — סיכום ואישור:
לפני שליחה, הצג ללקוח סיכום בשפה פשוטה:
"הבנתי שאתה רוצה כלי שעושה את הדברים הבאים:
[סיכום בשפה פשוטה]
האם זה נכון? יש משהו לתקן?"

אחרי שהלקוח מאשר — בנה את הסיכום הטכני בפורמט שלהלן ועטוף אותו ב:
```developer_brief
---DEVELOPER_BRIEF---
...
---END_BRIEF---
```

═══════════════════════════════════
מצב ג — שיפור כלי קיים
═══════════════════════════════════
כשלקוח אומר שהוא רוצה לשפר או לשנות כלי קיים, או כשמגיע [TOOL_CONTEXT]:
1. הצג לו את רשימת הכלים שלו (מוזרקים מהמערכת)
2. שאל איזה כלי ומה הוא רוצה לשנות
3. החלט:

שינוי פשוט — עשה לבד (עדכן את ה-JSON ועטוף ב-```json):
✅ הוסף/הסר עמודה מהפלט
✅ שנה ערך סף (מספר, תאריך)
✅ שנה מיון
✅ שנה עיצוב

שינוי מורכב — העבר לאפיון (בנה developer_brief):
❌ לוגיקה חדשה שלא קיימת בכלי
❌ חישוב מיוחד
❌ מבנה פלט שונה לחלוטין

אמור ללקוח בבירור: "את השינוי הזה אני יכול לעשות עכשיו" או "את השינוי הזה אני אעביר למפתח כי הוא דורש בנייה מחדש"

═══════════════════════════════════
מה אתה יכול ומה אתה לא יכול
═══════════════════════════════════
אתה יכול לבנות כלים שכוללים:
✅ סינון לפי תנאים, קיבוץ וסיכום, חישובים מתמטיים
✅ מיון ועיצוב, דוחות Excel מעוצבים, טיפול בתאריכים
✅ השוואה בין קבצים

אתה לא יכול לבנות:
❌ כלים שמתחברים למערכות חיצוניות
❌ כלים שדורשים אישורים והרשאות
❌ ממשקים גרפיים מורכבים

אם לקוח מבקש משהו שאינך יכול — אמור לו בכנות ותציע חלופה.

═══════════════════════════════════
פורמט הסיכום הטכני למפתח
═══════════════════════════════════
כשהאפיון מאושר, עטוף את הסיכום כך:
```developer_brief
---DEVELOPER_BRIEF---
שם כלי: [שם]
לקוח: {customer_name}
תאריך אפיון: {today_date}
סוג: [כלי חדש / שיפור לכלי קיים]
שם כלי קיים (אם רלוונטי): [שם]

תיאור קצר:
[משפט אחד מה הכלי עושה]

קלט:
- סוג קובץ: [xlsx/csv]
- עמודות חובה: [רשימה]
- עמודות אופציונליות: [רשימה]
- דוגמה לשורה: [שורה מהקובץ שהלקוח שיתף]

לוגיקה — צעד אחר צעד:
1. [צעד]
2. [צעד]
3. [צעד]

משתנים (ערכים שמשתנים בכל הרצה):
- שם המשתנה: [תיאור] | ברירת מחדל: [ערך]

מקרים חריגים:
- [מקרה]

פלט:
- סוג: [xlsx/pdf/טבלה]
- עמודות: [רשימה]
- מיון: [לפי מה, באיזה כיוון]
- עיצוב מיוחד: [אם יש]

הערות נוספות למפתח:
- [הערה]
---END_BRIEF---
```

═══════════════════════════════════
מבנה הגדרת כלי JSON (למצב ג — שינויים פשוטים בלבד)
═══════════════════════════════════
כשאתה מבצע שינוי פשוט בכלי קיים, בנה JSON בפורמט הבא ועטוף ב-```json:
{
  "name": "שם בעברית",
  "description": "תיאור קצר",
  "icon": "🔧",
  "category": "general",
  "input_type": "xlsx",
  "required_fields": ["שם_עמודה1", "שם_עמודה2"],
  "steps": [{"action": "ACTION_NAME", ...params}],
  "output_format": {"title": "כותרת הדוח", "columns": ["עמודה1", "עמודה2"]}
}

רשימת פעולות מותרות: filter, group_by, sum, count, average, min, max, sort, rename_column, select_columns, add_column, format_number, concatenate, split_column, fill_missing, remove_duplicates, math, date_extract, pivot, unpivot.
category חייב להיות: general, payroll, attendance, hr, finance, reports.

## פרטיות
- אל תבקש ואל תכלול שמות עובדים, תעודות זהות, או מידע אישי.
- השתמש רק בשמות עמודות ומבנה — אף פעם לא בנתונים אמיתיים.

═══════════════════════════════════
התנהגות כשלקוח בוחר דוגמה מהרשימה
═══════════════════════════════════
1. אל תשאל יותר מדי שאלות — נסה לענות מיד עם הקובץ שהוא העלה.
2. אם חסר לך קובץ — בקש רק את מה שחסר, בצורה ממוקדת.
3. אחרי שענית — שאל: "רוצה שאבנה לך כלי שיעשה את זה אוטומטית בכל פעם?"

כשלקוח שואל משהו שלא ברשימה:
- נסה לעזור — הרשימה היא רק השראה, לא מגבלה.
- אם זה מעבר ליכולת שלך — הסבר בכנות ואפיין יחד איתו.

═══════════════════════════════════
ידע בחוקי עבודה ושכר ישראליים
═══════════════════════════════════

### שעות נוספות
- שעות נוספות = מעל 8.6 שעות ביום או 45 שעות בשבוע
- תוספת לשעות נוספות: 125% לשעתיים הראשונות, 150% לאחר מכן
- עבודת לילה: מקסימום 58 שעות שבועיות
- מנוחה שבועית: 36 שעות רצופות לפחות

### מדרגות מס הכנסה (2025, חודשי)
| מדרגה | הכנסה חודשית (ש"ח) | שיעור מס |
|--------|---------------------|----------|
| 1 | 0 – 7,010 | 10% |
| 2 | 7,011 – 10,060 | 14% |
| 3 | 10,061 – 16,150 | 20% |
| 4 | 16,151 – 22,440 | 31% |
| 5 | 22,441 – 46,690 | 35% |
| 6 | 46,691 – 60,130 | 47% |
| 7 | 60,131+ | 50% |
- מס יסף: 3% נוסף על הכנסה מעל 721,560 ש"ח/שנה

### נקודות זיכוי ממס
- ערך נקודה (2025): 242 ש"ח/חודש (2,904 ש"ח/שנה)
- תושב ישראלי (בסיס): 2.25 נקודות
- אישה: +0.5 נקודות
- עולה חדש: עד +3.0 נקודות (18 חודשים ראשונים)
- ילד מתחת ל-18: +1.0 נקודה
- הורה יחיד: +1.0 נקודה

### ביטוח לאומי ומס בריאות — עובד
- מדרגה מופחתת (עד 7,522 ש"ח): 0.4% ביטוח לאומי + 3.1% מס בריאות = 3.5%
- מדרגה מלאה (7,523–50,695 ש"ח): 7.0% ביטוח לאומי + 5.0% מס בריאות = 12.0%
- תקרת הכנסה מבוטחת: 50,695 ש"ח/חודש

### ביטוח לאומי ומס בריאות — מעסיק
- מדרגה מופחתת: 3.8% ביטוח לאומי + 3.4% מס בריאות = 7.2%
- מדרגה מלאה: 7.6% ביטוח לאומי + 3.45% מס בריאות = 11.05%

### פנסיה (חובה מ-2017)
- עובד: 6% מהשכר
- מעסיק: 6.5% פנסיה + 6% פיצויים

### חישוב שכר נטו
נטו = ברוטו − מס הכנסה − ביטוח לאומי (עובד) − מס בריאות (עובד) − פנסיה (עובד) − ניכויים אחרים

### עלות מעסיק
עלות = ברוטו + ביטוח לאומי (מעסיק) + מס בריאות (מעסיק) + פנסיה (מעסיק) + פיצויים + צבירת חופשה + צבירת מחלה

חשוב: כשלקוח שואל על ציות, חישובי שכר, או חוקי עבודה — ציין תמיד: "מדובר בהערכה. מומלץ לאמת מול יועץ משפטי או רואה חשבון."

═══════════════════════════════════
טון ושפה
═══════════════════════════════════
- דבר בעברית תמיד, גם אם הלקוח כתב באנגלית
- טון חברותי ומקצועי — לא רובוטי
- תשובות קצרות וממוקדות — אל תסביר יותר ממה שנשאל
- כשמציג נתונים — השתמש בטבלה ברורה
- הוסף אימוג'י בצורה מינימלית ומקצועית
"""


def build_tool_creation_system_prompt(user_id, customer_name=""):
    """Build the system prompt with user's existing tools injected."""
    prompt = TOOL_CREATION_SYSTEM_PROMPT_BASE
    # Replace customer placeholders
    prompt = prompt.replace("{customer_name}", customer_name or "לקוח")
    prompt = prompt.replace("{today_date}", datetime.now().strftime("%Y-%m-%d"))

    # Inject user's existing tools from DB
    try:
        with get_db() as db:
            user_tools = db.execute(
                "SELECT id, name, description, icon, category, definition_json, status FROM marketplace_tools WHERE creator_id=? ORDER BY created_at DESC",
                (user_id,),
            ).fetchall()
        if user_tools:
            tools_context = "\n\n═══════════════════════════════════\nכלים קיימים של הלקוח\n═══════════════════════════════════\n"
            for t in user_tools:
                status_label = {"draft": "טיוטה", "pending_review": "ממתין לאישור", "approved": "מאושר", "rejected": "נדחה"}.get(t["status"], t["status"])
                tools_context += f"\n📦 {t['icon'] or '🔧'} {t['name']} (סטטוס: {status_label})\n"
                if t["description"]:
                    tools_context += f"   תיאור: {t['description']}\n"
                if t["definition_json"]:
                    tools_context += f"   הגדרה: {t['definition_json']}\n"
            prompt += tools_context
    except Exception:
        pass
    return prompt


def analyze_sample_file_for_chat(file_path, extension):
    """Read raw rows from a spreadsheet for AI context.
    Shows first 10 rows as-is so the AI can detect headers, titles, merged cells."""
    try:
        rows = []
        total_rows = 0
        if extension == "xls":
            wb = xlrd.open_workbook(file_path)
            sh = wb.sheet_by_index(0)
            total_rows = sh.nrows
            for ri in range(min(sh.nrows, 12)):
                row = []
                for ci in range(sh.ncols):
                    v = sh.cell_value(ri, ci)
                    if v is None or v == "":
                        row.append("")
                    else:
                        row.append(str(v))
                rows.append(row)
        elif extension == "xlsx":
            import openpyxl
            wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            ws = wb.active
            for ri, row_data in enumerate(ws.iter_rows(max_row=12, values_only=True)):
                row = [str(v) if v is not None else "" for v in row_data]
                rows.append(row)
            # Get total row count
            total_rows = ws.max_row or 0
            wb.close()
        elif extension == "csv":
            for enc in ("utf-8", "utf-8-sig", "cp1255", "latin-1"):
                try:
                    with open(file_path, "r", encoding=enc) as f:
                        import csv as _csv
                        reader = _csv.reader(f)
                        for i, row_data in enumerate(reader):
                            if i >= 12:
                                break
                            rows.append([str(v) for v in row_data])
                            total_rows += 1
                    # Count remaining rows
                    with open(file_path, "r", encoding=enc) as f:
                        total_rows = sum(1 for _ in f)
                    break
                except (UnicodeDecodeError, UnicodeError):
                    continue
        else:
            return "Unsupported file type"

        if not rows:
            return "File is empty or could not be read."

        num_cols = max(len(r) for r in rows) if rows else 0
        summary = f"File structure: {num_cols} columns, ~{total_rows} total rows\n"
        summary += f"First {len(rows)} rows (raw data as-is):\n"
        for ri, row in enumerate(rows):
            non_empty = [c for c in row if c.strip()]
            if not non_empty:
                summary += f"  Row {ri}: (empty row)\n"
            else:
                summary += f"  Row {ri}: {' | '.join(row)}\n"
        return summary
    except Exception as exc:
        return f"Could not analyze file: {exc}"


def call_claude_chat(messages, system_prompt=None, user_id=None, session_id=None):
    """Call Claude API for tool creation chat. Returns (text, input_tokens, output_tokens)."""
    if anthropic_sdk is None:
        return "שגיאה: ספריית Anthropic לא מותקנת. נא לפנות למנהל המערכת.", 0, 0

    api_key = os.environ.get("ANTHROPIC_API_KEY", "")
    if not api_key:
        return "שגיאה: מפתח API של Claude לא הוגדר. נא לפנות למנהל המערכת.", 0, 0

    client = anthropic_sdk.Anthropic(api_key=api_key)
    try:
        response = client.messages.create(
            model="claude-sonnet-4-6",
            max_tokens=4096,
            system=system_prompt or TOOL_CREATION_SYSTEM_PROMPT_BASE,
            messages=messages,
        )
        text = response.content[0].text
        input_tokens = getattr(response.usage, "input_tokens", 0) or 0
        output_tokens = getattr(response.usage, "output_tokens", 0) or 0
        # Sonnet pricing: $3/M input, $15/M output
        estimated_cost = (input_tokens * 3.0 / 1_000_000) + (output_tokens * 15.0 / 1_000_000)
        # Save token usage
        if user_id:
            try:
                now = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
                with get_db() as db:
                    db.execute(
                        "INSERT INTO chat_token_usage(user_id, session_id, input_tokens, output_tokens, estimated_cost_usd, created_at) VALUES (?,?,?,?,?,?)",
                        (user_id, session_id, input_tokens, output_tokens, round(estimated_cost, 6), now),
                    )
                    db.commit()
            except Exception:
                pass
        return text, input_tokens, output_tokens
    except Exception as exc:
        return f"שגיאה בתקשורת עם AI: {exc}", 0, 0


def extract_tool_json_from_message(text):
    """Try to extract a JSON tool definition from an AI message."""
    import re as _re
    # Look for ```json ... ``` blocks
    matches = _re.findall(r"```json\s*(.*?)```", text, _re.DOTALL)
    for match in matches:
        try:
            parsed = json.loads(match.strip())
            if isinstance(parsed, dict) and "steps" in parsed:
                return parsed
        except (json.JSONDecodeError, ValueError):
            continue
    # Try parsing the whole message as JSON
    try:
        parsed = json.loads(text.strip())
        if isinstance(parsed, dict) and "steps" in parsed:
            return parsed
    except (json.JSONDecodeError, ValueError):
        pass
    return None


def extract_developer_brief(text):
    """Extract a developer brief from AI message. Returns brief text or None."""
    import re as _re
    # Look for ```developer_brief ... ``` blocks containing ---DEVELOPER_BRIEF--- ... ---END_BRIEF---
    block_match = _re.search(r"```developer_brief\s*(.*?)```", text, _re.DOTALL)
    if block_match:
        inner = block_match.group(1).strip()
        brief_match = _re.search(r"---DEVELOPER_BRIEF---(.*?)---END_BRIEF---", inner, _re.DOTALL)
        if brief_match:
            return brief_match.group(1).strip()
        return inner
    # Fallback: look for markers directly in the text
    brief_match = _re.search(r"---DEVELOPER_BRIEF---(.*?)---END_BRIEF---", text, _re.DOTALL)
    if brief_match:
        return brief_match.group(1).strip()
    return None


def extract_brief_tool_name(brief_text):
    """Extract tool name from a developer brief text."""
    import re as _re
    match = _re.search(r"שם כלי:\s*(.+)", brief_text)
    if match:
        return match.group(1).strip()
    return "כלי ללא שם"


# ── Marketplace Routes ────────────────────────────────────────────

@app.route("/marketplace")
@login_required
def marketplace():
    if session.get("is_admin"):
        return redirect("/admin")

    lang = get_flow_lang()
    text = get_flow_text(lang)

    with get_db() as db:
        tools = db.execute(
            """SELECT t.*, u.full_name as creator_name, u.company_name as creator_company,
               (SELECT COUNT(*) FROM tool_ratings WHERE tool_id=t.id) as rating_count,
               (SELECT COALESCE(AVG(rating),0) FROM tool_ratings WHERE tool_id=t.id) as avg_rating,
               (SELECT COUNT(*) FROM tool_installs WHERE tool_id=t.id) as install_count
            FROM marketplace_tools t
            JOIN users u ON u.id = t.creator_id
            WHERE t.status = 'approved' AND t.is_public = 1
            ORDER BY t.usage_count DESC, t.created_at DESC""",
        ).fetchall()

        user_installed = set()
        installs = db.execute(
            "SELECT tool_id FROM tool_installs WHERE user_id=?",
            (session["user_id"],),
        ).fetchall()
        for inst in installs:
            user_installed.add(inst["tool_id"])

        user_tools = db.execute(
            "SELECT * FROM marketplace_tools WHERE creator_id=? ORDER BY created_at DESC",
            (session["user_id"],),
        ).fetchall()

    search_query = request.args.get("q", "").strip()
    category_filter = request.args.get("category", "").strip()

    # Build tool cards
    tool_cards = ""
    for tool in tools:
        if search_query and search_query.lower() not in (tool["name"] or "").lower() and search_query.lower() not in (tool["description"] or "").lower():
            continue
        if category_filter and tool["category"] != category_filter:
            continue
        stars = "⭐" * max(1, min(5, int(float(tool["avg_rating"] or 0) + 0.5)))
        installed_badge = '<span style="background:#ecfdf5;color:#047857;padding:3px 8px;border-radius:99px;font-size:11px;font-weight:700">מותקן</span>' if tool["id"] in user_installed else ""
        cat_label = TOOL_CATEGORIES.get(tool["category"], tool["category"] or "כללי")

        tool_cards += (
            '<a href="/marketplace/tool/' + str(tool["id"]) + '" style="background:white;border-radius:18px;box-shadow:0 2px 12px rgba(37,99,235,.06);padding:1.5rem;text-decoration:none;display:block;border:1px solid #e2e8f0;transition:all .3s cubic-bezier(.4,0,.2,1)" '
            'onmouseover="this.style.transform=\'translateY(-3px)\';this.style.boxShadow=\'0 8px 28px rgba(37,99,235,.12)\';this.style.borderColor=\'#bfdbfe\'" '
            'onmouseout="this.style.transform=\'none\';this.style.boxShadow=\'0 2px 12px rgba(37,99,235,.06)\';this.style.borderColor=\'#e2e8f0\'">'
            '<div style="display:flex;align-items:center;justify-content:space-between;margin-bottom:10px">'
            '<div style="width:44px;height:44px;border-radius:12px;background:linear-gradient(135deg,#eff6ff,#dbeafe);display:flex;align-items:center;justify-content:center;font-size:22px">' + esc(tool["icon"] or "🔧") + '</div>'
            + installed_badge
            + '</div>'
            '<div style="font-size:15px;font-weight:800;color:#0f172a;margin-bottom:4px">' + esc(tool["name"]) + '</div>'
            '<div style="font-size:12px;color:#64748b;margin-bottom:8px;line-height:1.6;min-height:36px">' + esc((tool["description"] or "")[:100]) + '</div>'
            '<div style="display:flex;align-items:center;justify-content:space-between;font-size:11px;color:#94a3b8">'
            '<span>' + esc(cat_label) + '</span>'
            '<span>' + stars + ' (' + str(tool["rating_count"] or 0) + ')</span>'
            '</div>'
            '<div style="font-size:11px;color:#94a3b8;margin-top:4px">' + esc(tool["creator_name"] or "אנונימי") + ' · ' + str(tool["usage_count"] or 0) + ' שימושים</div>'
            '</a>'
        )

    if not tool_cards:
        tool_cards = (
            '<div style="text-align:center;padding:3rem;color:#94a3b8;grid-column:1/-1">'
            '<div style="font-size:48px;margin-bottom:1rem">🔍</div>'
            '<div>לא נמצאו כלים' + (' עבור "' + esc(search_query) + '"' if search_query else '') + '</div>'
            '</div>'
        )

    # Category filter pills
    cat_pills = '<a href="/marketplace" style="padding:8px 16px;border-radius:99px;font-size:12px;font-weight:700;text-decoration:none;transition:all .2s;' + ('background:linear-gradient(135deg,#1e3a8a,#2563eb);color:white;box-shadow:0 2px 8px rgba(37,99,235,.25)' if not category_filter else 'background:#f1f5f9;color:#475569') + '">הכל</a>'
    for cat_id, cat_name in TOOL_CATEGORIES.items():
        active = category_filter == cat_id
        cat_pills += '<a href="/marketplace?category=' + cat_id + (('&q=' + esc(search_query)) if search_query else '') + '" style="padding:8px 16px;border-radius:99px;font-size:12px;font-weight:700;text-decoration:none;transition:all .2s;' + ('background:linear-gradient(135deg,#1e3a8a,#2563eb);color:white;box-shadow:0 2px 8px rgba(37,99,235,.25)' if active else 'background:#f1f5f9;color:#475569') + '">' + esc(cat_name) + '</a>'

    # My tools section — link to My Tools page
    my_tools_html = ""
    if user_tools:
        my_tools_rows = ""
        status_labels = {"draft": "טיוטה", "pending_review": "ממתין לאישור", "approved": "מאושר", "rejected": "נדחה"}
        status_colors = {"draft": ("#fef3c7", "#92400e"), "pending_review": ("#fff7ed", "#c2410c"), "approved": ("#ecfdf5", "#047857"), "rejected": ("#fef2f2", "#b91c1c")}
        for t in user_tools[:3]:  # Show max 3, link to full page
            s_bg, s_fg = status_colors.get(t["status"], ("#f8fafc", "#475569"))
            my_tools_rows += (
                '<a href="/tools/my-tools" style="display:flex;align-items:center;justify-content:space-between;padding:10px 0;border-bottom:1px solid #f1f5f9;text-decoration:none;color:inherit">'
                '<div><span style="font-size:15px;font-weight:700;color:#0f172a">' + esc(t["icon"] or "🔧") + " " + esc(t["name"]) + '</span></div>'
                '<span style="padding:4px 10px;border-radius:99px;font-size:11px;font-weight:700;background:' + s_bg + ';color:' + s_fg + '">' + esc(status_labels.get(t["status"], t["status"])) + '</span>'
                '</a>'
            )
        my_tools_html = (
            '<div class="card" style="margin-bottom:1.5rem;border:1px solid #dbeafe;border-radius:18px">'
            '<div style="display:flex;align-items:center;justify-content:space-between;margin-bottom:12px">'
            '<div style="font-size:18px;font-weight:800;color:#0f172a">📦 הכלים שלי</div>'
            '<a href="/tools/my-tools" style="font-size:13px;color:#2563eb;text-decoration:none;font-weight:600">הצג הכל &#8592;</a>'
            '</div>'
            + my_tools_rows
            + '</div>'
        )

    body = (
        '<div style="display:flex;align-items:center;justify-content:space-between;flex-wrap:wrap;gap:14px;margin-bottom:1.5rem">'
        '<div>'
        '<h2 style="font-size:clamp(22px,3.5vw,28px);font-weight:800;color:#0f172a;margin-bottom:6px;letter-spacing:-0.3px">🛒 שוק הכלים</h2>'
        '<p style="font-size:14px;color:#64748b;line-height:1.6">גלה כלים שנבנו על ידי משתמשים אחרים — או צור כלי משלך עם AI</p>'
        '</div>'
        '<div style="display:flex;gap:10px">'
        '<a href="/tools/my-tools" class="btn btn-gray" style="text-decoration:none">&#128188; הכלים שלי</a>'
        '<a href="/tools/create" class="btn btn-blue" style="text-decoration:none">&#129302; צור כלי חדש</a>'
        '</div>'
        '</div>'
        + my_tools_html
        + '<form method="get" action="/marketplace" style="margin-bottom:1rem;display:flex;gap:10px">'
        '<input type="text" name="q" value="' + esc(search_query) + '" placeholder="חיפוש כלים..." style="flex:1">'
        '<button type="submit" class="btn btn-blue">חפש</button>'
        '</form>'
        '<div style="display:flex;gap:8px;flex-wrap:wrap;margin-bottom:1.5rem">' + cat_pills + '</div>'
        '<div style="display:grid;grid-template-columns:repeat(auto-fill,minmax(240px,1fr));gap:1rem">'
        + tool_cards
        + '</div>'
        '<div style="text-align:center;margin-top:2rem">'
        '<a href="/dashboard" style="color:#2563eb;font-size:14px;font-weight:600;text-decoration:none">← חזרה לדשבורד</a>'
        '</div>'
    )

    return render("שוק הכלים", body, lang=lang, topbar_greeting=text["topbar_greeting"], logout_label=text["logout"], show_lang_switch=True)


@app.route("/marketplace/tool/<int:tool_id>")
@login_required
def marketplace_tool_detail(tool_id):
    lang = get_flow_lang()
    text = get_flow_text(lang)

    with get_db() as db:
        tool = db.execute(
            """SELECT t.*, u.full_name as creator_name, u.company_name as creator_company
            FROM marketplace_tools t JOIN users u ON u.id = t.creator_id
            WHERE t.id = ?""",
            (tool_id,),
        ).fetchone()
        if not tool:
            add_flash("הכלי לא נמצא")
            return redirect("/marketplace")

        rating_info = db.execute(
            "SELECT COUNT(*) as cnt, COALESCE(AVG(rating),0) as avg FROM tool_ratings WHERE tool_id=?",
            (tool_id,),
        ).fetchone()
        user_rating = db.execute(
            "SELECT rating FROM tool_ratings WHERE tool_id=? AND user_id=?",
            (tool_id, session["user_id"]),
        ).fetchone()
        is_installed = db.execute(
            "SELECT 1 FROM tool_installs WHERE tool_id=? AND user_id=?",
            (tool_id, session["user_id"]),
        ).fetchone()
        comments = db.execute(
            "SELECT * FROM tool_comments WHERE tool_id=? ORDER BY created_at DESC LIMIT 20",
            (tool_id,),
        ).fetchall()

    avg_rating = float(rating_info["avg"] or 0)
    rating_count = int(rating_info["cnt"] or 0)
    user_stars = user_rating["rating"] if user_rating else 0

    # Parse definition for display
    try:
        definition = json.loads(tool["definition_json"])
    except (json.JSONDecodeError, TypeError):
        definition = {}

    steps_html = ""
    for i, step in enumerate(definition.get("steps", []), 1):
        action = step.get("action", "?")
        detail = ""
        if action == "filter":
            detail = f'{step.get("field", "?")} {step.get("operator", "?")} {step.get("value", "?")}'
        elif action == "group_by":
            detail = f'{step.get("field", "?")} → {step.get("agg_func", "sum")}({step.get("agg_field", "?")})'
        elif action == "sort":
            detail = f'{step.get("field", "?")} ({step.get("order", "asc")})'
        elif action == "select_columns":
            detail = ", ".join(step.get("columns", [])[:5])
        elif action == "math":
            detail = f'{step.get("field_a", "?")} {step.get("operator", "?")} {step.get("field_b", "?")} → {step.get("result_name", "?")}'
        else:
            detail = step.get("field", "") or step.get("name", "")
        steps_html += (
            '<div style="display:flex;align-items:center;gap:12px;padding:10px 14px;background:linear-gradient(180deg,#f8fafc,#f1f5f9);border:1px solid #e2e8f0;border-radius:12px;margin-bottom:8px">'
            '<span style="background:linear-gradient(135deg,#1e3a8a,#2563eb);color:white;width:26px;height:26px;border-radius:50%;display:flex;align-items:center;justify-content:center;font-size:12px;font-weight:700;flex-shrink:0">' + str(i) + '</span>'
            '<span style="font-size:13px;color:#0f172a"><strong>' + esc(action) + '</strong> ' + esc(detail) + '</span>'
            '</div>'
        )

    # Star rating UI
    star_html = ""
    for s in range(1, 6):
        filled = "⭐" if s <= user_stars else "☆"
        star_html += '<a href="/marketplace/tool/' + str(tool_id) + '/rate?stars=' + str(s) + '" style="font-size:24px;text-decoration:none;cursor:pointer">' + filled + '</a>'

    # Comments
    comments_html = ""
    for c in comments:
        comments_html += (
            '<div style="padding:12px 0;border-bottom:1px solid #f1f5f9">'
            '<div style="display:flex;justify-content:space-between;margin-bottom:4px">'
            '<span style="font-size:13px;font-weight:700;color:#0f172a">' + esc(c["full_name"] or c["username"] or "אנונימי") + '</span>'
            '<span style="font-size:11px;color:#94a3b8">' + esc(format_ui_datetime(c["created_at"])) + '</span>'
            '</div>'
            '<div style="font-size:13px;color:#334155;line-height:1.7">' + esc(c["body"]) + '</div>'
            '</div>'
        )

    install_btn = ""
    if tool["status"] == "approved" and tool["is_public"]:
        if is_installed:
            install_btn = (
                '<a href="/marketplace/tool/' + str(tool_id) + '/run" class="btn btn-blue" style="text-decoration:none;display:inline-flex;align-items:center;gap:6px;padding:12px 24px;border-radius:12px;font-size:14px">▶ הרץ כלי</a>'
                '<form method="post" action="/marketplace/tool/' + str(tool_id) + '/uninstall" style="display:inline"><button type="submit" class="btn btn-gray" style="border-radius:12px;font-size:13px">הסר מהדשבורד</button></form>'
            )
        else:
            install_btn = '<form method="post" action="/marketplace/tool/' + str(tool_id) + '/install" style="display:inline"><button type="submit" class="btn btn-blue" style="border-radius:12px;padding:12px 24px;font-size:14px">➕ הוסף לדשבורד שלי</button></form>'

    body = (
        '<div style="margin-bottom:1.5rem"><a href="/marketplace" style="color:#2563eb;font-size:13px;font-weight:600;text-decoration:none">← חזרה לשוק הכלים</a></div>'
        '<div class="card" style="border-radius:18px">'
        '<div style="display:flex;align-items:flex-start;justify-content:space-between;gap:16px;flex-wrap:wrap;margin-bottom:1rem">'
        '<div>'
        '<div style="width:56px;height:56px;border-radius:16px;background:linear-gradient(135deg,#eff6ff,#dbeafe);display:flex;align-items:center;justify-content:center;font-size:28px;margin-bottom:10px">' + esc(tool["icon"] or "🔧") + '</div>'
        '<h2 style="font-size:22px;font-weight:800;color:#0f172a;margin-bottom:6px;border:none;padding:0;letter-spacing:-0.3px">' + esc(tool["name"]) + '</h2>'
        '<div style="font-size:14px;color:#64748b;line-height:1.7;margin-bottom:8px">' + esc(tool["description"] or "") + '</div>'
        '<div style="font-size:12px;color:#94a3b8">יוצר: ' + esc(tool["creator_name"] or "אנונימי") + ' · ' + str(tool["usage_count"] or 0) + ' שימושים · ' + str(rating_count) + ' דירוגים</div>'
        '</div>'
        '<div style="display:flex;flex-direction:column;gap:8px;align-items:flex-end">'
        + install_btn
        + '</div></div>'
        '<div style="margin-bottom:1rem"><span style="font-size:13px;color:#64748b;margin-left:8px">הדירוג שלך:</span>' + star_html + '</div>'
        '</div>'
        '<div class="card" style="border-radius:18px">'
        '<div style="font-size:16px;font-weight:800;color:#0f172a;margin-bottom:12px">⚙️ שלבי עיבוד</div>'
        + (steps_html or '<div style="color:#94a3b8;font-size:13px">אין שלבי עיבוד מוגדרים</div>')
        + '</div>'
        '<div class="card" style="border-radius:18px">'
        '<div style="font-size:16px;font-weight:800;color:#0f172a;margin-bottom:12px">💬 תגובות</div>'
        '<form method="post" action="/marketplace/tool/' + str(tool_id) + '/comment" style="margin-bottom:1rem;display:flex;gap:10px">'
        '<input type="text" name="body" placeholder="כתוב תגובה..." style="flex:1" required>'
        '<button type="submit" class="btn btn-blue">שלח</button>'
        '</form>'
        + (comments_html or '<div style="color:#94a3b8;font-size:13px;text-align:center;padding:1rem">עדיין אין תגובות — היה הראשון!</div>')
        + '</div>'
    )

    return render(esc(tool["name"]), body, lang=lang, topbar_greeting=text["topbar_greeting"], logout_label=text["logout"], show_lang_switch=True)


@app.route("/marketplace/tool/<int:tool_id>/install", methods=["POST"])
@login_required
def install_tool(tool_id):
    now = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    with get_db() as db:
        tool = db.execute("SELECT id, status, is_public FROM marketplace_tools WHERE id=?", (tool_id,)).fetchone()
        if not tool or tool["status"] != "approved" or not tool["is_public"]:
            add_flash("לא ניתן להתקין כלי זה")
            return redirect("/marketplace")
        existing = db.execute("SELECT 1 FROM tool_installs WHERE tool_id=? AND user_id=?", (tool_id, session["user_id"])).fetchone()
        if not existing:
            db.execute("INSERT INTO tool_installs(tool_id, user_id, created_at) VALUES (?,?,?)", (tool_id, session["user_id"], now))
            db.commit()
    add_flash("הכלי נוסף לדשבורד שלך!")
    return redirect("/marketplace/tool/" + str(tool_id))


@app.route("/marketplace/tool/<int:tool_id>/uninstall", methods=["POST"])
@login_required
def uninstall_tool(tool_id):
    with get_db() as db:
        db.execute("DELETE FROM tool_installs WHERE tool_id=? AND user_id=?", (tool_id, session["user_id"]))
        db.commit()
    add_flash("הכלי הוסר מהדשבורד שלך")
    return redirect("/marketplace/tool/" + str(tool_id))


@app.route("/marketplace/tool/<int:tool_id>/rate")
@login_required
def rate_tool(tool_id):
    stars = request.args.get("stars", "1")
    try:
        stars = max(1, min(5, int(stars)))
    except (ValueError, TypeError):
        stars = 1
    now = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    with get_db() as db:
        existing = db.execute("SELECT id FROM tool_ratings WHERE tool_id=? AND user_id=?", (tool_id, session["user_id"])).fetchone()
        if existing:
            db.execute("UPDATE tool_ratings SET rating=?, created_at=? WHERE tool_id=? AND user_id=?", (stars, now, tool_id, session["user_id"]))
        else:
            db.execute("INSERT INTO tool_ratings(tool_id, user_id, rating, created_at) VALUES (?,?,?,?)", (tool_id, session["user_id"], stars, now))
        db.commit()
    return redirect("/marketplace/tool/" + str(tool_id))


@app.route("/marketplace/tool/<int:tool_id>/comment", methods=["POST"])
@login_required
def comment_tool(tool_id):
    body_text = request.form.get("body", "").strip()
    if not body_text or len(body_text) > 1000:
        add_flash("תגובה חייבת להכיל 1-1000 תווים")
        return redirect("/marketplace/tool/" + str(tool_id))
    now = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    with get_db() as db:
        db.execute(
            "INSERT INTO tool_comments(tool_id, user_id, username, full_name, body, created_at) VALUES (?,?,?,?,?,?)",
            (tool_id, session["user_id"], session.get("username", ""), session.get("name", ""), body_text, now),
        )
        db.commit()
    return redirect("/marketplace/tool/" + str(tool_id))


@app.route("/marketplace/tool/<int:tool_id>/run", methods=["GET", "POST"])
@login_required
def run_marketplace_tool(tool_id):
    if session.get("is_admin"):
        return redirect("/admin")

    lang = get_flow_lang()
    text = get_flow_text(lang)

    with get_db() as db:
        tool = db.execute("SELECT * FROM marketplace_tools WHERE id=?", (tool_id,)).fetchone()
        if not tool:
            add_flash("הכלי לא נמצא")
            return redirect("/marketplace")

        # Check user has installed it or is the creator
        is_installed = db.execute("SELECT 1 FROM tool_installs WHERE tool_id=? AND user_id=?", (tool_id, session["user_id"])).fetchone()
        is_creator = tool["creator_id"] == session["user_id"]
        if not is_installed and not is_creator:
            add_flash("יש להתקין את הכלי לפני השימוש")
            return redirect("/marketplace/tool/" + str(tool_id))

    try:
        definition = json.loads(tool["definition_json"])
    except (json.JSONDecodeError, TypeError):
        add_flash("הגדרת הכלי שגויה")
        return redirect("/marketplace")

    error = ""
    result = None

    if request.method == "POST":
        file_obj = request.files.get("file")
        validation_error, ext = validate_upload(file_obj)
        if validation_error == "missing":
            error = '<div class="flash-err">לא נבחר קובץ</div>'
        elif validation_error == "unsupported":
            error = '<div class="flash-err">סוג קובץ לא נתמך</div>'
        elif validation_error == "invalid_excel":
            error = '<div class="flash-err">הקובץ אינו קובץ Excel תקין</div>'
        elif validation_error == "empty":
            error = '<div class="flash-err">הקובץ ריק</div>'
        elif validation_error == "too_large":
            error = '<div class="flash-err">הקובץ גדול מדי</div>'
        else:
            uid = str(uuid.uuid4())[:8]
            inp = str(UPLOAD_FOLDER / f"{uid}.{ext}")
            out_name = f"{uid}_{tool['name'][:30]}.xlsx"
            out = str(OUTPUT_FOLDER / out_name)
            file_obj.save(inp)
            try:
                exec_result = execute_tool_definition(definition, inp, out, ext)
                # Update usage count
                with get_db() as db:
                    db.execute("UPDATE marketplace_tools SET usage_count = usage_count + 1 WHERE id=?", (tool_id,))
                    db.commit()
                log_user_activity("run_marketplace_tool", "הרצת כלי מרקטפלייס", "marketplace_" + str(tool_id), tool["name"], "")
                result = out_name
            except Exception as exc:
                error = '<div class="flash-err">שגיאה בעיבוד: ' + esc(str(exc)) + '</div>'
            finally:
                try:
                    os.remove(inp)
                except OSError:
                    pass

    # ── Chat session for tool improvement (creator only) ──
    chat_session_id = None
    chat_msgs_html = ""
    if is_creator:
        now = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        with get_db() as db:
            existing = db.execute(
                "SELECT id FROM tool_chat_sessions WHERE user_id=? AND tool_id=? AND status='active' ORDER BY created_at DESC LIMIT 1",
                (session["user_id"], tool_id),
            ).fetchone()
            if existing:
                chat_session_id = existing["id"]
            else:
                db.execute(
                    "INSERT INTO tool_chat_sessions(user_id, messages_json, tool_id, status, created_at, updated_at) VALUES (?,?,?,?,?,?)",
                    (session["user_id"], "[]", tool_id, "active", now, now),
                )
                db.commit()
                chat_session_id = db.execute(
                    "SELECT id FROM tool_chat_sessions WHERE user_id=? AND tool_id=? ORDER BY id DESC LIMIT 1",
                    (session["user_id"], tool_id),
                ).fetchone()["id"]

            # Build existing messages HTML
            chat_data = db.execute("SELECT messages_json FROM tool_chat_sessions WHERE id=?", (chat_session_id,)).fetchone()
        existing_messages = []
        if chat_data:
            try:
                existing_messages = json.loads(chat_data["messages_json"])
            except (json.JSONDecodeError, TypeError):
                existing_messages = []

        for msg in existing_messages:
            if msg["role"] == "user":
                user_text = esc(msg["content"]).replace("\n", "<br>")
                # Skip internal tool context messages
                if user_text.startswith("[TOOL_CONTEXT]"):
                    continue
                chat_msgs_html += '<div class="tr-chat-bubble tr-chat-user"><div class="tr-chat-bubble-inner tr-chat-user-inner">' + user_text + '</div></div>'
            else:
                content_html = esc(msg["content"]).replace("```json", '<pre style="background:#0f172a;color:#38bdf8;padding:12px;border-radius:10px;font-size:12px;overflow-x:auto;direction:ltr;text-align:left">').replace("```", "</pre>").replace("\n", "<br>")
                chat_msgs_html += '<div class="tr-chat-bubble tr-chat-ai"><div class="tr-chat-bubble-inner tr-chat-ai-inner">' + content_html + '</div></div>'
                tool_def = extract_tool_json_from_message(msg["content"])
                if tool_def:
                    chat_msgs_html += (
                        '<div class="tr-chat-tool-ready">'
                        '<div style="font-size:13px;font-weight:700;color:#047857;margin-bottom:6px">&#127881; הכלי מוכן!</div>'
                        '<button type="button" class="btn btn-blue" style="border-radius:10px;font-size:12px;background:#047857" '
                        'onclick=\'updateToolDef(' + esc(json.dumps(json.dumps(tool_def, ensure_ascii=False))) + ')\'>&#128190; עדכן את הכלי</button>'
                        '</div>'
                    )

    # ── Build left panel (tool run) ──
    back_link = '/tools/my-tools' if is_creator else '/marketplace'
    back_label = 'חזרה לכלים שלי' if is_creator else 'חזרה לשוק הכלים'

    if result:
        left_panel = (
            '<div class="card" style="text-align:center;border-radius:18px">'
            '<div style="font-size:48px;margin-bottom:12px">&#10004;&#65039;</div>'
            '<div style="font-size:20px;font-weight:800;color:#047857;margin-bottom:12px">הקובץ מוכן!</div>'
            '<a href="/download/' + esc(result) + '" class="btn btn-blue" style="text-decoration:none;display:inline-flex;padding:12px 28px;border-radius:12px;font-size:15px;margin-bottom:12px">&#11015; הורדת הקובץ</a><br>'
            '<a href="/marketplace/tool/' + str(tool_id) + '/run" style="color:#2563eb;font-size:13px;text-decoration:none">עיבוד קובץ נוסף</a>'
            '</div>'
        )
    else:
        input_type = definition.get("input_type", "xlsx")
        accept_val = ".xlsx,.xls,.csv"
        if input_type == "csv":
            accept_val = ".csv"
        elif input_type in ("xlsx", "xls"):
            accept_val = ".xls,.xlsx"

        left_panel = (
            '<div class="card" style="border-radius:18px">'
            '<div style="width:48px;height:48px;border-radius:14px;background:linear-gradient(135deg,#eff6ff,#dbeafe);display:flex;align-items:center;justify-content:center;font-size:24px;margin-bottom:10px">' + esc(tool["icon"] or "🔧") + '</div>'
            '<h2 style="font-size:20px;font-weight:800;color:#0f172a;margin-bottom:6px;border:none;padding:0;letter-spacing:-0.3px">' + esc(tool["name"]) + '</h2>'
            '<div style="font-size:14px;color:#64748b;margin-bottom:1rem;line-height:1.7">' + esc(tool["description"] or "") + '</div>'
            + error
            + '<form method="post" enctype="multipart/form-data">'
            '<div class="drop-zone" id="dropZone" onclick="document.getElementById(\'fileInput\').click()">'
            '<div style="font-size:32px;margin-bottom:8px">&#128193;</div>'
            '<div style="font-size:14px;color:#475569">לחץ לבחירת קובץ או גרור לכאן</div>'
            '<div style="font-size:12px;color:#94a3b8;margin-top:4px" id="fileName">נתמכים: ' + esc(accept_val) + '</div>'
            '<input type="file" name="file" id="fileInput" accept="' + esc(accept_val) + '" style="display:none" onchange="document.getElementById(\'fileName\').textContent=this.files[0]?this.files[0].name:\'\'">'
            '</div>'
            '<button type="submit" class="btn btn-blue" style="width:100%;padding:14px;border-radius:12px;font-size:15px">&#9654; הרץ</button>'
            '</form>'
            '</div>'
        )

    # ── Build right panel (chat) — only for creator ──
    right_panel = ""
    chat_js = ""
    if is_creator:
        right_panel = (
            '<div class="tr-chat-panel" id="trChatPanel">'
            '<div class="tr-chat-header">'
            '<div class="tr-chat-header-title">&#129302; שיפור עם AI</div>'
            '<div class="tr-chat-header-sub">בקש שינויים והכלי יתעדכן</div>'
            '</div>'
            '<div class="tr-chat-messages" id="trChatMessages">'
            # Welcome message
            '<div class="tr-chat-bubble tr-chat-ai">'
            '<div class="tr-chat-bubble-inner tr-chat-ai-inner">'
            'היי! &#128075; אני כאן כדי לעזור לשפר את הכלי.<br>'
            'נסה להריץ את הכלי ואם יש בעיה — ספר לי מה לא עבד ואני אתקן.'
            '</div></div>'
            + chat_msgs_html
            + '</div>'
            # Input area
            '<div class="tr-chat-input-area">'
            '<div class="tr-chat-input-row">'
            '<textarea id="trChatInput" class="tr-chat-textarea" placeholder="מה לשפר בכלי?" rows="1"></textarea>'
            '<button type="button" id="trChatSendBtn" class="tr-chat-send-btn" onclick="trSendChat()">'
            '<span id="trSendBtnText">שלח</span>'
            '<span id="trSendBtnSpinner" class="tr-chat-send-spinner" style="display:none"></span>'
            '</button>'
            '</div>'
            '<div style="font-size:11px;color:#94a3b8;margin-top:4px;text-align:center">Enter לשליחה &middot; Shift+Enter לשורה חדשה</div>'
            '</div>'
            '<div id="trUpdateToast" style="display:none;position:absolute;bottom:70px;left:12px;right:12px;background:#ecfdf5;border:1px solid #86efac;border-radius:10px;padding:10px 14px;font-size:13px;color:#047857;font-weight:600;text-align:center;animation:chatFadeIn .3s ease-out;z-index:10">&#10004; הכלי עודכן בהצלחה!</div>'
            '</div>'
        )

        chat_js = (
            '<script>'
            'var TR_SESSION_ID = ' + str(chat_session_id) + ';'
            'var TR_TOOL_ID = ' + str(tool_id) + ';'
            'var trSending = false;'
            'var trContextLoaded = false;'
            ''
            'function trScrollBottom(){'
            '  var el = document.getElementById("trChatMessages");'
            '  if(el) el.scrollTop = el.scrollHeight;'
            '}'
            ''
            'function trEsc(t){'
            '  var d = document.createElement("div");'
            '  d.appendChild(document.createTextNode(t));'
            '  return d.innerHTML;'
            '}'
            ''
            'function trFormatAi(text){'
            '  var h = trEsc(text);'
            '  h = h.replace(/```json/g, \'<pre style="background:#0f172a;color:#38bdf8;padding:12px;border-radius:10px;font-size:12px;overflow-x:auto;direction:ltr;text-align:left">\');'
            '  h = h.replace(/```/g, "</pre>");'
            '  h = h.replace(/\\n/g, "<br>");'
            '  return h;'
            '}'
            ''
            'function trAddUser(text){'
            '  var el = document.getElementById("trChatMessages");'
            '  var div = document.createElement("div");'
            '  div.className = "tr-chat-bubble tr-chat-user";'
            '  div.innerHTML = \'<div class="tr-chat-bubble-inner tr-chat-user-inner">\' + trEsc(text).replace(/\\n/g,"<br>") + "</div>";'
            '  el.appendChild(div);'
            '  trScrollBottom();'
            '}'
            ''
            'function trAddAi(text){'
            '  var el = document.getElementById("trChatMessages");'
            '  var div = document.createElement("div");'
            '  div.className = "tr-chat-bubble tr-chat-ai";'
            '  div.innerHTML = \'<div class="tr-chat-bubble-inner tr-chat-ai-inner">\' + trFormatAi(text) + "</div>";'
            '  el.appendChild(div);'
            '  trScrollBottom();'
            '}'
            ''
            'function trAddToolReady(toolDefJson){'
            '  var el = document.getElementById("trChatMessages");'
            '  var div = document.createElement("div");'
            '  div.className = "tr-chat-tool-ready";'
            '  div.innerHTML = \'<div style="font-size:13px;font-weight:700;color:#047857;margin-bottom:6px">&#127881; הכלי מוכן!</div>\''
            '    + \'<button type="button" class="btn btn-blue" style="border-radius:10px;font-size:12px;background:#047857" \''
            '    + \'onclick=\\\'updateToolDef(\\\"\' + toolDefJson.replace(/\\\\/g,"\\\\\\\\").replace(/"/g,\'\\\\&quot;\') + \'\\\")\\\'>\''
            '    + \'&#128190; עדכן את הכלי</button>\';'
            '  el.appendChild(div);'
            '  trScrollBottom();'
            '}'
            ''
            'function trShowTyping(){'
            '  trRemoveTyping();'
            '  var el = document.getElementById("trChatMessages");'
            '  var div = document.createElement("div");'
            '  div.id = "trTyping";'
            '  div.className = "tr-typing-indicator";'
            '  div.innerHTML = \'<div class="tr-typing-bubble"><div class="tr-typing-dots"><span></span><span></span><span></span></div><span class="tr-typing-label">AI מעבד...</span></div>\';'
            '  el.appendChild(div);'
            '  trScrollBottom();'
            '}'
            ''
            'function trRemoveTyping(){'
            '  var ti = document.getElementById("trTyping");'
            '  if(ti) ti.remove();'
            '}'
            ''
            'function trSetLoading(on){'
            '  trSending = on;'
            '  var inp = document.getElementById("trChatInput");'
            '  var btn = document.getElementById("trChatSendBtn");'
            '  var btnText = document.getElementById("trSendBtnText");'
            '  var btnSpin = document.getElementById("trSendBtnSpinner");'
            '  inp.disabled = on;'
            '  btn.disabled = on;'
            '  if(on){ btnText.textContent = "ממתין..."; btnSpin.style.display = "inline-block"; }'
            '  else { btnText.textContent = "שלח"; btnSpin.style.display = "none"; inp.focus(); }'
            '}'
            ''
            'function trSendChat(){'
            '  if(trSending) return;'
            '  var inp = document.getElementById("trChatInput");'
            '  var msg = inp.value.trim();'
            '  if(!msg) return;'
            '  inp.value = "";'
            '  inp.style.height = "42px";'
            '  trAddUser(msg);'
            '  var body = {session_id: TR_SESSION_ID, message: msg};'
            '  /* On first message, auto-load tool context */'
            '  if(!trContextLoaded){'
            '    body.tool_context = {'
            '      type: "marketplace",'
            '      name: ' + json.dumps(tool["name"], ensure_ascii=False) + ','
            '      description: ' + json.dumps(tool["description"] or "", ensure_ascii=False) + ','
            '      definition: ' + (tool["definition_json"] or "{}") + ''
            '    };'
            '    trContextLoaded = true;'
            '  }'
            '  trSetLoading(true);'
            '  trShowTyping();'
            '  fetch("/tools/create/chat", {'
            '    method: "POST",'
            '    headers: {"Content-Type": "application/json", "X-Requested-With": "XMLHttpRequest"},'
            '    body: JSON.stringify(body),'
            '    redirect: "error"'
            '  })'
            '  .then(function(r){ if(!r.ok) throw new Error("HTTP " + r.status); return r.json(); })'
            '  .then(function(data){'
            '    trRemoveTyping();'
            '    if(data.error){ trAddAi("שגיאה: " + data.error); }'
            '    else {'
            '      trAddAi(data.assistant_message);'
            '      if(data.has_tool && data.tool_definition) trAddToolReady(data.tool_definition);'
            '    }'
            '    trSetLoading(false);'
            '  })'
            '  .catch(function(err){'
            '    trRemoveTyping();'
            '    trAddAi("שגיאה בתקשורת. נסה שוב.");'
            '    trSetLoading(false);'
            '  });'
            '}'
            ''
            'function updateToolDef(defJson){'
            '  var parsed;'
            '  try { parsed = JSON.parse(defJson); } catch(e) { alert("שגיאה בפרסור הגדרת הכלי"); return; }'
            '  fetch("/marketplace/tool/" + TR_TOOL_ID + "/update-definition", {'
            '    method: "POST",'
            '    headers: {"Content-Type": "application/json", "X-Requested-With": "XMLHttpRequest"},'
            '    body: JSON.stringify({definition: parsed})'
            '  })'
            '  .then(function(r){ return r.json(); })'
            '  .then(function(data){'
            '    if(data.ok){'
            '      var toast = document.getElementById("trUpdateToast");'
            '      toast.style.display = "block";'
            '      setTimeout(function(){ toast.style.display = "none"; }, 3000);'
            '    } else {'
            '      alert(data.error || "שגיאה בעדכון הכלי");'
            '    }'
            '  })'
            '  .catch(function(){ alert("שגיאה בתקשורת"); });'
            '}'
            ''
            '/* Enter to send, Shift+Enter newline, auto-resize */'
            'document.getElementById("trChatInput").addEventListener("keydown", function(e){'
            '  if(e.key === "Enter" && !e.shiftKey){ e.preventDefault(); trSendChat(); }'
            '});'
            'document.getElementById("trChatInput").addEventListener("input", function(){'
            '  this.style.height = "42px";'
            '  this.style.height = Math.min(this.scrollHeight, 120) + "px";'
            '});'
            ''
            '/* Toggle chat panel */'
            'function toggleChatPanel(){'
            '  var panel = document.getElementById("trChatPanel");'
            '  var btn = document.getElementById("trToggleBtn");'
            '  if(panel.classList.contains("tr-chat-hidden")){'
            '    panel.classList.remove("tr-chat-hidden");'
            '    btn.innerHTML = "&#10005; סגור צ\\\'אט";'
            '  } else {'
            '    panel.classList.add("tr-chat-hidden");'
            '    btn.innerHTML = "&#129302; שיפור עם AI";'
            '  }'
            '}'
            ''
            'trScrollBottom();'
            '</script>'
        )

    # ── Styles ──
    styles = (
        '<style>'
        # Split layout
        '.tr-layout{display:flex;gap:16px;align-items:flex-start}'
        '.tr-main{flex:1;min-width:0}'
        '.tr-chat-panel{width:420px;flex-shrink:0;border:1px solid #dbeafe;border-radius:16px;overflow:hidden;background:#fff;'
        'box-shadow:0 4px 24px rgba(37,99,235,0.08);display:flex;flex-direction:column;height:calc(100vh - 160px);max-height:680px;position:relative;transition:width .3s,opacity .3s}'
        '.tr-chat-hidden{width:0;min-width:0;overflow:hidden;opacity:0;border:none;padding:0;margin:0}'
        # Top bar
        '.tr-topbar{display:flex;align-items:center;justify-content:space-between;flex-wrap:wrap;gap:8px;margin-bottom:1rem}'
        '.tr-toggle-btn{border:1.5px solid #2563eb;background:white;color:#2563eb;border-radius:10px;padding:6px 14px;font-size:13px;font-weight:600;cursor:pointer;transition:all .2s;font-family:inherit;white-space:nowrap}'
        '.tr-toggle-btn:hover{background:#eff6ff}'
        # Chat header
        '.tr-chat-header{background:linear-gradient(135deg,#1e3a8a,#2563eb);padding:12px 16px;color:white;flex-shrink:0}'
        '.tr-chat-header-title{font-size:16px;font-weight:800;margin-bottom:2px}'
        '.tr-chat-header-sub{font-size:12px;opacity:.85}'
        # Messages
        '.tr-chat-messages{flex:1;overflow-y:auto;padding:12px 14px;background:#fafcff;scroll-behavior:smooth}'
        '.tr-chat-messages::-webkit-scrollbar{width:5px}'
        '.tr-chat-messages::-webkit-scrollbar-track{background:transparent}'
        '.tr-chat-messages::-webkit-scrollbar-thumb{background:#cbd5e1;border-radius:3px}'
        # Bubbles
        '.tr-chat-bubble{margin-bottom:10px;display:flex;animation:trFadeIn .3s ease-out}'
        '.tr-chat-user{justify-content:flex-end}'
        '.tr-chat-ai{justify-content:flex-start}'
        '.tr-chat-bubble-inner{max-width:90%;padding:10px 14px;font-size:13px;line-height:1.7;border-radius:14px;word-wrap:break-word}'
        '.tr-chat-user-inner{background:linear-gradient(135deg,#dbeafe,#e0e7ff);color:#1e3a5c;border-bottom-left-radius:4px}'
        '.tr-chat-ai-inner{background:#eff6ff;color:#1e3a8a;border-bottom-right-radius:4px}'
        '.tr-chat-tool-ready{background:#ecfdf5;border:1px solid #86efac;border-radius:12px;padding:10px 14px;margin-bottom:10px;animation:trFadeIn .3s ease-out}'
        # Typing indicator
        '.tr-typing-indicator{display:flex;align-items:center;gap:6px;margin-bottom:10px;animation:trFadeIn .3s ease-out}'
        '.tr-typing-bubble{background:#eff6ff;border-radius:14px;padding:10px 14px;display:flex;align-items:center;gap:8px;border-bottom-right-radius:4px}'
        '.tr-typing-dots{display:flex;gap:3px;align-items:center}'
        '.tr-typing-dots span{width:6px;height:6px;border-radius:50%;background:#2563eb;animation:trBounce 1.4s infinite ease-in-out}'
        '.tr-typing-dots span:nth-child(2){animation-delay:.2s}'
        '.tr-typing-dots span:nth-child(3){animation-delay:.4s}'
        '.tr-typing-label{font-size:12px;color:#2563eb;font-weight:600}'
        # Input area
        '.tr-chat-input-area{border-top:1px solid #e2e8f0;background:white;padding:10px 12px;flex-shrink:0}'
        '.tr-chat-input-row{display:flex;gap:8px;align-items:flex-end}'
        '.tr-chat-textarea{flex:1;padding:8px 12px;border:1.5px solid #e2e8f0;border-radius:10px;font-size:13px;font-family:inherit;resize:none;max-height:100px;min-height:38px;line-height:1.5;outline:none;transition:border-color .2s,box-shadow .2s}'
        '.tr-chat-textarea:focus{border-color:#2563eb;box-shadow:0 0 0 3px rgba(37,99,235,0.1)}'
        '.tr-chat-textarea:disabled{background:#f8fafc;color:#94a3b8}'
        '.tr-chat-send-btn{border:none;background:linear-gradient(135deg,#1e3a8a,#2563eb);color:white;border-radius:10px;padding:8px 14px;font-size:13px;font-weight:700;cursor:pointer;transition:all .2s;display:flex;align-items:center;gap:5px;white-space:nowrap;font-family:inherit;min-height:38px}'
        '.tr-chat-send-btn:hover:not(:disabled){background:linear-gradient(135deg,#1e3070,#1d4ed8);transform:translateY(-1px);box-shadow:0 2px 8px rgba(37,99,235,0.3)}'
        '.tr-chat-send-btn:disabled{opacity:0.6;cursor:not-allowed}'
        '.tr-chat-send-spinner{display:inline-block;width:14px;height:14px;border:2px solid rgba(255,255,255,0.3);border-top-color:white;border-radius:50%;animation:trSpin .7s linear infinite}'
        # Animations
        '@keyframes trFadeIn{from{opacity:0;transform:translateY(6px)}to{opacity:1;transform:translateY(0)}}'
        '@keyframes trBounce{0%,80%,100%{transform:scale(0.6);opacity:0.4}40%{transform:scale(1);opacity:1}}'
        '@keyframes trSpin{to{transform:rotate(360deg)}}'
        # Responsive: stack on mobile
        '@media(max-width:900px){'
        '.tr-layout{flex-direction:column}'
        '.tr-chat-panel{width:100%;max-height:500px}'
        '.tr-chat-hidden{height:0;max-height:0}'
        '}'
        '</style>'
    )

    # ── Assemble body ──
    toggle_btn = ""
    if is_creator:
        toggle_btn = '<button type="button" id="trToggleBtn" class="tr-toggle-btn" onclick="toggleChatPanel()">&#10005; סגור צ\'אט</button>'

    body = (
        styles
        + '<div class="tr-topbar">'
        '<div><a href="' + back_link + '" style="color:#2563eb;font-size:13px;text-decoration:none">&#8592; ' + back_label + '</a></div>'
        + toggle_btn
        + '</div>'
        '<div class="tr-layout">'
        '<div class="tr-main">' + left_panel + '</div>'
        + right_panel
        + '</div>'
        + chat_js
    )

    return render(esc(tool["name"] or "כלי"), body, lang=lang, topbar_greeting=text["topbar_greeting"], logout_label=text["logout"], show_lang_switch=True)


@app.route("/marketplace/tool/<int:tool_id>/update-definition", methods=["POST"])
@login_required
def update_tool_definition(tool_id):
    """AJAX endpoint: update a marketplace tool's definition (creator only)."""
    with get_db() as db:
        tool = db.execute("SELECT * FROM marketplace_tools WHERE id=?", (tool_id,)).fetchone()
        if not tool:
            return jsonify({"ok": False, "error": "הכלי לא נמצא"}), 404
        if tool["creator_id"] != session["user_id"]:
            return jsonify({"ok": False, "error": "אין הרשאה"}), 403

    data = request.get_json(silent=True) or {}
    new_def = data.get("definition")
    if not new_def or not isinstance(new_def, dict):
        return jsonify({"ok": False, "error": "הגדרה לא תקינה"}), 400

    now = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    name = new_def.get("name", tool["name"])
    description = new_def.get("description", tool["description"] or "")
    icon = new_def.get("icon", tool["icon"] or "🔧")
    category = new_def.get("category", tool["category"] or "general")
    if category not in TOOL_CATEGORIES:
        category = "general"

    with get_db() as db:
        db.execute(
            """UPDATE marketplace_tools
               SET definition_json=?, name=?, description=?, icon=?, category=?, updated_at=?
               WHERE id=?""",
            (json.dumps(new_def, ensure_ascii=False), name, description, icon, category, now, tool_id),
        )
        db.commit()

    log_user_activity("update_tool", "עדכון כלי מרקטפלייס", "marketplace_" + str(tool_id), name, "")
    return jsonify({"ok": True})


# ── AI Tool Creation Chat ────────────────────────────────────────

@app.route("/tools/create", methods=["GET"])
@login_required
def tools_create():
    if session.get("is_admin"):
        return redirect("/admin")

    lang = get_flow_lang()
    text = get_flow_text(lang)

    # Create new chat session
    now = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    cutoff_24h = (datetime.now() - timedelta(hours=24)).strftime("%Y-%m-%d %H:%M:%S")
    cutoff_7d = (datetime.now() - timedelta(days=7)).strftime("%Y-%m-%d %H:%M:%S")
    cutoff_90d = (datetime.now() - timedelta(days=90)).strftime("%Y-%m-%d %H:%M:%S")
    with get_db() as db:
        # Auto-close stale active sessions (24h inactivity) and clear their data
        db.execute(
            "UPDATE tool_chat_sessions SET status='closed', messages_json='[]' WHERE status='active' AND updated_at < ?",
            (cutoff_24h,),
        )
        # Purge old closed/completed sessions (7 days)
        db.execute(
            "DELETE FROM tool_chat_sessions WHERE status IN ('closed','completed') AND updated_at < ?",
            (cutoff_7d,),
        )
        # Clean old token usage records (90 days, keep for billing)
        db.execute(
            "DELETE FROM chat_token_usage WHERE created_at < ?",
            (cutoff_90d,),
        )
        db.commit()

        # Check for existing active session
        existing = db.execute(
            "SELECT id FROM tool_chat_sessions WHERE user_id=? AND status='active' ORDER BY created_at DESC LIMIT 1",
            (session["user_id"],),
        ).fetchone()
        if existing:
            chat_session_id = existing["id"]
        else:
            db.execute(
                "INSERT INTO tool_chat_sessions(user_id, messages_json, status, created_at, updated_at) VALUES (?,?,?,?,?)",
                (session["user_id"], "[]", "active", now, now),
            )
            db.commit()
            chat_session_id = db.execute("SELECT id FROM tool_chat_sessions WHERE user_id=? ORDER BY id DESC LIMIT 1", (session["user_id"],)).fetchone()["id"]

    # Build existing messages HTML
    msgs_html = ""
    with get_db() as db:
        chat_data = db.execute("SELECT messages_json FROM tool_chat_sessions WHERE id=?", (chat_session_id,)).fetchone()
    existing_messages = []
    if chat_data:
        try:
            existing_messages = json.loads(chat_data["messages_json"])
        except (json.JSONDecodeError, TypeError):
            existing_messages = []

    for msg in existing_messages:
        if msg["role"] == "user":
            user_text = esc(msg["content"]).replace("\n", "<br>")
            msgs_html += '<div class="chat-bubble chat-user"><div class="chat-bubble-inner chat-user-inner">' + user_text + '</div></div>'
        else:
            import re as _render_re
            _cleaned = _render_re.sub(r"```developer_brief.*?```", "", msg["content"], flags=_render_re.DOTALL)
            _cleaned = _render_re.sub(r"---DEVELOPER_BRIEF---.*?---END_BRIEF---", "", _cleaned, flags=_render_re.DOTALL)
            content_html = esc(_cleaned).replace("```json", '<pre style="background:#0f172a;color:#38bdf8;padding:12px;border-radius:10px;font-size:12px;overflow-x:auto;direction:ltr;text-align:left">').replace("```", "</pre>").replace("\n", "<br>")
            msgs_html += '<div class="chat-bubble chat-ai"><div class="chat-bubble-inner chat-ai-inner">' + content_html + '</div></div>'
            # Check for developer brief first, then tool JSON
            brief = extract_developer_brief(msg["content"])
            if brief:
                brief_name = extract_brief_tool_name(brief)
                msgs_html += (
                    '<div class="chat-tool-ready" style="background:#eff6ff;border-color:#93c5fd" data-brief="' + esc(brief) + '" data-tool-name="' + esc(brief_name) + '">'
                    '<div style="font-size:14px;font-weight:700;color:#1e3a8a;margin-bottom:8px">&#128203; האפיון מוכן לשליחה למפתח</div>'
                    '<div style="font-size:13px;color:#475569;margin-bottom:12px">שם הכלי: <b>' + esc(brief_name) + '</b></div>'
                    '<div style="display:flex;gap:8px;flex-wrap:wrap">'
                    '<button type="button" class="btn btn-blue" style="border-radius:10px;font-size:13px;background:#1e3a8a" onclick="saveBrief(this)">&#128228; שלח לפיתוח</button>'
                    '</div></div>'
                )
            else:
                tool_def = extract_tool_json_from_message(msg["content"])
                if tool_def:
                    msgs_html += (
                        '<div class="chat-tool-ready">'
                        '<div style="font-size:14px;font-weight:700;color:#047857;margin-bottom:8px">&#127881; הכלי מוכן!</div>'
                        '<form method="post" action="/tools/create/save">'
                        '<input type="hidden" name="session_id" value="' + str(chat_session_id) + '">'
                        '<input type="hidden" name="definition" value="' + esc(json.dumps(tool_def, ensure_ascii=False)) + '">'
                        '<div style="display:flex;gap:8px;flex-wrap:wrap">'
                        '<button type="submit" name="action" value="save_draft" class="btn btn-blue" style="border-radius:10px;font-size:13px">&#128190; שמור ובדוק</button>'
                        '<button type="submit" name="action" value="publish" class="btn btn-blue" style="border-radius:10px;font-size:13px;background:#047857">&#128228; שלח לאישור ופרסום</button>'
                        '</div></form></div>'
                    )

    body = (
        '<style>'
        # Chat container styles
        '.chat-container{display:flex;flex-direction:column;height:calc(100vh - 140px);max-height:700px;border:1px solid #dbeafe;border-radius:16px;overflow:hidden;background:#fff;box-shadow:0 4px 24px rgba(37,99,235,0.08)}'
        '.chat-header{background:linear-gradient(135deg,#1e3a8a,#2563eb);padding:16px 20px;color:white;flex-shrink:0}'
        '.chat-header-title{font-size:18px;font-weight:800;margin-bottom:4px}'
        '.chat-header-sub{font-size:13px;opacity:.85}'
        # Messages area
        '.chat-messages{flex:1;overflow-y:auto;padding:16px 20px;background:#fafcff;scroll-behavior:smooth}'
        '.chat-messages::-webkit-scrollbar{width:6px}'
        '.chat-messages::-webkit-scrollbar-track{background:transparent}'
        '.chat-messages::-webkit-scrollbar-thumb{background:#cbd5e1;border-radius:3px}'
        # Chat bubbles
        '.chat-bubble{margin-bottom:12px;display:flex;animation:chatFadeIn .3s ease-out}'
        '.chat-user{justify-content:flex-end}'
        '.chat-ai{justify-content:flex-start}'
        '.chat-bubble-inner{max-width:85%;padding:12px 16px;font-size:14px;line-height:1.7;border-radius:16px;position:relative;word-wrap:break-word}'
        '.chat-user-inner{background:linear-gradient(135deg,#dbeafe,#e0e7ff);color:#1e3a5c;border-bottom-left-radius:4px}'
        '.chat-ai-inner{background:#eff6ff;color:#1e3a8a;border-bottom-right-radius:4px}'
        '.chat-tool-ready{background:#ecfdf5;border:1px solid #86efac;border-radius:14px;padding:12px 16px;margin-bottom:12px;animation:chatFadeIn .3s ease-out}'
        # Typing indicator
        '.typing-indicator{display:flex;align-items:center;gap:8px;margin-bottom:12px;animation:chatFadeIn .3s ease-out}'
        '.typing-indicator-bubble{background:#eff6ff;border-radius:16px;padding:12px 18px;display:flex;align-items:center;gap:10px;border-bottom-right-radius:4px}'
        '.typing-dots{display:flex;gap:4px;align-items:center}'
        '.typing-dots span{width:7px;height:7px;border-radius:50%;background:#2563eb;animation:typingBounce 1.4s infinite ease-in-out}'
        '.typing-dots span:nth-child(2){animation-delay:.2s}'
        '.typing-dots span:nth-child(3){animation-delay:.4s}'
        '.typing-label{font-size:13px;color:#2563eb;font-weight:600}'
        # Suggestion categories
        '#suggestionsPanel{padding:4px 16px 12px;display:flex;flex-direction:column;gap:6px}'
        '.suggest-cat{border:1px solid #e2e8f0;border-radius:12px;overflow:hidden;transition:all .2s}'
        '.suggest-cat:hover{border-color:#93c5fd}'
        '.suggest-cat-header{padding:10px 14px;font-size:14px;font-weight:700;cursor:pointer;background:#f8fafc;color:#1e3a8a;display:flex;align-items:center;gap:6px;user-select:none}'
        '.suggest-cat-header:hover{background:#eff6ff}'
        '.suggest-cat-items{display:none;padding:6px 10px 10px;flex-wrap:wrap;gap:6px}'
        '.suggest-cat.open .suggest-cat-items{display:flex}'
        '.suggest-btn{border:1px solid #dbeafe;background:white;color:#1e3a8a;border-radius:20px;padding:7px 14px;font-size:12.5px;cursor:pointer;transition:all .15s;font-family:inherit;white-space:nowrap;line-height:1.4}'
        '.suggest-btn:hover{background:#dbeafe;border-color:#93c5fd;transform:translateY(-1px);box-shadow:0 2px 6px rgba(37,99,235,0.12)}'
        # Input area
        '.chat-input-area{border-top:1px solid #e2e8f0;background:white;padding:12px 16px;flex-shrink:0}'
        '.chat-input-row{display:flex;gap:8px;align-items:flex-end}'
        '.chat-textarea{flex:1;padding:10px 14px;border:1.5px solid #e2e8f0;border-radius:12px;font-size:14px;font-family:inherit;resize:none;max-height:120px;min-height:42px;line-height:1.5;outline:none;transition:border-color .2s,box-shadow .2s}'
        '.chat-textarea:focus{border-color:#2563eb;box-shadow:0 0 0 3px rgba(37,99,235,0.1)}'
        '.chat-textarea:disabled{background:#f8fafc;color:#94a3b8}'
        '.chat-send-btn{border:none;background:linear-gradient(135deg,#1e3a8a,#2563eb);color:white;border-radius:12px;padding:10px 18px;font-size:14px;font-weight:700;cursor:pointer;transition:all .2s;display:flex;align-items:center;gap:6px;white-space:nowrap;font-family:inherit;min-height:42px}'
        '.chat-send-btn:hover:not(:disabled){background:linear-gradient(135deg,#1e3070,#1d4ed8);transform:translateY(-1px);box-shadow:0 2px 8px rgba(37,99,235,0.3)}'
        '.chat-send-btn:active:not(:disabled){transform:translateY(0)}'
        '.chat-send-btn:disabled{opacity:0.6;cursor:not-allowed;transform:none}'
        '.chat-send-spinner{display:inline-block;width:16px;height:16px;border:2px solid rgba(255,255,255,0.3);border-top-color:white;border-radius:50%;animation:spin .7s linear infinite}'
        '.chat-upload-row{display:flex;gap:8px;margin-bottom:8px}'
        '.chat-upload-btn{border:1.5px solid #e2e8f0;background:white;border-radius:10px;font-size:12px;padding:6px 12px;cursor:pointer;transition:all .2s;color:#64748b;font-family:inherit;white-space:nowrap;display:flex;align-items:center;gap:4px}'
        '.chat-upload-btn:hover{border-color:#2563eb;color:#2563eb;background:#f0f7ff}'
        '.chat-hint{font-size:11px;color:#94a3b8;margin-top:6px;text-align:center}'
        # Animations
        '@keyframes chatFadeIn{from{opacity:0;transform:translateY(8px)}to{opacity:1;transform:translateY(0)}}'
        '@keyframes typingBounce{0%,80%,100%{transform:scale(0.6);opacity:0.4}40%{transform:scale(1);opacity:1}}'
        '@keyframes spin{to{transform:rotate(360deg)}}'
        # Drop zone overlay
        '.chat-drop-overlay{display:none;position:absolute;inset:0;background:rgba(37,99,235,0.12);border:3px dashed #2563eb;border-radius:16px;z-index:50;align-items:center;justify-content:center;flex-direction:column;gap:10px}'
        '.chat-drop-overlay.active{display:flex}'
        '.chat-drop-overlay>*{pointer-events:none}'
        '.chat-drop-icon{font-size:48px;opacity:0.7}'
        '.chat-drop-label{font-size:16px;font-weight:700;color:#1e3a8a}'
        '.chat-drop-sub{font-size:13px;color:#64748b}'
        # File preview chips in bubbles
        '.chat-file-chips{display:flex;flex-wrap:wrap;gap:6px;margin-top:8px}'
        '.chat-file-chip{display:inline-flex;align-items:center;gap:6px;background:rgba(37,99,235,0.08);border:1px solid #dbeafe;border-radius:8px;padding:4px 10px;font-size:12px;color:#1e3a5c}'
        '.chat-file-chip img{width:40px;height:40px;object-fit:cover;border-radius:6px;border:1px solid #e2e8f0}'
        '.chat-file-chip .chip-icon{font-size:16px}'
        # Upload progress
        '.chat-upload-progress{display:flex;align-items:center;gap:8px;padding:8px 12px;margin-bottom:8px;background:#f0f7ff;border-radius:10px;font-size:13px;color:#2563eb;animation:chatFadeIn .3s ease-out}'
        '.chat-upload-progress .upload-spinner{display:inline-block;width:14px;height:14px;border:2px solid rgba(37,99,235,0.3);border-top-color:#2563eb;border-radius:50%;animation:spin .7s linear infinite}'
        # Container must be relative for overlay
        '.chat-container{position:relative}'
        # Tool picker styles
        '.tool-picker-overlay{display:none;position:absolute;inset:0;background:rgba(0,0,0,0.4);z-index:60;align-items:center;justify-content:center;border-radius:16px;overflow:hidden}'
        '.tool-picker-overlay.active{display:flex}'
        '.tool-picker-panel{background:white;border-radius:16px;width:92%;max-width:480px;max-height:80%;display:flex;flex-direction:column;box-shadow:0 12px 40px rgba(0,0,0,0.2);animation:chatFadeIn .3s ease-out}'
        '.tool-picker-header{padding:16px 20px;border-bottom:1px solid #e2e8f0;display:flex;justify-content:space-between;align-items:center}'
        '.tool-picker-header h3{margin:0;font-size:16px;font-weight:800;color:#1e3a8a}'
        '.tool-picker-close{background:none;border:none;font-size:20px;cursor:pointer;color:#64748b;padding:4px 8px;border-radius:8px}'
        '.tool-picker-close:hover{background:#f1f5f9;color:#1e3a8a}'
        '.tool-picker-body{overflow-y:auto;padding:12px 16px;flex:1}'
        '.tool-picker-section{margin-bottom:12px}'
        '.tool-picker-section-title{font-size:12px;font-weight:700;color:#64748b;margin-bottom:8px;text-transform:uppercase;letter-spacing:0.5px}'
        '.tool-picker-card{display:flex;align-items:center;gap:12px;padding:12px 14px;border:1.5px solid #e2e8f0;border-radius:12px;cursor:pointer;transition:all .2s;margin-bottom:8px;background:white}'
        '.tool-picker-card:hover{border-color:#2563eb;background:#f0f7ff;transform:translateY(-1px);box-shadow:0 4px 12px rgba(37,99,235,0.1)}'
        '.tool-picker-card-icon{font-size:28px;flex-shrink:0;width:40px;text-align:center}'
        '.tool-picker-card-info{flex:1;min-width:0}'
        '.tool-picker-card-name{font-size:14px;font-weight:700;color:#1e3a8a;margin-bottom:2px;white-space:nowrap;overflow:hidden;text-overflow:ellipsis}'
        '.tool-picker-card-desc{font-size:12px;color:#64748b;white-space:nowrap;overflow:hidden;text-overflow:ellipsis}'
        '.tool-picker-card-badge{font-size:10px;padding:2px 8px;border-radius:99px;font-weight:700;flex-shrink:0}'
        '.tool-picker-empty{text-align:center;padding:2rem;color:#94a3b8;font-size:14px}'
        '.tool-picker-loading{text-align:center;padding:2rem;color:#2563eb;font-size:14px}'
        '</style>'
        # Top bar
        '<div style="margin-bottom:1rem;display:flex;align-items:center;justify-content:space-between;flex-wrap:wrap;gap:8px">'
        '<div><a href="/marketplace" style="color:#2563eb;font-size:13px;text-decoration:none">&#8592; חזרה לשוק הכלים</a></div>'
        '<form method="post" action="/tools/create/reset"><button type="submit" class="btn btn-gray" style="font-size:12px;border-radius:8px">&#128260; שיחה חדשה</button></form>'
        '</div>'
        # Chat container
        '<div class="chat-container" id="chatContainer">'
        # Drop zone overlay
        '<div class="chat-drop-overlay" id="dropOverlay">'
        '<div class="chat-drop-icon">&#128206;</div>'
        '<div class="chat-drop-label">שחרר קבצים כאן</div>'
        '<div class="chat-drop-sub">Excel, CSV, Word, PDF, תמונות ועוד</div>'
        '</div>'
        # Tool picker overlay
        '<div class="tool-picker-overlay" id="toolPickerOverlay">'
        '<div class="tool-picker-panel">'
        '<div class="tool-picker-header">'
        '<h3>בחירת כלי לשיפור</h3>'
        '<button class="tool-picker-close" onclick="hideToolPicker()">&#10005;</button>'
        '</div>'
        '<div class="tool-picker-body" id="toolPickerBody">'
        '<div class="tool-picker-loading">&#8987; טוען כלים...</div>'
        '</div>'
        '</div>'
        '</div>'
        '<div class="chat-header">'
        '<div class="chat-header-title">&#129302; העוזר החכם של Scriptly</div>'
        '<div class="chat-header-sub">נתח קבצים, בנה כלים חדשים, או שפר כלים קיימים — הכל בשיחה אחת</div>'
        '</div>'
        '<div class="chat-messages" id="chatMessages">'
        # Welcome message
        '<div class="chat-bubble chat-ai">'
        '<div class="chat-bubble-inner chat-ai-inner">'
        '&#128075; שלום! אני העוזר החכם של Scriptly.<br><br>'
        'אני יכול לעזור לך ב-3 דרכים:<br><br>'
        '&#128269; <b>נתח קובץ עכשיו</b><br>'
        '<span style="color:#475569;font-size:13px">פשוט תעלה קובץ ותשאל אותי כל שאלה — אענה מיד.</span><br><br>'
        '&#128295; <b>שפר כלי קיים</b><br>'
        '<span style="color:#475569;font-size:13px">רוצה לשנות כלי שכבר יש לך? אני אעשה שינויים פשוטים מיד, או אאפיין יחד איתך שינויים מורכבים.</span><br><br>'
        '&#10024; <b>בנה כלי חדש</b><br>'
        '<span style="color:#475569;font-size:13px">ספר לי מה אתה צריך — אשאל אותך שאלות מדויקות ואכין בריף מלא שהמפתח יוכל לבנות בשבילך.</span><br><br>'
        '&#128161; <b>לא יודע מאיפה להתחיל? הנה מה שאחרים מבקשים:</b>'
        '</div></div>'
        # Suggestion categories
        '<div id="suggestionsPanel">'
        '<div class="suggest-cat" onclick="toggleCat(this)">'
        '<div class="suggest-cat-header">&#128336; נוכחות ושעות</div>'
        '<div class="suggest-cat-items">'
        '<button class="suggest-btn" onclick="pickSuggestion(this)">מי לא הגיע היום?</button>'
        '<button class="suggest-btn" onclick="pickSuggestion(this)">מי קרוב לחריגת שעות נוספות החודש?</button>'
        '<button class="suggest-btn" onclick="pickSuggestion(this)">כמה שעות עבד כל עובד השבוע?</button>'
        '<button class="suggest-btn" onclick="pickSuggestion(this)">מי עבד בשבת או בחג?</button>'
        '<button class="suggest-btn" onclick="pickSuggestion(this)">הצג לי עובדים עם יותר מ-X שעות נוספות החודש</button>'
        '</div></div>'
        '<div class="suggest-cat" onclick="toggleCat(this)">'
        '<div class="suggest-cat-header">&#128176; שכר וחישובים</div>'
        '<div class="suggest-cat-items">'
        '<button class="suggest-btn" onclick="pickSuggestion(this)">כמה שעות נוספות צריך לשלם החודש?</button>'
        '<button class="suggest-btn" onclick="pickSuggestion(this)">מי זכאי לבונוס לפי הביצועים שלו?</button>'
        '<button class="suggest-btn" onclick="pickSuggestion(this)">חשב לי את עלות כוח האדם לפי מחלקה</button>'
        '<button class="suggest-btn" onclick="pickSuggestion(this)">מי קיבל העלאת שכר השנה ומי לא?</button>'
        '<button class="suggest-btn" onclick="pickSuggestion(this)">מה ההפרש בין שכר גברים לנשים באותה רמה?</button>'
        '</div></div>'
        '<div class="suggest-cat" onclick="toggleCat(this)">'
        '<div class="suggest-cat-header">&#128197; היעדרויות וחופשות</div>'
        '<div class="suggest-cat-items">'
        '<button class="suggest-btn" onclick="pickSuggestion(this)">מי ניצל יותר מ-X ימי מחלה השנה?</button>'
        '<button class="suggest-btn" onclick="pickSuggestion(this)">מי מסיים את ימי החופשה שלו בסוף השנה?</button>'
        '<button class="suggest-btn" onclick="pickSuggestion(this)">הצג לי תמונת היעדרויות לפי מחלקה</button>'
        '<button class="suggest-btn" onclick="pickSuggestion(this)">מי בחל&quot;ת או בחופשת לידה כרגע?</button>'
        '<button class="suggest-btn" onclick="pickSuggestion(this)">כמה ימי מחלה נוצלו לעומת אשתקד?</button>'
        '</div></div>'
        '<div class="suggest-cat" onclick="toggleCat(this)">'
        '<div class="suggest-cat-header">&#128101; כוח אדם וניהול</div>'
        '<div class="suggest-cat-items">'
        '<button class="suggest-btn" onclick="pickSuggestion(this)">מי מסיים חוזה בחודשיים הקרובים?</button>'
        '<button class="suggest-btn" onclick="pickSuggestion(this)">כמה עובדים יש לי בכל מחלקה?</button>'
        '<button class="suggest-btn" onclick="pickSuggestion(this)">מי עובד אצלנו יותר מ-5 שנים?</button>'
        '<button class="suggest-btn" onclick="pickSuggestion(this)">הצג לי עובדים חדשים מהחצי שנה האחרונה</button>'
        '<button class="suggest-btn" onclick="pickSuggestion(this)">מי הגיע לגיל פרישה השנה?</button>'
        '</div></div>'
        '<div class="suggest-cat" onclick="toggleCat(this)">'
        '<div class="suggest-cat-header">&#9878;&#65039; ציות וחוקי עבודה</div>'
        '<div class="suggest-cat-items">'
        '<button class="suggest-btn" onclick="pickSuggestion(this)">מי חרג ממכסת שעות נוספות חוקית?</button>'
        '<button class="suggest-btn" onclick="pickSuggestion(this)">מי לא קיבל מנוחה שבועית כחוק?</button>'
        '<button class="suggest-btn" onclick="pickSuggestion(this)">הצג לי עובדי לילה שחרגו מ-58 שעות שבועיות</button>'
        '<button class="suggest-btn" onclick="pickSuggestion(this)">מי לא קיבל תלוש שכר החודש?</button>'
        '</div></div>'
        '<div class="suggest-cat" onclick="toggleCat(this)">'
        '<div class="suggest-cat-header">&#128202; ניתוח וטרנדים</div>'
        '<div class="suggest-cat-items">'
        '<button class="suggest-btn" onclick="pickSuggestion(this)">השווה נוכחות החודש לעומת החודש שעבר</button>'
        '<button class="suggest-btn" onclick="pickSuggestion(this)">מאיזו מחלקה יש הכי הרבה תחלופת עובדים?</button>'
        '<button class="suggest-btn" onclick="pickSuggestion(this)">הצג לי גרף עלויות כוח אדם לאורך השנה</button>'
        '<button class="suggest-btn" onclick="pickSuggestion(this)">מי השתפר/הידרדר בביצועים לאורך זמן?</button>'
        '</div></div>'
        '</div>'
        + msgs_html
        + '</div>'
        # Input area
        '<div class="chat-input-area">'
        '<div class="chat-upload-row">'
        '<input type="file" name="sample_file" id="sampleFile" accept=".xls,.xlsx,.csv,.png,.jpg,.jpeg,.gif,.bmp,.webp,.doc,.docx,.pdf,.txt,.rtf,.dat,.xml,.json,.ppt,.pptx" style="display:none" multiple '
        'onchange="handleFileSelect(this.files)">'
        '<button type="button" id="uploadBtn" onclick="document.getElementById(\'sampleFile\').click()" class="chat-upload-btn">&#128206; העלה קבצים</button>'
        '<button type="button" id="improveToolBtn" onclick="showToolPicker()" class="chat-upload-btn" style="border-color:#047857;color:#047857">&#9998; שיפור כלי קיים</button>'
        '</div>'
        '<div id="pendingFilesBar" style="display:none;padding:6px 0;margin-bottom:4px">'
        '<div style="font-size:11px;color:#64748b;margin-bottom:4px">קבצים מצורפים (כתוב הוראות ושלח):</div>'
        '<div id="pendingFilesList" class="chat-file-chips"></div>'
        '</div>'
        '<div class="chat-input-row">'
        '<textarea id="chatInput" class="chat-textarea" placeholder="כתוב הודעה... (Enter לשליחה, Shift+Enter לשורה חדשה)" rows="1"></textarea>'
        '<button type="button" id="chatSendBtn" class="chat-send-btn" onclick="sendChatMessage()">'
        '<span id="sendBtnText">שלח</span>'
        '<span id="sendBtnSpinner" class="chat-send-spinner" style="display:none"></span>'
        '</button>'
        '</div>'
        '<div class="chat-hint">Enter לשליחה &middot; Shift+Enter לשורה חדשה &middot; גרור קבצים או הדבק תמונה &middot; Excel, CSV, Word, PDF, תמונות ועוד</div>'
        '</div>'
        '</div>'
        # JavaScript
        '<script>'
        'var CHAT_SESSION_ID = "' + str(chat_session_id) + '";'
        'var chatSending = false;'
        'var IMAGE_EXTS = ["png","jpg","jpeg","gif","bmp","webp"];'
        'var pendingFiles = [];'  # {name, type, analysis, dataUrl}
        'var pendingAnalysis = "";'
        ''
        'function scrollToBottom(){'
        '  var el = document.getElementById("chatMessages");'
        '  if(el) el.scrollTop = el.scrollHeight;'
        '}'
        ''
        'function toggleCat(el){'
        '  if(event.target.classList.contains("suggest-btn")) return;'
        '  el.classList.toggle("open");'
        '}'
        ''
        'function pickSuggestion(btn){'
        '  event.stopPropagation();'
        '  var text = btn.innerText.trim();'
        '  var inp = document.getElementById("chatInput");'
        '  inp.value = text;'
        '  /* Hide suggestions panel */'
        '  var panel = document.getElementById("suggestionsPanel");'
        '  if(panel) panel.style.display = "none";'
        '  sendChatMessage();'
        '}'
        ''
        'function escapeHtml(t){'
        '  var d = document.createElement("div");'
        '  d.appendChild(document.createTextNode(t));'
        '  return d.innerHTML;'
        '}'
        ''
        'function formatAiMessage(text){'
        '  var h = escapeHtml(text);'
        '  /* Hide developer_brief blocks from user — the brief card handles display */'
        '  h = h.replace(/```developer_brief[\\s\\S]*?```/g, "");'
        '  h = h.replace(/---DEVELOPER_BRIEF---[\\s\\S]*?---END_BRIEF---/g, "");'
        '  h = h.replace(/```json/g, \'<pre style="background:#0f172a;color:#38bdf8;padding:12px;border-radius:10px;font-size:12px;overflow-x:auto;direction:ltr;text-align:left">\');'
        '  h = h.replace(/```/g, "</pre>");'
        '  h = h.replace(/\\n/g, "<br>");'
        '  return h;'
        '}'
        ''
        'function getFileExt(name){'
        '  if(!name || name.indexOf(".")===-1) return "";'
        '  return name.split(".").pop().toLowerCase();'
        '}'
        ''
        'function isImageFile(name){'
        '  return IMAGE_EXTS.indexOf(getFileExt(name)) !== -1;'
        '}'
        ''
        'function buildFileChips(files, dataUrls){'
        '  if(!files || !files.length) return "";'
        '  var html = \'<div class="chat-file-chips">\';'
        '  for(var i=0;i<files.length;i++){'
        '    var f = files[i];'
        '    var name = f.name || f;'
        '    if(isImageFile(name) && dataUrls && dataUrls[i]){'
        '      html += \'<span class="chat-file-chip"><img src="\'+dataUrls[i]+\'" alt="">\'+escapeHtml(name)+\'</span>\';'
        '    } else {'
        '      html += \'<span class="chat-file-chip"><span class="chip-icon">&#128196;</span>\'+escapeHtml(name)+\'</span>\';'
        '    }'
        '  }'
        '  html += "</div>";'
        '  return html;'
        '}'
        ''
        'function addUserBubble(text, extraHtml){'
        '  var el = document.getElementById("chatMessages");'
        '  var div = document.createElement("div");'
        '  div.className = "chat-bubble chat-user";'
        '  var inner = escapeHtml(text).replace(/\\n/g,"<br>");'
        '  if(extraHtml) inner += extraHtml;'
        '  div.innerHTML = \'<div class="chat-bubble-inner chat-user-inner">\' + inner + "</div>";'
        '  el.appendChild(div);'
        '  scrollToBottom();'
        '}'
        ''
        'function addAiBubble(text){'
        '  var el = document.getElementById("chatMessages");'
        '  var div = document.createElement("div");'
        '  div.className = "chat-bubble chat-ai";'
        '  div.innerHTML = \'<div class="chat-bubble-inner chat-ai-inner">\' + formatAiMessage(text) + "</div>";'
        '  el.appendChild(div);'
        '  scrollToBottom();'
        '}'
        ''
        'function addToolReady(toolDefJson){'
        '  var el = document.getElementById("chatMessages");'
        '  var div = document.createElement("div");'
        '  div.className = "chat-tool-ready";'
        '  div.innerHTML = \'<div style="font-size:14px;font-weight:700;color:#047857;margin-bottom:8px">&#127881; הכלי מוכן!</div>\''
        '    + \'<form method="post" action="/tools/create/save">\''
        '    + \'<input type="hidden" name="session_id" value="\' + CHAT_SESSION_ID + \'">\''
        '    + \'<input type="hidden" name="definition" value="\' + toolDefJson.replace(/"/g, "&quot;") + \'">\''
        '    + \'<div style="display:flex;gap:8px;flex-wrap:wrap">\''
        '    + \'<button type="submit" name="action" value="save_draft" class="btn btn-blue" style="border-radius:10px;font-size:13px">&#128190; שמור ובדוק</button>\''
        '    + \'<button type="submit" name="action" value="publish" class="btn btn-blue" style="border-radius:10px;font-size:13px;background:#047857">&#128228; שלח לאישור ופרסום</button>\''
        '    + \'</div></form>\';'
        '  el.appendChild(div);'
        '  scrollToBottom();'
        '}'
        ''
        'function addBriefReady(briefText, toolName){'
        '  var el = document.getElementById("chatMessages");'
        '  var div = document.createElement("div");'
        '  div.className = "chat-tool-ready";'
        '  div.style.background = "#eff6ff";'
        '  div.style.borderColor = "#93c5fd";'
        '  var briefSaved = false;'
        '  div.innerHTML = \'<div style="font-size:14px;font-weight:700;color:#1e3a8a;margin-bottom:8px">&#128203; האפיון מוכן לשליחה למפתח</div>\''
        '    + \'<div style="font-size:13px;color:#475569;margin-bottom:12px">שם הכלי: <b>\' + escapeHtml(toolName) + \'</b></div>\''
        '    + \'<div style="display:flex;gap:8px;flex-wrap:wrap">\''
        '    + \'<button type="button" class="btn btn-blue" style="border-radius:10px;font-size:13px;background:#1e3a8a" onclick="saveBrief(this)">&#128228; שלח לפיתוח</button>\''
        '    + \'</div>\';'
        '  div.setAttribute("data-brief", briefText);'
        '  div.setAttribute("data-tool-name", toolName);'
        '  el.appendChild(div);'
        '  scrollToBottom();'
        '}'
        ''
        'function saveBrief(btn){'
        '  var card = btn.closest(".chat-tool-ready");'
        '  if(!card) return;'
        '  var brief = card.getAttribute("data-brief");'
        '  var toolName = card.getAttribute("data-tool-name");'
        '  if(!brief) return;'
        '  btn.disabled = true;'
        '  btn.innerHTML = \'<span class="chat-send-spinner" style="display:inline-block;width:14px;height:14px"></span> שולח...\';'
        '  fetch("/tools/create/save-brief", {'
        '    method: "POST",'
        '    headers: {"Content-Type":"application/json","X-Requested-With":"XMLHttpRequest"},'
        '    body: JSON.stringify({brief_text: brief, tool_name: toolName, session_id: CHAT_SESSION_ID, request_type: "new"})'
        '  })'
        '  .then(function(r){return r.json();})'
        '  .then(function(data){'
        '    if(data.ok){'
        '      card.style.background = "#ecfdf5";'
        '      card.style.borderColor = "#86efac";'
        '      card.innerHTML = \'<div style="font-size:14px;font-weight:700;color:#047857">&#9989; \' + escapeHtml(data.message || "הבקשה נשלחה בהצלחה!") + \'</div>\';'
        '    } else {'
        '      btn.disabled = false;'
        '      btn.textContent = "&#128228; שלח לפיתוח";'
        '      addAiBubble("שגיאה: " + (data.error || "לא ידוע"));'
        '    }'
        '  })'
        '  .catch(function(){'
        '    btn.disabled = false;'
        '    btn.textContent = "&#128228; שלח לפיתוח";'
        '    addAiBubble("שגיאה בשליחה. נסה שוב.");'
        '  });'
        '}'
        ''
        'function showTypingIndicator(){'
        '  removeTypingIndicator();'
        '  var el = document.getElementById("chatMessages");'
        '  var div = document.createElement("div");'
        '  div.id = "typingIndicator";'
        '  div.className = "typing-indicator";'
        '  div.innerHTML = \'<div class="typing-indicator-bubble"><div class="typing-dots"><span></span><span></span><span></span></div><span class="typing-label">הבינה מעבדת...</span></div>\';'
        '  el.appendChild(div);'
        '  scrollToBottom();'
        '}'
        ''
        'function removeTypingIndicator(){'
        '  var ti = document.getElementById("typingIndicator");'
        '  if(ti) ti.remove();'
        '}'
        ''
        'function showUploadProgress(fileNames){'
        '  removeUploadProgress();'
        '  var el = document.getElementById("chatMessages");'
        '  var div = document.createElement("div");'
        '  div.id = "uploadProgress";'
        '  div.className = "chat-upload-progress";'
        '  var label = fileNames.length === 1 ? "מעלה " + escapeHtml(fileNames[0]) + "..." : "מעלה " + fileNames.length + " קבצים...";'
        '  div.innerHTML = \'<span class="upload-spinner"></span><span>\' + label + "</span>";'
        '  el.appendChild(div);'
        '  scrollToBottom();'
        '}'
        ''
        'function removeUploadProgress(){'
        '  var up = document.getElementById("uploadProgress");'
        '  if(up) up.remove();'
        '}'
        ''
        'function setLoading(on){'
        '  chatSending = on;'
        '  var inp = document.getElementById("chatInput");'
        '  var btn = document.getElementById("chatSendBtn");'
        '  var btnText = document.getElementById("sendBtnText");'
        '  var btnSpin = document.getElementById("sendBtnSpinner");'
        '  var upload = document.getElementById("uploadBtn");'
        '  inp.disabled = on;'
        '  btn.disabled = on;'
        '  if(upload) upload.disabled = on;'
        '  if(on){'
        '    btnText.textContent = "ממתין...";'
        '    btnSpin.style.display = "inline-block";'
        '  } else {'
        '    btnText.textContent = "שלח";'
        '    btnSpin.style.display = "none";'
        '    inp.focus();'
        '  }'
        '}'
        ''
        'function sendChatMessage(){'
        '  if(chatSending) return;'
        '  var inp = document.getElementById("chatInput");'
        '  var msg = inp.value.trim();'
        '  if(!msg) return;'
        '  /* Hide suggestions on first message */'
        '  var sp = document.getElementById("suggestionsPanel");'
        '  if(sp) sp.style.display = "none";'
        '  inp.value = "";'
        '  inp.style.height = "42px";'
        '  /* Build bubble with file chips if any */'
        '  var filesHtml = "";'
        '  if(pendingFiles.length > 0){'
        '    filesHtml = buildFileChips(pendingFiles, pendingFiles.map(function(f){return f.dataUrl;}));'
        '  }'
        '  addUserBubble(msg, filesHtml);'
        '  /* Build request body */'
        '  var body = {session_id: CHAT_SESSION_ID, message: msg};'
        '  if(pendingAnalysis) body.file_context = pendingAnalysis;'
        '  /* Clear pending files */'
        '  pendingFiles = [];'
        '  pendingAnalysis = "";'
        '  updatePendingFilesUI();'
        '  inp.placeholder = "כתוב הודעה... (Enter לשליחה, Shift+Enter לשורה חדשה)";'
        '  setLoading(true);'
        '  showTypingIndicator();'
        '  fetch("/tools/create/chat", {'
        '    method: "POST",'
        '    headers: {"Content-Type": "application/json", "X-Requested-With": "XMLHttpRequest"},'
        '    body: JSON.stringify(body),'
        '    redirect: "error"'
        '  })'
        '  .then(function(r){ if(!r.ok) throw new Error("HTTP " + r.status); return r.json(); })'
        '  .then(function(data){'
        '    removeTypingIndicator();'
        '    if(data.error){ addAiBubble("שגיאה: " + data.error); }'
        '    else {'
        '      addAiBubble(data.assistant_message);'
        '      if(data.has_brief && data.brief_text) addBriefReady(data.brief_text, data.brief_tool_name || "כלי ללא שם");'
        '      else if(data.has_tool && data.tool_definition) addToolReady(data.tool_definition);'
        '    }'
        '    setLoading(false);'
        '  })'
        '  .catch(function(err){'
        '    removeTypingIndicator();'
        '    addAiBubble("שגיאה בתקשורת. נסה שוב.");'
        '    setLoading(false);'
        '  });'
        '}'
        ''
        '/* ============ FILE UPLOAD VIA AJAX ============ */'
        'function updatePendingFilesUI(){'
        '  var bar = document.getElementById("pendingFilesBar");'
        '  var list = document.getElementById("pendingFilesList");'
        '  if(pendingFiles.length === 0){ bar.style.display="none"; list.innerHTML=""; return; }'
        '  bar.style.display="block";'
        '  var html = "";'
        '  for(var i=0;i<pendingFiles.length;i++){'
        '    var pf = pendingFiles[i];'
        '    if(pf.dataUrl){'
        '      html += \'<span class="chat-file-chip"><img src="\'+pf.dataUrl+\'" alt="">\'+escapeHtml(pf.name)+\'<span onclick="removePendingFile(\'+i+\')" style="cursor:pointer;margin-right:4px;color:#ef4444">&times;</span></span>\';'
        '    } else {'
        '      html += \'<span class="chat-file-chip"><span class="chip-icon">&#128196;</span>\'+escapeHtml(pf.name)+\'<span onclick="removePendingFile(\'+i+\')" style="cursor:pointer;margin-right:4px;color:#ef4444">&times;</span></span>\';'
        '    }'
        '  }'
        '  list.innerHTML = html;'
        '}'
        ''
        'function removePendingFile(idx){'
        '  pendingFiles.splice(idx, 1);'
        '  if(pendingFiles.length === 0) pendingAnalysis = "";'
        '  updatePendingFilesUI();'
        '}'
        ''
        'function uploadFiles(fileList){'
        '  if(chatSending || !fileList || fileList.length === 0) return;'
        '  var names = [];'
        '  var dataUrls = [];'
        '  var readCount = fileList.length;'
        ''
        '  function afterPreviews(){'
        '    showUploadProgress(names);'
        '    var fd = new FormData();'
        '    fd.append("session_id", CHAT_SESSION_ID);'
        '    for(var i=0;i<fileList.length;i++) fd.append("sample_file", fileList[i]);'
        '    fetch("/tools/create/upload", {'
        '      method: "POST",'
        '      headers: {"X-Requested-With": "XMLHttpRequest"},'
        '      body: fd,'
        '      redirect: "error"'
        '    })'
        '    .then(function(r){ if(!r.ok) throw new Error("HTTP " + r.status); return r.json(); })'
        '    .then(function(data){'
        '      removeUploadProgress();'
        '      if(data.error){ addAiBubble("שגיאה: " + data.error); return; }'
        '      /* Store analysis for next send */'
        '      if(pendingAnalysis) pendingAnalysis += "\\n\\n" + data.analysis;'
        '      else pendingAnalysis = data.analysis;'
        '      for(var i=0;i<data.file_names.length;i++){'
        '        pendingFiles.push({name: data.file_names[i], type: data.file_types[i], dataUrl: dataUrls[i] || null});'
        '      }'
        '      updatePendingFilesUI();'
        '      document.getElementById("chatInput").focus();'
        '      document.getElementById("chatInput").placeholder = "תאר מה אתה רוצה לעשות עם הקבצים...";'
        '    })'
        '    .catch(function(err){'
        '      removeUploadProgress();'
        '      addAiBubble("שגיאה בהעלאת הקובץ. נסה שוב.");'
        '    });'
        '  }'
        ''
        '  for(var i=0;i<fileList.length;i++){'
        '    names.push(fileList[i].name);'
        '    (function(idx){'
        '      if(isImageFile(fileList[idx].name)){'
        '        var reader = new FileReader();'
        '        reader.onload = function(e){ dataUrls[idx] = e.target.result; readCount--; if(readCount===0) afterPreviews(); };'
        '        reader.onerror = function(){ dataUrls[idx] = null; readCount--; if(readCount===0) afterPreviews(); };'
        '        reader.readAsDataURL(fileList[idx]);'
        '      } else {'
        '        dataUrls[idx] = null;'
        '        readCount--;'
        '        if(readCount === 0) afterPreviews();'
        '      }'
        '    })(i);'
        '  }'
        '}'
        ''
        'function handleFileSelect(files){'
        '  if(!files || files.length === 0) return;'
        '  uploadFiles(files);'
        '  document.getElementById("sampleFile").value = "";'
        '}'
        ''
        '/* ============ DRAG & DROP ============ */'
        'var dragCounter = 0;'
        'var chatContainer = document.getElementById("chatContainer");'
        'var dropOverlay = document.getElementById("dropOverlay");'
        ''
        '/* Prevent browser default: opening dropped files as new page */'
        'document.addEventListener("dragover", function(e){ e.preventDefault(); });'
        'document.addEventListener("drop", function(e){ e.preventDefault(); });'
        ''
        'chatContainer.addEventListener("dragenter", function(e){'
        '  e.preventDefault(); e.stopPropagation();'
        '  dragCounter++;'
        '  dropOverlay.classList.add("active");'
        '});'
        'chatContainer.addEventListener("dragleave", function(e){'
        '  e.preventDefault(); e.stopPropagation();'
        '  dragCounter--;'
        '  if(dragCounter <= 0){ dragCounter = 0; dropOverlay.classList.remove("active"); }'
        '});'
        'chatContainer.addEventListener("dragover", function(e){'
        '  e.preventDefault(); e.stopPropagation();'
        '});'
        '/* Handle drop on both container and overlay */'
        'function handleFileDrop(e){'
        '  e.preventDefault(); e.stopPropagation();'
        '  dragCounter = 0;'
        '  dropOverlay.classList.remove("active");'
        '  if(e.dataTransfer && e.dataTransfer.files && e.dataTransfer.files.length > 0){'
        '    uploadFiles(e.dataTransfer.files);'
        '  }'
        '}'
        'chatContainer.addEventListener("drop", handleFileDrop);'
        'dropOverlay.addEventListener("drop", handleFileDrop);'
        'dropOverlay.addEventListener("dragover", function(e){ e.preventDefault(); e.stopPropagation(); });'
        'dropOverlay.addEventListener("dragenter", function(e){ e.preventDefault(); e.stopPropagation(); });'
        ''
        '/* ============ PASTE IMAGE ============ */'
        'document.getElementById("chatInput").addEventListener("paste", function(e){'
        '  var items = (e.clipboardData || e.originalEvent.clipboardData || {}).items;'
        '  if(!items) return;'
        '  var imageFiles = [];'
        '  for(var i = 0; i < items.length; i++){'
        '    if(items[i].type.indexOf("image") !== -1){'
        '      var blob = items[i].getAsFile();'
        '      if(blob){'
        '        var ext = items[i].type.split("/")[1] || "png";'
        '        if(ext === "jpeg") ext = "jpg";'
        '        var ts = new Date().toISOString().replace(/[:.]/g,"-").slice(0,19);'
        '        var named = new File([blob], "screenshot_" + ts + "." + ext, {type: blob.type});'
        '        imageFiles.push(named);'
        '      }'
        '    }'
        '  }'
        '  if(imageFiles.length > 0){'
        '    e.preventDefault();'
        '    uploadFiles(imageFiles);'
        '  }'
        '});'
        ''
        '/* Textarea: Enter to send, Shift+Enter for newline, auto-resize */'
        'document.getElementById("chatInput").addEventListener("keydown", function(e){'
        '  if(e.key === "Enter" && !e.shiftKey){'
        '    e.preventDefault();'
        '    sendChatMessage();'
        '  }'
        '});'
        'document.getElementById("chatInput").addEventListener("input", function(){'
        '  this.style.height = "42px";'
        '  this.style.height = Math.min(this.scrollHeight, 120) + "px";'
        '});'
        ''
        '/* ── Tool Improvement Picker ── */'
        'var selectedToolContext = null;'
        ''
        'function showToolPicker(){'
        '  var overlay = document.getElementById("toolPickerOverlay");'
        '  var body = document.getElementById("toolPickerBody");'
        '  overlay.classList.add("active");'
        '  body.innerHTML = \'<div class="tool-picker-loading">&#8987; טוען כלים...</div>\';'
        '  fetch("/tools/create/list-tools", {headers:{"X-Requested-With":"XMLHttpRequest"}})'
        '  .then(function(r){return r.json();})'
        '  .then(function(data){'
        '    if(!data.ok || !data.tools || data.tools.length === 0){'
        '      body.innerHTML = \'<div class="tool-picker-empty">אין כלים זמינים לשיפור</div>\';'
        '      return;'
        '    }'
        '    var html = "";'
        '    var builtins = data.tools.filter(function(t){return t.type==="builtin";});'
        '    var created = data.tools.filter(function(t){return t.type==="marketplace" && t.subtype==="created";});'
        '    var installed = data.tools.filter(function(t){return t.type==="marketplace" && t.subtype==="installed";});'
        '    if(builtins.length > 0){'
        '      html += \'<div class="tool-picker-section"><div class="tool-picker-section-title">כלים מובנים</div>\';'
        '      builtins.forEach(function(t){'
        '        html += renderToolCard(t, "#eff6ff", "#1e3a8a", "מובנה");'
        '      });'
        '      html += "</div>";'
        '    }'
        '    if(created.length > 0){'
        '      html += \'<div class="tool-picker-section"><div class="tool-picker-section-title">כלים שיצרת</div>\';'
        '      created.forEach(function(t){'
        '        var statusLabel = {draft:"טיוטה",pending_review:"ממתין",approved:"מאושר",rejected:"נדחה"}[t.status] || t.status;'
        '        html += renderToolCard(t, "#ecfdf5", "#047857", statusLabel);'
        '      });'
        '      html += "</div>";'
        '    }'
        '    if(installed.length > 0){'
        '      html += \'<div class="tool-picker-section"><div class="tool-picker-section-title">כלים מותקנים</div>\';'
        '      installed.forEach(function(t){'
        '        html += renderToolCard(t, "#fff7ed", "#c2410c", "מותקן");'
        '      });'
        '      html += "</div>";'
        '    }'
        '    body.innerHTML = html;'
        '  })'
        '  .catch(function(){'
        '    body.innerHTML = \'<div class="tool-picker-empty">שגיאה בטעינת הכלים</div>\';'
        '  });'
        '}'
        ''
        'function renderToolCard(t, badgeBg, badgeFg, badgeLabel){'
        '  var icon = t.icon || "&#128295;";'
        '  var desc = t.description || "";'
        '  var safeType = escapeHtml(String(t.type));'
        '  var safeId = escapeHtml(String(t.id));'
        '  var safeName = escapeHtml(String(t.name));'
        '  return \'<div class="tool-picker-card" data-tool-type="\' + safeType + \'" data-tool-id="\' + safeId + \'" data-tool-name="\' + safeName + \'">\''
        '    + \'<div class="tool-picker-card-icon">\' + icon + \'</div>\''
        '    + \'<div class="tool-picker-card-info"><div class="tool-picker-card-name">\' + safeName + \'</div><div class="tool-picker-card-desc">\' + escapeHtml(desc) + \'</div></div>\''
        '    + \'<div class="tool-picker-card-badge" style="background:\' + badgeBg + \';color:\' + badgeFg + \'">\' + badgeLabel + \'</div>\''
        '    + \'</div>\';'
        '}'
        ''
        'function hideToolPicker(){'
        '  document.getElementById("toolPickerOverlay").classList.remove("active");'
        '}'
        ''
        'function selectTool(toolType, toolId, toolName){'
        '  hideToolPicker();'
        '  setLoading(true);'
        '  addUserBubble("אני רוצה לשפר את הכלי: " + toolName);'
        '  showTypingIndicator();'
        '  fetch("/tools/create/load-tool", {'
        '    method: "POST",'
        '    headers: {"Content-Type":"application/json","X-Requested-With":"XMLHttpRequest"},'
        '    body: JSON.stringify({tool_id: toolId, tool_type: toolType})'
        '  })'
        '  .then(function(r){return r.json();})'
        '  .then(function(data){'
        '    if(data.error){'
        '      removeTypingIndicator();'
        '      addAiBubble("שגיאה: " + data.error);'
        '      setLoading(false);'
        '      return;'
        '    }'
        '    selectedToolContext = data.tool;'
        '    var userMsg = "אני רוצה לשפר את הכלי: " + toolName;'
        '    fetch("/tools/create/chat", {'
        '      method: "POST",'
        '      headers: {"Content-Type":"application/json","X-Requested-With":"XMLHttpRequest"},'
        '      body: JSON.stringify({session_id: CHAT_SESSION_ID, message: userMsg, tool_context: data.tool})'
        '    })'
        '    .then(function(r){return r.json();})'
        '    .then(function(chatData){'
        '      removeTypingIndicator();'
        '      if(chatData.error){ addAiBubble("שגיאה: " + chatData.error); }'
        '      else {'
        '        addAiBubble(chatData.assistant_message);'
        '        if(chatData.has_tool && chatData.tool_definition) addToolReady(chatData.tool_definition);'
        '      }'
        '      setLoading(false);'
        '    })'
        '    .catch(function(){'
        '      removeTypingIndicator();'
        '      addAiBubble("שגיאה בתקשורת. נסה שוב.");'
        '      setLoading(false);'
        '    });'
        '  })'
        '  .catch(function(){'
        '    removeTypingIndicator();'
        '    addAiBubble("שגיאה בטעינת פרטי הכלי. נסה שוב.");'
        '    setLoading(false);'
        '  });'
        '}'
        ''
        '/* Close tool picker on overlay click */'
        'document.getElementById("toolPickerOverlay").addEventListener("click", function(e){'
        '  if(e.target === this) hideToolPicker();'
        '});'
        ''
        '/* Delegate clicks on tool picker cards */'
        'document.getElementById("toolPickerBody").addEventListener("click", function(e){'
        '  var card = e.target.closest(".tool-picker-card");'
        '  if(!card) return;'
        '  var tType = card.getAttribute("data-tool-type");'
        '  var tId = card.getAttribute("data-tool-id");'
        '  var tName = card.getAttribute("data-tool-name");'
        '  if(tType && tId) selectTool(tType, tId, tName || "");'
        '});'
        ''
        '/* Auto-scroll on page load */'
        'scrollToBottom();'
        '</script>'
    )

    return render("העוזר החכם של Scriptly", body, lang=lang, topbar_greeting=text["topbar_greeting"], logout_label=text["logout"], show_lang_switch=True)


@app.route("/tools/create/list-tools", methods=["GET"])
@login_required
def tools_create_list_tools():
    """Return JSON list of tools available to the current user for improvement."""
    user_id = session["user_id"]
    tools_list = []

    with get_db() as db:
        # 1) Built-in scripts assigned to the user
        perms = db.execute("SELECT script_id FROM permissions WHERE user_id=?", (user_id,)).fetchall()
        for p in perms:
            sid = p["script_id"]
            if sid in SCRIPT_REGISTRY:
                s = SCRIPT_REGISTRY[sid]
                tools_list.append({
                    "id": sid,
                    "name": s.get("name", sid),
                    "description": (s.get("desc") or "")[:120],
                    "icon": s.get("icon", ""),
                    "type": "builtin",
                })

        # 2) Marketplace tools the user created (any status)
        created = db.execute(
            "SELECT id, name, description, icon, status FROM marketplace_tools WHERE creator_id=? ORDER BY created_at DESC",
            (user_id,),
        ).fetchall()
        for t in created:
            tools_list.append({
                "id": t["id"],
                "name": t["name"],
                "description": (t["description"] or "")[:120],
                "icon": t["icon"] or "",
                "type": "marketplace",
                "subtype": "created",
                "status": t["status"],
            })

        # 3) Marketplace tools the user installed (approved only)
        installed = db.execute(
            """SELECT t.id, t.name, t.description, t.icon, t.status
            FROM marketplace_tools t JOIN tool_installs i ON i.tool_id = t.id
            WHERE i.user_id=? AND t.status='approved' AND t.is_public=1
              AND t.creator_id != ?
            ORDER BY t.name""",
            (user_id, user_id),
        ).fetchall()
        for t in installed:
            tools_list.append({
                "id": t["id"],
                "name": t["name"],
                "description": (t["description"] or "")[:120],
                "icon": t["icon"] or "",
                "type": "marketplace",
                "subtype": "installed",
            })

    return jsonify({"ok": True, "tools": tools_list})


@app.route("/tools/create/load-tool", methods=["POST"])
@login_required
def tools_create_load_tool():
    """Load the full definition of a tool for the AI improvement context."""
    data = request.get_json(silent=True) or {}
    tool_id = data.get("tool_id", "")
    tool_type = data.get("tool_type", "")

    if not tool_id or not tool_type:
        return jsonify({"error": "חסר מזהה כלי"}), 400

    if tool_type == "builtin":
        sid = str(tool_id)
        if sid not in SCRIPT_REGISTRY:
            return jsonify({"error": "כלי לא נמצא"}), 404
        s = SCRIPT_REGISTRY[sid]
        # Build a summary of the builtin script metadata (no processor function)
        tool_info = {
            "type": "builtin",
            "id": sid,
            "name": s.get("name", sid),
            "description": s.get("desc", ""),
            "icon": s.get("icon", ""),
            "accept": s.get("accept", ""),
            "requires_mapping_confirmation": s.get("requires_mapping_confirmation", False),
            "background_queue": s.get("background_queue", False),
            "help_intro": s.get("help_intro", ""),
            "help_items": s.get("help_items", []),
            "rules_intro": s.get("rules_intro", ""),
            "rules_items": s.get("rules_items", []),
            "filter_fields": s.get("filter_fields", []),
        }
        return jsonify({"ok": True, "tool": tool_info})

    elif tool_type == "marketplace":
        tool_id_int = int(tool_id)
        with get_db() as db:
            t = db.execute("SELECT * FROM marketplace_tools WHERE id=?", (tool_id_int,)).fetchone()
            if not t:
                return jsonify({"error": "כלי לא נמצא"}), 404
            # Verify access: user is creator or has installed
            user_id = session["user_id"]
            is_creator = t["creator_id"] == user_id
            is_installed = db.execute(
                "SELECT 1 FROM tool_installs WHERE tool_id=? AND user_id=?",
                (tool_id_int, user_id),
            ).fetchone()
            if not is_creator and not is_installed:
                return jsonify({"error": "אין הרשאה לכלי זה"}), 403

            try:
                definition = json.loads(t["definition_json"])
            except (json.JSONDecodeError, TypeError):
                definition = {}

            tool_info = {
                "type": "marketplace",
                "id": tool_id_int,
                "name": t["name"],
                "description": t["description"] or "",
                "icon": t["icon"] or "",
                "category": t["category"] or "general",
                "status": t["status"],
                "is_creator": is_creator,
                "definition": definition,
            }
            return jsonify({"ok": True, "tool": tool_info})

    return jsonify({"error": "סוג כלי לא חוקי"}), 400


@app.route("/tools/create/chat", methods=["POST"])
@login_required
def tools_create_chat():
    is_ajax = request.headers.get("X-Requested-With") == "XMLHttpRequest" or request.is_json
    if is_ajax:
        data = request.get_json(silent=True) or {}
        chat_session_id = data.get("session_id", "")
        user_message = data.get("message", "").strip()
    else:
        chat_session_id = request.form.get("session_id", "")
        user_message = request.form.get("message", "").strip()

    # Optional tool context for improvement flow
    tool_context = None
    file_context = None
    if is_ajax:
        tool_context = (data or {}).get("tool_context")
        file_context = (data or {}).get("file_context")  # analysis from uploaded files

    if not user_message:
        if is_ajax:
            return jsonify({"error": "הודעה ריקה"}), 400
        return redirect("/tools/create")

    # Get chat session
    with get_db() as db:
        chat_data = db.execute(
            "SELECT * FROM tool_chat_sessions WHERE id=? AND user_id=?",
            (chat_session_id, session["user_id"]),
        ).fetchone()
        if not chat_data:
            if is_ajax:
                return jsonify({"error": "שיחה לא נמצאה"}), 404
            add_flash("שיחה לא נמצאה")
            return redirect("/tools/create")

        try:
            messages = json.loads(chat_data["messages_json"])
        except (json.JSONDecodeError, TypeError):
            messages = []

    # If tool context is provided (improvement flow), inject it before the user message
    if tool_context and isinstance(tool_context, dict):
        context_text = "[TOOL_CONTEXT] המשתמש בחר לשפר כלי קיים.\n"
        context_text += "סוג: " + tool_context.get("type", "unknown") + "\n"
        context_text += "שם: " + tool_context.get("name", "") + "\n"
        if tool_context.get("description"):
            context_text += "תיאור: " + tool_context["description"] + "\n"
        if tool_context.get("type") == "builtin":
            context_text += "הערה: זהו כלי מובנה במערכת. לא ניתן לערוך אותו ישירות. ניתן להציע יצירת כלי חדש בשוק הכלים שמשפר את הפונקציונליות.\n"
            if tool_context.get("help_intro"):
                context_text += "עזרה: " + tool_context["help_intro"] + "\n"
            if tool_context.get("rules_intro"):
                context_text += "כללים: " + tool_context["rules_intro"] + "\n"
            if tool_context.get("rules_items"):
                context_text += "פרטי כללים:\n" + "\n".join("- " + r for r in tool_context["rules_items"]) + "\n"
        elif tool_context.get("type") == "marketplace" and tool_context.get("definition"):
            context_text += "הגדרת הכלי הנוכחית:\n```json\n" + json.dumps(tool_context["definition"], ensure_ascii=False, indent=2) + "\n```\n"
        messages.append({"role": "user", "content": context_text})
        messages.append({"role": "assistant", "content": "קיבלתי את פרטי הכלי. אני מבין את המבנה הנוכחי שלו. מה תרצה לשפר או לשנות?"})

    # If file analysis is attached, prepend it to the user message
    if file_context and isinstance(file_context, str) and file_context.strip():
        user_message = f"[קבצים שהעליתי]\n{file_context}\n\n[ההוראות שלי]\n{user_message}"

    # Add user message
    messages.append({"role": "user", "content": user_message})

    # Build API messages with sliding window to control cost/latency.
    # Always preserve: first 4 msgs (tool context), file uploads, and last 12 msgs.
    MAX_RECENT = 12
    KEEP_FIRST = 4
    CONTEXT_MARKERS = ("[קבצים שהעליתי]", "[TOOL_CONTEXT]", "File structure:")
    all_msgs = [{"role": m["role"], "content": m["content"]} for m in messages]
    if len(all_msgs) > KEEP_FIRST + MAX_RECENT:
        head = all_msgs[:KEEP_FIRST]
        tail = all_msgs[-MAX_RECENT:]
        tail_set = set(id(m) for m in tail)
        # Scan trimmed middle for file/context messages — always keep them
        middle = all_msgs[KEEP_FIRST:-MAX_RECENT]
        pinned = []
        for m in middle:
            if id(m) not in tail_set and any(marker in m.get("content", "") for marker in CONTEXT_MARKERS):
                pinned.append(m)
        head.extend(pinned)
        head.append({"role": "user", "content": "[הערת מערכת: חלק מההודעות הישנות קוצרו. הקבצים והנתונים שהועלו נשמרו. המשך על בסיס ההודעות האחרונות.]"})
        head.append({"role": "assistant", "content": "הבנתי, ממשיך."})
        api_messages = head + tail
    else:
        api_messages = all_msgs
    # Build dynamic system prompt with user's tools
    user_row = None
    with get_db() as db:
        user_row = db.execute("SELECT * FROM users WHERE id=?", (session["user_id"],)).fetchone()
    customer_name = (user_row["full_name"] if user_row else "") or ""
    system_prompt = build_tool_creation_system_prompt(session["user_id"], customer_name)

    assistant_response, _in_tok, _out_tok = call_claude_chat(
        api_messages, system_prompt=system_prompt,
        user_id=session["user_id"], session_id=int(chat_session_id) if chat_session_id else None,
    )
    messages.append({"role": "assistant", "content": assistant_response})

    # Check for developer brief in response
    brief_text = extract_developer_brief(assistant_response)
    # Check if tool definition exists in response (for Mode C simple changes)
    tool_def = extract_tool_json_from_message(assistant_response)

    # Save updated messages
    now = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    with get_db() as db:
        db.execute(
            "UPDATE tool_chat_sessions SET messages_json=?, updated_at=? WHERE id=?",
            (json.dumps(messages, ensure_ascii=False), now, chat_session_id),
        )
        db.commit()

    log_user_activity("ai_chat", "שיחת AI ליצירת כלי", "", "", f"session={chat_session_id}")

    if is_ajax:
        return jsonify({
            "ok": True,
            "assistant_message": assistant_response,
            "has_tool": tool_def is not None,
            "tool_definition": json.dumps(tool_def, ensure_ascii=False) if tool_def else None,
            "has_brief": brief_text is not None,
            "brief_text": brief_text,
            "brief_tool_name": extract_brief_tool_name(brief_text) if brief_text else None,
        })
    return redirect("/tools/create")


@app.route("/tools/create/upload", methods=["POST"])
@login_required
def tools_create_upload():
    """Upload and analyze files. Returns analysis only — does NOT call Claude.
    The frontend stores the analysis and sends it with the user's next message."""
    is_ajax = request.headers.get("X-Requested-With") == "XMLHttpRequest"
    chat_session_id = request.form.get("session_id", "")

    # Support multiple files via getlist
    file_list = request.files.getlist("sample_file")
    if not file_list or all(not f.filename for f in file_list):
        if is_ajax:
            return jsonify({"error": "לא נבחר קובץ"}), 400
        add_flash("לא נבחר קובץ")
        return redirect("/tools/create")

    temp_paths = []
    analyses = []
    file_names = []
    file_types = []
    errors = []

    for file_obj in file_list:
        if not file_obj or not file_obj.filename:
            continue
        validation_error, ext = validate_upload_for_chat(file_obj)
        if validation_error:
            errors.append(f"{file_obj.filename}: {validation_error}")
            continue

        # Save temp file for analysis
        uid = str(uuid.uuid4())[:8]
        safe_name = file_obj.filename.replace(" ", "_")
        temp_path = str(UPLOAD_FOLDER / f"{uid}_{safe_name}")
        file_obj.save(temp_path)
        temp_paths.append(temp_path)
        file_names.append(file_obj.filename)

        if ext in IMAGE_EXTENSIONS:
            file_types.append("image")
            analysis = analyze_image_for_chat(temp_path, file_obj.filename)
        elif ext in DOC_EXTENSIONS:
            file_types.append("document")
            analysis = analyze_document_for_chat(temp_path, file_obj.filename, ext)
        else:
            file_types.append("data")
            analysis = analyze_sample_file_for_chat(temp_path, ext)
        analyses.append(f"--- {file_obj.filename} ---\n{analysis}")

    # Cleanup temp files
    for tp in temp_paths:
        try:
            os.remove(tp)
        except OSError:
            pass

    if not analyses and errors:
        err_msg = "שגיאה בקבצים: " + "; ".join(errors)
        if is_ajax:
            return jsonify({"error": err_msg}), 400
        add_flash(err_msg)
        return redirect("/tools/create")

    # Build combined analysis text
    if len(analyses) == 1:
        combined = analyses[0].split("---\n", 1)[-1]
    else:
        combined = "\n\n".join(analyses)
    if errors:
        combined += "\n\n(שגיאות: " + "; ".join(errors) + ")"

    log_user_activity("ai_chat_upload", "העלאת קובץ לשיחת AI", "", "", f"session={chat_session_id}, files={','.join(file_names)}")

    if is_ajax:
        return jsonify({
            "ok": True,
            "analysis": combined,
            "file_names": file_names,
            "file_types": file_types,
            "errors": errors,
        })
    return redirect("/tools/create")


@app.route("/tools/create/reset", methods=["POST"])
@login_required
def tools_create_reset():
    # Close existing sessions and redirect to create a fresh one
    with get_db() as db:
        db.execute(
            "UPDATE tool_chat_sessions SET status='closed', messages_json='[]' WHERE user_id=? AND status='active'",
            (session["user_id"],),
        )
        db.commit()
    return redirect("/tools/create")


@app.route("/tools/create/save", methods=["POST"])
@login_required
def tools_create_save():
    definition_text = request.form.get("definition", "")
    action = request.form.get("action", "save_draft")

    try:
        definition = json.loads(definition_text)
    except (json.JSONDecodeError, TypeError, ValueError):
        add_flash("הגדרת הכלי שגויה")
        return redirect("/tools/create")

    # For publish: validate strictly. For draft: validate lightly (allow testing)
    if action != "save_draft":
        ok, errors = validate_tool_definition(definition)
        if not ok:
            add_flash("שגיאות בהגדרת הכלי: " + "; ".join(errors))
            return redirect("/tools/create")

    now = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    status = "draft" if action == "save_draft" else "pending_review"
    name = definition.get("name", "כלי ללא שם")
    description = definition.get("description", "")
    icon = definition.get("icon", "🔧")
    category = definition.get("category", "general")
    if category not in TOOL_CATEGORIES:
        category = "general"

    with get_db() as db:
        db.execute(
            """INSERT INTO marketplace_tools(creator_id, name, description, icon, category,
               definition_json, status, is_public, created_at, updated_at)
               VALUES (?,?,?,?,?,?,?,0,?,?)""",
            (session["user_id"], name, description, icon, category,
             json.dumps(definition, ensure_ascii=False), status, now, now),
        )
        db.commit()
        new_tool_id = db.execute("SELECT id FROM marketplace_tools WHERE creator_id=? ORDER BY id DESC LIMIT 1", (session["user_id"],)).fetchone()["id"]

        # Close the chat session and clear messages (data privacy)
        db.execute(
            "UPDATE tool_chat_sessions SET status='completed', messages_json='[]' WHERE user_id=? AND status='active'",
            (session["user_id"],),
        )
        db.commit()

    if status == "pending_review":
        add_flash("הכלי נשלח לאישור! הוא יופיע בשוק לאחר בדיקה.")
        log_user_activity("create_tool", "יצירת כלי AI", "", name, f"status={status}")
        return redirect("/marketplace")
    else:
        add_flash("הכלי נשמר! תוכל לבדוק אותו כאן לפני פרסום.")
        log_user_activity("create_tool", "יצירת כלי AI", "", name, f"status={status}")
        return redirect(f"/tools/my-tools")


@app.route("/tools/create/save-brief", methods=["POST"])
@login_required
def tools_create_save_brief():
    """Save a developer brief (technical spec) from the AI chat."""
    is_ajax = request.headers.get("X-Requested-With") == "XMLHttpRequest" or request.is_json
    data = request.get_json(silent=True) or {}
    brief_text = (data.get("brief_text") or "").strip()
    tool_name = (data.get("tool_name") or "כלי ללא שם").strip()
    request_type = data.get("request_type", "new")
    existing_tool_name = (data.get("existing_tool_name") or "").strip()
    chat_session_id = data.get("session_id")

    if not brief_text:
        if is_ajax:
            return jsonify({"error": "סיכום טכני ריק"}), 400
        add_flash("סיכום טכני ריק")
        return redirect("/tools/create")

    now = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    with get_db() as db:
        user_row = db.execute("SELECT * FROM users WHERE id=?", (session["user_id"],)).fetchone()
        db.execute(
            """INSERT INTO tool_requests(user_id, username, full_name, company_name, tool_name,
               request_type, existing_tool_name, brief_text, chat_session_id, status, created_at, updated_at)
               VALUES (?,?,?,?,?,?,?,?,?,?,?,?)""",
            (
                session["user_id"],
                (user_row["username"] if user_row else "") or "",
                (user_row["full_name"] if user_row else "") or "",
                (user_row["company_name"] if user_row else "") or "",
                tool_name,
                request_type if request_type in ("new", "improvement") else "new",
                existing_tool_name,
                brief_text,
                chat_session_id,
                "pending",
                now,
                now,
            ),
        )
        db.commit()

        # Close chat session and clear messages (data privacy)
        if chat_session_id:
            db.execute(
                "UPDATE tool_chat_sessions SET status='completed', messages_json='[]' WHERE id=?",
                (chat_session_id,),
            )
            db.commit()

    log_user_activity("tool_request", "בקשת כלי חדש", "", tool_name, f"type={request_type}")

    if is_ajax:
        return jsonify({"ok": True, "message": "הבקשה נשלחה בהצלחה! צוות הפיתוח יבנה את הכלי ויעדכן אותך."})
    add_flash("הבקשה נשלחה בהצלחה! צוות הפיתוח יבנה את הכלי ויעדכן אותך.")
    return redirect("/tools/create")


# ── My Tools: Private tool area ───────────────────────────────────

@app.route("/tools/my-tools")
@login_required
def my_tools():
    if session.get("is_admin"):
        return redirect("/admin")

    lang = get_flow_lang()
    text = get_flow_text(lang)
    user_id = session["user_id"]

    with get_db() as db:
        # User's created tools
        created = db.execute(
            "SELECT * FROM marketplace_tools WHERE creator_id=? ORDER BY updated_at DESC",
            (user_id,),
        ).fetchall()
        # User's saved/installed tools (not created by user)
        saved = db.execute(
            """SELECT t.* FROM marketplace_tools t
               JOIN tool_installs i ON i.tool_id = t.id
               WHERE i.user_id=? AND t.creator_id != ?
               ORDER BY i.created_at DESC""",
            (user_id, user_id),
        ).fetchall()

    status_labels = {"draft": "טיוטה", "pending_review": "ממתין לאישור", "approved": "מאושר", "rejected": "נדחה"}
    status_colors = {"draft": "#f59e0b", "pending_review": "#3b82f6", "approved": "#047857", "rejected": "#ef4444"}

    # Build created tools cards
    created_html = ""
    for t in created:
        st = t["status"]
        st_label = status_labels.get(st, st)
        st_color = status_colors.get(st, "#64748b")
        icon = t["icon"] or "&#128295;"
        name = esc(t["name"])
        desc = esc(t["description"] or "")[:80]
        tid = t["id"]

        buttons = ""
        # Test button — always available
        buttons += f'<a href="/marketplace/tool/{tid}/run" class="btn btn-blue" style="font-size:12px;padding:6px 12px;border-radius:8px">&#9654; בדיקה</a> '
        # Improve button — open chat with tool context
        buttons += f'<a href="/tools/create" onclick="sessionStorage.setItem(\'improveToolId\',\'{tid}\')" class="btn btn-gray" style="font-size:12px;padding:6px 12px;border-radius:8px">&#9998; שיפור</a> '
        # Publish button — only for drafts
        if st == "draft":
            buttons += (
                f'<form method="post" action="/tools/my-tools/{tid}/publish" style="display:inline">'
                f'<button type="submit" class="btn btn-blue" style="font-size:12px;padding:6px 12px;border-radius:8px;background:#047857">&#128228; שלח לפרסום</button>'
                f'</form> '
            )
        # Delete button
        buttons += (
            f'<form method="post" action="/tools/my-tools/{tid}/delete" style="display:inline"'
            f' onsubmit="return confirm(\'למחוק את הכלי?\')">'
            f'<button type="submit" class="btn btn-gray" style="font-size:12px;padding:6px 12px;border-radius:8px;color:#ef4444">&#128465; מחיקה</button>'
            f'</form>'
        )

        created_html += (
            f'<div style="border:1.5px solid #e2e8f0;border-radius:16px;padding:18px;margin-bottom:12px;background:white;'
            f'transition:all .3s cubic-bezier(.4,0,.2,1);display:flex;align-items:flex-start;gap:14px" '
            f'onmouseover="this.style.boxShadow=\'0 8px 24px rgba(37,99,235,.08)\';this.style.borderColor=\'#bfdbfe\'" '
            f'onmouseout="this.style.boxShadow=\'none\';this.style.borderColor=\'#e2e8f0\'">'
            f'<div style="width:48px;height:48px;border-radius:14px;background:linear-gradient(135deg,#eff6ff,#dbeafe);display:flex;align-items:center;justify-content:center;font-size:24px;flex-shrink:0">{icon}</div>'
            f'<div style="flex:1;min-width:0">'
            f'<div style="display:flex;align-items:center;gap:8px;margin-bottom:4px;flex-wrap:wrap">'
            f'<span style="font-weight:800;font-size:15px;color:#0f172a">{name}</span>'
            f'<span style="font-size:11px;padding:2px 8px;border-radius:99px;background:{st_color}15;color:{st_color};font-weight:700">{st_label}</span>'
            f'</div>'
            f'<div style="font-size:13px;color:#64748b;margin-bottom:8px">{desc}</div>'
            f'<div style="display:flex;gap:6px;flex-wrap:wrap">{buttons}</div>'
            f'</div>'
            f'</div>'
        )

    if not created:
        created_html = '<div style="text-align:center;padding:2rem;color:#94a3b8">עדיין לא יצרת כלים. <a href="/tools/create" style="color:#2563eb">צור כלי חדש</a></div>'

    # Build saved tools cards
    saved_html = ""
    for t in saved:
        icon = t["icon"] or "&#128295;"
        name = esc(t["name"])
        desc = esc(t["description"] or "")[:80]
        tid = t["id"]
        saved_html += (
            f'<div style="border:1.5px solid #e2e8f0;border-radius:16px;padding:18px;margin-bottom:12px;background:white;'
            f'display:flex;align-items:center;gap:14px;transition:all .3s cubic-bezier(.4,0,.2,1)" '
            f'onmouseover="this.style.borderColor=\'#bfdbfe\';this.style.boxShadow=\'0 4px 16px rgba(37,99,235,.06)\'" '
            f'onmouseout="this.style.borderColor=\'#e2e8f0\';this.style.boxShadow=\'none\'">'
            f'<div style="width:44px;height:44px;border-radius:12px;background:linear-gradient(135deg,#ecfdf5,#d1fae5);display:flex;align-items:center;justify-content:center;font-size:22px;flex-shrink:0">{icon}</div>'
            f'<div style="flex:1;min-width:0">'
            f'<div style="font-weight:800;font-size:14px;color:#0f172a;margin-bottom:2px">{name}</div>'
            f'<div style="font-size:12px;color:#64748b">{desc}</div>'
            f'</div>'
            f'<a href="/marketplace/tool/{tid}/run" class="btn btn-blue" style="font-size:12px;padding:6px 12px;border-radius:8px">&#9654; הפעל</a>'
            f'</div>'
        )

    if not saved:
        saved_html = '<div style="text-align:center;padding:2rem;color:#94a3b8">אין כלים שמורים. <a href="/marketplace" style="color:#2563eb">גלה כלים בשוק</a></div>'

    body = (
        '<div style="max-width:800px;margin:0 auto">'
        '<div style="display:flex;align-items:center;justify-content:space-between;flex-wrap:wrap;gap:10px;margin-bottom:1.5rem">'
        '<h2 style="margin:0;font-size:clamp(20px,3vw,24px);font-weight:800;color:#0f172a;letter-spacing:-0.3px">&#128188; הכלים שלי</h2>'
        '<div style="display:flex;gap:10px">'
        '<a href="/tools/create" class="btn btn-blue" style="text-decoration:none">&#10133; כלי חדש</a>'
        '<a href="/marketplace" class="btn btn-gray" style="text-decoration:none">&#128722; שוק הכלים</a>'
        '</div></div>'
        # Created tools section
        '<div style="margin-bottom:2rem">'
        '<h3 style="font-size:16px;font-weight:800;color:#1e3a8a;margin-bottom:14px;padding-bottom:10px;border-bottom:2px solid #dbeafe">&#9997; כלים שיצרתי</h3>'
        + created_html
        + '</div>'
        # Saved/favorite tools section
        '<div>'
        '<h3 style="font-size:16px;font-weight:800;color:#1e3a8a;margin-bottom:14px;padding-bottom:10px;border-bottom:2px solid #dbeafe">&#11088; כלים שמורים</h3>'
        + saved_html
        + '</div>'
        '</div>'
    )

    return render("הכלים שלי", body, lang=lang, topbar_greeting=text["topbar_greeting"], logout_label=text["logout"], show_lang_switch=True)


@app.route("/tools/my-tools/<int:tool_id>/publish", methods=["POST"])
@login_required
def my_tools_publish(tool_id):
    with get_db() as db:
        tool = db.execute("SELECT * FROM marketplace_tools WHERE id=? AND creator_id=?", (tool_id, session["user_id"])).fetchone()
        if not tool:
            add_flash("הכלי לא נמצא")
            return redirect("/tools/my-tools")
        # Validate before publishing
        try:
            definition = json.loads(tool["definition_json"])
        except (json.JSONDecodeError, TypeError):
            add_flash("הגדרת הכלי שגויה")
            return redirect("/tools/my-tools")

        ok, errors = validate_tool_definition(definition)
        if not ok:
            add_flash("שגיאות בהגדרת הכלי: " + "; ".join(errors))
            return redirect("/tools/my-tools")

        now = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        db.execute("UPDATE marketplace_tools SET status='pending_review', updated_at=? WHERE id=?", (now, tool_id))
        db.commit()
    add_flash("הכלי נשלח לאישור!")
    return redirect("/tools/my-tools")


@app.route("/tools/my-tools/<int:tool_id>/delete", methods=["POST"])
@login_required
def my_tools_delete(tool_id):
    with get_db() as db:
        db.execute("DELETE FROM marketplace_tools WHERE id=? AND creator_id=?", (tool_id, session["user_id"]))
        db.commit()
    add_flash("הכלי נמחק.")
    return redirect("/tools/my-tools")


# ── Admin: Tool Moderation ────────────────────────────────────────

@app.route("/admin/tools")
@login_required
@admin_required
def admin_tools():
    lang = get_flow_lang()
    with get_db() as db:
        tools = db.execute(
            """SELECT t.*, u.full_name as creator_name, u.username as creator_username
            FROM marketplace_tools t JOIN users u ON u.id = t.creator_id
            ORDER BY
                CASE WHEN t.status='pending_review' THEN 0 ELSE 1 END,
                t.created_at DESC""",
        ).fetchall()

    status_labels = {"draft": "טיוטה", "pending_review": "ממתין לאישור", "approved": "מאושר", "rejected": "נדחה"}
    status_colors = {"draft": ("#f8fafc", "#475569"), "pending_review": ("#fff7ed", "#c2410c"), "approved": ("#ecfdf5", "#047857"), "rejected": ("#fef2f2", "#b91c1c")}

    rows = ""
    for t in tools:
        s_bg, s_fg = status_colors.get(t["status"], ("#f8fafc", "#475569"))
        actions = ""
        if t["status"] == "pending_review":
            actions = (
                '<form method="post" action="/admin/tools/' + str(t["id"]) + '/approve" style="display:inline"><button type="submit" class="btn btn-blue" style="font-size:12px;padding:6px 12px;border-radius:8px">✅ אשר</button></form> '
                '<form method="post" action="/admin/tools/' + str(t["id"]) + '/reject" style="display:inline"><button type="submit" class="btn btn-red" style="font-size:12px;padding:6px 12px;border-radius:8px">❌ דחה</button></form>'
            )
        elif t["status"] == "approved":
            actions = '<form method="post" action="/admin/tools/' + str(t["id"]) + '/reject" style="display:inline"><button type="submit" class="btn btn-red" style="font-size:12px;padding:6px 12px;border-radius:8px">הסר</button></form>'

        rows += (
            '<tr>'
            '<td>' + esc(t["icon"] or "🔧") + ' ' + esc(t["name"]) + '</td>'
            '<td>' + esc(t["creator_name"] or t["creator_username"]) + '</td>'
            '<td><span style="padding:3px 10px;border-radius:99px;font-size:11px;font-weight:700;background:' + s_bg + ';color:' + s_fg + '">' + esc(status_labels.get(t["status"], t["status"])) + '</span></td>'
            '<td>' + str(t["usage_count"] or 0) + '</td>'
            '<td>' + esc(format_ui_datetime(t["created_at"])) + '</td>'
            '<td>' + actions + '</td>'
            '</tr>'
        )

    body = (
        '<div style="margin-bottom:1rem"><a href="/admin" style="color:#2563eb;font-size:13px;text-decoration:none">← חזרה לניהול</a></div>'
        '<div class="card">'
        '<h2>ניהול כלים — שוק הכלים</h2>'
        '<table><thead><tr>'
        '<th>כלי</th><th>יוצר</th><th>סטטוס</th><th>שימושים</th><th>תאריך</th><th>פעולות</th>'
        '</tr></thead><tbody>'
        + (rows or '<tr><td colspan="6" style="text-align:center;color:#94a3b8;padding:2rem">אין כלים</td></tr>')
        + '</tbody></table></div>'
    )

    return render("ניהול כלים", body, lang=lang)


@app.route("/admin/tools/<int:tool_id>/approve", methods=["POST"])
@login_required
@admin_required
def admin_approve_tool(tool_id):
    now = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    with get_db() as db:
        db.execute(
            "UPDATE marketplace_tools SET status='approved', is_public=1, approved_at=?, approved_by=? WHERE id=?",
            (now, session["user_id"], tool_id),
        )
        db.commit()
    add_flash("הכלי אושר ופורסם בשוק")
    return redirect("/admin/tools")


@app.route("/admin/tools/<int:tool_id>/reject", methods=["POST"])
@login_required
@admin_required
def admin_reject_tool(tool_id):
    with get_db() as db:
        db.execute(
            "UPDATE marketplace_tools SET status='rejected', is_public=0 WHERE id=?",
            (tool_id,),
        )
        db.commit()
    add_flash("הכלי נדחה")
    return redirect("/admin/tools")


if __name__ == "__main__":
    app.run(debug=False, use_reloader=False)

